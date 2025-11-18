'''
SPDX-License-Identifier: Apache-2.0

Copyright 2023 Cisco Systems, Inc. and its affiliates

Licensed under the Apache License, Version 2.0 (the "License");
you may not use this file except in compliance with the License.
You may obtain a copy of the License at

http://www.apache.org/licenses/LICENSE-2.0

Unless required by applicable law or agreed to in writing, software
distributed under the License is distributed on an "AS IS" BASIS,
WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
See the License for the specific language governing permissions and
limitations under the License.
'''

import ns_def
import yaml
import random
import math
import re
import shlex
import csv
import sys
from pathlib import Path
from collections import defaultdict
from typing import Optional, List, Any, Set, Dict, Tuple
from ciscoconfparse import CiscoConfParse
import networkx as nx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


class ns_option_convert_to_master_csv():
    def rename_port(self):
        print('--- Rename port ---')
        self.csv_l1_rename_cmd_list = []
        self.script_path = "dummy"  # Path to the network_sketcher.py script

        generated_commands = set()  # Set to prevent duplicate commands
        used_random_numbers = set()  # Set to track used random numbers

        # --- Helper Functions (normally separate class methods) ---
        def _get_unique_random_number() -> int:
            """
            Generates a unique random number between 9000-9999 that hasn't been used before.
            Raises an exception if all numbers have been used.
            """
            available_numbers = set(range(9000, 10000)) - used_random_numbers

            if not available_numbers:
                raise ValueError("All random numbers (9000-9999) have been exhausted.")

            random_number = random.choice(list(available_numbers))
            used_random_numbers.add(random_number)
            return random_number

        def _format_interface_with_space(name: str) -> str:
            """
            Ensures that interface names have a space before the final digit sequence.
            Treats '/', '.', and digits as part of the number sequence.
            Also removes all spaces except the one before the final digit sequence.
            Returns the modified name if it doesn't end with a digit (adds random number).
            Handles cases like:
            """
            if not name:
                return None

            # First, remove all existing spaces
            name_no_spaces = name.replace(' ', '')

            # Check if it ends with a digit
            if not name_no_spaces or not name_no_spaces[-1].isdigit():
                # If name doesn't end with digit, add space and unique random number (9000-9999)
                random_number = _get_unique_random_number()
                return f"{name_no_spaces} {random_number}"

            # Then, add a space before the final digit sequence
            # Pattern to match: any character that is not a digit, '/', '.', or space,
            # followed by a digit (which may be followed by more digits, '/', or '.')
            formatted = re.sub(r'([^\d/.\s])(\d[/.\d]*)$', r'\1 \2', name_no_spaces)
            return formatted

        def _process_interface_name(name: str) -> str:
            """
            Processes the interface name:
            1. Trims whitespace
            2. If it doesn't end with a digit, adds a space and unique random number (9000-9999)
            3. Ensures proper formatting with space before the final number
            """
            cleaned = name.strip()

            # Check if the interface name ends with a digit
            if not cleaned:
                return None

            if not cleaned[-1].isdigit():
                # Remove all spaces and add unique random number
                cleaned_no_spaces = cleaned.replace(' ', '')
                random_number = _get_unique_random_number()
                cleaned = f"{cleaned_no_spaces} {random_number}"

            # Format with proper spacing
            return _format_interface_with_space(cleaned)

        def _is_valid_peer_interface(interface_name: str) -> bool:
            """
            Checks if a string appears to be a valid network interface name.
            Now accepts all interface names (including those ending with non-digits after processing).
            """
            if not interface_name:
                return False
            # Accepts interface names containing numbers, slashes, dots, or "Port" keyword
            # Also accepts names that have been processed with random numbers
            return bool(re.search(r'[\d/.]|Port', interface_name, re.IGNORECASE))

        def _parse_connection_field(connection_str: str) -> tuple:
            """
            Parses the 'Connection' column to extract peer device and interface.

            Handles two formats:
            1. Single line: Just the peer device name (e.g., "Level3 MPLS")
               Returns: (peer_device, None)

            2. Multi-line format:
               - Line 1: Peer device name (e.g., 'dmi01-yonkers-sw01', 'Main Switch')
               - Line 2-N: May contain whitespace or empty lines
               - Last line: Peer interface name (e.g., 'GigabitEthernet1/0/2', 'mgmt', '1/1/49')
               Returns: (peer_device, peer_interface)

            Examples:
            Input: "Level3 MPLS"
            Output: ('Level3 MPLS', None)

            Input: "dmi01-yonkers-sw01\n    \n  \n  GigabitEthernet1/0/2"
            Output: ('dmi01-yonkers-sw01', 'GigabitEthernet1/0/2')

            Input: "Main Switch\n    \n  \n  mgmt"
            Output: ('Main Switch', 'mgmt')
            """
            if not connection_str:
                return None, None

            # Split by newline and filter out empty/whitespace-only lines
            lines = [line.strip() for line in connection_str.split('\n') if line.strip()]

            if len(lines) == 0:
                return None, None
            elif len(lines) == 1:
                # Single line - only device name, no interface
                # This is the case for connections like "Level3 MPLS"
                return lines[0], None
            else:
                # Multi-line format
                # First non-empty line is the peer device name
                peer_device = lines[0]

                # Last non-empty line is the peer interface name
                peer_interface = lines[-1]

                return peer_device, peer_interface

        # --- Main Logic ---

        # 1. Parse self.full_filepath_csv (master data) to create a mapping of correct connection info
        # Key: (device1, device2) (sorted), Value: (interface of device1, interface of device2)
        master_link_map: Dict[Tuple[str, str], Tuple[str, str]] = {}

        try:
            with open(self.full_filepath_csv, 'r', newline='', encoding='utf-8') as f:
                reader = csv.DictReader(f)

                for row in reader:
                    # Only use 'Device', 'Interface', and 'Connection' columns
                    device = row.get('Device', '').strip()
                    interface_raw = row.get('Interface', '').strip()
                    connection_str = row.get('Connection', '').strip()

                    # Skip if any required field is missing
                    if not device or not interface_raw or not connection_str:
                        continue

                    interface = _process_interface_name(interface_raw)

                    # Skip if interface is invalid
                    if interface is None:
                        continue

                    # Parse the 'Connection' column to extract peer device and interface
                    peer_device, peer_interface_raw = _parse_connection_field(connection_str)

                    # For connections like "Level3 MPLS", peer_device will be "Level3 MPLS"
                    # and peer_interface_raw will be None
                    if peer_device:
                        if peer_interface_raw:
                            # Multi-line format with both device and interface
                            # Clean up peer interface name (remove parentheses content if any)
                            # Example: "fa0/1 (fa0/1)" -> "fa0/1"
                            peer_interface_cleaned = re.sub(r'\s*\([^)]*\)$', '', peer_interface_raw).strip()

                            # Process and format peer interface name
                            peer_interface = _process_interface_name(peer_interface_cleaned)

                            # Skip if peer interface is invalid
                            if peer_interface is None:
                                continue

                            # Check validity after processing
                            if _is_valid_peer_interface(peer_interface):
                                # Normalize device pair to create a consistent key (sorted alphabetically)
                                # Use device names from 'Device' column and parsed from 'Connection' column
                                if device < peer_device:
                                    key = (device, peer_device)
                                    value = (interface, peer_interface)
                                else:
                                    key = (peer_device, device)
                                    value = (peer_interface, interface)

                                # Store the correct interface information in the map
                                # This ensures we're using device names like 'Main Switch', not 'DOQHD947'
                                master_link_map[key] = value
                        else:
                            # Single-line format: only device name, no interface
                            # Example: "Level3 MPLS"
                            # We still need to store this connection but without peer interface
                            # Create a special marker for connections without peer interfaces
                            if device < peer_device:
                                key = (device, peer_device)
                                value = (interface, None)
                            else:
                                key = (peer_device, device)
                                value = (None, interface)

                            # Only store if this key doesn't already exist
                            # (to avoid overwriting valid interface connections)
                            if key not in master_link_map:
                                master_link_map[key] = value

        except FileNotFoundError:
            print(f"Error: Master CSV file not found at '{self.full_filepath_csv}'. Cannot generate commands.")
            return  # Exit if master file is not found
        except ValueError as e:
            print(f"Error: {e}")
            return  # Exit if all random numbers are exhausted
        except Exception as e:
            print(f"Error reading master CSV file '{self.full_filepath_csv}': {e}")
            return  # Exit on other file reading errors

        # 2. Process self.l1_connection_csv and generate commands
        try:
            for connection in self.l1_connection_csv:
                if len(connection) != 2:
                    # Skip malformed connection entries
                    print(f"Warning: Malformed connection entry skipped: {connection}")
                    continue

                devA, ifA = connection[0]
                devB, ifB = connection[1]

                # Keep original device names, only trim whitespace
                devA_clean = devA.strip()
                ifA_clean = _process_interface_name(ifA)
                devB_clean = devB.strip()
                ifB_clean = _process_interface_name(ifB)

                # Skip if either interface is invalid
                if ifA_clean is None or ifB_clean is None:
                    continue

                # Generate the key for lookup in the master map
                key_for_lookup = tuple(sorted((devA_clean, devB_clean)))

                if key_for_lookup in master_link_map:
                    correct_if_tuple = master_link_map[key_for_lookup]

                    correct_ifA_master = None
                    correct_ifB_master = None

                    # Determine the correct interface for each device based on the sorted key
                    if key_for_lookup[0] == devA_clean:
                        correct_ifA_master = correct_if_tuple[0]
                        correct_ifB_master = correct_if_tuple[1]
                    else:  # key_for_lookup[0] == devB_clean
                        correct_ifA_master = correct_if_tuple[1]
                        correct_ifB_master = correct_if_tuple[0]

                    # Generate command for Device A's port if it needs updating
                    # Skip if correct interface is None (single-line connection format)
                    if correct_ifA_master and ifA_clean != correct_ifA_master:
                        # Enclose device and port names with spaces in double quotes
                        device_name = f'"{devA_clean}"' if ' ' in devA_clean else devA_clean
                        original_port = f'"{ifA_clean}"' if ' ' in ifA_clean else ifA_clean
                        updated_port = f'"{correct_ifA_master}"' if ' ' in correct_ifA_master else correct_ifA_master
                        cmd = f"python {self.script_path} rename port {device_name} {original_port} {updated_port} --master {self.full_filepath}"
                        if cmd not in generated_commands:
                            self.csv_l1_rename_cmd_list.append(cmd)
                            generated_commands.add(cmd)

                    # Generate command for Device B's port if it needs updating
                    # Skip if correct interface is None (single-line connection format)
                    if correct_ifB_master and ifB_clean != correct_ifB_master:
                        device_name = f'"{devB_clean}"' if ' ' in devB_clean else devB_clean
                        original_port = f'"{ifB_clean}"' if ' ' in ifB_clean else ifB_clean
                        updated_port = f'"{correct_ifB_master}"' if ' ' in correct_ifB_master else correct_ifB_master
                        cmd = f"python {self.script_path} rename port {device_name} {original_port} {updated_port} --master {self.full_filepath}"
                        if cmd not in generated_commands:
                            self.csv_l1_rename_cmd_list.append(cmd)
                            generated_commands.add(cmd)

        except ValueError as e:
            print(f"Error: {e}")
            return  # Exit if all random numbers are exhausted

        print('Total interfaces ' + str(len(self.csv_l1_rename_cmd_list)))
        count_interfaces = len(self.csv_l1_rename_cmd_list)

        for cmd in self.csv_l1_rename_cmd_list:
            # Split command string into parts, handling quotes correctly
            parts = shlex.split(cmd)

            # Find --master option position
            master_index = None
            for i, part in enumerate(parts):
                if part == '--master':
                    master_index = i
                    break

            if master_index is None:
                print("Error: --master option not found")
                sys.exit(1)

            # Get master file path
            master_filepath = parts[master_index + 1]

            # Find rename position and extract arguments
            rename_index = parts.index('rename')
            args = parts[rename_index:master_index]

            # Process arguments
            processed_args = []
            for i, arg in enumerate(args):
                if i == 0:
                    # 'rename'
                    processed_args.append(arg)
                elif i == 1:
                    # 'port'
                    processed_args.append(arg)
                elif i == 2:
                    # Device name - keep as is, no cleaning
                    processed_args.append(arg)
                elif i in [3, 4]:
                    # Interface names - apply interface formatting and ensure space before final digits
                    converted = re.sub(r'([a-zA-Z]+)[-]?([/\d].*)', r'\1 \2', arg)
                    converted = re.sub(r'\s*\([^)]*\)$', '', converted).strip()
                    # Ensure space before final digit sequence (including / and .)
                    converted = _format_interface_with_space(converted)
                    if converted is None:
                        # Skip this command if interface is invalid
                        break
                    processed_args.append(converted)
                else:
                    processed_args.append(arg)

            # Skip if we broke out of the loop due to invalid interface
            if len(processed_args) < 5:
                continue

            # Generate arguments string for output
            args_str = ', '.join([f"'{arg}'" for arg in processed_args])

            # Generate final output
            output = f"network_sketcher_cli.ns_cli_run.cli_rename(self, self.full_filepath, [{args_str}])"
            parts_with_quotes = args_str.split(', ')
            stripped_parts = [part.strip("'") for part in parts_with_quotes]
            desired_parts = stripped_parts[2:]

            if ns_def.get_if_value(desired_parts[2]) != -1 and ns_def.get_if_value(desired_parts[1]) != -1:
                print(str(count_interfaces) + '.   ' + desired_parts[0] + '    ' + desired_parts[1] + ' --> ' +
                      desired_parts[2])
                count_interfaces -= 1

                # run the command
                import network_sketcher_cli
                exec(output)
            else:
                print('[ERROR]   ' + str(count_interfaces) + '.   ' + desired_parts[0] + '    ' + desired_parts[
                    1] + ' --> ' + desired_parts[2])
                count_interfaces -= 1

        return

    def __init__(self):
        """
        Network topology visualizer that converts CSV network data to PowerPoint diagrams.

        Usage:
            import ns_option_convert_to_master
            converter = ns_option_convert_to_master.ns_option_convert_to_master_csv()
            converter.set_output_path("output.pptx")
            converter.process("input.csv")
        """

        """Initialize the network topology visualizer."""
        self.tmp_pptx_file_path = None
        self.G = None
        self.clusters = {}
        self.node_tiers = {}
        self.pos = {}

        # Configuration parameters
        self.min_spacing = 3.0
        self.rect_width = 1.5
        self.rect_height = 0.8
        self.show_legend = True
        self.show_stats = True
        self.max_iterations = 1000

        # Color mapping for tiers
        self.color_map = {
            0: RGBColor(231, 76, 60),  # Red
            1: RGBColor(230, 126, 34),  # Orange
            2: RGBColor(243, 156, 18),  # Yellow
            3: RGBColor(155, 89, 182),  # Purple
            4: RGBColor(46, 204, 113),  # Green
            5: RGBColor(52, 152, 219),  # Blue
            6: RGBColor(26, 188, 156),  # Teal
            7: RGBColor(149, 165, 166),  # Gray
        }

        # Tier names
        self.tier_names = {
            0: "Internet/WAN",
            1: "Edge Router",
            2: "Core Layer",
            3: "Internal Router",
            4: "Distribution Layer",
            5: "Aggregation",
            6: "Access Layer",
            7: "Endpoint"
        }

        print("Network Topology Visualizer initialized")

    def set_output_path(self, output_path: str):
        """
        Set the output PowerPoint file path.

        Args:
            output_path: Path where the PowerPoint file will be saved
        """
        self.tmp_pptx_file_path = output_path

    def set_configuration(self, min_spacing: float = 3.0, rect_width: float = 1.5,
                          rect_height: float = 0.8, show_legend: bool = True,
                          show_stats: bool = True, max_iterations: int = 1000):
        """
        Set configuration parameters for the layout algorithm.

        Args:
            min_spacing: Minimum horizontal spacing between nodes
            rect_width: Width of device rectangles
            rect_height: Height of device rectangles
            show_legend: Whether to show legend in output
            show_stats: Whether to show statistics in output
            max_iterations: Maximum iterations for optimization
        """
        self.min_spacing = min_spacing
        self.rect_width = rect_width
        self.rect_height = rect_height
        self.show_legend = show_legend
        self.show_stats = show_stats
        self.max_iterations = max_iterations
        print(f"Configuration updated: spacing={min_spacing}, iterations={max_iterations}")

    @staticmethod
    def normalize_text(text: Optional[str]) -> str:
        """Normalize line breaks and consecutive spaces to a single space."""
        if text is None:
            return ""
        return re.sub(r"\s+", " ", str(text)).strip()

    @staticmethod
    def parse_connection_field(connection_str: str) -> Tuple[str, Optional[str]]:
        """
        Parse the 'Connection' column to extract peer device and interface.

        Handles two formats:
        1. Single line: Just the peer device name (e.g., "Level3 MPLS")
           Returns: (peer_device, None)

        2. Multi-line format:
           - Line 1: Peer device name (e.g., 'dmi01-yonkers-sw01', 'Main Switch')
           - Line 2-N: May contain whitespace or empty lines
           - Last line: Peer interface name (e.g., 'GigabitEthernet1/0/2', 'mgmt')
           Returns: (peer_device, peer_interface)

        Examples:
            Input: "Level3 MPLS"
            Output: ('Level3 MPLS', None)

            Input: "dmi01-yonkers-sw01\\n    \\n  \\n  GigabitEthernet1/0/2"
            Output: ('dmi01-yonkers-sw01', 'GigabitEthernet1/0/2')
        """
        if not connection_str:
            return "", None

        # Split by newline and filter out empty/whitespace-only lines
        lines = [line.strip() for line in connection_str.split('\n') if line.strip()]

        if len(lines) == 0:
            return "", None
        elif len(lines) == 1:
            # Single line - only device name, no interface
            return lines[0], None
        else:
            # Multi-line format
            peer_device = lines[0]
            peer_interface = lines[-1]

            # Clean up peer interface (remove parentheses content if any)
            peer_interface = re.sub(r'\s*\([^)]*\)$', '', peer_interface).strip()

            return peer_device, peer_interface

    @staticmethod
    def str_to_bool(s: Optional[str]) -> bool:
        """Convert True/False string to bool."""
        if s is None:
            return False
        return ns_option_convert_to_master_csv.normalize_text(s).lower() in \
            {"true", "1", "yes", "y"}

    @staticmethod
    def is_same_link(link: Dict[str, Any], a_dev: str, a_port: str,
                     b_dev: str, b_port: str) -> bool:
        """Determine if existing link and new link are the same physical link."""
        return (
                (link["src_device"] == a_dev and link["dst_device"] == b_dev and
                 link["src_port"] == a_port and link["dst_port"] == b_port)
                or
                (link["src_device"] == b_dev and link["dst_device"] == a_dev and
                 link["src_port"] == b_port and link["dst_port"] == a_port)
        )

    @staticmethod
    def extract_device_number(device_name: str) -> int:
        """Extract number from device name for sorting."""
        numbers = re.findall(r'\d+', device_name)
        if numbers:
            return int(numbers[-1])
        return 0

    @staticmethod
    def extract_device_base_name(device_name: str) -> str:
        """Extract base name without numbers for clustering detection."""
        base_name = re.sub(r'[-_]?\d+$', '', device_name)
        return base_name.lower()

    @staticmethod
    def get_redundant_keywords() -> Set[str]:
        """
        Return set of keywords that indicate redundant device types.
        Organized by network hierarchy and device categories.
        """
        return {
            # WAN/Internet/External Network keywords (Tier 0)
            'mpls', 'wan', 'internet', 'isp', 'carrier', 'provider',
            'transit', 'peering', 'upstream', 'backbone', 'metro',

            # Edge/Border keywords (Tier 1)
            'edge', 'border', 'bdr', 'dmz', 'perimeter',

            # Core network layer keywords (Tier 2)
            'core', 'spine', 'dcn',

            # Distribution/Aggregation keywords (Tier 3-4)
            'dist', 'distribution', 'agg', 'aggregation', 'distrib',

            # Access layer keywords (Tier 5-6)
            'acc', 'access', 'leaf', 'tor', 'floor', 'closet', 'idf',

            # Device type keywords
            'sw', 'switch', 'rtr', 'router', 'rt', 'gw', 'gateway',
            'fw', 'firewall', 'fwall', 'lb', 'loadbalancer', 'balancer',
            'vpn', 'wlc', 'wireless', 'ctl', 'ctrl', 'controller',
            'ap', 'accesspoint', 'srv', 'server', 'host',

            # Security keywords
            'sec', 'security', 'ids', 'ips', 'asa', 'ngfw', 'utm',

            # Redundancy/High Availability keywords
            'pri', 'primary', 'sec', 'secondary', 'std', 'standby',
            'actv', 'active', 'prim', 'back', 'backup', 'ha', 'red', 'redundant',
            'main', 'spare', 'hot', 'cold', 'warm',

            # Position/Location keywords
            'top', 'bot', 'bottom', 'left', 'right', 'rght',
            'east', 'west', 'north', 'south', 'nth', 'sth',
            'a', 'b', 'c', 'd', 'side',

            # MPLS/BGP/Routing keywords
            'abr', 'asbr', 'pe', 'ce', 'p', 'rr', 'reflector', 'bgp',

            # Data Center keywords
            'dc', 'datacenter', 'pod', 'rack', 'row', 'zone',

            # Campus/Building keywords
            'bldg', 'building', 'campus', 'site', 'branch', 'remote',

            # Stack/Cluster keywords
            'stack', 'cluster', 'vss', 'vpc', 'mlag', 'lag',

            # Service/Application keywords
            'dmz', 'mgmt', 'management', 'oob', 'console', 'admin',
            'prod', 'production', 'dev', 'test', 'staging', 'lab'
        }

    def get_wan_keywords(self) -> Set[str]:
        """
        Return set of keywords that specifically identify WAN/Internet/External connections.
        These devices should be placed at Tier 0 (topmost layer).
        """
        return {
            'mpls', 'wan', 'internet', 'isp', 'carrier', 'provider',
            'transit', 'peering', 'upstream', 'backbone', 'metro',
            'l3vpn', 'vpls', 'evpn', 'sdwan', 'overlay'
        }

    def detect_device_clusters(self) -> Dict[str, Set[str]]:
        """Detect device clusters for redundant pairs only."""
        print("\n=== Detecting Device Clusters ===")

        if self.G.number_of_nodes() == 0:
            print("  No nodes in graph")
            return {}

        base_name_groups = defaultdict(set)
        for node in self.G.nodes():
            base_name = self.extract_device_base_name(node)
            base_name_groups[base_name].add(node)

        clusters = {}
        cluster_id = 0
        redundant_keywords = self.get_redundant_keywords()
        wan_keywords = self.get_wan_keywords()

        for base_name, devices in base_name_groups.items():
            if len(devices) < 2:
                continue

            devices_list = list(devices)

            # Skip WAN/MPLS devices from clustering (they should remain independent at Tier 0)
            if any(wan_kw in base_name.lower() for wan_kw in wan_keywords):
                print(f"  Skipping WAN device clustering: {base_name}")
                continue

            # Check if devices have matching redundant keywords
            matching_keyword = None
            for keyword in redundant_keywords:
                if all(keyword in d.lower() for d in devices_list):
                    matching_keyword = keyword
                    break

            if not matching_keyword:
                continue

            # Verify connectivity pattern
            degrees = [self.G.degree(d) for d in devices_list]
            avg_degree = sum(degrees) / len(degrees) if degrees else 0

            similar_devices = set()
            for device in devices_list:
                if avg_degree == 0 or abs(self.G.degree(device) - avg_degree) <= max(avg_degree * 0.5, 2):
                    similar_devices.add(device)

            if len(similar_devices) >= 2:
                clusters[f"cluster_{cluster_id}"] = similar_devices
                cluster_id += 1
                print(f"  Cluster: {base_name} ({len(similar_devices)} devices)")

        print(f"Total clusters: {len(clusters)}")
        return clusters

    def calculate_tier_by_centrality(self) -> Dict[str, int]:
        """Calculate tier based on centrality metrics."""
        if self.G.number_of_nodes() == 0:
            return {}

        degree_centrality = nx.degree_centrality(self.G)
        betweenness_centrality = nx.betweenness_centrality(self.G)
        closeness_centrality = nx.closeness_centrality(self.G)

        try:
            eigenvector_centrality = nx.eigenvector_centrality(self.G, max_iter=1000)
        except:
            eigenvector_centrality = {node: 0 for node in self.G.nodes()}

        node_scores = {}
        wan_keywords = self.get_wan_keywords()

        for node in self.G.nodes():
            # Force WAN/MPLS nodes to highest priority (Tier 0)
            if any(wan_kw in node.lower() for wan_kw in wan_keywords):
                node_scores[node] = 1.0  # Maximum score
                continue

            score = (
                    degree_centrality[node] * 0.3 +
                    betweenness_centrality[node] * 0.4 +
                    closeness_centrality[node] * 0.2 +
                    eigenvector_centrality[node] * 0.1
            )
            node_scores[node] = score

        sorted_nodes = sorted(node_scores.items(), key=lambda x: x[1], reverse=True)

        node_tiers = {}
        num_nodes = len(sorted_nodes)

        for i, (node, score) in enumerate(sorted_nodes):
            # WAN/MPLS devices always at Tier 0
            if any(wan_kw in node.lower() for wan_kw in wan_keywords):
                node_tiers[node] = 0
                continue

            percentile = i / num_nodes if num_nodes > 0 else 0

            if percentile < 0.05:
                tier = 1
            elif percentile < 0.15:
                tier = 2
            elif percentile < 0.30:
                tier = 3
            elif percentile < 0.50:
                tier = 4
            elif percentile < 0.70:
                tier = 5
            elif percentile < 0.85:
                tier = 6
            else:
                tier = 7

            node_tiers[node] = tier

        return node_tiers

    def calculate_tier_by_device_role(self, node: str) -> int:
        """Calculate tier based on device role inferred from name and connectivity."""
        degree = self.G.degree(node)
        name_lower = node.lower()

        # Define keyword sets for each tier
        # Tier 0: Internet/WAN - Highest priority for external networks
        wan_keywords = {'mpls', 'wan', 'inet', 'isp', 'carr', 'prov',
                        'trans', 'peer', 'upstr', 'bckbn',
                        'extnet', 'pubcld', 'cldgw', 'egress', 'brdwan'}

        # Tier 1: Edge Router - Routers connected to WAN/MPLS
        edge_router_base_keywords = {'rtr', 'router', 'rt',
                                     'brdrtr', 'bdrrtr', 'perimr', 'branch', 'siter',
                                     'cer', 'vpnedg', 'inetedg', 'dmzrtr', 'ingres',
                                     'egres', 'mplsed'}

        # Tier 2: Core Layer
        core_keywords = {'core',
                         'spine', 'fabric', 'bckbn', 'centr', 'dccore',
                         'mainc', 'hspeed', 'routco', 'superc', 'netwco',
                         'enterc', 'distco', 'nxos', 'iosxe'}

        # Tier 4: Distribution Layer
        distribution_keywords = {'dist', 'agg', 'distrib',
                                 'agglay', 'campdi', 'bldist', 'flrdst', 'accagg',
                                 'l3dist', 'l3d', 'routdi', 'ivlr', 'policy',
                                 'catalyst', 'nexus'}

        # Tier 5: Aggregation switches
        aggregation_sw_base_keywords = {'sw', 'switch',
                                        'aggsw', 'stkmst', 'stkmem', 'clstagg', 'idfagg',
                                        'mdfagg', 'uplink', 'intsw', 'bldsw', 'flrsw',
                                        'dataagg', 'netagg', 'mlagg'}

        # Tier 2: Security/Gateway devices
        security_keywords = {'fw', 'firewall', 'lb', 'vpn', 'gw', 'gateway',
                             'secapp', 'utm', 'ids', 'ips', 'proxy',
                             'wlc', 'asa', 'webgw', 'emailgw'}

        # Tier 6: Access Layer
        access_keywords = {'edge', 'acc', 'access',
                           'accsw', 'client', 'user', 'wkstat', 'poe',
                           'desksw', 'flracc', 'clstac', 'portac', 'endpt',
                           'datapt', 'cubicl'}

        # Tier 3: Internal Router
        internal_router_keywords = {'intrt', 'camprt', 'ospf', 'eigrp', 'bgpint',
                                    'vlanrt', 'l3int', 'dcint', 'segmnt', 'zone',
                                    'privat', 'infra', 'servic', 'corefb', 'introu'}

        # Tier 7: Endpoints (low connectivity)
        endpoint_keywords = {'host', 'server', 'pc', 'wkstat', 'print',
                             'ipphon', 'camera', 'iot', 'ap', 'client',
                             'term', 'vm', 'vmhost', 'sensor', 'hmi'}

        # Tier 0: Internet/WAN - Highest priority for external networks
        if any(keyword in name_lower for keyword in wan_keywords):
            return 0

        # Tier 1: Edge Router - Routers connected to WAN/MPLS
        if any(keyword in name_lower for keyword in edge_router_base_keywords):
            neighbors = list(self.G.neighbors(node))
            has_external = any(
                any(wan_kw in str(self.G[node][n].get('links', [{}])[0].get('connection', '')).lower()
                    for wan_kw in wan_keywords)
                for n in neighbors if self.G.has_edge(node, n)
            )
            if has_external:
                return 1
            return 3

        # Tier 2: Core Layer
        if any(keyword in name_lower for keyword in core_keywords):
            return 2

        # Tier 4: Distribution Layer
        if any(keyword in name_lower for keyword in distribution_keywords):
            return 4

        # Tier 5: Aggregation switches
        if any(keyword in name_lower for keyword in aggregation_sw_base_keywords):
            if any(keyword in name_lower for keyword in access_keywords):
                return 6
            elif degree <= 3:
                return 6
            return 5

        # Tier 2: Security/Gateway devices
        if any(kw in name_lower for kw in security_keywords):
            return 2

        # Tier 7: Endpoints (low connectivity)
        if degree == 1 or any(keyword in name_lower for keyword in endpoint_keywords):
            return 7
        elif degree >= 8:
            return 3
        elif degree >= 4:
            return 4
        else:
            if any(keyword in name_lower for keyword in access_keywords):
                return 6
            return 6

    def calculate_integrated_tier(self) -> Dict[str, int]:
        """Calculate tiers with emphasis on separating different device types."""
        print("\n=== Calculating Tiers ===")

        if self.G.number_of_nodes() == 0:
            return {}

        tier_by_centrality = self.calculate_tier_by_centrality()

        node_tiers = {}

        for node in self.G.nodes():
            role_tier = self.calculate_tier_by_device_role(node)
            centrality_tier = tier_by_centrality.get(node, 6)
            combined_tier = int(role_tier * 0.7 + centrality_tier * 0.3)
            node_tiers[node] = combined_tier

        # Apply cluster constraint
        for cluster_id, cluster_devices in self.clusters.items():
            if not cluster_devices:
                continue

            cluster_tiers = [node_tiers[node] for node in cluster_devices if node in node_tiers]
            if cluster_tiers:
                cluster_tier = min(cluster_tiers)
                for node in cluster_devices:
                    node_tiers[node] = cluster_tier

        # Print tier distribution
        tier_counts = defaultdict(int)
        for tier in node_tiers.values():
            tier_counts[tier] += 1

        print("Tier Distribution:")
        for tier in sorted(tier_counts.keys()):
            print(f"  Tier {tier}: {tier_counts[tier]} nodes")

        return node_tiers

    @staticmethod
    def point_to_segment_distance(px: float, py: float, x1: float, y1: float,
                                  x2: float, y2: float) -> float:
        """Calculate minimum distance from point to line segment."""
        length_sq = (x2 - x1) ** 2 + (y2 - y1) ** 2

        if length_sq < 1e-10:
            return math.sqrt((px - x1) ** 2 + (py - y1) ** 2)

        t = max(0, min(1, ((px - x1) * (x2 - x1) + (py - y1) * (y2 - y1)) / length_sq))

        closest_x = x1 + t * (x2 - x1)
        closest_y = y1 + t * (y2 - y1)

        return math.sqrt((px - closest_x) ** 2 + (py - closest_y) ** 2)

    @staticmethod
    def do_segments_intersect(x1, y1, x2, y2, x3, y3, x4, y4) -> bool:
        """Determine if two line segments intersect."""

        def ccw(ax, ay, bx, by, cx, cy):
            return (cy - ay) * (bx - ax) > (by - ay) * (cx - ax)

        eps = 0.001
        if (abs(x1 - x3) < eps and abs(y1 - y3) < eps) or \
                (abs(x1 - x4) < eps and abs(y1 - y4) < eps) or \
                (abs(x2 - x3) < eps and abs(y2 - y3) < eps) or \
                (abs(x2 - x4) < eps and abs(y2 - y4) < eps):
            return False

        return (ccw(x1, y1, x3, y3, x4, y4) != ccw(x2, y2, x3, y3, x4, y4) and
                ccw(x1, y1, x2, y2, x3, y3) != ccw(x1, y1, x2, y2, x4, y4))

    def count_horizontal_connections(self, y_tolerance: float = 0.1) -> int:
        """Count horizontal connections that are NOT within valid clusters."""
        horizontal_count = 0

        for u, v in self.G.edges():
            if u not in self.pos or v not in self.pos:
                continue

            _, uy = self.pos[u]
            _, vy = self.pos[v]

            if abs(uy - vy) < y_tolerance:
                in_same_cluster = False
                for cluster_devices in self.clusters.values():
                    if u in cluster_devices and v in cluster_devices:
                        in_same_cluster = True
                        break

                if not in_same_cluster:
                    horizontal_count += 1

        return horizontal_count

    def count_edge_crossings(self) -> int:
        """Count number of edge crossings in the current layout."""
        edges = list(self.G.edges())
        crossings = 0

        for i, (u1, v1) in enumerate(edges):
            for u2, v2 in edges[i + 1:]:
                if u1 in self.pos and v1 in self.pos and u2 in self.pos and v2 in self.pos:
                    x1, y1 = self.pos[u1]
                    x2, y2 = self.pos[v1]
                    x3, y3 = self.pos[u2]
                    x4, y4 = self.pos[v2]

                    if self.do_segments_intersect(x1, y1, x2, y2, x3, y3, x4, y4):
                        crossings += 1

        return crossings

    def calculate_layout_score(self) -> Tuple[int, int, int]:
        """Calculate layout quality score."""
        horizontal_conn = self.count_horizontal_connections()
        crossings = self.count_edge_crossings()
        nodes_on_edges = 0

        return horizontal_conn, crossings, nodes_on_edges

    def calculate_vertical_alignment_score(self) -> float:
        """
        Calculate how well nodes are vertically aligned with their neighbors.
        Lower score is better (0 = perfect vertical alignment).

        Returns:
            Total horizontal distance between connected nodes
        """
        total_distance = 0.0
        edge_count = 0

        for u, v in self.G.edges():
            if u not in self.pos or v not in self.pos:
                continue

            ux, uy = self.pos[u]
            vx, vy = self.pos[v]

            # Only consider edges between different tiers (vertical connections)
            if abs(uy - vy) > 0.1:
                horizontal_distance = abs(ux - vx)
                total_distance += horizontal_distance
                edge_count += 1

        return total_distance / max(edge_count, 1)

    def calculate_combined_score(self, alpha=0.2, beta=0.2, gamma=0.6) -> Tuple[float, Dict[str, float]]:
        """
        Calculate combined score with PRIORITY on vertical alignment.

        Args:
            alpha: Weight for crossings (default 0.2) - reduced
            beta: Weight for horizontal connections (default 0.2) - reduced
            gamma: Weight for vertical alignment (default 0.6) - HIGHEST priority

        Returns:
            Tuple of (combined_score, score_details)
        """
        horiz, crossings, on_edges = self.calculate_layout_score()
        vertical_alignment = self.calculate_vertical_alignment_score()

        # Dramatically increase alignment score weight
        crossing_score = crossings * 20000  # Reduced from 30000
        horizontal_score = horiz * 30000  # Reduced from 50000
        alignment_score = vertical_alignment * 50000  # GREATLY increased from 15000

        combined = alpha * crossing_score + beta * horizontal_score + gamma * alignment_score

        details = {
            'crossings': crossings,
            'horizontal': horiz,
            'alignment': vertical_alignment,
            'crossing_score': crossing_score,
            'horizontal_score': horizontal_score,
            'alignment_score': alignment_score,
            'combined': combined
        }

        return combined, details

    def optimize_tier_positions(self, tier_nodes: List[str], base_y: float):
        """
        Optimize positions - VERTICAL ALIGNMENT PRIORITY.
        Places child nodes directly below their parents whenever possible.

        Args:
            tier_nodes: List of nodes in the tier
            base_y: Base Y coordinate for the tier
        """
        if len(tier_nodes) <= 1:
            return

        # Step 1: Identify clusters
        cluster_members = set()
        for cluster_devices in self.clusters.values():
            cluster_members.update(cluster_devices & set(tier_nodes))

        same_tier_cluster_edges = []
        for u, v in self.G.edges():
            if (u in tier_nodes and v in tier_nodes and
                    u in cluster_members and v in cluster_members):
                for cluster_devices in self.clusters.values():
                    if u in cluster_devices and v in cluster_devices:
                        same_tier_cluster_edges.append((u, v))
                        break

        subgroups = []
        processed = set()

        if same_tier_cluster_edges:
            cluster_subgraph = nx.Graph()
            cluster_subgraph.add_nodes_from(cluster_members)
            cluster_subgraph.add_edges_from(same_tier_cluster_edges)

            for component in nx.connected_components(cluster_subgraph):
                if len(component) > 1:
                    subgroups.append(component)
                    processed.update(component)

        standard_nodes = [n for n in tier_nodes if n not in processed]

        # Step 2: Assign Y coordinates
        y_offset = 0.5
        node_y_assignments = {}

        if len(subgroups) > 1:
            for i, group in enumerate(subgroups):
                y_adjustment = (i - (len(subgroups) - 1) / 2) * y_offset
                for node in group:
                    node_y_assignments[node] = base_y + y_adjustment
        elif len(subgroups) == 1:
            for node in subgroups[0]:
                node_y_assignments[node] = base_y + y_offset / 2

        for node in standard_nodes:
            node_y_assignments[node] = base_y

        # Step 3: Calculate target positions based on vertical neighbors
        def get_vertical_neighbor_avg(node):
            """Get average X position of vertical neighbors."""
            neighbors = [n for n in self.G.neighbors(node)
                         if n not in tier_nodes and n in self.pos]
            if not neighbors:
                return None

            weighted_sum = 0
            total_weight = 0
            for neighbor in neighbors:
                edge_data = self.G[node][neighbor]
                weight = len(edge_data.get('links', [1]))
                weighted_sum += self.pos[neighbor][0] * weight
                total_weight += weight

            return weighted_sum / total_weight if total_weight > 0 else None

        # Group by Y-level
        y_groups = defaultdict(list)
        for node in tier_nodes:
            y = node_y_assignments.get(node, base_y)
            y_groups[y].append(node)

        # Step 4: DIRECT PLACEMENT - Place each node at its target position
        print(f"      Direct vertical alignment...")

        for y_level, nodes_at_level in y_groups.items():
            # Calculate ideal X position for each node
            node_ideal_positions = []

            for node in nodes_at_level:
                target_x = get_vertical_neighbor_avg(node)

                if target_x is not None:
                    node_ideal_positions.append((node, target_x, True))  # Has vertical neighbor
                else:
                    current_x = self.pos.get(node, (0, y_level))[0]
                    node_ideal_positions.append((node, current_x, False))  # No vertical neighbor

            # Sort by ideal X position
            node_ideal_positions.sort(key=lambda x: x[1])

            # Place nodes with proper spacing, preserving relative order
            total_width = (len(node_ideal_positions) - 1) * (self.rect_width + self.min_spacing)
            start_x = -total_width / 2

            for i, (node, ideal_x, has_neighbor) in enumerate(node_ideal_positions):
                x = start_x + i * (self.rect_width + self.min_spacing)
                self.pos[node] = (x, y_level)

        # Step 5: Calculate how well we're aligned
        alignments_before = {}
        for node in tier_nodes:
            target_x = get_vertical_neighbor_avg(node)
            if target_x is not None:
                current_x = self.pos[node][0]
                alignments_before[node] = abs(current_x - target_x)

        # Step 6: Try to improve alignment by shifting entire groups
        print(f"      Adjusting group positions for better alignment...")

        best_layout = self.pos.copy()
        best_alignment_score = self.calculate_vertical_alignment_score()
        best_crossings = self.count_edge_crossings()

        # Try different shift amounts for the entire tier
        for shift in [-9.0, -6.75, -4.5, -2.25, 0, 2.25, 4.5, 6.75, 9.0]:
            test_pos = self.pos.copy()

            # Shift all nodes in this tier
            for node in tier_nodes:
                x, y = test_pos[node]
                test_pos[node] = (x + shift, y)

            # Evaluate
            self.pos = test_pos
            new_alignment = self.calculate_vertical_alignment_score()
            new_crossings = self.count_edge_crossings()

            # Accept if alignment improves and crossings don't increase too much
            if new_alignment < best_alignment_score and new_crossings <= best_crossings + 2:
                best_layout = test_pos.copy()
                best_alignment_score = new_alignment
                best_crossings = new_crossings
                print(f"        Shift {shift:+.2f}: Alignment improved to {new_alignment:.2f}")

        self.pos = best_layout.copy()

        # Step 7: Try individual node shifts
        print(f"      Fine-tuning individual node positions...")

        for node in tier_nodes:
            target_x = get_vertical_neighbor_avg(node)
            if target_x is None:
                continue

            current_x, y = self.pos[node]
            y_level = y
            nodes_at_level = [n for n in tier_nodes if abs(self.pos[n][1] - y_level) < 0.01]

            # Try moving this node closer to target
            for shift_factor in [0.9, 0.7, 0.5, 0.3]:  # Try different amounts
                shifted_x = current_x * (1 - shift_factor) + target_x * shift_factor

                # Create test position
                test_pos = self.pos.copy()
                test_pos[node] = (shifted_x, y)

                # Re-sort and enforce spacing
                nodes_with_x = [(n, test_pos[n][0]) for n in nodes_at_level]
                nodes_with_x.sort(key=lambda x: x[1])

                total_width = (len(nodes_with_x) - 1) * (self.rect_width + self.min_spacing)
                start_x = -total_width / 2

                for i, (n, _) in enumerate(nodes_with_x):
                    x = start_x + i * (self.rect_width + self.min_spacing)
                    test_pos[n] = (x, y_level)

                # Evaluate
                self.pos = test_pos
                new_alignment = self.calculate_vertical_alignment_score()
                new_crossings = self.count_edge_crossings()

                # Accept if improves alignment without adding crossings
                if new_alignment < best_alignment_score and new_crossings <= best_crossings:
                    best_layout = test_pos.copy()
                    best_alignment_score = new_alignment
                    best_crossings = new_crossings
                    print(f"        {node}: Moved {shift_factor * 100:.0f}% toward parent")
                    break  # Found improvement for this node

        self.pos = best_layout.copy()

        # Step 8: Final score
        final_score, final_details = self.calculate_combined_score()

        print(f"    Final - Crossings: {final_details['crossings']}, "
              f"Horizontal: {final_details['horizontal']}, "
              f"Alignment: {final_details['alignment']:.2f}")

        # Debug output
        print(f"    Vertical alignment details:")
        improvements_count = 0
        for node in tier_nodes:
            target_x = get_vertical_neighbor_avg(node)
            if target_x is not None:
                node_x = self.pos[node][0]
                offset = abs(node_x - target_x)
                before_offset = alignments_before.get(node, offset)

                if offset > 0.5:
                    improvement = before_offset - offset
                    if improvement > 0.1:
                        print(f"      {node}: offset={offset:.2f} (improved by {improvement:.2f})")
                        improvements_count += 1
                    else:
                        print(f"      {node}: offset={offset:.2f}")

        if improvements_count > 0:
            print(f"    Improved alignment for {improvements_count} nodes")

        if same_tier_cluster_edges:
            print(f"    Found {len(same_tier_cluster_edges)} same-tier cluster connections")


    def detect_network_groups(self) -> List[Set[str]]:
        """
        Detect separate network groups (connected components).

        Returns:
            List of sets, where each set contains device names in a group
        """
        print("\n=== Detecting Network Groups ===")

        if self.G.number_of_nodes() == 0:
            print("  No nodes in graph")
            return []

        # Find connected components
        components = list(nx.connected_components(self.G))

        print(f"Found {len(components)} network group(s):")
        for i, component in enumerate(components, 1):
            print(f"  Group {i}: {len(component)} devices")
            sample_devices = list(component)[:5]
            print(f"    Sample devices: {', '.join(sample_devices)}")
            if len(component) > 5:
                print(f"    ... and {len(component) - 5} more")

        return components

    def calculate_positions_for_group(self, group_nodes: Set[str]):
        """
        Calculate hierarchical positions using BOTTOM-UP approach.
        Places parent nodes above their children for better vertical alignment.

        Args:
            group_nodes: Set of node names in the group
        """
        print(f"\n=== Calculating Layout for Group ({len(group_nodes)} nodes) ===")

        if len(group_nodes) == 0:
            print("  ERROR: No nodes in group")
            return

        # Create subgraph for this group
        subgraph = self.G.subgraph(group_nodes).copy()

        # Temporarily replace self.G with subgraph for calculations
        original_graph = self.G
        self.G = subgraph

        # Detect clusters within this group
        self.clusters = self.detect_device_clusters()

        # Calculate tiers for this group
        self.node_tiers = self.calculate_integrated_tier()

        # Group nodes by tier
        tiers = defaultdict(list)
        for node, tier in self.node_tiers.items():
            tiers[tier].append(node)

        for tier in tiers:
            tiers[tier].sort(key=lambda x: (self.extract_device_number(x), x))

        self.pos = {}
        y_spacing = 3.0
        max_tier = max(tiers.keys()) if tiers else 0

        print("  Using BOTTOM-UP placement for better vertical alignment...")

        tier_order = sorted(tiers.keys())

        # Helper function to get child positions
        def get_child_x_position(node):
            """Get average X position of already-placed child nodes."""
            neighbors = [n for n in self.G.neighbors(node) if n in self.pos]
            if not neighbors:
                return None

            weighted_sum = 0
            total_weight = 0
            for neighbor in neighbors:
                edge_data = self.G[node][neighbor]
                weight = len(edge_data.get('links', [1]))
                weighted_sum += self.pos[neighbor][0] * weight
                total_weight += weight

            return weighted_sum / total_weight if total_weight > 0 else None

        # STRATEGY: Place bottom tiers first, then place parents above their children

        # Find the bottom-most tier (highest tier number)
        bottom_tier = max(tier_order)

        # Place bottom tier evenly
        nodes = tiers[bottom_tier]
        y = (max_tier - bottom_tier) * y_spacing

        num_nodes = len(nodes)
        if num_nodes == 1:
            self.pos[nodes[0]] = (0, y)
        else:
            total_width = (num_nodes - 1) * (self.rect_width + self.min_spacing)
            start_x = -total_width / 2
            for i, node in enumerate(nodes):
                x = start_x + i * (self.rect_width + self.min_spacing)
                self.pos[node] = (x, y)
                print(f"    {node}: placed at x={x:.2f} (bottom tier)")

        # Now place remaining tiers from bottom to top
        for tier in reversed(tier_order[:-1]):  # Exclude the bottom tier we just placed
            nodes = tiers[tier]
            y = (max_tier - tier) * y_spacing

            # For each node, try to place it above its children
            node_positions = []

            for node in nodes:
                child_x = get_child_x_position(node)
                if child_x is not None:
                    node_positions.append((node, child_x, True))
                    print(f"    {node}: target (above children) = {child_x:.2f}")
                else:
                    # No children placed yet - use center
                    node_positions.append((node, 0, False))

            # Sort by target position
            node_positions.sort(key=lambda x: x[1])

            # Place with spacing
            num_nodes = len(node_positions)
            total_width = (num_nodes - 1) * (self.rect_width + self.min_spacing)
            start_x = -total_width / 2

            for i, (node, target_x, has_children) in enumerate(node_positions):
                x = start_x + i * (self.rect_width + self.min_spacing)
                self.pos[node] = (x, y)

                if has_children:
                    actual_offset = abs(x - target_x)
                    print(f"    {node}: placed at x={x:.2f}, offset from children={actual_offset:.2f}")

        # Calculate initial score
        initial_horiz, initial_crossings, initial_on_edges = self.calculate_layout_score()
        initial_alignment = self.calculate_vertical_alignment_score()
        print(f"\n  After bottom-up placement:")
        print(f"    Crossings: {initial_crossings}, Horizontal: {initial_horiz}, Alignment: {initial_alignment:.2f}")

        # Refinement: Try to improve alignment further
        print("\n=== Refinement Optimization ===")

        best_pos = self.pos.copy()
        best_score, best_details = self.calculate_combined_score()

        for refinement_pass in range(3):
            print(f"\n  Refinement Pass {refinement_pass + 1}/3:")
            improved_this_pass = False

            # Go through each tier (except top and bottom)
            for tier in tier_order[1:-1]:
                nodes = tiers[tier]
                if len(nodes) <= 1:
                    continue

                y = (max_tier - tier) * y_spacing

                for node in nodes:
                    child_x = get_child_x_position(node)
                    if child_x is None:
                        continue

                    current_x, current_y = self.pos[node]
                    current_offset = abs(current_x - child_x)

                    # Skip if already well-aligned
                    if current_offset < 1.0:
                        continue

                    # Try moving toward children
                    for factor in [0.9, 0.7, 0.5]:
                        new_x = current_x * (1 - factor) + child_x * factor

                        test_pos = self.pos.copy()
                        test_pos[node] = (new_x, current_y)

                        # Re-sort tier to maintain spacing
                        tier_nodes_with_x = [(n, test_pos[n][0]) for n in nodes]
                        tier_nodes_with_x.sort(key=lambda x: x[1])

                        total_width = (len(tier_nodes_with_x) - 1) * (self.rect_width + self.min_spacing)
                        start_x = -total_width / 2

                        for i, (n, _) in enumerate(tier_nodes_with_x):
                            x = start_x + i * (self.rect_width + self.min_spacing)
                            test_pos[n] = (x, current_y)

                        # Evaluate
                        self.pos = test_pos
                        new_score, new_details = self.calculate_combined_score()

                        # Accept if improves
                        if new_score < best_score:
                            best_pos = test_pos.copy()
                            best_score = new_score
                            best_details = new_details
                            improved_this_pass = True
                            new_offset = abs(test_pos[node][0] - child_x)
                            print(f"    {node}: Improved alignment {current_offset:.2f} -> {new_offset:.2f}")
                            break
                        else:
                            self.pos = best_pos.copy()

            if not improved_this_pass:
                print(f"    No improvements in pass {refinement_pass + 1}")
                break

        # Apply best layout
        self.pos = best_pos.copy()

        # Final score
        final_score, final_details = self.calculate_combined_score()

        print(f"\n  Final Results:")
        print(f"    Crossings: {final_details['crossings']} (was {initial_crossings})")
        print(f"    Horizontal connections: {final_details['horizontal']} (was {initial_horiz})")
        print(f"    Vertical alignment: {final_details['alignment']:.2f} (was {initial_alignment:.2f})")

        alignment_improvement = initial_alignment - final_details['alignment']
        if alignment_improvement > 0:
            print(
                f"    Alignment improved by: {alignment_improvement:.2f} ({alignment_improvement / max(initial_alignment, 0.01) * 100:.1f}%)")
        elif alignment_improvement < 0:
            print(f"    Alignment decreased by: {-alignment_improvement:.2f}")

        # Show final alignment details for significant offsets
        print(f"\n  Final vertical alignment (showing offsets > 2.0):")

        def get_neighbor_avg(node):
            neighbors = [n for n in self.G.neighbors(node) if n in self.pos and n not in tiers[self.node_tiers[node]]]
            if not neighbors:
                return None
            return sum(self.pos[n][0] for n in neighbors) / len(neighbors)

        significant_offsets = []
        for tier in tier_order:
            for node in tiers[tier]:
                neighbor_avg = get_neighbor_avg(node)
                if neighbor_avg is not None:
                    node_x = self.pos[node][0]
                    offset = abs(node_x - neighbor_avg)
                    if offset > 2.0:
                        significant_offsets.append((node, offset, neighbor_avg, node_x))

        if significant_offsets:
            for node, offset, target, actual in significant_offsets:
                print(f"    {node}: offset={offset:.2f} (target={target:.2f}, actual={actual:.2f})")
        else:
            print(f"    All nodes well-aligned (offsets < 2.0)")

        # Restore original graph
        self.G = original_graph

    def calculate_positions(self):
        """Calculate hierarchical positions for all nodes."""
        print("\n=== Calculating Layout ===")

        if self.G.number_of_nodes() == 0:
            print("  ERROR: No nodes in graph")
            return

        self.clusters = self.detect_device_clusters()
        self.node_tiers = self.calculate_integrated_tier()

        # Group nodes by tier
        tiers = defaultdict(list)
        for node, tier in self.node_tiers.items():
            tiers[tier].append(node)

        for tier in tiers:
            tiers[tier].sort(key=lambda x: (self.extract_device_number(x), x))

        self.pos = {}
        y_spacing = 3.0
        max_tier = max(tiers.keys()) if tiers else 0

        # Initial placement
        for tier, nodes in tiers.items():
            y = (max_tier - tier) * y_spacing

            num_nodes = len(nodes)

            if num_nodes == 1:
                self.pos[nodes[0]] = (0, y)
            else:
                total_width = (num_nodes - 1) * (self.rect_width + self.min_spacing)
                start_x = -total_width / 2

                for i, node in enumerate(nodes):
                    x = start_x + i * (self.rect_width + self.min_spacing)
                    self.pos[node] = (x, y)

        # Optimize each tier
        print("\n=== Optimizing Positions ===")
        for tier in sorted(tiers.keys()):
            nodes = tiers[tier]
            if len(nodes) > 1:
                y = (max_tier - tier) * y_spacing
                print(f"  Optimizing tier {tier}...")
                self.optimize_tier_positions(nodes, y)

    def build_graph_from_csv(self, csv_path: str):
        """Read CSV and build network graph using only Device, Interface, and Connection columns."""
        print(f"\nReading CSV: {csv_path}")

        self.G = nx.Graph()
        rows_processed = 0
        links_created = 0

        try:
            with open(csv_path, "r", encoding="utf-8", newline="") as f:
                reader = csv.DictReader(f)

                fieldnames = reader.fieldnames
                if not fieldnames:
                    print("  ERROR: Empty CSV")
                    return

                for row in reader:
                    rows_processed += 1

                    device = self.normalize_text(row.get("Device", ""))
                    if not device:
                        continue

                    interface = self.normalize_text(row.get("Interface", ""))
                    connection_str = row.get("Connection", "")

                    if not connection_str:
                        continue

                    peer_device, peer_port = self.parse_connection_field(connection_str)

                    if not peer_device:
                        continue

                    self.G.add_node(device)
                    self.G.add_node(peer_device)

                    reachable = self.str_to_bool(row.get("Reachable", ""))
                    idval = self.normalize_text(row.get("ID", ""))
                    cable = self.normalize_text(row.get("Cable", ""))
                    cable_color = self.normalize_text(row.get("Cable Color", ""))

                    link_detail = {
                        "src_device": device,
                        "src_port": interface,
                        "dst_device": peer_device,
                        "dst_port": peer_port or "",
                        "ids": [idval] if idval else [],
                        "reachable": reachable,
                        "connection": self.normalize_text(connection_str),
                        "cable": cable,
                        "cable_color": cable_color,
                    }

                    if self.G.has_edge(device, peer_device):
                        links: List[Dict[str, Any]] = self.G[device][peer_device].setdefault("links", [])
                        merged = False
                        for l in links:
                            if self.is_same_link(l, device, interface, peer_device, peer_port or ""):
                                if idval and idval not in l["ids"]:
                                    l["ids"].append(idval)
                                l["reachable"] = l["reachable"] or reachable
                                merged = True
                                break
                        if not merged:
                            links.append(link_detail)
                            links_created += 1
                    else:
                        self.G.add_edge(device, peer_device, links=[link_detail])
                        links_created += 1

            print(f"  Loaded: {self.G.number_of_nodes()} nodes, {self.G.number_of_edges()} edges")

        except Exception as e:
            print(f"  ERROR: {str(e)}")

    def create_powerpoint(self):
        """Create PowerPoint presentation with separate slides for each network group."""
        print("\n=== Creating PowerPoint ===")

        if self.G.number_of_nodes() == 0:
            print("  ERROR: No nodes")
            return False

        if not self.tmp_pptx_file_path:
            print("  ERROR: Output path not set")
            return False

        # Detect network groups
        network_groups = self.detect_network_groups()

        if not network_groups:
            print("  ERROR: No network groups found")
            return False

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        blank_layout = prs.slide_layouts[6]

        # Process each group
        for group_idx, group_nodes in enumerate(network_groups, 1):
            print(f"\n=== Processing Group {group_idx}/{len(network_groups)} ===")

            # Calculate positions for this group
            self.calculate_positions_for_group(group_nodes)

            if not self.pos:
                print(f"  WARNING: No positions calculated for group {group_idx}")
                continue

            # Create slide for this group
            slide = prs.slides.add_slide(blank_layout)

            # Calculate transformation for this group
            group_pos = {node: pos for node, pos in self.pos.items() if node in group_nodes}

            if not group_pos:
                continue

            all_x = [x for x, y in group_pos.values()]
            all_y = [y for x, y in group_pos.values()]

            min_x, max_x = min(all_x), max(all_x)
            min_y, max_y = min(all_y), max(all_y)

            margin_x = self.rect_width * 2
            margin_y = self.rect_height * 2

            range_x = max_x - min_x + margin_x * 2
            range_y = max_y - min_y + margin_y * 2

            slide_width_in = prs.slide_width / 914400.0
            slide_height_in = prs.slide_height / 914400.0

            scale_in = min(slide_width_in / range_x, slide_height_in / range_y) * 0.85

            if range_x > 20 or range_y > 15:
                scale_in = 0.4

            def to_ppt_coords_in(x, y):
                px_in = (x - min_x + margin_x) * scale_in
                py_in = (max_y - y + margin_y) * scale_in
                return px_in, py_in

            # Draw nodes for this group
            for node in group_nodes:
                if node not in self.pos:
                    continue

                x, y = self.pos[node]
                px_in, py_in = to_ppt_coords_in(x, y)

                shape_w_in = self.rect_width * scale_in
                shape_h_in = self.rect_height * scale_in

                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(px_in - shape_w_in / 2),
                    Inches(py_in - shape_h_in / 2),
                    Inches(shape_w_in),
                    Inches(shape_h_in)
                )

                tier = self.node_tiers.get(node, 7)
                fill_color = self.color_map.get(tier, RGBColor(52, 73, 94))
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_color

                shape.line.color.rgb = RGBColor(44, 62, 80)
                shape.line.width = Pt(2)

                text_frame = shape.text_frame
                text_frame.clear()
                text_frame.word_wrap = True
                text_frame.margin_left = Pt(5)
                text_frame.margin_right = Pt(5)
                text_frame.margin_top = Pt(5)
                text_frame.margin_bottom = Pt(5)

                p = text_frame.paragraphs[0]
                p.text = node
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER

                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Draw edges for this group
            for u, v, data in self.G.edges(data=True):
                # Only draw edges where both nodes are in this group
                if u not in group_nodes or v not in group_nodes:
                    continue

                if u not in self.pos or v not in self.pos:
                    continue

                ux, uy = self.pos[u]
                vx, vy = self.pos[v]

                px1_in, py1_in = to_ppt_coords_in(ux, uy)
                px2_in, py2_in = to_ppt_coords_in(vx, vy)

                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(px1_in),
                    Inches(py1_in),
                    Inches(px2_in),
                    Inches(py2_in)
                )

                links = data.get("links", [])
                reachable = any(link.get("reachable", False) for link in links)

                line = connector.line
                if reachable:
                    line.color.rgb = RGBColor(39, 174, 96)
                    line.dash_style = MSO_LINE_DASH_STYLE.SOLID
                else:
                    line.color.rgb = RGBColor(231, 76, 60)
                    line.dash_style = MSO_LINE_DASH_STYLE.DASH

                line_width_pt = min(1.5 + len(links) * 0.5, 5)
                line.width = Pt(line_width_pt)

            # Add title to slide
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                Inches(10), Inches(0.5)
            )
            text_frame = title_box.text_frame
            p = text_frame.paragraphs[0]

            # Create descriptive title based on group devices
            sample_devices = sorted(list(group_nodes))[:3]
            if len(network_groups) > 1:
                p.text = f"Network Topology - Group {group_idx} ({len(group_nodes)} devices)"
            else:
                p.text = "Network Topology Diagram"

            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(44, 62, 80)

            # Add group info subtitle
            if len(network_groups) > 1:
                subtitle_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.7),
                    Inches(10), Inches(0.3)
                )
                text_frame = subtitle_box.text_frame
                p = text_frame.paragraphs[0]
                p.text = f"Devices: {', '.join(sample_devices)}" + (
                    f" and {len(group_nodes) - 3} more" if len(group_nodes) > 3 else ""
                )
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(100, 100, 100)

        # Save presentation
        try:
            prs.save(self.tmp_pptx_file_path)
            #print(f"\n  Saved: {self.tmp_pptx_file_path}")
            print(f"  Total slides created: {len(network_groups)}")
            return True
        except Exception as e:
            print(f"  ERROR: {str(e)}")
            return False

    def process(self, csv_path: str, output_path: str = None):
        """
        Main processing method.

        Args:
            csv_path: Path to input CSV file
            output_path: Path to output PowerPoint file (optional)

        Returns:
            bool: True if successful, False otherwise
        """
        print("\n" + "=" * 70)
        print("Network Topology Visualizer")
        print("=" * 70)

        # Set output path
        if output_path:
            self.set_output_path(output_path)
        elif not self.tmp_pptx_file_path:
            csv_file = Path(csv_path)
            self.set_output_path(str(csv_file.parent / f"{csv_file.stem}_topology.pptx"))

        # Build graph
        self.build_graph_from_csv(csv_path)

        if self.G.number_of_nodes() == 0:
            print("\nERROR: No nodes loaded")
            return False

        # Create PowerPoint (now handles multiple groups automatically)
        success = self.create_powerpoint()

        if success:
            print("\n" + "=" * 70)
            print("Process Complete! (Phase1)")
            print("=" * 70)

        return success





class  ns_option_convert_to_master_svg():
    def __init__(self):
        #convert svg to rough sketch with pptx
        import re
        import math
        from xml.etree import ElementTree as ET
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
        from pptx.dml.color import RGBColor

        # Prefer svg.path for robust path parsing (first Line.start and last Line.end).
        try:
            from svg.path import parse_path as svg_parse_path, Line as SvgLine
            SVG_PATH_AVAILABLE = True
        except Exception:
            SVG_PATH_AVAILABLE = False

        full_filepath = self.full_filepath

        '''other types support'''
        import xml.etree.ElementTree as ET
        import os

        if not os.path.exists(full_filepath):
            print(f"Error: The file '{full_filepath}' does not exist.")
            return

        # Parse the SVG file
        try:
            tree = ET.parse(full_filepath)
            root = tree.getroot()
        except Exception as e:
            print(f"ERROR: Failed to parse SVG: {e}")
            return

        # List to store target groups for transformation
        groups_to_transform = []

        def get_attr_ignore_namespace(element, attr_name):
            """Get attribute value ignoring namespace"""
            for key, value in element.attrib.items():
                if key.endswith('}' + attr_name) or key == attr_name:
                    return value
            return None

        def find_target_groups(current_element, parent_element=None):
            # Process only 'g' elements
            if current_element.tag.endswith('}g') or current_element.tag == 'g':
                group_context = get_attr_ignore_namespace(current_element, 'groupContext')

                if group_context == 'group':
                    #print(f"  ✓ Match found!")
                    groups_to_transform.append((parent_element, current_element))
                    # Remove 'return' to continue searching nested groups
                    # return  # <-- REMOVE THIS LINE

            # Continue searching in child elements
            for child in current_element:
                find_target_groups(child, current_element)

        # Start searching from root element
        find_target_groups(root)

        if groups_to_transform:
            #print(f"\nFound {len(groups_to_transform)} groups to transform.")

            for parent_of_target_group, target_group in groups_to_transform:
                # Extract original attributes
                original_id = target_group.get('id')
                original_transform = target_group.get('transform')
                original_v_mID = get_attr_ignore_namespace(target_group, 'mID')

                # Get title text
                original_title_element = None
                for child in target_group:
                    if child.tag.endswith('}title') or child.tag == 'title':
                        original_title_element = child
                        break
                original_title_text = original_title_element.text if original_title_element is not None else "Unknown Sheet"

                # Prepare attributes for new <g> element
                new_g_attrib = {
                    'id': original_id.replace('group', 'shape') if original_id else 'shape-transformed',
                }

                # OPTION 1: Remove transform attribute completely
                # This will place shapes at their default positions
                # (Do not add 'transform' to new_g_attrib)

                # OPTION 2: Keep transform but adjust Y coordinate
                # Parse and modify the transform value
                if original_transform:
                    import re
                    match = re.search(r'translate\(([-\d.]+),([-\d.]+)\)', original_transform)
                    if match:
                        x_val = float(match.group(1))
                        y_val = float(match.group(2))

                        # Adjust Y coordinate (add offset to move down)
                        # 10.2 inches = 734.4 points (at 72 DPI)
                        # Adjust this value based on your needs
                        y_offset = 734.4  # 10.2 inches in points
                        adjusted_y = y_val + y_offset

                        new_transform = f'translate({x_val},{adjusted_y})'
                        new_g_attrib['transform'] = new_transform
                        #print(f"  Adjusted transform: {original_transform} → {new_transform}")
                    else:
                        # If transform format is different, keep original
                        new_g_attrib['transform'] = original_transform

                # Add v:mID attribute (with namespace)
                if original_v_mID:
                    for key in target_group.attrib.keys():
                        if 'mID' in key:
                            new_g_attrib[key] = original_v_mID
                            break

                # Add v:groupContext attribute
                for key in target_group.attrib.keys():
                    if 'groupContext' in key:
                        new_g_attrib[key] = 'shape'
                        break

                # Create new <g> element (use original tag name)
                g_tag = None
                for child in root:
                    if child.tag.endswith('}g'):
                        g_tag = child.tag
                        break
                if g_tag is None:
                    g_tag = '{http://www.w3.org/2000/svg}g'

                new_g = ET.Element(g_tag, attrib=new_g_attrib)

                # Add <title> element
                title_tag = original_title_element.tag if original_title_element is not None else '{http://www.w3.org/2000/svg}title'
                new_title = ET.SubElement(new_g, title_tag)
                new_title.text = original_title_text

                # Add <desc> element (same as original structure)
                desc_tag = None
                for child in target_group:
                    if child.tag.endswith('}desc') or child.tag == 'desc':
                        desc_tag = child.tag
                        break
                if desc_tag:
                    new_desc = ET.SubElement(new_g, desc_tag)
                    new_desc.text = original_title_text

                # Add v:textBlock element
                for child in target_group:
                    if 'textBlock' in child.tag:
                        textblock_tag = child.tag
                        new_textblock = ET.SubElement(new_g, textblock_tag)
                        # Copy margins attribute if exists
                        margins = get_attr_ignore_namespace(child, 'margins')
                        if margins:
                            for key in child.attrib.keys():
                                if 'margins' in key:
                                    new_textblock.set(key, margins)
                                    break
                        break

                # Add v:textRect element
                for child in target_group:
                    if 'textRect' in child.tag:
                        textrect_tag = child.tag
                        new_textrect = ET.SubElement(new_g, textrect_tag)
                        new_textrect.set('cx', '50')
                        new_textrect.set('cy', '25')
                        new_textrect.set('width', '90')
                        new_textrect.set('height', '40')
                        break

                # Add <rect> element (with placeholder values)
                rect_tag = '{http://www.w3.org/2000/svg}rect'
                new_rect = ET.SubElement(new_g, rect_tag)
                new_rect.set('x', '0')
                new_rect.set('y', '0')
                new_rect.set('width', '100')
                new_rect.set('height', '50')
                new_rect.set('class', 'st1')

                # Add <text> element
                text_tag = '{http://www.w3.org/2000/svg}text'
                new_text = ET.SubElement(new_g, text_tag)
                new_text.set('x', '50')
                new_text.set('y', '30')
                new_text.set('class', 'st2')

                # Copy v:langID if exists
                for child in target_group:
                    if child.tag.endswith('}text') or child.tag == 'text':
                        langID = get_attr_ignore_namespace(child, 'langID')
                        if langID:
                            for key in child.attrib.keys():
                                if 'langID' in key:
                                    new_text.set(key, langID)
                                    break
                        break

                new_text.text = original_title_text

                # Add v:paragraph inside text element
                for child in target_group:
                    if child.tag.endswith('}text') or child.tag == 'text':
                        for subchild in child:
                            if 'paragraph' in subchild.tag:
                                paragraph_tag = subchild.tag
                                new_paragraph = ET.SubElement(new_text, paragraph_tag)
                                # Copy horizAlign attribute
                                horizAlign = get_attr_ignore_namespace(subchild, 'horizAlign')
                                if horizAlign:
                                    for key in subchild.attrib.keys():
                                        if 'horizAlign' in key:
                                            new_paragraph.set(key, horizAlign)
                                            break
                            elif 'tabList' in subchild.tag:
                                tablist_tag = subchild.tag
                                ET.SubElement(new_text, tablist_tag)
                        break

                # Replace original group with new group
                if parent_of_target_group is not None:
                    index = list(parent_of_target_group).index(target_group)
                    parent_of_target_group.remove(target_group)
                    parent_of_target_group.insert(index, new_g)

            # Output the transformed SVG
            ET.indent(tree, space="    ")
            transformed_svg_content = ET.tostring(root, encoding='unicode', xml_declaration=True)

        '''end of additional section'''


        # -------- Namespace-agnostic helpers --------
        def local_name(tag):
            # '{uri}name' -> 'name', or plain 'name' -> 'name'
            return tag.split('}')[-1] if tag else ''

        def get_attr(el, name, default=None):
            # Find attribute by local-name, ignoring namespace.
            for k, v in el.attrib.items():
                if k.split('}')[-1] == name:
                    return v
            return default

        def first_child(el, name):
            # Return first direct child with matching local-name.
            for c in el:
                if local_name(c.tag) == name:
                    return c
            return None

        def all_descendants(el, name):
            # Return all descendants (including el itself children) with matching local-name.
            return [e for e in el.iter() if local_name(e.tag) == name]

        # Visio layerMember check without namespace dependency.
        def in_connector_layer(el, parent_map):
            e = el
            while e is not None:
                if local_name(e.tag) == 'g':
                    if get_attr(e, 'layerMember') == '0':
                        return True
                e = parent_map.get(e)
            return False

        def pt_to_inches(pt):
            # Treat SVG coordinates as points (72 points per inch).
            return pt / 72.0

        # --- 2D transform helpers (translate/scale/matrix) ---
        def mat_identity():
            return [[1.0, 0.0, 0.0],
                    [0.0, 1.0, 0.0],
                    [0.0, 0.0, 1.0]]

        def mat_mul(A, B):
            C = [[0.0, 0.0, 0.0] for _ in range(3)]
            for i in range(3):
                for j in range(3):
                    C[i][j] = A[i][0] * B[0][j] + A[i][1] * B[1][j] + A[i][2] * B[2][j]
            return C

        def mat_translate(dx, dy):
            return [[1.0, 0.0, dx],
                    [0.0, 1.0, dy],
                    [0.0, 0.0, 1.0]]

        def mat_scale(sx, sy):
            return [[sx, 0.0, 0.0],
                    [0.0, sy, 0.0],
                    [0.0, 0.0, 1.0]]

        def mat_matrix(a, b, c, d, e, f):
            # SVG matrix(a b c d e f):
            # [a c e]
            # [b d f]
            # [0 0 1]
            return [[a, c, e],
                    [b, d, f],
                    [0.0, 0.0, 1.0]]

        # ADD THIS NEW FUNCTION
        def mat_rotate(cos_a, sin_a):
            """Create rotation matrix from cos and sin values"""
            return [[cos_a, -sin_a, 0.0],
                    [sin_a, cos_a, 0.0],
                    [0.0, 0.0, 1.0]]

        def parse_transform_to_matrix(s):
            # Supports translate, scale, rotate, and matrix. Skew is ignored for simplicity.
            if not s:
                return mat_identity()
            M = mat_identity()
            for func, args in re.findall(r'([a-zA-Z]+)\s*\(([^)]*)\)', s):
                parts = [p for p in re.split(r'[,\s]+', args.strip()) if p]
                try:
                    if func.lower() == 'translate':
                        dx = float(parts[0]) if len(parts) >= 1 else 0.0
                        dy = float(parts[1]) if len(parts) >= 2 else 0.0
                        M = mat_mul(M, mat_translate(dx, dy))
                    elif func.lower() == 'scale':
                        sx = float(parts[0]) if len(parts) >= 1 else 1.0
                        sy = float(parts[1]) if len(parts) >= 2 else sx
                        M = mat_mul(M, mat_scale(sx, sy))
                    elif func.lower() == 'rotate':
                        # rotate(angle) or rotate(angle, cx, cy)
                        angle_deg = float(parts[0]) if len(parts) >= 1 else 0.0
                        angle_rad = math.radians(angle_deg)
                        cos_a = math.cos(angle_rad)
                        sin_a = math.sin(angle_rad)

                        if len(parts) >= 3:
                            # rotate around point (cx, cy)
                            cx = float(parts[1])
                            cy = float(parts[2])
                            # Translate to origin, rotate, translate back
                            M = mat_mul(M, mat_translate(cx, cy))
                            M = mat_mul(M, mat_rotate(cos_a, sin_a))
                            M = mat_mul(M, mat_translate(-cx, -cy))
                        else:
                            # rotate around origin
                            M = mat_mul(M, mat_rotate(cos_a, sin_a))
                    elif func.lower() == 'matrix' and len(parts) >= 6:
                        a, b, c, d, e, f = map(float, parts[:6])
                        M = mat_mul(M, mat_matrix(a, b, c, d, e, f))
                except ValueError:
                    continue
            return M

        def apply_mat(M, x, y):
            x2 = M[0][0] * x + M[0][1] * y + M[0][2]
            y2 = M[1][0] * x + M[1][1] * y + M[1][2]
            return x2, y2

        # Build a parent map to accumulate transforms from ancestors.
        parent_map = {child: parent for parent in root.iter() for child in parent}

        def cumulative_matrix(el):
            # Accumulate transforms from root to the element.
            chain = []
            e = el
            while e is not None:
                chain.append(e)
                e = parent_map.get(e)
            chain.reverse()
            M = mat_identity()
            for node in chain:
                t = get_attr(node, 'transform')
                if t is None:
                    t = node.attrib.get('transform')  # non-namespaced attr fallback
                if t:
                    M = mat_mul(M, parse_transform_to_matrix(t))
            return M

        # --- Path endpoint extraction (prefer straight Line segments) ---
        def _manual_path_endpoints(d):
            if not d or not d.strip():
                return None, None
            s = re.sub(r'[,\s]+', ' ', d.strip())
            tokens = re.findall(r'([MLHVmlhv])|([-\d.]+)', s)
            current = None
            start = None
            last_cmd = None
            nums_buffer = []

            def flush(cmd, nums):
                nonlocal start, current
                if not cmd or not nums:
                    return
                if cmd in ('M', 'L'):
                    for i in range(0, len(nums) - 1, 2):
                        x = float(nums[i]); y = float(nums[i + 1])
                        current = (x, y)
                        if cmd == 'M' and start is None:
                            start = current
                    return
                if cmd in ('m', 'l'):
                    if current is None and len(nums) >= 2:
                        x = float(nums[0]); y = float(nums[1])
                        current = (x, y)
                        if cmd == 'm' and start is None:
                            start = current
                        for i in range(2, len(nums) - 1, 2):
                            dx = float(nums[i]); dy = float(nums[i + 1])
                            cx, cy = current
                            current = (cx + dx, cy + dy)
                    else:
                        for i in range(0, len(nums) - 1, 2):
                            dx = float(nums[i]); dy = float(nums[i + 1])
                            cx, cy = current
                            current = (cx + dx, cy + dy)
                    return
                if cmd == 'H':
                    for x in map(float, nums):
                        current = (x, current[1] if current else 0.0)
                    return
                if cmd == 'h':
                    for dx in map(float, nums):
                        if current is None:
                            current = (dx, 0.0)
                        else:
                            current = (current[0] + dx, current[1])
                    return
                if cmd == 'V':
                    for y in map(float, nums):
                        current = (current[0] if current else 0.0, y)
                    return
                if cmd == 'v':
                    for dy in map(float, nums):
                        if current is None:
                            current = (0.0, dy)
                        else:
                            current = (current[0], current[1] + dy)
                    return

            for cmd_token, num_token in tokens:
                if cmd_token:
                    if last_cmd and nums_buffer:
                        flush(last_cmd, nums_buffer)
                        nums_buffer = []
                    last_cmd = cmd_token
                else:
                    nums_buffer.append(num_token)
            if last_cmd and nums_buffer:
                flush(last_cmd, nums_buffer)

            if start and current and start != current:
                return start, current

            nums = [float(n) for n in re.findall(r'[-\d.]+', d)]
            if len(nums) >= 4:
                a = (nums[0], nums[1]); b = (nums[-2], nums[-1])
                if a != b:
                    return a, b
            return None, None

        def path_endpoints_prefer_lines(d):
            if not d or not d.strip():
                return None, None
            if SVG_PATH_AVAILABLE:
                try:
                    path_obj = svg_parse_path(d)
                    start = None
                    end = None
                    for seg in path_obj:
                        if isinstance(seg, SvgLine):
                            if start is None:
                                start = (seg.start.real, seg.start.imag)
                            end = (seg.end.real, seg.end.imag)
                    if start and end and start != end:
                        return start, end
                except Exception:
                    pass
            return _manual_path_endpoints(d)

        def canonical_line_key(x1, y1, x2, y2, ndigits=0):
            r = lambda v: round(v, ndigits)
            a = (r(x1), r(y1)); b = (r(x2), r(y2))
            return tuple(sorted((a, b)))

        def set_text_black(text_frame):
            p = text_frame.paragraphs[0]
            p.font.color.rgb = RGBColor(0, 0, 0)
            for run in p.runs:
                run.font.color.rgb = RGBColor(0, 0, 0)

        def is_valid_label(text):
            if not text:
                return False
            t = text.strip()
            if t == "" or t == ".local":
                return False
            if not re.search(r'[A-Za-z0-9]', t):
                return False
            return True

        def extract_scale(M):
            sx = math.sqrt(M[0][0] ** 2 + M[1][0] ** 2)
            sy = math.sqrt(M[0][1] ** 2 + M[1][1] ** 2)
            if not math.isfinite(sx) or sx == 0.0:
                sx = 1.0
            if not math.isfinite(sy) or sy == 0.0:
                sy = 1.0
            return sx, sy

        # Create presentation with page size matching a Letter-style slide.
        prs = Presentation()
        prs.slide_width = Inches(8.5)
        prs.slide_height = Inches(11)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout.

        # Initialize existing_names set before any shape processing
        existing_names = set()

        '''foreignObject support - convert to shapes'''
        # Find all <g> elements containing <switch><foreignObject> structure
        processed_foreign = set()

        for g in root.iter():
            if not (g.tag.endswith('}g') or g.tag == 'g'):
                continue

            # Look for <switch> in descendants (not just direct children)
            switch_element = None
            for desc in g.iter():
                if desc.tag.endswith('}switch') or desc.tag == 'switch':
                    switch_element = desc
                    break

            if switch_element is None:
                continue

            # Avoid processing the same switch multiple times
            if id(switch_element) in processed_foreign:
                continue
            processed_foreign.add(id(switch_element))

            # Look for <foreignObject> within <switch>
            foreign_object = None
            for child in switch_element:
                if child.tag.endswith('}foreignObject') or child.tag == 'foreignObject':
                    foreign_object = child
                    break

            if foreign_object is None:
                continue

            # Extract text from nested <div> elements
            text_content = None
            for div in foreign_object.iter():
                if div.tag.endswith('}div') or div.tag == 'div':
                    if div.text and div.text.strip():
                        text_content = div.text.strip()
                        break

            if not text_content or not is_valid_label(text_content):
                print(f"  Skipped: invalid label")
                continue

            # Find the root <g> that contains both the shape and the foreignObject
            # This is the <g> with data-cell-id attribute
            root_g = g
            temp = g
            while temp is not None:
                if get_attr(temp, 'data-cell-id') is not None:
                    root_g = temp
                    break
                temp = parent_map.get(temp)

            #print(f"  Root <g> id: {get_attr(root_g, 'data-cell-id')}")

            # Find <rect> with fill="none" within root_g (this is the bounding box)
            bbox_element = None
            for desc in root_g.iter():
                if local_name(desc.tag) == 'rect':
                    # Look for rect with fill="none" - this is the bounding box
                    fill_attr = desc.attrib.get('fill', '')
                    if fill_attr == 'none':
                        bbox_element = desc
                        #print(f"  Found bounding rect with fill='none'")
                        break

            if bbox_element is None:
                #print(f"  No bounding <rect> found")
                continue

            # Extract bounding box from rect
            try:
                bbox_x = float(bbox_element.attrib.get('x', '0'))
                bbox_y = float(bbox_element.attrib.get('y', '0'))
                bbox_w = float(bbox_element.attrib.get('width', '0'))
                bbox_h = float(bbox_element.attrib.get('height', '0'))
            except (ValueError, TypeError):
                print(f"  Failed to parse rect attributes")
                continue

            if not bbox_w or not bbox_h:
                print(f"  Invalid dimensions: w={bbox_w}, h={bbox_h}")
                continue

            #print(f"  Bounding box: x={bbox_x}, y={bbox_y}, w={bbox_w}, h={bbox_h}")

            # Get cumulative transform for the rect element
            M = cumulative_matrix(bbox_element)
            absx, absy = apply_mat(M, bbox_x, bbox_y)
            sx, sy = extract_scale(M)
            eff_w = bbox_w * sx
            eff_h = bbox_h * sy

            # Apply translate(0.5, 0.5) offset to match lines
            # Lines have <g transform="translate(0.5,0.5)"> wrapper
            absx += 0.5
            absy += 0.5

            # Skip if already processed
            if text_content in existing_names:
                print(f"  Skipped: already processed")
                continue

            # Create shape in PowerPoint
            left = Inches(pt_to_inches(absx))
            top = Inches(pt_to_inches(absy))
            width = Inches(pt_to_inches(eff_w))
            height = Inches(pt_to_inches(eff_h))

            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(0, 0, 0)

            tf = shape.text_frame
            tf.clear()
            tf.paragraphs[0].text = text_content
            set_text_black(tf)

            shape.name = text_content
            existing_names.add(text_content)

        '''end of foreignObject support'''

        # Find groups matching Visio 'groupContext' == 'shape' without namespace dependency.
        groups = [g for g in root.iter() if local_name(g.tag) == 'g' and get_attr(g, 'groupContext') == 'shape']
        ''' end of additional part 2 '''

        # Create presentation with page size matching a Letter-style slide.
        prs = Presentation()
        prs.slide_width = Inches(8.5)
        prs.slide_height = Inches(11)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout.

        # Initialize existing_names set before any shape processing
        existing_names = set()
        processed_texts = set()  # Initialize early for use in groups loop
        processed_images = set()  # Initialize early for use in groups loop

        '''foreignObject support - convert to shapes'''
        # Find all <g> elements containing <switch><foreignObject> structure
        processed_foreign = set()

        for g in root.iter():
            if not (g.tag.endswith('}g') or g.tag == 'g'):
                continue

            # Look for <switch> in descendants (not just direct children)
            switch_element = None
            for desc in g.iter():
                if desc.tag.endswith('}switch') or desc.tag == 'switch':
                    switch_element = desc
                    break

            if switch_element is None:
                continue

            # Avoid processing the same switch multiple times
            if id(switch_element) in processed_foreign:
                continue
            processed_foreign.add(id(switch_element))

            # Look for <foreignObject> within <switch>
            foreign_object = None
            for child in switch_element:
                if child.tag.endswith('}foreignObject') or child.tag == 'foreignObject':
                    foreign_object = child
                    break

            if foreign_object is None:
                continue

            # Extract text from nested <div> elements
            text_content = None
            for div in foreign_object.iter():
                if div.tag.endswith('}div') or div.tag == 'div':
                    if div.text and div.text.strip():
                        text_content = div.text.strip()
                        break

            if not text_content or not is_valid_label(text_content):
                print(f"  Skipped: invalid label")
                continue

            # Find the root <g> that contains both the shape and the foreignObject
            # This is the <g> with data-cell-id attribute
            root_g = g
            temp = g
            while temp is not None:
                if get_attr(temp, 'data-cell-id') is not None:
                    root_g = temp
                    break
                temp = parent_map.get(temp)

            # Find <rect> with fill="none" within root_g (this is the bounding box)
            bbox_element = None
            for desc in root_g.iter():
                if local_name(desc.tag) == 'rect':
                    # Look for rect with fill="none" - this is the bounding box
                    fill_attr = desc.attrib.get('fill', '')
                    if fill_attr == 'none':
                        bbox_element = desc
                        break

            if bbox_element is None:
                #print(f"  No bounding <rect> found")
                continue

            # Extract bounding box from rect
            try:
                bbox_x = float(bbox_element.attrib.get('x', '0'))
                bbox_y = float(bbox_element.attrib.get('y', '0'))
                bbox_w = float(bbox_element.attrib.get('width', '0'))
                bbox_h = float(bbox_element.attrib.get('height', '0'))
            except (ValueError, TypeError):
                print(f"  Failed to parse rect attributes")
                continue

            if not bbox_w or not bbox_h:
                print(f"  Invalid dimensions: w={bbox_w}, h={bbox_h}")
                continue

            # Get cumulative transform for the rect element
            M = cumulative_matrix(bbox_element)
            absx, absy = apply_mat(M, bbox_x, bbox_y)
            sx, sy = extract_scale(M)
            eff_w = bbox_w * sx
            eff_h = bbox_h * sy

            # Apply translate(0.5, 0.5) offset to match lines
            # Lines have <g transform="translate(0.5,0.5)"> wrapper
            absx += 0.5
            absy += 0.5

            # Skip if already processed
            if text_content in existing_names:
                print(f"  Skipped: already processed")
                continue

            # Create shape in PowerPoint
            left = Inches(pt_to_inches(absx))
            top = Inches(pt_to_inches(absy))
            width = Inches(pt_to_inches(eff_w))
            height = Inches(pt_to_inches(eff_h))

            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(0, 0, 0)

            tf = shape.text_frame
            tf.clear()
            tf.paragraphs[0].text = text_content
            set_text_black(tf)

            shape.name = text_content
            existing_names.add(text_content)

        '''end of foreignObject support'''

        # Find groups matching Visio 'groupContext' == 'shape' without namespace dependency.
        groups = [g for g in root.iter() if local_name(g.tag) == 'g' and get_attr(g, 'groupContext') == 'shape']

        for g in groups:
            rect = first_child(g, 'rect')
            if rect is None:
                continue

            # ============== STROKE CHECK ==============
            # Check if rect has visible stroke
            stroke_attr = rect.attrib.get('stroke')
            class_attr = rect.attrib.get('class', '')

            # Skip if stroke="none"
            if stroke_attr == 'none':
                print(f"  → SKIP: stroke=none")
                # Mark text as processed
                text_el = first_child(g, 'text')
                if text_el is not None:
                    processed_texts.add(text_el)
                continue

            # Skip if class has stroke:none
            if class_attr:
                print(f"  → Checking CSS for class: {class_attr}")
                style_elements = root.findall('.//{http://www.w3.org/2000/svg}style')
                print(f"  → Found {len(style_elements)} style elements")
                skip_shape = False
                for style_elem in style_elements:
                    if style_elem.text:
                        style_text = style_elem.text
                        print(f"  → Style text length: {len(style_text)}")
                        class_pattern = rf'\.{re.escape(class_attr)}\s*\{{([^}}]+)\}}'
                        match = re.search(class_pattern, style_text)
                        if match:
                            style_content = match.group(1)
                            print(f"  → Found CSS for .{class_attr}: {style_content[:100]}")
                            if 'stroke:none' in style_content or 'stroke: none' in style_content:
                                print(f"  → SKIP: class {class_attr} has stroke:none")

                                # Mark text as processed to prevent standalone rendering
                                text_el = first_child(g, 'text')
                                if text_el is not None:
                                    #print(f"  → Text element id: {id(text_el)}, text: {text_el.text}")  # ADD DEBUG
                                    processed_texts.add(text_el)
                                    print(f"  → Added text element to processed_texts")

                                # ADD: Also add all text descendants
                                for text_desc in g.iter():
                                    if local_name(text_desc.tag) == 'text':
                                        #print(f"  → Also adding text descendant id: {id(text_desc)}")  # ADD DEBUG
                                        processed_texts.add(text_desc)

                                skip_shape = True
                                break
                        else:
                            print(f"  → No CSS match for .{class_attr}")
                if skip_shape:
                    continue
            # ============== END OF STROKE CHECK ==============

            name = None
            desc_el = first_child(g, 'desc')
            if desc_el is not None and desc_el.text:
                name = desc_el.text.strip()
            if not name:
                text_el = first_child(g, 'text')
                if text_el is not None:
                    text_content = "".join(text_el.itertext()).strip()
                    if text_content:
                        name = text_content

            print(f"  → Shape name: {name}")

            if not is_valid_label(name):
                print(f"  → SKIP: invalid label")
                continue

            tx_ty = get_attr(g, 'transform')
            if tx_ty is None:
                tx_ty = g.attrib.get('transform')
            M_g = parse_transform_to_matrix(tx_ty) if tx_ty else mat_identity()

            x = float(rect.attrib.get('x', '0'))
            y = float(rect.attrib.get('y', '0'))

            x, y = apply_mat(M_g, x, y)
            w = float(rect.attrib.get('width', '0'))
            h = float(rect.attrib.get('height', '0'))

            left = Inches(pt_to_inches(x))
            top = Inches(pt_to_inches(y))
            width = Inches(pt_to_inches(w))
            height = Inches(pt_to_inches(h))

            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(0, 0, 0)

            tf = shape.text_frame
            tf.clear()
            tf.paragraphs[0].text = name
            set_text_black(tf)

            shape.name = name
            existing_names.add(name)

        # ... (rest of the code continues with standalone <rect> processing, etc.)

        # ============================================================
        # NEW SECTION: Process draw.io shapes from data-cell-id groups
        # Handles ellipse (WAN) and complex paths (Router-line2)
        # ============================================================

        # Find all <g> elements with data-cell-id (draw.io format)
        drawio_cell_groups = []
        for g in root.iter():
            if local_name(g.tag) == 'g':
                cell_id = get_attr(g, 'data-cell-id')
                if cell_id and cell_id not in ('0', '1'):
                    # Skip if already processed via groupContext='shape'
                    if get_attr(g, 'groupContext') == 'shape':
                        continue
                    drawio_cell_groups.append((cell_id, g))

        if drawio_cell_groups:
            print(f"\n Processing {len(drawio_cell_groups)} draw.io cell groups...")

        for cell_id, cell_group in drawio_cell_groups:
            # Extract label from foreignObject
            label = None
            for elem in cell_group.iter():
                if local_name(elem.tag) == 'foreignObject':
                    for div in elem.iter():
                        if div.tag.endswith('}div') or div.tag == 'div':
                            if div.text and div.text.strip():
                                label = div.text.strip()
                                break
                    if label:
                        break

            if not label or not is_valid_label(label):
                continue

            if label in existing_names:
                continue

            # Find bounding box
            bbox_x = bbox_y = bbox_w = bbox_h = None
            shape_type = MSO_AUTO_SHAPE_TYPE.RECTANGLE
            transform_elem = None

            # Check for ellipse
            for elem in cell_group.iter():
                if local_name(elem.tag) == 'ellipse':
                    try:
                        cx = float(elem.attrib.get('cx', 0))
                        cy = float(elem.attrib.get('cy', 0))
                        rx = float(elem.attrib.get('rx', 0))
                        ry = float(elem.attrib.get('ry', 0))
                        bbox_x = cx - rx
                        bbox_y = cy - ry
                        bbox_w = rx * 2
                        bbox_h = ry * 2
                        shape_type = MSO_AUTO_SHAPE_TYPE.OVAL
                        transform_elem = elem
                        break
                    except:
                        pass

            # Check for path (complex shapes like Router-line2)
            if bbox_x is None:
                all_x = []
                all_y = []

                for elem in cell_group.iter():
                    if local_name(elem.tag) == 'path':
                        d = elem.attrib.get('d', '')
                        if not d:
                            continue

                        stroke = elem.attrib.get('stroke', '')
                        fill = elem.attrib.get('fill', '')

                        # Accept if has stroke or fill
                        if (stroke and stroke != 'none') or (fill and fill != 'none'):
                            nums = re.findall(r'[-\d.]+', d)
                            try:
                                coords = [float(n) for n in nums]
                                if len(coords) >= 2:
                                    xs = coords[0::2]
                                    ys = coords[1::2]
                                    all_x.extend(xs)
                                    all_y.extend(ys)
                            except:
                                pass

                if all_x and all_y:
                    bbox_x = min(all_x)
                    bbox_y = min(all_y)
                    bbox_w = max(all_x) - min(all_x)
                    bbox_h = max(all_y) - min(all_y)
                    transform_elem = cell_group

            # Skip if no geometry
            if bbox_x is None or bbox_w is None or bbox_h is None:
                continue

            if bbox_w <= 0 or bbox_h <= 0:
                continue

            # Apply transforms
            M = mat_identity()
            if transform_elem is not None:
                M = cumulative_matrix(transform_elem)

            absx, absy = apply_mat(M, bbox_x, bbox_y)
            sx, sy = extract_scale(M)
            eff_w = bbox_w * sx
            eff_h = bbox_h * sy

            # Create shape
            left = Inches(pt_to_inches(absx))
            top = Inches(pt_to_inches(absy))
            width = Inches(pt_to_inches(eff_w))
            height = Inches(pt_to_inches(eff_h))

            try:
                shape = slide.shapes.add_shape(shape_type, left, top, width, height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                shape.line.color.rgb = RGBColor(0, 0, 0)

                tf = shape.text_frame
                tf.clear()
                tf.paragraphs[0].text = label
                tf.paragraphs[0].alignment = 1  # Center
                set_text_black(tf)

                shape.name = label
                existing_names.add(label)

                print(f"    draw.io shape: {label}")

            except Exception as e:
                print(f"    Failed: {label} - {e}")

        # ============================================================
        # END OF NEW SECTION
        # ============================================================

        # ============================================================
        # NEW SECTION: Process standalone <rect> elements (like Router-line1)
        # This section handles rectangles that are not part of the groupContext='shape' groups
        # ============================================================

        # Helper function to check if rect has visible stroke
        def has_visible_stroke(rect_elem):
            """Check if rect has a visible stroke (not 'none')"""
            # Check direct stroke attribute
            stroke_attr = rect_elem.attrib.get('stroke')
            if stroke_attr == 'none':
                return False

            # Check class attribute for stroke:none
            class_attr = rect_elem.attrib.get('class', '')
            if class_attr:
                # Parse style definitions from <style> element
                style_elements = root.findall('.//{http://www.w3.org/2000/svg}style')
                for style_elem in style_elements:
                    if style_elem.text:
                        style_text = style_elem.text
                        # Look for class definition
                        class_pattern = rf'\.{re.escape(class_attr)}\s*{{([^}}]+)}}'
                        match = re.search(class_pattern, style_text)
                        if match:
                            style_content = match.group(1)
                            # Check if stroke:none is defined
                            if 'stroke:none' in style_content or 'stroke: none' in style_content:
                                return False

            return True

        # Get all <rect> elements that have a visible stroke
        standalone_rects = []
        for rect_elem in root.iter():
            if not (rect_elem.tag.endswith('}rect') or rect_elem.tag == 'rect'):
                continue

            # Check if rect has visible stroke (not stroke="none" or class with stroke:none)
            if not has_visible_stroke(rect_elem):
                #print(f"DEBUG: Skipping rect with no visible stroke (stroke={rect_elem.attrib.get('stroke', 'N/A')})")

                # ADD: Mark associated text elements as processed
                current = rect_elem
                cell_group = None
                while current is not None:
                    if local_name(current.tag) == 'g':
                        cell_id = get_attr(current, 'data-cell-id')
                        if cell_id:
                            cell_group = current
                            break
                    current = parent_map.get(current)

                if cell_group:
                    # Add all text elements in this group to processed_texts
                    for text_elem in cell_group.iter():
                        if local_name(text_elem.tag) == 'text':
                            processed_texts.add(text_elem)
                            print(f"  → Added text element id: {id(text_elem)} to processed_texts")

                continue

            # Check fill attribute to exclude bounding boxes
            fill_attr = rect_elem.attrib.get('fill')
            if fill_attr == 'none':
                continue

            # Check if this rect is already processed (part of groupContext='shape')
            is_already_processed = False
            temp = rect_elem
            while temp is not None:
                if local_name(temp.tag) == 'g' and get_attr(temp, 'groupContext') == 'shape':
                    is_already_processed = True
                    break
                temp = parent_map.get(temp)

            if is_already_processed:
                continue

            standalone_rects.append(rect_elem)

        #print(f"DEBUG: Found {len(standalone_rects)} standalone rects with visible strokes")

        # Process each standalone rect
        for rect_elem in standalone_rects:
            # Get rectangle attributes
            try:
                x = float(rect_elem.attrib.get('x', '0'))
                y = float(rect_elem.attrib.get('y', '0'))
                w = float(rect_elem.attrib.get('width', '0'))
                h = float(rect_elem.attrib.get('height', '0'))
            except (ValueError, TypeError):
                continue

            if w <= 0 or h <= 0:
                continue

            # Get cumulative transform
            M = cumulative_matrix(rect_elem)
            absx, absy = apply_mat(M, x, y)
            sx, sy = extract_scale(M)
            eff_w = w * sx
            eff_h = h * sy

            # Find associated text label - IMPROVED STRATEGY
            label = None

            # Navigate up to find the containing <g> with data-cell-id or v:mID
            current = rect_elem
            cell_group = None
            while current is not None:
                if local_name(current.tag) == 'g':
                    cell_id = get_attr(current, 'data-cell-id')
                    v_mID = get_attr(current, 'mID')
                    if cell_id or v_mID:
                        cell_group = current
                        break
                current = parent_map.get(current)

            if cell_group:
                cell_id = get_attr(cell_group, 'data-cell-id') or get_attr(cell_group, 'mID')
                #print(f"DEBUG: Found cell group with id: {cell_id}")

                # Search for <text> elements (Visio style)
                for child in cell_group.iter():
                    if local_name(child.tag) == 'text':
                        text_content = "".join(child.itertext()).strip()
                        if is_valid_label(text_content):
                            label = text_content
                            #print(f"DEBUG: Found label in text: {label}")
                            processed_texts.add(child)
                            break

                # Search for <desc> element (Visio style)
                if not label:
                    for child in cell_group:
                        if local_name(child.tag) == 'desc':
                            if child.text and child.text.strip():
                                text_content = child.text.strip()
                                if is_valid_label(text_content):
                                    label = text_content
                                    #print(f"DEBUG: Found label in desc: {label}")
                                    break

                # Search for <switch><foreignObject><div> within the cell group (draw.io style)
                if not label:
                    for child in cell_group.iter():
                        if local_name(child.tag) == 'switch':
                            for switch_child in child:
                                if local_name(switch_child.tag) == 'foreignObject':
                                    # Extract text from divs
                                    for div in switch_child.iter():
                                        if div.tag.endswith('}div') or div.tag == 'div':
                                            if div.text and div.text.strip():
                                                text_content = div.text.strip()
                                                if is_valid_label(text_content):
                                                    label = text_content
                                                    #print(f"DEBUG: Found label in foreignObject: {label}")
                                                    break
                                    if label:
                                        break
                            if label:
                                break

            #print(f"DEBUG: Rect at ({x}, {y}) size ({w}, {h}) with label: {label}")

            # Skip if no valid label found or already processed
            if not label or not is_valid_label(label):
                #print(f"DEBUG: Skipping - invalid label")
                continue

            if label in existing_names:
                #print(f"DEBUG: Skipping - already exists: {label}")
                continue

            # Create shape in PowerPoint
            left = Inches(pt_to_inches(absx))
            top = Inches(pt_to_inches(absy))
            width = Inches(pt_to_inches(eff_w))
            height = Inches(pt_to_inches(eff_h))

            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)

            # Set fill color (default white, or parse from rect's fill attribute)
            fill_attr = rect_elem.attrib.get('fill')
            if fill_attr and fill_attr != 'none':
                try:
                    if fill_attr.startswith('#') and len(fill_attr) in (4, 7):
                        if len(fill_attr) == 7:
                            r = int(fill_attr[1:3], 16)
                            g = int(fill_attr[3:5], 16)
                            b = int(fill_attr[5:7], 16)
                        else:
                            r = int(fill_attr[1] * 2, 16)
                            g = int(fill_attr[2] * 2, 16)
                            b = int(fill_attr[3] * 2, 16)
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(r, g, b)
                    else:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                except Exception:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 255, 255)

            # Set stroke color
            stroke_attr = rect_elem.attrib.get('stroke')
            if stroke_attr and stroke_attr != 'none':
                try:
                    if stroke_attr.startswith('#') and len(stroke_attr) in (4, 7):
                        if len(stroke_attr) == 7:
                            r = int(stroke_attr[1:3], 16)
                            g = int(stroke_attr[3:5], 16)
                            b = int(stroke_attr[5:7], 16)
                        else:
                            r = int(stroke_attr[1] * 2, 16)
                            g = int(stroke_attr[2] * 2, 16)
                            b = int(stroke_attr[3] * 2, 16)
                        shape.line.color.rgb = RGBColor(r, g, b)
                    else:
                        shape.line.color.rgb = RGBColor(0, 0, 0)
                except Exception:
                    shape.line.color.rgb = RGBColor(0, 0, 0)
            else:
                # No stroke attribute means default black stroke
                shape.line.color.rgb = RGBColor(0, 0, 0)

            # Add text
            tf = shape.text_frame
            tf.clear()
            tf.paragraphs[0].text = label
            set_text_black(tf)

            shape.name = label
            existing_names.add(label)
            #print(f"DEBUG: Created shape: {label}")

        # ============================================================
        # END OF NEW SECTION
        # ============================================================

        # Consolidate group images + texts into a single "big square with text" shape.

        for g in [e for e in root.iter() if local_name(e.tag) == 'g']:
            images = [img for img in g if local_name(img.tag) == 'image']
            texts = [t for t in g if local_name(t.tag) == 'text']

            if not images or not texts:
                continue

            label = None
            for t in texts:
                content = "".join(t.itertext()).strip()
                if is_valid_label(content):
                    label = content
                    break
            if label is None:
                continue

            best = None  # (absx, absy, eff_w, eff_h, area, element)
            for img in images:
                w_attr = img.attrib.get('width')
                h_attr = img.attrib.get('height')
                x_attr = img.attrib.get('x')
                y_attr = img.attrib.get('y')

                try:
                    w = float(w_attr) if w_attr is not None else None
                    h = float(h_attr) if h_attr is not None else None
                    x = float(x_attr) if x_attr is not None else 0.0
                    y = float(y_attr) if y_attr is not None else 0.0
                except ValueError:
                    w = h = None
                    x = y = 0.0

                if w is None or h is None or w <= 0 or h <= 0:
                    continue

                M = cumulative_matrix(img)
                absx, absy = apply_mat(M, x, y)
                sx, sy = extract_scale(M)
                eff_w = w * sx
                eff_h = h * sy
                area = eff_w * eff_h

                if not best or area > best[4]:
                    best = (absx, absy, eff_w, eff_h, area, img)

            if not best:
                continue

            absx, absy, eff_w, eff_h, _, big_img_el = best

            left = Inches(pt_to_inches(absx))
            top = Inches(pt_to_inches(absy))
            width = Inches(pt_to_inches(eff_w))
            height = Inches(pt_to_inches(eff_h))

            uni = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
            uni.fill.solid()
            uni.fill.fore_color.rgb = RGBColor(255, 255, 255)
            uni.line.color.rgb = RGBColor(0, 0, 0)

            tf = uni.text_frame
            tf.clear()
            tf.paragraphs[0].text = label
            set_text_black(tf)

            uni.name = label
            existing_names.add(label)

            for img in images:
                processed_images.add(img)
            for t in texts:
                processed_texts.add(t)

        # Draw single straight connectors for ALL <path> elements by consolidating endpoints.
        # Filter out lines that are completely inside shapes
        seen = set()
        connector_count = 0
        filtered_count = 0

        # Build a list of all created shapes with their bounding boxes
        shape_bounds = []
        for shape_name in existing_names:
            # Find the shape by name
            for shape in slide.shapes:
                if hasattr(shape, 'name') and shape.name == shape_name:
                    shape_bounds.append({
                        'name': shape_name,
                        'left': shape.left,
                        'top': shape.top,
                        'right': shape.left + shape.width,
                        'bottom': shape.top + shape.height
                    })
                    break

        def point_inside_shape(x, y, shape_bound):
            """Check if point (x, y) is inside shape bounds"""
            # Convert point coordinates to EMU (914400 EMU = 1 inch)
            x_emu = x / 72.0 * 914400
            y_emu = y / 72.0 * 914400

            return (shape_bound['left'] <= x_emu <= shape_bound['right'] and
                    shape_bound['top'] <= y_emu <= shape_bound['bottom'])

        def both_endpoints_inside_any_shape(x1, y1, x2, y2, shape_bounds_list):
            """Check if both endpoints of a line (shortened by 10%) are inside the same shape at ver 2.6.0f """
    
            # Calculate line vector
            dx = x2 - x1
            dy = y2 - y1
    
            # Shorten line by 10% (5% from each end)
            # Move start point 5% toward end point
            shortened_x1 = x1 + dx * 0.05
            shortened_y1 = y1 + dy * 0.05
    
            # Move end point 5% toward start point
            shortened_x2 = x2 - dx * 0.05
            shortened_y2 = y2 - dy * 0.05
    
            for shape_bound in shape_bounds_list:
                if (point_inside_shape(shortened_x1, shortened_y1, shape_bound) and
                        point_inside_shape(shortened_x2, shortened_y2, shape_bound)):
                    return True, shape_bound['name']
            return False, None

        for p in [e for e in root.iter() if local_name(e.tag) == 'path']:
            d = p.attrib.get('d', '')

            # CHECK: Skip all closed paths (these are shapes, not connectors)
            has_closure = bool(re.search(r'[Zz]', d))

            if has_closure:
                # Closed paths are shapes (icons, decorations), not connectors
                continue

            start, end = path_endpoints_prefer_lines(d)
            if not start or not end:
                continue

            M = cumulative_matrix(p)
            x1, y1 = apply_mat(M, start[0], start[1])
            x2, y2 = apply_mat(M, end[0], end[1])

            if not all(map(math.isfinite, [x1, y1, x2, y2])):
                continue
            if (x1, y1) == (x2, y2):
                continue

            # NEW: Check if both endpoints are inside the same shape
            is_inside, shape_name = both_endpoints_inside_any_shape(x1, y1, x2, y2, shape_bounds)
            if is_inside:
                filtered_count += 1
                print(f"  -> Filtered line inside {shape_name}: ({x1:.0f},{y1:.0f}) to ({x2:.0f},{y2:.0f})")
                continue

            key = canonical_line_key(x1, y1, x2, y2)
            if key in seen:
                continue
            seen.add(key)

            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(pt_to_inches(x1)),
                Inches(pt_to_inches(y1)),
                Inches(pt_to_inches(x2)),
                Inches(pt_to_inches(y2)),
            )

            stroke = p.attrib.get('stroke')
            if stroke:
                try:
                    if stroke.startswith('#') and len(stroke) in (4, 7):
                        if len(stroke) == 7:
                            r = int(stroke[1:3], 16)
                            g = int(stroke[3:5], 16)
                            b = int(stroke[5:7], 16)
                        else:
                            r = int(stroke[1] * 2, 16)
                            g = int(stroke[2] * 2, 16)
                            b = int(stroke[3] * 2, 16)
                        conn.line.color.rgb = RGBColor(r, g, b)
                    elif stroke.startswith('rgb'):
                        nums = [int(v) for v in re.findall(r'\d+', stroke)[:3]]
                        if len(nums) == 3:
                            conn.line.color.rgb = RGBColor(nums[0], nums[1], nums[2])
                        else:
                            conn.line.color.rgb = RGBColor(0, 0, 0)
                    else:
                        conn.line.color.rgb = RGBColor(0, 0, 0)
                except Exception:
                    conn.line.color.rgb = RGBColor(0, 0, 0)
            else:
                conn.line.color.rgb = RGBColor(0, 0, 0)

            sw = p.attrib.get('stroke-width')
            if sw:
                try:
                    conn.line.width = Pt(float(sw))
                except Exception:
                    pass

            connector_count += 1

        print(f"\n--- Connector Summary ---")
        print(f"Total connectors created: {connector_count}")
        print(f"Lines filtered (inside shapes): {filtered_count}")

        # Draw shapes from standalone <text> elements (labels) not consolidated in groups.
        for t in [e for e in root.iter() if local_name(e.tag) == 'text']:
            if t in processed_texts:
                continue

            content = "".join(t.itertext()).strip()


            if not is_valid_label(content):
                continue
            if content in existing_names:
                continue

            x_attr = t.attrib.get('x')
            y_attr = t.attrib.get('y')
            try:
                x = float(x_attr) if x_attr is not None else None
                y = float(y_attr) if y_attr is not None else None
            except ValueError:
                x = y = None

            if x is None or y is None or not all(map(math.isfinite, [x, y])):
                continue

            M = cumulative_matrix(t)
            absx, absy = apply_mat(M, x, y)

            w_pt = 130.5
            h_pt = 40.0

            left = Inches(pt_to_inches(absx))
            top = Inches(pt_to_inches(absy))
            width = Inches(pt_to_inches(w_pt))
            height = Inches(pt_to_inches(h_pt))

            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(0, 0, 0)

            tf = shape.text_frame
            tf.clear()
            tf.paragraphs[0].text = content
            set_text_black(tf)

            shape.name = content
            existing_names.add(content)

        out_path = self.tmp_pptx_file_path
        prs.save(out_path)



        '''
        Handles the removal of duplicate Shapes
        '''
        from pptx import Presentation
        from collections import defaultdict
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        def has_text(shape):
            """
            Function to check if a Shape has text

            Args:
                shape: PowerPoint Shape

            Returns:
                bool: True if shape has text
            """
            try:
                # Check if shape has text_frame attribute
                if hasattr(shape, 'text_frame'):
                    # Check if text exists and is not empty
                    if shape.text_frame.text and shape.text_frame.text.strip():
                        return True
                # Check if shape has text attribute
                if hasattr(shape, 'text'):
                    if shape.text and shape.text.strip():
                        return True
            except:
                pass
            return False

        def is_valid_shape_for_processing(shape):
            """
            Function to check if a Shape is valid for processing
            Excludes Lines and Connectors

            Args:
                shape: PowerPoint Shape

            Returns:
                bool: True if shape should be processed
            """
            try:
                # Exclude Lines and Connectors
                if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    return False
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and hasattr(shape, 'connector_format'):
                    return False

                # Only process shapes with text
                return has_text(shape)
            except:
                return False

        def shapes_overlap(shape1, shape2, tolerance=100000):
            """
            Function to determine if two Shapes overlap or are very close

            Args:
                shape1, shape2: Two shapes to compare
                tolerance: Tolerance value for overlap detection (in EMU units)
                           Default 100000 EMU ≈ 0.35 cm

            Returns:
                bool: True if shapes overlap or are within tolerance distance
            """
            # Get coordinates and sizes
            left1, top1 = shape1.left, shape1.top
            right1, bottom1 = left1 + shape1.width, top1 + shape1.height

            left2, top2 = shape2.left, shape2.top
            right2, bottom2 = left2 + shape2.width, top2 + shape2.height

            # Check for overlap with tolerance
            # Shapes overlap if they are within tolerance distance
            horizontal_overlap = not (right1 < left2 - tolerance or right2 < left1 - tolerance)
            vertical_overlap = not (bottom1 < top2 - tolerance or bottom2 < top1 - tolerance)

            overlap = horizontal_overlap and vertical_overlap

            return overlap

        def get_shape_volume(shape):
            """
            Function to calculate the volume (area) of a Shape

            Args:
                shape: PowerPoint Shape

            Returns:
                int: Area (width * height)
            """
            return shape.width * shape.height

        def find_overlapping_groups(shapes, tolerance=100000):
            """
            Function to detect groups of overlapping Shapes
            Only processes shapes with text, excludes Lines and Connectors
            Uses Union-Find algorithm for better grouping

            Args:
                shapes: List of Shapes
                tolerance: Tolerance value for overlap detection (in EMU units)

            Returns:
                list: List of groups containing overlapping Shapes
            """
            # Filter shapes: only include shapes with text, exclude Lines and Connectors
            valid_shapes = [shape for shape in shapes if is_valid_shape_for_processing(shape)]

            if len(valid_shapes) == 0:
                return []

            # Create overlap matrix
            n = len(valid_shapes)
            parent = list(range(n))

            def find(x):
                if parent[x] != x:
                    parent[x] = find(parent[x])
                return parent[x]

            def union(x, y):
                px, py = find(x), find(y)
                if px != py:
                    parent[px] = py

            # Check all pairs for overlap
            for i in range(n):
                for j in range(i + 1, n):
                    if shapes_overlap(valid_shapes[i], valid_shapes[j], tolerance):
                        union(i, j)

            # Group shapes by parent
            groups_dict = defaultdict(list)
            for i in range(n):
                groups_dict[find(i)].append(valid_shapes[i])

            # Return only groups with more than one shape
            groups = [group for group in groups_dict.values() if len(group) > 1]

            return groups

        def remove_overlapping_shapes(prs, tolerance=100000, area_threshold_percent=50):
            """
            Process PowerPoint presentation and delete all but the largest overlapping Shape
            Only processes shapes with text, excludes Lines and Connectors
            Only deletes shapes whose area is less than the specified percentage of the largest shape

            Args:
                prs: Presentation object
                tolerance: Tolerance value for overlap detection (in EMU units)
                           Default 100000 EMU ≈ 0.35 cm
                area_threshold_percent: Percentage threshold for deletion
                                        Only shapes with area < (largest_area * threshold / 100) will be deleted
                                        Default 50 (50%)

            Returns:
                Presentation: Modified presentation object
            """
            # Process each slide
            for slide_idx, slide in enumerate(prs.slides):
                print(f"\n=== Processing Slide {slide_idx + 1} ===")
                print(f"Tolerance setting: {tolerance} EMU ({tolerance / 914400:.2f} inches)")
                print(f"Area threshold: {area_threshold_percent}% of largest shape")

                # Get all Shapes in the slide
                shapes = list(slide.shapes)

                # Count valid shapes (with text, excluding Lines and Connectors)
                valid_shape_count = sum(1 for shape in shapes if is_valid_shape_for_processing(shape))
                print(f"Total shapes: {len(shapes)}, Valid shapes with text: {valid_shape_count}")

                # Detect groups of overlapping Shapes
                overlapping_groups = find_overlapping_groups(shapes, tolerance)

                print(f"Number of overlapping groups detected: {len(overlapping_groups)}")

                # Process each group
                for group_idx, group in enumerate(overlapping_groups):
                    #print(f"\n--- Group {group_idx + 1} (Total: {len(group)} shapes) ---")

                    # Display information and area for each Shape
                    shape_info = []
                    for shape in group:
                        volume = get_shape_volume(shape)
                        name = shape.name if hasattr(shape, 'name') else "Unknown"
                        text = shape.text if hasattr(shape, 'text') else ""
                        shape_type = shape.shape_type
                        shape_info.append((shape, volume, name))
                        #print(f"  Shape: {name}")
                        #print(f"    Text: '{text}'")
                        #print(f"    Area: {volume:,} EMU²")
                        #print(f"    Type: {shape_type}")
                        #print(f"    Position: ({shape.left}, {shape.top})")
                        #print(f"    Size: {shape.width} x {shape.height}")

                    # Identify the Shape with maximum area
                    max_shape, max_volume, max_name = max(shape_info, key=lambda x: x[1])
                    threshold_area = max_volume * area_threshold_percent / 100

                    print(f"\n  → Keep: {max_name} (Area: {max_volume:,} EMU²)")
                    print(f"  → Deletion threshold: {threshold_area:,} EMU² ({area_threshold_percent}% of largest)")

                    # Delete shapes that meet BOTH conditions:
                    # 1. Not the largest shape
                    # 2. Area is less than threshold_percent of the largest shape
                    shapes_to_delete = []
                    shapes_to_keep = []

                    for shape, volume, name in shape_info:
                        if shape != max_shape:
                            area_percentage = (volume / max_volume) * 100
                            if volume < threshold_area:
                                shapes_to_delete.append((shape, name, volume, area_percentage))
                                print(f"  → Delete: {name} (Area: {volume:,} EMU², {area_percentage:.1f}% of largest)")
                            else:
                                shapes_to_keep.append((shape, name, volume, area_percentage))
                                print(
                                    f"  → Keep: {name} (Area: {volume:,} EMU², {area_percentage:.1f}% of largest - above threshold)")

                    # Execute deletion
                    for shape, name, volume, percentage in shapes_to_delete:
                        text = shape.text if hasattr(shape, 'text') else ""
                        print(f"  → Deleting: {name} (Text: '{text}')")

                        # Delete Shape
                        sp = shape.element
                        sp.getparent().remove(sp)

                    if len(shapes_to_delete) == 0:
                        print(f"  → No shapes to delete (all shapes are above {area_threshold_percent}% threshold)")

            return prs

        # Your existing code continues here:
        out_path = self.tmp_pptx_file_path
        prs.save(out_path)

        # Add overlapping shape removal process
        print("\n=== Starting overlapping shape removal process ===")
        prs = Presentation(out_path)  # Reload the saved presentation
        prs = remove_overlapping_shapes(prs, tolerance=100000, area_threshold_percent=50)  # Remove overlapping shapes
        prs.save(out_path)  # Save again with overlapping shapes removed
        #print(f"\nProcessing complete: Saved to {out_path}")




class  ns_option_convert_to_master_yaml():
    def __init__(self):

        #parameter
        tmp_ppt_width = 30  # inches
        tmp_ppt_hight = 15  # inches

        path_array_inches = []
        device_array = []
        path_strings = []

        ### read the yaml file
        with open(str(self.full_filepath), 'r') as yml:
            config = yaml.safe_load(yml)
            #print(config)

        tmp_min_x = 999999
        tmp_min_y = 999999
        tmp_max_x = 0
        tmp_max_y = 0
        for tmp_array in config['nodes']:
            #print(tmp_array['label'],tmp_array['x'],tmp_array['y'],tmp_array['id'])
            if tmp_array['x'] < tmp_min_x:
                tmp_min_x = tmp_array['x']
            if tmp_array['y'] < tmp_min_y:
                tmp_min_y = tmp_array['y']

        for tmp_array in config['nodes']:
            #print(tmp_array['label'],tmp_array['x'],tmp_array['y'])
            device_array.append([tmp_array['label'], int (tmp_array['x'] - tmp_min_x ), int(tmp_array['y'] - tmp_min_y),tmp_array['id']])
            if int (tmp_array['x'] - tmp_min_x ) > tmp_max_x:
                tmp_max_x = int (tmp_array['x'] - tmp_min_x )
            if int(tmp_array['y'] - tmp_min_y) > tmp_max_y:
                tmp_max_y = int(tmp_array['y'] - tmp_min_y)

        #print(device_array)

        for tmp_array in config['links']:
            for i in device_array:
                if i[3] == tmp_array['n1']:
                    tmp_tmp_label_start = i[0]
                if i[3] == tmp_array['n2']:
                    tmp_tmp_label_end = i[0]

            path_strings.append([tmp_array['id'],tmp_array['i1'],tmp_array['i2'],tmp_tmp_label_start,tmp_tmp_label_end])
        #print(path_strings)

        '''create ppt'''
        ppt = Presentation()
        ppt.slide_width = Inches(tmp_ppt_width)
        ppt.slide_height = Inches(tmp_ppt_hight)

        slide_layout_5 = ppt.slide_layouts[5]
        slide = ppt.slides.add_slide(slide_layout_5)

        shapes = slide.shapes

        ### adjust ratio of ppt size
        ppt_ratio_x = float(tmp_ppt_width) / float(tmp_max_x)
        ppt_ratio_y = float(tmp_ppt_hight)/ float(tmp_max_y)

        '''create the path array'''
        for path_string in path_strings:
            start_x = 0
            start_y = 0
            end_x = 0
            end_y = 0

            for d in device_array:
                if path_string[3] == d[0]:
                    start_x = d[1]
                    start_y = d[2]

                if path_string[4] == d[0]:
                    end_x = d[1]
                    end_y = d[2]


            path_array_inches.append([start_x * ppt_ratio_x  * 0.5 + 0.5, start_y * ppt_ratio_y  * 0.5 + 0.5, end_x * ppt_ratio_x  * 0.5 + 0.5, end_y * ppt_ratio_y  * 0.5 + 0.5])
        #print(path_array_inches)

        '''add connectors'''
        for tmp_path_array_inches in path_array_inches:
            shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(tmp_path_array_inches[0]), Inches(tmp_path_array_inches[1]), \
                             Inches(tmp_path_array_inches[2]), Inches(tmp_path_array_inches[3]))

        '''add shapes'''
        for tmp_device_array in device_array:
            tmp_1st_array = []
            tmp_3rd_array = []
            tmp_1st_array.append([tmp_device_array[1],tmp_device_array[2]])
            tmp_3rd_array.append([5,30])
            #print(tmp_1st_array[0])
            shape_left = float(tmp_1st_array[0][0]) * ppt_ratio_x * 0.5 + 0.5
            shape_top = float(tmp_1st_array[0][1]) * ppt_ratio_y * 0.5 + 0.5
            #print(tmp_3rd_array[0] , ppt_ratio_x)
            shape_width = abs(tmp_3rd_array[0][0] * ppt_ratio_x)
            shape_hight = abs(tmp_3rd_array[0][1] * ppt_ratio_y)


            shapes = slide.shapes
            shapes = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(shape_left), Inches(shape_top), Inches(shape_width), Inches(shape_hight))
            shapes.adjustments[0] = 0.0
            shapes.text = tmp_device_array[0]

        '''add root folder'''
        shapes = slide.shapes
        shapes = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0), Inches(0), Inches(tmp_ppt_width), Inches(tmp_ppt_hight))
        shapes.adjustments[0] = 0.0
        shapes.text = config['lab']['title']

        ppt.save(self.tmp_pptx_file_path)


class  ns_overwrite_line_to_master_yaml():
    def __init__(self):
        import ns_def
        path_strings = []
        device_array = []

        ### read the yaml file
        with open(str(self.full_filepath), 'r') as yml:
            config = yaml.safe_load(yml)
            #print(config)

        #get node data
        for tmp_array in config['nodes']:
            #print(tmp_array)
            device_array.append([tmp_array['label'],tmp_array['id'],tmp_array['interfaces']])
        #print(device_array)

        #get line data and map to path_strings array
        for tmp_array in config['links']:
            for i in device_array:
                if i[1] == tmp_array['n1']:
                    tmp_tmp_label_start = i[0]

                    #get interface name
                    for b in i[2]:
                        if b['id'] == tmp_array['i1']:
                            tmp_tmp_int_start = b['label']

                if i[1] == tmp_array['n2']:
                    tmp_tmp_label_end = i[0]

                    #get interface name
                    for b in i[2]:
                        if b['id'] == tmp_array['i2']:
                            tmp_tmp_int_end = b['label']

            path_strings.append([tmp_array['id'],tmp_tmp_int_start,tmp_tmp_int_end,tmp_tmp_label_start,tmp_tmp_label_end])
        #print(path_strings)

        #adjust CML's interface format to NS
        new_line_array = []
        tmp_i = 1
        for tmp_path in path_strings:
            #get string part and number part from path_strings
            new_line_array.append([ns_def.adjust_portname(tmp_path[1]),ns_def.adjust_portname(tmp_path[2]),tmp_path[3],tmp_path[4],tmp_i ])
            tmp_i+=1

        #print(new_line_array)

        '''overwrite Master Data File'''
        new_master_line_array = []
        converted_tuple = {}
        tmp_used_array = []
        ppt_meta_file = str(self.excel_file_path)
        master_line_array = ns_def.convert_master_to_array('Master_Data', ppt_meta_file, '<<POSITION_LINE>>')

        for tmp_master_line_array in master_line_array:
            if tmp_master_line_array[0] >= 3:
                #print (tmp_master_line_array)
                ### Change values at each line ###
                tmp_master_line_array[1][13] = 'Unknown'
                tmp_master_line_array[1][14] = 'Unknown'
                tmp_master_line_array[1][15] = 'Unknown'
                tmp_master_line_array[1][17] = 'Unknown'
                tmp_master_line_array[1][18] = 'Unknown'
                tmp_master_line_array[1][19] = 'Unknown'
                #print(tmp_master_line_array)
                for tmp_new_line_array in new_line_array:
                    if tmp_master_line_array[1][0] == tmp_new_line_array[2] and tmp_master_line_array[1][1] == tmp_new_line_array[3] and tmp_new_line_array[4] not in tmp_used_array:
                        tmp_master_line_array[1][2] = str(tmp_new_line_array[0][0] + ' ' + str(tmp_new_line_array[0][2]))
                        tmp_master_line_array[1][3] = str(tmp_new_line_array[1][0] + ' ' + str(tmp_new_line_array[1][2]))
                        tmp_master_line_array[1][12] = str(tmp_new_line_array[0][1])
                        tmp_master_line_array[1][16] = str(tmp_new_line_array[1][1])
                        tmp_used_array.append(tmp_new_line_array[4])
                        #print(tmp_used_array)
                        break
                #print(tmp_master_line_array)
                new_master_line_array.append(tmp_master_line_array)


        converted_tuple = ns_def.convert_array_to_tuple(new_master_line_array)
        ns_def.overwrite_excel_meta(converted_tuple, self.excel_file_path, 'Master_Data', '<<POSITION_LINE>>', 0,0)


class  ns_l3_config_to_master_yaml():
    def __init__(self):
        # parameter
        l3_table_ws_name = 'Master_Data_L3'
        l3_table_file = self.full_filepath
        target_node_definition_ios = ['iosv','csr1000v','iosvl2','cat8000v']
        target_node_definition_asa = ['asav']
        target_node_definition_iosxr = ['iosxrv9000']

        #get L3 Table Excel file
        l3_table_array = []
        l3_table_array = ns_def.convert_excel_to_array(l3_table_ws_name, l3_table_file, 3)
        #print('--- l3_table_array ---')
        #print(l3_table_array)

        '''get L3 ipaddress from yaml'''
        ### read the yaml file
        with open(str(self.yaml_full_filepath), 'r') as yml:
            config = yaml.safe_load(yml)
        # print(config)

        config_array = []
        last_int_array = []
        overwrite_l3_table_array = []

        print('--- [label, node_definition, id, configuration] ---')
        for tmp_config in config['nodes']:
            config_array.append([tmp_config['label'], tmp_config['node_definition'], tmp_config['id'], tmp_config['configuration']])
            # print(tmp_config['configuration'])
            if tmp_config['node_definition'] in target_node_definition_ios or tmp_config['node_definition'] in target_node_definition_asa or tmp_config['node_definition'] in target_node_definition_iosxr:
                '''
                CiscoConfParse
                '''
                CONFIG = tmp_config['configuration']

                if tmp_config['node_definition'] in target_node_definition_ios:
                    parse = CiscoConfParse(CONFIG.splitlines(), syntax='ios', factory=True)
                elif tmp_config['node_definition'] in target_node_definition_asa:
                    parse = CiscoConfParse(CONFIG.splitlines(), syntax='asa', factory=True)
                elif tmp_config['node_definition'] in target_node_definition_iosxr:
                    parse = CiscoConfParse(CONFIG.splitlines(), syntax='iosxr', factory=True)


                int_array = [[tmp_config['label'], tmp_config['node_definition'], tmp_config['id']]]
                dummy_array = []

                for tmp_parse in parse.find_objects(r'^interface\s'):
                    int_char = list(str(tmp_parse.interface_object))
                    int_char_2 = str(tmp_parse.interface_object)

                    for i, tmp_char in enumerate(str(tmp_parse.interface_object)):
                        if re.fullmatch('[0-9]+', tmp_char):
                            # print(str(i),tmp_char)
                            int_char.insert(i, ' ')
                            int_char_2 = str("".join(int_char))
                            break

                    if str(tmp_parse.ipv4_addr) != '':
                        dummy_array.append([int_char_2, tmp_parse.ipv4_addr + '/' + str(tmp_parse.ipv4_masklength)])

                int_array.append(dummy_array)
                #print(int_array)
                last_int_array.append(int_array)

        for tmp_l3_table_array in l3_table_array:
            #print(tmp_l3_table_array[1])
            for tmp_last_int_array in last_int_array:
                if tmp_l3_table_array[1][1] == tmp_last_int_array[0][0]:
                    #print(tmp_last_int_array[0][0] , tmp_l3_table_array[1][1])
                    if tmp_last_int_array[0][0] == tmp_l3_table_array[1][1]:
                        flag_l3_exist = False
                        for tmp_tmp_last_int_array in tmp_last_int_array[1]:
                            if tmp_tmp_last_int_array[0] == tmp_l3_table_array[1][2]:
                                #print('--- L3 address match ---   ' + str(tmp_last_int_array) , str(tmp_l3_table_array))
                                if len(tmp_l3_table_array[1]) == 3:
                                    tmp_l3_table_array[1].append('')
                                tmp_l3_table_array[1].append(tmp_tmp_last_int_array[1])
                                overwrite_l3_table_array.append(tmp_l3_table_array)
                                flag_l3_exist = True
                        if flag_l3_exist == False:
                            overwrite_l3_table_array.append(tmp_l3_table_array)

        print('--- overwrite_l3_table_array ---')
        #print(overwrite_l3_table_array)


        # write to master file
        last_overwrite_l3_table_tuple = {}
        last_overwrite_l3_table_tuple = ns_def.convert_array_to_tuple(overwrite_l3_table_array)
        print('--- last_overwrite_l3_table_tuple ---')
        #print(last_overwrite_l3_table_tuple)

        master_excel_meta = last_overwrite_l3_table_tuple
        excel_file_path = self.full_filepath
        worksheet_name = l3_table_ws_name
        section_write_to = '<<L3_TABLE>>'
        offset_row = 2
        offset_column = 0
        ns_def.overwrite_excel_meta(master_excel_meta, excel_file_path, worksheet_name, section_write_to, offset_row, offset_column)



