'''
SPDX-License-Identifier: Apache-2.0

Copyright 2023 Cisco Systems, Inc. and its affiliates

Licensed under the Apache License, Version 2.0 (the "License");
you may not use this file except in compliance with the License.
You may obtain a copy of the License at

http://www.apache.org/licenses/LICENSE-2.0

Unless required by applicable law or agreed to in writing, software
distributed under the License is distributed on an "AS IS" BASIS,
WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
See the License for the specific language governing permissions and
limitations under the License.
'''
import ast
import copy

from pptx import *
import sys, os, re, shutil
import numpy as np
import math
import nsm_def , nsm_ddx_figure
import openpyxl
from pptx import Presentation
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from pptx.dml.line import LineFormat
from pptx.shapes.connector import Connector
from pptx.util import Inches, Cm, Pt
from collections import defaultdict
import time

def _ts():
    return time.perf_counter()

def log_elapsed(label, t0):
    print(f"[TIME] {label}: {(_ts() - t0):.3f}s")

def get_text_wh_cached(ctx, font_size, text):
    """
    ctx: ns_front_run ã® selfï¼ˆæœ¬ãƒ„ãƒ¼ãƒ«ã§ã¯ nsm_l3_diagram_create.__init__/l3_area_create ã«æ¸¡ã•ã‚Œã‚‹ selfï¼‰
    """
    if not hasattr(ctx, "_text_size_cache") or ctx._text_size_cache is None:
        ctx._text_size_cache = {}

    if text is None:
        text = ""
    else:
        text = str(text)

    key = (font_size, text)
    v = ctx._text_size_cache.get(key)
    if v is None:
        v = nsm_def.get_description_width_hight(font_size, text)
        ctx._text_size_cache[key] = v
    return v

class  nsm_l3_diagram_create():
    def __init__(self):
        import time
        _t0 = time.perf_counter()
        #print('--- nsm_l3_diagram_create ---')
        '''
        STEP0 get values of Master Data
        '''
        #parameter
        ws_name = 'Master_Data'
        ws_l2_name = 'Master_Data_L2'
        ws_l3_name = 'Master_Data_L3'
        excel_maseter_file = self.inFileTxt_L3_3_1.get()

        self._text_size_cache = {}
        self._l3_ipset_cache = {}  # key: (device, if) -> list (remake_array)

        if self.click_value_l3 == 'L3-4-1':
            excel_maseter_file = self.outFileTxt_L3_3_5_1.get()

        _cache = getattr(self, '_l3_data_cache', None)
        if _cache is not None:
            self.result_get_l2_broadcast_domains = _cache['result']
            self.update_l2_table_array = _cache['update_l2']
            self.target_l2_broadcast_group_array = _cache['target_groups']
            self.position_folder_array = _cache['pos_folder_arr']
            self.position_shape_array = _cache['pos_shape_arr']
            self.position_line_array = _cache['pos_line_arr']
            self.position_style_shape_array = _cache['pos_style_arr']
            self.position_tag_array = _cache['pos_tag_arr']
            self.root_folder_array = _cache['root_folder_arr']
            self.position_folder_tuple = _cache['pos_folder_t']
            self.position_shape_tuple = _cache['pos_shape_t']
            self.position_line_tuple = _cache['pos_line_t']
            self.position_style_shape_tuple = _cache['pos_style_t']
            self.position_tag_tuple = _cache['pos_tag_t']
            self.root_folder_tuple = _cache['root_folder_t']
            self.folder_wp_name_array = _cache['folder_wp']
            self.l2_table_array = _cache.get('l2_table_arr', [])
            self.l3_table_array = _cache.get('l3_table_arr', [])
            self.device_list_array = _cache.get('device_list', [])
            self.wp_list_array = _cache.get('wp_list', [])
            self.all_shape_list_array = _cache.get('all_shape_list', [])
            self.update_l3_table_array = _cache['update_l3']
            self.l3_rows_by_device = _cache['l3_by_dev']
            self.l3_rows_by_device_if = _cache['l3_by_dev_if']
            self.l2_rows_by_area = _cache['l2_by_area']
            self.groups_by_member = _cache['grp_by_member']
            self.members_by_group = _cache['members_by_grp']
            self.device_l2_boradcast_domain_array = _cache['result'][1]
            self.device_l2_directly_l3vport_array = _cache['result'][2]
            self.device_l2_other_array = _cache['result'][3]
            self.marged_l2_broadcast_group_array = _cache['result'][4]
        else:
            self.result_get_l2_broadcast_domains = nsm_def.get_l2_broadcast_domains.run(self, excel_maseter_file)

        # ============================================================
        # BUILD INDEXES for faster lookup
        # ============================================================
        import time
        _idx_start = time.time()

        if _cache is None:
            self.l3_table_array = nsm_def.convert_master_to_array(ws_l3_name, excel_maseter_file, '<<L3_TABLE>>')
            self.update_l3_table_array = []
            for tmp_l3_table_array in self.l3_table_array:
                if tmp_l3_table_array[0] != 1 and tmp_l3_table_array[0] != 2:
                    tmp_l3_table_array[1].extend(['', '', '', ''])
                    del tmp_l3_table_array[1][6:]
                    work_new_l3_table_array = tmp_l3_table_array[1][4].split(',')
                    add_ip_address_set_array = []
                    for tmp_tmp_new_l3_table_array in work_new_l3_table_array:
                        if nsm_def.check_ip_format(tmp_tmp_new_l3_table_array) == 'IPv4':
                            change_tmp_ip_address = str(tmp_tmp_new_l3_table_array).replace('[', '').replace(']', '').replace('\'', '').replace(' ', '')
                            ip_address_set_array = nsm_def.get_ip_address_set(change_tmp_ip_address)
                            add_ip_address_set_array.append(ip_address_set_array)
                    tmp_l3_table_array[1].append(str(add_ip_address_set_array))
                    self.update_l3_table_array.append(tmp_l3_table_array[1])

            self.l3_rows_by_device = {}
            for row in self.update_l3_table_array:
                dev = row[1]
                self.l3_rows_by_device.setdefault(dev, []).append(row)

            self.l3_rows_by_device_if = {}
            for row in self.update_l3_table_array:
                key = (row[1], row[2])
                self.l3_rows_by_device_if.setdefault(key, []).append(row)

            self.l2_rows_by_area = {}
            for row in self.update_l2_table_array:
                self.l2_rows_by_area.setdefault(row[0], []).append(row)

            self.groups_by_member = {}
            self.members_by_group = {}
            for gid, group in enumerate(self.target_l2_broadcast_group_array):
                self.members_by_group[gid] = group[1]
                for member in group[1]:
                    self.groups_by_member.setdefault(tuple(member), []).append(gid)

            if not (self.click_value_l3 == 'L3-4-1' and self.flag_re_create == True):
                print(f"  Areas: {len(self.l2_rows_by_area)}")
                print(f"  Devices: {len(self.l3_rows_by_device)}")
                print(f"  Broadcast groups: {len(self.members_by_group)}")

            self.device_l2_boradcast_domain_array = self.result_get_l2_broadcast_domains[1]
            self.device_l2_directly_l3vport_array = self.result_get_l2_broadcast_domains[2]
            self.device_l2_other_array = self.result_get_l2_broadcast_domains[3]
            self.marged_l2_broadcast_group_array = self.result_get_l2_broadcast_domains[4]



        # GET way point with folder tuple
        self.wp_with_folder_tuple = {}
        for tmp_wp_folder_name in self.folder_wp_name_array[1]:
            current_row = 1
            flag_start_row = False
            flag_end_row = False

            while flag_end_row == False:
                if str(self.position_shape_tuple[current_row, 1]) == tmp_wp_folder_name:
                    start_row = current_row
                    flag_start_row = True
                if flag_start_row == True and str(self.position_shape_tuple[current_row, 1]) == '<END>':
                    flag_end_row = True
                    end_row = current_row - 1
                current_row += 1
            # print(tmp_wp_folder_name,start_row,end_row)

            for i in range(start_row, end_row + 1):
                flag_start_column = False
                current_column = 2
                while flag_start_column == False:
                    if str(self.position_shape_tuple[i, current_column]) != '<END>':
                        self.wp_with_folder_tuple[self.position_shape_tuple[i, current_column]] = tmp_wp_folder_name
                    else:
                        flag_start_column = True
                    current_column += 1

        #print('---- wp_with_folder_tuple ----')
        #print(self.wp_with_folder_tuple)
        print(f"[TIME] create_master_file_one_area: {time.perf_counter() - _t0:.3f}s")

        '''
        Create per area l3 ppt
        '''
        if self.click_value == 'L3-3-2':
            '''GET SIZE'''
            self.page_size_array = []
            self.slide_width = 0.0
            self.slide_hight = 0.0
            for tmp_new_position_folder_array in self.folder_wp_name_array[0]:
                action_type = 'GET_SIZE'
                offset_x = 0.0 #inches
                offset_y = 0.0 #inches
                self.page_size_array.append( nsm_l3_diagram_create.l3_area_create(self, tmp_new_position_folder_array , action_type ,offset_x ,offset_y))

            for tmp_page_size_array in self.page_size_array:
                if self.slide_width < tmp_page_size_array[3]:
                    self.slide_width = tmp_page_size_array[3]
                if self.slide_hight < tmp_page_size_array[4]:
                    self.slide_hight = tmp_page_size_array[4]

            #add page margin
            self.slide_width += 1.0 * 2  #page margin
            self.slide_hight += 1.0 * 2  #page margin

            #print('--- self.page_size_array ,self.slide_width ,self.slide_hight ---  [outline_shape_type, outline_shape_left, outline_shape_top, outline_shape_width, outline_shape_hight, folder_shape_text] ,self.slide_width ,self.slide_hight')
            #print(self.page_size_array,self.slide_width ,self.slide_hight)

            # Calculate area offset for one area at ver 2.3.3
            if self.click_value_l3 == 'L3-4-1':
                create_master_file_one_area.calculate_area_offset(self)

                # Compute per-device rightward shifts to avoid L3 segment
                # connector lines passing over other devices. Runs only in
                # 2nd pass (after 1st pass populated _collision_lines_per_area
                # and _collision_rows_per_area). The shifts are applied during
                # the device draw via the per-device hook below; this pass
                # also bumps calculated_max_right_edge so slide_width adjusts.
                if self.flag_re_create == True and getattr(self, '_l3_render_quality', 3) >= 3:
                    # Pass 1 (cross-area): seed device_extra_shift and detect
                    # over-spread. If over-spread, immediately re-pack with
                    # same-area-only blocker checking + left-pack relaxation so
                    # the 2nd pass CREATE below already draws the packed layout
                    # (previously this packing ran as a separate 4th pass after
                    # rendering & discarding the spread layout). The packed
                    # calculated_max_right_edge then drives slide_width below.
                    self._l3_packed_applied = False
                    compute_l3_device_shifts(self)
                    if _detect_l3_overspread(self):
                        self._l3_same_area_only = True
                        try:
                            compute_l3_device_shifts(self)
                            self._l3_packed_applied = True
                        finally:
                            self._l3_same_area_only = False
                elif self.flag_re_create == True:
                    # Quality 2: skip the entire 3rd pass (device collision-shift).
                    # Devices keep their natural left-aligned positions from the
                    # 1st/2nd pass; connector lines are NOT routed around devices.
                    # device_extra_shift stays empty so the per-device draw hook
                    # (left_offset += device_extra_shift.get(name, 0.0)) adds 0.
                    self.device_extra_shift = {}
                    self._l3_packed_applied = False

                # Recalculate slide_width based on actual area positions after calculate_area_offset
                # calculated_max_right_edge is computed in calculate_area_offset() and contains
                # the maximum (area_start_x + area_width) across all areas
                #
                # Note: area_width from calculate_area_offset is based on device positions only.
                # The actual outline extends beyond devices with additional margins:
                # - between_shape_column (0.5") on each side of devices
                # - Outline has between_shape_column * 2 (1.0") additional margin
                # Total additional margin needed: approximately 1.0" - 1.5" per area edge
                if hasattr(self, 'calculated_max_right_edge') and self.calculated_max_right_edge > 0:
                    # Add margins:
                    # - 1.0" left margin (standard)
                    # - 1.0" right margin (standard)
                    # - 1.0" additional for outline shape margins (between_shape_column * 2)
                    outline_margin = 1.0  # Additional margin for outline shape beyond device positions
                    required_width = self.calculated_max_right_edge + 2.0 + outline_margin
                    # Apply maximum slide size limit (56 inches as per PowerPoint limitation)
                    if required_width > 56.0:
                        required_width = 56.0
                    if required_width > self.slide_width:
                        print(f"[L3 All Areas] Adjusting slide_width: {self.slide_width:.2f} -> {required_width:.2f} inches (max_right_edge={self.calculated_max_right_edge:.2f})")
                        self.slide_width = required_width

                # Snap WayPoints in `_wp_`-only Area Position rows AFTER the
                # slide_width adjustment above so the overflow check uses the
                # final diagram width (not the pre-bump GET_SIZE value).
                # Otherwise WayPoints anchored near the rightmost device would
                # spuriously exceed overflow_limit and fall back to natural
                # left (observed on Sample.figure5 NAGAINAMAE-TEST-WAN-fix001
                # which should snap above DC-TOP3 / FW-12 but was skipped
                # because slide_width was still the pre-bump 24.01 instead of
                # the final 27.58).
                if self.flag_re_create == True:
                    _compute_wp_x_snap_targets(self)

            '''CREATE L3 DIAGRAM (reuse cached result from __init__)'''
            #self.result_get_l2_broadcast_domains = nsm_def.get_l2_broadcast_domains.run(self, excel_maseter_file)  ## 'self.update_l2_table_array, device_l2_boradcast_domain_array, device_l2_directly_l3vport_array, device_l2_other_array, marged_l2_broadcast_group_array'
            #self.active_ppt = Presentation()  # define target ppt object


            if os.path.exists(self.output_ppt_file) and self.flag_second_page == True:
                self.active_ppt = Presentation(self.output_ppt_file)
                slide_layout = self.active_ppt.slide_layouts[5]  # Blank layout
                #self.active_ppt.slides.add_slide(slide_layout)
            else:
                self.active_ppt = Presentation()

            # Snapshot mutable state before 2nd pass CREATE so the conditional
            # 3rd pass (overlap-fix re-render) can restore baseline accurately.
            # This is only needed in 2nd pass (flag_re_create=True) AllAreas.
            # `svg_capture_list_len` lets the SVG path truncate captured
            # entries before re-render (set in nsm_l3_svg_create._run_all_areas).
            # `position_shape_array` / `position_folder_array` are deepcopied
            # because l3_area_create mutates inner lists (del at L394-397).
            if self.flag_re_create == True and self.click_value_l3 == 'L3-4-1':
                import copy as _copy
                _svg_cap = getattr(self, '_svg_capture_list', None)
                self._state_snapshot = {
                    'y_grid_segment_array': list(self.y_grid_segment_array),
                    'add_shape_array': list(self.add_shape_array),
                    'add_shape_write_array': list(getattr(self, 'add_shape_write_array', [])),
                    'per_index2_after_array': list(self.per_index2_after_array),
                    'add_shape_array_len': len(self.add_shape_array),
                    'svg_capture_list_len': (len(_svg_cap) if _svg_cap is not None else None),
                    'position_shape_array': _copy.deepcopy(self.position_shape_array),
                    'position_folder_array': _copy.deepcopy(getattr(self, 'position_folder_array', [])),
                }

            for tmp_new_position_folder_array in self.folder_wp_name_array[0]:
                action_type = 'CREATE'
                offset_x = 0.0 #inches
                offset_y = 0.0 #inches

                for tmp_page_size_array in self.page_size_array:
                    if tmp_page_size_array[5] == tmp_new_position_folder_array:
                        offset_x = tmp_page_size_array[1]
                        offset_y = tmp_page_size_array[2]
                        break

                nsm_l3_diagram_create.l3_area_create(self, tmp_new_position_folder_array , action_type,offset_x ,offset_y)

            ### save pptx file
            # Skip save on 1st pass (only for All Areas L3-4-1 mode)
            if self.flag_re_create == False and self.click_value_l3 == 'L3-4-1':
                #print("[L3 Diagram] 1st pass done (skip save)")
                return

            # Runaway-device pull-back (quality-independent): a device connected
            # by a long L3 segment can land far to the right of its area's
            # device cluster (e.g. SRV_* at the end of a cross-diagram segment),
            # which makes the area frame overrun and overlap the next area's
            # frame. The grid-column layout does not cause this (it is driven by
            # L3 connectivity), so it cannot be fixed by compacting the master.
            # Here we measure the drawn positions, and for any area whose
            # rightmost device is separated from the rest of the cluster by a
            # large horizontal gap, we pull that device back to align with the
            # area's rightmost cluster device via device_extra_shift, then
            # re-render. Runs before the segment lane-separation retry so the
            # latter operates on the final X positions.
            if self.flag_re_create == True and self.click_value_l3 == 'L3-4-1' and hasattr(self, '_state_snapshot'):
                import copy as _pb_copy
                PULLBACK_MAX_RETRIES = 8
                # Phase A - horizontal compaction to convergence. Each area's
                # devices are swept and any oversized empty band (a device
                # cluster dragged far right next to its cross-area peers) is
                # closed by pulling the right group left. This is monotonic-left
                # and bounded by the area's natural packed width, so it always
                # converges, and it is the primary tool for over-spread areas
                # (long inter-area connector lines are accepted as the trade-off
                # for removing the overrun). The resulting state is snapshotted
                # as the safe fallback: it keeps the page width sane even if the
                # following separation dance fails to settle.
                for _pb_retry in range(PULLBACK_MAX_RETRIES):
                    if _pull_back_runaway_devices(self):
                        _restore_state_and_rerender(self)
                        continue
                    break
                _pb_shift_baseline = _pb_copy.deepcopy(getattr(self, 'device_extra_shift', {}) or {})
                _pb_slide_baseline = getattr(self, 'slide_width', 0.0)
                # Phase B - interleaved compaction + separation. Separation
                # pushes overlapping area frames apart (right); a bridging device
                # that belongs to the left area but is wired into the right area
                # chases that push, so compaction re-fires to pull it back,
                # shrinking the left area until the frames clear. This dance is
                # what resolves side-by-side area overlaps. It is guarded against
                # divergence (dense masters where areas interleave within rows so
                # the chase never settles): if the separation shift stops
                # shrinking we roll back to the Phase-A compaction-only fallback
                # rather than ballooning the page width.
                _pb_prev_sep = None
                _pb_diverged = False
                for _pb_retry in range(PULLBACK_MAX_RETRIES):
                    if _pull_back_runaway_devices(self):
                        _restore_state_and_rerender(self)
                        continue
                    _pb_sep = _separate_overlapping_areas(self)
                    if _pb_sep > 1e-6:
                        if _pb_prev_sep is not None and _pb_sep > _pb_prev_sep - 1e-6:
                            _pb_diverged = True
                            break
                        _pb_prev_sep = _pb_sep
                        _restore_state_and_rerender(self)
                        continue
                    break
                if _pb_diverged:
                    self.device_extra_shift = _pb_shift_baseline
                    self.slide_width = _pb_slide_baseline
                    _restore_state_and_rerender(self)

            # Conditional 3rd pass: detect same-Y X-overlap among segment
            # bars in 2nd pass entries and retry CREATE with adjusted
            # optimize_y_grid_array until no overlap or MAX_RETRIES reached.
            # The device packing for over-spread layouts now happens BEFORE the
            # 2nd pass CREATE (see the compute_l3_device_shifts block after
            # calculate_area_offset), so the 2nd pass already drew the packed
            # layout. This lane-separation retry therefore operates directly on
            # the packed bars. When the packed path ran, the packed layout needs
            # more lane-separation iterations (shorter bars create new same-Y
            # X-overlaps that must be pushed onto free lanes), so the retry
            # budget is raised; otherwise the original budget is kept so
            # non-over-spread diagrams (5site/flow) behave exactly as before.
            if self.flag_re_create == True and self.click_value_l3 == 'L3-4-1' and hasattr(self, '_state_snapshot') and getattr(self, '_l3_render_quality', 3) >= 3:
                MAX_RETRIES = 120 if getattr(self, '_l3_packed_applied', False) else 50
                for _retry in range(MAX_RETRIES):
                    had_overlap = _detect_and_fix_2nd_pass_overlaps(self)
                    if not had_overlap:
                        break
                    _restore_state_and_rerender(self)
                else:
                    print(f"[L3 overlap fix] Warning: did not converge after {MAX_RETRIES} retries")

            ### save pptx file
            self.active_ppt.save(self.output_ppt_file)

        '''
        Modify style of device for All Areas at ver 2.3.0
        '''
        if self.click_value_l3 == 'L3-4-1':
            from pptx.dml.color import RGBColor
            prs = Presentation(self.output_ppt_file)

            for slide in prs.slides:
                shapes_to_process = list(slide.shapes)

                while shapes_to_process:
                    shape = shapes_to_process.pop()
                    if shape.has_text_frame:
                        for text in shape.text.splitlines():
                            if text in self.global_wp_array:
                                shape.fill.solid()
                                shape.fill.fore_color.rgb = RGBColor(220, 230, 242)
                                ### apply attribute color to shape at ver 2.4.0
                                tmp_rgp_color = self.attribute_tuple1_1[text]
                                shape.fill.fore_color.rgb = RGBColor(tmp_rgp_color[0], tmp_rgp_color[1],tmp_rgp_color[2])

                                if shape.adjustments:
                                    shape.adjustments[0] = 0.2002

            #print(self.output_ppt_file)
            prs.save(self.output_ppt_file)

    def l3_area_create(self, target_folder_name, action_type,offset_x ,offset_y):
        #print('--- l3_area_create -',action_type,' - ',target_folder_name,'---')
        import os as _os
        if _os.environ.get('NS_DEBUG_L3', '0') == '1':
            try:
                with open('c:/work_local/dbg_outline.txt', 'a', encoding='utf-8') as _f:
                    _f.write(f'[START] l3_area_create {target_folder_name} {action_type}\n')
            except OSError:
                pass
        t0 = _ts()

        ### add for y-grid optimize at ve 2.4.1
        if self.flag_re_create == True and self.flag_second_page == False and action_type == 'GET_SIZE':
            self.optimize_y_grid_array = get_optimize_y_grid_array(self)
            #print(self.optimize_y_grid_array)

        self.used_l3segment_array = []
        self._used_l3segment_set = set()
        ### get l3segment in the target folder
        target_all_device_array = []

        # Use pre-filtered rows by area
        area_l2_rows = self.l2_rows_by_area.get(target_folder_name, [])
        for tmp_update_l2_table_array in area_l2_rows:
            if 'L3' in tmp_update_l2_table_array[2]:
                target_all_device_array.append([tmp_update_l2_table_array[1], tmp_update_l2_table_array[3]])
            if 'L3' in tmp_update_l2_table_array[4]:
                target_all_device_array.append([tmp_update_l2_table_array[1], tmp_update_l2_table_array[5]])

        target_all_device_array = nsm_def.get_l2_broadcast_domains.get_unique_list(target_all_device_array)
        # Build a set of tuples for O(1) membership tests (list elements are lists, so convert to tuples)
        target_all_device_set = {tuple(x) for x in target_all_device_array}
        #print(target_all_device_array)

        update_l2_broadcast_group_array = []
        for tmp_target_l2_broadcast_group_array in self.target_l2_broadcast_group_array:
            for tmp_tmp_target_l2_broadcast_group_array in tmp_target_l2_broadcast_group_array[1]:
                if tuple(tmp_tmp_target_l2_broadcast_group_array) in target_all_device_set:
                    #print(tmp_tmp_target_l2_broadcast_group_array)
                    update_l2_broadcast_group_array.append(tmp_target_l2_broadcast_group_array)
                    break

        ### add at ver 2.4.1  remove no l3if device from self.update_l3_table_array
        l3_device_array = [sub_array[1] for sub_array in self.update_l3_table_array]
        # Removing duplicates while preserving order
        l3_device_array = list(dict.fromkeys(l3_device_array))
        # Build set for O(1) membership tests in the filter below
        l3_device_set = set(l3_device_array)
        #print(l3_device_array)
        # Create the filtered array
        l3only_position_shape_array = []

        for index, sub_array in self.position_shape_array:
            if index == 1:
                # Append the original array for index 1
                l3only_position_shape_array.append([index, sub_array])
            else:
                first_item = sub_array[0]
                last_item = sub_array[-1]
                # Filter devices present in l3_device_set, keeping '_AIR_' and the first/last items
                filtered_devices = [
                    item for item in sub_array
                    if item in l3_device_set or item == '_AIR_' or item == first_item or item == last_item
                ]
                l3only_position_shape_array.append([index, filtered_devices])

        # set the result
        #print(l3only_position_shape_array)
        local_position_shape_array = l3only_position_shape_array

        target_position_shape_array = []
        flag_match_folder = False
        for tmp_position_shape_array in local_position_shape_array:
            if tmp_position_shape_array[1][0] == target_folder_name:
                del tmp_position_shape_array[1][0]
                del tmp_position_shape_array[1][-1]
                target_position_shape_array.append(tmp_position_shape_array[1])
                flag_match_folder = True

            elif flag_match_folder == True and tmp_position_shape_array[1][0] == '':
                del tmp_position_shape_array[1][0]
                del tmp_position_shape_array[1][-1]
                target_position_shape_array.append(tmp_position_shape_array[1])
            elif flag_match_folder == True and tmp_position_shape_array[1][0] != '':
                break

        # for bug fix. add at ver 2.5.1d
        target_position_shape_array = [item for item in target_position_shape_array if not all(element == '_AIR_' for element in item)]

        #print('--- target_position_shape_array ---')
        #print(target_position_shape_array,len(target_position_shape_array))

        ### Kyuusai if len(target_position_shape_array) == 1
        if len(target_position_shape_array) == 1:
            target_position_shape_array.append(['_AIR_9999']) # add dummy shape to second row

        ### get WP of <<POSITION_FOLDER>>in MASTER EXCEL
        wp_exist_array = [[], [], [], []]  # up/down/left/right
        target_folder_row = 1

        # check left/right wp
        for tmp_position_folder_array in self.position_folder_array:
            if tmp_position_folder_array[1][0] != '<<POSITION_FOLDER>>' and tmp_position_folder_array[1][0] != '<SET_WIDTH>':
                #print(tmp_position_folder_array)
                for index, tmp_folder_name in enumerate(tmp_position_folder_array[1]):
                    if tmp_folder_name == target_folder_name:
                        target_folder_row = tmp_position_folder_array[0]
                        if len(tmp_position_folder_array[1]) >= (index + 2):
                            if '_wp_' in str(tmp_position_folder_array[1][index+1]):
                                #print('right ', tmp_position_folder_array[1][index+1])
                                wp_exist_array[3].extend([tmp_position_folder_array[1][index+1]])
                        if index >= 1:
                            if '_wp_' in str(tmp_position_folder_array[1][index - 1]):
                                #print('left ', tmp_position_folder_array[1][index - 1])
                                wp_exist_array[2].extend([tmp_position_folder_array[1][index - 1]])

        #check up/down wp
        for tmp_position_folder_array in self.position_folder_array:
            if tmp_position_folder_array[0]  == target_folder_row -2:
                for tmp_tmp_position_folder_array in tmp_position_folder_array[1]:
                    if '_wp_' in str(tmp_tmp_position_folder_array):
                        wp_exist_array[0].extend([tmp_tmp_position_folder_array])

            if tmp_position_folder_array[0]  == target_folder_row +2:
                for tmp_tmp_position_folder_array in tmp_position_folder_array[1]:
                    if '_wp_' in str(tmp_tmp_position_folder_array):
                        wp_exist_array[1].extend([tmp_tmp_position_folder_array])

        #print('--- wp_exist_array ---')
        #print(wp_exist_array,'  #up/down/left/right')

        #convert _wp_ folder name to shape name
        new_wp_exist_array = [[], [], [], []]

        for index,tmp_wp_exist_array in enumerate(wp_exist_array):
            for tmp_tmp_wp_exist_array in tmp_wp_exist_array:
                for tmp_wp_with_folder_tuple in self.wp_with_folder_tuple:
                    if tmp_tmp_wp_exist_array == str(self.wp_with_folder_tuple[tmp_wp_with_folder_tuple]):
                        #check if wp is connected to own folder

                        #check wp name
                        for tmp_update_l2_broadcast_group_array in update_l2_broadcast_group_array:
                            for tmp_tmp_update_l2_broadcast_group_array in tmp_update_l2_broadcast_group_array[1]:
                                #print(tmp_tmp_target_l2_broadcast_group_array )
                                if tmp_tmp_update_l2_broadcast_group_array[0] == str(tmp_wp_with_folder_tuple):
                                    # print(index,tmp_wp_with_folder_tuple)
                                    new_wp_exist_array[index].extend([tmp_wp_with_folder_tuple])

                        #check broadcast domain number
                        for tmp_device_l2_boradcast_domain_array in self.device_l2_boradcast_domain_array:
                            for tmp_tmp_device_l2_boradcast_domain_array in tmp_device_l2_boradcast_domain_array:
                                #print(tmp_tmp_device_l2_boradcast_domain_array)
                                if tmp_tmp_device_l2_boradcast_domain_array[1] == str(tmp_wp_with_folder_tuple):
                                    for tmp_update_l2_broadcast_group_array in update_l2_broadcast_group_array:
                                        for tmp_tmp_update_l2_broadcast_group_array in tmp_update_l2_broadcast_group_array:
                                            if  tmp_tmp_device_l2_boradcast_domain_array[0] in tmp_tmp_update_l2_broadcast_group_array :
                                                #print(tmp_tmp_device_l2_boradcast_domain_array[0],tmp_tmp_device_l2_boradcast_domain_array[1])
                                                new_wp_exist_array[index].extend([tmp_wp_with_folder_tuple])

        new_wp_exist_array = [list(set(new_wp_exist_array [0])),list(set(new_wp_exist_array [1])),list(set(new_wp_exist_array [2])),list(set(new_wp_exist_array[3]))]

        #print('--- new_wp_exist_array ---')
        #print(new_wp_exist_array ,'  #up/down/left/right')
        self.new_wp_exist_array = new_wp_exist_array

        marge_target_position_shape_array = target_position_shape_array
        if new_wp_exist_array[0] != []:
            marge_target_position_shape_array.insert(0,new_wp_exist_array[0])
        if new_wp_exist_array[1] != []:
            marge_target_position_shape_array.append(new_wp_exist_array[1])

        #print('--- marge_target_position_shape_array ---')
        #print(marge_target_position_shape_array)

        wp_marge_target_position_shape_array = marge_target_position_shape_array

        #left_right_wp_row_num = math.floor(len(marge_target_position_shape_array) * 0.5 - 1)
        left_right_wp_row_num = math.floor((len(marge_target_position_shape_array) - 1) * 0.5)     # updated

        # write wp left
        if new_wp_exist_array[2] != []:
            for tmp_i in new_wp_exist_array[2]:
                shape_text = tmp_i

                # insert wp to edge of left side
                wp_marge_target_position_shape_array[left_right_wp_row_num].insert(0, shape_text)

        # write wp right
        if new_wp_exist_array[3] != []:
            for tmp_i in new_wp_exist_array[3]:
                shape_text = tmp_i

                # insert wp to edge of right side
                wp_marge_target_position_shape_array[left_right_wp_row_num].extend([shape_text])

        #print('--- wp_marge_target_position_shape_array---')
        #print(wp_marge_target_position_shape_array)

        #get up down l3if array
        self.up_down_l3if_array = get_up_down_l3if_count(self,wp_marge_target_position_shape_array)

        # Rebuild device-interface lookup caches to match the newly computed up_down_l3if_array.
        # These caches must be refreshed each time l3_area_create is called for a different area;
        # without invalidation the stale first-area cache would be reused for all subsequent areas.
        self._up_by_device = {}
        for item in self.up_down_l3if_array[0]:
            self._up_by_device.setdefault(item[0], []).append(item)
        self._down_by_device = {}
        for item in self.up_down_l3if_array[1]:
            self._down_by_device.setdefault(item[0], []).append(item)

        #print('--- self.up_down_l3if_array ---  # up/down')
        #print(self.up_down_l3if_array[0])
        #print(self.up_down_l3if_array[1])

        ### get index for target_position_shape_array
        #get index_1  -> target_position_shape_array
        self.index_1_array = []
        for self.index_1,tmp_target_position_shape_array in enumerate(target_position_shape_array):
            #print(self.index_1,tmp_target_position_shape_array )
            self.index_1_array.append(self.index_1)

        #get index_11  -> marge_target_position_shape_array
        self.index_11_array = []
        for self.index_11,tmp_target_position_shape_array in enumerate(marge_target_position_shape_array):
            #print(self.index_11,tmp_target_position_shape_array )
            self.index_11_array.append(self.index_1)


        ### GET identify L3IF that is connected to l3 segment
        self.l3_if_has_l3_segment_array = []
        self._l3_if_has_l3_segment_set = set()
        return_get_l3_segment_array = []

        self.used_l3segment_array = []
        self._used_l3segment_set = set()
        self._combined_device_set = set(item for sublist in target_position_shape_array for item in sublist)
        for self.index_2, tmp_target_position_shape_array in enumerate(target_position_shape_array):
            for return_get_l3_segment_num in get_l3_segment_num(self,tmp_target_position_shape_array,target_position_shape_array)[1]:
                return_get_l3_segment_array.append(return_get_l3_segment_num)

                for tmp_target_l2_broadcast_group_array in self.target_l2_broadcast_group_array:
                    if return_get_l3_segment_num in tmp_target_l2_broadcast_group_array[1]:
                        self.l3_if_has_l3_segment_array.extend(tmp_target_l2_broadcast_group_array[1])
                        for _m in tmp_target_l2_broadcast_group_array[1]:
                            self._l3_if_has_l3_segment_set.add(tuple(_m))

        ### GET L3 instance
        self.defalut_l3_instance_name = 'Defalut'
        self.l3_instance_array = []
        self.update_l3_instance_array = []
        for tmp_update_l3_table_array in self.update_l3_table_array:
            if tmp_update_l3_table_array[3].replace(' ','') != '':
                self.l3_instance_array.append([tmp_update_l3_table_array[1],tmp_update_l3_table_array[3]])

        self.l3_instance_array = nsm_def.get_l2_broadcast_domains.get_unique_list(self.l3_instance_array)
        #print('--- self.l3_instance_array ---')
        #print(self.l3_instance_array)

        for tmp_update_l3_table_array in self.update_l3_table_array:
            for tmp_l3_instance_array  in self.l3_instance_array:
                if tmp_l3_instance_array[0] == tmp_update_l3_table_array[1]:
                    if tmp_update_l3_table_array[3].replace(' ','') != '':
                        self.update_l3_instance_array.append([tmp_update_l3_table_array,tmp_update_l3_table_array[3]])
                        #print([tmp_update_l3_table_array,tmp_update_l3_table_array[3]])
                    else:
                        self.update_l3_instance_array.append([tmp_update_l3_table_array, self.defalut_l3_instance_name])
                        #print([tmp_update_l3_table_array, self.defalut_l3_instance_name])

        #print('--- self.update_l3_instance_array ---')
        #print(self.update_l3_instance_array)


        '''
        make ppt diagram
        '''
        ### create slide
        if action_type == 'CREATE':
            # 1st pass (All Areas): only measure y-grid (skip most shapes/tags/lines)
            minimal_1st_pass = (self.click_value_l3 == 'L3-4-1' and
                               self.flag_re_create == False and
                               self.flag_second_page == False)

            self.active_ppt.slide_width = Inches(10.0)
            self.active_ppt.slide_height = Inches(5.0)

            ### adjust to large size , ver 2.1
            if self.slide_width > 56.0:
                self.slide_width = 56.0
            if self.slide_hight > 56.0:
                self.slide_hight = 56.0

            #input from get_size
            if self.active_ppt.slide_width < Inches(self.slide_width):
                self.active_ppt.slide_width = Inches(self.slide_width)
            if self.active_ppt.slide_height < Inches(self.slide_hight):
                self.active_ppt.slide_height = Inches(self.slide_hight)

            self.title_only_slide_layout = self.active_ppt.slide_layouts[5]
            self.slide = self.active_ppt.slides.add_slide(self.title_only_slide_layout)
            self.slide.shapes.title.left = Inches(0.0)
            self.slide.shapes.title.top = Inches(0.0)
            self.slide.shapes.title.width = Inches(10.0)
            self.slide.shapes.title.height = Inches(1.0)

            self.shape = self.slide.shapes
            self.shape.title.text = '[L3] ' + target_folder_name
            if self.flag_second_page == True and self.click_value_l3 == 'L3-4-1':
                self.shape.title.text = '[L3] ' + target_folder_name +'  <Focus on Connectivity>'
            elif self.flag_second_page == False and self.click_value_l3 == 'L3-4-1':
                self.shape.title.text = '[L3] ' + target_folder_name + '  <Focus on Area>'

            if self.click_value_VPN == 'VPN-1-3':  # add ver 2.3.2
                if self.flag_second_page == True and self.click_value_l3 == 'L3-4-1':
                    self.shape.title.text = '[VPNs on L3] <Focus on Connectivity>'
                elif self.flag_second_page == False and self.click_value_l3 == 'L3-4-1':
                    self.shape.title.text = '[VPNs on L3] <Focus on Area>'

        ### parameter
        self.left_margin = 1.0 # Inches
        self.top_margin = 1.0 # Inches

        # input from get_size
        if action_type == 'CREATE':
            self.left_margin = self.left_margin - offset_x + self.left_margin # Inches
            self.top_margin = self.top_margin - offset_y + self.top_margin  # Inches

        top_offset = 0.0 #Inches
        left_offset = 0.0 #Inches

        self.folder_font_type = 'Calibri'
        self.folder_font_size = 10  # Pt
        self.shape_font_type = 'Calibri'
        self.shae_font_size = 6.0  # Pt
        self.shae_font_large_size = 8.0  # Pt
        self.tag_font_large_size = 4.0 # Pt

        self.between_shape_column = 0.5 #inches
        between_shape_row = 0.25  # inches
        self.between_l3if = 0.25  #inches

        l3_segment_up_down_offset = 0.25

        if self.click_value_l3 == 'L3-4-1': # inches  #changed at ver 2.3.3
            l3_segment_up_down_offset = 0.40

        min_between_line = 0.075  # inches
        min_shape_width = 1.0 #inches

        l3_segment_hight_ratio = 1.75 # ratio

        ''' 
        main loop 
        '''
        shape_left = self.left_margin
        shape_top = self.top_margin

        '''
        loop write device
        '''
        l3segment_line_array = []
        self.connected_l3if_key_array = []
        self.all_l3if_tag_array = []
        self.mark_multi_ip_array = []
        self.size_l3_instance_array = []
        self.used_l3segment_array = []
        self._used_l3segment_set = set()
        self._combined_device_set = None
        self.area_position_array = [9999.0,0.0,0.0,0.0,target_folder_name]  # shape_left, shape_top, shape_width, shape_hight, shape_text
        self.outline_position_array = [9999.0, 0.0, 0.0, 0.0]  # shape_left, shape_top, shape_width, shape_hight

        max_offset_x = 0.0
        end_l3_seg_inche_x = 0.0
        self.mark_wp_top = self.top_margin + top_offset
        flag_first_colmun = True

        ### ver 2.3.3 make self.target_offset_shape_array ###
        # get shape name in the folder and sort
        ws_name = 'Master_Data'
        ppt_meta_file = str(self.inFileTxt_11_1.get())
        # Cache POSITION_SHAPE read: l3_area_create is called multiple times per All Areas pass,
        # so avoid re-reading the same file on every invocation.
        _ori_ps_key = (ws_name, ppt_meta_file, '<<POSITION_SHAPE>>')
        if not hasattr(self, '_ori_pos_shape_cache'):
            self._ori_pos_shape_cache = {}
        if _ori_ps_key not in self._ori_pos_shape_cache:
            self._ori_pos_shape_cache[_ori_ps_key] = nsm_def.convert_master_to_array(
                ws_name, ppt_meta_file, '<<POSITION_SHAPE>>')
        self.ori_position_shape_array = self._ori_pos_shape_cache[_ori_ps_key]
        ori_position_shape_tuple = nsm_def.convert_array_to_tuple(self.ori_position_shape_array)
        self.shape_folder_tuple = nsm_def.get_shape_folder_tuple(ori_position_shape_tuple)
        # print(self.shape_folder_tuple)

        self.target_offset_shape_array = []

        def process_elements(array):
            target_array = []
            for element in array:
                if element[0] >= 2:
                    items = element[1]
                    for item in items[1:]:
                        if '_AIR_' not in item and '<END>' not in item:
                            target_array.append(item)
                            break  # Stop after appending the first valid item
            return target_array

        # Process the array and print the result
        self.target_offset_shape_array = process_elements(self.ori_position_shape_array)
        #print('--- self.target_offset_shape_array ---')
        #print(self.target_offset_shape_array)
        ################################################

        current_segment_count = 0  # add at ver 2.4.1
        for self.index_2,tmp_target_position_shape_array in enumerate(target_position_shape_array):
            start_l3_seg_inche_x = self.left_margin + left_offset
            sum_delta_y_grid = 0.0  # at ver 2.4.1

            ''' write device and wp(up/down)'''
            self.flag_area_equel_left = True
            self.second_area_offset = 0.0

            for tmp_tmp_target_position_shape_array in tmp_target_position_shape_array:

                ''' OFFSET Ver 2.3.3 '''
                if self.click_value_l3 == 'L3-4-1':
                    if action_type == 'CREATE' and tmp_tmp_target_position_shape_array not in '_AIR_' and self.flag_second_page == False:
                        shape_name = tmp_tmp_target_position_shape_array
                        left_offset += create_master_file_one_area.get_l3_shape_offset(self,shape_name ,left_offset)
                        # Add the per-device collision-avoidance shift computed
                        # by compute_l3_device_shifts(). The cumulative cascade
                        # within a row is handled automatically by the existing
                        # `left_offset += shape_width + between_shape_column`
                        # accumulation, so we only inject this device's own
                        # additional delta (0 for non-shifted devices).
                        if hasattr(self, 'device_extra_shift'):
                            left_offset += self.device_extra_shift.get(shape_name, 0.0)

                tmp_left_array = []
                tmp_right_array = []
                for i in new_wp_exist_array[2]:
                    tmp_left_array.append(i)
                for k in new_wp_exist_array[3]:
                    tmp_left_array.append(k)

                if tmp_tmp_target_position_shape_array not in tmp_left_array  and tmp_tmp_target_position_shape_array not in tmp_right_array and '_AIR_' not in tmp_tmp_target_position_shape_array:  # except left/right wp in writing pre device. If you need _AIR_ empty space, delete '_AIR_ not in tmp_tmp_target_position_shape_array'
                    shape_text  = tmp_tmp_target_position_shape_array
                    shape_type  = 'DEVICE_NORMAL'
                    shape_left  = self.left_margin + left_offset
                    # WayPoint X-snap override: in L3-4-1 All Areas 2nd pass
                    # CREATE, snap WayPoints in `_wp_`-only Area Position rows
                    # to their precomputed target X (center-aligned over the
                    # leftmost connected device). left_offset is intentionally
                    # left untouched so downstream shapes in the same row
                    # continue their natural cascade. See
                    # _compute_wp_x_snap_targets for selection criteria.
                    if (self.click_value_l3 == 'L3-4-1'
                            and action_type == 'CREATE'
                            and self.flag_re_create == True):
                        _wp_target = getattr(self, '_wp_x_target', None)
                        if _wp_target and shape_text in _wp_target:
                            shape_left = _wp_target[shape_text]
                    shape_top   = self.top_margin + top_offset
                    shape_width_hight_array = get_text_wh_cached(self,self.shae_font_large_size,tmp_tmp_target_position_shape_array)
                    shape_width = shape_width_hight_array[0]
                    shape_hight = shape_width_hight_array[1] * 5

                    ### Add wp change at ver 2.3.0 ####
                    if self.click_value_l3 == 'L3-4-1' and len(self.wp_list_array) != 0:
                        self.global_wp_array = copy.deepcopy(self.wp_list_array)
                        self.wp_list_array = []
                    ###################################

                    if shape_text in self.wp_list_array:
                        shape_type = 'WAY_POINT_NORMAL'

                    self.shape_width_if_array = get_shape_width_if_array(self,tmp_tmp_target_position_shape_array)  ### return_shape_width,tmp_up_array,tmp_down_array,full_ip_address_width_array
                    tmp_shpae_width = self.shape_width_if_array[0]
                    if shape_width < tmp_shpae_width:
                        shape_width = tmp_shpae_width

                    if min_shape_width > shape_width:
                        shape_width = min_shape_width

                    #check l3 instance exist
                    for tmp_l3_instance_array in self.l3_instance_array:
                        if tmp_l3_instance_array[0] == shape_text:
                            shape_type = 'DEVICE_L3_INSTANCE'

                    #write l3 instance
                    self.between_l3instance = min_between_line * 2
                    tmp_l3_add_shape_array = []

                    if shape_type == 'DEVICE_L3_INSTANCE':
                        tmp_l3_instance_array = []
                        for tmp_update_l3_instance_array in self.update_l3_instance_array:
                            if tmp_update_l3_instance_array[0][1] == shape_text:
                                tmp_l3_instance_array.append(tmp_update_l3_instance_array[1])

                        tmp_l3_instance_array = nsm_def.get_l2_broadcast_domains.get_unique_list(tmp_l3_instance_array)
                        offset_l3_instance = get_text_wh_cached(self,self.shae_font_size,shape_text)[0] + self.between_l3instance

                        calc_need_device_width = get_text_wh_cached(self,self.shae_font_size,shape_text)[0] + offset_l3_instance
                        for tmp_tmp_l3_instance_array in tmp_l3_instance_array:
                            l3_shape_text = tmp_tmp_l3_instance_array
                            l3_shape_width = get_text_wh_cached(self,self.shae_font_size,tmp_tmp_l3_instance_array)[0]
                            l3_shape_hight = get_text_wh_cached(self,self.shae_font_size,tmp_tmp_l3_instance_array)[1]

                            if (min_shape_width * 0.3) > l3_shape_width:
                                l3_shape_width = (min_shape_width * 0.3)

                            l3_shape_top = shape_top + (shape_hight * 0.5) - l3_shape_hight * 0.5
                            l3_shape_left = shape_left + offset_l3_instance
                            l3_shape_type = 'L3_INSTANCE'

                            calc_need_device_width += l3_shape_width + self.between_l3instance

                            if action_type == 'CREATE' and not minimal_1st_pass:
                                #self.shape = self.slide.shapes
                                #nsm_ddx_figure.extended.add_shape(self, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight, l3_shape_text)
                                tmp_l3_add_shape_array.append([l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight, l3_shape_text])

                            offset_l3_instance += l3_shape_width + self.between_l3instance
                            #self.size_l3_instance_array.append([shape_text, l3_shape_text, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight])

                        if calc_need_device_width > shape_width:
                            ### extend distance of shape

                            shape_width = calc_need_device_width
                            for tmp_tmp_l3_add_shape_array in tmp_l3_add_shape_array:
                                l3_shape_type = tmp_tmp_l3_add_shape_array[0]
                                l3_shape_left = tmp_tmp_l3_add_shape_array[1]
                                l3_shape_top = tmp_tmp_l3_add_shape_array[2]
                                l3_shape_width = tmp_tmp_l3_add_shape_array[3]
                                l3_shape_hight = tmp_tmp_l3_add_shape_array[4]
                                l3_shape_text = tmp_tmp_l3_add_shape_array[5]

                                # add at ver 2.4.1
                                if self.flag_re_create == True and self.flag_second_page == False and action_type == 'CREATE' and len(self.per_index2_before_array) > self.index_2:
                                    l3_shape_top -= (self.per_index2_before_array[self.index_2] - self.per_index2_after_array[self.index_2])

                                self.shape = self.slide.shapes
                                nsm_ddx_figure.extended.add_shape(self, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight, l3_shape_text)

                                self.size_l3_instance_array.append([shape_text, l3_shape_text, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight])

                        else:
                            ### not extend distance of shape
                            tmp_l3_instance_array = []
                            for tmp_update_l3_instance_array in self.update_l3_instance_array:
                                if tmp_update_l3_instance_array[0][1] == shape_text:
                                    tmp_l3_instance_array.append(tmp_update_l3_instance_array[1])

                            tmp_l3_instance_array = nsm_def.get_l2_broadcast_domains.get_unique_list(tmp_l3_instance_array)
                            tmp_plus_width = shape_width / (len(tmp_l3_instance_array) + 1)
                            offset_l3_instance = 0.0

                            for tmp_tmp_l3_instance_array in tmp_l3_instance_array:
                                l3_shape_text = tmp_tmp_l3_instance_array
                                l3_shape_width = get_text_wh_cached(self,self.shae_font_size, tmp_tmp_l3_instance_array)[0]
                                l3_shape_hight = get_text_wh_cached(self,self.shae_font_size, tmp_tmp_l3_instance_array)[1]

                                if (min_shape_width * 0.3) > l3_shape_width:
                                    l3_shape_width = (min_shape_width * 0.3)

                                l3_shape_top = shape_top + (shape_hight * 0.5) - l3_shape_hight * 0.5
                                l3_shape_left = shape_left + tmp_plus_width - l3_shape_width * 0.5 + offset_l3_instance
                                l3_shape_type = 'L3_INSTANCE'

                                if action_type == 'CREATE' and not minimal_1st_pass:
                                    # add at ver 2.4.1
                                    if self.flag_re_create == True and self.flag_second_page == False and len(self.per_index2_before_array) > self.index_2:
                                        l3_shape_top -= (self.per_index2_before_array[self.index_2] - self.per_index2_after_array[self.index_2])
                                    self.shape = self.slide.shapes
                                    nsm_ddx_figure.extended.add_shape(self, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight, l3_shape_text)

                                if self.click_value_l3 == 'L3-4-1':
                                    self.add_shape_array.append([shape_type, shape_left, shape_top, shape_width, shape_hight,shape_text])  # add ver 2.3.3

                                offset_l3_instance += tmp_plus_width
                                self.size_l3_instance_array.append([shape_text, l3_shape_text, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight])


                    if '_AIR_' not in shape_text:
                        if action_type == 'CREATE' and not minimal_1st_pass:
                            # add at ver 2.4.1
                            #print(self.index_2,self.per_index2_before_array,self.per_index2_after_array)
                            if self.flag_re_create == True and self.flag_second_page == False and len(self.per_index2_before_array) > self.index_2:
                                shape_top -= (self.per_index2_before_array[self.index_2] - self.per_index2_after_array[self.index_2])

                            self.shape = self.slide.shapes
                            nsm_ddx_figure.extended.add_shape(self, shape_type, shape_left, shape_top, shape_width, shape_hight, shape_text)
                            self.slide.shapes._spTree.remove(self.shape._element)  # move shape to back layer
                            self.slide.shapes._spTree.insert(2, self.shape._element)  # move shape to back layer

                        if self.click_value_l3 == 'L3-4-1':
                            self.add_shape_array.append([shape_type, shape_left, shape_top , shape_width, shape_hight, shape_text])  # add ver 2.3.3

                        # Capture per-device positions in 1st pass CREATE for
                        # the post-pass collision-avoidance shift analysis.
                        # Use the post-get_l3_shape_offset coordinates so that
                        # device positions, tag positions, and segment line
                        # natural_x are all in the same coordinate frame.
                        if self.click_value_l3 == 'L3-4-1' and self.flag_re_create == False and action_type == 'CREATE':
                            if shape_type in ('DEVICE_NORMAL', 'DEVICE_L3_INSTANCE', 'WAY_POINT_NORMAL'):
                                if not hasattr(self, '_collision_devs_per_area'):
                                    self._collision_devs_per_area = {}
                                area_devs = self._collision_devs_per_area.setdefault(target_folder_name, {})
                                if shape_text != '_AIR_' and shape_text not in area_devs:
                                    area_devs[shape_text] = {'left': shape_left, 'top': shape_top, 'width': shape_width, 'height': shape_hight}

                        '''GET Folder and Outline position'''
                        # get folder left
                        if (shape_type == 'DEVICE_L3_INSTANCE' or shape_type == 'DEVICE_NORMAL' or shape_type == 'WAY_POINT_NORMAL') and self.area_position_array[0] + self.between_shape_column  > shape_left:
                            self.area_position_array[0] = shape_left - self.between_shape_column
                            self.outline_position_array[0] = shape_left - self.between_shape_column * 2

                        # get folder top
                        if (shape_type == 'DEVICE_L3_INSTANCE' or shape_type == 'DEVICE_NORMAL') and flag_first_colmun == True:
                            self.area_position_array[1] = shape_top - between_shape_row
                            if new_wp_exist_array[0] == []:
                                # 0.04 = area_margin_y_top (0.34) - base area_margin_y (0.30):
                                # extend outline further up to preserve the gap with the area frame.
                                self.outline_position_array[1] = shape_top - between_shape_row * 2 - 0.04
                            flag_first_colmun = False

                        # get folder width
                        if (shape_type == 'DEVICE_L3_INSTANCE' or shape_type == 'DEVICE_NORMAL' or shape_type == 'WAY_POINT_NORMAL') and (self.outline_position_array[0] + self.outline_position_array[2]) < shape_left + shape_width:
                            self.outline_position_array[2] = shape_left + shape_width + self.between_shape_column * 2 - self.outline_position_array[0]
                        if (shape_type == 'DEVICE_L3_INSTANCE' or shape_type == 'DEVICE_NORMAL' ) and (self.area_position_array[0] + self.area_position_array[2]) < shape_left + shape_width:
                            self.area_position_array[2] = shape_left + shape_width + self.between_shape_column - self.area_position_array[0]

                        # get folder hight
                        if shape_type == 'DEVICE_L3_INSTANCE' or shape_type == 'DEVICE_NORMAL':
                            self.area_position_array[3] = shape_top + shape_hight + between_shape_row - self.area_position_array[1]
                            self.outline_position_array[3] = shape_top + shape_hight + between_shape_row * 3 - self.outline_position_array[1]

                        # adjust shape_hight if downside wp exist
                        if new_wp_exist_array[1] != []: #up/down/left/right
                            self.area_position_array[3] = (shape_top - between_shape_row) - self.area_position_array[1]


                        #adjust shape top of outline, if wp
                        if shape_type == 'WAY_POINT_NORMAL':
                            if shape_text in  new_wp_exist_array[0]:
                                self.outline_position_array[1] = shape_top - between_shape_row

                            if shape_text in new_wp_exist_array[1]:
                                self.outline_position_array[3] = (shape_top + shape_hight + between_shape_row) - self.outline_position_array[1]

                    left_offset += shape_width + self.between_shape_column

                    #print('--- self.size_l3_instance_array ---  shape_text, l3_shape_text, l3_shape_type, l3_shape_left, l3_shape_top, l3_shape_width, l3_shape_hight')
                    #print(self.size_l3_instance_array)

                    '''write l3 if '''

                    #print('### self.shape_width_if_array[1], self.shape_width_if_array[2]  ',self.shape_width_if_array[1], self.shape_width_if_array[2])
                    tag_up_offset_x = self.between_l3if
                    tag_down_offset_x = self.between_l3if
                    # Use index instead of full scan
                    device_l3_rows = self.l3_rows_by_device.get(shape_text, [])
                    for tmp_update_l3_table_array in device_l3_rows:
                            '''write upside l3 if'''
                            for up_shape_width_if_array in self.shape_width_if_array[1]:
                                if up_shape_width_if_array[0] == tmp_update_l3_table_array[2]:
                                    #print('##UP   ',shape_text,up_shape_width_if_array[1])
                                    shape_width_hight_array = get_text_wh_cached(self,self.tag_font_large_size, up_shape_width_if_array[1]) # width, hight
                                    tag_shape_type  = 'TAG_NORMAL'
                                    tag_shape_left  = shape_left + tag_up_offset_x
                                    tag_shape_top   = shape_top - shape_width_hight_array[1] * 0.5
                                    tag_shape_width = shape_width_hight_array[0]
                                    tag_shape_hight = shape_width_hight_array[1]
                                    tag_shape_text  = up_shape_width_if_array[1]

                                    if action_type == 'CREATE' and not minimal_1st_pass:
                                        self.shape = self.slide.shapes
                                        nsm_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text)

                                    self.all_l3if_tag_array.append([tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text,up_shape_width_if_array[0],shape_text])

                                    #reflect description ip address name distance-1
                                    flag_match_shape_width_if_array = False
                                    for tmp_shape_width_if_array in self.shape_width_if_array[3]:
                                        if [shape_text ,up_shape_width_if_array[0]] == tmp_shape_width_if_array[0]:
                                            tag_up_offset_x += tmp_shape_width_if_array[1]
                                            flag_match_shape_width_if_array = True
                                            tmp_add_width = tmp_shape_width_if_array[1]

                                    if flag_match_shape_width_if_array == False:
                                        tag_up_offset_x += tag_shape_width + self.between_l3if
                                        tmp_add_width = tag_shape_width + self.between_l3if

                                    '''mark ip address(up side)'''
                                    offset_ipaddress = 0.0 #inches
                                    key_if = (tmp_update_l3_table_array[1], tmp_update_l3_table_array[2])
                                    remake_array = self._l3_ipset_cache.get(key_if)
                                    if remake_array is None:
                                        try:
                                            remake_array = ast.literal_eval(tmp_update_l3_table_array[6])
                                        except (ValueError, SyntaxError):
                                            remake_array = []
                                        self._l3_ipset_cache[key_if] = remake_array

                                    if remake_array:
                                        #print('##mark ip address', len(remake_array),remake_array)
                                        for tmp_remake_array in remake_array:
                                            tag_shape_type = 'IP_ADDRESS_TAG'
                                            tag_ip_width = get_text_wh_cached(self,self.shae_font_size,tmp_remake_array[2])[0]
                                            self.mark_multi_ip_array.append([tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top - tag_shape_hight - offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2],tmp_remake_array,shape_text,len(remake_array),[shape_text,up_shape_width_if_array[0]],tag_shape_left ])

                                            offset_ipaddress += tag_shape_hight

                                            # reflect description ip address name distance-2
                                            #if tag_ip_width + tag_shape_width * 0.5 > tmp_add_width:
                                            #    tag_up_offset_x += (tag_ip_width + tag_shape_width * 0.5) - tmp_add_width


                            #write downside l3 if
                            for down_shape_width_if_array in self.shape_width_if_array[2]:
                                if down_shape_width_if_array[0] == tmp_update_l3_table_array[2]:
                                    #print('##DOWN ',shape_text,down_shape_width_if_array[1])
                                    shape_width_hight_array = get_text_wh_cached(self,self.tag_font_large_size, down_shape_width_if_array[1]) # width, hight
                                    tag_shape_type  = 'TAG_NORMAL'
                                    tag_shape_left  = shape_left + tag_down_offset_x
                                    tag_shape_top   = shape_top + shape_hight  - shape_width_hight_array[1] * 0.5
                                    tag_shape_width = shape_width_hight_array[0]
                                    tag_shape_hight = shape_width_hight_array[1]
                                    tag_shape_text  = down_shape_width_if_array[1]

                                    if action_type == 'CREATE' and not minimal_1st_pass:
                                        self.shape = self.slide.shapes
                                        nsm_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text)

                                    self.all_l3if_tag_array.append([tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text,down_shape_width_if_array[0],shape_text])

                                    # reflect description ip address name distance-1
                                    flag_match_shape_width_if_array = False
                                    for tmp_shape_width_if_array in self.shape_width_if_array[3]:
                                        if [shape_text, down_shape_width_if_array[0]] == tmp_shape_width_if_array[0]:
                                            tag_down_offset_x += tmp_shape_width_if_array[1]
                                            flag_match_shape_width_if_array = True
                                            tmp_add_width = tmp_shape_width_if_array[1]

                                    if flag_match_shape_width_if_array == False:
                                        tag_down_offset_x += tag_shape_width + self.between_l3if
                                        tmp_add_width = tag_shape_width + self.between_l3if

                                    '''mark ip address(down side)'''
                                    offset_ipaddress = 0.0 #inches
                                    key_if = (tmp_update_l3_table_array[1], tmp_update_l3_table_array[2])
                                    remake_array = self._l3_ipset_cache.get(key_if)
                                    if remake_array is None:
                                        try:
                                            remake_array = ast.literal_eval(tmp_update_l3_table_array[6])
                                        except (ValueError, SyntaxError):
                                            remake_array = []
                                        self._l3_ipset_cache[key_if] = remake_array

                                    if remake_array:
                                        #print('##mark ip address', len(remake_array),remake_array)
                                        for tmp_remake_array in remake_array:
                                            tag_shape_type = 'IP_ADDRESS_TAG'
                                            tag_ip_width = get_text_wh_cached(self,self.shae_font_size,tmp_remake_array[2])[0]
                                            self.mark_multi_ip_array.append([tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top + tag_shape_hight + offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2],tmp_remake_array,shape_text,len(remake_array),[shape_text,down_shape_width_if_array[0]],tag_shape_left])

                                            offset_ipaddress += tag_shape_hight

                                            # reflect description ip address name distance-2
                                            #if tag_ip_width + tag_shape_width * 0.5 > tmp_add_width:
                                            #    tag_up_offset_x += (tag_ip_width + tag_shape_width * 0.5) - tmp_add_width

            ''' pre mark wp(left/right)'''
            if (new_wp_exist_array[2] != [] or new_wp_exist_array[3] != []) and left_right_wp_row_num == self.index_2:
                self.mark_wp_top = self.top_margin + top_offset

            '''add top_offset for folder outline'''
            #upside
            if shape_text in self.wp_list_array:
                top_offset += between_shape_row * 2

            #downside
            if len(target_position_shape_array) > self.index_2 + 1:
                for down_target_position_shape_array in target_position_shape_array[self.index_2 + 1]:
                    if down_target_position_shape_array in self.wp_list_array and len(target_position_shape_array) == self.index_2 +2:
                        top_offset += between_shape_row * 2
                        break

            ### check end l3 segment point
            if end_l3_seg_inche_x < shape_left + shape_width:
                end_l3_seg_inche_x = shape_left + shape_width

            ### write broadcast domain line
            top_device_name_array = tmp_target_position_shape_array

            #get count l3_segment
            return_get_l3_segment_num = get_l3_segment_num(self,top_device_name_array ,target_position_shape_array)
            count_l3segment = return_get_l3_segment_num[0]
            self.connected_l3if_key_array.append(return_get_l3_segment_num[1])
            #print('##  return_get_l3_segment_num ',return_get_l3_segment_num ,top_device_name_array )
            tmp_l3segment_y_array = []
            if count_l3segment != 0:
                tmp_line_offset = l3_segment_up_down_offset
                #current_segment_count = 0
                for tmp_count_l3segment in range(count_l3segment):
                    #add distance upside or downside of device shape
                    if tmp_count_l3segment == 0:
                        top_offset += between_shape_row + l3_segment_up_down_offset
                    else:
                        top_offset += between_shape_row
                    if tmp_count_l3segment +1 == count_l3segment:
                        top_offset += l3_segment_up_down_offset

                    ### add at ver 2.4.1 ###
                    if self.flag_re_create == True and self.flag_second_page == False:
                        #print('self.y_grid_segment_array,self.optimize_y_grid_array  ---> ', self.y_grid_segment_array[current_segment_count][1][1],self.optimize_y_grid_array[current_segment_count][1][1])
                        delta_y_grid = self.y_grid_segment_array[current_segment_count][1][1] - self.optimize_y_grid_array[current_segment_count][1][1] - sum_delta_y_grid
                        sum_delta_y_grid += delta_y_grid
                        if delta_y_grid != 0:
                            #print('delta_y_grid -->', delta_y_grid)
                            tmp_line_offset -= delta_y_grid
                        current_segment_count += 1

                    tmp_l3segment_y_array.append(shape_top + tmp_line_offset  + shape_hight + between_shape_row)
                    tmp_line_offset += between_shape_row

            l3segment_line_array.append([[start_l3_seg_inche_x,end_l3_seg_inche_x],tmp_l3segment_y_array,return_get_l3_segment_num[1]])

            ### add at ver 2.4.1
            if self.flag_re_create == True and self.flag_second_page == False and tmp_l3segment_y_array != [] and action_type == 'CREATE':
                self.per_index2_after_array.append(max(tmp_l3segment_y_array))
                #print(self.per_index2_before_array)
                #print(self.per_index2_after_array)

            elif self.flag_re_create == False and self.flag_second_page == False and tmp_l3segment_y_array != [] and action_type == 'CREATE':
                self.per_index2_before_array.append(max(tmp_l3segment_y_array))

            elif self.flag_re_create == True and self.flag_second_page == False and tmp_l3segment_y_array == [] and action_type == 'CREATE' and self.per_index2_after_array != []:
                self.per_index2_after_array.append(self.per_index2_after_array[-1] + 0.5) #0.5 is the offset value in the case of WP; this offset value is uniformly set to 0.5 because it is difficult to determine whether it is WP or not.

            elif self.flag_re_create == False and self.flag_second_page == False and tmp_l3segment_y_array == [] and action_type == 'CREATE' and self.per_index2_before_array != []:
                self.per_index2_before_array.append(self.per_index2_before_array[-1])

            '''change offset  check_move_to_right '''
            top_offset += shape_hight + between_shape_row

            now_offset_x = end_l3_seg_inche_x - self.left_margin
            if max_offset_x < now_offset_x :
                max_offset_x = now_offset_x

            _move_right = check_move_to_right(self,top_device_name_array,target_position_shape_array)
            if _move_right:
                left_offset = max_offset_x + self.left_margin # add 1.0 at ver 2.3.4
            else:
                left_offset = start_l3_seg_inche_x - self.left_margin


        #print('## end_l3_seg_inche_x ', end_l3_seg_inche_x)
        #print('--- l3segment_line_array ---')
        #print(l3segment_line_array)


        ''' 
        write wp(left/right)
        '''
        shape_top = self.mark_wp_top
        shape_width_hight_array = get_text_wh_cached(self,self.shae_font_large_size, str(new_wp_exist_array[2]))
        shape_width = shape_width_hight_array[0]
        shape_hight = shape_width_hight_array[1] * 5


        ''' write wp left'''
        offset_shape_left = 0.0
        if new_wp_exist_array[2] != []:
            for tmp_i in new_wp_exist_array[2]:
                shape_text = tmp_i
                shape_type = 'WAY_POINT_NORMAL'

                self.shape_width_if_array = get_shape_width_if_array(self, shape_text)  # return_shape_width, tmp_up_array, tmp_down_array
                tmp_shpae_width = self.shape_width_if_array[0]
                if shape_width < tmp_shpae_width:
                    shape_width = tmp_shpae_width

                if min_shape_width > shape_width:
                    shape_width = min_shape_width

                #print('### WRITE LEFT WP  ', shape_text, shape_width,new_wp_exist_array[2])
                shape_left = self.left_margin - shape_width - self.between_shape_column * 3 + offset_shape_left

                if action_type == 'CREATE' and not minimal_1st_pass:
                    self.shape = self.slide.shapes
                    nsm_ddx_figure.extended.add_shape(self, shape_type, shape_left, shape_top, shape_width, shape_hight, shape_text)
                    if self.click_value_l3 == 'L3-4-1':
                        self.add_shape_array.append([shape_type, shape_left, shape_top, shape_width, shape_hight, shape_text])  # add ver 2.3.3

                #get left side folder and outline point
                self.area_position_array[0] = shape_left + shape_width + self.between_shape_column * 2 - offset_shape_left

                self.outline_position_array[0] = shape_left - self.between_shape_column * 2
                self.outline_position_array[2] += shape_width + self.between_shape_column * 3

                '''write wp_left l3 if '''
                # print('### self.shape_width_if_array[1], self.shape_width_if_array[2]  ',self.shape_width_if_array[1], self.shape_width_if_array[2])
                tag_up_offset_x = self.between_l3if
                tag_down_offset_x = self.between_l3if

                # Use index instead of full scan
                device_l3_rows = self.l3_rows_by_device.get(shape_text, [])
                for tmp_update_l3_table_array in device_l3_rows:
                        for up_shape_width_if_array in self.shape_width_if_array[1]:
                            #write up side l3 if
                            if up_shape_width_if_array[0] == tmp_update_l3_table_array[2]:
                                # print('##UP   ',shape_text,up_shape_width_if_array[1])
                                shape_width_hight_array = get_text_wh_cached(self,self.tag_font_large_size, up_shape_width_if_array[1])  # width, hight
                                tag_shape_type = 'TAG_NORMAL'
                                tag_shape_left = shape_left + tag_up_offset_x
                                tag_shape_top = shape_top - shape_width_hight_array[1] * 0.5
                                tag_shape_width = shape_width_hight_array[0]
                                tag_shape_hight = shape_width_hight_array[1]
                                tag_shape_text = up_shape_width_if_array[1]

                                if action_type == 'CREATE' and not minimal_1st_pass:
                                    self.shape = self.slide.shapes
                                    nsm_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text)

                                self.all_l3if_tag_array.append([tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text,up_shape_width_if_array[0],shape_text])

                                # reflect description ip address name distance-1
                                flag_match_shape_width_if_array = False
                                for tmp_shape_width_if_array in self.shape_width_if_array[3]:
                                    if [shape_text, up_shape_width_if_array[0]] == tmp_shape_width_if_array[0]:
                                        tag_up_offset_x += tmp_shape_width_if_array[1]
                                        flag_match_shape_width_if_array = True
                                        tmp_add_width = tmp_shape_width_if_array[1]

                                if flag_match_shape_width_if_array == False:
                                    tag_up_offset_x += tag_shape_width + self.between_l3if
                                    tmp_add_width = tag_shape_width + self.between_l3if


                                '''wp_left write ip address(up side)'''
                                offset_ipaddress = 0.0  # inches
                                key_if = (tmp_update_l3_table_array[1], tmp_update_l3_table_array[2])
                                remake_array = self._l3_ipset_cache.get(key_if)
                                if remake_array is None:
                                    try:
                                        remake_array = ast.literal_eval(tmp_update_l3_table_array[6])
                                    except (ValueError, SyntaxError):
                                        remake_array = []
                                    self._l3_ipset_cache[key_if] = remake_array

                                if remake_array:
                                    for tmp_remake_array in remake_array:
                                        #print('##wp_left write ip address', tmp_remake_array)
                                        tag_shape_type = 'IP_ADDRESS_TAG'
                                        tag_ip_width = get_text_wh_cached(self,self.shae_font_size, tmp_remake_array[2])[0]
                                        #self.shape = self.slide.shapes

                                        self.mark_multi_ip_array.append([tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top - tag_shape_hight - offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2], tmp_remake_array, shape_text, len(remake_array), [shape_text, up_shape_width_if_array[0]], tag_shape_left])
                                        offset_ipaddress += tag_shape_hight

                                        # reflect description ip address name distance-2
                                        #if tag_ip_width + tag_shape_width * 0.5 > tmp_add_width:
                                        #    tag_up_offset_x += (tag_ip_width + tag_shape_width * 0.5) - tmp_add_width

                        # write down side l3 if
                        for down_shape_width_if_array in self.shape_width_if_array[2]:
                            if down_shape_width_if_array[0] == tmp_update_l3_table_array[2]:
                                #print('##DOWN ',shape_text,down_shape_width_if_array[1])
                                shape_width_hight_array = get_text_wh_cached(self,self.tag_font_large_size, down_shape_width_if_array[1]) # width, hight
                                tag_shape_type  = 'TAG_NORMAL'
                                tag_shape_left  = shape_left + tag_down_offset_x
                                tag_shape_top   = shape_top + shape_hight  - shape_width_hight_array[1] * 0.5
                                tag_shape_width = shape_width_hight_array[0]
                                tag_shape_hight = shape_width_hight_array[1]
                                tag_shape_text  = down_shape_width_if_array[1]

                                if action_type == 'CREATE' and not minimal_1st_pass:
                                    self.shape = self.slide.shapes
                                    nsm_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text)

                                self.all_l3if_tag_array.append([tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text,down_shape_width_if_array[0],shape_text])

                                # reflect description ip address name distance-1
                                flag_match_shape_width_if_array = False
                                for tmp_shape_width_if_array in self.shape_width_if_array[3]:
                                    if [shape_text, down_shape_width_if_array[0]] == tmp_shape_width_if_array[0]:
                                        tag_up_offset_x += tmp_shape_width_if_array[1]
                                        flag_match_shape_width_if_array = True
                                        tmp_add_width = tmp_shape_width_if_array[1]

                                if flag_match_shape_width_if_array == False:
                                    tag_down_offset_x += tag_shape_width + self.between_l3if
                                    tmp_add_width = tag_shape_width + self.between_l3if

                                '''wp_left write ip address(down side)'''
                                offset_ipaddress = 0.0  # inches
                                key_if = (tmp_update_l3_table_array[1], tmp_update_l3_table_array[2])
                                remake_array = self._l3_ipset_cache.get(key_if)
                                if remake_array is None:
                                    try:
                                        remake_array = ast.literal_eval(tmp_update_l3_table_array[6])
                                    except (ValueError, SyntaxError):
                                        remake_array = []
                                    self._l3_ipset_cache[key_if] = remake_array

                                if remake_array:
                                    for tmp_remake_array in remake_array:
                                        #print('##wp_left write ip address', tmp_remake_array)
                                        tag_shape_type = 'IP_ADDRESS_TAG'
                                        tag_ip_width = get_text_wh_cached(self,self.shae_font_size, tmp_remake_array[2])[0]

                                        #self.shape = self.slide.shapes
                                        self.mark_multi_ip_array.append([tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top + tag_shape_hight + offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2], tmp_remake_array, shape_text, len(remake_array), [shape_text, down_shape_width_if_array[0]],tag_shape_left])

                                        offset_ipaddress += tag_shape_hight

                                        # reflect description ip address name distance-2
                                        #if tag_ip_width + tag_shape_width * 0.5 > tmp_add_width:
                                        #    tag_up_offset_x += (tag_ip_width + tag_shape_width * 0.5) - tmp_add_width

                offset_shape_left = offset_shape_left - shape_width - self.between_shape_column

        ''' write wp right'''
        offset_shape_right = 0.0
        if new_wp_exist_array[3] != []:

            for tmp_i in reversed(new_wp_exist_array[3]):
                shape_text = tmp_i

                shape_type = 'WAY_POINT_NORMAL'
                self.shape_width_if_array = get_shape_width_if_array(self, shape_text)  # return_shape_width, tmp_up_array, tmp_down_array

                tmp_shpae_width = self.shape_width_if_array[0]

                if shape_width < tmp_shpae_width:
                    shape_width = tmp_shpae_width

                if min_shape_width > shape_width:
                    shape_width = min_shape_width

                #print('### WRITE RIGHT WP  ', new_wp_exist_array[3], shape_width)
                #shape_left = end_l3_seg_inche_x + self.between_shape_column * 3
                shape_left = self.area_position_array[0] + self.area_position_array[2] + self.between_shape_column * 2  + offset_shape_right # updated
                if action_type == 'CREATE' and not minimal_1st_pass:
                    self.shape = self.slide.shapes
                    nsm_ddx_figure.extended.add_shape(self, shape_type, shape_left, shape_top, shape_width, shape_hight, shape_text)

                if self.click_value_l3 == 'L3-4-1':
                    self.add_shape_array.append([shape_type, shape_left, shape_top, shape_width, shape_hight, shape_text])  # add ver 2.3.3

                #get left side folder and outline point
                self.outline_position_array[2] += shape_width + self.between_shape_column * 3

                '''write wp_right l3 if '''
                tag_up_offset_x = self.between_l3if
                tag_down_offset_x = self.between_l3if

                # Use index instead of full scan
                device_l3_rows = self.l3_rows_by_device.get(shape_text, [])
                for tmp_update_l3_table_array in device_l3_rows:
                        for up_shape_width_if_array in self.shape_width_if_array[1]:
                            # write up side l3 if
                            if up_shape_width_if_array[0] == tmp_update_l3_table_array[2]:
                                # print('##UP   ',shape_text,up_shape_width_if_array[1])
                                shape_width_hight_array = get_text_wh_cached(self,self.tag_font_large_size, up_shape_width_if_array[1])  # width, hight
                                tag_shape_type = 'TAG_NORMAL'
                                tag_shape_left = shape_left + tag_up_offset_x
                                tag_shape_top = shape_top - shape_width_hight_array[1] * 0.5
                                tag_shape_width = shape_width_hight_array[0]
                                tag_shape_hight = shape_width_hight_array[1]
                                tag_shape_text = up_shape_width_if_array[1]

                                if action_type == 'CREATE' and not minimal_1st_pass:
                                    self.shape = self.slide.shapes
                                    nsm_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text)

                                self.all_l3if_tag_array.append([tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text,up_shape_width_if_array[0],shape_text])

                                # reflect description ip address name distance-1
                                flag_match_shape_width_if_array = False
                                for tmp_shape_width_if_array in self.shape_width_if_array[3]:
                                    if [shape_text, up_shape_width_if_array[0]] == tmp_shape_width_if_array[0]:
                                        tag_up_offset_x += tmp_shape_width_if_array[1]
                                        flag_match_shape_width_if_array = True
                                        tmp_add_width = tmp_shape_width_if_array[1]

                                if flag_match_shape_width_if_array == False:
                                    tag_up_offset_x += tag_shape_width + self.between_l3if
                                    tmp_add_width = tag_shape_width + self.between_l3if

                                #tag_up_offset_x += tag_shape_width + self.between_l3if

                                '''wp_right write ip address(up side)'''
                                offset_ipaddress = 0.0  # inches
                                key_if = (tmp_update_l3_table_array[1], tmp_update_l3_table_array[2])
                                remake_array = self._l3_ipset_cache.get(key_if)
                                if remake_array is None:
                                    try:
                                        remake_array = ast.literal_eval(tmp_update_l3_table_array[6])
                                    except (ValueError, SyntaxError):
                                        remake_array = []
                                    self._l3_ipset_cache[key_if] = remake_array

                                if remake_array:
                                    for tmp_remake_array in remake_array:
                                        #print('##wp_right write ip address', tmp_remake_array)
                                        tag_shape_type = 'IP_ADDRESS_TAG'
                                        tag_ip_width = get_text_wh_cached(self,self.shae_font_size, tmp_remake_array[2])[0]
                                        #self.shape = self.slide.shapes
                                        #nsm_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top - tag_shape_hight - offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_tmp_update_l3_table_array[2])
                                        self.mark_multi_ip_array.append([tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top - tag_shape_hight - offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2], tmp_remake_array, shape_text, len(remake_array), [shape_text, up_shape_width_if_array[0]],tag_shape_left])

                                        offset_ipaddress += tag_shape_hight

                                        # reflect description ip address name distance-2
                                        #if tag_ip_width + tag_shape_width * 0.5 > tmp_add_width:
                                        #    tag_up_offset_x += (tag_ip_width + tag_shape_width * 0.5) - tmp_add_width


                        # wite down side l3 if
                        for down_shape_width_if_array in self.shape_width_if_array[2]:
                            if down_shape_width_if_array[0] == tmp_update_l3_table_array[2]:
                                shape_width_hight_array = get_text_wh_cached(self,self.tag_font_large_size, down_shape_width_if_array[1])  # width, hight
                                tag_shape_type = 'TAG_NORMAL'
                                tag_shape_left = shape_left + tag_down_offset_x
                                tag_shape_top = shape_top + shape_hight - shape_width_hight_array[1] * 0.5
                                tag_shape_width = shape_width_hight_array[0]
                                tag_shape_hight = shape_width_hight_array[1]
                                tag_shape_text = down_shape_width_if_array[1]

                                if action_type == 'CREATE' and not minimal_1st_pass:
                                    self.shape = self.slide.shapes
                                    nsm_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text)

                                self.all_l3if_tag_array.append([tag_shape_type, tag_shape_left, tag_shape_top, tag_shape_width, tag_shape_hight, tag_shape_text,down_shape_width_if_array[0],shape_text])

                                # reflect description ip address name distance-1
                                flag_match_shape_width_if_array = False
                                for tmp_shape_width_if_array in self.shape_width_if_array[3]:
                                    if [shape_text, down_shape_width_if_array[0]] == tmp_shape_width_if_array[0]:
                                        tag_up_offset_x += tmp_shape_width_if_array[1]
                                        flag_match_shape_width_if_array = True
                                        tmp_add_width = tmp_shape_width_if_array[1]

                                if flag_match_shape_width_if_array == False:
                                    tag_down_offset_x += tag_shape_width + self.between_l3if
                                    tmp_add_width = tag_shape_width + self.between_l3if

                                #tag_down_offset_x += tag_shape_width + self.between_l3if

                                '''wp_right write ip address(down side)'''
                                offset_ipaddress = 0.0  # inches
                                key_if = (tmp_update_l3_table_array[1], tmp_update_l3_table_array[2])
                                remake_array = self._l3_ipset_cache.get(key_if)
                                if remake_array is None:
                                    try:
                                        remake_array = ast.literal_eval(tmp_update_l3_table_array[6])
                                    except (ValueError, SyntaxError):
                                        remake_array = []
                                    self._l3_ipset_cache[key_if] = remake_array

                                if remake_array:
                                    for tmp_remake_array in remake_array:
                                        #print('##wp_right write ip address', tmp_remake_array)
                                        tag_shape_type = 'IP_ADDRESS_TAG'
                                        tag_ip_width = get_text_wh_cached(self,self.shae_font_size, tmp_remake_array[2])[0]
                                        #self.shape = self.slide.shapes
                                        #nsm_ddx_figure.extended.add_shape(self, tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top + tag_shape_hight + offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_tmp_update_l3_table_array[2])
                                        self.mark_multi_ip_array.append( [tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top + tag_shape_hight + offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2], tmp_remake_array, shape_text, len(remake_array), [shape_text, down_shape_width_if_array[0]],tag_shape_left])

                                        offset_ipaddress += tag_shape_hight

                                        # reflect description ip address name distance-2
                                        if tag_ip_width + tag_shape_width * 0.5 > tmp_add_width:
                                            tag_up_offset_x += (tag_ip_width + tag_shape_width * 0.5) - tmp_add_width

                offset_shape_right = offset_shape_right + shape_width + self.between_shape_column

        '''
        write folder line and outline 
        '''
        #print('--- self.area_position_array ---')
        #print(self.area_position_array)

        ### write folder
        folder_shape_type = 'FOLDER_NORMAL'
        folder_shape_left = self.area_position_array[0]
        folder_shape_top = self.area_position_array[1]
        folder_shape_width = self.area_position_array[2]
        folder_shape_hight = self.area_position_array[3]
        folder_shape_text = self.area_position_array[4]
        max_folder_left_width = 0.0

        if action_type == 'CREATE' and self.click_value_l3 != 'L3-4-1':
            self.shape = self.slide.shapes
            nsm_ddx_figure.extended.add_shape(self, folder_shape_type, folder_shape_left, folder_shape_top, folder_shape_width, folder_shape_hight, folder_shape_text)
            self.slide.shapes._spTree.remove(self.shape._element)  # move shape to back layer
            self.slide.shapes._spTree.insert(2, self.shape._element)  # move shape to back layer

        elif action_type == 'CREATE' and self.click_value_l3 == 'L3-4-1' and self.flag_second_page == False:
            '''write folder line when l3 all areas'''
            #print(self.add_shape_write_array )    #shape_type, shape_left, shape_top, shape_width, shape_hight, shape_text
            #print(self.unique_area_device_array)

            area_groups = {}

            # Group the shapes by area_name
            for shape in self.add_shape_write_array:
                shape_area_name = None

                # Find the corresponding area_name for the shape's device_name
                for area_device in self.unique_area_device_array:
                    if shape[5] == area_device[1]:  # shape[5] corresponds to the device_name
                        shape_area_name = area_device[0]  # area_name is in area_device[0]
                        break

                if shape_area_name:
                    if shape_area_name not in area_groups:
                        area_groups[shape_area_name] = []
                    area_groups[shape_area_name].append(shape)

            # Now, for each area_name, calculate the new values
            area_outline_array = []

            for area_name, shapes in area_groups.items():
                # Initialize min/max variables
                min_shape_left = float('inf')
                min_shape_top = float('inf')
                max_right_edge = float('-inf')  # This will store the maximum (shape_left + shape_width)
                max_bottom_edge = float('-inf')  # This will store the maximum (shape_top + shape_height)

                # Loop through all shapes in the area and find the new min/max values
                for shape in shapes:
                    min_shape_left = min(min_shape_left, shape[1])  # shape[1] is shape_left
                    min_shape_top = min(min_shape_top, shape[2])  # shape[2] is shape_top
                    max_right_edge = max(max_right_edge, shape[1] + shape[3])  # shape[1] + shape[3] is the rightmost edge
                    max_bottom_edge = max(max_bottom_edge, shape[2] + shape[4])  # shape[2] + shape[4] is the bottommost edge

                # Calculate the final output values based on the new format
                width_difference = max_right_edge - min_shape_left  # Right edge minus left edge
                height_difference = max_bottom_edge - min_shape_top  # Bottom edge minus top edge

                # Append the results for this area
                area_outline_array.append([area_name, min_shape_left, min_shape_top, width_difference, height_difference])

            #print(area_outline_array)

            # write the area outline
            area_margin_x = 0.5 #inchi
            area_margin_y = 0.15 #inchi (bottom)
            area_margin_y_top = 0.34 #inchi (top; minimum that keeps Branch-2/IP gap >= 2px with label_offset=0.055)

            for tmp_area_outline_array in area_outline_array:
                if '_wp_' not in tmp_area_outline_array[0]:
                    self.shape = self.slide.shapes
                    folder_shape_left = tmp_area_outline_array[1] - area_margin_x
                    folder_shape_top = tmp_area_outline_array[2] - area_margin_y_top
                    folder_shape_width = tmp_area_outline_array[3] + (area_margin_x * 2)
                    folder_shape_hight = tmp_area_outline_array[4] + (area_margin_y_top + area_margin_y)
                    folder_shape_text = tmp_area_outline_array[0]

                    nsm_ddx_figure.extended.add_shape(self, folder_shape_type, folder_shape_left, folder_shape_top,folder_shape_width, folder_shape_hight, folder_shape_text)
                    self.slide.shapes._spTree.remove(self.shape._element)  # move shape to back layer
                    self.slide.shapes._spTree.insert(2, self.shape._element)  # move shape to back layer

                    if max_folder_left_width < (folder_shape_left + folder_shape_width):
                        max_folder_left_width = folder_shape_left + folder_shape_width

        ### write outline
        outline_shape_type = 'OUTLINE_NORMAL'
        outline_shape_left = self.outline_position_array[0]
        outline_shape_top = self.outline_position_array[1]
        outline_shape_width = self.outline_position_array[2]
        outline_shape_hight = self.outline_position_array[3]
        outline_shape_text = ''

        ### Kyuusai if len(target_position_shape_array) == 1
        outline_shape_hight += 0.2

        '''
        loop write l3segment
        '''
        all_l3_netowrk_list =[]
        all_l3segment_l3_netowrk_list = []
        self.all_written_if_line_array = []
        self.all_written_line_position_array = []
        self._written_line_x_buckets = {}
        multi_ip_address_num = 1
        for index_5,tmp_l3segment_line_array in enumerate(l3segment_line_array):
            #print(index_5,tmp_l3segment_line_array)
            if len(l3segment_line_array) > index_5 + 1:
                if l3segment_line_array[index_5][0][0] < l3segment_line_array[index_5 + 1][0][0]:
                    start_x = l3segment_line_array[index_5][0][0]
                else:
                    start_x = l3segment_line_array[index_5 + 1][0][0]

                if l3segment_line_array[index_5][0][1] > l3segment_line_array[index_5 + 1][0][1]:
                    end_x = l3segment_line_array[index_5][0][1]
                else:
                    end_x = l3segment_line_array[index_5 + 1][0][1]

                for index_55,tmp_tmp_l3segment_line_array in enumerate(tmp_l3segment_line_array[1]) :
                    '''pre check L3 SEGMENT'''
                    line_type = 'L3_SEGMENT'
                    inche_from_connect_x = start_x
                    inche_from_connect_y = tmp_tmp_l3segment_line_array
                    inche_to_connect_x = end_x
                    inche_to_connect_y = tmp_tmp_l3segment_line_array
                    #print(line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)
                    #nsm_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                    shape_text = 'xxx.xxx.xxx.xxx/xx'
                    shape_type = 'L3_SEGMENT_GRAY'
                    shape_width = inche_to_connect_x - inche_from_connect_x
                    shape_hight = get_text_wh_cached(self,self.tag_font_large_size, shape_text)[1] * l3_segment_hight_ratio  # l3 segment hight ratio
                    shape_left = inche_from_connect_x
                    shape_top = inche_from_connect_y - shape_hight * 0.5

                    #self.shape = self.slide.shapes
                    #nsm_ddx_figure.extended.add_shape(self, shape_type, shape_left, shape_top, shape_width, shape_hight, shape_text)
                    #self.slide.shapes._spTree.remove(self.shape._element)  # move shape to back layer
                    #self.slide.shapes._spTree.insert(2, self.shape._element)  # move shape to back layer

                    '''write lines between device and l3 segment'''
                    # identify key l3 if for current broadcast domain
                    #print('##self.connected_l3if_key_array[index_5][index_55]  ', self.connected_l3if_key_array[index_5][index_55])

                    ###get l3IF of the broadcast domain

                    edge_left_x = 999.0   # inches
                    edge_right_x = -999.0  # inches

                    l3segment_edge_array = []
                    l3_network_list = []
                    the_l3segment_l3_network_list = []

                    for tmp_target_l2_broadcast_group_array in self.target_l2_broadcast_group_array:
                        if self.connected_l3if_key_array[index_5][index_55] in tmp_target_l2_broadcast_group_array[1]:
                            #print('##MATCH  ',self.connected_l3if_key_array[index_5][index_55])
                            #print('##Target IFs  ', len(tmp_target_l2_broadcast_group_array[1]),tmp_target_l2_broadcast_group_array[1])
                            for tmp_all_l3if_tag_array in self.all_l3if_tag_array:
                                #print([tmp_all_l3if_tag_array[7],tmp_all_l3if_tag_array[6]] , tmp_target_l2_broadcast_group_array[1])
                                if [tmp_all_l3if_tag_array[7],tmp_all_l3if_tag_array[6]] in tmp_target_l2_broadcast_group_array[1]:
                                    #print('## Write line ', tmp_all_l3if_tag_array)

                                    line_type = 'L3_SEGMENT-L3IF'
                                    inche_from_connect_x = tmp_all_l3if_tag_array[1] + tmp_all_l3if_tag_array[3] * 0.5
                                    inche_from_connect_y = tmp_all_l3if_tag_array[2] + tmp_all_l3if_tag_array[4]
                                    inche_to_connect_x = tmp_all_l3if_tag_array[1] + tmp_all_l3if_tag_array[3] * 0.5
                                    inche_to_connect_y = tmp_tmp_l3segment_line_array - shape_hight * 0.5

                                    #for up
                                    if [tmp_all_l3if_tag_array[7],tmp_all_l3if_tag_array[6]] in self.up_down_l3if_array[0]:
                                        inche_from_connect_y = tmp_all_l3if_tag_array[2]
                                        inche_to_connect_y = tmp_tmp_l3segment_line_array + shape_hight * 0.5

                                    '''check to exist the near vertical line. if there is near line, make between size to about 0.075 inches'''

                                    def _norm_y(y1, y2):
                                        return (y1, y2) if y1 <= y2 else (y2, y1)

                                    def _overlap_y(a0, a1, b0, b1):
                                        # overlap if max(starts) <= min(ends)
                                        return max(a0, b0) <= min(a1, b1)

                                    def _is_nearly_same_line(x1, y1_start, y1_end, x2, y2_start, y2_end, tolerance=0.01):
                                        """
                                        Determines if two lines are at nearly the same position
                                        Returns True if both X-axis distance and Y-axis endpoints are within tolerance
                                        """
                                        y1_0, y1_1 = _norm_y(y1_start, y1_end)
                                        y2_0, y2_1 = _norm_y(y2_start, y2_end)

                                        x_distance = abs(x1 - x2)
                                        y_start_distance = abs(y1_0 - y2_0)
                                        y_end_distance = abs(y1_1 - y2_1)

                                        return (x_distance <= tolerance and
                                                y_start_distance <= tolerance and
                                                y_end_distance <= tolerance)

                                    def _has_conflict_bucketed(x_candidate, new_y0, new_y1, buckets, bucket_size,
                                                      current_x, current_y_start, current_y_end):
                                        bk = int(x_candidate / bucket_size)
                                        for b in (bk - 1, bk, bk + 1):
                                            for tmp_line in buckets.get(b, []):
                                                old_x = tmp_line[1]
                                                old_y0, old_y1 = _norm_y(tmp_line[2], tmp_line[4])
                                                if _is_nearly_same_line(current_x, current_y_start, current_y_end,
                                                                        old_x, old_y0, old_y1):
                                                    continue
                                                if not (x_candidate - bucket_size < old_x < x_candidate + bucket_size):
                                                    continue
                                                if _overlap_y(new_y0, new_y1, old_y0, old_y1):
                                                    return True
                                        return False

                                    new_y0, new_y1 = _norm_y(inche_from_connect_y, inche_to_connect_y)
                                    x0 = inche_from_connect_x
                                    step = min_between_line

                                    candidates = [x0, x0 + step, x0 - step]

                                    chosen_x = x0
                                    for x_candidate in candidates:
                                        if not _has_conflict_bucketed(
                                                x_candidate, new_y0, new_y1,
                                                self._written_line_x_buckets,
                                                min_between_line,
                                                inche_from_connect_x,
                                                inche_from_connect_y,
                                                inche_to_connect_y
                                        ):
                                            chosen_x = x_candidate
                                            break

                                    inche_from_connect_x = chosen_x
                                    inche_to_connect_x = chosen_x

                                    # Capture line geometry in 1st pass for the
                                    # post-pass device-shift collision analysis
                                    # (see compute_l3_device_shifts). natural_x
                                    # is the un-bucketed tag center; the device
                                    # name lets us anchor the line to the device
                                    # so that subsequent shifts move both
                                    # together.
                                    if action_type == 'CREATE' and self.click_value_l3 == 'L3-4-1' and self.flag_re_create == False:
                                        if not hasattr(self, '_collision_lines_per_area'):
                                            self._collision_lines_per_area = {}
                                        natural_x = tmp_all_l3if_tag_array[1] + tmp_all_l3if_tag_array[3] * 0.5
                                        self._collision_lines_per_area.setdefault(target_folder_name, []).append({
                                            'device_name': tmp_all_l3if_tag_array[7],
                                            'natural_x': natural_x,
                                            'y0': min(inche_from_connect_y, inche_to_connect_y),
                                            'y1': max(inche_from_connect_y, inche_to_connect_y),
                                        })

                                    # Rest of the code remains unchanged
                                    if self.click_value_VPN == 'VPN-1-3':
                                        if [tmp_all_l3if_tag_array[7], tmp_all_l3if_tag_array[6]] in self.vpn_hostname_if_list:
                                            line_type = 'L3_SEGMENT-VPN'

                                    skip_connectors = (self.click_value_l3 == 'L3-4-1' and self.flag_re_create == False and self.flag_second_page == False)

                                    if action_type == 'CREATE' and not skip_connectors:
                                        nsm_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                    _line_entry = [line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y]
                                    self.all_written_line_position_array.append(_line_entry)
                                    _bk = int(inche_from_connect_x / min_between_line)
                                    self._written_line_x_buckets.setdefault(_bk, []).append(_line_entry)

                                    if inche_from_connect_x < edge_left_x:
                                        edge_left_x = inche_from_connect_x
                                    if inche_from_connect_x > edge_right_x:
                                        edge_right_x = inche_from_connect_x

                                    #make l3 network address list
                                    key = (tmp_all_l3if_tag_array[7], tmp_all_l3if_tag_array[6])
                                    for tmp_update_l3_table_array in self.l3_rows_by_device_if.get(key, []):
                                        key_if = (tmp_update_l3_table_array[1], tmp_update_l3_table_array[2])
                                        remake_array = self._l3_ipset_cache.get(key_if)
                                        if remake_array is None:
                                            try:
                                                remake_array = ast.literal_eval(tmp_update_l3_table_array[6])
                                            except (ValueError, SyntaxError):
                                                remake_array = []
                                            self._l3_ipset_cache[key_if] = remake_array

                                        for tmp_remake_array in remake_array:
                                            l3_network_list.append(tmp_remake_array[1])
                                            all_l3_netowrk_list.append([tmp_remake_array[1],[tmp_all_l3if_tag_array[7],tmp_all_l3if_tag_array[6]]])
                                            the_l3segment_l3_network_list.append([tmp_remake_array[1],[tmp_all_l3if_tag_array[7],tmp_all_l3if_tag_array[6]]])


                                    #mark written if line for write ip address
                                    self.all_written_if_line_array.append([tmp_all_l3if_tag_array[7],tmp_all_l3if_tag_array[6]])

                    l3segment_edge_array = [edge_left_x,edge_right_x ]
                    l3_network_list = nsm_def.get_l2_broadcast_domains.get_unique_list(l3_network_list)

                    #print('--- l3_network_list ---')
                    #print(l3_network_list )



                    '''write l3 segment'''
                    #print('### l3segment_edge_array   ',l3segment_edge_array)

                    if len(l3_network_list) >= 2:
                        tmp_text = ''
                        for index_01, tmp_l3_network_list in enumerate(l3_network_list):
                            tmp_text += '(' + str(nsm_def.num2alpha(multi_ip_address_num)) + ')' + tmp_l3_network_list + '  '

                            # pre match to l3if ip address
                            for index_11,tmp_the_l3segment_l3_network_list in enumerate(the_l3segment_l3_network_list):
                                if tmp_the_l3segment_l3_network_list[0] == tmp_l3_network_list:
                                    the_l3segment_l3_network_list[index_11].append('(' + str(nsm_def.num2alpha(multi_ip_address_num)) + ')')

                            multi_ip_address_num += 1

                        shape_text = tmp_text


                    elif l3_network_list != []:
                        shape_text = str(l3_network_list[0])

                    else:
                        shape_text = ''


                    #print('--- the_l3segment_l3_network_list ---')
                    #print(the_l3segment_l3_network_list)

                    all_l3segment_l3_netowrk_list.extend(the_l3segment_l3_network_list)


                    shape_left = l3segment_edge_array[0]
                    shape_width = l3segment_edge_array[1] - l3segment_edge_array[0]

                    #check width text
                    tmp_text_width = get_text_wh_cached(self,self.tag_font_large_size, shape_text)[0]
                    if tmp_text_width > shape_width:
                        shape_width = tmp_text_width

                    #add left right inches
                    tmp_add_width = get_text_wh_cached(self,self.tag_font_large_size, 'aa')[0]
                    shape_left -= tmp_add_width
                    shape_width += tmp_add_width * 2

                    #check min width
                    tmp_char_width = get_text_wh_cached(self,self.tag_font_large_size,shape_text)[0]
                    if shape_width < tmp_char_width:
                        shape_width = tmp_char_width

                    if self.click_value_VPN == 'VPN-1-3': #add ver 2.3.2
                        #print('--- self.vpn_hostname_if_list ---')
                        #print(self.vpn_hostname_if_list)

                        list1 = self.vpn_hostname_if_list
                        list2 = the_l3segment_l3_network_list

                        # Logic to check for a complete match (both elements)
                        for item1 in list1:
                            match_found = False
                            # Extract the two elements of list1 (device name and VPN name)
                            device_name = item1[0]
                            vpn_name = item1[1]

                            # Loop over each entry in list2
                            for item2 in list2:
                                # Extract the second element of item2, which is a list with [device_name, interface]
                                interface_info = item2[1]

                                # Check if both the device names and the interface names match
                                if device_name == interface_info[0] and vpn_name == interface_info[1]:
                                    #print(f"--- VPN Segment found:  {interface_info}")
                                    match_found = True
                                    break  # Stop once a match is found

                            if match_found == True:
                                shape_type = 'L3_SEGMENT_VPN'

                    #write l3 segment
                    if action_type == 'CREATE':
                        self.shape = self.slide.shapes
                        nsm_ddx_figure.extended.add_shape(self, shape_type, shape_left, shape_top, shape_width, shape_hight, shape_text)
                        self.slide.shapes._spTree.remove(self.shape._element)  # move shape to back layer
                        self.slide.shapes._spTree.insert(3, self.shape._element)  # move shape to back layer

                        self.y_grid_segment_array.append([index_5,[shape_left, shape_top, shape_width, shape_hight]])


        #print('#### self.all_written_line_position_array  | line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y \n'  ,self.all_written_line_position_array , len(self.all_written_line_position_array))

        '''Write Outline was moved at ver 2.4.2a'''
        if action_type == 'CREATE':
            ### adjust outline width . add ver 2.3.3 ###
            if (max_folder_left_width + 0.5) > (outline_shape_left + outline_shape_width):
                outline_shape_width += max_folder_left_width - (outline_shape_left + outline_shape_width) + 0.5

            #### fix at case 13 at ver 2.4.2a###
            if (shape_hight + shape_top + 0.5) > (outline_shape_top + outline_shape_hight):
                outline_shape_hight = outline_shape_hight + ((shape_hight + shape_top + 0.5) - (outline_shape_top + outline_shape_hight))

            self.shape = self.slide.shapes
            nsm_ddx_figure.extended.add_shape(self, outline_shape_type, outline_shape_left, outline_shape_top, outline_shape_width, outline_shape_hight, outline_shape_text)

            # move shape to back layer when set 'OUTLINE_NORMAL' at ver 2.3.0
            if outline_shape_type == 'OUTLINE_NORMAL':
                self.slide.shapes._spTree.remove(self.shape._element)  # move shape to back layer
                self.slide.shapes._spTree.insert(2, self.shape._element)  # move shape to back layer

        '''write ip on L3 IF'''
        #print('--- self.mark_multi_ip_array ---  [tag_shape_type, tag_shape_left + tag_shape_width * 0.6, tag_shape_top - tag_shape_hight - offset_ipaddress, tag_ip_width, tag_shape_hight, tmp_remake_array[2],tmp_remake_array,shape_text ]')
        #print('--- self.mark_multi_ip_array ---')
        #print(self.mark_multi_ip_array)
        #print('--- self.all_written_if_line_array) ---')
        #print(self.all_written_if_line_array)
        #print('--- all_l3segment_l3_netowrk_list ---')
        #print(all_l3segment_l3_netowrk_list)

        # Track interfaces already drawn for no-L3-segment case to avoid duplicates
        drawn_no_l3seg_interfaces = []

        if self.mark_multi_ip_array != []:
            for tmp_mark_multi_ip_array in self.mark_multi_ip_array:
                if tmp_mark_multi_ip_array[8] == 1:
                    tag_ip_text = tmp_mark_multi_ip_array[5]
                    tag_ip_width = get_text_wh_cached(self,self.tag_font_large_size, tag_ip_text)[0]
                    tag_ip_text_2 = tmp_mark_multi_ip_array[6][0]
                    tag_ip_width_2 = get_text_wh_cached(self,self.tag_font_large_size, tag_ip_text_2 + str('x'))[0]
                    tag_ip_left_2 = tmp_mark_multi_ip_array[10]

                    #chack add (x)
                    for tmp_all_l3segment_l3_netowrk_list in all_l3segment_l3_netowrk_list:
                        if tmp_all_l3segment_l3_netowrk_list[1] == tmp_mark_multi_ip_array[9] and tmp_all_l3segment_l3_netowrk_list[0] == tmp_mark_multi_ip_array[6][1] and len(tmp_all_l3segment_l3_netowrk_list) == 3:
                            tag_ip_text = tmp_all_l3segment_l3_netowrk_list[2] + tmp_mark_multi_ip_array[5]
                            tag_ip_width = get_text_wh_cached(self,self.tag_font_large_size, tag_ip_text)[0]
                            tag_ip_text_2 = tmp_all_l3segment_l3_netowrk_list[2] + tmp_mark_multi_ip_array[6][0]
                            tag_ip_width_2 = get_text_wh_cached(self,self.tag_font_large_size, tag_ip_text_2 + str('x'))[0]

                    # check no line l3 if
                    if tmp_mark_multi_ip_array[9] in self.all_written_if_line_array:
                        if action_type == 'CREATE' and not minimal_1st_pass:
                            self.shape = self.slide.shapes
                            nsm_ddx_figure.extended.add_shape(self, tmp_mark_multi_ip_array[0], tmp_mark_multi_ip_array[1], tmp_mark_multi_ip_array[2], tag_ip_width, tmp_mark_multi_ip_array[4], tag_ip_text)
                    else:
                        if action_type == 'CREATE' and not minimal_1st_pass:
                            self.shape = self.slide.shapes
                            nsm_ddx_figure.extended.add_shape(self, tmp_mark_multi_ip_array[0], tag_ip_left_2, tmp_mark_multi_ip_array[2], tag_ip_width_2, tmp_mark_multi_ip_array[4], tag_ip_text_2)

                else:
                    # Flag to track if we found a matching L3 segment
                    found_l3segment_match = False

                    for tmp_all_l3segment_l3_netowrk_list in all_l3segment_l3_netowrk_list:
                        if tmp_all_l3segment_l3_netowrk_list[1] == tmp_mark_multi_ip_array[9] and tmp_all_l3segment_l3_netowrk_list[0] == tmp_mark_multi_ip_array[6][1]:
                            found_l3segment_match = True
                            # If the index does not exist, expand the list and provide an empty string at index 2. Bug fix at Ver 2.5.3a
                            if len(tmp_all_l3segment_l3_netowrk_list) <= 2:
                                tmp_all_l3segment_l3_netowrk_list.extend(
                                    [''] * (3 - len(tmp_all_l3segment_l3_netowrk_list)))
                            ################################################################

                            tag_ip_text = tmp_all_l3segment_l3_netowrk_list[2] + tmp_mark_multi_ip_array[5]
                            tag_ip_width = get_text_wh_cached(self,self.tag_font_large_size, tag_ip_text)[0]
                            tag_ip_text_2 = tmp_all_l3segment_l3_netowrk_list[2]  + tmp_mark_multi_ip_array[6][0]
                            tag_ip_width_2 = get_text_wh_cached(self,self.tag_font_large_size, tag_ip_text_2 + str('x'))[0]
                            tag_ip_left_2 = tmp_mark_multi_ip_array[10]

                            # print('### tag_ip_text, tag_ip_width',tag_ip_text,tag_ip_width)

                            # check no line l3 if
                            if tmp_mark_multi_ip_array[9] in self.all_written_if_line_array:
                                if action_type == 'CREATE' and not minimal_1st_pass:
                                    self.shape = self.slide.shapes
                                    nsm_ddx_figure.extended.add_shape(self, tmp_mark_multi_ip_array[0], tmp_mark_multi_ip_array[1], tmp_mark_multi_ip_array[2], tag_ip_width, tmp_mark_multi_ip_array[4], tag_ip_text)
                            else:
                                if action_type == 'CREATE' and not minimal_1st_pass:
                                    self.shape = self.slide.shapes
                                    nsm_ddx_figure.extended.add_shape(self, tmp_mark_multi_ip_array[0], tag_ip_left_2, tmp_mark_multi_ip_array[2], tag_ip_width_2, tmp_mark_multi_ip_array[4], tag_ip_text_2)
                            break

                    # If no L3 segment match found, still draw the IP addresses (stacked vertically)
                    if not found_l3segment_match:
                        current_if_key = tmp_mark_multi_ip_array[9]
                        # Check if already drawn for this interface
                        if current_if_key not in drawn_no_l3seg_interfaces:
                            drawn_no_l3seg_interfaces.append(current_if_key)
                            if action_type == 'CREATE' and not minimal_1st_pass:
                                # Collect full entries to use their pre-calculated y positions.
                                # Each entry already has the correct y for UP (stacked upward)
                                # or DOWN (stacked downward) interfaces.
                                all_ip_entries_for_if = []
                                for check_ip_array in self.mark_multi_ip_array:
                                    if check_ip_array[9] == current_if_key:
                                        all_ip_entries_for_if.append(check_ip_array)

                                tag_ip_left_2 = tmp_mark_multi_ip_array[10]
                                for ip_entry in all_ip_entries_for_if:
                                    ip_addr = ip_entry[6][0]
                                    ip_y = ip_entry[2]  # Pre-calculated y (UP=upward, DOWN=downward)
                                    tag_ip_width = get_text_wh_cached(self, self.tag_font_large_size, ip_addr)[0]
                                    self.shape = self.slide.shapes
                                    nsm_ddx_figure.extended.add_shape(self, ip_entry[0], tag_ip_left_2, ip_y, tag_ip_width, ip_entry[4], ip_addr)
        '''
        Write line of L3 instance
        '''
        used_line_array = []
        for tmp_update_l3_instance_array in self.update_l3_instance_array:
            for tmp_all_l3if_tag_array in self.all_l3if_tag_array:
                if tmp_update_l3_instance_array[0][1] == tmp_all_l3if_tag_array[7] and tmp_update_l3_instance_array[0][2] == tmp_all_l3if_tag_array[6]:
                    line_type = 'L3_INSTANCE'
                    inche_from_connect_x = tmp_all_l3if_tag_array[1] + tmp_all_l3if_tag_array[3] * 0.5
                    inche_from_connect_y = tmp_all_l3if_tag_array[2]

                    for tmp_size_l3_instance_array in self.size_l3_instance_array:
                        if tmp_size_l3_instance_array[0] == tmp_update_l3_instance_array[0][1] and tmp_size_l3_instance_array[1] == tmp_update_l3_instance_array[1]:
                            inche_to_connect_x = tmp_size_l3_instance_array[3] + tmp_size_l3_instance_array[5] * 0.5
                            inche_to_connect_y = tmp_size_l3_instance_array[4] + tmp_size_l3_instance_array[6]

                            # for up
                            if [tmp_all_l3if_tag_array[7], tmp_all_l3if_tag_array[6]] in self.up_down_l3if_array[0]:
                                inche_from_connect_y = tmp_all_l3if_tag_array[2] + tmp_all_l3if_tag_array[4]
                                inche_to_connect_y = tmp_size_l3_instance_array[4]

                            # Skip connectors on 1st pass (All Areas mode only)
                            skip_connectors = (self.click_value_l3 == 'L3-4-1' and self.flag_re_create == False and self.flag_second_page == False)

                            if action_type == 'CREATE' and (not skip_connectors) and [line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y] not in used_line_array:
                                nsm_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)
                                used_line_array.append([line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y])
                                break

        log_elapsed(f"l3_area_create {action_type} {target_folder_name}", t0)
        if action_type == 'CREATE':
            self._t_get_l3_segment_num = 0.0

        # Capture per-area row layout for the post-pass device-shift collision
        # analysis. Per-device positions are captured directly inside the
        # device draw loop (above) using the post-get_l3_shape_offset values.
        if action_type == 'CREATE' and self.click_value_l3 == 'L3-4-1' and self.flag_re_create == False:
            if not hasattr(self, '_collision_rows_per_area'):
                self._collision_rows_per_area = {}
            self._collision_rows_per_area[target_folder_name] = [list(r) for r in target_position_shape_array]

        return ([outline_shape_type, outline_shape_left, outline_shape_top, outline_shape_width, outline_shape_hight, folder_shape_text])


'''
LOCAL DEF
'''
#add at ver 2.4.1
def get_optimize_y_grid_array(self):
    # Global lane packing (interval graph coloring): for each segment, pick
    # the smallest k such that lane Y = first_shape_top + k * 0.25 has no
    # horizontal overlap with any segment ALREADY placed at that Y across
    # ALL groups. This catches cross-group collisions that the previous
    # per-group algorithm missed (different index_5 groups whose
    # first_shape_top happens to coincide). The first segment of each group
    # still prefers k = 0 (its original Y) but is allowed to shift down on
    # cross-group conflict; same-group "fix_first" semantics are preserved
    # because group iteration order keeps the group's earliest seg first.
    print('--- optimize_y_grid_array ---')
    y_grid_segment_per_inches = 0.25
    x_grid_segment_buffer = 0.03

    shapes_by_index = {}
    for index_5, shape_data in self.y_grid_segment_array:
        if index_5 not in shapes_by_index:
            shapes_by_index[index_5] = []
        shapes_by_index[index_5].append(shape_data)

    global_lanes = {}  # rounded Y -> list of (left_with_buf, right_with_buf)
    done_y_grid_segment_array = []

    for idx_key in shapes_by_index:
        group = shapes_by_index[idx_key]
        if not group:
            continue

        first_shape_top = group[0][1]

        for num, seg in enumerate(group):
            seg_left = round(seg[0] - x_grid_segment_buffer, 3)
            seg_right = round(seg[0] + seg[2] + x_grid_segment_buffer, 3)

            k = 0
            Y = first_shape_top
            while k < 1000:
                Y = round(first_shape_top + k * y_grid_segment_per_inches, 3)
                conflict = False
                for (ll, lr) in global_lanes.get(Y, []):
                    if not (lr < seg_left or ll > seg_right):
                        conflict = True
                        break
                if not conflict:
                    break
                k += 1

            global_lanes.setdefault(Y, []).append((seg_left, seg_right))
            if num == 0 and k == 0:
                done_y_grid_segment_array.append([idx_key, seg])
            else:
                done_y_grid_segment_array.append([idx_key, [seg[0], Y, seg[2], seg[3]]])

    return done_y_grid_segment_array


def _detect_l3_overspread(self, threshold=3.0):
    # Return True if any device row's cumulative line-over-device shift exceeds
    # `threshold` inches. Legitimate single-row nudges are well under 1.1" in
    # practice (measured across bug_fix0605 / 5site_2DC_WAN / flow_topology),
    # while the multi-row misalignment cascade spreads a row by 30-44". A 3.0"
    # threshold cleanly separates the two, so the conditional 4th pass only
    # fires on the pathological case and leaves all other diagrams untouched.
    shifts = getattr(self, 'device_extra_shift', None)
    if not shifts:
        return False
    rows_per_area = getattr(self, '_collision_rows_per_area', {})
    devs_per_area = getattr(self, '_collision_devs_per_area', {})
    shape_folder_tuple = getattr(self, 'shape_folder_tuple', {})
    for area_name, rows in rows_per_area.items():
        devs = devs_per_area.get(area_name, {})
        if not devs:
            continue
        for row in rows:
            row_devs = [d for d in row
                        if isinstance(d, str) and d in devs and d != '_AIR_']
            if not row_devs:
                continue
            row_devs.sort(key=lambda d: devs[d]['left'])
            acc = 0.0
            last_area = None
            for d in row_devs:
                d_area = shape_folder_tuple.get(d)
                if d_area != last_area:
                    acc = 0.0
                    last_area = d_area
                acc += shifts.get(d, 0.0)
                if acc > threshold:
                    return True
    return False


def compute_l3_device_shifts(self):
    # Detect L3 segment connector lines that vertically pass over other devices
    # and compute per-device rightward shifts that resolve every collision.
    # Cascade within a row is automatic via the existing
    # `left_offset += shape_width + between_shape_column` accumulation in the
    # device draw loop, so we only store the *additional* delta that each
    # device needs on top of the cascaded movement of devices to its left.
    # Iteration: a fixpoint loop runs until no collisions remain or MAX_ITER
    # is reached. Each iteration only ever increases shifts (monotone).
    BUFFER = 0.20  # inches. Doubled from 0.10 per user request for more
                   # breathing room between L3 segment connectors and device
                   # edges (per Sever-13~1~ visual gap feedback). Still
                   # safely exceeds the ±0.075" line bucketing wiggle.
    MAX_ITER = 1000

    # Same-area blocker limiting (conditional L3 "4th pass"). Normally False
    # (original behaviour: a connector must clear EVERY device it vertically
    # passes over, across all areas). When the 4th pass sets
    # self._l3_same_area_only=True before re-invoking this function, the
    # collision check ignores blockers that belong to a DIFFERENT real area.
    # This removes the pathological over-shift where a row's devices cascade
    # rightward into the next area's column band and then shift further to
    # clear THAT area's devices (which the per-area expansion propagation
    # below moves out of the way anyway). With same-area-only blockers each
    # device only clears its own area's upper rows, so devices pack tightly
    # behind the first one that cleared the block (e.g. PC_NewY_4 sits just
    # right of PC_NewY_3 instead of jumping past Atlanta). Constraint A
    # (connector not over an upper device) still holds within the area, and
    # across areas the expansion propagation keeps the bands separated.
    same_area_only = getattr(self, '_l3_same_area_only', False)

    if not hasattr(self, '_collision_lines_per_area') or not self._collision_lines_per_area:
        return
    rows_per_area = getattr(self, '_collision_rows_per_area', {})
    devs_per_area = getattr(self, '_collision_devs_per_area', {})
    shape_folder_tuple = getattr(self, 'shape_folder_tuple', {})

    # Re-entry guard: this function may be called a second time by the
    # conditional 4th pass (with same_area_only set). The per-area expansion
    # propagation below mutates update_start_area_array in place
    # (entry[1] += add). On the normal call we snapshot the clean baseline
    # start_x values; on the 4th-pass re-call we restore them first so
    # additions are recomputed from the baseline instead of accumulating
    # (no double-shift). The snapshot is keyed to the normal call so it is
    # refreshed for every export.
    _usa = getattr(self, 'update_start_area_array', None)
    if _usa:
        if not same_area_only:
            self._update_start_area_orig = [
                (e[1] if isinstance(e, (list, tuple)) and len(e) >= 2 else None)
                for e in _usa
            ]
        elif hasattr(self, '_update_start_area_orig'):
            for e, orig in zip(_usa, self._update_start_area_orig):
                if orig is not None and isinstance(e, list) and len(e) >= 2:
                    e[1] = orig

    # Re-entry guard for calculated_max_right_edge (same rationale as the
    # update_start_area_array guard above). Step 7 below does
    # `calculated_max_right_edge += max_total_growth`. Now that the packing
    # path calls this function twice in succession BEFORE slide_width is
    # computed (cross-area then same-area), a naive second `+=` would
    # double-count the growth. Snapshot the clean post-calculate_area_offset
    # baseline on the normal call and restore it before the bump on the
    # same-area re-call, so each call bumps from the clean baseline.
    if hasattr(self, 'calculated_max_right_edge'):
        if not same_area_only:
            self._calc_max_right_edge_orig = self.calculated_max_right_edge
        elif hasattr(self, '_calc_max_right_edge_orig'):
            self.calculated_max_right_edge = self._calc_max_right_edge_orig

    self.device_extra_shift = {}

    def _per_area_cumul(row, shifts_self, shifts_iter=None):
        # Compute per-device cumulative shift, resetting at every area
        # boundary inside a row. This matches the actual draw-loop behavior
        # where get_l3_shape_offset resets left_offset to area_start_x at the
        # first device of each area's sub-row.
        out = {}
        acc = 0.0
        last_area = None
        for d in row:
            d_area = shape_folder_tuple.get(d)
            if d_area != last_area:
                acc = 0.0
                last_area = d_area
            acc += shifts_self.get(d, 0.0)
            if shifts_iter is not None:
                acc += shifts_iter.get(d, 0.0)
            out[d] = acc
        return out

    for area_name, lines in self._collision_lines_per_area.items():
        rows = rows_per_area.get(area_name, [])
        devs = devs_per_area.get(area_name, {})
        if not rows or not lines or not devs:
            continue

        # Per real-area top Y (min device top), used to group areas into
        # horizontal Y-bands. same_area_only ignores a cross-area blocker
        # ONLY when the two areas share a Y-band (side-by-side, separated in X
        # by expansion propagation). Areas in a different Y-band (vertically
        # stacked, e.g. Datacenter above NewYork) are NOT separated in X, so
        # their devices remain real blockers (an inter-area connector must
        # still avoid them).
        _Y_BAND_TOL = 0.5  # inches
        real_area_top = {}
        for _dn, _b in devs.items():
            _ra = shape_folder_tuple.get(_dn)
            if _ra is None:
                continue
            if _ra not in real_area_top or _b['top'] < real_area_top[_ra]:
                real_area_top[_ra] = _b['top']

        # Order devices in each row by their 1st-pass left position (this is
        # the order the draw loop will visit them and the order in which
        # cascade applies).
        row_order = []
        for row in rows:
            row_devs = [d for d in row if isinstance(d, str) and d in devs and d != '_AIR_']
            row_devs.sort(key=lambda d: devs[d]['left'])
            row_order.append(row_devs)

        lines_by_dev = {}
        for line in lines:
            lines_by_dev.setdefault(line['device_name'], []).append(line)

        # Fixpoint loop. Each iteration scans rows left-to-right, applying
        # cascade WITHIN each area's sub-row only.
        for _iteration in range(MAX_ITER):
            # Baseline cumulative shift from prior iterations (per-area).
            base_cumul = {}
            for row in row_order:
                base_cumul.update(_per_area_cumul(row, self.device_extra_shift))

            iter_extra = {}
            any_added = False

            for row in row_order:
                running = 0.0
                last_area = None
                for d in row:
                    d_area = shape_folder_tuple.get(d)
                    if d_area != last_area:
                        running = 0.0
                        last_area = d_area
                    cumul_d = base_cumul.get(d, 0.0) + running
                    d_lines = lines_by_dev.get(d, [])
                    if not d_lines:
                        continue
                    max_need = 0.0
                    for line in d_lines:
                        effective_x = line['natural_x'] + cumul_d
                        for blocker_name, b in devs.items():
                            if blocker_name == d:
                                continue
                            # 4th pass: ignore blockers in a different real
                            # area ONLY when that area shares a Y-band with the
                            # device's area (side-by-side; expansion separates
                            # them in X). Cross-area blockers in a different
                            # Y-band (vertically stacked) are kept, since X
                            # separation does not apply to them. This avoids
                            # the over-shift where a row cascades into the next
                            # side-by-side area's columns and jumps past it,
                            # while still preventing inter-area connectors from
                            # crossing a stacked area's devices.
                            if same_area_only:
                                _b_area = shape_folder_tuple.get(blocker_name)
                                if (_b_area != d_area and abs(
                                        real_area_top.get(_b_area, -1e9)
                                        - real_area_top.get(d_area, 1e9)) <= _Y_BAND_TOL):
                                    continue
                            b_left = b['left'] + base_cumul.get(blocker_name, 0.0) + iter_extra.get(blocker_name, 0.0)
                            b_right = b_left + b['width']
                            if not (b_left < effective_x < b_right):
                                continue
                            if line['y0'] < b['top'] + b['height'] and b['top'] < line['y1']:
                                need = b_right - effective_x + BUFFER
                                if need > max_need:
                                    max_need = need
                    if max_need > 0:
                        iter_extra[d] = max_need
                        running += max_need
                        any_added = True

            if not any_added:
                break
            for d, need in iter_extra.items():
                self.device_extra_shift[d] = self.device_extra_shift.get(d, 0.0) + need
        else:
            print(f"[L3 collision] Warning: shift fixpoint did not converge after {MAX_ITER} iterations for area '{area_name}'")

    # ===== Left-pack relaxation (4th pass / same_area_only) =====
    # The fixpoint above only ever ADDS rightward shift (monotone), so a
    # device can keep an early shift it no longer needs once the cascade has
    # moved it past the block (e.g. PC_NewY_4 ends far right of PC_NewY_3 even
    # though its connector would clear the upper rows just one pitch to the
    # right). Pull each device LEFT to the leftmost position whose connector
    # still clears all relevant blockers (same Y-band cross-area blockers are
    # ignored, identical to the fixpoint rule) and that keeps >= pitch from
    # its left neighbour. This is decrease-only, so constraint A (connector
    # not over an upper device) is preserved, and the per-area expansion
    # propagation below then runs on the smaller, packed shifts.
    if same_area_only and self.device_extra_shift:
        _Y_BAND_TOL_R = 0.5  # inches
        _GAP = getattr(self, 'between_shape_column', 0.5)
        for _aname in self._collision_lines_per_area:
            _devs = devs_per_area.get(_aname, {})
            _rows = rows_per_area.get(_aname, [])
            _lns = self._collision_lines_per_area.get(_aname, [])
            if not _devs or not _rows or not _lns:
                continue
            _rat = {}
            for _dn, _b in _devs.items():
                _ra = shape_folder_tuple.get(_dn)
                if _ra is None:
                    continue
                if _ra not in _rat or _b['top'] < _rat[_ra]:
                    _rat[_ra] = _b['top']
            _conn = {}
            for _l in _lns:
                _dn = _l['device_name']
                if _dn in _devs and _dn not in _conn:
                    _conn[_dn] = (_l['natural_x'] - _devs[_dn]['left'], _l['y0'], _l['y1'])
            # Current absolute X per device (post-fixpoint).
            _cur = {}
            for _row in _rows:
                _cur.update(_per_area_cumul(
                    [d for d in _row if isinstance(d, str) and d in _devs and d != '_AIR_'],
                    self.device_extra_shift))
            _cur_abs = {d: _devs[d]['left'] + c for d, c in _cur.items()}
            _final_abs = {}
            # Process rows top-to-bottom so upper-row final positions are used
            # as blockers for lower rows.
            _row_order = sorted(
                range(len(_rows)),
                key=lambda ri: min(
                    [_devs[d]['top'] for d in _rows[ri]
                     if isinstance(d, str) and d in _devs], default=1e9))
            for _ri in _row_order:
                _seg = [d for d in _rows[_ri]
                        if isinstance(d, str) and d in _devs and d != '_AIR_']
                _seg.sort(key=lambda d: _devs[d]['left'])
                _prev_right = {}
                _prev_cumt = {}
                for d in _seg:
                    a = shape_folder_tuple.get(d)
                    natl = _devs[d]['left']
                    w = _devs[d]['width']
                    cabs = _cur_abs.get(d, natl)
                    lo = _prev_right.get(a)
                    cand = natl if lo is None else max(natl, lo + _GAP)
                    if cand < cabs:
                        co = _conn.get(d)
                        if co is not None:
                            off, y0, y1 = co
                            for _it in range(400):
                                cx = cand + off
                                worst = None
                                for bn, b in _devs.items():
                                    if bn == d:
                                        continue
                                    barea = shape_folder_tuple.get(bn)
                                    if (barea != a and abs(_rat.get(barea, -1e9)
                                            - _rat.get(a, 1e9)) <= _Y_BAND_TOL_R):
                                        continue
                                    babs = _final_abs.get(bn, _cur_abs.get(bn, b['left']))
                                    if not (babs < cx < babs + b['width']):
                                        continue
                                    if y0 < b['top'] + b['height'] and b['top'] < y1:
                                        if worst is None or babs + b['width'] > worst:
                                            worst = babs + b['width']
                                if worst is None:
                                    break
                                cand = worst - off + BUFFER
                                if cand >= cabs:
                                    break
                    target = min(cand, cabs)
                    if lo is not None and target < lo + _GAP:
                        target = lo + _GAP
                    _final_abs[d] = target
                    _prev_right[a] = target + w
                    cumt = target - natl
                    pc = _prev_cumt.get(a, 0.0)
                    self.device_extra_shift[d] = cumt - pc
                    _prev_cumt[a] = cumt
    # ===== end left-pack relaxation =====

    # ===== Per-area expansion propagation =====
    # When devices inside an area shift right, the area's right edge expands.
    # Subsequent areas IN THE SAME HORIZONTAL ROW must shift right by the
    # cumulative expansion of all prior areas in that row to avoid frame
    # overlap. Areas in a different Y-band (e.g. data center row vs branch
    # row) are independent.
    if not self.device_extra_shift:
        return

    update_start_area_array = getattr(self, 'update_start_area_array', [])

    # 1. Per-real-area max expansion (max sub-row cumulative shift)
    area_expansion = {}
    for area_name, rows in rows_per_area.items():
        devs = devs_per_area.get(area_name, {})
        for row in rows:
            row_devs = [d for d in row if isinstance(d, str) and d in devs and d != '_AIR_']
            row_devs.sort(key=lambda d: devs[d]['left'])
            cumul_in_area = _per_area_cumul(row_devs, self.device_extra_shift)
            for d, c in cumul_in_area.items():
                d_area = shape_folder_tuple.get(d)
                if d_area is None:
                    continue
                if c > area_expansion.get(d_area, 0.0):
                    area_expansion[d_area] = c

    if not area_expansion:
        return

    # 2. Per-real-area top Y (= min device top across all devices in that area).
    # Used to group areas into horizontal Y-bands for cascade propagation.
    area_top = {}
    for area_devs in devs_per_area.values():
        for d_name, d_info in area_devs.items():
            d_real_area = shape_folder_tuple.get(d_name)
            if d_real_area is None:
                continue
            t = d_info['top']
            if d_real_area not in area_top or area_top[d_real_area] > t:
                area_top[d_real_area] = t

    # 3. Normalise update_start_area_array entries: smallest start_x per area.
    area_orig_start = {}
    for entry in update_start_area_array:
        if isinstance(entry, (list, tuple)) and len(entry) >= 2:
            a, sx = entry[0], entry[1]
            if a not in area_orig_start or area_orig_start[a] > sx:
                area_orig_start[a] = sx

    # 4. Group areas into Y-bands (areas with similar top are in the same
    # horizontal row of areas). Tolerance: 0.5" accounts for minor draw-time
    # shifts.
    Y_BAND_TOLERANCE = 0.5
    y_bands = []  # list of (representative_top, [area_names])
    for area_name, top in sorted(area_top.items(), key=lambda x: x[1]):
        placed = False
        for band in y_bands:
            if abs(band[0] - top) <= Y_BAND_TOLERANCE:
                band[1].append(area_name)
                placed = True
                break
        if not placed:
            y_bands.append((top, [area_name]))

    # 5. Within each Y-band, sort by start_x and propagate cumulative
    # expansion to subsequent areas.
    area_offset_addition = {}
    max_total_growth = 0.0
    for band_top, band_areas in y_bands:
        band_areas_sorted = sorted(
            (a for a in band_areas if a in area_orig_start),
            key=lambda a: area_orig_start[a],
        )
        cumul_exp = 0.0
        for area_name in band_areas_sorted:
            area_offset_addition[area_name] = cumul_exp
            cumul_exp += area_expansion.get(area_name, 0.0)
        if cumul_exp > max_total_growth:
            max_total_growth = cumul_exp

    # 6. Update update_start_area_array entries in-place.
    for entry in update_start_area_array:
        if isinstance(entry, (list, tuple)) and len(entry) >= 2:
            a = entry[0]
            add = area_offset_addition.get(a, 0.0)
            if add > 0:
                entry[1] += add

    # 7. Bump calculated_max_right_edge by max total growth across Y-bands.
    if hasattr(self, 'calculated_max_right_edge') and max_total_growth > 0:
        self.calculated_max_right_edge += max_total_growth


def _compute_wp_x_snap_targets(self):
    # Compute absolute target X (inches) for WayPoints whose Area Position
    # row contains only `_wp_`-suffixed area names (or blanks). Snaps each
    # such WayPoint horizontally near its leftmost connected device (center-
    # aligned). Falls back to natural cascade for WayPoints that fail any
    # check (no connection, no device position data, would overflow slide).
    #
    # Output: self._wp_x_target = {wp_name: absolute_target_left}
    # The draw loop consults this dict after computing the natural shape_left
    # and overrides if the entry exists. Designed for L3-4-1 All Areas mode,
    # 2nd pass only (after compute_l3_device_shifts has run). Operates on
    # _collision_devs_per_area populated in the 1st pass.
    #
    # IMPORTANT: In All Areas mode, create_master_file_one_area rewrites
    # POSITION_FOLDER and POSITION_SHAPE on the staging master to a single
    # 'All Areas' composite, which loses the `_wp_` folder identity. We must
    # therefore read these sections from the ORIGINAL master file
    # (self.inFileTxt_L3_3_1.get()), not the rewritten staging copy.
    self._wp_x_target = {}

    devs_per_area = getattr(self, '_collision_devs_per_area', None)
    if not devs_per_area:
        return

    # Merge per-area device geometry into a single dict for cross-area
    # lookup (All Areas mode uses a single composite key but stay defensive).
    all_devs = {}
    for area_devs in devs_per_area.values():
        all_devs.update(area_devs)
    if not all_devs:
        return

    # Load original (pre-rewrite) sections to retrieve `_wp_` folder identity.
    try:
        orig_master = self.inFileTxt_L3_3_1.get()
    except AttributeError:
        return
    if not orig_master:
        return
    try:
        orig_pos_folder = nsm_def.convert_master_to_array(
            'Master_Data', orig_master, '<<POSITION_FOLDER>>')
        orig_pos_shape = nsm_def.convert_master_to_array(
            'Master_Data', orig_master, '<<POSITION_SHAPE>>')
        orig_pos_line = nsm_def.convert_master_to_array(
            'Master_Data', orig_master, '<<POSITION_LINE>>')
    except Exception:
        return
    if not orig_pos_folder or not orig_pos_shape:
        return

    # 1. Identify target rows: every non-empty cell ends with '_wp_'.
    target_rows = []
    for entry in orig_pos_folder:
        if not entry or len(entry) < 2:
            continue
        cells = entry[1] if isinstance(entry[1], list) else []
        if not cells:
            continue
        if cells[0] in ('<<POSITION_FOLDER>>', '<SET_WIDTH>'):
            continue
        area_cells = [c for c in cells[1:] if isinstance(c, str)]
        non_empty = [c for c in area_cells if c != '']
        if not non_empty:
            continue
        if all(c.endswith('_wp_') for c in non_empty):
            target_rows.append(non_empty)
    if not target_rows:
        return

    # 2. Collect `_wp_` folder names (any folder appearing in any cell).
    wp_folders = set()
    for entry in orig_pos_folder:
        if not entry or len(entry) < 2:
            continue
        cells = entry[1] if isinstance(entry[1], list) else []
        if not cells or cells[0] in ('<<POSITION_FOLDER>>', '<SET_WIDTH>'):
            continue
        for c in cells[1:]:
            if isinstance(c, str) and c.endswith('_wp_'):
                wp_folders.add(c)
    if not wp_folders:
        return

    # 3. Rebuild WP_name -> wp_folder from original POSITION_SHAPE. Each row
    #    that contains a wp_folder name has its other (non-marker, non-folder)
    #    cells treated as WP device names belonging to that folder.
    wp_to_folder = {}
    for entry in orig_pos_shape:
        if not entry or len(entry) < 2:
            continue
        if entry[0] in (1, 2):  # header rows
            continue
        cells = entry[1] if isinstance(entry[1], list) else []
        if not cells:
            continue
        folder_in_row = None
        for c in cells:
            if isinstance(c, str) and c in wp_folders:
                folder_in_row = c
                break
        if folder_in_row is None:
            continue
        for c in cells:
            if not isinstance(c, str):
                continue
            if c == folder_in_row or c == '<END>' or c.startswith('<'):
                continue
            if c == '':
                continue
            wp_to_folder[c] = folder_in_row
    if not wp_to_folder:
        return

    # 4. Identify WPs placed in target rows.
    target_areas = set()
    for row in target_rows:
        target_areas.update(row)
    target_wps = {wp: f for wp, f in wp_to_folder.items() if f in target_areas}
    if not target_wps:
        return
    wp_set = set(wp_to_folder.keys())

    # 5. Build WayPoint -> connected device map from original POSITION_LINE.
    wp_to_devices = {}
    for item in orig_pos_line:
        if not item or len(item) < 2:
            continue
        if item[0] in (1, 2):
            continue
        cells = item[1] if isinstance(item[1], list) else []
        if len(cells) < 2:
            continue
        from_name, to_name = cells[0], cells[1]
        if not isinstance(from_name, str) or not isinstance(to_name, str):
            continue
        if from_name in wp_set and to_name not in wp_set:
            wp_to_devices.setdefault(from_name, set()).add(to_name)
        if to_name in wp_set and from_name not in wp_set:
            wp_to_devices.setdefault(to_name, set()).add(from_name)
    target_wp_to_row = {wp: row for wp in target_wps for row in target_rows if target_wps[wp] in row}

    # 6. Provisional target_x: center the WayPoint above the leftmost
    #    connected device. Skip when geometry is missing.
    #    `_collision_devs_per_area` captures 1st pass coordinates; the
    #    leftmost device may have moved rightward in 2nd pass due to
    #    compute_l3_device_shifts (line-over-device avoidance). We add the
    #    device's own extra_shift as a first-order correction. Cumulative
    #    cascade from prior devices in the same row is ignored for
    #    simplicity ("rough placement" accepted by spec).
    extra_shift = getattr(self, 'device_extra_shift', {}) or {}
    provisional = {}
    for wp_name in target_wps:
        wp_info = all_devs.get(wp_name)
        if not wp_info:
            continue
        connected = wp_to_devices.get(wp_name, set())
        if not connected:
            continue
        dev_info_list = []
        for dev in connected:
            d = all_devs.get(dev)
            if d:
                shifted_left = d['left'] + extra_shift.get(dev, 0.0)
                dev_info_list.append((shifted_left, d['width'], dev))
        if not dev_info_list:
            continue
        # `connected` is a set, so its iteration order is hash-seed dependent.
        # When two connected devices share the same leftmost X but differ in
        # width, sorting on x[0] alone leaves the tie to be broken by the
        # (nondeterministic) insertion order, which made dev_center -- and
        # thus the WayPoint snap target -- jitter by (width_diff / 2) between
        # runs. Break exact-left ties by device name so the leftmost choice
        # (and the whole diagram) is deterministic regardless of hash seed.
        dev_info_list.sort(key=lambda x: (x[0], x[2]))
        leftmost_left, leftmost_width = dev_info_list[0][0], dev_info_list[0][1]
        dev_center = leftmost_left + leftmost_width / 2.0
        target_left = dev_center - wp_info['width'] / 2.0
        if target_left < self.left_margin:
            target_left = self.left_margin
        provisional[wp_name] = target_left

    if not provisional:
        return

    # 6.5. Line-over-device safety check. Project each target WP's owned
    # connector lines to where they would land after snap
    # (line_x_new = line.natural_x + delta) and discard the snap when the
    # projected line passes vertically through any other device's post-shift
    # x-range. compute_l3_device_shifts resolved collisions based on the
    # 1st-pass (unshifted) line positions, so a WP snap that moves the line
    # rightward by `delta` can re-introduce a collision against a device
    # that had already been shifted. Observed on Sample.figure51 where WAN-1
    # snap moved its connector into the (now-shifted) WAN-1R~1~ x-range.
    lines_by_dev = {}
    for area_lines in getattr(self, '_collision_lines_per_area', {}).values():
        for ln in area_lines:
            lines_by_dev.setdefault(ln['device_name'], []).append(ln)

    safe_provisional = {}
    for wp_name, target_left in provisional.items():
        wp_info = all_devs.get(wp_name)
        if not wp_info:
            continue
        delta = target_left - wp_info['left']
        wp_lines = lines_by_dev.get(wp_name, [])
        collision = False
        for ln in wp_lines:
            line_new_x = ln['natural_x'] + delta
            y0, y1 = ln['y0'], ln['y1']
            for other_name, other_info in all_devs.items():
                if other_name == wp_name:
                    continue
                other_left = other_info['left'] + extra_shift.get(other_name, 0.0)
                other_right = other_left + other_info['width']
                other_top = other_info['top']
                other_bot = other_top + other_info['height']
                # x-containment uses ±0.075" tolerance for line bucketing
                # wiggle; y overlap is strict half-open interval.
                if (other_left - 0.075 < line_new_x < other_right + 0.075
                        and y0 < other_bot and other_top < y1):
                    collision = True
                    break
            if collision:
                break
        if not collision:
            safe_provisional[wp_name] = target_left

    provisional = safe_provisional
    if not provisional:
        return

    # 7. Per-row overlap avoidance and overflow check.
    between = getattr(self, 'between_shape_column', 0.5)
    # Use calculated_max_right_edge when available (post calculate_area_offset)
    # so the limit reflects the actual diagram width, not the unadjusted
    # slide_width. Fall back to slide_width - margin.
    overflow_limit = getattr(self, 'slide_width', 56.0) - 1.0
    edge = getattr(self, 'calculated_max_right_edge', 0.0)
    if edge > overflow_limit:
        overflow_limit = edge

    final = {}
    for row in target_rows:
        row_wps = []
        for area_name in row:
            for wp_name, folder in target_wps.items():
                if folder == area_name and wp_name in provisional:
                    geom = all_devs.get(wp_name)
                    if geom:
                        row_wps.append((wp_name, provisional[wp_name], geom['width']))
        row_wps.sort(key=lambda w: w[1])
        last_right = None
        for wp_name, target_x, width in row_wps:
            # Push right to clear prior WP in the same row.
            if last_right is not None and target_x < last_right + between:
                target_x = last_right + between
            # Overflow policy: skip (keep natural cascade) if it would not
            # fit within the diagram.
            if target_x + width > overflow_limit:
                continue
            final[wp_name] = target_x
            last_right = target_x + width

    self._wp_x_target = final


def _detect_and_fix_2nd_pass_overlaps(self):
    # Detect same-Y X-overlap among 2nd pass segment bar entries in
    # y_grid_segment_array (= entries appended after the snapshot length).
    # When found, push the LATER segment's optimize Y down by 0.25" so the
    # next CREATE pass places it on the next lane. Returns True if any
    # overlaps were detected (re-render needed), False otherwise.
    snap = getattr(self, '_state_snapshot', None)
    if snap is None:
        return False
    snapshot_len = len(snap['y_grid_segment_array'])
    new_entries = self.y_grid_segment_array[snapshot_len:]
    if not new_entries:
        return False

    x_buf = 0.03
    by_y = {}
    for local_idx, entry in enumerate(new_entries):
        _, seg = entry
        Y = round(seg[1], 3)
        by_y.setdefault(Y, []).append((local_idx, seg))

    overlap_local_idxs = set()
    for Y, items in by_y.items():
        if len(items) < 2:
            continue
        for i, (ia, sa) in enumerate(items):
            la = sa[0] - x_buf
            ra = sa[0] + sa[2] + x_buf
            for ib, sb in items[i + 1:]:
                lb = sb[0] - x_buf
                rb = sb[0] + sb[2] + x_buf
                if not (ra < lb or rb < la):
                    # Push the LATER segment (higher index) down
                    overlap_local_idxs.add(max(ia, ib))

    if not overlap_local_idxs:
        return False

    # Apply +0.25" Y shift to corresponding entries in optimize_y_grid_array.
    # 2nd pass CREATE appends segments in the same order as 1st pass, so the
    # local index inside new_entries equals the global index in
    # optimize_y_grid_array.
    y_step = 0.25
    opt = getattr(self, 'optimize_y_grid_array', None)
    if opt is None:
        return False
    for li in overlap_local_idxs:
        if li >= len(opt):
            continue
        opt[li][1][1] = round(opt[li][1][1] + y_step, 3)

    return True


def _pull_back_runaway_devices(self, gap_threshold=2.0):
    # Compact each real area horizontally by closing oversized empty bands.
    # Connectivity-driven placement can drag a device - or a whole sub-cluster
    # of devices (e.g. servers pulled next to their cross-area peers) - far to
    # the right of the rest of the area, leaving a wide empty vertical band of
    # X between a left group and a right group. That overruns the area frame
    # and makes it overlap the neighbouring area. Here we sweep each area's
    # devices left to right and, wherever a clean empty band (no device of any
    # row spans it) is wider than `gap_threshold` inches, pull the entire right
    # group - and everything further right - left so the band shrinks to one
    # normal column gap. The shift is a negative delta on device_extra_shift[
    # name] (same hook the 3rd-pass rightward collision shift uses), which moves
    # the device AND its connected segment geometry on the next CREATE pass.
    # Resolving the area overlap is prioritised, so the inter-area connector
    # lines of the pulled-back devices may become long and cross other content.
    # The move is purely leftward and bounded by the area's natural packed
    # width, so it converges (unlike the area-separation step, which can
    # diverge). A single runaway device is just a right group of size one, so
    # this generalises - and preserves - the previous single-device behaviour
    # (it lands one column to the right of the left group).
    # Returns True when at least one device was moved (caller must re-render);
    # False when nothing changed.
    write_arr = getattr(self, 'add_shape_write_array', None)
    uad = getattr(self, 'unique_area_device_array', [])
    if not write_arr or not uad:
        return False

    dev_area = {}
    for ad in uad:
        if isinstance(ad, (list, tuple)) and len(ad) >= 2:
            dev_area[ad[1]] = ad[0]

    # area_name -> { device_name: [min_left, max_right] } aggregated over every
    # shape carrying that device's name (icon + tags), matching the shapes the
    # drawn area frame is built from.
    area_devs = {}
    for sh in write_arr:
        if not isinstance(sh, (list, tuple)) or len(sh) < 6:
            continue
        name = sh[5]
        a = dev_area.get(name)
        if a is None or '_wp_' in a:
            continue
        left = sh[1]
        right = sh[1] + sh[3]
        top = sh[2]
        ext = area_devs.setdefault(a, {}).get(name)
        if ext is None:
            area_devs[a][name] = [left, right, top]
        else:
            if left < ext[0]:
                ext[0] = left
            if right > ext[1]:
                ext[1] = right
            if top < ext[2]:
                ext[2] = top

    if not hasattr(self, 'device_extra_shift') or self.device_extra_shift is None:
        self.device_extra_shift = {}

    spacing = getattr(self, 'between_shape_column', 0.5)

    # Step 1 - per-area target shift. Sweep each area left to right and close
    # every clean empty band (no device of any row reaches into it) wider than
    # gap_threshold by accruing a cumulative leftward `pull`. Devices at/after a
    # band inherit the accrued pull, so their target shift is -pull; devices
    # before the first band keep target 0.
    target = {}
    for a, devs in area_devs.items():
        if len(devs) < 2:
            continue
        order = sorted(devs.items(), key=lambda kv: kv[1][0])  # by min_left
        pull = 0.0
        prev_max_right = None
        for name, (dleft, dright, _dtop) in order:
            eff_left = dleft - pull
            if prev_max_right is not None and (eff_left - prev_max_right) > gap_threshold:
                pull += (eff_left - prev_max_right) - spacing
                eff_left = dleft - pull
            if pull > 1e-6:
                target[name] = -pull
            eff_right = dright - pull
            if prev_max_right is None or eff_right > prev_max_right:
                prev_max_right = eff_right

    if not target:
        return False

    # Step 2 - apply targets through the row cascade. device_extra_shift is
    # summed into the per-row running left_offset, so a device's net move equals
    # the sum of deltas injected at or before it in its row. Walking each row
    # left to right we inject only the INCREMENT between consecutive devices'
    # targets. This is the crux of holding everything else still: when a pulled
    # device (target -pull) is followed by a device of ANOTHER area (target 0),
    # that next device receives +pull, cancelling the cascade so the neighbour
    # area stays exactly where it was. Only the over-spread devices move left;
    # the long inter-area connector lines that result are the accepted
    # trade-off. (Injecting -pull on every pulled device instead would both sum
    # per row and drag the trailing areas left.)
    row_dev = {}
    for a, devs in area_devs.items():
        for name, (dleft, dright, dtop) in devs.items():
            row_dev.setdefault(round(dtop, 2), []).append((dleft, name))

    changed = False
    for _key, lst in row_dev.items():
        lst.sort()
        prev_target = 0.0
        for _l, name in lst:
            tgt = target.get(name, 0.0)
            inc = tgt - prev_target
            if abs(inc) > 1e-9:
                self.device_extra_shift[name] = self.device_extra_shift.get(name, 0.0) + inc
                changed = True
            prev_target = tgt

    return changed


def _separate_overlapping_areas(self, area_margin_x=0.5, area_margin_y=0.34, desired_gap=0.5):
    # Detect drawn area frames that overlap horizontally (while also overlapping
    # vertically, i.e. they are genuinely side by side on the page) and push the
    # right-hand area - and every area further right - to the right so the
    # frames clear with `desired_gap` between them. In All-Areas mode every area
    # is drawn inside ONE combined page, so an area cannot be moved via a
    # per-area page offset; instead the shift is added to device_extra_shift for
    # every device of the area (the per-device draw hook moves the device AND its
    # connected segment geometry), which rigidly translates the whole area and
    # its frame to the right while preserving its internal layout. Only areas
    # that actually overlap are moved; non-overlapping layouts are left
    # untouched, so this is a no-op (returns 0.0) for normal masters. The area
    # frame is reproduced exactly as it is drawn at L1648-1684: bounding box of
    # every shape in the area (devices + segment bars + labels) expanded by
    # area_margin_x / area_margin_y, with '_wp_' areas excluded (they draw no
    # frame). Returns the largest per-area shift applied this round (0.0 when
    # nothing was shifted); the caller uses the magnitude to detect divergence.
    write_arr = getattr(self, 'add_shape_write_array', None)
    uad = getattr(self, 'unique_area_device_array', [])
    if not write_arr or not uad:
        return 0.0

    dev_area = {}
    for ad in uad:
        if isinstance(ad, (list, tuple)) and len(ad) >= 2:
            dev_area[ad[1]] = ad[0]

    # area_name -> [min_left, min_top, max_right, max_bottom] over all shapes
    bbox = {}
    for sh in write_arr:
        if not isinstance(sh, (list, tuple)) or len(sh) < 6:
            continue
        a = dev_area.get(sh[5])
        if a is None or '_wp_' in a:
            continue
        l = sh[1]
        t = sh[2]
        r = sh[1] + sh[3]
        b = sh[2] + sh[4]
        f = bbox.get(a)
        if f is None:
            bbox[a] = [l, t, r, b]
        else:
            f[0] = min(f[0], l)
            f[1] = min(f[1], t)
            f[2] = max(f[2], r)
            f[3] = max(f[3], b)

    if len(bbox) < 2:
        return 0.0

    # Expand bounding boxes by the drawn frame margins.
    frames = {a: [f[0] - area_margin_x, f[1] - area_margin_y,
                  f[2] + area_margin_x, f[3] + area_margin_y]
              for a, f in bbox.items()}

    # Left-to-right sweep: each area is pushed right just enough to clear every
    # previously placed area it overlaps vertically. Only rightward shifts are
    # ever produced, and an area that already clears its predecessors keeps its
    # original position (shift 0).
    order = sorted(frames.items(), key=lambda kv: kv[1][0])
    shifts = {}
    for i in range(len(order)):
        iname, iframe = order[i]
        required_left = iframe[0]
        for j in range(i):
            jname, jframe = order[j]
            j_right = jframe[2] + shifts.get(jname, 0.0)
            y_overlap = (jframe[1] < iframe[3]) and (iframe[1] < jframe[3])
            if y_overlap and required_left < j_right + desired_gap:
                required_left = j_right + desired_gap
        sh = required_left - iframe[0]
        if sh > 1e-6:
            shifts[iname] = sh

    if not shifts:
        return 0.0

    if not hasattr(self, 'device_extra_shift') or self.device_extra_shift is None:
        self.device_extra_shift = {}

    # Apply the per-area target shift through the row cascade. device_extra_shift
    # is summed into the per-row running left_offset, so a device's net move is
    # the sum of every delta injected at or before it in its row. Injecting the
    # full target on every device would therefore sum per row and overshoot (the
    # cause of the earlier runaway page width); instead, walking each row left to
    # right, we inject only the INCREMENT between consecutive devices' targets.
    # This both moves an area rigidly and lets areas further right inherit the
    # push for free (their increment is measured relative to the area in front
    # of them), so the net shift of every device equals its area's target.
    dev_pos = {}  # device_name -> [min_left, min_top]
    for sh in write_arr:
        if not isinstance(sh, (list, tuple)) or len(sh) < 6:
            continue
        name = sh[5]
        a = dev_area.get(name)
        if a is None or '_wp_' in a:
            continue
        rec = dev_pos.get(name)
        if rec is None:
            dev_pos[name] = [sh[1], sh[2]]
        else:
            if sh[1] < rec[0]:
                rec[0] = sh[1]
            if sh[2] < rec[1]:
                rec[1] = sh[2]

    rows = {}
    for name, (l, t) in dev_pos.items():
        rows.setdefault(round(t, 2), []).append((l, name))

    for _key, lst in rows.items():
        lst.sort()
        prev_target = 0.0
        for _l, name in lst:
            target = shifts.get(dev_area.get(name), 0.0)
            inc = target - prev_target
            if abs(inc) > 1e-9:
                self.device_extra_shift[name] = self.device_extra_shift.get(name, 0.0) + inc
            prev_target = target

    max_shift = max(shifts.values(), default=0.0)
    if max_shift <= 1e-6:
        return 0.0

    # Grow the page so the shifted areas are not clipped.
    self.slide_width = getattr(self, 'slide_width', 0.0) + max_shift

    return max_shift


def _restore_state_and_rerender(self):
    # Restore mutable state arrays to their pre-2nd-pass-CREATE snapshot and
    # re-run the CREATE pass with the updated optimize_y_grid_array. Called
    # from the conditional 3rd pass loop after _detect_and_fix_2nd_pass_overlaps
    # returns True.
    # Note: position_shape_array / position_folder_array are deepcopied
    # from the snapshot each retry because l3_area_create mutates their
    # inner lists (L394-397). Without this, the 3rd pass area-matching
    # fails and most devices/segments are not rendered.
    # IMPORTANT: do NOT `from pptx import Presentation` locally here -- the
    # SVG path monkeypatches `nsm_l3_diagram_create.Presentation = _MockPresentation`
    # at module level. A local import re-binds the name to the real
    # Presentation, which causes the 3rd pass to create real lxml shapes
    # that the mocked add_shape cannot handle (TypeError: expected
    # lxml.etree._Element, got _MockElement). Reference the module-level
    # `Presentation` instead so the SVG path uses the mock.
    import copy
    snap = self._state_snapshot
    self.y_grid_segment_array = list(snap['y_grid_segment_array'])
    self.add_shape_array = list(snap['add_shape_array'])
    self.add_shape_write_array = list(snap['add_shape_write_array'])
    self.per_index2_after_array = list(snap['per_index2_after_array'])
    self.position_shape_array = copy.deepcopy(snap['position_shape_array'])
    self.position_folder_array = copy.deepcopy(snap['position_folder_array'])

    # SVG path: truncate captured entries back to the 2nd pass baseline so
    # the re-render captures the corrected shapes only.
    svg_cap_len = snap.get('svg_capture_list_len')
    svg_cap = getattr(self, '_svg_capture_list', None)
    if svg_cap is not None and svg_cap_len is not None:
        del svg_cap[svg_cap_len:]

    # Fresh Presentation; the CREATE pass below populates it.
    self.active_ppt = Presentation()

    for tmp_new_position_folder_array in self.folder_wp_name_array[0]:
        action_type = 'CREATE'
        offset_x = 0.0
        offset_y = 0.0
        for tmp_page_size_array in self.page_size_array:
            if tmp_page_size_array[5] == tmp_new_position_folder_array:
                offset_x = tmp_page_size_array[1]
                offset_y = tmp_page_size_array[2]
                break
        nsm_l3_diagram_create.l3_area_create(self, tmp_new_position_folder_array, action_type, offset_x, offset_y)


def get_shape_width_if_array(self,device_name):
    if not hasattr(self, '_up_by_device'):
        self._up_by_device = {}
        for item in self.up_down_l3if_array[0]:
            self._up_by_device.setdefault(item[0], []).append(item)
        self._down_by_device = {}
        for item in self.up_down_l3if_array[1]:
            self._down_by_device.setdefault(item[0], []).append(item)

    tmp_up_array = []
    tmp_down_array = []
    full_ip_address_width_array = []

    for item in self._up_by_device.get(device_name, []):
        tmp_return = nsm_def.get_tag_name_from_full_name(device_name, item[1], self.position_line_tuple)
        if tmp_return == '_NO_MATCH_':
            tmp_if_array = nsm_def.adjust_portname(item[1])
            tmp_up_array.append([item[1], str(tmp_if_array[0] + ' ' + tmp_if_array[2]).replace('  ', ' ')])
        else:
            tmp_up_array.append([item[1], tmp_return])

    for item in self._down_by_device.get(device_name, []):
        tmp_return = nsm_def.get_tag_name_from_full_name(device_name, item[1], self.position_line_tuple)
        if tmp_return == '_NO_MATCH_':
            tmp_if_array = nsm_def.adjust_portname(item[1])
            tmp_down_array.append([item[1], str(tmp_if_array[0] + ' ' + tmp_if_array[2]).replace('  ', ' ')])
        else:
            tmp_down_array.append([item[1], tmp_return])

    tmp_up_width = self.between_l3if
    tmp_down_width = self.between_l3if
    full_ip_address = 'xxx.xxx.xxx.xxx/xxxx'
    distance_full_ip_address = get_text_wh_cached(self, self.tag_font_large_size, full_ip_address)[0]

    l3seg_set = getattr(self, '_l3_if_has_l3_segment_set', set())

    for tmp_tmp_up_array in tmp_up_array:
        shape_width_hight_array = get_text_wh_cached(self, self.tag_font_large_size, tmp_tmp_up_array[1])
        first_up_width = shape_width_hight_array[0] + self.between_l3if

        tmp_tmp_up_width = distance_full_ip_address
        if tmp_tmp_up_width > first_up_width and (device_name, tmp_tmp_up_array[0]) not in l3seg_set:
            first_up_width = tmp_tmp_up_width
            full_ip_address_width_array.append([[device_name, tmp_tmp_up_array[0]], tmp_tmp_up_width])

        tmp_up_width += first_up_width

    for tmp_tmp_down_array in tmp_down_array:
        shape_width_hight_array = get_text_wh_cached(self, self.tag_font_large_size, tmp_tmp_down_array[1])
        first_down_width = shape_width_hight_array[0] + self.between_l3if

        tmp_tmp_down_width = distance_full_ip_address
        if tmp_tmp_down_width > first_down_width and (device_name, tmp_tmp_down_array[0]) not in l3seg_set:
            first_down_width = tmp_tmp_down_width
            full_ip_address_width_array.append([[device_name, tmp_tmp_down_array[0]], tmp_tmp_down_width])

        tmp_down_width += first_down_width

    return_shape_width = max(tmp_up_width, tmp_down_width)
    return ([return_shape_width, tmp_up_array, tmp_down_array, full_ip_address_width_array])

def get_up_down_l3if_count(self,marge_target_position_shape_array):
    buttom_array = []
    top_array = []

    row_device_sets = [set(dev for dev in row) for row in marge_target_position_shape_array]

    group_members_by_device = {}
    for grp in self.target_l2_broadcast_group_array:
        member_name_set = set(m[0] for m in grp[1])
        for m in grp[1]:
            group_members_by_device.setdefault(m[0], []).append((m, member_name_set))

    for index_8, tmp_marge_target_position_shape_array in enumerate(marge_target_position_shape_array):
        for device_name in tmp_marge_target_position_shape_array:

            if '_AIR_' in device_name:
                continue
            tmp_buttom_array = []
            tmp_top_array = []

            for member, member_name_set in group_members_by_device.get(device_name, []):
                flag_buttom_exist = False
                for i in range(index_8 + 1, len(marge_target_position_shape_array)):
                    if not member_name_set.isdisjoint(row_device_sets[i]):
                        flag_buttom_exist = True
                        break

                flag_top_exist = False
                for i in range(0, index_8):
                    if not member_name_set.isdisjoint(row_device_sets[i]):
                        flag_top_exist = True
                        break

                if self.new_wp_exist_array[0] == [] and index_8 == 0:
                    tmp_buttom_array.append(member)
                elif flag_buttom_exist and not flag_top_exist:
                    tmp_buttom_array.append(member)
                else:
                    tmp_top_array.append(member)

            buttom_array.extend(tmp_buttom_array)
            top_array.extend(tmp_top_array)

    buttom_array = nsm_def.get_l2_broadcast_domains.get_unique_list(buttom_array)
    top_array = nsm_def.get_l2_broadcast_domains.get_unique_list(top_array)
    return ([top_array,buttom_array])


def check_move_to_right(self,top_device_name_array,target_position_shape_array):
    if self.index_2 >= 1 and self.index_2 <= len(self.index_1_array) - 2:
        source_set = set(target_position_shape_array[self.index_2 + 1])
        for i in range(self.index_2):
            target_set = set(target_position_shape_array[i])

            for grp in self.target_l2_broadcast_group_array:
                member_names = set(m[0] for m in grp[1] if m != [])
                if not member_names.isdisjoint(source_set) and not member_names.isdisjoint(target_set):
                    return True

    return False

def get_l3_segment_num(self,top_device_name_array,target_position_shape_array):
    import time
    _t0 = time.perf_counter()
    count_l3segment = 0
    connected_l3if_key_array = []
    tmp_used_l3segment_array = []

    top_set = set(top_device_name_array)
    used_set = self._used_l3segment_set
    tmp_used_set = set()

    if getattr(self, '_combined_device_set', None) is None:
        self._combined_device_set = set(item for sublist in target_position_shape_array for item in sublist)
    combined_set = self._combined_device_set

    def _hk(v):
        return tuple(v) if isinstance(v, list) else v

    index_1_set = set(self.index_1_array)

    if self.index_2 + 1 in index_1_set:
        for i in range(self.index_2 + 1, len(self.index_1_array)):
            buttom_set = set(target_position_shape_array[i])
            for grp in self.target_l2_broadcast_group_array:
                gid = _hk(grp[0])
                if gid in used_set:
                    continue
                members = grp[1]
                for m in members:
                    if m[0] in top_set and m not in connected_l3if_key_array:
                        for m2 in members:
                            if m2[0] in buttom_set and gid not in tmp_used_set:
                                count_l3segment += 1
                                connected_l3if_key_array.append(m)
                                tmp_used_set.add(gid)
                                tmp_used_l3segment_array.append(grp[0])
                                break

        used_set.update(tmp_used_set)

    if self.index_2 + 1 in index_1_set:
        buttom_device_name_array = target_position_shape_array[self.index_2 + 1]
        buttom_set = set(buttom_device_name_array)
        for grp in self.target_l2_broadcast_group_array:
            members = grp[1]
            tmp_count = sum(1 for m in members if m[0] in buttom_set)

            updated_members = [m for m in members if m[0] in combined_set]

            if tmp_count == len(updated_members) and tmp_count >= 2:
                count_l3segment += 1
                connected_l3if_key_array.append(updated_members[-1])

    if self.index_2 == 0 and self.new_wp_exist_array[0] == []:
        row0_set = set(target_position_shape_array[0])
        for grp in self.target_l2_broadcast_group_array:
            gid = _hk(grp[0])
            tmp_count = 0
            if gid not in used_set:
                members = grp[1]
                for m in members:
                    if m[0] in top_set and m not in connected_l3if_key_array:
                        for m2 in members:
                            if m2[0] in row0_set and gid not in tmp_used_set:
                                tmp_count += 1
                                if tmp_count >= 2:
                                    count_l3segment += 1
                                    connected_l3if_key_array.insert(0, m)
                                    tmp_used_set.add(gid)
                                    tmp_used_l3segment_array.append(grp[0])

        used_set.update(tmp_used_set)

    self._t_get_l3_segment_num = getattr(self, "_t_get_l3_segment_num", 0.0) + (time.perf_counter() - _t0)
    return ([count_l3segment,connected_l3if_key_array])



class  create_master_file_one_area():
    def __init__(self):
        print('--- create_master_file_one_area--- ')

        #copy master file
        shutil.copy(self.inFileTxt_L3_3_1.get(), self.excel_maseter_file_backup)

        #GET backup master file parameter
        # parameter
        ws_name = 'Master_Data'
        tmp_ws_name = '_tmp_'
        ppt_meta_file = str(self.excel_maseter_file_backup)

        # convert from master to array and convert to tuple
        self.position_folder_array = nsm_def.convert_master_to_array(ws_name, ppt_meta_file, '<<POSITION_FOLDER>>')
        self.position_shape_array = nsm_def.convert_master_to_array(ws_name, ppt_meta_file, '<<POSITION_SHAPE>>')
        #self.root_folder_array = nsm_def.convert_master_to_array(ws_name, ppt_meta_file, '<<ROOT_FOLDER>>')
        self.position_folder_tuple = nsm_def.convert_array_to_tuple(self.position_folder_array)
        self.position_shape_tuple = nsm_def.convert_array_to_tuple(self.position_shape_array)
        #self.root_folder_tuple = nsm_def.convert_array_to_tuple(self.root_folder_array)

        #print('---- self.position_folder_tuple ----')
        #print(self.position_folder_tuple)
        #print('---- self.position_folder_array ----')
        #print(self.position_folder_array)
        #print('---- self.position_shape_tuple ----')
        #print(self.position_shape_tuple)
        #print('---- self.position_shape_array ----')
        #print(self.position_shape_array)

        # GET Folder and wp name List
        folder_wp_name_array = nsm_def.get_folder_wp_array_from_master(ws_name, ppt_meta_file)
        #print('---- folder_wp_name_array ----')
        #print(folder_wp_name_array)

        #SET new <<POSITION_FOLDER>>
        self.new_position_folder_tuple = {}
        self.new_position_folder_tuple = {(1, 1): '<<POSITION_FOLDER>>', (1, 2): 1.0, (2, 1): 1, (2, 2): 'All Areas'}

        write_to_section = '<<POSITION_FOLDER>>'
        offset_row = 0
        offset_column = 0
        nsm_def.clear_section_sheet('Master_Data', self.excel_maseter_file_backup, self.position_folder_tuple)
        nsm_def.write_excel_meta(self.new_position_folder_tuple , self.excel_maseter_file_backup, 'Master_Data', write_to_section, offset_row, offset_column)


        '''
        make one area shape <<POSITION_SHAPE>>
        '''
        # Create a new dictionary to store the filtered and renumbered tuples
        original_dict = self.position_folder_tuple
        new_dict = {}
        new_x = 1

        for (x, y), value in original_dict.items():
            if y != 1 and value != '' and isinstance(value, str):  # Exclude if y is 1, value is empty string, or value is not a string
                if new_x not in new_dict:
                    new_dict[new_x] = {}
                new_dict[new_x][y] = value
                if y == max(k[1] for k in original_dict.keys() if k[0] == x):
                    new_x += 1

        # Renumber the (x, y) pairs to ensure no gaps
        renumbered_dict = {}
        new_x = 1
        for x in sorted(new_dict.keys()):
            new_y = 1
            for y in sorted(new_dict[x].keys()):
                renumbered_dict[(new_x, new_y)] = new_dict[x][y]
                new_y += 1
            new_x += 1

        '''make one area from per area''' # renumbered_dict is folder summary

        input_position_shape_tuple = self.position_shape_tuple
        area_start_x = 2
        area_start_y = 1
        area_max_x = 2
        area_max_y = 1
        master_x = 2
        master_y = 1
        new_tuple = {}
        pre_xx = 2
        pre_yy = 1

        for (yy, xx), now_area in sorted(renumbered_dict.items()):
            #print(yy,xx,now_area)

            if yy > pre_yy:
                master_y = area_max_y + 1
                master_x = 2
            elif xx > pre_xx:
                master_y = area_start_y
                master_x = area_max_x

            flag_inarea = False
            flag_first_inarea = True
            for (y, x), value in sorted(input_position_shape_tuple.items()):
                if x == 1 and value == now_area:
                    flag_inarea = True

                if flag_inarea == True:
                    if x == 1:
                        if value == '<END>':
                            pre_yy = yy
                            pre_xx = xx

                            if area_max_y < (master_y - 1):
                                area_max_y = master_y - 1
                            break
                    else:
                        if value != '<END>':
                            if flag_first_inarea == True:
                                area_start_y = master_y
                                area_start_x = master_x
                                flag_first_inarea = False

                            new_tuple[(master_y,master_x)] = value
                            master_x += 1

                        if value == '<END>':
                            master_y += 1
                            area_max_x = master_x
                            master_x = area_start_x

        #print(new_tuple)

        # Find the maximum x for each y
        max_x_for_y = {}
        for (y, x) in new_tuple.keys():
            if y not in max_x_for_y:
                max_x_for_y[y] = x
            else:
                if x > max_x_for_y[y]:
                    max_x_for_y[y] = x

        # Add '<END>' to the new_tuple
        all_max_y = 1
        for y, max_x in max_x_for_y.items():
            new_tuple[(y, max_x + 1)] = '<END>'
            all_max_y = y

        # Find the maximum x value for each y
        max_x_per_y = {}
        for (y, x) in new_tuple.keys():
            if y not in max_x_per_y:
                max_x_per_y[y] = x
            else:
                max_x_per_y[y] = max(max_x_per_y[y], x)

        # Add missing (y, x) pairs with '_AIR_'
        for y, max_x in max_x_per_y.items():
            for x in range(2, max_x + 1):
                if (y, x) not in new_tuple:
                    new_tuple[(y, x)] = '_AIR_'

        #last input
        new_tuple[(1, 1)] = 'All Areas'
        new_tuple[(all_max_y + 1, 1)] = '<END>' #change y axis at ver 2.3.4

        # SET new <<POSITION_SHAPE>>
        write_to_section = '<<POSITION_SHAPE>>'
        offset_row = 1
        offset_column = 0
        nsm_def.clear_section_sheet('Master_Data', self.excel_maseter_file_backup, self.position_shape_tuple)
        nsm_def.write_excel_meta(new_tuple, self.excel_maseter_file_backup, 'Master_Data',write_to_section, offset_row, offset_column)

    def calculate_area_offset(self):
        #print(self.add_shape_array)
        ''' get DEVICE_NORMAL'''
        device_normal_array = []
        for tmp_add_shape_array in self.add_shape_array:
            if tmp_add_shape_array[0] == 'DEVICE_NORMAL' or tmp_add_shape_array[0] == 'DEVICE_L3_INSTANCE':
                device_normal_array.append(tmp_add_shape_array[1:])

        #print('shape_left, shape_top, shape_width, shape_hight,shape_text')
        #print(device_normal_array)

        # List to store the transformed data
        transformed_data = []

        # Processing each sublist
        data = device_normal_array
        self.shape_left_inchi_array = []
        for sublist in data:
            # Add the 1st (index 0), 3rd (index 2), and 4th (index 3) elements together
            sum_value = sublist[0] + sublist[2]
            # Create a new sublist with the last element (string) and the calculated sum
            # shape_text,  shape_top, shape_left,shape_left + shape_width
            transformed_data.append([sublist[4], sublist[1],sublist[0] , sum_value])
            self.shape_left_inchi_array.append([sublist[4],sublist[0]])

        # Print the result
        #print(transformed_data)

        ''' GET Area-Device array'''
        data = self.l2_table_array
        # Filtering data based on condition: first item >= 3, then extracting the first two items
        filtered_data = []
        for item in data:
            if item[0] >= 3:
                filtered_data.append(item[1][:2])  # Extracting the first two items from the second list

        # Removing duplicates
        unique_data = []
        seen = set()
        for entry in filtered_data:
            if entry[1] not in seen:
                unique_data.append(entry)
                seen.add(entry[1])

        # Display the results
        #print(unique_data)

        # change 'N/A' to own folder name
        for item in unique_data:
            if item[0] == 'N/A':
                corresponding_value = self.shape_folder_tuple.get(item[1], None)
                if corresponding_value:
                    item[0] = corresponding_value

        # Display the results
        #print(unique_data)
        self.unique_area_device_array = unique_data

        # Dictionary to group by the first element
        grouped_data = defaultdict(list)

        # Group the data by the first element
        data = unique_data
        for entry in data:
            grouped_data[entry[0]].append(entry[1])

        # Create the result where the second element is an array of devices
        self.result_area_device_array = [[key, value] for key, value in grouped_data.items()]

        # Output the result
        #print(self.result_area_device_array)

        ''' Calculate the distance to the right of each area'''
        # Build a dict for O(1) lookup: device_name -> [shape_top, shape_left, shape_left+width]
        transformed_data_dict = {data[0]: data[1:] for data in transformed_data}

        # Iterate through the self.result_area_device_array to match and add corresponding data
        for area in self.result_area_device_array:
            area_name = area[0]  # Site name
            device_list = area[1]  # List of devices related to the site

            # Create a new list to store the devices with their matched data
            new_device_list = []

            # For each device in the device list, look up in the dict (O(1) vs. O(n) linear scan)
            for device in device_list:
                if device in transformed_data_dict:
                    new_device_list.append([device, *transformed_data_dict[device]])

            # Update the area[1] with the newly created list with the added data
            area[1] = new_device_list

        # Print the updated self.result_area_device_array
        #print('--- self.result_area_device_array ---')
        #print(self.result_area_device_array)

        ''' GET max width per area '''
        area_min_max_diff_array = []
        # Process the self.result_area_device_array to find the required differences
        for area in self.result_area_device_array:
            area_name = area[0]  # Area name, e.g., 'DC-TOP1'
            device_data_list = area[1]  # List of devices and their data

            # Create a dictionary to group devices by their second value (data[1] value)
            grouped_data = {}

            # Group devices by the second element (data[1])
            for device_data in device_data_list:
                key = device_data[1]  # The second item, which we group by (e.g., 3.02, 11.04, etc.)
                if key not in grouped_data:
                    grouped_data[key] = []
                grouped_data[key].append(device_data)

            # Process each group and calculate the required difference
            for key, group in grouped_data.items():
                # Find the smallest third element and the largest fourth element in the group
                min_third = min([item[2] for item in group])  # Smallest value in the third position
                max_fourth = max([item[3] for item in group])  # Largest value in the fourth position

                # Calculate the difference
                difference = max_fourth - min_third

                # Print the result for this group
                area_min_max_diff_array.append([area_name,min_third,max_fourth,difference])
                #print(f"Area: {area_name}, Group: {key}, Min Third: {min_third}, Max Fourth: {max_fourth}, Difference: {difference}")

        #print('--- area_min_max_diff_array ---')
        #print(area_min_max_diff_array)

        ''' Create a defaultdict to group the entries by the category (first element) '''
        # A defaultdict allows us to easily append items to lists without checking if the key exists.
        grouped = defaultdict(list)

        # Iterate through each entry in the input data
        for entry in area_min_max_diff_array:
            category = entry[0]  # Get the category (first element)
            fourth_value = entry[3]  # Get the fourth value (index 3)
            grouped[category].append(
                (fourth_value, entry))  # Append the fourth value and the full entry to the corresponding category

        # List to store the final results
        result_area_width = []

        # Iterate over the grouped categories to find the entry with the maximum fourth value
        for category, values in grouped.items():
            # Use the max function to find the entry with the largest fourth value within each category
            # 'values' is a list of tuples, where the first item in each tuple is the fourth value (index 3)
            max_value_entry = max(values, key=lambda x: x[0])

            # Append the category and the maximum fourth value to the result list
            result_area_width.append([category, max_value_entry[0]])

        # Print the result: for each category, output the category name and the maximum fourth value
        #print('--- result_area_width ---')
        #print(result_area_width)

        ''' GET area location '''
        ws_name = 'Master_Data'
        ppt_meta_file = str(self.inFileTxt_11_1.get())
        ori_position_folder_array = nsm_def.convert_master_to_array(ws_name, ppt_meta_file, '<<POSITION_FOLDER>>')
        #print(ori_position_folder_array)

        # List of strings to exclude
        exclude_strings = ['<<POSITION_FOLDER>>', '<SET_WIDTH>']

        # Initialize an empty list to store the processed results
        processed_array = []

        # Iterate through all elements in ori_position_folder_array
        for element in ori_position_folder_array:
            # Extract the second element (the list) of the current element
            second_element = element[1]

            # Remove numbers and exclude the strings specified in exclude_strings
            cleaned = [item for item in second_element if
                       isinstance(item, str) and item not in exclude_strings and item != '']

            # Only add non-empty lists to processed_array
            if cleaned:
                processed_array.append(cleaned)

        # Print the processed results
        #print(processed_array)

        ''' get start left per area '''
        start_left_inchi = self.left_margin
        between_area_inchi = 1.5

        # Create a dictionary for quick lookup of result_area_width
        result_dict = {item[0]: item[1] for item in result_area_width}

        # Initialize start_area_array
        start_area_array = []

        # Iterate over processed_array to compute start_area_array
        for sublist in processed_array:
            sublist_result = []
            for index, item in enumerate(sublist):
                if item in result_dict:
                    if index == 0:
                        # For the first item, start from 0 + start_left_inchi
                        start_value = start_left_inchi
                    else:
                        # Calculate the value based on conditions for subsequent items
                        start_value = sum(result_dict[prev_item] for prev_item in sublist[:index])
                        start_value += start_left_inchi
                        start_value += between_area_inchi * index
                    sublist_result.append([item, start_value])
            start_area_array.append(sublist_result)

        # Output the start_area_array
        #self.update_start_area_array = []
        for sublist in start_area_array:
            for item in sublist:
                self.update_start_area_array.append(item)

        # Calculate and store the maximum right edge for slide width adjustment
        # Each area's right edge = start_x + area_width
        max_right_edge = 0.0
        for area_info in self.update_start_area_array:
            area_name = area_info[0]
            area_start_x = area_info[1]
            if area_name in result_dict:
                area_width = result_dict[area_name]
                area_right_edge = area_start_x + area_width
                if area_right_edge > max_right_edge:
                    max_right_edge = area_right_edge

        self.calculated_max_right_edge = max_right_edge
        #print(f"[calculate_area_offset] max_right_edge = {max_right_edge}")

        #print(start_area_array)
        #print(self.update_start_area_array)

    def get_l3_shape_offset(self,shape_name,left_offset):
        target_folder_name = self.shape_folder_tuple.get(shape_name)
        area_value = None
        for item in self.update_start_area_array:
            if item[0] == target_folder_name:
                area_value = item[1]
                break

        if shape_name in self.target_offset_shape_array:
            #print(shape_name, target_folder_name, area_value)
            for item in self.shape_left_inchi_array:
                if item[0] == shape_name:
                    left_value = item[1]
                    break

            offset_value = area_value - left_offset

        else:
            offset_value = 0.0

        return offset_value
