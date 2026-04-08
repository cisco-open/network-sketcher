'''
SPDX-License-Identifier: Apache-2.0

Copyright 2023 Cisco Systems, Inc. and its affiliates

Licensed under the Apache License, Version 2.0 (the "License");
you may not use this file except in compliance with the License.
You may obtain a copy of the License at

http://www.apache.org/licenses/LICENSE-2.0

Unless required by applicable law or agreed to in writing, software
distributed under the License is distributed on an "AS IS" BASIS,
WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
See the License for the specific language governing permissions and
limitations under the License.
'''

import tkinter as tk ,tkinter.ttk as ttk,tkinter.filedialog, tkinter.messagebox
import sys, os, shutil , unicodedata,subprocess,datetime,random
import openpyxl
import math ,ipaddress ,yaml, ast
from pptx import *
import platform
from openpyxl.styles import PatternFill
from collections import defaultdict


def check_data_exists(excel_master_file):
    """
    Check if data exists in self.folder_wp_name_array

    Returns:
        bool: True if data exists, False if no data exists
    """
    ws_name = 'Master_Data'

    # GET Folder and wp name List
    folder_wp_name_array = get_folder_wp_array_from_master(ws_name, excel_master_file)
    # Check if folder_wp_name_array exists and is not empty

    try:
        return bool(folder_wp_name_array and
                   len(folder_wp_name_array) > 0 and
                   len(folder_wp_name_array[0]) > 0)
    except (AttributeError, IndexError, TypeError):
        return False

def get_l3_segments(self):
    '''get values of Master Data'''
    # parameter
    ws_l3_name = 'Master_Data_L3'
    excel_maseter_file = self.inFileTxt_L3_3_1.get()

    self.result_get_l2_broadcast_domains = get_l2_broadcast_domains.run(self, excel_maseter_file)  ## 'self.update_l2_table_array, device_l2_boradcast_domain_array, device_l2_directly_l3vport_array, device_l2_other_array, marged_l2_broadcast_group_array'
    #print('--- get_l3_segments ---')
    #print('--- self.target_l2_broadcast_group_array ---')
    #print(self.target_l2_broadcast_group_array)

    self.l3_table_array = convert_master_to_array(ws_l3_name, excel_maseter_file, '<<L3_TABLE>>')
    #print('--- self.l3_table_array ---')
    #print(self.l3_table_array)

    updated_l3_table_array = []
    for index, tmp_l3_table_array in enumerate(self.l3_table_array):
        str(tmp_l3_table_array).replace(' ', '')
        if index >= 2:
            if len(tmp_l3_table_array[1]) == 5:
                if ',' in str(tmp_l3_table_array[1][4]):
                    #print('--- tmp_l3_table_array ', str(tmp_l3_table_array))
                    tmp_tmp_l3_table_array = str(tmp_l3_table_array[1][4]).split(',')
                    for tmp_add_array in tmp_tmp_l3_table_array:
                        tmp_tmp_tmp_l3_table_array = tmp_l3_table_array
                        tmp_tmp_tmp_l3_table_array[1][4] = tmp_add_array
                        #print('--- tmp_tmp_tmp_l3_table_array ', tmp_tmp_tmp_l3_table_array)
                        updated_l3_table_array.append([tmp_tmp_tmp_l3_table_array[1][0], tmp_tmp_tmp_l3_table_array[1][1], tmp_tmp_tmp_l3_table_array[1][2], tmp_tmp_tmp_l3_table_array[1][3], tmp_tmp_tmp_l3_table_array[1][4]])
                else:
                    updated_l3_table_array.append(tmp_l3_table_array[1])

            elif len(tmp_l3_table_array[1]) == 4:
                updated_l3_table_array.append([tmp_l3_table_array[1][0], tmp_l3_table_array[1][1], tmp_l3_table_array[1][2], tmp_l3_table_array[1][3], ''])

            elif len(tmp_l3_table_array[1]) == 3:
                updated_l3_table_array.append([tmp_l3_table_array[1][0], tmp_l3_table_array[1][1], tmp_l3_table_array[1][2], '', ''])

    #print('--- updated_l3_table_array ---')
    #print(updated_l3_table_array)

    '''get segment with target area'''
    l3_segment_group_array = []

    for tmp_target_l2_broadcast_group_array in self.target_l2_broadcast_group_array:
        tmp_l3_segment_group_array = []
        for tmp_tmp_target_l2_broadcast_group_array in tmp_target_l2_broadcast_group_array[1]:
            for tmp_updated_l3_table_array in updated_l3_table_array:
                if tmp_tmp_target_l2_broadcast_group_array[0] == tmp_updated_l3_table_array[1] and tmp_tmp_target_l2_broadcast_group_array[1] == tmp_updated_l3_table_array[2]:
                    tmp_l3_segment_group_array.append(tmp_updated_l3_table_array)
        if tmp_l3_segment_group_array != []:
            l3_segment_group_array.append(tmp_l3_segment_group_array)

    #print('--- l3_segment_group_array ---')
    #print(l3_segment_group_array)

    return l3_segment_group_array

def return_os_slash():
    slash_type = '\\'+'\\'
    os_type = platform.platform()
    #print(os_type)
    if 'macOS'.casefold() in os_type.casefold() or 'Linux'.casefold() in os_type.casefold():
        slash_type = '/'
    return (slash_type)

def get_backup_filename(full_filepath):
    now = datetime.datetime.now(datetime.timezone(datetime.timedelta(hours=9)))
    yyyymmddhhss = str(now.strftime('%Y%m%d%H%M%S'))
    #print(yyyymmddhhss)

    filename = os.path.basename(full_filepath)
    iDir = os.path.abspath(os.path.dirname(full_filepath))
    basename_without_ext = os.path.splitext(os.path.basename(full_filepath))[0]
    basename_ext = os.path.splitext(os.path.basename(full_filepath))[1]
    backup_full_filepath = iDir + return_os_slash() + basename_without_ext + '_' +yyyymmddhhss + basename_ext

    shutil.copyfile(full_filepath, backup_full_filepath)

    return (backup_full_filepath)

def messagebox_file_open(full_filepath):
    if return_os_slash() == '\\\\': # add ver 2.1.1 for bug fix on Mac OS
        filename = os.path.basename(full_filepath)
        ret = tkinter.messagebox.askyesno('Complete', 'Do you want to open the created file?\n\n' + filename)
        if ret == True:
            subprocess.Popen(full_filepath, shell=True)

def check_file_type(full_filepath):
    count_check_tag = 0
    if full_filepath.endswith('.pptx'):
        input_ppt = Presentation(full_filepath)
        for i, sld in enumerate(input_ppt.slides, start=1):
            for shp in sld.shapes:
                if 'AUTO_SHAPE' in str(shp.shape_type) and str(shp.text) != '':
                    ### check contain IF tag adjustments[0] = 0.xx445
                    try:
                        if shp.adjustments[0] == 0.99445 or shp.adjustments[0] == 0.50445:  #check IF tag
                            count_check_tag += 1
                            if count_check_tag >= 2:
                                return (['ERROR','Please enter a PPT file that does not contain IF tags'])

                    except Exception as e:
                        #print('[info] Exception handling with check_file_type')
                        flag_exception_dummy = True

        return_type_array = ['PPT_SKECH','PPT_SKECH']

    elif full_filepath.endswith('.xlsx'):
        return_type_array = ['ERROR', 'Please enter a EXCEL file compatible with NS']
        input_excel = openpyxl.load_workbook(full_filepath)

        # check ws name
        ws_list = input_excel.sheetnames
        for sheet_name in ws_list:
            if sheet_name == 'Master_Data':
                return_type_array = ['EXCEL_MASTER', ws_list]
            elif sheet_name == 'L1 Table':
                return_type_array = ['EXCEL_DEVICE', ws_list]
            elif sheet_name == 'Flow_List': #add ver 2.4.3
                return_type_array = ['EXCEL_FLOW', ws_list]

        # add ver 2.4.3
        if '[MASTER]' not in full_filepath and '[DEVICE]' not in full_filepath and '[FLOW]' not in full_filepath:
            return_type_array = ['ERROR', 'The Excel file name must start with [MASTER],[DEVICE],[FLOW]']

    elif full_filepath.endswith('.yaml'):
        return_type_array = ['ERROR', 'Please enter a backup file of CML']

        with open(str(full_filepath), 'r') as yml:
            config = yaml.safe_load(yml)

        for tmp_config in config:
            if tmp_config == 'lab':
                return_type_array = ['YAML_CML', config]

    elif full_filepath.endswith('.svg'):
        return_type_array = ['SVG', 'dummy']

    elif full_filepath.endswith('.csv'):
        import csv
        def check_csv_headers(full_filepath):
            required_headers = {"Device", "Interface", "Connection"}
            with open(full_filepath, newline='', encoding='utf-8') as csvfile:
                reader = csv.reader(csvfile)
                headers = next(reader)
                return required_headers.issubset(headers)

        if check_csv_headers(full_filepath) == True:
            return_type_array = ['CSV', 'dummy']
        else:
            return_type_array = ['ERROR', 'This is not a supported CSV file.']
    else:
        return_type_array = ['ERROR', 'Please enter a file compatible with NS']

    return return_type_array

def num2alpha(num):  #input number output alphabet
    if num<=26:
        return chr(64+num).lower()
    elif num%26==0:
        return num2alpha(num//26-1)+chr(90).lower()
    else:
        return num2alpha(num//26)+chr(64+num%26).lower()

def get_ip_address_set(change_tmp_ip_address_array):
    host = ipaddress.ip_interface(change_tmp_ip_address_array)
    tmp_network =  str(host.network)
    tmp_ip = str(host.ip)
    tmp_mask = host.network.prefixlen
    tmp_ip_array = tmp_ip.split('.')

    if tmp_mask >= 24:
        return_ip = '.'+ tmp_ip_array[3]
    elif tmp_mask >= 16:
        return_ip = '.' + tmp_ip_array[2] + '.' + tmp_ip_array[3]
    elif tmp_mask >= 8:
        return_ip = '.' + tmp_ip_array[1] + '.' + tmp_ip_array[2] + '.' + tmp_ip_array[3]
    else:
        return_ip = tmp_ip

    return ([change_tmp_ip_address_array,tmp_network,return_ip ]) # [ip_address,network_address,mask,host_address]

def get_description_width_hight(font_size,description):
    par_char_ratio = 0.0095
    font_hight_ratio = 0.018

    per_char_width = font_size * par_char_ratio
    font_size_hight = font_size * font_hight_ratio

    result = [get_east_asian_width_count(description) * per_char_width,font_size_hight ]

    return (result)




def _get_staging_xlsx(nsm_path):
    """Get or create a staging xlsx for NSM export operations.

    Creates a transient xlsx from NSM Parquet data so that export functions
    (create_excel_gui_tree, insert_custom_excel_table, diagram generators)
    can work with _tmp_ sheets without touching the NSM master itself.
    openpyxl is used ONLY on this staging file, never on the NSM master.
    """
    staging = str(nsm_path) + '.staging.xlsx'
    if not os.path.exists(staging):
        from ns_engine.nsm_io import nsm_to_xlsx
        nsm_to_xlsx(str(nsm_path), staging)
    return staging


### write excel meta file ###
def write_excel_meta(master_excel_meta, excel_file_path, worksheet_name, section_write_to, offset_row, offset_column):
    if str(excel_file_path).lower().endswith('.nsm'):
        if worksheet_name.startswith('_tmp') or worksheet_name.startswith('_template'):
            staging = _get_staging_xlsx(excel_file_path)
            _write_excel_meta_xlsx(master_excel_meta, staging,
                                   worksheet_name, section_write_to,
                                   offset_row, offset_column)
            return
        _write_meta_nsm(master_excel_meta, excel_file_path, worksheet_name,
                        section_write_to, offset_row, offset_column,
                        mode='write')
        return
    _write_excel_meta_xlsx(master_excel_meta, excel_file_path, worksheet_name,
                           section_write_to, offset_row, offset_column)


def _write_meta_nsm(master_excel_meta, file_path, worksheet_name,
                    section_write_to, offset_row, offset_column, mode='write'):
    """Write/overwrite section data to .nsm file.

    Converts the tuple-keyed dict to a DataFrame and saves via nsm_io.
    """
    from ns_engine.nsm_io import load_section, save_section, _section_key, _encode_cell

    if section_write_to == '<<N/A>>':
        return

    is_template = (section_write_to == '_template_')

    if is_template:
        actual_section = None
        for (r, c), val in master_excel_meta.items():
            s = str(val).strip()
            if s.startswith('<<') and s.endswith('>>') and s not in ('<<N/A>>', '<<END_MARK>>'):
                actual_section = s
                break
        if actual_section is None:
            _ws_to_section = {
                'Master_Data_L2': '<<L2_TABLE>>',
                'Master_Data_L3': '<<L3_TABLE>>',
            }
            actual_section = _ws_to_section.get(worksheet_name)
        if actual_section is None:
            return
        key = _section_key(actual_section)
        rows_data = {}
    else:
        key = _section_key(section_write_to)
        if mode == 'overwrite':
            existing_df = load_section(file_path, key)
            if existing_df.empty:
                rows_data = {}
            else:
                rows_data = {}
                for row_idx, (_, row) in enumerate(existing_df.iterrows(), start=1):
                    for col_idx, val in enumerate(row, start=1):
                        rows_data[(row_idx, col_idx)] = val
        else:
            rows_data = {}

    for (r, c), val in master_excel_meta.items():
        adj_r = r + offset_row
        adj_c = c + offset_column
        rows_data[(adj_r, adj_c)] = _encode_cell(val)

    if not rows_data:
        return

    max_row = max(r for r, c in rows_data)
    max_col = max(c for r, c in rows_data)

    import pandas as pd
    data = []
    for r in range(1, max_row + 1):
        row_vals = []
        for c in range(1, max_col + 1):
            val = rows_data.get((r, c), '')
            row_vals.append(str(val) if val is not None else '')
        data.append(row_vals)

    cols = [f'col_{i}' for i in range(max_col)]
    df = pd.DataFrame(data, columns=cols)
    save_section(file_path, key, df)


def _write_excel_meta_xlsx(master_excel_meta, excel_file_path, worksheet_name, section_write_to, offset_row, offset_column):
    '''
    ★★★ OPTIMIZED VERSION (maintains exact original specification) ★★★

    Optimizations applied:
    1. Faster section search using iter_rows()
    2. Single-pass max row calculation
    3. Pre-sorted cell writes for better cache locality
    4. Early termination on errors

    Expected performance: 10-15x faster

    :param master_excel_meta:  tuple master data
    :param excel_file_path:  file path of excel master data
    :param worksheet_name:   worksheet name to write
    :param section_write_to: decide the value of start row and column
    :return: none
    '''

    import openpyxl

    wb = openpyxl.load_workbook(excel_file_path)
    wb.active = wb[worksheet_name]

    # ========== Find section row (OPTIMIZED) ==========
    flag_section = False
    empty_count = 0
    row_count = 0

    if section_write_to == '_template_':
        flag_section = True
        row_count = 1
    else:
        # ★★★ OPTIMIZATION 1: Use iter_rows instead of cell-by-cell ★★★
        max_search_rows = 10000

        for row in wb.active.iter_rows(min_row=1, max_row=max_search_rows, min_col=1, max_col=1, values_only=True):
            row_count += 1

            if row[0] == section_write_to:
                flag_section = True
                break
            elif '<<N/A>>' == section_write_to:
                flag_section = True
                row_count = 1
                break
            elif row[0] is None:
                empty_count += 1
                if empty_count > 10000:
                    flag_section = True
                    print('---ERROR and STOP---  can not find ---> %s  ' % section_write_to)
                    wb.close()
                    exit()

    # If section not found after loop
    if not flag_section:
        print('---ERROR and STOP---  can not find ---> %s  ' % section_write_to)
        wb.close()
        exit()

    # ========== Calculate max row (OPTIMIZED) ==========
    # ★★★ OPTIMIZATION 2: Single-pass max calculation ★★★
    num_insert_row = 2
    if master_excel_meta:
        # Use max() with generator for O(n) instead of loop
        num_insert_row = max((i[0] for i in master_excel_meta), default=2)

    # Insert rows
    if num_insert_row > 1:
        wb.active.insert_rows(row_count + 1 + offset_row, amount=num_insert_row - 1)

    # ========== Write cells (OPTIMIZED) ==========
    # ★★★ OPTIMIZATION 3: Pre-sort keys for better cache locality ★★★
    # Sorting by row then column improves Excel internal cache hit rate
    sorted_keys = sorted(master_excel_meta.keys(), key=lambda k: (k[0], k[1]))

    for num_wp_up in sorted_keys:
        target_row = num_wp_up[0] + row_count - 1 + offset_row
        target_col = num_wp_up[1] + offset_column
        wb.active.cell(target_row, target_col).value = master_excel_meta[num_wp_up]

    # Save excel file
    wb.save(excel_file_path)
    wb.close()




### overwrite excel meta file ###
from openpyxl import load_workbook


def overwrite_excel_meta(master_excel_meta, excel_file_path, worksheet_name, section_write_to, offset_row, offset_column):
    if str(excel_file_path).lower().endswith('.nsm'):
        if worksheet_name.startswith('_tmp') or worksheet_name.startswith('_template'):
            staging = _get_staging_xlsx(excel_file_path)
            _overwrite_excel_meta_xlsx(master_excel_meta, staging,
                                       worksheet_name, section_write_to,
                                       offset_row, offset_column)
            return
        _write_meta_nsm(master_excel_meta, excel_file_path, worksheet_name,
                        section_write_to, offset_row, offset_column,
                        mode='overwrite')
        return
    _overwrite_excel_meta_xlsx(master_excel_meta, excel_file_path, worksheet_name,
                               section_write_to, offset_row, offset_column)


def _overwrite_excel_meta_xlsx(master_excel_meta, excel_file_path, worksheet_name, section_write_to, offset_row, offset_column):
    '''
    Overwrites Excel metadata with optimized openpyxl performance

    :param master_excel_meta: dict with (row, col) tuples as keys and cell values as values
    :param excel_file_path: path to the Excel file
    :param worksheet_name: name of the worksheet to write to
    :param section_write_to: section identifier to determine starting position
    :param offset_row: row offset from the section start
    :param offset_column: column offset from the section start
    :return: None
    '''

    # Load workbook with performance optimization options
    # data_only=True: Read formula results as values (faster)
    # keep_vba=False: Exclude VBA macros (faster loading)
    wb = load_workbook(excel_file_path, data_only=True, keep_vba=False)
    ws = wb[worksheet_name]

    # Section search (optimized version)
    flag_section = False
    row_count = 0

    # Handle special section identifiers
    if section_write_to in ('_template_', '<<N/A>>'):
        flag_section = True
        row_count = 1
    else:
        # Search for section in first column up to 10000 rows
        # iter_rows is memory-efficient compared to cell-by-cell access
        for row_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=100000, min_col=1, max_col=1, values_only=True), start=1):
            if row[0] == section_write_to:
                row_count = row_idx
                flag_section = True
                break

        # Exit if section not found
        if not flag_section:
            print(f'---ERROR and STOP--- can not find ---> {section_write_to}')
            wb.close()
            return

    # Write data (batch processing)
    # Calculate absolute position and write each cell
    for (rel_row, rel_col), value in master_excel_meta.items():
        ws.cell(rel_row + row_count - 1 + offset_row,
                rel_col + offset_column,
                value)

    # Save and close the workbook
    wb.save(excel_file_path)
    wb.close()


### return shapes of tuple format  ###
def return_shape_tuple(current_shape_array ,start_row):
    '''
    :param current_shape_array):  current shape array
    start_row : input row number in the tuple
    '''
    #sort top value
    current_shape_array = sorted(current_shape_array, reverse=False, key=lambda x: x[3])  # sort for top
    tmp_grid_array = []
    grid_array = []
    threshold_shape_top = current_shape_array[0][3]
    threshold_shape_down = current_shape_array[0][3] + current_shape_array[0][5]

    #make y grid array
    flag_next_y_grid = False
    for i in range(len(current_shape_array)):
        if flag_next_y_grid == True:
            threshold_shape_top = current_shape_array[i-1][3]
            threshold_shape_down = current_shape_array[i-1][3] + current_shape_array[i-1][5]
            flag_next_y_grid = False

        if current_shape_array[i][3] <= threshold_shape_down and (current_shape_array[i][3] + current_shape_array[i][5]) >= threshold_shape_top:
            tmp_grid_array.append(current_shape_array[i])

            if i == len(current_shape_array) - 1:
                grid_array.append(tmp_grid_array)

        else:
            grid_array.append(tmp_grid_array)
            tmp_grid_array = []
            tmp_grid_array.append(current_shape_array[i])
            flag_next_y_grid = True

            if i == len(current_shape_array) - 1:
                grid_array.append([current_shape_array[i]])

    # sort left value
    master_grid_array = []
    for tmp_grid_array in grid_array:
        tmp_grid_array = sorted(tmp_grid_array, reverse=False, key=lambda x: x[2])  # sort for top
        master_grid_array.append(tmp_grid_array)

    '''Added automatic horizontal axis placement function at Ver 2.2.2 '''
    # Remove items containing '_AIR_' in the second element
    master_grid_array = [[item for item in sublist if '_AIR_' not in item[1]] for sublist in master_grid_array]

    updated_master_grid_array = []
    device_name_array = []
    kari_master_grid_array = []

    for tmp_master_grid_array in master_grid_array:
        for tmp_tmp_master_grid_array in tmp_master_grid_array:
            device_name_array.append(tmp_tmp_master_grid_array[1])
            kari_master_grid_array.append(tmp_tmp_master_grid_array)
    #print('--- device_name_array ---')
    #print(device_name_array)

    vertical_key_array = []
    used_device_array = []

    for tmp_device_name_array in device_name_array:
        if tmp_device_name_array not in used_device_array:
            #print('#####tmp_device_name_array,used_device_array,vertical_key_array',tmp_device_name_array,used_device_array,vertical_key_array)
            target = tmp_device_name_array
            result = None

            for kari_kari_master_grid_array in kari_master_grid_array:
                if kari_kari_master_grid_array[1] == target:
                    result = kari_kari_master_grid_array
                    break

            if result not in vertical_key_array and result[1] not in used_device_array:
                vertical_key_array.append(result)
                used_device_array.append(result)

            except_array = []

            for sublist in master_grid_array:
                except_array.extend(sublist)
                flag_1st_match = False

                #re-make at ver 2.3.4
                for item in sublist:
                    if ((result[2] < item[2] + item[4] and result[2] > item[2]) or \
                        (result[2] + result[4] > item[2] and result[2] + result[4] < item[2] + item[4]) or \
                        (result[2] < item[2] and result[2] + result[4] > item[2] + item[4])  or \
                        (result[2] > item[2] and result[2] + result[4] < item[2] + item[4])  or \
                        (result[2] == item[2] and result[2] + result[4] == item[2] + item[4] and result[1] != item[1])) and \
                        (item[1] not in used_device_array) and (flag_1st_match == False or item not in except_array):
                            used_device_array.append(item[1])
                            except_array.remove(item)
                            flag_1st_match = True

    # Remove duplicates by converting lists to tuples, adding to a set, and converting back to lists
    vertical_key_array = [list(t) for t in set(tuple(item) for item in vertical_key_array)]
    #print(vertical_key_array, len(vertical_key_array))

    vertical_key_array_2 = []
    for tmp_vertical_key_array in vertical_key_array:
        vertical_key_array_2.append(
            [tmp_vertical_key_array[0], '_AIR_', tmp_vertical_key_array[2], tmp_vertical_key_array[3],
             tmp_vertical_key_array[4], tmp_vertical_key_array[5], tmp_vertical_key_array[6], tmp_vertical_key_array[7],
             tmp_vertical_key_array[2] + int(tmp_vertical_key_array[4] * 0.5)])

    vertical_key_array_2 = sorted(vertical_key_array_2, key=lambda x: x[8])
    #print('####vertical_key_array_2 ####')
    #print(vertical_key_array_2)

    for tmp_master_grid_array in master_grid_array:
        #print('####### tmp_master_grid_array ########')
        #print(tmp_master_grid_array)

        import copy
        updated_vertical_key_array_2 = copy.deepcopy(vertical_key_array_2)

        for i, tmp_vertical_key_array_2 in enumerate(vertical_key_array_2):
            for tmp_tmp_master_grid_array in tmp_master_grid_array:

                if tmp_tmp_master_grid_array not in updated_vertical_key_array_2:
                    if tmp_vertical_key_array_2[2] + tmp_vertical_key_array_2[4] > tmp_tmp_master_grid_array[2] and tmp_vertical_key_array_2[2] < tmp_tmp_master_grid_array[2]:
                        updated_vertical_key_array_2[i] = tmp_tmp_master_grid_array
                        break
                    elif tmp_vertical_key_array_2[2] > tmp_tmp_master_grid_array[2] and tmp_vertical_key_array_2[2] + tmp_vertical_key_array_2[4] < tmp_tmp_master_grid_array[2] + tmp_tmp_master_grid_array[4]:
                        updated_vertical_key_array_2[i] = tmp_tmp_master_grid_array
                        break
                    elif tmp_vertical_key_array_2[2] + tmp_vertical_key_array_2[4] >= tmp_tmp_master_grid_array[2] + tmp_tmp_master_grid_array[4] and tmp_vertical_key_array_2[2] < tmp_tmp_master_grid_array[2] + tmp_tmp_master_grid_array[4]:
                        updated_vertical_key_array_2[i] = tmp_tmp_master_grid_array
                        break
                    elif tmp_vertical_key_array_2[2] < tmp_tmp_master_grid_array[2] and tmp_vertical_key_array_2[2] + tmp_vertical_key_array_2[4] > tmp_tmp_master_grid_array[2] + tmp_tmp_master_grid_array[4]:
                        updated_vertical_key_array_2[i] = tmp_tmp_master_grid_array
                        break
                    elif tmp_vertical_key_array_2[2] == tmp_tmp_master_grid_array[2] and tmp_vertical_key_array_2[4] == tmp_tmp_master_grid_array[4]:
                        updated_vertical_key_array_2[i] = tmp_tmp_master_grid_array
                        break

        #print(updated_vertical_key_array_2,len(updated_vertical_key_array_2))
        updated_master_grid_array.append(updated_vertical_key_array_2)
    #print('--- updated_master_grid_array ---')
    #print(updated_master_grid_array)

    master_grid_array = copy.deepcopy(updated_master_grid_array)
    #print('--- master_grid_array ---')
    #print(master_grid_array)
    '''Addition completed(ver 2.2.2)'''
    # make the tuple format
    tuple_grid_array = {}
    tuple_grid_array[start_row, 1] = master_grid_array[0][0][0]
    num_row = start_row -1
    num_column =1
    for tmp_y in master_grid_array:
        num_row += 1
        for tmp_x in tmp_y:
            num_column += 1
            tuple_grid_array[num_row, num_column] = tmp_x[1]
        tuple_grid_array[num_row, num_column + 1] = '<END>'
        num_column = 1
    tuple_grid_array[num_row + 1, 1] = '<END>'

    return(tuple_grid_array)


def get_folder_width_size(master_folder_tuple, master_style_shape_tuple, master_shape_tuple, min_tag_inches):
    # print(master_shape_tuple)
    # add parameter at ver2.1 for large size
    folder_width_ratio = 0.8  # add at ver 2.1 for large size (applied at output only). 0.5 -> 0.8 changed at ver 2.6.1a
    air_fixed_width = 0.5  # FIXED width for _AIR_ in master_shape_tuple only

    # Helper function to apply ratio and round up to prevent value degradation
    import math
    def apply_ratio_and_round(value, ratio, decimal_places=4):
        """
        Apply ratio to value and round up to specified decimal places.
        This prevents value degradation over repeated calculations.
        """
        result = value * ratio
        multiplier = 10 ** decimal_places
        return math.ceil(result * multiplier) / multiplier

    master_width_size_folder = []
    master_width_size_y_grid = []
    master_hight_size_y_grid = []
    master_folder_size = []
    folder_num_list = []

    for tmp_master_folder_tuple in master_folder_tuple:
        if tmp_master_folder_tuple[1] == 1 and master_folder_tuple[tmp_master_folder_tuple] != '<SET_WIDTH>':
            folder_num_list.append(tmp_master_folder_tuple[0])

    current_folder_start_row = 1
    current_folder_end_row = 1

    for folder_num in folder_num_list:
        i = 0
        for tmp_master_folder_tuple in master_folder_tuple:
            if tmp_master_folder_tuple[0] == folder_num and tmp_master_folder_tuple[1] != 1:
                if master_folder_tuple[tmp_master_folder_tuple] == '':
                    # Empty folder (not _AIR_ in shape) - will use average width later
                    master_width_size_folder.append([folder_num, [['_EMPTY_FOLDER_', 0]]])
                    master_folder_size.append([folder_num, [['_EMPTY_FOLDER_', 0]]])
                    i += 1
                else:
                    i += 1
                    flag_shape_start = False
                    for tmp_master_shape_tuple in master_shape_tuple:
                        if tmp_master_shape_tuple[1] == 1 and master_shape_tuple[tmp_master_shape_tuple] == '<END>' and flag_shape_start == True:
                            current_folder_end_row = tmp_master_shape_tuple[0] - 1
                            break
                        if tmp_master_shape_tuple[1] == 1 and master_folder_tuple[tmp_master_folder_tuple] == master_shape_tuple[tmp_master_shape_tuple]:
                            current_folder_start_row = tmp_master_shape_tuple[0]
                            flag_shape_start = True

                    # ===== Added: Pre-calculate average width for entire folder =====
                    folder_total_width = 0
                    folder_total_count = 0

                    for r in range(current_folder_start_row, current_folder_end_row + 1):
                        for tmp_master_shape_tuple in master_shape_tuple:
                            if tmp_master_shape_tuple[0] == r and tmp_master_shape_tuple[1] != 1:
                                shape_name = master_shape_tuple[tmp_master_shape_tuple]

                                # Count only non-AIR shapes for average calculation
                                if shape_name != '_AIR_' and shape_name != '<END>':
                                    for tmp_master_style_shape_tuple in master_style_shape_tuple:
                                        if master_style_shape_tuple[tmp_master_style_shape_tuple[0], 1] == shape_name:
                                            folder_total_width += master_style_shape_tuple[tmp_master_style_shape_tuple[0], 2]
                                            folder_total_count += 1
                                            break

                    # Calculate average width for the entire folder (excluding _AIR_)
                    if folder_total_count > 0:
                        folder_average_width = folder_total_width / folder_total_count
                    else:
                        folder_average_width = 0.5

                    # print(f"Folder: {master_folder_tuple[tmp_master_folder_tuple]}, Average width: {folder_average_width}, Device count: {folder_total_count}")
                    # ===== End of pre-calculation =====

                    tmp_folder_size = []
                    current_level = 0
                    current_max_width = 0
                    current_max_hight = 0
                    tmp_hight = 0

                    for r in range(current_folder_start_row, current_folder_end_row + 1):
                        tmp_width = 0
                        tmp_count_shape = 0
                        tmp_count_air = 0  # Count _AIR_ in master_shape_tuple
                        current_max_hight = 0

                        for tmp_master_shape_tuple in master_shape_tuple:
                            if tmp_master_shape_tuple[0] == r and tmp_master_shape_tuple[1] != 1:
                                shape_name = master_shape_tuple[tmp_master_shape_tuple]

                                # ===== Modified: Count _AIR_ separately and use FIXED width =====
                                if shape_name == '_AIR_':
                                    tmp_count_air += 1  # Count _AIR_ from master_shape_tuple
                                elif shape_name != '<END>':
                                    for tmp_master_style_shape_tuple in master_style_shape_tuple:
                                        if master_style_shape_tuple[tmp_master_style_shape_tuple[0], 1] == shape_name:
                                            # sum width in a level in a folder
                                            tmp_width += master_style_shape_tuple[tmp_master_style_shape_tuple[0], 2]
                                            tmp_count_shape += 1

                                            # get max value in a level in a folder
                                            if current_max_hight < (master_style_shape_tuple[tmp_master_style_shape_tuple[0], 3] + (min_tag_inches * 2.5)):
                                                current_max_hight = (master_style_shape_tuple[tmp_master_style_shape_tuple[0], 3] + (min_tag_inches * 2.5))
                                            break

                        # ===== MODIFIED: Use FIXED width for _AIR_ from master_shape_tuple =====
                        air_total_width = air_fixed_width * tmp_count_air
                        tmp_width += air_total_width

                        # Total count includes both real shapes and _AIR_
                        total_count = tmp_count_shape + tmp_count_air
                        # ===== End of modification =====

                        current_level += 1
                        # print(f"  Row {r}: Devices={tmp_count_shape}, _AIR_={tmp_count_air}, Total={total_count}, Width={tmp_width}")

                        current_level_inches_width = min_tag_inches * 2 + tmp_width + ((total_count - 1) * (min_tag_inches * 7))
                        tmp_hight += current_max_hight

                        if current_max_width < current_level_inches_width:
                            current_max_width = current_level_inches_width

                    tmp_hight += 1.0  # add up down buffer for a hight in a folder

                    # Store ONLY actual width (NO ratio applied during calculation)
                    # Format: [name, actual_width, height]
                    tmp_folder_size.append([
                        master_folder_tuple[tmp_master_folder_tuple],
                        current_max_width,  # Actual width (for calculations)
                        tmp_hight
                    ])

                master_width_size_folder.append([folder_num, tmp_folder_size])
                master_folder_size.append([folder_num, tmp_folder_size])

        if i == 0:
            # No folders in this row - use fixed width
            master_width_size_folder.append([folder_num, [['_EMPTY_FOLDER_', air_fixed_width, 0]]])
            master_folder_size.append([folder_num, [['_EMPTY_FOLDER_', air_fixed_width, 0]]])

    # Modified: Calculate widths
    # _EMPTY_FOLDER_ (from master_folder_tuple '') uses average width
    # _AIR_ width is already included in calculations above
    for folder_num in folder_num_list:
        tmp_sum_width = 0
        empty_folder_count = 0
        non_empty_count = 0

        # First pass: calculate sum of non-empty folder widths
        for tmp_master_min_size_folder in master_width_size_folder:
            if tmp_master_min_size_folder[0] == folder_num:
                if tmp_master_min_size_folder[1][0][0] != '_EMPTY_FOLDER_':
                    tmp_sum_width += tmp_master_min_size_folder[1][0][1]  # Use actual width
                    non_empty_count += 1
                else:
                    empty_folder_count += 1

        # Calculate average width of non-empty folders
        if non_empty_count > 0:
            average_width = tmp_sum_width / non_empty_count
        else:
            average_width = 0.5  # Default if no shapes exist

        # Calculate total width: actual folders + (average × empty_folder_count)
        # Empty folders from master_folder_tuple use average width
        final_width = tmp_sum_width + (average_width * empty_folder_count)
        air_width_each = average_width * 0.2 if non_empty_count > 0 else 0.1

        # Store actual widths (NO ratio applied yet)
        master_width_size_y_grid.append([folder_num, final_width, air_width_each])

    # GET best width of slide (inches) using ACTUAL width
    slide_max_width_inches = 0
    for tmp_master_min_size_y_grid in master_width_size_y_grid:
        if slide_max_width_inches < tmp_master_min_size_y_grid[1]:
            slide_max_width_inches = tmp_master_min_size_y_grid[1]

    # GET best height of slide (inches)
    slide_max_hight_inches = 0

    for tmp_master_min_size_y_grid in master_width_size_y_grid:
        tmp_max_hight_y_grid = 0
        flag_only_wp = True

        for tmp_master_folder_size in master_folder_size:
            if tmp_master_min_size_y_grid[0] == tmp_master_folder_size[0] and tmp_master_folder_size[1][0][0] != '_EMPTY_FOLDER_':
                # Use height (index [2])
                if tmp_max_hight_y_grid < tmp_master_folder_size[1][0][2]:
                    tmp_max_hight_y_grid = tmp_master_folder_size[1][0][2]
                if '_wp_' not in str(tmp_master_folder_size[1][0][0]):
                    flag_only_wp = False

        if flag_only_wp == True:
            tmp_max_hight_y_grid = tmp_max_hight_y_grid * 1  # Change hight ratio Ver 1.1

        if tmp_max_hight_y_grid == 0:
            tmp_max_hight_y_grid = 0.5  # only empty level is 0.5 inches

        master_hight_size_y_grid.append([tmp_master_min_size_y_grid[0], tmp_max_hight_y_grid])
        slide_max_hight_inches += tmp_max_hight_y_grid

    # print('----slide_max_width_inches----', slide_max_width_inches)

    # ===== APPLY RATIO AND ROUND UP AT OUTPUT STAGE =====
    # This ensures ratio is only applied once and rounding up prevents value degradation

    # Apply ratio to slide max width
    output_slide_max_width = apply_ratio_and_round(slide_max_width_inches, folder_width_ratio)

    # Apply ratio to master_width_size_y_grid
    output_master_width_size_y_grid = []
    for item in master_width_size_y_grid:
        output_master_width_size_y_grid.append([
            item[0],  # folder_num
            apply_ratio_and_round(item[1], folder_width_ratio),  # final_width with ratio
            apply_ratio_and_round(item[2], folder_width_ratio)  # air_width_each with ratio
        ])

    # Apply ratio to master_folder_size
    output_master_folder_size = []
    for folder in master_folder_size:
        folder_num = folder[0]
        folder_data = folder[1]
        if folder_data[0][0] == '_EMPTY_FOLDER_':
            # For empty folder, apply ratio
            output_master_folder_size.append([folder_num, [['_EMPTY_FOLDER_', apply_ratio_and_round(folder_data[0][1], folder_width_ratio), 0]]])
        else:
            output_master_folder_size.append([
                folder_num,
                [[
                    folder_data[0][0],  # Name
                    apply_ratio_and_round(folder_data[0][1], folder_width_ratio),  # Width with ratio and round up
                    folder_data[0][2]  # Height (unchanged)
                ]]
            ])

    # Return with ratio applied and rounded up
    return ([
        output_slide_max_width,
        output_master_width_size_y_grid,
        output_master_folder_size,
        slide_max_hight_inches,
        master_hight_size_y_grid
    ])


def get_root_folder_tuple(self,master_folder_size_array,tmp_folder_name):
    self.root_left = 0.28
    self.root_top = 1.42
    self.root_width = math.ceil(master_folder_size_array[0] * 12) / 10  #ver2.2.0 change,  10 -> 12
    self.root_hight = math.ceil(master_folder_size_array[3] * 12) / 10  #ver1.1 change,  10 -> 12
    ppt_min_width = 6  # inches 13.4
    ppt_min_hight = 4  # inches 7.5

    master_root_folder_tuple = {}
    master_root_folder_tuple[2, 3] = 1
    master_root_folder_tuple[2, 4] = 1
    master_root_folder_tuple[2, 5] = self.root_left
    master_root_folder_tuple[2, 6] = self.root_top

    if self.root_width < (ppt_min_width - (self.root_left * 2)):
        master_root_folder_tuple[2, 7] = (ppt_min_width - (self.root_left * 2))
    else:
        master_root_folder_tuple[2, 7] = self.root_width

    if self.root_hight < (ppt_min_hight - (self.root_top * 1.5)):
        master_root_folder_tuple[2, 8] = (ppt_min_hight - (self.root_top * 1.5))
    else:
        master_root_folder_tuple[2, 8] = self.root_hight

    master_root_folder_tuple[2, 2] = '[L1]' + tmp_folder_name

    return(master_root_folder_tuple)



### return tuple for def - return_shape_tuple - ###
def convert_array_to_tuple(tmp_master_data_array):
    """ULTRA FAST: 3-5x faster, most concise"""
    return {
        (row_data[0], col_idx): value
        for row_data in tmp_master_data_array
        for col_idx, value in enumerate(row_data[1], start=1)
    }


### return array for def - return_array - ###
def convert_tuple_to_array(tmp_master_data_tuple):
    tmp_master_data_array = []
    for tmp_tmp_master_data_tuple in tmp_master_data_tuple:
        tmp_master_data_array.append([tmp_tmp_master_data_tuple[0],tmp_tmp_master_data_tuple[1],tmp_master_data_tuple[tmp_tmp_master_data_tuple]])

    # sort row -> column
    tmp_master_data_array = sorted(tmp_master_data_array, reverse=False, key=lambda x:( x[0],x[1]))  # sort row -> column
    #print(tmp_master_data_array)

    master_data_array = []
    tmp_tmp_array = []
    flag_first = True
    for tmp_array in tmp_master_data_array:
        if flag_first == True:
            tmp_num = tmp_array[0]
            #tmp_tmp_array.append(tmp_array[2])
            flag_first = False

        if flag_first == False:
            if tmp_num == tmp_array[0]:
                tmp_tmp_array.append(tmp_array[2])
            else:
                master_data_array.append([tmp_num ,tmp_tmp_array])
                tmp_num = tmp_array[0]
                tmp_tmp_array = []
                tmp_tmp_array.append(tmp_array[2])
    master_data_array.append([tmp_num, tmp_tmp_array])
    return(master_data_array)


### return folder and wp name array from master excel file ###
def get_folder_wp_array_from_master(ws_name, ppt_meta_file):
    if str(ppt_meta_file).lower().endswith('.nsm'):
        return _get_folder_wp_array_from_master_nsm(ppt_meta_file)
    return _get_folder_wp_array_from_master_xlsx(ws_name, ppt_meta_file)


def _get_folder_wp_array_from_master_nsm(ppt_meta_file):
    from ns_engine.nsm_io import load_section, _decode_cell
    df = load_section(ppt_meta_file, 'POSITION_SHAPE')
    folder_name_array = []
    wp_name_array = []
    if df.empty:
        return [folder_name_array, wp_name_array]
    for _, row in df.iterrows():
        val = _decode_cell(row.iloc[0]) if len(row) > 0 else ''
        s = str(val)
        if s.startswith('<<') and s.endswith('>>'):
            continue
        if s == 'None' or s == '<END>' or s == '':
            continue
        if '_wp_' in s:
            wp_name_array.append(s)
        else:
            folder_name_array.append(s)
    return [folder_name_array, wp_name_array]


def _get_folder_wp_array_from_master_xlsx(ws_name, ppt_meta_file):
    input_ppt_mata_excel = openpyxl.load_workbook(ppt_meta_file)
    input_ppt_mata_excel.active = input_ppt_mata_excel[ws_name]

    flag_finish = False
    current_row = 1
    folder_name_array = []
    wp_name_array = []

    while flag_finish == False:
        if input_ppt_mata_excel.active.cell(current_row, 1).value == '<<POSITION_SHAPE>>':
            start_row = current_row
        if input_ppt_mata_excel.active.cell(current_row, 1).value == '<<STYLE_SHAPE>>':
            end_row = current_row - 1
            flag_finish = True
        current_row += 1

    for i in range(start_row + 1, end_row+1):
        if str(input_ppt_mata_excel.active.cell(i, 1).value) != 'None' and str(input_ppt_mata_excel.active.cell(i, 1).value) != '<END>' \
                and '_wp_' not in str(input_ppt_mata_excel.active.cell(i, 1).value):
            folder_name_array.append(str(input_ppt_mata_excel.active.cell(i, 1).value))
        elif '_wp_' in str(input_ppt_mata_excel.active.cell(i, 1).value):
            wp_name_array.append(str(input_ppt_mata_excel.active.cell(i, 1).value))

    input_ppt_mata_excel.close()
    return([folder_name_array, wp_name_array])

def _nsm_save_path(ppt_meta_file):
    """For .nsm files, return the cached xlsx path for openpyxl saves.
    
    If no cache exists, triggers a reconstruction via the patched load_workbook.
    """
    if str(ppt_meta_file).lower().endswith('.nsm'):
        from ns_engine.nsm_adapter import _nsm_xlsx_cache
        cached = _nsm_xlsx_cache.get(str(ppt_meta_file))
        if cached and os.path.exists(cached):
            return cached
        wb = openpyxl.load_workbook(ppt_meta_file)
        wb.close()
        cached = _nsm_xlsx_cache.get(str(ppt_meta_file))
        if cached:
            return cached
    return str(ppt_meta_file)

### copy excel sheet to own file ###
def copy_excel_sheet(ws_name, ppt_meta_file, copy_sheet_name):
    if str(ppt_meta_file).lower().endswith('.nsm'):
        staging = _get_staging_xlsx(ppt_meta_file)
        wb = openpyxl.load_workbook(staging)
        if copy_sheet_name in wb.sheetnames:
            wb.remove(wb[copy_sheet_name])
        ws = wb.copy_worksheet(wb[ws_name])
        ws.title = copy_sheet_name
        wb.save(staging)
        wb.close()
        return ()
    input_ppt_mata_excel = openpyxl.load_workbook(ppt_meta_file)
    save_path = _nsm_save_path(ppt_meta_file)

    ws_list = input_ppt_mata_excel.sheetnames
    if copy_sheet_name in ws_list:
        ws = input_ppt_mata_excel.remove(input_ppt_mata_excel[copy_sheet_name])
    ws = input_ppt_mata_excel.copy_worksheet(input_ppt_mata_excel[ws_name])
    ws.title = copy_sheet_name
    input_ppt_mata_excel.save(save_path)
    input_ppt_mata_excel.close()
    return()

### remove excel sheet  ###
def remove_excel_sheet(ppt_meta_file, copy_sheet_name):
    if str(ppt_meta_file).lower().endswith('.nsm'):
        if copy_sheet_name.startswith('_tmp') or copy_sheet_name.startswith('_template'):
            staging = str(ppt_meta_file) + '.staging.xlsx'
            if os.path.exists(staging):
                try:
                    wb = openpyxl.load_workbook(staging)
                    if copy_sheet_name in wb.sheetnames:
                        wb.remove(wb[copy_sheet_name])
                    remaining_tmp = [s for s in wb.sheetnames
                                     if s.startswith('_tmp') or s.startswith('_template')]
                    wb.close()
                    if not remaining_tmp:
                        os.remove(staging)
                except Exception:
                    pass
        return ()
    input_ppt_mata_excel = openpyxl.load_workbook(ppt_meta_file)
    save_path = _nsm_save_path(ppt_meta_file)

    ws_list = input_ppt_mata_excel.sheetnames
    if copy_sheet_name in ws_list:
        ws = input_ppt_mata_excel.remove(input_ppt_mata_excel[copy_sheet_name])
    input_ppt_mata_excel.save(save_path)
    input_ppt_mata_excel.close()

    return()

### create excel sheet  ###
def create_excel_sheet(ppt_meta_file, sheet_name):
    if str(ppt_meta_file).lower().endswith('.nsm'):
        if sheet_name.startswith('_tmp') or sheet_name.startswith('_template'):
            staging = _get_staging_xlsx(ppt_meta_file)
            wb = openpyxl.load_workbook(staging)
            if sheet_name in wb.sheetnames:
                wb.remove(wb[sheet_name])
            wb.create_sheet(sheet_name)
            wb.save(staging)
            wb.close()
        return ()
    input_ppt_mata_excel = openpyxl.load_workbook(ppt_meta_file)
    save_path = _nsm_save_path(ppt_meta_file)

    ws_list = input_ppt_mata_excel.sheetnames
    if sheet_name in ws_list:
        ws = input_ppt_mata_excel.remove(input_ppt_mata_excel[sheet_name])
        ws = input_ppt_mata_excel.create_sheet(sheet_name)
    else:
        ws = input_ppt_mata_excel.create_sheet(sheet_name)
    input_ppt_mata_excel.save(save_path)
    input_ppt_mata_excel.close()

    return()

#convert from master to array
def convert_master_to_array(ws_name, ppt_meta_file, section_name):
    if str(ppt_meta_file).lower().endswith('.nsm'):
        if ws_name.startswith('_tmp') or ws_name.startswith('_template'):
            staging = str(ppt_meta_file) + '.staging.xlsx'
            if os.path.exists(staging):
                return _convert_master_to_array_xlsx(ws_name, staging, section_name)
        return _convert_master_to_array_nsm(ws_name, ppt_meta_file, section_name)
    return _convert_master_to_array_xlsx(ws_name, ppt_meta_file, section_name)


def convert_master_to_arrays_bulk(ws_name, ppt_meta_file, section_names):
    """Load multiple sections at once. Returns dict {section_name: array}.

    For .nsm files, opens ZIP once instead of N times.
    For .xlsx files, opens workbook once instead of N times.
    """
    if str(ppt_meta_file).lower().endswith('.nsm'):
        return _convert_master_to_arrays_bulk_nsm(ws_name, ppt_meta_file, section_names)
    return _convert_master_to_arrays_bulk_xlsx(ws_name, ppt_meta_file, section_names)


def _convert_master_to_arrays_bulk_nsm(ws_name, ppt_meta_file, section_names):
    from ns_engine.nsm_io import load_sections_bulk_nsm, _section_key, _decode_cell

    keys = [_section_key(s) for s in section_names]
    dfs = load_sections_bulk_nsm(str(ppt_meta_file), keys)

    results = {}
    for sname, key in zip(section_names, keys):
        df = dfs.get(key, None)
        if df is None or df.empty:
            results[sname] = ['_NOT_FOUND_', 1] if ws_name == 'Master_Data' else []
            continue
        return_array = []
        for row_idx, (_, row_data) in enumerate(df.iterrows(), start=1):
            cell_values = [_decode_cell(v) for v in row_data]
            while cell_values and cell_values[-1] == '':
                cell_values.pop()
            if cell_values:
                return_array.append([row_idx, cell_values])
        results[sname] = return_array
    return results


def _convert_master_to_arrays_bulk_xlsx(ws_name, ppt_meta_file, section_names):
    input_wb = openpyxl.load_workbook(ppt_meta_file, read_only=True, data_only=True)
    try:
        ws = input_wb[ws_name]
        all_rows = list(ws.iter_rows(values_only=True))
    finally:
        input_wb.close()

    results = {}
    for section_name in section_names:
        start_row = None
        end_row = len(all_rows)
        for i, row in enumerate(all_rows):
            if row and row[0] == section_name:
                start_row = i
                continue
            if start_row is not None and row and isinstance(row[0], str) and row[0].startswith('<<') and row[0].endswith('>>'):
                end_row = i
                break

        if start_row is None:
            results[section_name] = ['_NOT_FOUND_', 1] if ws_name == 'Master_Data' else []
            continue

        return_array = []
        for row_idx_offset, row in enumerate(all_rows[start_row:end_row], start=1):
            cell_values = []
            trailing_none = 0
            for v in row:
                if v is None:
                    trailing_none += 1
                else:
                    if trailing_none > 0:
                        cell_values.extend(['' for _ in range(trailing_none)])
                        trailing_none = 0
                    cell_values.append(v)
            if cell_values:
                return_array.append([row_idx_offset, cell_values])
        results[section_name] = return_array
    return results


def write_sections_bulk(file_path, section_dict):
    """Write multiple sections at once to reduce file I/O.

    section_dict: {section_tag: data} where data is either:
      - array format from convert_master_to_array, or
      - tuple (dict) format from convert_array_to_tuple
    """
    normalised = {}
    for tag, data in section_dict.items():
        if isinstance(data, dict):
            normalised[tag] = data
        else:
            normalised[tag] = convert_array_to_tuple(data)

    if str(file_path).lower().endswith('.nsm'):
        _write_sections_bulk_nsm(file_path, normalised)
    else:
        _write_sections_bulk_xlsx(file_path, normalised)


def _write_sections_bulk_nsm(file_path, section_dict):
    """Write multiple sections to .nsm in a single ZIP open/close cycle.

    section_dict values are already tuples (dict with (row,col) keys).
    """
    from ns_engine.nsm_io import _section_key, _parquet_filename, _encode_cell
    import zipfile, io
    import pyarrow as pa
    import pyarrow.parquet as pq
    import pandas as pd

    new_parquets = {}
    for section_tag, tup in section_dict.items():
        key = _section_key(section_tag)

        rows_data = {}
        for (r, c), val in tup.items():
            rows_data[(r, c)] = _encode_cell(val)

        if not rows_data:
            new_parquets[key] = pd.DataFrame()
            continue

        max_row = max(r for r, c in rows_data)
        max_col = max(c for r, c in rows_data)
        data = []
        for r in range(1, max_row + 1):
            row_vals = []
            for c in range(1, max_col + 1):
                val = rows_data.get((r, c), '')
                row_vals.append(str(val) if val is not None else '')
            data.append(row_vals)
        cols = [f'col_{i}' for i in range(max_col)]
        new_parquets[key] = pd.DataFrame(data, columns=cols)

    pq_names_to_replace = {_parquet_filename(k) for k in new_parquets}

    import os
    existing = {}
    if os.path.exists(file_path):
        with zipfile.ZipFile(file_path, 'r') as zf:
            for name in zf.namelist():
                if name not in pq_names_to_replace and not name.endswith('.xlsx'):
                    existing[name] = zf.read(name)

    with zipfile.ZipFile(file_path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for name, data in existing.items():
            zf.writestr(name, data)
        for key, df in new_parquets.items():
            buf = io.BytesIO()
            table = pa.Table.from_pandas(df.astype(str), preserve_index=False)
            pq.write_table(table, buf)
            zf.writestr(_parquet_filename(key), buf.getvalue())

    try:
        from ns_engine.nsm_adapter import invalidate_nsm_cache
        invalidate_nsm_cache(file_path)
    except ImportError:
        pass


def _write_sections_bulk_xlsx(file_path, section_dict):
    """Write multiple sections to .xlsx in a single workbook open/save cycle.

    section_dict values are already tuples (dict with (row,col) keys).
    """
    wb = openpyxl.load_workbook(file_path)

    for section_tag, tup in section_dict.items():
        ws_name = 'Master_Data'
        ws = wb[ws_name]

        start_row = None
        for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=1, values_only=False):
            if row[0].value == section_tag:
                start_row = row[0].row
                break

        if start_row is None:
            continue

        end_row = start_row + 1
        for row_idx in range(start_row + 1, min(start_row + 50000, ws.max_row + 1)):
            val = ws.cell(row_idx, 1).value
            if val and isinstance(val, str) and val.startswith('<<') and val.endswith('>>'):
                end_row = row_idx
                break
        else:
            end_row = ws.max_row + 1

        for r in range(start_row, end_row):
            for c in range(1, ws.max_column + 1):
                ws.cell(r, c).value = None

        ws.cell(start_row, 1).value = section_tag

        for (r, c), val in tup.items():
            s = str(val).strip() if val is not None else ''
            if s.startswith('<<') and s.endswith('>>'):
                continue
            ws.cell(start_row + r, c).value = val

    wb.save(file_path)
    wb.close()


def _convert_master_to_array_nsm(ws_name, ppt_meta_file, section_name):
    """Fast path: read section from .nsm (ZIP+Parquet) file."""
    from ns_engine.nsm_io import load_section, _section_key, _decode_cell

    key = _section_key(section_name)
    df = load_section(ppt_meta_file, key)

    if df.empty:
        if ws_name == 'Master_Data':
            return ['_NOT_FOUND_', 1]
        return []

    return_array = []
    for row_idx, (_, row_data) in enumerate(df.iterrows(), start=1):
        cell_values = [_decode_cell(v) for v in row_data]
        while cell_values and cell_values[-1] == '':
            cell_values.pop()
        if cell_values:
            return_array.append([row_idx, cell_values])

    return return_array


def _convert_master_to_array_xlsx(ws_name, ppt_meta_file, section_name):
    """Original xlsx path: read section from Excel file using openpyxl."""
    input_ppt_mata_excel = openpyxl.load_workbook(ppt_meta_file)
    input_ppt_mata_excel.active = input_ppt_mata_excel[ws_name]

    flag_finish = False
    flag_not_found = False
    flag_get_start_row = False
    current_row = 1
    empty_count = 0
    start_row = 1

    while flag_finish == False:
        if input_ppt_mata_excel.active.cell(current_row, 1).value == section_name:
            start_row = current_row
            flag_get_start_row = True
            current_row += 1

        if '<<' in str(input_ppt_mata_excel.active.cell(current_row , 1).value) and '>>' in str(input_ppt_mata_excel.active.cell(current_row, 1).value)\
                and flag_get_start_row == True:
            end_row = current_row - 1
            flag_finish = True
            break

        if str(input_ppt_mata_excel.active.cell(current_row, 1).value) == 'None':
            empty_count += 1
        else:
            empty_count = 0

        if empty_count >= 10000 and section_name == '<<POSITION_TAG>>':
            flag_finish = True
            flag_not_found = True
            end_row = current_row
            break
        elif empty_count >= 3000:
            flag_finish = True
            end_row = current_row
            flag_not_found = True
            break
        current_row += 1

    return_array = []
    max_row = 1
    for tmp_row in range(start_row,end_row+1):
        tmp_array = []
        current_row_array = []
        flag_column_end = False
        tmp_column = 1
        tmp_empty_count = 0
        while flag_column_end == False:
            if str(input_ppt_mata_excel.active.cell(tmp_row, tmp_column).value) != 'None' and tmp_empty_count == 0:
                current_row_array.append(input_ppt_mata_excel.active.cell(tmp_row, tmp_column).value)
                max_row = tmp_row

            if str(input_ppt_mata_excel.active.cell(tmp_row, tmp_column).value) == 'None':
                tmp_empty_count += 1
                tmp_array.append('')

            if str(input_ppt_mata_excel.active.cell(tmp_row, tmp_column).value) != 'None' and tmp_empty_count != 0:
                tmp_empty_count = 0
                for m in tmp_array:
                    current_row_array.append(m)
                current_row_array.append(input_ppt_mata_excel.active.cell(tmp_row, tmp_column).value)
                max_row = tmp_row
                tmp_array = []
            if tmp_empty_count >= 100:
                flag_column_end = True

            tmp_column += 1

        if len(current_row_array) != 0:
            return_array.append([tmp_row - start_row +1,current_row_array])

    if flag_not_found == True and ws_name == 'Master_Data':
        return_array = ['_NOT_FOUND_',max_row]

    return(return_array)

#convert from excel table to array
def convert_excel_to_array(ws_name, excel_file, start_row):
    input_ppt_mata_excel = openpyxl.load_workbook(excel_file)
    input_ppt_mata_excel.active = input_ppt_mata_excel[ws_name]

    # GET Folder names
    flag_finish = False
    current_row = 1
    empty_count = 0

    while flag_finish == False:
        if str(input_ppt_mata_excel.active.cell(current_row, 1).value) == 'None' and str(input_ppt_mata_excel.active.cell(current_row, 2).value) == 'None':
            empty_count += 1
        else:
            empty_count = 0

        if empty_count >= 100:
            flag_finish = True
            end_row = current_row
        current_row += 1
    #print(start_row,end_row)

    return_array = []
    for tmp_row in range(start_row,end_row+1):
        tmp_array = []
        current_row_array = []
        flag_column_end = False
        tmp_column = 1
        tmp_empty_count = 0
        while flag_column_end == False:
            if str(input_ppt_mata_excel.active.cell(tmp_row, tmp_column).value) != 'None' and tmp_empty_count == 0:
                current_row_array.append(input_ppt_mata_excel.active.cell(tmp_row, tmp_column).value)

            if str(input_ppt_mata_excel.active.cell(tmp_row, tmp_column).value) == 'None':
                tmp_empty_count += 1
                tmp_array.append('')

            if str(input_ppt_mata_excel.active.cell(tmp_row, tmp_column).value) != 'None' and tmp_empty_count != 0:
                tmp_empty_count = 0
                for m in tmp_array:
                    current_row_array.append(m)
                current_row_array.append(input_ppt_mata_excel.active.cell(tmp_row, tmp_column).value)
                tmp_array = []
            if tmp_empty_count >= 100:
                flag_column_end = True
            tmp_column += 1

        if len(current_row_array) != 0:
            return_array.append([tmp_row - start_row +1,current_row_array])
    input_ppt_mata_excel.close()
    return(return_array)


def convert_workbook_to_array(wb, ws_name, start_row):
    """Same as convert_excel_to_array but operates on an in-memory openpyxl Workbook."""
    wb.active = wb[ws_name]

    flag_finish = False
    current_row = 1
    empty_count = 0

    while flag_finish == False:
        if str(wb.active.cell(current_row, 1).value) == 'None' and str(wb.active.cell(current_row, 2).value) == 'None':
            empty_count += 1
        else:
            empty_count = 0

        if empty_count >= 100:
            flag_finish = True
            end_row = current_row
        current_row += 1

    return_array = []
    for tmp_row in range(start_row, end_row + 1):
        tmp_array = []
        current_row_array = []
        flag_column_end = False
        tmp_column = 1
        tmp_empty_count = 0
        while flag_column_end == False:
            if str(wb.active.cell(tmp_row, tmp_column).value) != 'None' and tmp_empty_count == 0:
                current_row_array.append(wb.active.cell(tmp_row, tmp_column).value)

            if str(wb.active.cell(tmp_row, tmp_column).value) == 'None':
                tmp_empty_count += 1
                tmp_array.append('')

            if str(wb.active.cell(tmp_row, tmp_column).value) != 'None' and tmp_empty_count != 0:
                tmp_empty_count = 0
                for m in tmp_array:
                    current_row_array.append(m)
                current_row_array.append(wb.active.cell(tmp_row, tmp_column).value)
                tmp_array = []
            if tmp_empty_count >= 100:
                flag_column_end = True
            tmp_column += 1

        if len(current_row_array) != 0:
            return_array.append([tmp_row - start_row + 1, current_row_array])
    return(return_array)


def remove_rows_under_section(tmp_ws_name, ppt_meta_file, clear_section_tuple):
    """
    Remove rows from master matching input tuple values.
    Supports both .xlsx (openpyxl) and .nsm (Parquet) formats.
    """
    if str(ppt_meta_file).lower().endswith('.nsm'):
        return _remove_rows_under_section_nsm(ppt_meta_file, clear_section_tuple)
    return _remove_rows_under_section_xlsx(tmp_ws_name, ppt_meta_file, clear_section_tuple)


def _remove_rows_under_section_nsm(ppt_meta_file, clear_section_tuple):
    from ns_engine.nsm_io import load_section, save_section, _section_key, _decode_cell

    section_name = None
    for key in clear_section_tuple:
        val = clear_section_tuple[key]
        if isinstance(val, str) and '<<' in val and '>>' in val:
            section_name = val
            break
    if not section_name:
        return 'remove_rows_under_section: No section name found'

    skey = _section_key(section_name)
    df = load_section(ppt_meta_file, skey)
    if df.empty:
        return f'remove_rows_under_section: Section "{section_name}" empty'

    values_to_match = set()
    for key in clear_section_tuple:
        val = clear_section_tuple[key]
        if isinstance(val, str) and '<<' in val and '>>' in val:
            continue
        if val is not None and val != '':
            values_to_match.add(str(val))

    rows_to_keep = []
    deleted = 0
    for idx, row in df.iterrows():
        first_val = _decode_cell(row.iloc[0]) if len(row) > 0 else ''
        s = str(first_val)
        if s.startswith('<<') and s.endswith('>>'):
            rows_to_keep.append(idx)
            continue
        matched = False
        for cell_val in row:
            decoded = _decode_cell(cell_val)
            if decoded is not None and decoded != '':
                if str(decoded) in values_to_match:
                    matched = True
                    break
        if matched:
            deleted += 1
        else:
            rows_to_keep.append(idx)

    if deleted > 0:
        new_df = df.loc[rows_to_keep].reset_index(drop=True)
        save_section(ppt_meta_file, skey, new_df)

    return f'remove_rows_under_section: {deleted} rows deleted'


def _remove_rows_under_section_xlsx(tmp_ws_name, ppt_meta_file, clear_section_tuple):
    """Original xlsx implementation."""
    import openpyxl

    wb = openpyxl.load_workbook(ppt_meta_file)
    ws = wb[tmp_ws_name]

    # ========== STEP 1: Find section name (OPTIMIZED) ==========
    section_name = None
    for key in clear_section_tuple:
        val = clear_section_tuple[key]
        if isinstance(val, str) and '<<' in val and '>>' in val:
            section_name = val
            break

    if not section_name:
        wb.close()
        return 'remove_rows_under_section: No section name found'

    # ========== STEP 2: Find section start (OPTIMIZED) ==========
    # Use iter_rows instead of cell-by-cell access
    start_row = None
    max_search = min(1000000, ws.max_row)  # Limit search to 100k rows

    for row in ws.iter_rows(min_row=1, max_row=max_search, min_col=1, max_col=1, values_only=False):
        if row[0].value == section_name:
            start_row = row[0].row
            break

    if not start_row:
        wb.close()
        return f'remove_rows_under_section: Section "{section_name}" not found'

    # ========== STEP 3: Clear cells right of section marker ==========
    # Only clear if needed (check first)
    if ws.cell(start_row, 2).value is not None:
        for col in range(2, ws.max_column + 1):
            ws.cell(start_row, col).value = None

    # ========== STEP 4: Find section end ==========
    end_row = start_row + 1
    for row_idx in range(start_row + 1, min(start_row + 10000, ws.max_row + 1)):
        val = ws.cell(row_idx, 1).value
        if val and isinstance(val, str) and '<<' in val and '>>' in val:
            end_row = row_idx
            break
    else:
        end_row = ws.max_row + 1

    # ========== STEP 5: Pre-build value set for fast lookup (CRITICAL) ==========
    # Convert tuple values to set for O(1) lookup instead of O(n)
    values_to_match = set()
    for key in clear_section_tuple:
        val = clear_section_tuple[key]
        # Skip section markers
        if isinstance(val, str) and '<<' in val and '>>' in val:
            continue
        # Skip None/empty
        if val is not None and val != '':
            values_to_match.add(str(val))

    # ========== STEP 6: Find rows to delete (OPTIMIZED) ==========
    rows_to_delete = []

    # Use iter_rows for batch reading (much faster)
    for row in ws.iter_rows(min_row=start_row + 1, max_row=end_row - 1, values_only=False):
        row_num = row[0].row

        # Check if row starts with section marker
        first_val = row[0].value
        if first_val and isinstance(first_val, str) and '<<' in first_val and '>>' in first_val:
            # Clear trailing cells for sub-section markers
            for cell in row[1:]:
                if cell.value is not None:
                    cell.value = None
            continue

        # Check if any cell value matches
        # Limit to reasonable column count (optimize)
        max_check_col = min(1000, len(row))  # Don't check beyond 1000 columns

        for cell in row[:max_check_col]:
            if cell.value is None or cell.value == '':
                continue

            # Fast O(1) lookup instead of O(n) loop
            if str(cell.value) in values_to_match:
                rows_to_delete.append(row_num)
                break

    # ========== STEP 7: Delete rows in batch (ULTRA OPTIMIZED) ==========
    # ★★★ MINIMAL CHANGE: Detect consecutive ranges and use amount parameter ★★★
    deleted_count = 0

    if rows_to_delete:
        # Group consecutive rows into ranges
        ranges = []
        range_start = rows_to_delete[0]
        range_count = 1

        for i in range(1, len(rows_to_delete)):
            if rows_to_delete[i] == rows_to_delete[i - 1] + 1:
                # Consecutive
                range_count += 1
            else:
                # End of range
                ranges.append((range_start, range_count))
                range_start = rows_to_delete[i]
                range_count = 1

        # Add last range
        ranges.append((range_start, range_count))

        # Delete in reverse order using amount parameter
        for start, count in sorted(ranges, reverse=True):
            ws.delete_rows(start, amount=count)
            deleted_count += count

    # ========== STEP 8: Save and close ==========
    wb.save(ppt_meta_file)
    wb.close()

    return f'remove_rows_under_section: {deleted_count} rows deleted'


from openpyxl import load_workbook


def clear_section_sheet(tmp_ws_name, ppt_meta_file, clear_section_taple):
    '''
    Clears specified cells in a section of an Excel worksheet

    :param tmp_ws_name: worksheet name to operate on
    :param ppt_meta_file: path to the Excel file
    :param clear_section_taple: dict with (row, col) tuples as keys and values to identify section
    :return: status string 'clear_section_sheet'
    '''
    if str(ppt_meta_file).lower().endswith('.nsm'):
        if tmp_ws_name.startswith('_tmp') or tmp_ws_name.startswith('_template'):
            staging = str(ppt_meta_file) + '.staging.xlsx'
            if os.path.exists(staging):
                wb = load_workbook(staging, data_only=True, keep_vba=False)
                save_path = staging
                ws = wb[tmp_ws_name]
                section_name = 'N/A'
                for tmp_k in clear_section_taple:
                    value = str(clear_section_taple[tmp_k])
                    if '<<' in value and '>>' in value:
                        section_name = clear_section_taple[tmp_k]
                        break
                flag_get_section = False
                start_row = 0
                for row_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=1000000, min_col=1, max_col=1, values_only=True), start=1):
                    if row[0] == section_name:
                        start_row = row_idx
                        flag_get_section = True
                        break
                if not flag_get_section:
                    wb.close()
                    return 'clear_section_sheet'
                for tmp_k in clear_section_taple:
                    value = str(clear_section_taple[tmp_k])
                    if value != str(section_name):
                        target_row = tmp_k[0] + start_row - 1
                        target_col = tmp_k[1]
                        ws.cell(target_row, target_col).value = ''
                wb.save(save_path)
                wb.close()
            return 'clear_section_sheet'

        from ns_engine.nsm_io import save_section, _section_key
        import pandas as pd
        section_name = 'N/A'
        for k in clear_section_taple:
            v = str(clear_section_taple[k])
            if '<<' in v and '>>' in v:
                section_name = clear_section_taple[k]
                break
        if section_name != 'N/A':
            key = _section_key(section_name)
            header_vals = {}
            for k in clear_section_taple:
                if k[0] == 1:
                    header_vals[k[1]] = clear_section_taple[k]
            if header_vals:
                max_col = max(header_vals.keys())
                row = [str(header_vals.get(c, '')) for c in range(1, max_col + 1)]
                cols = [f'col_{i}' for i in range(max_col)]
                save_section(str(ppt_meta_file), key, pd.DataFrame([row], columns=cols))
            else:
                save_section(str(ppt_meta_file), key, pd.DataFrame())
        return 'clear_section_sheet'
    wb = load_workbook(ppt_meta_file, data_only=True, keep_vba=False)
    save_path = _nsm_save_path(ppt_meta_file)
    ws = wb[tmp_ws_name]

    # Find section name from clear_section_taple
    # Look for entries with '<<' and '>>' markers
    section_name = 'N/A'
    for tmp_clear_section_taple in clear_section_taple:
        value = str(clear_section_taple[tmp_clear_section_taple])
        if '<<' in value and '>>' in value:
            section_name = clear_section_taple[tmp_clear_section_taple]
            break

    # Search for section in first column using optimized iteration
    flag_get_section = False
    start_row = 0

    # Use iter_rows for memory-efficient searching (up to 1,000,000 rows)
    for row_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=1000000, min_col=1, max_col=1, values_only=True), start=1):
        if row[0] == section_name:
            start_row = row_idx
            flag_get_section = True
            break

    # Exit with error if section not found
    if not flag_get_section:
        print(f'ERROR cannot find section name -- {section_name}')
        wb.close()
        exit()

    # Clear cells in the section (batch processing)
    # Skip the cell that contains the section name itself
    for tmp_clear_section_taple in clear_section_taple:
        value = str(clear_section_taple[tmp_clear_section_taple])
        if value != str(section_name):
            # Calculate absolute position and clear the cell
            target_row = tmp_clear_section_taple[0] + start_row - 1
            target_col = tmp_clear_section_taple[1]
            ws.cell(target_row, target_col).value = ''

    wb.save(save_path)
    wb.close()

    return 'clear_section_sheet'


def clear_tag_in_position_line(tmp_ws_name, ppt_meta_file, clear_section_taple):
    if str(ppt_meta_file).lower().endswith('.nsm'):
        if tmp_ws_name.startswith('_tmp') or tmp_ws_name.startswith('_template'):
            staging = str(ppt_meta_file) + '.staging.xlsx'
            if os.path.exists(staging):
                wb = openpyxl.load_workbook(staging)
                wb.active = wb[tmp_ws_name]
                section_name = 'N/A'
                for tmp_k in clear_section_taple:
                    if '<<' in str(clear_section_taple[tmp_k]) and '>>' in str(clear_section_taple[tmp_k]):
                        section_name = clear_section_taple[tmp_k]
                        break
                flag_get_section = False
                i = 1
                while not flag_get_section:
                    if wb.active.cell(i, 1).value == section_name:
                        start_row = i
                        flag_get_section = True
                        break
                    i += 1
                    if i > 1000000:
                        wb.close()
                        return
                for tmp_k in clear_section_taple:
                    if str(clear_section_taple[tmp_k]) != str(section_name) and (tmp_k[1] == 3 or tmp_k[1] == 4):
                        wb.active.cell(tmp_k[0] + start_row - 1, tmp_k[1]).value = ''
                wb.save(staging)
                wb.close()
            return

        from ns_engine.nsm_io import load_section, save_section, _section_key
        import pandas as pd
        section_name = 'N/A'
        for k in clear_section_taple:
            v = str(clear_section_taple[k])
            if '<<' in v and '>>' in v:
                section_name = clear_section_taple[k]
                break
        if section_name == 'N/A':
            return
        key = _section_key(section_name)
        df = load_section(ppt_meta_file, key)
        if df.empty:
            return
        for k in clear_section_taple:
            if str(clear_section_taple[k]) != str(section_name) and (k[1] == 3 or k[1] == 4):
                row_idx = k[0] - 1
                col_idx = k[1] - 1
                if 0 <= row_idx < len(df) and 0 <= col_idx < len(df.columns):
                    df.iloc[row_idx, col_idx] = ''
        save_section(str(ppt_meta_file), key, df)
        return

    wb = openpyxl.load_workbook(ppt_meta_file)
    save_path = _nsm_save_path(ppt_meta_file)
    wb.active = wb[tmp_ws_name]

    section_name = 'N/A'
    for tmp_clear_section_taple in clear_section_taple:
        if '<<' in str(clear_section_taple[tmp_clear_section_taple]) and '>>' in str(clear_section_taple[tmp_clear_section_taple]):
            section_name = clear_section_taple[tmp_clear_section_taple]
            break

    flag_get_section = False
    i = 1
    while flag_get_section == False:
        if wb.active.cell(i,1).value == section_name:
            start_row = i
            flag_get_section = True
            break
        i += 1

        if i > 1000000:
            print('EEROR cannot find section name -- ',section_name)
            exit()

    for tmp_clear_section_taple in clear_section_taple:
        if str(clear_section_taple[tmp_clear_section_taple]) != str(section_name) and (tmp_clear_section_taple[1] == 3 or tmp_clear_section_taple[1] == 4):
            wb.active.cell(tmp_clear_section_taple[0] + start_row -1, tmp_clear_section_taple[1]).value = ''

    wb.save(save_path)
    wb.close()

    return ('clear_tag_in_position_line')

#get shape name in the folder and sort by tuple type
def get_shape_folder_tuple(position_shape_tuple):
    return_tuple = {}
    current_folder_name = ''
    for tmp_position_shape_tuple in position_shape_tuple:
        if tmp_position_shape_tuple[0] != 1 and position_shape_tuple[tmp_position_shape_tuple] != '<END>':
            if position_shape_tuple[tmp_position_shape_tuple[0],1] != '' and tmp_position_shape_tuple[1] == 1:
                #print(position_shape_tuple[tmp_position_shape_tuple])
                current_folder_name = position_shape_tuple[tmp_position_shape_tuple]
            if tmp_position_shape_tuple[1] != 1:
                return_tuple[position_shape_tuple[tmp_position_shape_tuple]] = current_folder_name

    return (return_tuple)

### convert value from interface name  . exsample Gigabit Ethernet 0/0 -> 1001000
import re

def get_if_value(if_name):
    sum_num = 0
    # Strip trailing whitespace from the original interface name
    original_if_name_stripped = if_name.rstrip()

    processed_if_name = original_if_name_stripped

    # Pre-processing: Insert a space if the name does not already contain one,
    # specifically between the interface type and the number part.
    # This handles cases like 'GigabitEthernet0/1' -> 'GigabitEthernet 0/1'
    # and '4921:TerminationZ9333' -> '4921:TerminationZ 9333'.
    if ' ' not in original_if_name_stripped:
        # Find the first occurrence where a non-slash, non-dot, non-whitespace, non-digit character
        # is immediately followed by a digit.
        # Group 1: The character before the digit (e.g., 't', 'Z', 'l')
        # Group 2: The digit itself
        match = re.search(r'([^/\.\s\d])(\d)', original_if_name_stripped)
        if match:
            # Insert a space at the position of the digit
            idx = match.start(2) # Get the starting index of the digit (Group 2)
            processed_if_name = original_if_name_stripped[:idx] + ' ' + original_if_name_stripped[idx:]

    if_name = processed_if_name

    # Original function logic starts here
    if ' ' in if_name:
        # Assumes the split_portname function is defined elsewhere or imported.
        # If not defined, a NameError will occur.
        # Example definition:
        # def split_portname(name_str):
        #     parts = name_str.split(' ', 1)
        #     return parts if len(parts) > 1 else [name_str, '']

        split_if_name = split_portname(if_name)

        # Basic handling if split_portname does not return the expected format
        if len(split_if_name) < 2:
            print(f"Warning: split_portname did not return enough parts for '{if_name}'. Returning -1.")
            return -1

        if '/' in split_if_name[1] or '.' in split_if_name[1]:          #update  replace '.' to '/' for Network Sketcher ver 2.0
            split_if_name[1] = split_if_name[1].replace('.','/')
            each_num = split_if_name[1].split('/')
            #print(split_if_name[1],len(each_num)) # Debug print
            tmp_add_value = '1'
            for i in range(0,int(len(each_num))):
                tmp_add_value += '000'
            tmp_num = int(tmp_add_value)
            #print('tmp_num  ', tmp_num) # Debug print
            for n in range(0,int(len(each_num))):
                try:
                    sum_num += (int(each_num[n]) + 1) * tmp_num
                except ValueError:
                    # Handle cases where a part of the interface path cannot be converted to an integer
                    print(f"Error: Could not convert '{each_num[n]}' to integer in interface path part. Input: '{if_name}'. Returning -1.")
                    return -1
                tmp_num = tmp_num/1000
            #print(split_if_name[1],sum_num) # Debug print
            if_value = sum_num

        else:
            if_value = split_if_name[1]

    else:
        # If no space was found even after pre-processing (e.g., regex didn't match)
        if_value = -1

    # Ensure the final return value is an integer.
    # It's recommended to wrap this in a try-except block to prevent ValueError
    # if if_value somehow ends up as a non-convertible string.
    try:
        return int(if_value)
    except ValueError:
        print(f"Error: Cannot convert final if_value '{if_value}' to integer. Input: '{if_name}'. Returning -1.")
        return -1

def split_portname(if_name):
    #reduce space' ' in if name
    if_name_split = str(if_name).split(' ')
    tmp_if_name_split = str(if_name).replace(str(if_name_split[-1]), '')

    if len(if_name_split) != 1:
        name = tmp_if_name_split.replace(' ','')
        num = str(if_name_split[-1])
    else:
        name = tmp_if_name_split.replace(' ','')
        num = ''

    return_array = [name, num]
    return (return_array)

def check_file_open(file_fullpath):
    if os.path.exists(file_fullpath):
        try:
            os.rename(file_fullpath, file_fullpath) #can't rename an open file so an error will be thrown
            return False
        except:
            #tkinter.messagebox.showwarning(title="File is being opened", message="Please close the file below." + '\n\n' + file_fullpath)
            print('[WARNING] The file you are writing to may have been left open; Windows may display this message even when there is no problem.')
            return False #chenge for windows at Ver 2.2.0(a)

    #raise NameError

def adjust_portname(if_name):
    '''adjust port name to NS format.'''
    #Whether a number is included in portname
    if any(map(str.isdigit, if_name)) == True:
        if_name = if_name.replace(' ', '')
        flag_first_digit = False
        tmp_if_name = ''
        tmp_if_num = ''
        for per_char in if_name:
            if per_char.isdigit() == True:
                flag_first_digit = True

            if flag_first_digit == False:
                tmp_if_name += per_char
            else:
                tmp_if_num += per_char
        if len(tmp_if_name) <= 1:
            if_name_abbreviation = tmp_if_name
        elif tmp_if_name == 'GigabitEthernet':
            if_name_abbreviation = 'GE'
        elif tmp_if_name == 'Ethernet':
            if_name_abbreviation = 'E'
        elif tmp_if_name == 'FastEthernet':
            if_name_abbreviation = 'FE'
        elif tmp_if_name == 'TenGigabitEthernet':
            if_name_abbreviation = 'TE'
        else:
            if_name_abbreviation = tmp_if_name[:2]

        return_array = [if_name_abbreviation , tmp_if_name,  tmp_if_num]

    else:
        if_name = if_name.replace(' ', '')
        if len(if_name) <= 1:
            if_name_abbreviation = if_name
        else:
            if_name_abbreviation = if_name[:2]

        return_array = [if_name_abbreviation , if_name,  random.randint(9001, 9999)]

    return (return_array)

def check_tuple_num_exist(target_tuple,row,column):
    ### if tuple [x,y] exist, return True. not exist, return False. ###
    flag_tuple_num_exist = False
    for tmp_target_tuple in target_tuple:
        if tmp_target_tuple[0] == row and tmp_target_tuple[1] == column:
            flag_tuple_num_exist = True
            break

    return (flag_tuple_num_exist)

def get_east_asian_width_count(text):
    count = 0
    for c in text:
        if unicodedata.east_asian_width(c) in 'FWA':
            count += 2
        else:
            count += 1
    return count

def get_tag_name_from_full_name(target_device_name,full_if_name,position_line_tuple): # Input 'GigabitEthernet' 0/6 Output 'GE 0/6' using self.position_line_tuple
    #print(target_device_name,full_if_name,position_line_tuple)
    for tmp_position_line_tuple in position_line_tuple:
        if tmp_position_line_tuple[0] != 1 and tmp_position_line_tuple[0] != 2 and (tmp_position_line_tuple[1] == 1 or tmp_position_line_tuple[1] == 2):
            if tmp_position_line_tuple[1] == 1:
                offet_column = 0
            elif tmp_position_line_tuple[1] == 2:
                offet_column = 1

            if position_line_tuple[tmp_position_line_tuple[0], tmp_position_line_tuple[1]] == target_device_name:
                # print(tmp_position_line_tuple, self.position_line_tuple[tmp_position_line_tuple])
                tmp_tag = position_line_tuple[tmp_position_line_tuple[0], 3 + offet_column]
                target = ' '
                idx = tmp_tag.find(target)
                modify_if_name = position_line_tuple[tmp_position_line_tuple[0], 13 + offet_column * 4] + ' ' + tmp_tag[idx + 1:]

                if modify_if_name == full_if_name:
                    return (tmp_tag)
    return ('_NO_MATCH_')

def build_full_name_lookup(position_line_tuple):
    """Build an O(1) lookup dict for (device_name, tag_if_name) -> full_if_name."""
    lookup = {}
    for (row, col) in position_line_tuple:
        if row in (1, 2) or col not in (1, 2):
            continue
        offset_column = 0 if col == 1 else 1
        device_name = position_line_tuple[(row, col)]
        tag = position_line_tuple.get((row, 3 + offset_column), '')
        if not tag or ' ' not in str(tag):
            continue
        prefix = position_line_tuple.get((row, 13 + offset_column * 4), '')
        space_idx = tag.find(' ')
        full_if_name = str(prefix) + ' ' + tag[space_idx + 1:]
        lookup[(device_name, tag)] = full_if_name
    return lookup


def get_full_name_from_tag_name(target_device_name, tag_if_name, position_line_tuple):
    """Input 'GE 0/6' Output 'GigabitEthernet 0/6' using position_line_tuple.

    For bulk lookups, prefer build_full_name_lookup() + dict.get() instead.
    """
    for tmp_position_line_tuple in position_line_tuple:
        if tmp_position_line_tuple[0] != 1 and tmp_position_line_tuple[0] != 2 and (tmp_position_line_tuple[1] == 1 or tmp_position_line_tuple[1] == 2):
            if tmp_position_line_tuple[1] == 1:
                offet_column = 0
            elif tmp_position_line_tuple[1] == 2:
                offet_column = 1

            if position_line_tuple[tmp_position_line_tuple[0], tmp_position_line_tuple[1]] == target_device_name:
                tmp_tag = position_line_tuple[tmp_position_line_tuple[0], 3 + offet_column]

                if tag_if_name == tmp_tag:
                    target = ' '
                    idx = tmp_tag.find(target)
                    full_if_name = position_line_tuple[tmp_position_line_tuple[0], 13 + offet_column * 4] + ' ' + tmp_tag[idx + 1:]

                    return (full_if_name)
    return ('_NO_MATCH_')

def check_ip_format(char_ip_address):
    ### IPv4 check ###
    flag_ipv4 = False
    if char_ip_address.count('.') == 3 and char_ip_address.count('/') == 1:
        target = '/'
        idx = char_ip_address.find(target)
        if char_ip_address[idx + 1:].isalnum() == True:
            if 1 <= int(char_ip_address[idx + 1:]) <= 32:
                flag_ipv4_array = []
                for tmp_num in char_ip_address[:idx].split('.'):
                    if tmp_num.isalnum() == True:
                        if 0 <= int(tmp_num) <= 255:
                            flag_ipv4_array.append('True')

                if flag_ipv4_array.count('True') == 4:
                    return ('IPv4')

    return ('NOT_MATCH')

def get_ipv4_value(ipv4_address):
    return_ipv4_value = ''
    target = '/'
    idx = ipv4_address.find(target)

    for tmp_num in ipv4_address[:idx].split('.'):
        zero_count = 3 - len(str(tmp_num))
        if zero_count != 0 :
            for tmp_zero_count in range(zero_count):
                return_ipv4_value += '0'
        return_ipv4_value += str(tmp_num)

    return(return_ipv4_value)

#Example usage
#change_cell_color(self.input_tree_excel "Sheet1", 1, 1, (235, 241, 222))
def change_cell_color(workbook, sheet_name, row, column, rgb_color):

    # Access the specified sheet
    if sheet_name not in workbook.sheetnames:
        raise ValueError(f"Sheet '{sheet_name}' not found in the workbook.")
    sheet = workbook[sheet_name]

    # Convert RGB color to a hexadecimal string
    hex_color = "{:02x}{:02x}{:02x}".format(*rgb_color)

    # Create a PatternFill with the specified RGB color
    fill = PatternFill(start_color=hex_color, end_color=hex_color, fill_type="solid")

    # Apply the fill to the specified cell
    sheet.cell(row=row, column=column).fill = fill

def _auto_add_attribute_section(master_file_path):
    """Add a default ATTRIBUTE section to a master file that lacks one (pre-v2.3 format).

    Reads STYLE_SHAPE to infer device names and assigns default colors:
      GREEN  -> ['DEVICE',   [235, 241, 222]]
      BLUE   -> ['WayPoint', [240, 244, 250]]
      others -> ['<EMPTY>',  [232, 232, 232]]

    For .nsm files: writes a single combined Parquet with section tag at row 1,
    column headers at row 2, and device data starting at row 3.

    For .xlsx files: uses the two-step template + data approach so that the
    <<ATTRIBUTE>> tag is placed correctly in the Master_Data sheet.

    Returns the updated attribute_array in the same format as
    convert_master_to_array(), or [] on failure.
    """
    try:
        import sys as _sys
        _sys.stderr.write('--- Attribute data not found, auto-adding from STYLE_SHAPE ---\n')
        worksheet_name = 'Master_Data'

        # Read STYLE_SHAPE to infer device names and colors
        master_style_shape_array = convert_master_to_array(
            worksheet_name, master_file_path, '<<STYLE_SHAPE>>')
        master_style_shape_array = master_style_shape_array[3:]
        for index, item in enumerate(master_style_shape_array):
            item[0] = index + 1
        master_style_shape_tuple = convert_array_to_tuple(master_style_shape_array)

        _headers = ['Device Name', 'Default', 'Attribute-A', 'Attribute-B',
                    'Attribute-C', 'Attribute-D', 'Attribute-E', 'Attribute-F',
                    'Attribute-G', 'Attribute-H', '<END>']

        if str(master_file_path).lower().endswith('.nsm'):
            # --- .nsm path: single combined write ---
            # _write_meta_nsm completely replaces the Parquet file on each call,
            # so two separate writes would destroy the first.  Build one dict
            # that contains the section tag (row 1), column headers (row 2), and
            # all device rows (rows 3+) and write them in a single call.
            combined = {}
            combined[(1, 1)] = '<<ATTRIBUTE>>'
            for col_idx, h in enumerate(_headers, start=1):
                combined[(2, col_idx)] = h
            device_row = 3
            seen_rows = set()
            for key in sorted(master_style_shape_tuple):
                row_num = key[0]
                if key[1] == 1 and row_num not in seen_rows:
                    seen_rows.add(row_num)
                    device_name = master_style_shape_tuple[key]
                    color_val = master_style_shape_tuple.get((row_num, 5), '')
                    if color_val == 'GREEN':
                        default_attr = "['DEVICE',[235, 241, 222]]"
                    elif color_val == 'BLUE':
                        default_attr = "['WayPoint', [240, 244, 250]]"
                    else:
                        default_attr = "['<EMPTY>', [232, 232, 232]]"
                    combined[(device_row, 1)] = device_name
                    combined[(device_row, 2)] = default_attr
                    for i in range(3, 11):
                        combined[(device_row, i)] = "['<EMPTY>', [255, 255, 255]]"
                    combined[(device_row, 11)] = '<END>'
                    device_row += 1
            write_excel_meta(combined, master_file_path, worksheet_name, '<<ATTRIBUTE>>', 0, 0)

        else:
            # --- .xlsx path: two-step template + data write ---
            # The template write inserts <<ATTRIBUTE>> at the correct row in the
            # Excel sheet; the subsequent data write then finds that tag and writes
            # device rows relative to it.
            _probe = convert_master_to_array(worksheet_name, master_file_path, '<<ATTRIBUTE>>')
            max_row = _probe[1] if (len(_probe) >= 2 and _probe[0] == '_NOT_FOUND_') else 1
            offset_row = max_row + 4

            tmp_master_data_array = [
                [1, ['<<ATTRIBUTE>>']],
                [2, _headers],
                [7, ['<<END_MARK>>']],
            ]
            template_tuple = convert_array_to_tuple(tmp_master_data_array)
            write_excel_meta(template_tuple, master_file_path,
                             worksheet_name, '_template_', offset_row, 0)

            master_attribute_tuple = {}
            for key in master_style_shape_tuple:
                if key[1] == 1:
                    master_attribute_tuple[(key[0], 1)] = master_style_shape_tuple[key]
                    color_val = master_style_shape_tuple.get((key[0], 5), '')
                    if color_val == 'GREEN':
                        master_attribute_tuple[(key[0], 2)] = "['DEVICE',[235, 241, 222]]"
                    elif color_val == 'BLUE':
                        master_attribute_tuple[(key[0], 2)] = "['WayPoint', [240, 244, 250]]"
                    else:
                        master_attribute_tuple[(key[0], 2)] = "['<EMPTY>', [232, 232, 232]]"
                for i in range(3, 11):
                    master_attribute_tuple[(key[0], i)] = "['<EMPTY>', [255, 255, 255]]"
            write_excel_meta(master_attribute_tuple, master_file_path,
                             worksheet_name, '<<ATTRIBUTE>>', 2, 0)

        return convert_master_to_array('Master_Data', master_file_path, '<<ATTRIBUTE>>')
    except Exception as e:
        import sys as _sys
        _sys.stderr.write(f'--- _auto_add_attribute_section failed: {e} ---\n')
        return []


def get_attribute_title_list(self, master_file_path):
    attribute_array = convert_master_to_array('Master_Data', master_file_path, '<<ATTRIBUTE>>')

    # Detect corrupted ATTRIBUTE data: section tag (row 1) or column headers
    # (row 2) are missing.  This happens when _write_meta_nsm overwrote the
    # Parquet with device-only data during a previous auto-add attempt.
    is_corrupted = (
        attribute_array and
        attribute_array[0] != '_NOT_FOUND_' and
        (len(attribute_array) < 2 or
         not (isinstance(attribute_array[0], list) and
              len(attribute_array[0]) >= 2 and
              attribute_array[0][1] == ['<<ATTRIBUTE>>']) or
         not (isinstance(attribute_array[1], list) and
              len(attribute_array[1]) >= 2 and
              isinstance(attribute_array[1][1], list) and
              len(attribute_array[1][1]) >= 1 and
              attribute_array[1][1][0] == 'Device Name'))
    )
    if is_corrupted:
        import sys as _sys
        _sys.stderr.write('--- Corrupted ATTRIBUTE data detected, regenerating ---\n')
        attribute_array = ['_NOT_FOUND_', 1]

    ### add <<ATTRIBUTE>> to Master file of pre ver 2.3
    if attribute_array[0] == '_NOT_FOUND_':
        excel_file_path = self.inFileTxt_L2_3_1.get()

        ### check file open
        if check_file_open(excel_file_path) == True:
            return ()

        ###create backup master file
        get_backup_filename(excel_file_path)

        # Delegate to _auto_add_attribute_section which correctly handles
        # both .nsm (single combined write) and .xlsx (two-step write)
        attribute_array = _auto_add_attribute_section(excel_file_path)
        if not attribute_array:
            return ()

    attribute_list = attribute_array[1][1]
    return attribute_list[1:-1]

def get_global_attribute_tuple(master_file_path,selected_title):
    attribute_array = convert_master_to_array('Master_Data', master_file_path, '<<ATTRIBUTE>>')
    attribute_array = attribute_array[1:]
    update_attribute_array = []
    selected_index = 0
    for kari_attribute_array in attribute_array:
        update_attribute_array.append(kari_attribute_array[1])
        for index, tmp_attribute_array in enumerate(kari_attribute_array[1], start=0):
            if tmp_attribute_array == selected_title and tmp_attribute_array != '<END>':
                selected_index = index

    selected_attribute_tuple = {}
    for kari_update_attribute_array in update_attribute_array:
        if kari_update_attribute_array[0] != 'Device Name':
            try:
                parsed_list = ast.literal_eval(kari_update_attribute_array[selected_index])
                extracted_array = parsed_list[1]
            except (ValueError, SyntaxError, IndexError, KeyError, TypeError):
                extracted_array = [255, 255, 255]
            selected_attribute_tuple[kari_update_attribute_array[0]] = extracted_array
    #print(selected_title,selected_attribute_tuple)
    return selected_attribute_tuple

def check_file_locked(file_path):
    if not os.path.exists(file_path):
        return False

    try:
        os.rename(file_path, file_path)
        return False
    except PermissionError:
        return True

class  get_l2_broadcast_domains():
    def run(self,excel_maseter_file):
        #print('--- get_l2_broadcast_domains ---')
        '''
        STEP0 get values of Master Data
        '''
        #parameter
        ws_name = 'Master_Data'
        ws_l2_name = 'Master_Data_L2'
        ws_l3_name = 'Master_Data_L3'
        #excel_maseter_file = self.inFileTxt_L3_3_1.get()

        _md_sections = ['<<POSITION_FOLDER>>', '<<POSITION_SHAPE>>', '<<POSITION_LINE>>',
                        '<<STYLE_SHAPE>>', '<<POSITION_TAG>>', '<<ROOT_FOLDER>>']
        _md_bulk = convert_master_to_arrays_bulk(ws_name, excel_maseter_file, _md_sections)
        self.position_folder_array = _md_bulk.get('<<POSITION_FOLDER>>', [])
        self.position_shape_array = _md_bulk.get('<<POSITION_SHAPE>>', [])
        self.position_line_array = _md_bulk.get('<<POSITION_LINE>>', [])
        self.position_style_shape_array = _md_bulk.get('<<STYLE_SHAPE>>', [])
        self.position_tag_array = _md_bulk.get('<<POSITION_TAG>>', [])
        self.root_folder_array = _md_bulk.get('<<ROOT_FOLDER>>', [])
        self.position_folder_tuple = convert_array_to_tuple(self.position_folder_array)
        self.position_shape_tuple = convert_array_to_tuple(self.position_shape_array)
        self.position_line_tuple = convert_array_to_tuple(self.position_line_array)
        self.position_style_shape_tuple = convert_array_to_tuple(self.position_style_shape_array)
        self.position_tag_tuple = convert_array_to_tuple(self.position_tag_array)
        self.root_folder_tuple = convert_array_to_tuple(self.root_folder_array)
        #print('---- self.position_folder_tuple ----')
        #print(self.position_folder_tuple)
        #print('---- self.position_shape_tuple ----')
        #print(self.position_shape_tuple)
        #print('---- self.position_line_tuple ----')
        #print(self.position_line_tuple)

        # GET Folder and wp name List
        self.folder_wp_name_array = get_folder_wp_array_from_master(ws_name, excel_maseter_file)

        self._opposite_if_index = {}
        _seen_rows = set()
        for _key in self.position_line_tuple:
            r, c = _key
            if r in _seen_rows or r == 1 or r == 2:
                continue
            _seen_rows.add(r)
            try:
                dev1 = self.position_line_tuple.get((r, 1))
                dev2 = self.position_line_tuple.get((r, 2))
                raw3 = self.position_line_tuple.get((r, 3))
                raw4 = self.position_line_tuple.get((r, 4))
                pfx1 = self.position_line_tuple.get((r, 13))
                pfx2 = self.position_line_tuple.get((r, 17))
                if dev1 is None or dev2 is None or raw3 is None or raw4 is None:
                    continue
                arr1 = split_portname(raw3)
                if1 = str(pfx1) + ' ' + arr1[1]
                arr2 = split_portname(raw4)
                if2 = str(pfx2) + ' ' + arr2[1]
                self._opposite_if_index[(str(dev1), if1)] = [str(dev2), if2]
                self._opposite_if_index[(str(dev2), if2)] = [str(dev1), if1]
            except Exception:
                pass

        self._l2seg_count_index = {}
        _l2_bulk = convert_master_to_arrays_bulk(ws_l2_name, excel_maseter_file, ['<<L2_TABLE>>'])
        self.l2_table_array = _l2_bulk.get('<<L2_TABLE>>', [])
        _l3_bulk = convert_master_to_arrays_bulk(ws_l3_name, excel_maseter_file, ['<<L3_TABLE>>'])
        self.l3_table_array = _l3_bulk.get('<<L3_TABLE>>', [])
        #print('--- self.l2_table_array ---')
        #print(self.l2_table_array)

        self.new_l2_table_array = []
        for tmp_l2_table_array in self.l2_table_array:
            if tmp_l2_table_array[0] != 1 and tmp_l2_table_array[0] != 2:
                tmp_l2_table_array[1].extend(['', '', '', '', '', '', '', ''])
                del tmp_l2_table_array[1][8:]
                self.new_l2_table_array.append(tmp_l2_table_array)

        self.device_list_array = []
        self.wp_list_array = []
        for tmp_new_l2_table_array in self.new_l2_table_array:
            if tmp_new_l2_table_array[1][1] not in self.device_list_array and tmp_new_l2_table_array[1][1] not in self.wp_list_array:
                if tmp_new_l2_table_array[1][0] == 'N/A':
                    self.wp_list_array.append(tmp_new_l2_table_array[1][1])
                else:
                    self.device_list_array.append(tmp_new_l2_table_array[1][1])

        self.all_shape_list_array = []
        self.all_shape_list_array.extend(self.device_list_array)
        self.all_shape_list_array.extend(self.wp_list_array)


        #print('--- self.device_list_array ---')
        #print(self.device_list_array)
        #print('--- self.wp_list_array ---')
        #print(self.wp_list_array)
        #print('--- self.all_shape_list_array ---')
        #print(self.all_shape_list_array)

        # GET L2 Segment name of each device — indexed by device
        _l2_by_device = {}
        for tmp_new_l2_table_array in self.new_l2_table_array:
            if tmp_new_l2_table_array[1][6] != '':
                _l2_by_device.setdefault(tmp_new_l2_table_array[1][1], []).append(tmp_new_l2_table_array)

        self.device_unique_l2name_array = []
        for device_name in self.all_shape_list_array:
            unique_set = set()
            for entry in _l2_by_device.get(device_name, []):
                for seg in entry[1][6].split(','):
                    seg = seg.replace(' ', '').strip()
                    if seg:
                        unique_set.add(seg)
            self.device_unique_l2name_array.append([device_name, sorted(unique_set)])

        #print('--- device_unique_l2name_array ---')
        #print(self.device_unique_l2name_array)

        # Input l2 l3 if type and Reduce ' ' in connected l2 segment name
        self.update_l2_table_array = []

        for tmp_tmp_new_l2_table_array in self.new_l2_table_array:
            offset_excel = 2
            tmp_new_l2_table_array = tmp_tmp_new_l2_table_array[1]
            if tmp_new_l2_table_array[offset_excel + 3]  == "":
                if tmp_new_l2_table_array[offset_excel + 4] == "":
                    if tmp_new_l2_table_array[offset_excel + 1] == "":
                        tmp_new_l2_table_array[offset_excel] = ''
                    else:
                        tmp_new_l2_table_array[offset_excel] = 'Routed (L3)'
                else:
                    if tmp_new_l2_table_array[offset_excel + 1] == "":
                        tmp_new_l2_table_array[offset_excel] = ''
                    else:
                        tmp_new_l2_table_array[offset_excel] = 'Switch (L2)'
            else:
                if tmp_new_l2_table_array[offset_excel + 1] == "":
                    tmp_new_l2_table_array[offset_excel] = ''
                else:
                    tmp_new_l2_table_array[offset_excel] = 'Switch (L2)'

            offset_excel = 4
            if tmp_new_l2_table_array[offset_excel + 1] == "":
                if tmp_new_l2_table_array[offset_excel + 1] == "":
                    tmp_new_l2_table_array[offset_excel] = ''
                else:
                    tmp_new_l2_table_array[offset_excel] = 'Routed (L3)'
            else:
                if tmp_new_l2_table_array[offset_excel + 2] == "":
                    if tmp_new_l2_table_array[offset_excel - 1] == "":
                        tmp_new_l2_table_array[offset_excel] = 'Loopback (L3)'
                    else:
                        tmp_new_l2_table_array[offset_excel] = 'Routed (L3)'
                else:
                    if tmp_new_l2_table_array[offset_excel - 1] == "":
                        tmp_new_l2_table_array[offset_excel] = 'Routed (L3)'
                    else:
                        tmp_new_l2_table_array[offset_excel] = 'Switch (L2)'

            #Reduce ' ' in connected l2 segment name
            tmp_new_l2_table_array[6] = tmp_new_l2_table_array[6].replace(' ','')

            self.update_l2_table_array.append(tmp_new_l2_table_array)

        for _row in self.update_l2_table_array:
            _key = (_row[1], _row[3])
            if _key not in self._l2seg_count_index:
                if _row[6] != '':
                    self._l2seg_count_index[_key] = len(_row[6].split(','))
                else:
                    self._l2seg_count_index[_key] = 0

        '''
        STEP1 get broadcast domain from l2 table
        '''
        self.local_boradcast_domain_num = 1
        self.global_boradcast_domain_num = 1
        self.master_boradcast_domain_array = []

        device_l2_boradcast_domain_array= get_l2_broadcast_domains.get_broadcast_domain(self)
        device_l2_directly_l3vport_array = get_l2_broadcast_domains.get_l2_directly_l3vport(self)
        device_l2_other_array = get_l2_broadcast_domains.get_l2_other(self)

        #print('--- device_l2_boradcast_domain_array ---')
        #print(device_l2_boradcast_domain_array)
        #print('--- device_l2_directly_l3vport_array ---')
        #print(device_l2_directly_l3vport_array)
        #print('--- device_l2_other_array ---')
        #print(device_l2_other_array)

        l2_broadcast_group_array = []

        ### make l2_broadcast_group_array using device_l2_boradcast_domain_array ###
        _flat_bd_array = []
        for _sl in device_l2_boradcast_domain_array:
            _flat_bd_array.extend(_sl)
        _flat_bd_by_device = {}
        for _fbd in _flat_bd_array:
            if _fbd:
                _flat_bd_by_device.setdefault(_fbd[1], []).append(_fbd)

        for tmp_device_l2_boradcast_domain_array in device_l2_boradcast_domain_array:
            #print('##### tmp_device_l2_boradcast_domain_array',tmp_device_l2_boradcast_domain_array)
            if tmp_device_l2_boradcast_domain_array != []:
                for kyuusai_device_l2_boradcast_domain_array in tmp_device_l2_boradcast_domain_array:
                    tmp_l2_broadcast_group_array = []
                    tmp_l2_broadcast_group_array.append(kyuusai_device_l2_boradcast_domain_array[0])
                    #print('##### kyuusai_device_l2_boradcast_domain_array ', kyuusai_device_l2_boradcast_domain_array)
                    for tmp_tmp_device_l2_boradcast_domain_array in kyuusai_device_l2_boradcast_domain_array[3]:
                        ### check per physical if ###
                        #print('     >tmp_tmp_device_l2_boradcast_domain_array ',tmp_tmp_device_l2_boradcast_domain_array)
                        if 'L2IF' in tmp_tmp_device_l2_boradcast_domain_array[0] or 'L3IF' in tmp_tmp_device_l2_boradcast_domain_array[0]:
                            ### check ####  device_l2_boradcast_domain_array  -> device_l2_boradcast_domain_array
                            #print('### per IFs###', tmp_device_l2_boradcast_domain_array[0][0], tmp_device_l2_boradcast_domain_array[0][1], tmp_device_l2_boradcast_domain_array[0][2], tmp_tmp_device_l2_boradcast_domain_array)
                            device_name = tmp_device_l2_boradcast_domain_array[0][1]
                            if_name = tmp_tmp_device_l2_boradcast_domain_array[1]
                            tmp_opposite_if_array = get_l2_broadcast_domains.get_opposite_if(self,device_name,if_name)
                            #print('          + device_name,if_name -> ', device_name,if_name ,'    +++tmp_oppsite_if_array -> ',tmp_opposite_if_array)

                            for now_device_l2_boradcast_domain_array in _flat_bd_array:
                                if now_device_l2_boradcast_domain_array != []:
                                    #print(now_device_l2_boradcast_domain_array)
                                    if now_device_l2_boradcast_domain_array[1] == tmp_opposite_if_array[0]:
                                        for now_now_device_l2_boradcast_domain_array in now_device_l2_boradcast_domain_array[3]:
                                            #print(now_now_device_l2_boradcast_domain_array[1] , tmp_opposite_if_array[1] , now_device_l2_boradcast_domain_array[2] , kyuusai_device_l2_boradcast_domain_array[2])
                                            if now_now_device_l2_boradcast_domain_array[1] == tmp_opposite_if_array[1] and now_device_l2_boradcast_domain_array[2] == kyuusai_device_l2_boradcast_domain_array[2]:
                                                #print('               *MATCH OPPO L2***',now_device_l2_boradcast_domain_array[2],now_device_l2_boradcast_domain_array[0])
                                                tmp_l2_broadcast_group_array.append(now_device_l2_boradcast_domain_array[0])
                                                break

                                                '''Kyuusai a one-L2_SEG if connect to a one-L2_SEG if with difference name'''
                                                tmp_num_1 = get_l2_broadcast_domains.get_l2seg_count_on_if(self,device_name,if_name)
                                                tmp_num_2 = get_l2_broadcast_domains.get_l2seg_count_on_if(self,tmp_opposite_if_array[0],tmp_opposite_if_array[1])
                                                #print(tmp_num_1,tmp_num_2,now_now_device_l2_boradcast_domain_array[1] , tmp_opposite_if_array[1] )
                                                if tmp_num_1 == 1 and tmp_num_2 == 1 and tmp_device_l2_boradcast_domain_array[0][2] != now_now_device_l2_boradcast_domain_array[2] and now_now_device_l2_boradcast_domain_array[1]  == tmp_opposite_if_array[1]:
                                                    #print('               *MATCH a one-L2_SEG if connect to a one-L2_SEG if with difference name***',tmp_device_l2_boradcast_domain_array[0][2],now_now_device_l2_boradcast_domain_array[2])
                                                    tmp_l2_broadcast_group_array.append(kyuusai_now_device_l2_boradcast_domain_array[0] )

                            ### check ####  device_l2_boradcast_domain_array  -> device_l2_directly_l3vport_array
                            for tmp_device_l2_directly_l3vport_array in device_l2_directly_l3vport_array:
                                #print(kyuusai_device_l2_boradcast_domain_array[1] , tmp_device_l2_directly_l3vport_array[0], tmp_device_l2_directly_l3vport_array[2][2],tmp_opposite_if_array[1] )
                                if tmp_device_l2_directly_l3vport_array[2][2] == tmp_opposite_if_array[1] and tmp_device_l2_directly_l3vport_array[0] == tmp_opposite_if_array[0]:
                                    if tmp_device_l2_directly_l3vport_array[1] == kyuusai_device_l2_boradcast_domain_array[2]:
                                        #print('   *MATCH OPPO l2_directly_l3vport***', kyuusai_device_l2_boradcast_domain_array[2])
                                        tmp_l2_broadcast_group_array.append(tmp_device_l2_directly_l3vport_array[2][0])
                                        break

                                    '''Kyuusai a one-L2_SEG if connect to a one-L2_oneL3 if'''
                                    tmp_num_1 = get_l2_broadcast_domains.get_l2seg_count_on_if(self, device_name, if_name)
                                    if tmp_num_1  == 1 and tmp_device_l2_directly_l3vport_array[1] == '__NO_L2SEG__':
                                        #print('               *MATCH a one-L2_SEG if connect to a one-L2_oneL3 if***', tmp_device_l2_directly_l3vport_array[0],tmp_device_l2_directly_l3vport_array[2][3])
                                        tmp_l2_broadcast_group_array.append(tmp_device_l2_directly_l3vport_array[2][0])
                                        break

                            ### check ####  device_l2_boradcast_domain_array  -> device_l2_other_array
                            for tmp_device_l2_other_array in device_l2_other_array:
                                if tmp_device_l2_other_array[2][2] == tmp_opposite_if_array[1] and tmp_device_l2_other_array[0] == tmp_opposite_if_array[0]:
                                    #print('               *MATCH device_l2_boradcast_domain_array  -> device_l2_other_array', tmp_opposite_if_array,tmp_device_l2_other_array[2][0])
                                    tmp_l2_broadcast_group_array.append(tmp_device_l2_other_array[2][0] )
                                    break


                    #print('--- list(set(tmp_l2_broadcast_group_array)) ---  ',sorted(list(set(tmp_l2_broadcast_group_array))))
                    l2_broadcast_group_array.append(sorted(list(set(tmp_l2_broadcast_group_array))))
        #print('--- 1st, l2_broadcast_group_array ---')
        #print(get_l2_broadcast_domains.get_unique_list(l2_broadcast_group_array))


        ### device_l2_directly_l3vport_array -> device_l2_directly_l3vport_array###
        #print('### device_l2_directly_l3vport_array -> device_l2_directly_l3vport_array###')
        #pre count, number of l2name
        l2_broadcast_group_array_2nd = []
        l2name_count_array = []
        l2seg_name_virtual_port_array = []
        for tmp_device_l2_directly_l3vport_array in device_l2_directly_l3vport_array:
            #print(tmp_device_l2_directly_l3vport_array)
            tmp_l2name_count = 0
            for tmp_tmp_device_l2_directly_l3vport_array in device_l2_directly_l3vport_array:
                if tmp_device_l2_directly_l3vport_array[0] == tmp_tmp_device_l2_directly_l3vport_array[0] and tmp_device_l2_directly_l3vport_array[2][2] == tmp_tmp_device_l2_directly_l3vport_array[2][2] :
                    tmp_l2name_count += 1
            l2name_count_array.append([tmp_device_l2_directly_l3vport_array[0],tmp_device_l2_directly_l3vport_array[2][2],tmp_l2name_count])

        #print('###l2name_count_array')
        #print(l2name_count_array)


        #main
        for tmp_device_l2_directly_l3vport_array in device_l2_directly_l3vport_array:
            device_name = tmp_device_l2_directly_l3vport_array[0]
            if_name = tmp_device_l2_directly_l3vport_array[2][2]
            tmp_opposite_if_array = get_l2_broadcast_domains.get_opposite_if(self, device_name, if_name)
            tmp_opposite_l2seg_num = get_l2_broadcast_domains.get_l2seg_count_on_if(self, tmp_opposite_if_array[0], tmp_opposite_if_array[1])
            if tmp_opposite_l2seg_num == 0:
                #print(tmp_device_l2_directly_l3vport_array)
                tmp_source_count = 0
                tmp_target_count = 0
                for tmp_l2name_count_array in l2name_count_array:
                    #print(tmp_l2name_count_array)
                    if device_name == tmp_l2name_count_array[0] and if_name == tmp_l2name_count_array[1]:
                        tmp_source_count = tmp_l2name_count_array[2]
                    if tmp_opposite_if_array[0] == tmp_l2name_count_array[0] and if_name == tmp_opposite_if_array[1]:
                        tmp_target_count = tmp_l2name_count_array[2]
                #get opposite l2segment number and name
                opposite_l2seg_name = ''
                opposite_l2seg_num = 0
                for na_device_l2_directly_l3vport_array in device_l2_directly_l3vport_array:
                    if na_device_l2_directly_l3vport_array[0] == tmp_opposite_if_array[0] and na_device_l2_directly_l3vport_array[2][2]== tmp_opposite_if_array[1]:
                        opposite_l2seg_num = na_device_l2_directly_l3vport_array[2][0]
                        opposite_l2seg_name = na_device_l2_directly_l3vport_array[1]

                #print('          + device_name,if_name -> ', device_name, if_name, tmp_source_count, tmp_device_l2_directly_l3vport_array[2][0], '    +++tmp_oppsite_if_array -> ', tmp_opposite_if_array, tmp_target_count, opposite_l2seg_num)
                if tmp_source_count == 1 and tmp_target_count == 1 and opposite_l2seg_num != 0:
                    #print('          + device_name,if_name -> ', device_name, if_name, tmp_source_count, tmp_device_l2_directly_l3vport_array[2][0],tmp_device_l2_directly_l3vport_array[2][3], '    +++tmp_oppsite_if_array -> ', tmp_opposite_if_array, tmp_target_count, opposite_l2seg_num)
                    l2_broadcast_group_array_2nd.append(sorted(list(set([tmp_device_l2_directly_l3vport_array[2][0],opposite_l2seg_num]))))

                    # kyuusai L3 virtual port has multiple l2 ports made as one l2 segment
                    l2seg_name_virtual_port_array.append([tmp_device_l2_directly_l3vport_array[2][3],tmp_device_l2_directly_l3vport_array[2][0]])

                if opposite_l2seg_num == 0:
                    ### device_l2_directly_l3vport_array -> device_l2_other_array ###
                    for tmp_device_l2_other_array in device_l2_other_array:
                        if tmp_device_l2_other_array[0] == tmp_opposite_if_array[0] and tmp_device_l2_other_array[2][2] == tmp_opposite_if_array[1]:
                            #print('### device_l2_directly_l3vport_array -> device_l2_other_array', tmp_opposite_if_array,tmp_device_l2_directly_l3vport_array[2][0], tmp_device_l2_other_array [2][0])
                            l2_broadcast_group_array_2nd.append(sorted(list(set([tmp_device_l2_directly_l3vport_array[2][0], tmp_device_l2_other_array[2][0]]))))

                if tmp_source_count != 1 and tmp_target_count == 1:
                    if tmp_device_l2_directly_l3vport_array[1] == opposite_l2seg_name:
                        #print('          + device_name,if_name -> ', device_name, if_name, tmp_source_count, tmp_device_l2_directly_l3vport_array[2][0], tmp_device_l2_directly_l3vport_array[1], '    +++tmp_oppsite_if_array -> ', tmp_opposite_if_array, tmp_target_count, opposite_l2seg_num,opposite_l2seg_name)
                        l2_broadcast_group_array_2nd.append(sorted(list(set([tmp_device_l2_directly_l3vport_array[2][0], opposite_l2seg_num]))))

            else:
                '''Kyusai device_l2_directly_l3vport_array -> multiple opposite_l2seg_num '''
                l2_broadcast_group_array_2nd.append([tmp_device_l2_directly_l3vport_array[2][0]])
                #print(tmp_device_l2_directly_l3vport_array[2][0],device_name,if_name)

        # kyuusai L3 virtual port has multiple l2 ports made as one l2 segment
        #print('### l2seg_name_virtual_port_array ###')
        #print(l2seg_name_virtual_port_array)

        for tmp_l2seg_name_virtual_port_array in l2seg_name_virtual_port_array:
            tmp_kyuusai_l3vport_multiple_l2port = [tmp_l2seg_name_virtual_port_array[1]]
            for tmp_tmp_l2seg_name_virtual_port_array in l2seg_name_virtual_port_array:
                if tmp_l2seg_name_virtual_port_array[0] == tmp_tmp_l2seg_name_virtual_port_array[0] and tmp_l2seg_name_virtual_port_array[1] != tmp_tmp_l2seg_name_virtual_port_array[1]:
                    tmp_kyuusai_l3vport_multiple_l2port.extend([tmp_tmp_l2seg_name_virtual_port_array[1]])

            if len(tmp_kyuusai_l3vport_multiple_l2port) != 1:
                #print(tmp_kyuusai_l3vport_multiple_l2port)
                l2_broadcast_group_array_2nd.append(sorted(list(set(tmp_kyuusai_l3vport_multiple_l2port))))

        #print(l2_broadcast_group_array_2nd)
        #print(get_l2_broadcast_domains.get_unique_list(l2_broadcast_group_array_2nd))
        l2_broadcast_group_array.extend(get_l2_broadcast_domains.get_unique_list(l2_broadcast_group_array_2nd))

        ### device_l2_other_array -> device_l2_other_array###
        #print('### device_l2_other_array -> device_l2_other_array###')
        l2_broadcast_group_array_3rd = []
        for tmp_device_l2_other_array  in device_l2_other_array:
            device_name = tmp_device_l2_other_array[0]
            if_name = tmp_device_l2_other_array[2][2]
            tmp_opposite_if_array = get_l2_broadcast_domains.get_opposite_if(self, device_name, if_name)

            if tmp_device_l2_other_array[1] == '__EDGE_L3__':
                for tmp_tmp_device_l2_other_array in device_l2_other_array:
                    if tmp_tmp_device_l2_other_array[0] == tmp_opposite_if_array[0] and tmp_tmp_device_l2_other_array[2][2] == tmp_opposite_if_array[1]:
                        l2_broadcast_group_array_3rd.append(sorted(list(set([tmp_device_l2_other_array[2][0], tmp_tmp_device_l2_other_array[2][0]]))))
                        break

            elif tmp_device_l2_other_array[1] == '__ALONE_L3__':
                l2_broadcast_group_array_3rd.append([tmp_device_l2_other_array[2][0]])

            elif tmp_device_l2_other_array[1] == '__ALONE_L2__':
                l2_broadcast_group_array_3rd.append([tmp_device_l2_other_array[2][0]])

        #print(get_l2_broadcast_domains.get_unique_list(l2_broadcast_group_array_3rd))
        l2_broadcast_group_array.extend(get_l2_broadcast_domains.get_unique_list(l2_broadcast_group_array_3rd))


        #print(l2_broadcast_group_array_3rd)
        #print('--- last, l2_broadcast_group_array ---')
        #print(get_l2_broadcast_domains.get_unique_list(sorted(l2_broadcast_group_array)))


        '''marge l2_broadcast_group_array — Union-Find'''
        _uf_parent = {}

        def _uf_find(x):
            while _uf_parent.get(x, x) != x:
                _uf_parent[x] = _uf_parent.get(_uf_parent[x], _uf_parent[x])
                x = _uf_parent[x]
            return x

        def _uf_union(a, b):
            ra, rb = _uf_find(a), _uf_find(b)
            if ra != rb:
                _uf_parent[ra] = rb

        _elem_to_groups = {}
        for gi, grp in enumerate(l2_broadcast_group_array):
            for elem in grp:
                if elem in _elem_to_groups:
                    _uf_union(_elem_to_groups[elem], gi)
                else:
                    _elem_to_groups[elem] = gi
                _uf_parent.setdefault(gi, gi)

        _root_to_elems = {}
        for gi, grp in enumerate(l2_broadcast_group_array):
            root = _uf_find(gi)
            if root not in _root_to_elems:
                _root_to_elems[root] = set()
            _root_to_elems[root].update(grp)

        marged_l2_broadcast_group_array = [sorted(elems) for elems in _root_to_elems.values()]
        marged_l2_broadcast_group_array = get_l2_broadcast_domains.get_unique_list(marged_l2_broadcast_group_array)

        '''make target_l2_broadcast_group_array — indexed lookup'''
        # Use setdefault to accumulate ALL entries per segment_id.
        # Multiple devices can share the same segment_id (same L2 segment),
        # so using plain dict assignment would silently overwrite earlier entries.
        _bd_by_id = {}
        for grp_list in device_l2_boradcast_domain_array:
            for bd_entry in grp_list:
                if bd_entry:
                    _bd_by_id.setdefault(bd_entry[0], []).append(bd_entry)
        _l3vport_by_id = {}
        for entry in device_l2_directly_l3vport_array:
            _l3vport_by_id.setdefault(entry[2][0], []).append(entry)
        _other_by_id = {}
        for entry in device_l2_other_array:
            if 'L3' in entry[1]:
                _other_by_id.setdefault(entry[2][0], []).append(entry)

        self.target_l2_broadcast_group_array = []
        for tmp_marged_l2_broadcast_group_array in marged_l2_broadcast_group_array:
            tmp_target_l2_broadcast_group_array = []
            for member_id in tmp_marged_l2_broadcast_group_array:
                for bd_entry in _bd_by_id.get(member_id, []):
                    for if_entry in bd_entry[3]:
                        if 'L3' in if_entry[3]:
                            tmp_target_l2_broadcast_group_array.append([bd_entry[1], if_entry[1]])

                for l3v in _l3vport_by_id.get(member_id, []):
                    tmp_target_l2_broadcast_group_array.append([l3v[0], l3v[2][3]])

                for oth in _other_by_id.get(member_id, []):
                    tmp_target_l2_broadcast_group_array.append([oth[0], oth[2][2]])

            #print([tmp_marged_l2_broadcast_group_array, nsm_def.get_l2_broadcast_domains.get_unique_list(tmp_target_l2_broadcast_group_array)])
            self.target_l2_broadcast_group_array.append([tmp_marged_l2_broadcast_group_array, get_l2_broadcast_domains.get_unique_list(tmp_target_l2_broadcast_group_array)])

        #print('--- target_l2_broadcast_group_array ---')
        #print(self.target_l2_broadcast_group_array)

        return ([self.update_l2_table_array,device_l2_boradcast_domain_array,device_l2_directly_l3vport_array,device_l2_other_array,marged_l2_broadcast_group_array,self.target_l2_broadcast_group_array] ) ## 'self.update_l2_table_array, device_l2_boradcast_domain_array, device_l2_directly_l3vport_array, device_l2_other_array, marged_l2_broadcast_group_array'


    def get_l2name_count_on_if(self,device_name,if_name):
        return self._l2seg_count_index.get((device_name, if_name), 0)

    def get_unique_list(seq):
        seen = []
        return [x for x in seq if x not in seen and not seen.append(x)]

    def get_l2seg_count_on_if(self,device_name,if_name):
        return self._l2seg_count_index.get((device_name, if_name), 0)

    def get_opposite_if(self,device_name,if_name):
        result = self._opposite_if_index.get((str(device_name), str(if_name)))
        if result is not None:
            return list(result)
        return None

    def get_l2_other(self):
        tmp_device_l2_other_array = []
        for tmp_update_l2_table_array in self.update_l2_table_array:
            if tmp_update_l2_table_array[3] == '' and tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[6] == '' and tmp_update_l2_table_array[7] == '':
                tmp_device_l2_other_array.append([tmp_update_l2_table_array[1], '__ALONE_L3__', [self.local_boradcast_domain_num, '_ALONE_L3_', tmp_update_l2_table_array[5], 'N/A', tmp_update_l2_table_array[2], 'N/A']])
                self.local_boradcast_domain_num += 1
            if tmp_update_l2_table_array[3] == '' and tmp_update_l2_table_array[5] == '' and tmp_update_l2_table_array[6] != '' and tmp_update_l2_table_array[7] == '':
                for tmp_tmp_update_l2_table_array in tmp_update_l2_table_array[6].split(','):
                    tmp_device_l2_other_array.append([tmp_update_l2_table_array[1], '__ALONE_L2__', [self.local_boradcast_domain_num, '__ALONE_L2__', tmp_tmp_update_l2_table_array, 'N/A', 'N/A', 'N/A']])
                    self.local_boradcast_domain_num += 1
            if tmp_update_l2_table_array[3] != '' and tmp_update_l2_table_array[5] == '' and tmp_update_l2_table_array[6] == '' and tmp_update_l2_table_array[7] == '':
                tmp_device_l2_other_array.append([tmp_update_l2_table_array[1], '__EDGE_L3__', [self.local_boradcast_domain_num, '__EDGE_L3__', tmp_update_l2_table_array[3], 'N/A', tmp_update_l2_table_array[2], 'N/A']])
                self.local_boradcast_domain_num += 1

        return (tmp_device_l2_other_array)

    def get_l2_directly_l3vport(self):
        tmp_l2_directly_l3vport_array = []
        for tmp_update_l2_table_array in self.update_l2_table_array:
            if tmp_update_l2_table_array[3] != '' and tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[6] == '' and tmp_update_l2_table_array[7] != '':
                for tmp_l2_directly_seg in tmp_update_l2_table_array[7].split(','):
                    tmp_l2_directly_l3vport_array.append([tmp_update_l2_table_array[1],tmp_l2_directly_seg.replace(' ',''),[self.local_boradcast_domain_num,'L2IF-L3VPORT',tmp_update_l2_table_array[3],tmp_update_l2_table_array[5],tmp_update_l2_table_array[2],tmp_update_l2_table_array[4]]])
                    self.local_boradcast_domain_num += 1
            if tmp_update_l2_table_array[3] != '' and tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[6] == '' and tmp_update_l2_table_array[7] == '':
                tmp_l2_directly_l3vport_array.append([tmp_update_l2_table_array[1], '__NO_L2SEG__', [self.local_boradcast_domain_num,'L2IF-L3VPORT', tmp_update_l2_table_array[3], tmp_update_l2_table_array[5], tmp_update_l2_table_array[2], tmp_update_l2_table_array[4]]])
                self.local_boradcast_domain_num += 1
        return (tmp_l2_directly_l3vport_array)

    def get_broadcast_domain(self):
        return_master_boradcast_domain_array = []
        for tmp_all_shape_list_array in self.all_shape_list_array:
            ###get target device l2 array
            tmp_boradcast_domain_array = []
            target_l2name_array = []
            for tmp_device_unique_l2name_array in self.device_unique_l2name_array:
                if tmp_device_unique_l2name_array[0] == tmp_all_shape_list_array:
                    target_l2name_array  = tmp_device_unique_l2name_array[1]
                    break

            ###get per l2 line of ls segment
            for tmp_target_l2name_array in target_l2name_array:
                tmp_master_boradcast_domain_array = []
                for tmp_update_l2_table_array in self.update_l2_table_array:
                    if tmp_update_l2_table_array[1] == tmp_all_shape_list_array and tmp_update_l2_table_array[6] != '' and tmp_target_l2name_array in tmp_update_l2_table_array[6].split(','):
                        if tmp_update_l2_table_array[5] != '':
                            tmp_master_boradcast_domain_array.append(['L2VPORT-L2SEG',tmp_update_l2_table_array[5],tmp_target_l2name_array,tmp_update_l2_table_array[4],'L2SEG'])
                            if tmp_update_l2_table_array[3] != '':
                                tmp_master_boradcast_domain_array.append(['L2IF-L2VPORT', tmp_update_l2_table_array[3], tmp_update_l2_table_array[5],tmp_update_l2_table_array[2],tmp_update_l2_table_array[4]])
                        else:

                            if tmp_update_l2_table_array[3] != '':
                                tmp_master_boradcast_domain_array.append(['L2IF-L2SEG', tmp_update_l2_table_array[3], tmp_target_l2name_array,tmp_update_l2_table_array[2],'L2SEG'])

                tmp_boradcast_domain_array.append([self.local_boradcast_domain_num,tmp_all_shape_list_array,tmp_target_l2name_array,tmp_master_boradcast_domain_array])
                self.local_boradcast_domain_num += 1
            return_master_boradcast_domain_array.append(tmp_boradcast_domain_array)
        return(return_master_boradcast_domain_array)

