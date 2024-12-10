'''
SPDX-License-Identifier: Apache-2.0

Copyright 2023 Cisco Systems, Inc. and its affiliates

Licensed under the Apache License, Version 2.0 (the "License");
you may not use this file except in compliance with the License.
You may obtain a copy of the License at

http://www.apache.org/licenses/LICENSE-2.0

Unless required by applicable law or agreed to in writing, software
distributed under the License is distributed on an "AS IS" BASIS,
WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
See the License for the specific language governing permissions and
limitations under the License.
'''

from pptx import Presentation
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN , MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE,MSO_CONNECTOR
from pptx.dml.color import RGBColor,MSO_THEME_COLOR
from pptx.dml.line import LineFormat
from pptx.shapes.connector import Connector
from pptx.util import Inches,Cm,Pt
import sys, os
import openpyxl
import math
import unicodedata

class  ns_ddx_figure_run():
    def __init__(self):
        #parameter
        ppt_min_width = 13.4  #  inches
        ppt_min_hight = 7.5     #  inches
        ppt_max_width = 56  #  inches
        ppt_max_hight = 56   #  inches

        ppt_meta_file = self.excel_file_path

        ### Shared system setting values
        self.ws_name_PPT_META = self.worksheet_name
        self.outline_margin_root_folder = 0.2  ##  inches
        self.outline_margin_sub_folder = 0.1  ##  inches

        ### shared open the files
        if os.path.isfile(self.output_ppt_file) == True:
            self.active_ppt = Presentation(self.output_ppt_file)
        else:
            self.active_ppt = Presentation()

        # width inches of slide master
        if (self.root_width + self.root_left * 2) < ppt_min_width:
            self.active_ppt.slide_width = Inches(ppt_min_width)
        elif (self.root_width + self.root_left * 2) > ppt_max_width:
            self.active_ppt.slide_width = Inches(ppt_max_width)
        else:
            self.active_ppt.slide_width = Inches(self.root_width + self.root_left * 2)

        # height inches of slide master
        if (self.root_hight + self.root_top * 1.5) < ppt_min_hight:
            self.active_ppt.slide_height = Inches(ppt_min_hight)
        elif (self.root_hight + self.root_top * 1.5) > ppt_max_hight:
            self.active_ppt.slide_height = Inches(ppt_max_hight)
        else:
            self.active_ppt.slide_height = Inches(self.root_hight + self.root_top * 1.5)

        self.input_ppt_mata_excel = openpyxl.load_workbook(ppt_meta_file)

        ### shared open setting values
        self.coord_list = []  # coord list of folders and shapes
        self.shapes_size_array = []  # add ver 2.0 . for writing l2 figure
        self.temp_list_num = 0
        self.folder_font_type = 'Calibri'
        self.folder_font_size = 10  # Pt
        self.shape_font_type = 'Calibri'
        self.shae_font_size = 6  # Pt

        '''main'''
        if self.ws_name_PPT_META != '<ALL_Worksheets>':
            ### select active worksheet
            self.input_ppt_mata_excel.active = self.input_ppt_mata_excel[self.ws_name_PPT_META]

            ### add new slide
            ns_ddx_figure_run.add_slide(self)

            ### add root folder and get location
            ns_ddx_figure_run.add_root_folder(self)

            ### add sub folders from excel also shapes
            ns_ddx_figure_run.add_sub_folder(self)

            ### add l2 material for l2 shape ###
            if self.click_value == 'L2-3-2':
                ns_ddx_figure_run.add_l2_material(self)
                ns_ddx_figure_run.add_l2_line(self)

            ### add line between folder and shape
            if self.click_value != 'L2-3-2':
                ns_ddx_figure_run.add_line(self)

            ### last, save_pptx
            ns_ddx_figure_run.save_pptx(self)

        else:
            print('<ALL_Worksheets>  selected')
            ws_list = self.input_ppt_mata_excel.get_sheet_names()

            for ws_name in ws_list:
                ### select active worksheet
                self.input_ppt_mata_excel.active = self.input_ppt_mata_excel[ws_name]

                ### clear coord_list
                self.coord_list = []  # coord list of folders and shapes

                ### add new slide
                ns_ddx_figure_run.add_slide(self)

                ### add root folder and get location
                ns_ddx_figure_run.add_root_folder(self)

                ### add sub folders from excel also shapes
                ns_ddx_figure_run.add_sub_folder(self)

                ### add line between folder and shape
                ns_ddx_figure_run.add_line(self)


            ### last, save_pptx
            ns_ddx_figure_run.save_pptx(self)

    def add_slide(self):
        self.title_only_slide_layout = self.active_ppt.slide_layouts[5]
        self.slide = self.active_ppt.slides.add_slide(self.title_only_slide_layout)


    def add_root_folder(self):
        temp_rootfolder_row = 0
        temp_rootfolder_flag = False
        temp_row = 1
        self.root_folder = [0.28, 1.42, 9.45, 5.75]  # Defalut setting.  left , top , width , hight(Inches)  / EMU is 914400‬

        while temp_rootfolder_row < 50000 and temp_rootfolder_flag == False:
            temp_rootfolder_row += 1
            if str(self.input_ppt_mata_excel.active.cell(temp_rootfolder_row, 1).value) == '<<ROOT_FOLDER>>':
                temp_rootfolder_flag = True
                temp_row = temp_rootfolder_row

        ### input page title text and input root holder values
        self.shape = self.slide.shapes
        self.shape.title.text = str(self.input_ppt_mata_excel.active.cell(temp_row+1, 2).value)
        self.root_folder[0] = float(self.input_ppt_mata_excel.active.cell(temp_row+1, 5).value)
        self.root_folder[1] = float(self.input_ppt_mata_excel.active.cell(temp_row+1, 6).value)
        self.root_folder[2] = float(self.input_ppt_mata_excel.active.cell(temp_row+1, 7).value) * float(self.input_ppt_mata_excel.active.cell(temp_row+1, 3).value)
        self.root_folder[3] = float(self.input_ppt_mata_excel.active.cell(temp_row+1, 8).value) * float(self.input_ppt_mata_excel.active.cell(temp_row+1, 4).value)

        ### change ppt title , add ver 2.1 ###
        if self.click_value == 'VPN-1-1':
            self.shape.title.text = '[VPNs on L1] All Areas'

        # style size,outline
        self.shape = self.shape.add_shape(MSO_SHAPE.RECTANGLE, Inches(self.root_folder[0]), Inches(self.root_folder[1]), Inches(self.root_folder[2]), Inches(self.root_folder[3]))
        #style fill
        shape_fill = self.shape.fill

        #chage outline coler to 255,255,255 at ver 2.3.0
        #shape_fill = shape_fill.background()
        shape_fill.solid()
        shape_fill.fore_color.rgb = RGBColor(255, 255, 255)

        # style line
        shape_line = self.shape.line
        shape_line.color.rgb = RGBColor(0, 0, 0)
        shape_line.color.brightness = 0.0
        shape_line.width = Pt(1.5)
        shape_line.fill.solid()
        self.shape.shadow.inherit = False

        ### adjust outline margin
        self.root_folder[0] = self.root_folder[0] + self.outline_margin_root_folder
        self.root_folder[1] = self.root_folder[1] + self.outline_margin_root_folder
        self.root_folder[2] = self.root_folder[2] - (self.outline_margin_root_folder*2)
        self.root_folder[3] = self.root_folder[3] - (self.outline_margin_root_folder*2)

    def add_sub_folder(self):
        ###get number of row from meta excel
        temp_row = 1
        temp_temp_row_folder = 0
        excel_row_sum = 0
        excel_col_sum = 0
        temp_count_set_width = 0
        self.sub_folder_list = [[],[]]  #  row list , column list
        self.margin_top_text = 0.0
        self.margin_bottom_text = -0.05
        temp_subfolder_row = 0
        temp_subfolder_flag = False

        while temp_subfolder_row < 50000 and temp_subfolder_flag == False:
            temp_subfolder_row += 1
            if str(self.input_ppt_mata_excel.active.cell(temp_subfolder_row, 1).value) == '<<POSITION_FOLDER>>':
                temp_subfolder_flag = True
                temp_row = temp_subfolder_row

        while self.input_ppt_mata_excel.active.cell(temp_row, 1).value != None:
            if self.input_ppt_mata_excel.active.cell(temp_row, 1).value != '<SET_WIDTH>' and str(self.input_ppt_mata_excel.active.cell(temp_row, 1).value).startswith('<<POSITION_FOLDER>>') != True:
                self.sub_folder_list[0].append(self.input_ppt_mata_excel.active.cell(temp_row, 1).value)
                excel_row_sum = excel_row_sum + float(self.input_ppt_mata_excel.active.cell(temp_row, 1).value)
            elif self.input_ppt_mata_excel.active.cell(temp_row, 1).value == '<SET_WIDTH>' or str(self.input_ppt_mata_excel.active.cell(temp_row, 1).value).startswith('<<POSITION_FOLDER>>') == True:
                temp_count_set_width +=1
            temp_row += 1

        #### Set initial value
        self.sub_folder = [self.root_folder[0],self.root_folder[1],self.root_folder[2],self.root_folder[3]]  # left , top , width , hight(Inches)  / EMU is 914400‬

        '''write sub folders'''
        for temp_row_folder in range(temp_subfolder_row, temp_row):
            if self.input_ppt_mata_excel.active.cell(temp_row_folder, 1).value == '<SET_WIDTH>' or str(self.input_ppt_mata_excel.active.cell(temp_row_folder, 1).value).startswith('<<POSITION_FOLDER>>') == True:
                ### get num col in each row
                excel_col_sum = 0
                temp_column = 2
                self.sub_folder_list[1] = []
                while self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_column).value != None:
                    self.sub_folder_list[1].append(self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_column).value)
                    excel_col_sum = excel_col_sum + float(self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_column).value)
                    temp_column +=1

            if self.input_ppt_mata_excel.active.cell(temp_row_folder, 1).value != '<SET_WIDTH>' and str(self.input_ppt_mata_excel.active.cell(temp_row_folder, 1).value).startswith('<<POSITION_FOLDER>>') != True:
                ###write folders on the row
                for temp_col_folder in range(0, len(self.sub_folder_list[1])):
                    ### get width and hight reflected ratio
                    self.sub_folder[2] = (self.sub_folder_list[1][temp_col_folder] / excel_col_sum) * self.root_folder[2] #col
                    self.sub_folder[3] = (self.sub_folder_list[0][temp_temp_row_folder] / excel_row_sum) * self.root_folder[3] #row

                    ### adjust margin
                    folder_left = self.sub_folder[0]
                    folder_top = self.sub_folder[1]
                    folder_width = self.sub_folder[2]
                    folder_hight = self.sub_folder[3]

                    ### write sub folder
                    self.shape = self.slide.shapes
                    self.shape = self.shape.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(folder_left), Inches(folder_top), Inches(folder_width), Inches(folder_hight))

                    ###add self.coord_list for line
                    self.temp_folder_text = 'dummy'
                    if self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_col_folder+2).value != None:
                        x_left = folder_left
                        x_middle = folder_left + (folder_width*0.5)
                        x_right = folder_left + folder_width
                        y_top = folder_top
                        y_middle = folder_top + (folder_hight*0.5)
                        y_down = folder_top + folder_hight

                        self.temp_folder_text =  str(self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_col_folder+2).value)
                        ### refrect tag function own shape
                        if '<' in self.temp_folder_text:
                            self.temp_folder_text = str('<') + str(self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_col_folder+2).value).split('<')[-1].split('>')[0] + str('>')

                        self.coord_list.append([self.temp_folder_text,x_left,x_middle,x_right,y_top,y_middle,y_down]) #coord of x_left,x_middle,x_right,y_top,y_middle,y_down
                        self.temp_list_num += 1

                    ### pass to add_shape()
                    self.shape_folder_left = self.sub_folder[0]
                    self.shape_folder_top = self.sub_folder[1]
                    self.shape_folder_width = self.sub_folder[2]
                    self.shape_folder_hight = self.sub_folder[3]

                    ### chenge left and top of start folder
                    self.sub_folder[0] = self.sub_folder[0] + self.sub_folder[2]

                    shape_fill = self.shape.fill
                    shape_fill = shape_fill.background()

                    shape_line = self.shape.line
                    shape_line.color.rgb = RGBColor(0, 0, 0)
                    shape_line.color.brightness = 0.0
                    shape_line.width = Pt(1.5)

                    self.shape.adjustments[0] = 0.1  # curve of ROUNDED_RECTANGLE 0.0~1.0
                    self.shape.shadow.inherit = False  # disalbe dealut shadow effect

                    '''change stlye from meta file'''
                    temp_style_row = 1
                    temp_style_flag = False
                    temp_match_flag = False
                    while temp_style_row < 50000 and temp_match_flag == False:
                        if str(self.input_ppt_mata_excel.active.cell(temp_style_row, 1).value) == '<<STYLE_FOLDER>>':
                            temp_style_flag = True
                        ## Default set###
                        if str(self.input_ppt_mata_excel.active.cell(temp_style_row, 1).value) == '<DEFAULT>':
                            temp_default_value = self.input_ppt_mata_excel.active.cell(temp_style_row, 2).value
                            if temp_default_value == 'YES':
                                shape_line.fill.solid()
                            elif temp_default_value == 'NO':
                                shape_line.fill.background()

                        ## empty valse set###
                        if temp_style_flag == True and self.input_ppt_mata_excel.active.cell(temp_style_row, 1).value == '<EMPTY>':
                            temp_empty_value = self.input_ppt_mata_excel.active.cell(temp_style_row, 2).value
                            if self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_col_folder+2).value == None:
                                if temp_empty_value == 'YES':
                                    shape_line.fill.solid()
                                elif temp_empty_value == 'NO':
                                    shape_line.fill.background()
                                    ''' delete shape for network sketcher'''
                                    sp = self.shape._element
                                    sp.getparent().remove(sp)
                                break

                        ## Outline and Adjust Margin Option set###
                        if temp_style_flag == True and str(self.input_ppt_mata_excel.active.cell(temp_style_row, 1).value) == str(self.temp_folder_text):
                            temp_match_flag = True
                            if self.input_ppt_mata_excel.active.cell(temp_style_row, 2).value == 'YES':
                                shape_line.fill.solid()
                            elif self.input_ppt_mata_excel.active.cell(temp_style_row, 2).value == 'NO':
                                shape_line.fill.background()
                                ''' delete shape for network sketcher'''
                                sp = self.shape._element
                                sp.getparent().remove(sp)

                            if self.input_ppt_mata_excel.active.cell(temp_style_row, 4).value != '<AUTO>':
                                self.shape_folder_top = self.shape_folder_top + self.input_ppt_mata_excel.active.cell(temp_style_row, 4).value
                                self.shape_folder_hight = self.shape_folder_hight - self.input_ppt_mata_excel.active.cell(temp_style_row, 4).value

                            if self.input_ppt_mata_excel.active.cell(temp_style_row, 5).value != '<AUTO>':
                                self.shape_folder_hight = self.shape_folder_hight - self.input_ppt_mata_excel.active.cell(temp_style_row, 5).value

                            ### Add Text into the sub folder, when UP
                            if self.input_ppt_mata_excel.active.cell(temp_style_row, 3).value == 'UP':
                                if '<' in str(self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_col_folder+2).value):
                                    self.shape.text = str(self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_col_folder+2).value).split('<')[0]
                                    self.shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
                                    self.shape.text_frame.margin_top = Inches(self.margin_top_text)

                                else:
                                    self.shape.text = str(self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_col_folder+2).value)
                                    self.shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
                                    self.shape.text_frame.margin_top = Inches(self.margin_top_text)

                            ### Add Text into the sub folder, when DOWN
                            if self.input_ppt_mata_excel.active.cell(temp_style_row, 3).value == 'DOWN':
                                if '<' in str(self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_col_folder+2).value):
                                    self.shape.text = str(self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_col_folder+2).value).split('<')[0]
                                    self.shape.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
                                    self.shape.text_frame.margin_bottom = Inches(self.margin_bottom_text)
                                else:
                                    self.shape.text = str(self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_col_folder+2).value)
                                    self.shape.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
                                    self.shape.text_frame.margin_bottom = Inches(self.margin_bottom_text)

                            if self.input_ppt_mata_excel.active.cell(temp_style_row, 3).value == 'UP' or self.input_ppt_mata_excel.active.cell(temp_style_row, 3).value == 'DOWN':
                                self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                                self.shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                                self.shape.text_frame.paragraphs[0].font.name = self.folder_font_type
                                self.shape.text_frame.paragraphs[0].font.size = Pt(self.folder_font_size)
                            break

                        ## empty value set ###
                        if temp_style_flag == True and self.input_ppt_mata_excel.active.cell(temp_style_row, 1).value == None:
                            break

                        temp_style_row += 1

                    '''run add shapes in own sub folder'''
                    self.ddx_meta_row_folder = temp_row_folder
                    self.ddx_meta_col_folder = temp_col_folder + 2
                    if self.input_ppt_mata_excel.active.cell(temp_row_folder, temp_col_folder + 2).value != None:
                        ns_ddx_figure_run.add_shape(self)

                temp_temp_row_folder +=1
                self.sub_folder[0] = self.root_folder[0]
                self.sub_folder[1] = self.sub_folder[1] + self.sub_folder[3]


    def add_shape(self):
        print('Write Folder ---> ' + str(self.input_ppt_mata_excel.active.cell(self.ddx_meta_row_folder, self.ddx_meta_col_folder).value))
        #print(self.shape_folder_left,self.shape_folder_top,self.shape_folder_width,self.shape_folder_hight)

        '''Define reflect margin in the sub folder'''
        folder_left = self.shape_folder_left + self.outline_margin_sub_folder
        folder_top = self.shape_folder_top + self.outline_margin_sub_folder
        folder_width = self.shape_folder_width - (self.outline_margin_sub_folder * 2)
        folder_hight = self.shape_folder_hight - (self.outline_margin_sub_folder * 2)
        ## to use shapes in subholder
        shape_left = folder_left
        shape_top = folder_top
        shape_width = folder_width
        shape_hight = folder_hight

        ### get row of <<POSITION_SHAPE>> ###
        temp_positon_row = 1
        temp_position_flag = False
        while temp_positon_row < 50000 and temp_position_flag == False:
            if str(self.input_ppt_mata_excel.active.cell(temp_positon_row, 1).value) == '<<POSITION_SHAPE>>':
                temp_position_flag = True
            temp_positon_row += 1

        ### Search row number of the sub_folder
        temp_subfolder_row = temp_positon_row
        temp_subfolder_flag = False
        temp_subfolder_exceed_flag = False

        while temp_subfolder_row < 50000 and temp_subfolder_flag == False:
            if str(self.input_ppt_mata_excel.active.cell(temp_subfolder_row, 1).value) == str(self.temp_folder_text):
                temp_subfolder_flag = True
            if str(self.input_ppt_mata_excel.active.cell(temp_subfolder_row, 1).value).startswith('<<') == True:
                temp_subfolder_exceed_flag = True
                #print(str(self.input_ppt_mata_excel.active.cell(self.ddx_meta_row_folder, self.ddx_meta_col_folder).value) + ' is not found under <<POSITION_SHAPE>> in excel')
            temp_subfolder_row += 1

        ### get num of row about the sub folder
        if temp_subfolder_flag == True and temp_subfolder_exceed_flag == False:
            temp_row = temp_subfolder_row -1
            excel_row_sum = 0
            self.shape_width_sum = 0.0
            self.shape_hight_sum = 0.0
            temp_subfolder_row_flag = False
            temp_subfolder_col_flag = False
            temp_subfolder_row_exceed_flag = False

            while self.input_ppt_mata_excel.active.cell(temp_row, 1).value != '<END>' and temp_subfolder_row_flag == False:
                if self.input_ppt_mata_excel.active.cell(temp_row, 1).value == '<END>':
                    temp_subfolder_row_flag = True
                if temp_row > 50000:
                    temp_subfolder_row_exceed_flag = True

                ### get sum of hight of shapes
                if temp_subfolder_row_flag == False:
                    self.shape_hight_sum = self.shape_hight_sum + ns_ddx_figure_run.get_shape_highest(self,temp_row)
                    #print('####CHECK#### ',self.shape.text,ns_ddx_figure_run.get_shape_highest(self,temp_row),temp_row)

                ### process of before while loop
                temp_row += 1
                excel_row_sum += 1

            ### write each shape into sub folder
            if temp_subfolder_row_exceed_flag == False:
                for temp_row_subfolder in range(temp_subfolder_row-1, temp_subfolder_row+excel_row_sum-1):
                    excel_col_sum = 0
                    self.shape_width_sum = 0
                    temp_col = 2

                    ### get num of col on the row
                    while self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_col).value != '<END>' and temp_subfolder_col_flag == False:
                        if self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_col).value == '<END>':
                            temp_subfolder_col_flag = True
                        if temp_col > 1000:
                            temp_subfolder_col_flag = True

                        ### get sum of width of shapes
                        if temp_subfolder_col_flag == False:
                            self.shape_width_sum += ns_ddx_figure_run.get_shape_width(self,self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_col).value)

                        ### process of before while loop
                        temp_col += 1
                        excel_col_sum += 1

                    '''write each shape in the own sub folder'''
                    ### get width and hight margin
                    temp_width_margin = (folder_width - self.shape_width_sum) / (excel_col_sum + 1)
                    temp_hight_margin = (folder_hight - self.shape_hight_sum) / (excel_row_sum + 1)
                    shape_top += temp_hight_margin

                    for temp_sub_col_count in range(2,excel_col_sum+2):
                        shape_left = shape_left + temp_width_margin
                        shape_degree = 0.0 # initial degree
                        shape_width = ns_ddx_figure_run.get_shape_width(self,self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_sub_col_count).value)
                        shape_hight = ns_ddx_figure_run.get_shape_hight(self,self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_sub_col_count).value)
                        segment_flag = False

                        ### check meant to segment as <SEGMENT>
                        if str(self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_sub_col_count).value).startswith('<SEGMENT>') == True:
                            segment_flag = True
                            ns_ddx_figure_run.add_segment_line(self,shape_left,shape_top,shape_width,shape_hight,folder_left,folder_top,folder_width,folder_hight,temp_sub_col_count,temp_row_subfolder,temp_width_margin,temp_hight_margin)

                        ###write normal shape
                        #add  and _AIR_ not included at ver 2.2.2(a)
                        if self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_sub_col_count).value != None and segment_flag == False and '_AIR_' not in str(self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_sub_col_count).value):
                            self.shape = self.slide.shapes
                            self.shape = self.shape.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(shape_left), Inches(shape_top), Inches(shape_width), Inches(shape_hight))
                            ## write text and reflected TAG function in the shape

                            if '<' in str(self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_sub_col_count).value):
                                self.shape.text = str(self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_sub_col_count).value).split('<')[0]
                            else:
                                self.shape.text = str(self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_sub_col_count).value)

                            self.shapes_size_array.append([self.shape.text,[shape_left, shape_top, shape_width, shape_hight]])  #add ver 2.0 for writing l2 shapes

                            ###add self.coord_list for line
                            x_left = shape_left
                            x_middle = shape_left + (shape_width * 0.5)
                            x_right = shape_left + shape_width
                            y_top = shape_top
                            y_middle = shape_top + (shape_hight * 0.5)
                            y_down = shape_top + shape_hight
                            temp_shape_text = str(self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_sub_col_count).value)
                            ### refrect tag function own shape
                            if '<' in temp_shape_text:
                                temp_shape_text = str('<')+str(self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_sub_col_count).value).split('<')[-1].split('>')[0]+str('>')

                            self.coord_list.append([temp_shape_text, x_left, x_middle, x_right, y_top, y_middle, y_down])  # coord of x_left,x_middle,x_right,y_top,y_middle,y_down
                            self.temp_list_num += 1

                            '''set to initial style'''
                            self.shape.adjustments[0] = 0.0

                            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                            self.shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            self.shape.text_frame.paragraphs[0].font.name = self.shape_font_type
                            self.shape.text_frame.paragraphs[0].font.size = Pt(self.shae_font_size)
                            self.shape.text_frame.margin_top = 0
                            self.shape.text_frame.margin_bottom = 0
                            self.shape.text_frame.margin_left = 0
                            self.shape.text_frame.margin_right = 0

                            shape_fill = self.shape.fill
                            shape_fill = shape_fill.background()

                            shape_line = self.shape.line
                            shape_line.color.rgb = RGBColor(0, 0, 0)
                            shape_line.color.brightness = 0.0
                            shape_line.width = Pt(1.0)

                            self.shape.adjustments[0] = shape_degree  # initial
                            self.shape.shadow.inherit = False  # disalbe dealut shadow effect

                            ### change style for _AIR_ shape###
                            if  '_AIR_' in temp_shape_text:
                                shape_line.color.rgb = RGBColor(255, 255, 255)
                                self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

                            '''change to specific style from excel meta'''
                            ### search row of <<STYLE_SHAPE>
                            temp_style_shape_row = 1
                            temp_style_shape_flag = False
                            while temp_style_shape_row < 50000 and temp_style_shape_flag == False:
                                if str(self.input_ppt_mata_excel.active.cell(temp_style_shape_row, 1).value) == '<<STYLE_SHAPE>>':
                                    temp_style_shape_flag = True
                                temp_style_shape_row += 1

                            ###get style values of shape from excel meta
                            temp_temp_style_shape_row = temp_style_shape_row
                            temp_temp_style_shape_flag = False
                            while temp_temp_style_shape_row < 50000 and temp_temp_style_shape_flag == False:
                                if self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value == None:
                                    temp_temp_style_shape_flag = True
                                    break
                                if str(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value) == str(temp_shape_text):
                                    temp_temp_style_shape_flag = True
                                    ### Degree of Roundness###
                                    self.shape.adjustments[0] = float(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 4).value) # curve of ROUNDED_RECTANGLE 0.0~1.0

                                    ### fill Color ###
                                    if self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 5).value == 'ORANGE':
                                        shape_fill = self.shape.fill
                                        shape_fill.solid()
                                        shape_fill.fore_color.rgb  = RGBColor(253, 234, 218)
                                    elif self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 5).value == 'BLUE':
                                        shape_fill = self.shape.fill
                                        shape_fill.solid()
                                        shape_fill.fore_color.rgb = RGBColor(220, 230, 242)
                                    elif self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 5).value == 'GREEN':
                                        shape_fill = self.shape.fill
                                        shape_fill.solid()
                                        shape_fill.fore_color.rgb = RGBColor(235, 241, 222)
                                    elif self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 5).value == 'GRAY':
                                        shape_fill = self.shape.fill
                                        shape_fill.solid()
                                        shape_fill.fore_color.rgb = RGBColor(242, 242, 242)

                                    ### change style for _AIR_ shape###
                                    if  '_AIR_' in temp_shape_text:
                                        shape_fill = self.shape.fill
                                        shape_fill.solid()
                                        shape_fill.fore_color.rgb = RGBColor(255, 255, 255)

                                    '''
                                    ### change style for L2 materials shape### ver 2.0 for device frame, wp frame
                                    '''
                                    ### DEVICE FRAME ###
                                    if self.click_value == 'L2-3-2' and self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 5).value == 'GREEN':
                                        self.shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                                        self.shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
                                        shape_fill.fore_color.rgb = RGBColor(250, 251, 247)
                                        self.shape.text_frame.paragraphs[0].font.size = Pt(self.folder_font_size)
                                        self.shape.text_frame.margin_left = Inches(0.1)
                                        self.shape.text_frame.margin_top = Inches(0.05)

                                    ### WAY POINT ###
                                    if self.click_value == 'L2-3-2' and self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 5).value == 'BLUE':
                                        shape_fill.fore_color.rgb = RGBColor(237, 242, 249)
                                        self.shape.adjustments[0] = 0.1
                                        self.shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                                        self.shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
                                        self.shape.text_frame.paragraphs[0].font.size = Pt(self.folder_font_size)
                                        self.shape.text_frame.margin_left = Inches(0.1)

                                temp_temp_style_shape_row += 1


                        ### process of before loop
                        shape_left = shape_left + shape_width

                    # process before change row
                    shape_top += ns_ddx_figure_run.get_shape_highest(self,temp_row_subfolder)
                    shape_left = folder_left

    def add_line(self):
        #print('### self.coord_list ###')
        #print(self.coord_list)

        ### get <<POSITION_LINE>> from excel
        temp_line_row = 1
        temp_line_flag = False
        while temp_line_row < 50000 and temp_line_flag == False:
            if str(self.input_ppt_mata_excel.active.cell(temp_line_row, 1).value) == '<<POSITION_LINE>>':
                temp_line_flag = True
            temp_line_row += 1

        ### read each row under <<POSITION_LINE>>
        temp_temp_line_row = temp_line_row + 1
        while temp_temp_line_row < 50000 and self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 1).value != None:
            ''' write each line From To'''
            From_name = self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 1).value
            To_name = self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 2).value
            from_connect_x = 2
            from_connect_y = 5
            to_connect_x = 2
            to_connect_y = 5

            self.coord_match_from_flag = False
            self.coord_match_to_flag = False
            ###GET From index type as list
            for y, row in enumerate(self.coord_list):
                try:
                    temp_from = (y, row.index(From_name))
                    self.coord_match_from_flag = True
                    break
                except ValueError:
                    pass
            ###GET To index
            for y, row in enumerate(self.coord_list):
                try:
                    temp_to = (y, row.index(To_name))
                    self.coord_match_to_flag = True
                    break
                except ValueError:
                    pass

            if self.coord_match_from_flag == True and self.coord_match_to_flag == True:
                ### compare the coord value of shape and deside connect side on the shape
                if self.coord_list[temp_from[0]][6] >= self.coord_list[temp_to[0]][4]:
                    from_connect_y = 4
                    to_connect_y = 6
                else:
                    from_connect_y = 6
                    to_connect_y = 4

                ### change from_side and to_side of each shape coord
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 5).value) == 'RIGHT':
                    from_connect_x = 3
                    from_connect_y = 5
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 5).value) == 'LEFT':
                    from_connect_x = 1
                    from_connect_y = 5
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 6).value) == 'RIGHT':
                    to_connect_x = 3
                    to_connect_y = 5
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 6).value) == 'LEFT':
                    to_connect_x = 1
                    to_connect_y = 5

                ### set coord values
                inche_from_connect_x = float(self.coord_list[temp_from[0]][from_connect_x])
                inche_from_connect_y = float(self.coord_list[temp_from[0]][from_connect_y])
                inche_to_connect_x = float(self.coord_list[temp_to[0]][to_connect_x])
                inche_to_connect_y = float(self.coord_list[temp_to[0]][to_connect_y])

                ### apply offset value to x y coord
                if self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 7).value != None: # From X
                    inche_from_connect_x = self.coord_list[temp_from[0]][from_connect_x] + self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 7).value

                if self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 8).value != None: # From Y
                    inche_from_connect_y = self.coord_list[temp_from[0]][from_connect_y] + self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 8).value

                if self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 9).value != None and self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 9).value != '<FROM_X>': # To X
                    inche_to_connect_x = self.coord_list[temp_to[0]][to_connect_x] + self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 9).value

                if self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 10).value != None and self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 10).value !=  '<FROM_Y>': # To Y
                    inche_to_connect_y = self.coord_list[temp_to[0]][to_connect_y] + self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 10).value

                '''if Offset To_X or Y cell is <FROM_X> below <<POSITION_LINE>> in excel '''
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 9).value) == '<FROM_X>':
                    inche_to_connect_x = inche_from_connect_x
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 10).value) == '<FROM_Y>':
                    inche_to_connect_y = inche_from_connect_y

                #write line

                self.shape = self.slide.shapes
                self.shape = self.shape.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(inche_from_connect_x), Inches(inche_from_connect_y), Inches(inche_to_connect_x), Inches(inche_to_connect_y))

                '''change style of line'''
                shape_line = self.shape.line
                shape_line.color.rgb = RGBColor(0, 0, 0)
                shape_line.color.brightness = 0.0
                shape_line.width = Pt(0.5)
                shape_line.fill.solid()
                self.shape.shadow.inherit = False

                if self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 12).value == 'NO':
                    shape_line.fill.background()

                '''write oval mark meant to channel '''
                if self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 11).value != None:
                    #calc degree
                    temp_y = inche_to_connect_y - inche_from_connect_y
                    temp_x = inche_to_connect_x - inche_from_connect_x
                    temp_degree = math.degrees(math.atan2(temp_y, temp_x))

                    #calc location
                    if self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 5).value != None:
                        temp_oval_width = 0.1
                        temp_oval_hight = float(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 11).value)
                        temp_oval_left = inche_from_connect_x + (((inche_to_connect_x - inche_from_connect_x) * 0.5) - (temp_oval_width * 0.5))
                        temp_oval_top  = inche_from_connect_y + (((inche_to_connect_y - inche_from_connect_y) * 0.5) - (temp_oval_hight * 0.5))
                    else:
                        temp_oval_width = 0.1
                        temp_oval_hight = float(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 11).value)
                        temp_oval_left = inche_from_connect_x + (((inche_to_connect_x - inche_from_connect_x) * 0.5) - (temp_oval_width * 0.5))
                        temp_oval_top  = inche_from_connect_y + (((inche_to_connect_y - inche_from_connect_y) * 0.5) - (temp_oval_hight * 0.5))

                    self.shape = self.slide.shapes
                    self.shape = self.shape.add_shape(MSO_SHAPE.OVAL, Inches(temp_oval_left), Inches(temp_oval_top), Inches(temp_oval_width), Inches(temp_oval_hight)) #Inches(shape_left), Inches(shape_top), Inches(shape_width), Inches(shape_hight)
                    self.shape.rotation = temp_degree

                    shape_fill = self.shape.fill
                    shape_fill = shape_fill.background()
                    shape_line = self.shape.line
                    shape_line.color.rgb = RGBColor(0, 0, 0)
                    shape_line.color.brightness = 0.0
                    shape_line.width = Pt(0.3)
                    self.shape.shadow.inherit = False  # disalbe dealut shadow effect

                '''Add line tag on the line'''
                ### From side tag
                if self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 3).value != None:
                    ns_ddx_figure_run.add_line_tag(self,'From', inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y,temp_temp_line_row,From_name)
                ### To side tag
                if self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 4).value != None:
                    ns_ddx_figure_run.add_line_tag(self,'To', inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y,temp_temp_line_row,To_name)

            '''last process, increment row num '''
            temp_temp_line_row += 1

    def add_line_tag(self,tag_side, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y,temp_temp_line_row,Target_name):
        self.tag_font_size = 6
        import ns_def
        ### calc tag_width
        per_char_width = 0.045   # tuning value
        font_size_hight = 0.014  # tuning value
        if tag_side == 'From':
            tag_width = ns_ddx_figure_run.get_east_asian_width_count(self,str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 3).value)) * per_char_width
            tag_hight = self.tag_font_size * font_size_hight
            #tag_width = ns_def.get_description_width_hight(self.shae_font_size,str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 3).value))[0]
            #tag_hight = ns_def.get_description_width_hight(self.shae_font_size,str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 3).value))[1]
        elif tag_side == 'To':
            tag_width = ns_ddx_figure_run.get_east_asian_width_count(self,str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 4).value)) * per_char_width
            tag_hight = self.tag_font_size * font_size_hight
            #tag_width = ns_def.get_description_width_hight(self.shae_font_size, str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 4).value))[0]
            #tag_hight = ns_def.get_description_width_hight(self.shae_font_size, str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 4).value))[1]
        elif tag_side == 'Segment':
            tag_width = ns_ddx_figure_run.get_east_asian_width_count(self,str(Target_name)) * per_char_width
            tag_hight = self.tag_font_size * font_size_hight
            #tag_width = ns_def.get_description_width_hight(self.shae_font_size, str(Target_name))[0]
            #tag_hight = ns_def.get_description_width_hight(self.shae_font_size, str(Target_name))[1]

        ###get <<POSITION_TAG>> from excel
        temp_line_row = 1
        temp_line_flag = False
        while temp_line_row < 50000 and temp_line_flag == False:
            if str(self.input_ppt_mata_excel.active.cell(temp_line_row, 1).value) == '<<POSITION_TAG>>':
                temp_line_flag = True
            temp_line_row += 1

        ### Set Default value
        tag_type = str(self.input_ppt_mata_excel.active.cell(temp_line_row, 2).value)
        temp_tag_degree = 0

        ### read each row under <<POSITION_LINE>>
        temp_temp_temp_line_row = temp_line_row + 1
        temp_tag_flag = False
        while temp_temp_temp_line_row < 50000 and self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 1).value != None and temp_tag_flag == False:
            if self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 2).value != None:
                tag_type = str(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 2).value)

            if tag_side == 'From' and tag_type == 'SHAPE':
                tag_left  = inche_from_connect_x - (tag_width * 0.5)
                tag_top   = inche_from_connect_y - (tag_hight * 0.5)
                temp_tag_text = str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 3).value)

                ### get position tag
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 1).value) == Target_name and self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 3).value != None:
                    tag_left += float(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 3).value)
                    temp_tag_flag = True
                elif str(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 1).value) == Target_name:
                    temp_tag_flag = True
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 1).value) == Target_name and self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 4).value != None:
                    tag_top += float(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 4).value)
                    temp_tag_flag = True
                elif str(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 1).value) == Target_name:
                    temp_tag_flag = True

            elif tag_side == 'To'and tag_type == 'SHAPE':
                tag_left  = inche_to_connect_x - (tag_width * 0.5)
                tag_top   = inche_to_connect_y - (tag_hight * 0.5)
                temp_tag_text = str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 4).value)

                ### get position tag
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 1).value) == Target_name and self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 3).value != None:
                    tag_left += float(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 3).value)
                    temp_tag_flag = True
                elif str(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 1).value) == Target_name:
                    temp_tag_flag = True
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 1).value) == Target_name and self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 4).value != None:
                    tag_top += float(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 4).value)
                    temp_tag_flag = True
                elif str(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 1).value) == Target_name:
                    temp_tag_flag = True

            elif tag_side == 'From' and tag_type == 'LINE':
                tag_left  = inche_from_connect_x - (tag_width * 0.5)
                tag_top   = inche_from_connect_y - (tag_hight * 0.5)
                temp_tag_text = str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 3).value)

                #### calc degree
                temp_y = inche_to_connect_y - inche_from_connect_y
                temp_x = inche_to_connect_x - inche_from_connect_x
                temp_degree = math.degrees(math.atan2(temp_y, temp_x))

                ### change position tag
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 1).value) == Target_name and self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 5).value != None:
                    temp_tag_flag = True
                    tag_left += (math.cos(math.radians(temp_degree)) * self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 5).value)
                    tag_top += (math.sin(math.radians(temp_degree)) * self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 5).value)

                if str(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 1).value) == Target_name and self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 6).value == 'YES':
                    temp_tag_degree = temp_degree
                    temp_tag_flag = True
                    if temp_x < 0:
                        temp_tag_degree += 180

            elif tag_side == 'To' and tag_type == 'LINE':
                tag_left  = inche_to_connect_x - (tag_width * 0.5)
                tag_top   = inche_to_connect_y - (tag_hight * 0.5)
                temp_tag_text = str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 4).value)

                #### calc degree
                temp_y = inche_to_connect_y - inche_from_connect_y
                temp_x = inche_to_connect_x - inche_from_connect_x
                temp_degree = math.degrees(math.atan2(temp_y, temp_x))

                ### change position tag
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 1).value) == Target_name and self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 5).value != None:
                    temp_tag_flag = True
                    tag_left -= (math.cos(math.radians(temp_degree)) * self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 5).value)
                    tag_top -= (math.sin(math.radians(temp_degree))*self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 5).value)

                if str(self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 1).value) == Target_name and self.input_ppt_mata_excel.active.cell(temp_temp_temp_line_row, 6).value == 'YES':
                    temp_tag_degree = temp_degree
                    temp_tag_flag = True
                    if temp_x < 0:
                        temp_tag_degree += 180

            elif tag_side == 'Segment':
                self.adjust_segment_char_y = 0.025
                self.segment_font_size = 8
                self.tag_font_size = self.segment_font_size
                tag_left  = inche_to_connect_x - (tag_width)
                tag_top   = inche_to_connect_y - (tag_hight) + self.adjust_segment_char_y
                temp_tag_text = Target_name

            '''last process, increment row num '''
            temp_temp_temp_line_row += 1

        ###Write Tag Name
        '''if From is <BULLET>'''
        if tag_side == 'From' and temp_tag_text == '<BULLET>':
            tag_left = inche_from_connect_x - (self.size_bullet * 0.5)
            tag_top = inche_from_connect_y - (self.size_bullet * 0.5)
            tag_width = self.size_bullet
            tag_hight = self.size_bullet

        '''if To is <BULLET>'''
        if tag_side == 'To' and temp_tag_text == '<BULLET>':
            tag_left = inche_to_connect_x - (self.size_bullet * 0.5)
            tag_top = inche_to_connect_y - (self.size_bullet * 0.5)
            tag_width = self.size_bullet
            tag_hight = self.size_bullet

        self.shape = self.slide.shapes
        self.shape = self.shape.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(tag_left), Inches(tag_top), Inches(tag_width), Inches(tag_hight))
        self.shape.text = temp_tag_text
        self.shape.rotation = temp_tag_degree

        '''set initial style'''
        self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        self.shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        self.shape.text_frame.paragraphs[0].font.name = self.shape_font_type
        self.shape.text_frame.paragraphs[0].font.size = Pt(self.tag_font_size)
        self.shape.text_frame.margin_top = 0
        self.shape.text_frame.margin_bottom = 0
        self.shape.text_frame.margin_left = 0
        self.shape.text_frame.margin_right = 0
        self.shape.text_frame.auto_size = True
        self.shape.text_frame.word_wrap = False

        shape_fill = self.shape.fill
        shape_fill.solid()
        shape_fill.fore_color.rgb = RGBColor(255, 255, 255)

        shape_line = self.shape.line
        shape_line.color.rgb = RGBColor(0, 0, 0)
        shape_line.color.brightness = 0.0
        shape_line.width = Pt(0.1)

        self.shape.adjustments[0] = 0.99445  # initial    #0.99445 is do not change for identify if tag
        self.shape.shadow.inherit = False  # disalbe dealut shadow effect

        '''if To or From is <BULLET>'''
        if temp_tag_text == '<BULLET>':
            self.shape.text = ''
            shape_fill = self.shape.fill
            shape_fill.solid()
            shape_fill.fore_color.rgb = RGBColor(0, 0, 0)
            shape_line = self.shape.line
            shape_line.fill.background()
            self.shape.adjustments[0] = 1.0  # initial
            self.shape.shadow.inherit = False  # disalbe dealut shadow effect

        '''set segment line and write start and end bullet'''
        if tag_side == 'Segment':
            self.shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            shape_fill = self.shape.fill
            shape_fill.background()
            shape_line = self.shape.line
            shape_line.fill.background()
            self.shape.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

            ###write bullet at start and end of the line
            self.size_bullet = 0.04
            self.shape = self.slide.shapes
            self.shape = self.shape.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(inche_from_connect_x-(self.size_bullet*0.5)), Inches(inche_from_connect_y-(self.size_bullet*0.5)), Inches(self.size_bullet), Inches(self.size_bullet))
            shape_fill = self.shape.fill
            shape_fill.solid()
            shape_fill.fore_color.rgb = RGBColor(0, 0, 0)
            shape_line = self.shape.line
            shape_line.fill.background()
            self.shape.adjustments[0] = 1.0  # initial
            self.shape.shadow.inherit = False  # disalbe dealut shadow effect

            self.shape = self.slide.shapes
            self.shape = self.shape.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(inche_to_connect_x-(self.size_bullet*0.5)), Inches(inche_to_connect_y-(self.size_bullet*0.5)), Inches(self.size_bullet), Inches(self.size_bullet))
            shape_fill = self.shape.fill
            shape_fill.solid()
            shape_fill.fore_color.rgb = RGBColor(0, 0, 0)
            shape_line = self.shape.line
            shape_line.fill.background()
            self.shape.adjustments[0] = 1.0  # initial
            self.shape.shadow.inherit = False  # disalbe dealut shadow effect

    def add_segment_line(self,shape_left,shape_top,shape_width,shape_hight,folder_left,folder_top,folder_width,folder_hight,temp_sub_col_count,temp_row_subfolder,temp_width_margin,temp_hight_margin):
        ###input line list to list type
        temp_cell_value = str(self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_sub_col_count).value)
        temp_cell_value = '["' + temp_cell_value.split('["')[-1].split('"]')[0] + '"]'
        segment_line_list = eval(temp_cell_value)

        ### get total hight for multi segment lines
        margin_segment_hight = (temp_hight_margin + shape_hight) / (len(segment_line_list)+1)

        ### set line coord
        inche_from_connect_x = folder_left
        inche_from_connect_y = shape_top - (temp_hight_margin*0.5)
        inche_to_connect_x = folder_left + folder_width
        inche_to_connect_y = shape_top - (temp_hight_margin*0.5)
        #print(segment_line_list)

        '''wirte each segment lines'''
        for temp_segment_hight in range(0, len(segment_line_list)):
            inche_from_connect_y += margin_segment_hight
            inche_to_connect_y += margin_segment_hight

            ### check left cell in excel
            if temp_sub_col_count >= 3:
                inche_from_connect_x = shape_left - (temp_width_margin * 0.5)
            ### check right cell in excel
            if self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_sub_col_count+1).value != '<END>':
                inche_to_connect_x = shape_left + shape_width + (temp_width_margin * 0.5)

            # write line
            self.shape = self.slide.shapes
            self.shape = self.shape.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(inche_from_connect_x), Inches(inche_from_connect_y), Inches(inche_to_connect_x), Inches(inche_to_connect_y))

            '''change style of line'''
            shape_line = self.shape.line
            shape_line.color.rgb = RGBColor(0, 0, 0)
            shape_line.color.brightness = 0.0
            shape_line.width = Pt(0.8)
            shape_line.fill.solid()
            self.shape.shadow.inherit = False

            '''add segment line tag'''
            temp_temp_line_row = '<TEMP_SEGMENT>'
            To_name = str(segment_line_list[temp_segment_hight]).split('<')[0]
            ns_ddx_figure_run.add_line_tag(self,'Segment', inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y, temp_temp_line_row, To_name)

            '''set coord list'''
            ### refrect tag function own shape
            if '<' in str(segment_line_list[temp_segment_hight]):
                temp_shape_text = str('<') + str(segment_line_list[temp_segment_hight]).split('<')[-1].split('>')[0] + str('>')

            x_left = inche_from_connect_x
            x_middle = inche_from_connect_x + ((inche_to_connect_x-inche_from_connect_x)*0.5)
            x_right = inche_to_connect_x
            y_top = inche_from_connect_y
            y_middle = inche_from_connect_y
            y_down = inche_from_connect_y
            self.coord_list.append([temp_shape_text, x_left, x_middle, x_right, y_top, y_middle, y_down])  # coord of x_left,x_middle,x_right,y_top,y_middle,y_down

    def get_shape_width(self,shape_name):
        ### search row of <<STYLE_SHAPE>
        temp_style_shape_row = 1
        temp_style_shape_flag = False
        shape_width = 1.0
        while temp_style_shape_row < 50000 and temp_style_shape_flag == False:
            if str(self.input_ppt_mata_excel.active.cell(temp_style_shape_row, 1).value) == '<<STYLE_SHAPE>>':
                temp_style_shape_flag = True
            temp_style_shape_row += 1
        ## get shape width from excel <<STYLE_SHAPE>>
        temp_temp_style_shape_row = temp_style_shape_row
        temp_temp_style_shape_flag = False
        temp_defalut_style_shape_flag = False

        ### if shape_name == None , change to <EMPTY>
        if shape_name == None:
            shape_name = '<EMPTY>'

        while temp_temp_style_shape_row < 50000 and temp_temp_style_shape_flag == False:
            if self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value == None:
                temp_temp_style_shape_flag = True
                temp_defalut_style_shape_flag = True
                break
            if str(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value) == str(shape_name):
                temp_temp_style_shape_flag = True
                ### Set Width
                shape_width = float(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 2).value)
            temp_temp_style_shape_row +=1

        #### if Default flag = True, select <DEFAULT>
        if temp_defalut_style_shape_flag == True:
            temp_temp_style_shape_row = temp_style_shape_row
            temp_temp_style_shape_flag = False
            while temp_temp_style_shape_row < 50000 and temp_temp_style_shape_flag == False:
                if self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value == None:
                    temp_temp_style_shape_flag = True
                    break
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value) == '<DEFAULT>':
                    temp_temp_style_shape_flag = True
                    ### Set Width
                    shape_width = float(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 2).value)
                temp_temp_style_shape_row += 1

        return (shape_width)

    def get_shape_hight(self,shape_name):
        ### search row of <<STYLE_SHAPE>
        temp_style_shape_row = 1
        temp_style_shape_flag = False
        shape_hight = 0.5

        while temp_style_shape_row < 50000 and temp_style_shape_flag == False:
            if str(self.input_ppt_mata_excel.active.cell(temp_style_shape_row, 1).value) == '<<STYLE_SHAPE>>':
                temp_style_shape_flag = True
            temp_style_shape_row += 1

        ### if shape_name == None , change to <EMPTY>
        if shape_name == None:
            shape_name = '<EMPTY>'

        ## get shape hight from excel <<STYLE_SHAPE>>
        temp_temp_style_shape_row = temp_style_shape_row
        temp_temp_style_shape_flag = False
        temp_defalut_style_shape_flag = False
        while temp_temp_style_shape_row < 50000 and temp_temp_style_shape_flag == False:
            if self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value == None:
                temp_temp_style_shape_flag = True
                temp_defalut_style_shape_flag = True
                break
            if str(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value) == str(shape_name):
                temp_temp_style_shape_flag = True
                ### Set hight
                shape_hight = float(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 3).value)
            temp_temp_style_shape_row +=1

        #### if Default flag = True, select <DEFAULT>
        if temp_defalut_style_shape_flag == True:
            temp_temp_style_shape_row = temp_style_shape_row
            temp_temp_style_shape_flag = False
            while temp_temp_style_shape_row < 50000 and temp_temp_style_shape_flag == False:
                if self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value == None:
                    temp_temp_style_shape_flag = True
                    break
                if str(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value) == '<DEFAULT>':
                    temp_temp_style_shape_flag = True
                    ### Set Hight
                    shape_hight = float(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 3).value)
                temp_temp_style_shape_row += 1

        return (shape_hight)

    def get_shape_highest(self,row_value):
        temp_col = 2
        temp_row_subfolder = row_value
        temp_subfolder_col_flag = False
        temp_temp_no_match_flag = False
        shape_highest = 0.0
        default_shape_hight = 0.0
        empty_shape_hight = 0.0
        ### get num of col on the row
        while self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_col).value != '<END>' and temp_subfolder_col_flag == False:
            if self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_col).value == '<END>':
                temp_subfolder_col_flag = True
            if temp_col > 1000:
                temp_subfolder_col_flag = True

            '''get highest hight '''
            if temp_subfolder_col_flag == False:
                self.shape_width_sum += ns_ddx_figure_run.get_shape_width(self,self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_col).value)

                ### search row of <<STYLE_SHAPE>
                temp_style_shape_row = 1
                temp_style_shape_flag = False

                while temp_style_shape_row < 50000 and temp_style_shape_flag == False:
                    if str(self.input_ppt_mata_excel.active.cell(temp_style_shape_row, 1).value) == '<<STYLE_SHAPE>>':
                        temp_style_shape_flag = True
                    temp_style_shape_row += 1

                ## get shape hight from excel <<STYLE_SHAPE>>
                temp_temp_style_shape_row = temp_style_shape_row
                temp_temp_style_shape_flag = False
                while temp_temp_style_shape_row < 50000 and temp_temp_style_shape_flag == False:
                    #keep the default value
                    if self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value == '<DEFAULT>':
                        default_shape_hight = float(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 3).value)
                    # keep the empty value
                    if self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value == '<EMPTY>':
                        empty_shape_hight = float(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 3).value)

                    if self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value == None:
                        temp_temp_no_match_flag = True
                        break
                    if str(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value) == str(self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_col).value):
                        temp_temp_style_shape_flag = True
                        ### Set highest
                        if shape_highest < float(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 3).value):
                            shape_highest = float(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 3).value)


                    ##empty process
                    if str(self.input_ppt_mata_excel.active.cell(temp_row_subfolder, temp_col).value) == '<EMPTY>':
                        temp_temp_style_shape_flag = True
                        ### Set highest
                        if shape_highest < empty_shape_hight:
                            shape_highest = empty_shape_hight

                    temp_temp_style_shape_row += 1

                ### Default set when no match
                if temp_temp_no_match_flag == True and default_shape_hight > shape_highest:  # modify for NS 2.0
                    #print(str(self.input_ppt_mata_excel.active.cell(temp_temp_style_shape_row, 1).value), shape_highest)
                    shape_highest = default_shape_hight

                    ### process of before while loop
            temp_col += 1
        return (shape_highest)

    def get_east_asian_width_count(self,text):
        count = 0
        for c in text:
            if unicodedata.east_asian_width(c) in 'FWA':
                count += 2
            else:
                count += 1
        return count

    def save_pptx(self):
        self.active_ppt.save(self.output_ppt_file)


    def add_l2_material(self):
        #print('--- self.shapes_size_array ---', self.shapes_size_array)  # from ns_ddx_figure
        import ns_l2_diagram_create , ns_def , ns_ddx_figure

        '''write materials of l2 shape'''
        for tmp_shapes_size_array in self.shapes_size_array:
            for tmp_all_device_l2_size_array in self.all_device_l2_size_array:
                if tmp_all_device_l2_size_array[0] == tmp_shapes_size_array[0]:
                    target_device_name = tmp_shapes_size_array[0]  # device_name
                    device_size_array = tmp_all_device_l2_size_array
                    action_type = 'WRITE_DEVICE_L2'  # 'RETURN_DEVICE_SIZE' - > return array[width, hight] , 'WRITE_DEVICE_L2' -> write device l2 materials
                    write_left_top_array = [tmp_shapes_size_array[1][0], tmp_shapes_size_array[1][1], device_size_array[1]]  # [left , top , [offset_left, offset_top , right , left]]

                    '''
                    :param action_type: RETURN_DEVICE_SIZE' - > return array[left, top , width, hight] , 'WRITE_DEVICE_L2' -> write device l2 materials
                    :param input_device_name: target device_name
                    :param write_left_top_array: [left , top , [offset_left, offset_top , right , left]] or [left , top , [device_size_array]]
                    :return: RETURN_DEVICE_SIZE' - > return array[left, top , width, hight]
                    '''

                    offset_left_master = 0.0
                    offset_top_master = 0.0
                    if action_type == 'WRITE_DEVICE_L2':
                        #print(write_left_top_array[0], write_left_top_array[2][0])
                        offset_left_master = write_left_top_array[0] - write_left_top_array[2][0]
                        offset_top_master = write_left_top_array[1] - write_left_top_array[2][1]

                    #self.title_only_slide_layout = self.active_ppt.slide_layouts[5]
                    #self.slide = self.active_ppt.slides.add_slide(self.title_only_slide_layout)
                    self.shape = self.slide.shapes

                    ### default parameter ###
                    self.folder_font_type = 'Calibri'
                    self.folder_font_size = 10  # Pt
                    self.shape_font_type = 'Calibri'
                    self.shae_font_size = 6  # Pt

                    self.roundness = 0.0  # curve of ROUNDED_RECTANGLE 0.0-1.0 * 100(%)
                    self.shape_width_min = 0.5  # in <<STYLE_SHAPE>> inches
                    self.shape_hight_min = 0.1  # in <<STYLE_SHAPE>> inches
                    self.per_char_inchi = 0.1  # inches of per char count in shape

                    shape_left = 0
                    shape_top = 0
                    shape_type = 'L2_SEGMENT'
                    shape_text = 'Dummy'
                    shape_interval_width_ratio = 0.75  # ratio of interval shapes(width)
                    shape_interval_hight_ratio = 0.75  # ratio of interval shapes(hight)
                    between_tag = 0.2  # distance between tags
                    l2seg_size_margin = 0.7  # inches   between l2seg and if, l2seg and vport
                    l2seg_size_margin_left_right_add = 0.4  # add inches l2seg and if, l2seg and vport. and right or left

                    offset_left_shape = offset_left_master
                    offset_top_shape = offset_top_master

                    self.shape.title.text = '[L2] ' + self.l2_folder_name
                    self.slide.shapes.title.left = Inches(0.0)
                    self.slide.shapes.title.top = Inches(0.0)
                    self.slide.shapes.title.width = Inches(14.0)
                    self.slide.shapes.title.height = Inches(1.0)

                    #if self.click_value == 'L2-3-3':
                    #    self.shape.title.text = '[L2] ' + target_device_name

                    '''
                    STEP1.1 define functions
                    '''
                    new_l2_table_array = []
                    for tmp_l2_table_array in self.l2_table_array:
                        if tmp_l2_table_array[0] != 1 and tmp_l2_table_array[0] != 2:
                            tmp_l2_table_array[1].extend(['', '', '', '', '', '', '', ''])
                            del tmp_l2_table_array[1][8:]
                            new_l2_table_array.append(tmp_l2_table_array)

                    # print('---- new_l2_table_array ----')
                    # print(new_l2_table_array)


                    new_l2_table_tuple = ns_def.convert_array_to_tuple(new_l2_table_array)
                    # print('---- new_l2_table_tuple ----')
                    # print(new_l2_table_tuple)

                    # input l2 l3 if type
                    update_l2_table_array = []

                    for tmp_tmp_new_l2_table_array in new_l2_table_array:
                        offset_excel = 2
                        tmp_new_l2_table_array = tmp_tmp_new_l2_table_array[1]
                        if tmp_new_l2_table_array[offset_excel + 3] == "":
                            if tmp_new_l2_table_array[offset_excel + 4] == "":
                                if tmp_new_l2_table_array[offset_excel + 1] == "":
                                    tmp_new_l2_table_array[offset_excel] = ''
                                else:
                                    tmp_new_l2_table_array[offset_excel] = 'Routed (L3)'
                            else:
                                if tmp_new_l2_table_array[offset_excel + 1] == "":
                                    tmp_new_l2_table_array[offset_excel] = ''
                                else:
                                    tmp_new_l2_table_array[offset_excel] = 'Switch (L2)'
                        else:
                            if tmp_new_l2_table_array[offset_excel + 1] == "":
                                tmp_new_l2_table_array[offset_excel] = ''
                            else:
                                tmp_new_l2_table_array[offset_excel] = 'Switch (L2)'

                        offset_excel = 4
                        if tmp_new_l2_table_array[offset_excel + 1] == "":
                            if tmp_new_l2_table_array[offset_excel + 1] == "":
                                tmp_new_l2_table_array[offset_excel] = ''
                            else:
                                tmp_new_l2_table_array[offset_excel] = 'Routed (L3)'
                        else:
                            if tmp_new_l2_table_array[offset_excel + 2] == "":
                                if tmp_new_l2_table_array[offset_excel - 1] == "":
                                    tmp_new_l2_table_array[offset_excel] = 'Loopback (L3)'
                                else:
                                    tmp_new_l2_table_array[offset_excel] = 'Routed (L3)'
                            else:
                                if tmp_new_l2_table_array[offset_excel - 1] == "":
                                    tmp_new_l2_table_array[offset_excel] = 'Routed (L3)'
                                else:
                                    tmp_new_l2_table_array[offset_excel] = 'Switch (L2)'

                        update_l2_table_array.append(tmp_new_l2_table_array)  # print(tmp_new_l2_table_array)

                    #print('--- update_l2_table_array ---')
                    #print(update_l2_table_array)

                    # GET L2 Segment name of each device
                    device_l2name_array = []
                    unique_l2name_array = []
                    for tmp_new_l2_table_array in new_l2_table_array:
                        if tmp_new_l2_table_array[1][6] != '':
                            tmp_l2seg = []
                            for tmp_char in tmp_new_l2_table_array[1][6].split(','):
                                tmp_char = tmp_char.replace(' ', '')  # [Temporary setting] replace ' ' to '' in l2 segment name
                                tmp_l2seg.append(tmp_char.strip())
                                for tmp_tmp_char in tmp_l2seg:
                                    if tmp_tmp_char not in unique_l2name_array:
                                        unique_l2name_array.append(tmp_tmp_char)

                            device_l2name_array.append([tmp_new_l2_table_array[1][1], tmp_l2seg])

                    unique_l2name_array.sort()

                    # print('--- device_l2name_array ---')
                    # print(device_l2name_array)
                    # print('--- unique_l2name_array ---')
                    # print(unique_l2name_array)

                    # get direction of phygical port of each device  (UP/DOWN/RIGHT/LEFT)
                    device_list_array = []
                    self.device_list_2_array = []
                    wp_list_array = []
                    shape_list_array = []
                    for tmp_new_l2_table_array in new_l2_table_array:
                        if tmp_new_l2_table_array[1][1] not in device_list_array and tmp_new_l2_table_array[1][1] not in wp_list_array:
                            if tmp_new_l2_table_array[1][0] == 'N/A':
                                wp_list_array.append(tmp_new_l2_table_array[1][1])
                            else:
                                device_list_array.append(tmp_new_l2_table_array[1][1])


                    shape_list_array = device_list_array
                    shape_list_array.extend(wp_list_array)

                    # print('--- shape_list_array ---')
                    # print(shape_list_array)
                    # print('--- device_list_array ---')
                    # print(device_list_array)
                    # print('--- wp_list_array ---')
                    # print(wp_list_array)

                    shape_if_array = []
                    for tmp_shape_list_array in shape_list_array:
                        tmp_shape_if_array = []
                        for tmp_new_l2_table_array in new_l2_table_array:
                            if tmp_shape_list_array == tmp_new_l2_table_array[1][1] and tmp_new_l2_table_array[1][3] != '':
                                tmp_shape_if_array.append(tmp_new_l2_table_array[1][3])
                        shape_if_array.append([tmp_shape_list_array, tmp_shape_if_array])

                    # print('--- shape_if_array ---')
                    # print(shape_if_array)

                    # create modify_position_shape_array for decide up / down of shape
                    modify_position_shape_array = []
                    for tmp_position_shape_array in self.position_shape_array:

                        if tmp_position_shape_array[0] != 1 and tmp_position_shape_array[1][0] != '<END>':
                            if tmp_position_shape_array[1][0] != '':
                                tmp_folder_name = tmp_position_shape_array[1][0]
                            else:
                                tmp_position_shape_array[1][0] = tmp_folder_name
                            # print(tmp_position_shape_array)
                            modify_position_shape_array.append(tmp_position_shape_array)
                    #print('--- modify_position_shape_array ---')
                    #print(modify_position_shape_array)

                    # create modify_position_shape_array for decide up / down of shape
                    modify_position_folder_array = []
                    for tmp_position_folder_array in self.position_folder_array:

                        if tmp_position_folder_array[0] != 1 and tmp_position_folder_array[1][0] != '<SET_WIDTH>':
                            tmp_position_folder_array[1][0] = ''
                            modify_position_folder_array.append(tmp_position_folder_array)
                    # print('--- modify_position_folder_array ---')
                    # print(modify_position_folder_array)

                    #### decide up/down/right/left ####
                    tmp_device_line_array = []
                    direction_if_array = []
                    for tmp_shape_if_array in shape_if_array:
                        tmp_direction_if_array = [tmp_shape_if_array[0], [], [], [], []]  # UP/DOWN/RIGHT/LEFT
                        # print('#########', tmp_shape_if_array[0], '#########')
                        for tmp_tmp_shape_if_array in tmp_shape_if_array[1]:
                            # print(tmp_tmp_shape_if_array)

                            # get direction of if
                            for tmp_position_line_tuple in self.position_line_tuple:
                                if tmp_position_line_tuple[0] != 1 and tmp_position_line_tuple[0] != 2 and (tmp_position_line_tuple[1] == 1 or tmp_position_line_tuple[1] == 2):
                                    if tmp_position_line_tuple[1] == 1:
                                        offet_column = 0
                                    elif tmp_position_line_tuple[1] == 2:
                                        offet_column = 1

                                    if self.position_line_tuple[tmp_position_line_tuple[0], tmp_position_line_tuple[1]] == tmp_shape_if_array[0]:
                                        # print(tmp_position_line_tuple, self.position_line_tuple[tmp_position_line_tuple])
                                        tmp_tag = self.position_line_tuple[tmp_position_line_tuple[0], 3 + offet_column]
                                        target = ' '
                                        idx = tmp_tag.find(target)
                                        modify_if_name = self.position_line_tuple[tmp_position_line_tuple[0], 13 + offet_column * 4] + ' ' + tmp_tag[idx + 1:]

                                        # print(modify_if_name)

                                        if tmp_tmp_shape_if_array == modify_if_name:
                                            if self.position_line_tuple[tmp_position_line_tuple[0], 5 + offet_column] == 'RIGHT':
                                                if self.position_line_tuple[tmp_position_line_tuple[0], 8 + offet_column * 2] == '':
                                                    tmp_tag_offset = 0.0
                                                else:
                                                    tmp_tag_offset = self.position_line_tuple[tmp_position_line_tuple[0], 8 + offet_column * 2]
                                                tmp_direction_if_array[3].extend([[modify_if_name, tmp_tag_offset]])  # print(modify_if_name, '  RIGHT')
                                            elif self.position_line_tuple[tmp_position_line_tuple[0], 5 + offet_column] == 'LEFT':
                                                if self.position_line_tuple[tmp_position_line_tuple[0], 8 + offet_column * 2] == '':
                                                    tmp_tag_offset = 0.0
                                                else:
                                                    tmp_tag_offset = self.position_line_tuple[tmp_position_line_tuple[0], 8 + offet_column * 2]
                                                tmp_direction_if_array[4].extend([[modify_if_name, tmp_tag_offset]])  # print(modify_if_name, '  LEFT')
                                            else:
                                                if offet_column == 0:
                                                    opposite_device_name = self.position_line_tuple[tmp_position_line_tuple[0], tmp_position_line_tuple[1] + 1]  # print(tmp_shape_if_array[0], '  ',opposite_device_name)

                                                else:
                                                    opposite_device_name = self.position_line_tuple[tmp_position_line_tuple[0], tmp_position_line_tuple[1] - 1]  # print(tmp_shape_if_array[0], '  ', opposite_device_name)

                                                ### 'TOP or DOWN'####
                                                if tmp_shape_if_array[0] in wp_list_array or opposite_device_name in wp_list_array:
                                                    # include wp case
                                                    # print(tmp_shape_if_array[0], '  ',opposite_device_name)

                                                    origin_folder_name = ''
                                                    opposite_folder_name = ''
                                                    for tmp_modify_position_shape_array in modify_position_shape_array:
                                                        update_tmp_modify_position_shape_array = tmp_modify_position_shape_array[1]

                                                        for index_31, tmp_update_tmp_modify_position_shape_array in enumerate(update_tmp_modify_position_shape_array):
                                                            if index_31 != 0:
                                                                # print(tmp_update_tmp_modify_position_shape_array)
                                                                if tmp_shape_if_array[0] == tmp_update_tmp_modify_position_shape_array:
                                                                    origin_folder_name = update_tmp_modify_position_shape_array[0]
                                                                if opposite_device_name == tmp_update_tmp_modify_position_shape_array:
                                                                    opposite_folder_name = update_tmp_modify_position_shape_array[0]

                                                    for tmp_modify_position_folder_array in modify_position_folder_array:
                                                        if origin_folder_name in tmp_modify_position_folder_array[1]:
                                                            origin_folder_num = tmp_modify_position_folder_array[0]
                                                        if opposite_folder_name in tmp_modify_position_folder_array[1]:
                                                            opposite_folder_num = tmp_modify_position_folder_array[0]

                                                    # print(origin_folder_name,origin_folder_num,'    ' , opposite_folder_name,opposite_folder_num)
                                                    if origin_folder_num > opposite_folder_num:
                                                        if self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2] == '':
                                                            tmp_tag_offset = 0.0
                                                        else:
                                                            tmp_tag_offset = self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2]

                                                        tmp_direction_if_array[1].extend([[modify_if_name, tmp_tag_offset]])  # print(modify_if_name, '  UP')
                                                    else:
                                                        if self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2] == '':
                                                            tmp_tag_offset = 0.0
                                                        else:
                                                            tmp_tag_offset = self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2]

                                                        tmp_direction_if_array[2].extend([[modify_if_name, tmp_tag_offset]])  # print(modify_if_name, '  DOWN')

                                                else:
                                                    # NOT include wp case
                                                    for tmp_modify_position_shape_array in modify_position_shape_array:
                                                        if tmp_shape_if_array[0] in tmp_modify_position_shape_array[1]:
                                                            origin_device_num = tmp_modify_position_shape_array[0]
                                                        if opposite_device_name in tmp_modify_position_shape_array[1]:
                                                            opposite_device_num = tmp_modify_position_shape_array[0]
                                                    # print(origin_device_num,opposite_device_num)

                                                    if origin_device_num > opposite_device_num:
                                                        if self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2] == '':
                                                            tmp_tag_offset = 0.0
                                                        else:
                                                            tmp_tag_offset = self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2]

                                                        tmp_direction_if_array[1].extend([[modify_if_name, tmp_tag_offset]])  # print(modify_if_name, '  UP')
                                                    else:
                                                        if self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2] == '':
                                                            tmp_tag_offset = 0.0
                                                        else:
                                                            tmp_tag_offset = self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2]

                                                        tmp_direction_if_array[2].extend([[modify_if_name, tmp_tag_offset]])  # print(modify_if_name, '  DOWN')

                        direction_if_array.append(tmp_direction_if_array)

                    #print('--- direction_if_array ---')
                    #print(direction_if_array)

                    ### sort if location
                    new_direction_if_array = []
                    for tmp_direction_if_array in direction_if_array:
                        sorted_direction_if_array = []
                        # print('tmp_direction_if_array  ', tmp_direction_if_array)
                        for i in range(0, 5):
                            # print('tmp_direction_if_array [ ] ' ,str(i), ' ', tmp_direction_if_array[i])
                            if i == 0 or tmp_direction_if_array[i] == []:
                                sorted_direction_if_array.extend([tmp_direction_if_array[i]])
                            else:
                                # print(len(tmp_direction_if_array[i]),' ',tmp_direction_if_array[i])

                                if len(tmp_direction_if_array[i]) == 1:
                                    del tmp_direction_if_array[i][0][-1]
                                    sorted_direction_if_array.extend([tmp_direction_if_array[i][0]])
                                else:
                                    # print(tmp_direction_if_array[i])
                                    sorted_data = sorted(tmp_direction_if_array[i], key=lambda x: (x[1]), reverse=False)
                                    # print(sorted_data)

                                    sorted_data_array = []
                                    for tmp_sorted_data in sorted_data:
                                        sorted_data_array.append(tmp_sorted_data[0])

                                    # print(sorted_data_array)
                                    sorted_direction_if_array.extend([sorted_data_array])

                        #print('--- sorted_direction_if_array ---')
                        #print(sorted_direction_if_array)

                        new_direction_if_array.append(sorted_direction_if_array)

                    #print('--- new_direction_if_array ---')
                    #print(new_direction_if_array)
                    self.new_direction_if_2_array = new_direction_if_array

                    '''
                    STEP1.2 locate materials in shape
                    '''
                    # get target_device_l2_array
                    target_device_l2_array = []

                    for tmp_device_l2name_array in device_l2name_array:
                        if target_device_name == tmp_device_l2name_array[0]:
                            for tmp_char in tmp_device_l2name_array[1]:
                                # print(tmp_char,tmp_device_l2name_array[1],target_device_l2_array)
                                if tmp_char not in target_device_l2_array:
                                    target_device_l2_array.extend([tmp_char])
                    target_device_l2_array.sort()

                    flag_l2_segment_empty = False
                    if target_device_l2_array == []:
                        flag_l2_segment_empty = True  # have not l2 segment on the device
                        target_device_l2_array.extend(['_DummyL2Segment_'])

                    # print('--- target_device_l2_array ---')
                    # print(target_device_name,target_device_l2_array)

                    '''write l2 segment of shape'''
                    count_l2name_array = 0
                    pre_shape_width = 0
                    pre_offset_left_shape = 0
                    l2seg_size_array = []
                    for tmp_target_device_l2_array in target_device_l2_array:
                        shape_text = tmp_target_device_l2_array
                        self.shape = self.slide.shapes
                        shape_width = self.shape_width_min
                        shape_hight = ns_def.get_description_width_hight(self.shae_font_size,shape_text)[1]

                        if ns_def.get_description_width_hight(self.shae_font_size,shape_text)[0] > self.shape_width_min:
                            shape_width = ns_def.get_description_width_hight(self.shae_font_size,shape_text)[0]
                        else:
                            shape_width = self.shape_width_min

                        if flag_l2_segment_empty == True:
                            shape_width = 0.01
                            shape_hight = 0.01

                        if count_l2name_array > 0:
                            offset_left_shape -= pre_shape_width
                            offset_left_shape += pre_shape_width * shape_interval_width_ratio
                            offset_top_shape += shape_hight * shape_interval_hight_ratio

                            if pre_offset_left_shape + pre_shape_width + (shape_width * shape_interval_width_ratio) > (offset_left_shape + shape_width):
                                offset_left_shape += ((pre_offset_left_shape + pre_shape_width + (shape_width * shape_interval_width_ratio)) - (offset_left_shape + shape_width))

                        pre_offset_left_shape = offset_left_shape
                        pre_shape_width = shape_width

                        l2_segment_only_array = []
                        if flag_l2_segment_empty == False:
                            shape_type = 'L2_SEGMENT'

                            for tmp_update_l2_table_array in update_l2_table_array:
                                if target_device_name == tmp_update_l2_table_array[1] and tmp_update_l2_table_array[3] == '' and tmp_update_l2_table_array[5] == '':
                                    tmp_l2seg = []
                                    tmp_char = tmp_update_l2_table_array[6].replace(' ', '')  # [Temporary setting] replace ' ' to '' in l2 segment name
                                    tmp_l2seg = tmp_char.split(',')

                                    if shape_text in tmp_l2seg:
                                        shape_type = 'L2_SEGMENT_GRAY'
                                        break

                            ns_ddx_figure.extended.add_shape(self, shape_type, shape_left + offset_left_shape, shape_top + offset_top_shape, shape_width, shape_hight, shape_text)

                        l2seg_size_array.append([shape_left + offset_left_shape, shape_top + offset_top_shape, shape_width, shape_hight, shape_text])

                        offset_left_shape += shape_width
                        offset_top_shape += shape_hight

                        count_l2name_array += 1

                    # get virtual port of shape
                    target_device_vport_array = []
                    target_device_vport_if_array = []
                    for tmp_new_l2_table_array in new_l2_table_array:
                        if tmp_new_l2_table_array[1][1] == target_device_name and tmp_new_l2_table_array[1][5] != '':
                            if tmp_new_l2_table_array[1][5] not in target_device_vport_array:
                                target_device_vport_array.append(tmp_new_l2_table_array[1][5])
                                target_device_vport_if_array.append([tmp_new_l2_table_array[1][5], [tmp_new_l2_table_array[1][3]]])
                            else:
                                for tmp_target_device_vport_if_array in target_device_vport_if_array:
                                    if tmp_target_device_vport_if_array[0] == tmp_new_l2_table_array[1][5]:
                                        tmp_target_device_vport_if_array[1].extend([tmp_new_l2_table_array[1][3]])

                    # target_device_vport_array = sorted(target_device_vport_array, key=lambda x: (x[0]), reverse=False)
                    # target_device_vport_if_array = sorted(target_device_vport_if_array, key=lambda x: (x[0]), reverse=False)
                    # print('--- target_device_vport_array ---')
                    # print(target_device_vport_array)
                    # print('--- target_device_vport_if_array ---')
                    # print(target_device_vport_if_array)

                    ### set temporarily device size
                    # print('--- l2seg_size_array (left, top , width, hight)---')
                    # print(l2seg_size_array)  #left, top , width, hight, text

                    device_size_array = [l2seg_size_array[0][0], l2seg_size_array[0][1], l2seg_size_array[-1][0] + l2seg_size_array[-1][2] - l2seg_size_array[0][0] ,l2seg_size_array[-1][1] + l2seg_size_array[-1][3] - l2seg_size_array[0][1]  ]  # left, top  width, hight

                    device_size_array = [device_size_array[0] - l2seg_size_margin, device_size_array[1 ] -l2seg_size_margin ,device_size_array[2] + l2seg_size_margin * 2,device_size_array[3] + l2seg_size_margin * 2]

                    #print('--- device_size_array (left, top , width, hight) at 1st ---  ')
                    #print(device_size_array)

                    ### reflect if and vpc to device size
                    current_direction_if_array = []
                    for tmp_new_direction_if_array in new_direction_if_array:
                        if tmp_new_direction_if_array[0] == target_device_name:
                            current_direction_if_array = tmp_new_direction_if_array
                            # print('--- current_direction_if_array ---  up/down/right/left', current_direction_if_array)
                            break

                    # remove duplicate Physical interface
                    sec_current_direction_if_array = []
                    for tmp_num in range (0,5):
                        if tmp_num != 0:
                            if len(current_direction_if_array[tmp_num]) != 0:
                                sec_current_direction_if_array.append(sorted(set(current_direction_if_array[tmp_num]), key=current_direction_if_array[tmp_num].index))
                            else:
                                sec_current_direction_if_array.append([])
                        else:
                            sec_current_direction_if_array.append(current_direction_if_array[0])

                    # print('--- sec_current_direction_if_array --- ' , sec_current_direction_if_array)
                    current_direction_if_array = sec_current_direction_if_array
                    # print('--- current_direction_if_array (up/down/right/left) ---  ')
                    # print(current_direction_if_array)

                    # check exist IF or Vport on up/down/right/left
                    flag_exist_if_vport_array = [False,False],[False,False],[False,False],[False,False]
                    exit_if_vport_num_array = [0, 0, 0,0]
                    exit_if_vport_num_l3_only_array = [ 0, 0, 0,0]
                    # print('### check exist IF or Vport on up/down/right/left ###')
                    if current_direction_if_array[1] != []:
                        # print('Exsit IF UP',current_direction_if_array[1])
                        flag_exist_if_vport_array[0][0] = True
                        for tmp_current_direction_if_array in current_direction_if_array[1]:
                            for tmp_target_device_vport_if_array in target_device_vport_if_array:
                                if tmp_current_direction_if_array in tmp_target_device_vport_if_array[1]:
                                    # print('Exsit Vport UP  ', tmp_current_direction_if_array, ' ',tmp_target_device_vport_if_array[0])
                                    flag_exist_if_vport_array[0][1] = True
                                    exit_if_vport_num_array[0] += 1

                                    for tmp_update_l2_table_array in update_l2_table_array:
                                        if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] == tmp_target_device_vport_if_array[0] \
                                                and tmp_update_l2_table_array[4] != 'Switch (L2)':
                                            exit_if_vport_num_l3_only_array[0] += 1

                    if current_direction_if_array[2] != []:
                        # print('Exsit IF DOWN  ',current_direction_if_array[2])
                        flag_exist_if_vport_array[1][0] = True
                        for tmp_current_direction_if_array in current_direction_if_array[2]:
                            for tmp_target_device_vport_if_array in target_device_vport_if_array:
                                if tmp_current_direction_if_array in tmp_target_device_vport_if_array[1]:
                                    # print('Exsit Vport DOWN  ', tmp_current_direction_if_array, ' ',tmp_target_device_vport_if_array)
                                    flag_exist_if_vport_array[1][1] = True
                                    exit_if_vport_num_array[1] += 1

                                    for tmp_update_l2_table_array in update_l2_table_array:
                                        if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] == tmp_target_device_vport_if_array[0] \
                                                and tmp_update_l2_table_array[4] != 'Switch (L2)':
                                            exit_if_vport_num_l3_only_array[1] += 1

                    if current_direction_if_array[3] != []:
                        # print('Exsit IF RIGHT ',current_direction_if_array[3])
                        flag_exist_if_vport_array[2][0] = True
                        for tmp_current_direction_if_array in current_direction_if_array[3]:
                            for tmp_target_device_vport_if_array in target_device_vport_if_array:
                                if tmp_current_direction_if_array in tmp_target_device_vport_if_array[1]:
                                    # print('Exsit Vport RIGHT  ', tmp_current_direction_if_array, ' ',tmp_target_device_vport_if_array[0])
                                    flag_exist_if_vport_array[2][1] = True
                                    exit_if_vport_num_array[2] += 1

                                    for tmp_update_l2_table_array in update_l2_table_array:
                                        if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] == tmp_target_device_vport_if_array[0] \
                                                and tmp_update_l2_table_array[4] != 'Switch (L2)':
                                            exit_if_vport_num_l3_only_array[2] += 1

                    if current_direction_if_array[4] != []:
                        # print('Exsit IF LEFT ',current_direction_if_array[4])
                        flag_exist_if_vport_array[3][0] = True
                        for tmp_current_direction_if_array in current_direction_if_array[4]:
                            for tmp_target_device_vport_if_array in target_device_vport_if_array:
                                if tmp_current_direction_if_array in tmp_target_device_vport_if_array[1]:
                                    # print('Exsit Vport LEFT  ', tmp_current_direction_if_array, ' ',tmp_target_device_vport_if_array)
                                    flag_exist_if_vport_array[3][1] = True
                                    exit_if_vport_num_array[3] += 1

                                    for tmp_update_l2_table_array in update_l2_table_array:
                                        if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] == tmp_target_device_vport_if_array[0] \
                                                and tmp_update_l2_table_array[4] != 'Switch (L2)':
                                            exit_if_vport_num_l3_only_array[3] += 1

                    #print('--- flag_exist_if_vport_array ,exit_if_vport_num_array ,exit_if_vport_num_l3_only_array up/down/right/left ---')
                    #print(flag_exist_if_vport_array,exit_if_vport_num_array,exit_if_vport_num_l3_only_array)

                    ''' count Virtual port that has not physical IF(include loopback)'''
                    count_other_if = 0
                    other_if_array = []
                    for tmp_target_device_vport_if_array in target_device_vport_if_array:
                        if tmp_target_device_vport_if_array[1] == ['']:
                            count_other_if += 1
                            other_if_array.append(tmp_target_device_vport_if_array[0])
                    #print('--- target_device_name, other_if_array,count_other_if ---')
                    #print(target_device_name, other_if_array,count_other_if)

                    # extend device frame distance up/down/right/left by vport exist
                    if flag_exist_if_vport_array[0][1] == True or count_other_if != 0:
                        device_size_array[1] -= l2seg_size_margin
                        device_size_array[3] += l2seg_size_margin
                    if flag_exist_if_vport_array[1][1] == True:
                        device_size_array[3] += l2seg_size_margin
                    if flag_exist_if_vport_array[2][1] == True:
                        device_size_array[2] += (l2seg_size_margin + l2seg_size_margin_left_right_add)
                    if flag_exist_if_vport_array[3][1] == True:
                        device_size_array[0] -= (l2seg_size_margin + l2seg_size_margin_left_right_add)
                        device_size_array[2] += (l2seg_size_margin + l2seg_size_margin_left_right_add)

                    # print('--- device_size_array (left, top , width, hight) at 2nd ---  ')
                    # print(device_size_array)

                    ''' extend device_size_array by left number of if and vport to downside'''
                    need_top_distance_leftside = len(current_direction_if_array[4]) * (self.shape_hight_min + between_tag) + l2seg_size_array[0][1] + l2seg_size_array[0][3]
                    need_top_distance_rightside = l2seg_size_array[-1][1] + l2seg_size_array[-1][3] + exit_if_vport_num_l3_only_array[2] * (self.shape_hight_min + between_tag) - (l2seg_size_margin * 0.75)
                    downside_keep_distance = device_size_array[1] + device_size_array[3] - l2seg_size_margin - (self.shape_hight_min * 1.5)
                    # print('--- need_top_distance_leftside,need_top_distance_rightside,downside_keep_distance --- ', need_top_distance_leftside,need_top_distance_rightside,downside_keep_distance)

                    if need_top_distance_leftside > need_top_distance_rightside:
                        if downside_keep_distance < need_top_distance_leftside:
                            device_size_array[3] += (need_top_distance_leftside - downside_keep_distance)
                    else:
                        if downside_keep_distance < need_top_distance_rightside:
                            device_size_array[3] += (need_top_distance_rightside - downside_keep_distance)

                    ''' extend device_size_array by right number of if and vport to upside'''
                    need_top_distance_rightside = len(current_direction_if_array[3]) * (self.shape_hight_min + between_tag) + l2seg_size_margin
                    if flag_exist_if_vport_array[0][1] == True:
                        need_top_distance_rightside += l2seg_size_margin

                    need_top_distance_leftside = (l2seg_size_array[-1][1] - l2seg_size_array[0][1]) + exit_if_vport_num_l3_only_array[3] * (self.shape_hight_min + between_tag) + (l2seg_size_margin * 0.75)
                    upside_keep_distance = l2seg_size_array[-1][1] - device_size_array[1]
                    # print('--- need_top_distance_rightside,need_top_distance_leftside,upside_keep_distance --- ', need_top_distance_rightside,need_top_distance_leftside,upside_keep_distance)

                    if need_top_distance_leftside > need_top_distance_rightside:
                        if upside_keep_distance < need_top_distance_leftside:
                            device_size_array[1] -= (need_top_distance_leftside - upside_keep_distance)
                            device_size_array[3] += (need_top_distance_leftside - upside_keep_distance)
                    else:
                        if upside_keep_distance < need_top_distance_rightside:
                            device_size_array[1] -= (need_top_distance_rightside - upside_keep_distance)
                            device_size_array[3] += (need_top_distance_rightside - upside_keep_distance)

                    # print('--- device_size_array (left, top , width, hight) at 3rd ---  ')
                    # print(device_size_array)

                    '''
                    write physical if and vport tag
                    '''
                    tag_size_array = []
                    '''write physical if of target shape [UP]'''
                    tmp_up_tag_distance_sum = l2seg_size_array[0][0] + (l2seg_size_array[0][2] * 0.5)
                    last_tag_left = tmp_up_tag_distance_sum
                    last_tag_width = 0.0

                    if flag_exist_if_vport_array[0][0] == True  : # upside IF TAG
                        # write up tag
                        for tmp_current_direction_if_array in current_direction_if_array[1]:
                            tag_type = ''
                            for tmp_update_l2_table_array in update_l2_table_array:
                                if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[3] == tmp_current_direction_if_array:
                                    if 'L2' in str(tmp_update_l2_table_array[2]):
                                        tag_type = 'L2_TAG'
                                        break
                                    elif 'L3' in str(tmp_update_l2_table_array[2]):
                                        tag_type = 'L3_TAG'
                                        break

                            tmp_if_array = ns_def.adjust_portname(tmp_current_direction_if_array)
                            tmp_if_name = ns_def.get_tag_name_from_full_name(target_device_name, tmp_current_direction_if_array, self.position_line_tuple)
                            tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                            tag_left = tmp_up_tag_distance_sum
                            tag_top = device_size_array[1] - ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1] * 0.5
                            tag_width = tmp_if_distance
                            tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                            tag_name = tmp_if_name

                            self.shape = self.slide.shapes
                            ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                            tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name , tmp_current_direction_if_array])

                            tmp_up_tag_distance_sum += tmp_if_distance + between_tag

                        last_tag_left = tag_left
                        last_tag_width = tag_width
                        # adjust device size width
                        l2seg_rightside = l2seg_size_array[-1][0] + l2seg_size_array[-1][2] + l2seg_size_margin
                        if_tag_leftside = tag_left + tag_width + l2seg_size_margin

                        if if_tag_leftside > l2seg_rightside:
                            device_size_array[2] += (if_tag_leftside-l2seg_rightside)

                    ''' write virtual if of target shape [UP]'''
                    used_vport_name_array = []
                    vport_with_l2seg_array = []
                    offset_vport_L3 = 0.0
                    if flag_exist_if_vport_array[0][1] == True or count_other_if != 0:
                        tag_type = ''
                        for tmp_update_l2_table_array in update_l2_table_array:
                            if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[5] not in used_vport_name_array and \
                                    (tmp_update_l2_table_array[3] in current_direction_if_array[1] or tmp_update_l2_table_array[3] == ''):
                                if 'L2' in str(tmp_update_l2_table_array[4]):
                                    tag_type = 'L2_TAG'
                                elif 'Routed (L3)' in str(tmp_update_l2_table_array[4]):
                                    tag_type = 'L3_TAG'
                                else:
                                    tag_type = 'GRAY_TAG'

                                # print(tmp_update_l2_table_array)
                                tmp_if_array = ns_def.adjust_portname(tmp_update_l2_table_array[5])
                                tmp_if_name = str(tmp_if_array[0]) + ' ' + str(tmp_if_array[2])
                                tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                                for tmp_tag_size_array in tag_size_array:
                                    if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                        if tag_type == 'L2_TAG':
                                            tag_left = tmp_tag_size_array[0] + (tmp_tag_size_array[2] - tmp_if_distance) * 0.5
                                            used_vport_name_array.append(tmp_update_l2_table_array[5])
                                            vport_with_l2seg_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_update_l2_table_array[5], tmp_update_l2_table_array[6], tmp_update_l2_table_array[7]])
                                            break

                                        if tag_type == 'L3_TAG' or tag_type == 'GRAY_TAG':
                                            tag_left = last_tag_left + last_tag_width + offset_vport_L3 + (between_tag * 0.5)
                                            used_vport_name_array.append(tmp_update_l2_table_array[5])
                                            offset_vport_L3 += tmp_if_distance + (between_tag * 0.5)  # device_size_array[0] = device_size_array[0] - tmp_if_distance - (between_tag * 0.5)  # device_size_array[2] = device_size_array[2] + tmp_if_distance + (between_tag * 0.5)

                                ### for other_if_array
                                for tmp_other_if_array in other_if_array:
                                    if tmp_other_if_array == tmp_update_l2_table_array[5]:
                                        tag_left = last_tag_left + last_tag_width + offset_vport_L3 + (between_tag * 0.5)
                                        used_vport_name_array.append(tmp_other_if_array)
                                        offset_vport_L3 += tmp_if_distance + (between_tag * 0.5)

                                tag_top = device_size_array[1] + l2seg_size_margin - ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1] * 0.5
                                tag_width = tmp_if_distance
                                tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                                tag_name = tmp_if_name

                                self.shape = self.slide.shapes
                                ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                                tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_update_l2_table_array[5]])
                                vport_with_l2seg_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_update_l2_table_array[5], tmp_update_l2_table_array[6], tmp_update_l2_table_array[7]])

                                ''' write Directory L2 Segment name under virtual port tag [UP]'''
                                if tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[6] == '' and tmp_update_l2_table_array[7] != '':
                                    tmp_l2seg = []
                                    tmp_char = tmp_update_l2_table_array[7].replace(' ', '')  # [Temporary setting] replace ' ' to '' in l2 segment name
                                    tmp_l2seg = tmp_char.split(',')

                                    offset_hight = 0.0
                                    offset_left = 0.05

                                    for tmp_tmp_l2seg in tmp_l2seg:
                                        offset_hight += tag_hight
                                        # offset_left += 0.05
                                        tag_width = ns_def.get_description_width_hight(self.shae_font_size,tmp_tmp_l2seg)[0]
                                        self.shape = self.slide.shapes
                                        ns_ddx_figure.extended.add_shape(self, 'L2SEG_TEXT', tag_left + offset_left, tag_top + offset_hight, tag_width, tag_hight, tmp_tmp_l2seg)


                        # adjust device size left + width
                        if (tag_left + tag_width + l2seg_size_margin) > (device_size_array[0] + device_size_array[2]):
                            device_size_array[2] += ((tag_left + tag_width + l2seg_size_margin) - (device_size_array[0] + device_size_array[2]))

                    '''write physical if of target shape [DOWN]'''
                    tmp_down_tag_distance_sum = l2seg_size_array[0][0]

                    if flag_exist_if_vport_array[1][0] == True: # Downside IF TAG
                        # set left start point of if tag
                        for tmp_current_direction_if_array in current_direction_if_array[2]:
                            tmp_if_name = ns_def.get_tag_name_from_full_name(target_device_name, tmp_current_direction_if_array, self.position_line_tuple)
                            tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]
                            tag_left = tmp_down_tag_distance_sum
                            tmp_down_tag_distance_sum += tmp_if_distance + between_tag

                        # offset start point
                        tag_offset = 0.0

                        if tmp_down_tag_distance_sum > l2seg_size_array[-1][0]:
                            # print(tmp_down_tag_distance_sum , l2seg_size_array[-1][0])
                            tag_offset = tmp_down_tag_distance_sum - l2seg_size_array[-1][0]
                            tmp_down_tag_distance_sum = l2seg_size_array[0][0] - tag_offset
                            device_size_array[0] -= tag_offset
                            device_size_array[2] += tag_offset
                        else:
                            tmp_down_tag_distance_sum = l2seg_size_array[0][0]

                        # write down tag
                        flag_down_tag_left = False
                        if flag_exist_if_vport_array[1][0] == True:
                            for tmp_current_direction_if_array in current_direction_if_array[2]:
                                tag_type = ''
                                for tmp_update_l2_table_array in update_l2_table_array:
                                    if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[3] == tmp_current_direction_if_array:
                                        if 'L2' in str(tmp_update_l2_table_array[2]):
                                            tag_type = 'L2_TAG'
                                            break
                                        elif 'L3' in str(tmp_update_l2_table_array[2]):
                                            tag_type = 'L3_TAG'
                                            break

                                tmp_if_array = ns_def.adjust_portname(tmp_current_direction_if_array)
                                tmp_if_name = ns_def.get_tag_name_from_full_name(target_device_name, tmp_current_direction_if_array, self.position_line_tuple)
                                tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                                tag_left = tmp_down_tag_distance_sum
                                tag_top = device_size_array[1] + device_size_array[3] - ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1] * 0.5
                                tag_width = tmp_if_distance
                                tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                                tag_name = tmp_if_name

                                self.shape = self.slide.shapes
                                ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                                tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name , tmp_current_direction_if_array])

                                if flag_down_tag_left == False:
                                    down_tag_left_edge = tag_left
                                    flag_down_tag_left = True

                                tmp_down_tag_distance_sum += tmp_if_distance + between_tag

                    ''' write virtual if of target shape [DOWN]'''
                    vport_with_l2seg_array =[]
                    offset_vport_L3 = 0.0
                    if flag_exist_if_vport_array[1][1] == True:
                        tag_type = ''
                        for tmp_update_l2_table_array in  reversed(update_l2_table_array):
                            if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[5] not in used_vport_name_array \
                                    and tmp_update_l2_table_array[3] in current_direction_if_array[2]:
                                if 'L2' in str(tmp_update_l2_table_array[4]):
                                    tag_type = 'L2_TAG'
                                elif 'Routed (L3)' in str(tmp_update_l2_table_array[4]):
                                    tag_type = 'L3_TAG'
                                else:
                                    tag_type = 'GRAY_TAG'

                                # print(tmp_update_l2_table_array)
                                tmp_if_array = ns_def.adjust_portname(tmp_update_l2_table_array[5])
                                tmp_if_name = str(tmp_if_array[0]) + ' ' + str(tmp_if_array[2])
                                tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                                for tmp_tag_size_array in tag_size_array:
                                    if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                        if tag_type == 'L2_TAG':
                                            tag_left = tmp_tag_size_array[0] + (tmp_tag_size_array[2] - tmp_if_distance) * 0.5
                                            used_vport_name_array.append(tmp_update_l2_table_array[5])
                                            vport_with_l2seg_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_update_l2_table_array[5], tmp_update_l2_table_array[6], tmp_update_l2_table_array[7]])
                                            break

                                if tag_type == 'L3_TAG' or tag_type == 'GRAY_TAG':
                                    tag_left = down_tag_left_edge - offset_vport_L3 - tmp_if_distance - (between_tag * 0.5)
                                    used_vport_name_array.append(tmp_update_l2_table_array[5])
                                    offset_vport_L3 += tmp_if_distance + (between_tag * 0.5)
                                    device_size_array[0] = device_size_array[0] - tmp_if_distance - (between_tag * 0.5)
                                    device_size_array[2] = device_size_array[2] + tmp_if_distance + (between_tag * 0.5)

                                tag_top = device_size_array[1] + device_size_array[3] - 0.1 - l2seg_size_margin
                                tag_width = tmp_if_distance
                                tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                                tag_name = tmp_if_name

                                self.shape = self.slide.shapes
                                ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                                tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_update_l2_table_array[5]])
                                vport_with_l2seg_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_update_l2_table_array[5], tmp_update_l2_table_array[6], tmp_update_l2_table_array[7]])

                                ''' write Directory L2 Segment name under virtual port tag [DOWN]'''
                                if tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[6] == '' and tmp_update_l2_table_array[7] != '':
                                    # print(tag_name,tmp_update_l2_table_array[7])
                                    tmp_l2seg = []
                                    tmp_char = tmp_update_l2_table_array[7].replace(' ', '')  # [Temporary setting] replace ' ' to '' in l2 segment name
                                    tmp_l2seg = tmp_char.split(',')

                                    offset_hight = 0.0
                                    offset_left = 0.05

                                    for tmp_tmp_l2seg in tmp_l2seg:
                                        offset_hight += tag_hight
                                        # offset_left += 0.05
                                        tag_width = ns_def.get_description_width_hight(self.shae_font_size,tmp_tmp_l2seg)[0]
                                        self.shape = self.slide.shapes
                                        ns_ddx_figure.extended.add_shape(self, 'L2SEG_TEXT', tag_left + offset_left, tag_top + offset_hight, tag_width, tag_hight, tmp_tmp_l2seg)

                    '''write physical if of target shape [RIGHT]'''
                    if flag_exist_if_vport_array[2][0] == True:
                        offset_hight = 0.0
                        for tmp_current_direction_if_array in reversed(current_direction_if_array[3]):
                            tmp_if_array = ns_def.adjust_portname(tmp_current_direction_if_array)
                            tmp_if_name = ns_def.get_tag_name_from_full_name(target_device_name, tmp_current_direction_if_array, self.position_line_tuple)
                            tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                            tag_left = device_size_array[0] + device_size_array[2] - (tmp_if_distance * 0.5)
                            tag_top = l2seg_size_array[-1][1] - (l2seg_size_array[0][3] * 2) + offset_hight
                            tag_width = tmp_if_distance
                            tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                            tag_name = tmp_if_name

                            tag_type = 'GRAY_TAG'
                            for tmp_update_l2_table_array in update_l2_table_array:
                                if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[3] == tmp_current_direction_if_array:
                                    if 'Routed (L3)' == tmp_update_l2_table_array[2]:
                                        tag_type = 'L3_TAG'
                                    elif 'Switch (L2)' == tmp_update_l2_table_array[2]:
                                        tag_type = 'L2_TAG'
                                    else:
                                        tag_type = 'GRAY_TAG'

                            self.shape = self.slide.shapes
                            ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                            tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_current_direction_if_array])

                            offset_hight -= (self.shape_hight_min + between_tag)

                    ''' write virtual if of target shape [RIGHT]'''
                    offset_hight = 0
                    if flag_exist_if_vport_array[2][1] == True:
                        for tmp_update_l2_table_array in  update_l2_table_array:
                            if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[5] not in used_vport_name_array \
                                    and tmp_update_l2_table_array[3] in current_direction_if_array[3]:
                                tmp_if_array = []
                                for tmp_target_device_vport_if_array in target_device_vport_if_array:
                                    if tmp_update_l2_table_array[3] in tmp_target_device_vport_if_array[1] and tmp_target_device_vport_if_array[0] not in used_vport_name_array:
                                        used_vport_name_array.append(tmp_target_device_vport_if_array[0])

                                        tmp_if_array = ns_def.adjust_portname(tmp_target_device_vport_if_array[0])
                                        tmp_if_name = str(tmp_if_array[0]) + ' ' + str(tmp_if_array[2])
                                        tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                                        tag_left = device_size_array[0] + device_size_array[2] - (tmp_if_distance * 0.5) - l2seg_size_margin - l2seg_size_margin_left_right_add
                                        tag_top = l2seg_size_array[-1][1] + (l2seg_size_array[-1][3] * 2 ) + offset_hight
                                        tag_width = tmp_if_distance
                                        tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                                        tag_name = tmp_if_name

                                        if 'Routed (L3)' == tmp_update_l2_table_array[4]:
                                            tag_type = 'L3_TAG'
                                            break
                                        elif 'Switch (L2)' == tmp_update_l2_table_array[4]:
                                            tag_type = 'L2_TAG'
                                            for tmp_tag_size_array in tag_size_array:
                                                if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                                    tag_top = tmp_tag_size_array[1]
                                                    break
                                            break
                                        elif 'Loopback (L3)' == tmp_update_l2_table_array[4]:
                                            tag_type = 'GRAY_TAG'
                                            break

                                self.shape = self.slide.shapes
                                ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                                tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_target_device_vport_if_array[0]])

                                if tag_type != 'L2_TAG':
                                    offset_hight += (ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1] + between_tag)

                                    ''' write Directory L2 Segment name under virtual port tag [RIGHT]'''
                                    if tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[6] == '' and tmp_update_l2_table_array[7] != '':
                                        tmp_l2seg = []
                                        tmp_char = tmp_update_l2_table_array[7].replace(' ', '')  # [Temporary setting] replace ' ' to '' in l2 segment name
                                        tmp_l2seg = tmp_char.split(',')



                                        half_num = math.floor(len(tmp_l2seg) * 0.5)
                                        flag_over_half = False
                                        half_count = 0
                                        upside_l2seg = ''
                                        downside_l2seg = ''
                                        for tmp_tmp_l2seg in tmp_l2seg:
                                            if half_num > half_count:
                                                upside_l2seg += (tmp_tmp_l2seg + ' ')
                                                half_count += 1
                                            else:
                                                flag_over_half = True
                                                downside_l2seg += (tmp_tmp_l2seg + ' ')

                                        offset_hight_l2seg = tag_hight
                                        offset_RIGHT_l2seg  = 0.05

                                        if upside_l2seg == '' and downside_l2seg != '':
                                            upside_l2seg = downside_l2seg
                                            downside_l2seg = ''

                                        self.shape = self.slide.shapes
                                        tag_width = ns_def.get_description_width_hight(self.shae_font_size,upside_l2seg)[0]
                                        ns_ddx_figure.extended.add_shape(self, 'L2SEG_TEXT', tag_left + offset_RIGHT_l2seg, tag_top + offset_hight_l2seg, tag_width, tag_hight, upside_l2seg)

                                        offset_hight_l2seg += tag_hight
                                        tag_width = ns_def.get_description_width_hight(self.shae_font_size,downside_l2seg)[0]
                                        self.shape = self.slide.shapes
                                        ns_ddx_figure.extended.add_shape(self, 'L2SEG_TEXT', tag_left + offset_RIGHT_l2seg, tag_top + offset_hight_l2seg, tag_width, tag_hight, downside_l2seg)


                    '''write physical if of target shape [LEFT]'''
                    if flag_exist_if_vport_array[3][0] == True:
                        offset_hight = 0.0
                        for tmp_current_direction_if_array in current_direction_if_array[4]:
                            tmp_if_array = ns_def.adjust_portname(tmp_current_direction_if_array)
                            tmp_if_name = ns_def.get_tag_name_from_full_name(target_device_name, tmp_current_direction_if_array, self.position_line_tuple)
                            tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                            tag_left = device_size_array[0] - (tmp_if_distance * 0.5)
                            tag_top = l2seg_size_array[0][1] + (l2seg_size_array[0][3] * 2) + offset_hight
                            tag_width = tmp_if_distance
                            tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                            tag_name = tmp_if_name

                            tag_type = 'GRAY_TAG'
                            for tmp_update_l2_table_array in update_l2_table_array:
                                if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[3] == tmp_current_direction_if_array:
                                    if 'Routed (L3)' == tmp_update_l2_table_array[2]:
                                        tag_type = 'L3_TAG'
                                    elif 'Switch (L2)' == tmp_update_l2_table_array[2]:
                                        tag_type = 'L2_TAG'
                                    else:
                                        tag_type = 'GRAY_TAG'

                            self.shape = self.slide.shapes
                            ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                            tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_current_direction_if_array])

                            offset_hight += (ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1] + between_tag)

                    ''' write virtual if of target shape [LEFT]'''
                    offset_hight = 0.0
                    if flag_exist_if_vport_array[3][1] == True:
                        for tmp_update_l2_table_array in  update_l2_table_array:
                            if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[5] not in used_vport_name_array \
                                    and tmp_update_l2_table_array[3] in current_direction_if_array[4]:
                                tmp_if_array = []
                                for tmp_target_device_vport_if_array in target_device_vport_if_array:
                                    if tmp_update_l2_table_array[3] in tmp_target_device_vport_if_array[1] and tmp_target_device_vport_if_array[0] not in used_vport_name_array:
                                        used_vport_name_array.append(tmp_target_device_vport_if_array[0])

                                        tmp_if_array = ns_def.adjust_portname(tmp_target_device_vport_if_array[0])
                                        tmp_if_name = str(tmp_if_array[0]) + ' ' + str(tmp_if_array[2])
                                        tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                                        tag_left = device_size_array[0] - (tmp_if_distance * 0.5) + l2seg_size_margin + l2seg_size_margin_left_right_add
                                        tag_top = l2seg_size_array[0][1] - ((exit_if_vport_num_l3_only_array[3]) * (between_tag + l2seg_size_array[0][3])) + offset_hight
                                        tag_width = tmp_if_distance
                                        tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                                        tag_name = tmp_if_name

                                        if 'Routed (L3)' == tmp_update_l2_table_array[4]:
                                            tag_type = 'L3_TAG'
                                            break
                                        elif 'Switch (L2)' == tmp_update_l2_table_array[4]:
                                            tag_type = 'L2_TAG'
                                            for tmp_tag_size_array in tag_size_array:
                                                if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                                    tag_top = tmp_tag_size_array[1]
                                                    break
                                            break
                                        elif 'Loopback (L3)' == tmp_update_l2_table_array[4]:
                                            tag_type = 'GRAY_TAG'
                                            break

                                self.shape = self.slide.shapes
                                ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                                tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_target_device_vport_if_array[0]])

                                if tag_type != 'L2_TAG':
                                    offset_hight += (ns_def.get_description_width_hight(self.shae_font_size,tag_name)[1] + between_tag)

                                    ''' write Directory L2 Segment name under virtual port tag [LEFT]'''
                                    if tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[6] == '' and tmp_update_l2_table_array[7] != '':
                                        tmp_l2seg = []
                                        tmp_char = tmp_update_l2_table_array[7].replace(' ', '')  # [Temporary setting] replace ' ' to '' in l2 segment name
                                        tmp_l2seg = tmp_char.split(',')



                                        half_num = math.floor(len(tmp_l2seg) * 0.5)
                                        flag_over_half = False
                                        half_count = 0
                                        upside_l2seg = ''
                                        downside_l2seg = ''
                                        for tmp_tmp_l2seg in tmp_l2seg:
                                            if half_num > half_count:
                                                upside_l2seg += (tmp_tmp_l2seg + ' ')
                                                half_count += 1
                                            else:
                                                flag_over_half = True
                                                downside_l2seg += (tmp_tmp_l2seg + ' ')

                                        offset_hight_l2seg = tag_hight
                                        offset_left_l2seg  = 0.05

                                        if upside_l2seg == '' and downside_l2seg != '':
                                            upside_l2seg = downside_l2seg
                                            downside_l2seg = ''

                                        self.shape = self.slide.shapes
                                        tag_width = ns_def.get_description_width_hight(self.shae_font_size,upside_l2seg)[0]
                                        ns_ddx_figure.extended.add_shape(self, 'L2SEG_TEXT', tag_left + offset_left_l2seg, tag_top + offset_hight_l2seg, tag_width, tag_hight, upside_l2seg)

                                        offset_hight_l2seg += ns_def.get_description_width_hight(self.shae_font_size,downside_l2seg)[1]
                                        tag_width = ns_def.get_description_width_hight(self.shae_font_size,downside_l2seg)[0]
                                        self.shape = self.slide.shapes
                                        ns_ddx_figure.extended.add_shape(self, 'L2SEG_TEXT', tag_left + offset_left_l2seg, tag_top + offset_hight_l2seg, tag_width, tag_hight, downside_l2seg)

                    '''
                    write lines
                    '''
                    used_vport_name_array = []
                    inche_from_connect_x = 0.0
                    inche_from_connect_y = 0.0
                    inche_to_connect_x = 0.0
                    inche_to_connect_y = 0.0
                    # print('--- tag_size_array 2nd  (left, top , width, hight)--- ')
                    # print(tag_size_array)
                    for tmp_update_l2_table_array in update_l2_table_array:
                        if tmp_update_l2_table_array[1] == target_device_name:
                            ### Physical IF to Virtual port
                            if tmp_update_l2_table_array[3] != '' and tmp_update_l2_table_array[5] != '':
                                # print('#LINE#  [From] ', tmp_update_l2_table_array[3], ' [To] ', tmp_update_l2_table_array[5])
                                if tmp_update_l2_table_array[3] in current_direction_if_array[1]:
                                    # UP
                                    for tmp_tag_size_array in tag_size_array:
                                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                            inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                            inche_from_connect_y = tmp_tag_size_array[1] + tmp_tag_size_array[3]
                                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                            inche_to_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                            inche_to_connect_y = tmp_tag_size_array[1]

                                elif tmp_update_l2_table_array[3] in current_direction_if_array[2]:
                                    # DOWN
                                    for tmp_tag_size_array in tag_size_array:
                                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                            inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                            inche_from_connect_y = tmp_tag_size_array[1]
                                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                            inche_to_connect_x  = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                            inche_to_connect_y = tmp_tag_size_array[1] + tmp_tag_size_array[3]

                                elif tmp_update_l2_table_array[3] in current_direction_if_array[3]:
                                    # RIHGT
                                    for tmp_tag_size_array in tag_size_array:
                                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                            inche_from_connect_x = tmp_tag_size_array[0]
                                            inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)
                                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                            inche_to_connect_x = tmp_tag_size_array[0] + tmp_tag_size_array[2]
                                            inche_to_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)

                                elif tmp_update_l2_table_array[3] in current_direction_if_array[4]:
                                    # LEFT
                                    for tmp_tag_size_array in tag_size_array:
                                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                            inche_from_connect_x = tmp_tag_size_array[0] + tmp_tag_size_array[2]
                                            inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)
                                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                            inche_to_connect_x  = tmp_tag_size_array[0]
                                            inche_to_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)

                                line_type = 'NORMAL'
                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                            ### Virtual port to L2 Segment
                            if tmp_update_l2_table_array[6] != '' and tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[5] not in used_vport_name_array:
                                # print('#LINE#  [From] ', tmp_update_l2_table_array[3], ' [To] ', tmp_update_l2_table_array[5])
                                used_vport_name_array.append(tmp_update_l2_table_array[5])

                                # UP for other's L3 Virtual port (include loopback)
                                if tmp_update_l2_table_array[5] in other_if_array:
                                    for tmp_tag_size_array in tag_size_array:
                                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                            inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                            inche_from_connect_y = tmp_tag_size_array[1] + tmp_tag_size_array[3]

                                    modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                                    for tmp_l2seg_size_array in l2seg_size_array:
                                        for tmp_modify_update_l2_table_array in modify_update_l2_table_array:
                                            if tmp_l2seg_size_array[4] == tmp_modify_update_l2_table_array:
                                                inche_to_connect_x = tmp_l2seg_size_array[0] + tmp_l2seg_size_array[2] - (tmp_l2seg_size_array[2] * 0.1)
                                                inche_to_connect_y = tmp_l2seg_size_array[1]
                                                line_type = 'NORMAL'
                                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                if tmp_update_l2_table_array[3] in current_direction_if_array[1]:
                                    # UP
                                    for tmp_tag_size_array in tag_size_array:
                                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                            inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                            inche_from_connect_y = tmp_tag_size_array[1] + tmp_tag_size_array[3]

                                    modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                                    for tmp_l2seg_size_array in l2seg_size_array:
                                        for tmp_modify_update_l2_table_array in modify_update_l2_table_array:
                                            if tmp_l2seg_size_array[4] == tmp_modify_update_l2_table_array:
                                                inche_to_connect_x = tmp_l2seg_size_array[0] + tmp_l2seg_size_array[2] - (tmp_l2seg_size_array[2] * 0.1)
                                                inche_to_connect_y = tmp_l2seg_size_array[1]
                                                line_type = 'NORMAL'
                                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                elif tmp_update_l2_table_array[3] in current_direction_if_array[2]:
                                    # DOWN
                                    for tmp_tag_size_array in tag_size_array:
                                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                            inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                            inche_from_connect_y = tmp_tag_size_array[1]

                                    modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                                    for tmp_l2seg_size_array in l2seg_size_array:
                                        for tmp_modify_update_l2_table_array in modify_update_l2_table_array:
                                            if tmp_l2seg_size_array[4] == tmp_modify_update_l2_table_array:
                                                inche_to_connect_x = tmp_l2seg_size_array[0] + (tmp_l2seg_size_array[2] * 0.1)
                                                inche_to_connect_y = tmp_l2seg_size_array[1] + tmp_l2seg_size_array[3]
                                                # used_vport_name_array.append(tmp_update_l2_table_array[5])
                                                line_type = 'NORMAL'
                                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                elif tmp_update_l2_table_array[3] in current_direction_if_array[3]:
                                    # RIHGT
                                    for tmp_tag_size_array in tag_size_array:
                                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                            inche_from_connect_x = tmp_tag_size_array[0]
                                            inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)

                                    modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                                    for tmp_l2seg_size_array in l2seg_size_array:
                                        for tmp_modify_update_l2_table_array in modify_update_l2_table_array:
                                            if tmp_l2seg_size_array[4] == tmp_modify_update_l2_table_array:
                                                inche_to_connect_x = tmp_l2seg_size_array[0] + tmp_l2seg_size_array[2]
                                                inche_to_connect_y = tmp_l2seg_size_array[1] + (tmp_l2seg_size_array[3] * 0.5)
                                                # used_vport_name_array.append(tmp_update_l2_table_array[5])
                                                line_type = 'NORMAL'
                                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)


                                elif tmp_update_l2_table_array[3] in current_direction_if_array[4]:
                                    # LEFT
                                    for tmp_tag_size_array in tag_size_array:
                                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                            inche_from_connect_x = tmp_tag_size_array[0] + tmp_tag_size_array[2]
                                            inche_from_connect_y = tmp_tag_size_array[1]+ (tmp_tag_size_array[3] * 0.5)

                                    modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                                    for tmp_l2seg_size_array in l2seg_size_array:
                                        for tmp_modify_update_l2_table_array in modify_update_l2_table_array:
                                            if tmp_l2seg_size_array[4] == tmp_modify_update_l2_table_array:
                                                inche_to_connect_x = tmp_l2seg_size_array[0]
                                                inche_to_connect_y = tmp_l2seg_size_array[1] + (tmp_l2seg_size_array[3] * 0.5)
                                                # used_vport_name_array.append(tmp_update_l2_table_array[5])
                                                line_type = 'NORMAL'
                                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                            ### Physical IF to L2 Segment
                            if tmp_update_l2_table_array[6] != '' and tmp_update_l2_table_array[5] == '' and tmp_update_l2_table_array[3] != '':
                                if tmp_update_l2_table_array[3] in current_direction_if_array[1]:
                                    # UP
                                    if flag_exist_if_vport_array[0][1] == True:
                                        for tmp_tag_size_array in tag_size_array:
                                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                                inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                                inche_from_connect_y = tmp_tag_size_array[1] + tmp_tag_size_array[3]
                                                inche_to_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                                inche_to_connect_y = device_size_array[1] + l2seg_size_margin
                                                line_type = 'NORMAL'
                                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                                inche_from_connect_x = inche_to_connect_x
                                                inche_from_connect_y = inche_to_connect_y

                                    else:
                                        for tmp_tag_size_array in tag_size_array:
                                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                                inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                                inche_from_connect_y = tmp_tag_size_array[1] + tmp_tag_size_array[3]

                                    modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                                    for tmp_l2seg_size_array in l2seg_size_array:
                                        for tmp_update_l2_table_array in modify_update_l2_table_array:
                                            if tmp_l2seg_size_array[4] == tmp_update_l2_table_array:
                                                # print('tmp_tmp_update_l2_table_array   ',tmp_tmp_update_l2_table_array )
                                                inche_to_connect_x = tmp_l2seg_size_array[0] + tmp_l2seg_size_array[2] - (tmp_l2seg_size_array[2] * 0.1)
                                                inche_to_connect_y = tmp_l2seg_size_array[1]
                                                line_type = 'NORMAL'
                                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                elif tmp_update_l2_table_array[3] in current_direction_if_array[2]:
                                    # DOWN
                                    if flag_exist_if_vport_array[1][1] == True:
                                        for tmp_tag_size_array in tag_size_array:
                                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                                inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                                inche_from_connect_y = tmp_tag_size_array[1]
                                                inche_to_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                                inche_to_connect_y = device_size_array[1] + device_size_array[3] - 0.1 - l2seg_size_margin
                                                line_type = 'NORMAL'
                                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                                inche_from_connect_x = inche_to_connect_x
                                                inche_from_connect_y = inche_to_connect_y

                                    else:
                                        for tmp_tag_size_array in tag_size_array:
                                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                                inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                                inche_from_connect_y = tmp_tag_size_array[1]

                                    modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                                    for tmp_l2seg_size_array in l2seg_size_array:
                                        for tmp_update_l2_table_array in modify_update_l2_table_array:
                                            if tmp_l2seg_size_array[4] == tmp_update_l2_table_array:
                                                # print('tmp_tmp_update_l2_table_array   ',tmp_tmp_update_l2_table_array )
                                                inche_to_connect_x = tmp_l2seg_size_array[0] + (tmp_l2seg_size_array[2] * 0.1)
                                                inche_to_connect_y = tmp_l2seg_size_array[1] + tmp_l2seg_size_array[3]
                                                line_type = 'NORMAL'
                                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                elif tmp_update_l2_table_array[3] in current_direction_if_array[3]:
                                    # RIHGT
                                    if flag_exist_if_vport_array[2][1] == True:
                                        for tmp_tag_size_array in tag_size_array:
                                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                                inche_from_connect_x = tmp_tag_size_array[0]
                                                inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)
                                                inche_to_connect_x = tmp_tag_size_array [0]  - l2seg_size_margin - l2seg_size_margin_left_right_add
                                                inche_to_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)
                                                line_type = 'NORMAL'
                                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                                inche_from_connect_x = inche_to_connect_x
                                                inche_from_connect_y = inche_to_connect_y

                                    else:
                                        for tmp_tag_size_array in tag_size_array:
                                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                                inche_from_connect_x = tmp_tag_size_array[0]
                                                inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)

                                    modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                                    for tmp_l2seg_size_array in l2seg_size_array:
                                        for tmp_update_l2_table_array in modify_update_l2_table_array:
                                            if tmp_l2seg_size_array[4] == tmp_update_l2_table_array:
                                                # print('tmp_tmp_update_l2_table_array   ',tmp_tmp_update_l2_table_array )
                                                inche_to_connect_x = tmp_l2seg_size_array[0] + tmp_l2seg_size_array[2]
                                                inche_to_connect_y = tmp_l2seg_size_array[1] + (tmp_l2seg_size_array[3] * 0.5)
                                                line_type = 'NORMAL'
                                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                elif tmp_update_l2_table_array[3] in current_direction_if_array[4]:
                                    # LEFT
                                    if flag_exist_if_vport_array[3][1] == True:
                                        for tmp_tag_size_array in tag_size_array:
                                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                                inche_from_connect_x = tmp_tag_size_array[0] + tmp_tag_size_array[2]
                                                inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)
                                                inche_to_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5) + l2seg_size_margin +l2seg_size_margin_left_right_add + 0.3
                                                inche_to_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)
                                                line_type = 'NORMAL'
                                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                                inche_from_connect_x = inche_to_connect_x
                                                inche_from_connect_y = inche_to_connect_y

                                    else:
                                        for tmp_tag_size_array in tag_size_array:
                                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                                inche_from_connect_x = tmp_tag_size_array[0] + tmp_tag_size_array[2]
                                                inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)

                                    modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                                    for tmp_l2seg_size_array in l2seg_size_array:
                                        for tmp_update_l2_table_array in modify_update_l2_table_array:
                                            if tmp_l2seg_size_array[4] == tmp_update_l2_table_array:
                                                # print('tmp_tmp_update_l2_table_array   ',tmp_tmp_update_l2_table_array )
                                                inche_to_connect_x = tmp_l2seg_size_array[0]
                                                inche_to_connect_y = tmp_l2seg_size_array[1] + (tmp_l2seg_size_array[3] * 0.5)
                                                line_type = 'NORMAL'
                                                ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)


                    '''write device frame'''
                    #self.shape = self.slide.shapes
                    #ns_ddx_figure.extended.add_shape(self, 'DEVICE_FRAME', device_size_array[0], device_size_array[1], device_size_array[2], device_size_array[3], target_device_name)
                    #self.slide.shapes._spTree.remove(self.shape._element)   # move shape to back layer
                    #self.slide.shapes._spTree.insert(2, self.shape._element)  # move shape to back layer

                    #if action_type == 'RETURN_DEVICE_SIZE':
                    #    return ([device_size_array[0], device_size_array[1],device_size_array[2],device_size_array[3]])

                    self.all_tag_size_array.append([target_device_name,tag_size_array])


    def add_l2_line(self):
        import ns_def
        #print('### self.all_tag_size_line_array [left, top , width, hight] ###')
        #print(self.all_tag_size_array)
        #print('######### self.shapes_size_array')
        #print(self.shapes_size_array)

        check_current_device_array = []
        for tmp_shapes_size_array in self.shapes_size_array:
            check_current_device_array.append(tmp_shapes_size_array[0])

        #print('### check_current_device_array ###')
        #print(check_current_device_array)

        new_all_tag_size_array = []
        for tmp_all_tag_size_array in self.all_tag_size_array:
            for tmp_shapes_size_array in self.shapes_size_array:
                if tmp_all_tag_size_array[0] == tmp_shapes_size_array[0]:
                    new_all_tag_size_array.append(tmp_all_tag_size_array)
                    break

        #print('#### new_all_tag_size_array ####')
        #print(new_all_tag_size_array)
        #print('#### self.new_direction_if_2_array up/down/right/left ####')
        #print(self.new_direction_if_2_array)

        ### get <<POSITION_LINE>> from excel
        temp_line_row = 1
        temp_line_flag = False
        while temp_line_row < 50000 and temp_line_flag == False:
            if str(self.input_ppt_mata_excel.active.cell(temp_line_row, 1).value) == '<<POSITION_LINE>>':
                temp_line_flag = True
            temp_line_row += 1

        ### read each row under <<POSITION_LINE>>
        temp_temp_line_row = temp_line_row + 1
        while temp_temp_line_row < 50000 and self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 1).value != None:
            ''' write each line From To'''
            ### from_side and to_side line data.  [device_name , tag_name , right/left]

            tmp_from_if_array = ns_def.split_portname(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 3).value)
            tmp_from_if_fullname = str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 13).value)+ ' ' + str(tmp_from_if_array [1])
            tmp_to_if_array = ns_def.split_portname(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 4).value)
            tmp_to_if_fullname = str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 17).value) + ' ' + str(tmp_to_if_array[1])

            line_from_array = [self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 1).value, tmp_from_if_fullname,str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 5).value)]
            line_to_array = [self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 2).value, tmp_to_if_fullname, str(self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 6).value)]

            for tmp_all_tag_size_array in new_all_tag_size_array:
                if line_from_array[0] == tmp_all_tag_size_array[0]:
                    for tmp_tmp_all_tag_size_array in tmp_all_tag_size_array[1]:
                        #print(line_from_array[1],tmp_tmp_all_tag_size_array[5])
                        if line_from_array[1] == tmp_tmp_all_tag_size_array[5]:
                            inche_from_connect_x = tmp_tmp_all_tag_size_array[0] + (tmp_tmp_all_tag_size_array[2] * 0.5)
                            inche_from_connect_y = tmp_tmp_all_tag_size_array[1] + (tmp_tmp_all_tag_size_array[3] * 0.5)
                            break

            for tmp_all_tag_size_array in new_all_tag_size_array:
                if line_to_array[0] == tmp_all_tag_size_array[0]:
                    for tmp_tmp_all_tag_size_array in tmp_all_tag_size_array[1]:
                        #print(line_to_array[1], tmp_tmp_all_tag_size_array[5])
                        if line_to_array[1] == tmp_tmp_all_tag_size_array[5]:
                            inche_to_connect_x = tmp_tmp_all_tag_size_array[0] + (tmp_tmp_all_tag_size_array[2] * 0.5)
                            inche_to_connect_y = tmp_tmp_all_tag_size_array[1] + (tmp_tmp_all_tag_size_array[3] * 0.5)
                            break

            #write line
            if line_to_array[0] in check_current_device_array and line_from_array[0] in check_current_device_array:
                #adjust line point x,y
                for tmp_new_direction_if_2_array in self.new_direction_if_2_array:
                    if tmp_new_direction_if_2_array[0] == line_from_array[0]:
                        if line_from_array[1] in tmp_new_direction_if_2_array[1]:
                            inche_from_connect_y -= (tmp_tmp_all_tag_size_array[3] * 0.5)
                        elif line_from_array[1] in tmp_new_direction_if_2_array[2]:
                            inche_from_connect_y += (tmp_tmp_all_tag_size_array[3] * 0.5)
                        elif line_from_array[1] in tmp_new_direction_if_2_array[3]:
                            inche_from_connect_x += (tmp_tmp_all_tag_size_array[2] * 0.5)
                        elif line_from_array[1] in tmp_new_direction_if_2_array[4]:
                            inche_from_connect_x -= (tmp_tmp_all_tag_size_array[2] * 0.5)

                    if tmp_new_direction_if_2_array[0] == line_to_array[0]:
                        if line_to_array[1] in tmp_new_direction_if_2_array[1]:
                            inche_to_connect_y -= (tmp_tmp_all_tag_size_array[3] * 0.5)
                        elif line_to_array[1] in tmp_new_direction_if_2_array[2]:
                            inche_to_connect_y += (tmp_tmp_all_tag_size_array[3] * 0.5)
                        elif line_to_array[1] in tmp_new_direction_if_2_array[3]:
                            inche_to_connect_x += (tmp_tmp_all_tag_size_array[2] * 0.5)
                        elif line_to_array[1] in tmp_new_direction_if_2_array[4]:
                            inche_to_connect_x -= (tmp_tmp_all_tag_size_array[2] * 0.5)

                self.shape = self.slide.shapes
                self.shape = self.shape.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(inche_from_connect_x), Inches(inche_from_connect_y), Inches(inche_to_connect_x), Inches(inche_to_connect_y))

                '''change style of line'''
                shape_line = self.shape.line
                shape_line.color.rgb = RGBColor(0, 0, 0)
                shape_line.color.brightness = 0.0
                shape_line.width = Pt(0.5)
                shape_line.fill.solid()
                self.shape.shadow.inherit = False

                if self.input_ppt_mata_excel.active.cell(temp_temp_line_row, 12).value == 'NO':
                    shape_line.fill.background()


            '''last process, increment row num '''
            temp_temp_line_row += 1


if __name__ == '__main__':
    ns_ddx_figure_run()


class extended():
    def __init__(self):
        print('ns_ddx_figure_extended()')

    def add_shape(self,shape_type,shape_left, shape_top, shape_width, shape_hight,shape_text):

        self.add_shape_write_array.append([shape_type,shape_left, shape_top, shape_width, shape_hight,shape_text])

        '''Default Parameter'''
        self.shape = self.shape.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(shape_left), Inches(shape_top), Inches(shape_width), Inches(shape_hight))

        # fill parameter
        shape_fill = self.shape.fill
        shape_fill.solid()
        shape_fill.fore_color.rgb = RGBColor(255, 255, 255)

        # line parameter
        shape_line = self.shape.line
        shape_line.color.rgb = RGBColor(0, 0, 0)
        shape_line.color.brightness = 0.0
        shape_line.width = Pt(1.0)
        shape_line.fill.solid()
        self.shape.shadow.inherit = False  # disalbe dealut shadow effect
        self.shape.adjustments[0] = float(0.0)  # curve of ROUNDED_RECTANGLE 0.0~1.0

        # text parameter
        self.shape.text = shape_text
        self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        self.shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        self.shape.text_frame.paragraphs[0].font.name = self.shape_font_type
        self.shape.text_frame.paragraphs[0].font.size = Pt(self.shae_font_size)
        self.shape.text_frame.margin_top = 0
        self.shape.text_frame.margin_bottom = 0
        self.shape.text_frame.margin_left = 0
        self.shape.text_frame.margin_right = 0


        if shape_type == 'L2_SEGMENT':
            #fill paramter
            shape_fill.fore_color.rgb = RGBColor(247,245,249)
            #line paramter
            shape_line.color.rgb = RGBColor(112, 48, 160)
            shape_line.width = Pt(0.5)
            self.shape.adjustments[0] = float(0.15)  # curve of ROUNDED_RECTANGLE 0.0~1.0
            # text parameter
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(112, 48, 160)

        elif shape_type == 'FOLDER_NORMAL':
            #fill paramter
            shape_fill.background()
            # line paramter
            shape_line.width = Pt(1.0)
            shape_line.color.rgb = RGBColor(205, 205, 205)
            self.shape.adjustments[0] = float(0.015)  # curve of ROUNDED_RECTANGLE 0.0~1.0
            # text parameter
            #self.shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            self.shape.text_frame.paragraphs[0].font.size = Pt(self.folder_font_size )
            self.shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        elif shape_type == 'OUTLINE_NORMAL':
            #fill paramter change color to 255,255,255 at ver 2.3.0
            #shape_fill.background()
            shape_fill.fore_color.rgb = RGBColor(255, 255, 255)
            # line paramter
            shape_line.width = Pt(1.0)
            shape_line.color.rgb = RGBColor(0, 0, 0)
            self.shape.adjustments[0] = float(0.0)  # curve of ROUNDED_RECTANGLE 0.0~1.0

        elif shape_type == 'IP_ADDRESS_TAG':
            #fill paramter
            shape_fill.background()
            # line paramter
            shape_line.fill.background()
            shape_line.width = Pt(0.75)
            self.shape.adjustments[0] = float(0.0)  # curve of ROUNDED_RECTANGLE 0.0~1.0
            # text parameter
            #self.shape.text_frame.paragraphs[0].font.size = Pt(self.tag_font_large_size)
            self.shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            #self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        elif shape_type == 'L2_SEGMENT_GRAY':
            #fill paramter
            shape_fill.fore_color.rgb = RGBColor(255, 255, 255)
            # line paramter
            shape_line.color.rgb = RGBColor(127, 127, 127)
            shape_line.width = Pt(0.5)
            self.shape.adjustments[0] = float(0.15)  # curve of ROUNDED_RECTANGLE 0.0~1.0
            # text parameter
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(127, 127, 127)

        elif shape_type == 'L3_SEGMENT_GRAY':
            #fill paramter
            shape_fill.fore_color.rgb = RGBColor(249, 249, 249)
            # line paramter
            shape_line.color.rgb = RGBColor(0, 0, 0)
            shape_line.width = Pt(0.75)
            self.shape.adjustments[0] = float(0.3)  # curve of ROUNDED_RECTANGLE 0.0~1.0
            # text parameter
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        elif shape_type == 'L3_SEGMENT_VPN':
            #fill paramter
            shape_fill.fore_color.rgb = RGBColor(251, 243, 243)
            # line paramter
            shape_line.color.rgb = RGBColor(192, 0, 0)
            shape_line.width = Pt(0.75)
            self.shape.adjustments[0] = float(0.3)  # curve of ROUNDED_RECTANGLE 0.0~1.0
            # text parameter
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        elif shape_type == 'DEVICE_FRAME':
            '''Note, not used'''
            #fill paramter
            shape_fill.fore_color.rgb = RGBColor(250, 251, 247)
            #line paramter
            shape_line.color.rgb = RGBColor(0, 0, 0)
            shape_line.width = Pt(1.0)
            self.shape.adjustments[0] = float(0.0)  # curve of ROUNDED_RECTANGLE 0.0~1.0
            # text parameter
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            self.shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            self.shape.text_frame.margin_top = Inches(0.05)
            self.shape.text_frame.margin_left = Inches(0.1)
            self.shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
            self.shape.text_frame.paragraphs[0].font.size = Pt(16.0)

        elif shape_type == 'DEVICE_NORMAL':
            #fill paramter
            shape_fill.fore_color.rgb = RGBColor(235, 241, 222)
            #line paramter
            shape_line.color.rgb = RGBColor(0, 0, 0)
            shape_line.width = Pt(1.0)
            self.shape.adjustments[0] = float(0.0)  # curve of ROUNDED_RECTANGLE 0.0~1.0
            # text parameter
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            self.shape.text_frame.paragraphs[0].font.size = Pt(self.shae_font_large_size)

        elif shape_type == 'DEVICE_L3_INSTANCE':
            #fill paramter
            shape_fill.fore_color.rgb = RGBColor(250, 251, 247)
            #line paramter
            shape_line.color.rgb = RGBColor(0, 0, 0)
            shape_line.width = Pt(0.5)
            self.shape.adjustments[0] = float(0.0)  # curve of ROUNDED_RECTANGLE 0.0~1.0
            # text parameter
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            self.shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            self.shape.text_frame.margin_left = Inches(0.1)
            #self.shape.text_frame.paragraphs[0].font.size = Pt(16.0)

        elif shape_type == 'L3_INSTANCE':
            #fill paramter
            shape_fill.fore_color.rgb = RGBColor(230, 224, 236)
            #line paramter
            shape_line.color.rgb = RGBColor(0, 0, 0)
            shape_line.width = Pt(1.0)
            self.shape.adjustments[0] = float(0.2)  # curve of ROUNDED_RECTANGLE 0.0~1.0
            # text parameter
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            #self.shape.text_frame.paragraphs[0].font.size = Pt(16.0)

        elif shape_type == 'L2_TAG':
            #fill paramter
            shape_fill.fore_color.rgb = RGBColor(255,255,255)
            #line paramter
            shape_line.color.rgb = RGBColor(0, 112, 192)
            shape_line.width = Pt(0.5)
            self.shape.adjustments[0] = float(0.50445)  # curve of ROUNDED_RECTANGLE 0.0~1.0  #0.50445 is do not change for identify if tag
            # text parameter
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)

        elif shape_type == 'L3_TAG':
            #fill paramter
            shape_fill.fore_color.rgb = RGBColor(255,255,255)
            #line paramter
            shape_line.color.rgb = RGBColor(192, 0, 0)
            shape_line.width = Pt(0.5)
            self.shape.adjustments[0] = float(0.50445)  # curve of ROUNDED_RECTANGLE 0.0~1.0  #0.50445 is do not change for identify if tag
            # text parameter
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)

        elif shape_type == 'TAG_NORMAL':
            #fill paramter
            shape_fill.fore_color.rgb = RGBColor(255,255,255)
            #line paramter
            shape_line.color.rgb = RGBColor(0,0,0)
            shape_line.width = Pt(0.5)
            self.shape.adjustments[0] = float(0.50445)  # curve of ROUNDED_RECTANGLE 0.0~1.0 #0.50445 is do not change for identify if tag
            self.shape.text_frame.paragraphs[0].font.size = Pt(self.tag_font_large_size)
            # text parameter
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)

        elif shape_type == 'GRAY_TAG':
            # fill paramter
            shape_fill.fore_color.rgb = RGBColor(255, 255, 255)
            # line paramter
            shape_line.color.rgb = RGBColor(127, 127, 127)
            shape_line.width = Pt(0.5)
            self.shape.adjustments[0] = float(0.50445)  # curve of ROUNDED_RECTANGLE 0.0~1.0 #0.50445 is do not change for identify if tag
            # text parameter
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(127, 127, 127)

        elif shape_type == 'L2SEG_TEXT':
            shape_fill.background()
            shape_line.fill.background()
            self.shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            self.shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(163, 101, 209)

        elif shape_type == 'WAY_POINT':
            '''Note, not used'''
            shape_fill.fore_color.rgb = RGBColor(237, 242, 249)
            self.shape.adjustments[0] = 0.2
            self.shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            self.shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
            self.shape.text_frame.paragraphs[0].font.size = Pt(self.shae_font_size)
            self.shape.text_frame.margin_left = Inches(0.1)

        elif shape_type == 'WAY_POINT_NORMAL':
            shape_fill.fore_color.rgb = RGBColor(220, 230, 242)
            self.shape.adjustments[0] = 0.2
            self.shape.text_frame.paragraphs[0].font.size = Pt(self.shae_font_large_size)

        else:
            print('### not defined shape_type ###')

    def add_line(self,line_type,inche_from_connect_x,inche_from_connect_y,inche_to_connect_x,inche_to_connect_y):
        from pptx.oxml import parse_xml

        self.shape = self.slide.shapes
        if line_type == 'VPN_curve':
            #self.shape = self.shape.add_connector(MSO_CONNECTOR.CURVE, Inches(inche_from_connect_x), Inches(inche_from_connect_y), Inches(inche_to_connect_x), Inches(inche_to_connect_y))
            self.shape = self.shape.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(inche_from_connect_x), Inches(inche_from_connect_y), Inches(inche_to_connect_x), Inches(inche_to_connect_y))
        else:
            self.shape = self.shape.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(inche_from_connect_x), Inches(inche_from_connect_y), Inches(inche_to_connect_x), Inches(inche_to_connect_y))

        '''change style of line'''
        #defalut , l2 material
        shape_line = self.shape.line
        shape_line.color.rgb = RGBColor(0, 112, 192)
        shape_line.color.brightness = 0.0
        shape_line.width = Pt(0.75)
        shape_line.fill.solid()
        self.shape.shadow.inherit = False

        if line_type == 'L3_SEGMENT':
            shape_line.color.rgb = RGBColor(0, 0, 0)
            line_elem = self.shape.line._get_or_add_ln()
            shape_line.width = Pt(2.5)
            line_elem.append(parse_xml("""
                    <a:headEnd type="diamond" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>
             """))
            line_elem.append(parse_xml("""
                    <a:tailEnd type="diamond" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>
             """))


        if line_type == 'L3_SEGMENT-L3IF':
            shape_line.color.rgb = RGBColor(0, 0, 0)
            line_elem = self.shape.line._get_or_add_ln()
            shape_line.width = Pt(0.7)
            line_elem.append(parse_xml("""
                    <a:tailEnd type="diamond" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>
             """))

        if line_type == 'L3_SEGMENT-VPN':
            shape_line.color.rgb = RGBColor(192, 0, 0)
            line_elem = self.shape.line._get_or_add_ln()
            shape_line.width = Pt(0.7)
            line_elem.append(parse_xml("""
                    <a:tailEnd type="diamond" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>
             """))

        if line_type == 'L3_INSTANCE':
            shape_line.color.rgb = RGBColor(96, 74, 123)
            shape_line.width = Pt(0.7)

        if line_type == 'VPN':
            self.shape.shadow.inherit = True
            shape_line.color.brightness = 0.3
            shape_line.color.rgb = RGBColor(255, 0, 0)
            shape_line.width = Pt(4.0)
            line_elem = self.shape.line._get_or_add_ln()
            line_elem.append(parse_xml("""
                    <a:headEnd type="diamond" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>
             """))
            line_elem.append(parse_xml("""
                    <a:tailEnd type="diamond" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>
             """))

        if line_type == 'VPN_curve':
            self.shape.shadow.inherit = True
            shape_line.color.brightness = 0.3
            shape_line.color.rgb = RGBColor(255, 0, 0)
            shape_line.width = Pt(4.0)
            #shape_line.dash_style = MSO_LINE.SQUARE_DOT
            line_elem = self.shape.line._get_or_add_ln()
            line_elem.append(parse_xml("""
                    <a:headEnd type="diamond" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>
             """))


    def l2_device_materials(self,action_type,input_device_name,write_left_top_array,wp_list_array):
        import ns_def
        '''
        :param action_type: RETURN_DEVICE_SIZE' - > return array[left, top , width, hight] , 'WRITE_DEVICE_L2' -> write device l2 materials
        :param input_device_name: target device_name
        :param write_left_top_array: [left , top , [offset_left, offset_top , right , left]] or [left , top , [device_size_array]]
        :return: RETURN_DEVICE_SIZE' - > return array[left, top , width, hight]
        '''
        target_device_name = input_device_name

        offset_left_master = 0.0
        offset_top_master = 0.0
        if action_type == 'WRITE_DEVICE_L2':
            #print(write_left_top_array[0],write_left_top_array[2][0])
            offset_left_master = write_left_top_array[0] - write_left_top_array[2][0]
            offset_top_master = write_left_top_array[1] - write_left_top_array[2][1]

        self.active_ppt.slide_width = Inches(self.ppt_width + self.ppt_edge_margin * 2)
        self.active_ppt.slide_height = Inches(self.ppt_hight + self.ppt_edge_margin * 2)

        self.title_only_slide_layout = self.active_ppt.slide_layouts[5]
        self.slide = self.active_ppt.slides.add_slide(self.title_only_slide_layout)
        self.slide.shapes.title.left = Inches(0.0)
        self.slide.shapes.title.top = Inches(0.0)
        self.slide.shapes.title.width = Inches(14.0)
        self.slide.shapes.title.height = Inches(1.0)
        self.shape = self.slide.shapes

        ### default parameter ###
        self.folder_font_type = 'Calibri'
        self.folder_font_size = 10  # Pt
        self.shape_font_type = 'Calibri'
        self.shae_font_size = 6  # Pt
        self.shae_font_large_size = 8  # Pt

        self.roundness = 0.0  # curve of ROUNDED_RECTANGLE 0.0-1.0 * 100(%)
        self.shape_width_min = 0.3  # in <<STYLE_SHAPE>> inches
        self.shape_hight_min = 0.1  # in <<STYLE_SHAPE>> inches
        self.per_char_inchi = 0.1  # inches of per char count in shape

        shape_left = 0
        shape_top = 0
        shape_type = 'L2_SEGMENT'
        shape_text = 'Dummy'
        shape_interval_width_ratio = 0.75  # ratio of interval shapes(width)
        shape_interval_hight_ratio = 0.75  # ratio of interval shapes(hight)
        between_tag = 0.2  # distance between tags
        l2seg_size_margin = 0.7  # inches   between l2seg and if, l2seg and vport
        l2seg_size_margin_left_right_add = 0.4  # add inches l2seg and if, l2seg and vport. and right or left

        offset_left_shape = offset_left_master
        offset_top_shape = offset_top_master

        self.shape.title.text = input_device_name

        if self.click_value == 'L2-3-3':
            self.shape.title.text = '[L2] ' + input_device_name

        '''
        STEP1.1 define functions
        '''
        new_l2_table_array = []
        for tmp_l2_table_array in self.l2_table_array:
            if tmp_l2_table_array[0] != 1 and tmp_l2_table_array[0] != 2:
                tmp_l2_table_array[1].extend(['','','','','','','',''])
                del tmp_l2_table_array[1][8:]
                new_l2_table_array.append(tmp_l2_table_array)

        #print('---- new_l2_table_array ----')
        #print(new_l2_table_array)

        new_l2_table_tuple = ns_def.convert_array_to_tuple(new_l2_table_array)
        #print('---- new_l2_table_tuple ----')
        #print(new_l2_table_tuple)

        # input l2 l3 if type
        update_l2_table_array = []

        for tmp_tmp_new_l2_table_array in new_l2_table_array:
            offset_excel = 2
            tmp_new_l2_table_array = tmp_tmp_new_l2_table_array[1]
            if tmp_new_l2_table_array[offset_excel + 3]  == "":
                if tmp_new_l2_table_array[offset_excel + 4] == "":
                    if tmp_new_l2_table_array[offset_excel + 1] == "":
                        tmp_new_l2_table_array[offset_excel] = ''
                    else:
                        tmp_new_l2_table_array[offset_excel] = 'Routed (L3)'
                else:
                    if tmp_new_l2_table_array[offset_excel + 1] == "":
                        tmp_new_l2_table_array[offset_excel] = ''
                    else:
                        tmp_new_l2_table_array[offset_excel] = 'Switch (L2)'
            else:
                if tmp_new_l2_table_array[offset_excel + 1] == "":
                    tmp_new_l2_table_array[offset_excel] = ''
                else:
                    tmp_new_l2_table_array[offset_excel] = 'Switch (L2)'

            offset_excel = 4
            if tmp_new_l2_table_array[offset_excel + 1] == "":
                if tmp_new_l2_table_array[offset_excel + 1] == "":
                    tmp_new_l2_table_array[offset_excel] = ''
                else:
                    tmp_new_l2_table_array[offset_excel] = 'Routed (L3)'
            else:
                if tmp_new_l2_table_array[offset_excel + 2] == "":
                    if tmp_new_l2_table_array[offset_excel - 1] == "":
                        tmp_new_l2_table_array[offset_excel] = 'Loopback (L3)'
                    else:
                        tmp_new_l2_table_array[offset_excel] = 'Routed (L3)'
                else:
                    if tmp_new_l2_table_array[offset_excel - 1] == "":
                        tmp_new_l2_table_array[offset_excel] = 'Routed (L3)'
                    else:
                        tmp_new_l2_table_array[offset_excel] = 'Switch (L2)'

            update_l2_table_array.append(tmp_new_l2_table_array)
            #print(tmp_new_l2_table_array)

        #print('--- update_l2_table_array ---')
        #print(update_l2_table_array)

        # GET L2 Segment name of each device
        device_l2name_array = []
        unique_l2name_array = []
        for tmp_new_l2_table_array in new_l2_table_array:
            if tmp_new_l2_table_array[1][6] != '':
                tmp_l2seg = []
                for tmp_char in tmp_new_l2_table_array[1][6].split(','):
                    tmp_char = tmp_char.replace(' ','') #[Temporary setting] replace ' ' to '' in l2 segment name
                    tmp_l2seg.append(tmp_char.strip())
                    for tmp_tmp_char in tmp_l2seg:
                        if tmp_tmp_char not in unique_l2name_array:
                            unique_l2name_array.append(tmp_tmp_char)

                device_l2name_array.append([tmp_new_l2_table_array[1][1],tmp_l2seg])

        unique_l2name_array.sort()

        #print('--- device_l2name_array ---')
        #print(device_l2name_array)
        #print('--- unique_l2name_array ---')
        #print(unique_l2name_array)

        #get direction of phygical port of each device  (UP/DOWN/RIGHT/LEFT)
        device_list_array = []
        wp_list_array = []
        shape_list_array = []
        for tmp_new_l2_table_array in new_l2_table_array:
            if tmp_new_l2_table_array[1][1] not in device_list_array and tmp_new_l2_table_array[1][1] not in wp_list_array:
                if tmp_new_l2_table_array[1][0] == 'N/A':
                    wp_list_array.append(tmp_new_l2_table_array[1][1])
                else:
                    device_list_array.append(tmp_new_l2_table_array[1][1])

        shape_list_array= device_list_array
        shape_list_array.extend(wp_list_array)

        #print('--- shape_list_array ---')
        #print(shape_list_array)
        #print('--- device_list_array ---')
        #print(device_list_array)
        #print('--- wp_list_array ---')
        #print(wp_list_array)


        shape_if_array = []
        for tmp_shape_list_array in shape_list_array:
            tmp_shape_if_array = []
            for tmp_new_l2_table_array in new_l2_table_array:
                if tmp_shape_list_array == tmp_new_l2_table_array[1][1] and tmp_new_l2_table_array[1][3] != '':
                    tmp_shape_if_array.append(tmp_new_l2_table_array[1][3])
            shape_if_array.append([tmp_shape_list_array, tmp_shape_if_array])

        #print('--- shape_if_array ---')
        #print(shape_if_array)

        #create modify_position_shape_array for decide up / down of shape
        modify_position_shape_array = []
        for tmp_position_shape_array in self.position_shape_array:

            if tmp_position_shape_array[0] != 1 and tmp_position_shape_array[1][0] != '<END>':
                if tmp_position_shape_array[1][0] != '':
                    tmp_folder_name = tmp_position_shape_array[1][0]
                else:
                    tmp_position_shape_array[1][0] = tmp_folder_name
                #print(tmp_position_shape_array)
                modify_position_shape_array.append(tmp_position_shape_array)
        #print('--- modify_position_shape_array ---')
        #print(modify_position_shape_array)

        #create modify_position_shape_array for decide up / down of shape
        modify_position_folder_array = []
        for tmp_position_folder_array in self.position_folder_array:

            if tmp_position_folder_array[0] != 1 and tmp_position_folder_array[1][0] != '<SET_WIDTH>':
                tmp_position_folder_array[1][0] = ''
                modify_position_folder_array.append(tmp_position_folder_array)
        #print('--- modify_position_folder_array ---')
        #print(modify_position_folder_array)

        #### decide up/down/right/left ####
        tmp_device_line_array = []
        direction_if_array = []
        for tmp_shape_if_array in shape_if_array:
            tmp_direction_if_array = [tmp_shape_if_array[0], [], [], [], []] # UP/DOWN/RIGHT/LEFT
            #print('#########', tmp_shape_if_array[0], '#########')
            for tmp_tmp_shape_if_array in tmp_shape_if_array[1]:
                #print(tmp_tmp_shape_if_array)

                #get direction of if
                for tmp_position_line_tuple in self.position_line_tuple:
                    if tmp_position_line_tuple[0] != 1 and tmp_position_line_tuple[0] != 2 and (tmp_position_line_tuple[1] == 1 or tmp_position_line_tuple[1] == 2):
                        if tmp_position_line_tuple[1] == 1:
                            offet_column = 0
                        elif tmp_position_line_tuple[1] == 2:
                            offet_column = 1

                        if self.position_line_tuple[tmp_position_line_tuple[0],tmp_position_line_tuple[1]] == tmp_shape_if_array[0]:
                            #print(tmp_position_line_tuple, self.position_line_tuple[tmp_position_line_tuple])
                            tmp_tag = self.position_line_tuple[tmp_position_line_tuple[0],3 + offet_column]
                            target = ' '
                            idx = tmp_tag.find(target)
                            modify_if_name = self.position_line_tuple[tmp_position_line_tuple[0],13 + offet_column * 4] + ' ' + tmp_tag[idx + 1:]

                            #print(modify_if_name)

                            if tmp_tmp_shape_if_array == modify_if_name:

                                if self.position_line_tuple[tmp_position_line_tuple[0],5 + offet_column] == 'RIGHT':
                                    if self.position_line_tuple[tmp_position_line_tuple[0],8 + offet_column * 2] == '':
                                        tmp_tag_offset = 0.0
                                    else:
                                        tmp_tag_offset = self.position_line_tuple[tmp_position_line_tuple[0],8 + offet_column * 2]
                                    tmp_direction_if_array[3].extend([[modify_if_name,tmp_tag_offset]])
                                    #print(modify_if_name, '  RIGHT')
                                elif self.position_line_tuple[tmp_position_line_tuple[0],5 + offet_column] == 'LEFT':
                                    if self.position_line_tuple[tmp_position_line_tuple[0],8 + offet_column * 2] == '':
                                        tmp_tag_offset = 0.0
                                    else:
                                        tmp_tag_offset = self.position_line_tuple[tmp_position_line_tuple[0],8 + offet_column * 2]
                                    tmp_direction_if_array[4].extend([[modify_if_name,tmp_tag_offset]])
                                    #print(modify_if_name, '  LEFT')
                                else:
                                    if offet_column == 0:
                                        opposite_device_name = self.position_line_tuple[tmp_position_line_tuple[0],tmp_position_line_tuple[1] + 1]
                                        #print(tmp_shape_if_array[0], '  ',opposite_device_name)

                                    else:
                                        opposite_device_name = self.position_line_tuple[tmp_position_line_tuple[0],tmp_position_line_tuple[1] - 1]
                                        #print(tmp_shape_if_array[0], '  ', opposite_device_name)

                                    ### 'TOP or DOWN'####
                                    if tmp_shape_if_array[0] in wp_list_array or opposite_device_name in wp_list_array:
                                        # include wp case
                                        #print(tmp_shape_if_array[0], '  ',opposite_device_name)
                                        origin_folder_name = ''
                                        opposite_folder_name = ''
                                        for tmp_modify_position_shape_array in modify_position_shape_array:
                                            update_tmp_modify_position_shape_array = tmp_modify_position_shape_array[1]

                                            for index_31,tmp_update_tmp_modify_position_shape_array in enumerate(update_tmp_modify_position_shape_array):
                                                if index_31 != 0:
                                                    #print(tmp_update_tmp_modify_position_shape_array)
                                                    if tmp_shape_if_array[0] == tmp_update_tmp_modify_position_shape_array:
                                                        origin_folder_name = update_tmp_modify_position_shape_array[0]
                                                    if opposite_device_name == tmp_update_tmp_modify_position_shape_array:
                                                        opposite_folder_name = update_tmp_modify_position_shape_array[0]

                                        for tmp_modify_position_folder_array in modify_position_folder_array:
                                            if origin_folder_name in tmp_modify_position_folder_array[1]:
                                                origin_folder_num = tmp_modify_position_folder_array[0]
                                            if opposite_folder_name in tmp_modify_position_folder_array[1]:
                                                opposite_folder_num = tmp_modify_position_folder_array[0]

                                        #print(origin_folder_name,origin_folder_num,'    ' , opposite_folder_name,opposite_folder_num)

                                        if origin_folder_num > opposite_folder_num:
                                            if self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2] == '':
                                                tmp_tag_offset = 0.0
                                            else:
                                                tmp_tag_offset = self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2]

                                            tmp_direction_if_array[1].extend([[modify_if_name, tmp_tag_offset]])
                                            #print(modify_if_name, '  UP')
                                        else:
                                            if self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2] == '':
                                                tmp_tag_offset = 0.0
                                            else:
                                                tmp_tag_offset = self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2]

                                            tmp_direction_if_array[2].extend([[modify_if_name, tmp_tag_offset]])
                                            #print(modify_if_name, '  DOWN')

                                    else:
                                        # NOT include wp case
                                        origin_device_num = 0
                                        opposite_device_num = 0
                                        for tmp_modify_position_shape_array in modify_position_shape_array:
                                            if tmp_shape_if_array[0] in tmp_modify_position_shape_array[1]:
                                                origin_device_num = tmp_modify_position_shape_array[0]
                                            if opposite_device_name in tmp_modify_position_shape_array[1]:
                                                opposite_device_num = tmp_modify_position_shape_array[0]

                                        if origin_device_num > opposite_device_num:
                                            if self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2] == '':
                                                tmp_tag_offset = 0.0
                                            else:
                                                tmp_tag_offset = self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2]

                                            tmp_direction_if_array[1].extend([[modify_if_name, tmp_tag_offset]])
                                            #print(modify_if_name, '  UP')
                                        else:
                                            if self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2] == '':
                                                tmp_tag_offset = 0.0
                                            else:
                                                tmp_tag_offset = self.position_line_tuple[tmp_position_line_tuple[0], 7 + offet_column * 2]

                                            tmp_direction_if_array[2].extend([[modify_if_name, tmp_tag_offset]])
                                            #print(modify_if_name, '  DOWN')

            direction_if_array.append(tmp_direction_if_array)

        #print('--- direction_if_array ---')
        #print(direction_if_array)


        ### sort if location
        new_direction_if_array = []
        for tmp_direction_if_array in direction_if_array:
            sorted_direction_if_array = []
            #print('tmp_direction_if_array  ', tmp_direction_if_array)
            for i in range(0,5):
                #print('tmp_direction_if_array [ ] ' ,str(i), ' ', tmp_direction_if_array[i])
                if i == 0 or tmp_direction_if_array[i] == []:
                    sorted_direction_if_array.extend([tmp_direction_if_array[i]])
                else:
                    #print(len(tmp_direction_if_array[i]),' ',tmp_direction_if_array[i])

                    if len(tmp_direction_if_array[i]) == 1:
                        del tmp_direction_if_array[i][0][-1]
                        sorted_direction_if_array.extend([tmp_direction_if_array[i][0]])
                    else:
                        #print(tmp_direction_if_array[i])
                        sorted_data = sorted(tmp_direction_if_array[i], key=lambda x: (x[1]), reverse=False)
                        #print(sorted_data)

                        sorted_data_array = []
                        for tmp_sorted_data in sorted_data:
                            sorted_data_array.append(tmp_sorted_data[0])

                        #print(sorted_data_array)
                        sorted_direction_if_array.extend([sorted_data_array])

            #print('--- sorted_direction_if_array ---')
            #print(sorted_direction_if_array)

            new_direction_if_array.append(sorted_direction_if_array)

        #print('--- new_direction_if_array ---')
        #print(new_direction_if_array)

        '''
        STEP1.2 locate materials in shape
        '''
        import ns_ddx_figure

        #get target_device_l2_array
        target_device_l2_array = []

        for tmp_device_l2name_array in device_l2name_array:
            if target_device_name  == tmp_device_l2name_array[0]:
                for tmp_char in tmp_device_l2name_array[1]:
                    #print(tmp_char,tmp_device_l2name_array[1],target_device_l2_array)
                    if tmp_char not in target_device_l2_array:
                        target_device_l2_array.extend([tmp_char])
        target_device_l2_array.sort()

        flag_l2_segment_empty = False
        if target_device_l2_array == []:
            flag_l2_segment_empty = True # have not l2 segment on the device
            target_device_l2_array.extend(['_DummyL2Segment_'])

        #print('--- target_device_l2_array ---')
        #print(target_device_name,target_device_l2_array)

        '''write l2 segment of shape'''
        count_l2name_array = 0
        pre_shape_width = 0
        pre_offset_left_shape = 0
        l2seg_size_array = []
        for tmp_target_device_l2_array in target_device_l2_array:
            shape_text = tmp_target_device_l2_array
            self.shape = self.slide.shapes
            shape_width = self.shape_width_min
            shape_hight = ns_def.get_description_width_hight(self.shae_font_size,shape_text)[1]


            if ns_def.get_description_width_hight(self.shae_font_size,shape_text)[0] > self.shape_width_min:
                shape_width = ns_def.get_description_width_hight(self.shae_font_size,shape_text)[0]
            else:
                shape_width = self.shape_width_min

            if flag_l2_segment_empty == True:
                shape_width = 0.01
                shape_hight = 0.01

            if count_l2name_array > 0 :
                offset_left_shape -= pre_shape_width
                offset_left_shape += pre_shape_width * shape_interval_width_ratio
                offset_top_shape += shape_hight * shape_interval_hight_ratio

                if pre_offset_left_shape + pre_shape_width + (shape_width * shape_interval_width_ratio) > (offset_left_shape  + shape_width):
                    offset_left_shape += ((pre_offset_left_shape + pre_shape_width + (shape_width * shape_interval_width_ratio)) - (offset_left_shape  + shape_width))

            pre_offset_left_shape = offset_left_shape
            pre_shape_width = shape_width

            l2_segment_only_array = []
            if flag_l2_segment_empty == False:
                shape_type = 'L2_SEGMENT'

                for tmp_update_l2_table_array in update_l2_table_array:
                    if target_device_name == tmp_update_l2_table_array[1] and tmp_update_l2_table_array[3] == '' and tmp_update_l2_table_array[5] == '':
                        tmp_l2seg=[]
                        tmp_char = tmp_update_l2_table_array[6].replace(' ', '')  # [Temporary setting] replace ' ' to '' in l2 segment name
                        tmp_l2seg = tmp_char.split(',')

                        if shape_text in tmp_l2seg:
                            shape_type = 'L2_SEGMENT_GRAY'
                            break

                ns_ddx_figure.extended.add_shape(self, shape_type, shape_left + offset_left_shape, shape_top + offset_top_shape, shape_width, shape_hight, shape_text)

            l2seg_size_array.append([shape_left + offset_left_shape, shape_top + offset_top_shape, shape_width, shape_hight, shape_text])

            offset_left_shape += shape_width
            offset_top_shape += shape_hight

            count_l2name_array += 1

        #get virtual port of shape
        target_device_vport_array = []
        target_device_vport_if_array = []
        for tmp_new_l2_table_array in new_l2_table_array:
            if tmp_new_l2_table_array[1][1] == target_device_name and tmp_new_l2_table_array[1][5] != '':
                if tmp_new_l2_table_array[1][5] not in target_device_vport_array:
                    target_device_vport_array.append(tmp_new_l2_table_array[1][5])
                    target_device_vport_if_array.append([tmp_new_l2_table_array[1][5],[tmp_new_l2_table_array[1][3]]])
                else:
                    for tmp_target_device_vport_if_array in target_device_vport_if_array:
                        if tmp_target_device_vport_if_array[0] == tmp_new_l2_table_array[1][5]:
                            tmp_target_device_vport_if_array[1].extend([tmp_new_l2_table_array[1][3]])

        #target_device_vport_array = sorted(target_device_vport_array, key=lambda x: (x[0]), reverse=False)
        #target_device_vport_if_array = sorted(target_device_vport_if_array, key=lambda x: (x[0]), reverse=False)
        #print('--- target_device_vport_array ---')
        #print(target_device_vport_array)
        #print('--- target_device_vport_if_array ---')
        #print(target_device_vport_if_array)

        ### set temporarily device size
        #print('--- l2seg_size_array (left, top , width, hight)---')
        #print(l2seg_size_array)  #left, top , width, hight, text

        device_size_array = [l2seg_size_array[0][0],l2seg_size_array[0][1],l2seg_size_array[-1][0] + l2seg_size_array[-1][2] - l2seg_size_array[0][0] \
            ,l2seg_size_array[-1][1] + l2seg_size_array[-1][3] - l2seg_size_array[0][1]] #left, top  width, hight

        device_size_array = [device_size_array[0] - l2seg_size_margin, device_size_array[1] -l2seg_size_margin ,device_size_array[2] + l2seg_size_margin * 2,device_size_array[3] + l2seg_size_margin * 2]

        #print('--- device_size_array (left, top , width, hight) at 1st ---  ')
        #print(device_size_array)

        ### reflect if and vpc to device size
        current_direction_if_array = []
        for tmp_new_direction_if_array in new_direction_if_array:
            if tmp_new_direction_if_array[0] == target_device_name:
                current_direction_if_array = tmp_new_direction_if_array
                #print('--- current_direction_if_array ---  up/down/right/left', current_direction_if_array)
                break

        #remove duplicate Physical interface
        sec_current_direction_if_array = []
        for tmp_num in range(0,5):
            if tmp_num != 0:
                if len(current_direction_if_array[tmp_num]) != 0:
                    sec_current_direction_if_array.append(sorted(set(current_direction_if_array[tmp_num]), key=current_direction_if_array[tmp_num].index))
                else:
                    sec_current_direction_if_array.append([])
            else:
                sec_current_direction_if_array.append(current_direction_if_array[0])

        #print('--- sec_current_direction_if_array --- ' , sec_current_direction_if_array)
        current_direction_if_array = sec_current_direction_if_array
        #print('--- current_direction_if_array (up/down/right/left) ---  ')
        #print(current_direction_if_array)

        #check exist IF or Vport on up/down/right/left
        flag_exist_if_vport_array = [False,False],[False,False],[False,False],[False,False]
        exit_if_vport_num_array = [0,0,0,0]
        exit_if_vport_num_l3_only_array = [0,0,0,0]
        #print('### check exist IF or Vport on up/down/right/left ###')
        if current_direction_if_array[1] != []:
            #print('Exsit IF UP',current_direction_if_array[1])
            flag_exist_if_vport_array[0][0] = True
            for tmp_current_direction_if_array in current_direction_if_array[1]:
                for tmp_target_device_vport_if_array in target_device_vport_if_array:
                    if tmp_current_direction_if_array in tmp_target_device_vport_if_array[1]:
                        #print('Exsit Vport UP  ', tmp_current_direction_if_array, ' ',tmp_target_device_vport_if_array[0])
                        flag_exist_if_vport_array[0][1] = True
                        exit_if_vport_num_array[0] += 1

                        for tmp_update_l2_table_array in update_l2_table_array:
                            if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] == tmp_target_device_vport_if_array[0] \
                                and tmp_update_l2_table_array[4] != 'Switch (L2)':
                                exit_if_vport_num_l3_only_array[0] += 1

        if current_direction_if_array[2] != []:
            #print('Exsit IF DOWN  ',current_direction_if_array[2])
            flag_exist_if_vport_array[1][0] = True
            for tmp_current_direction_if_array in current_direction_if_array[2]:
                for tmp_target_device_vport_if_array in target_device_vport_if_array:
                    if tmp_current_direction_if_array in tmp_target_device_vport_if_array[1]:
                        #print('Exsit Vport DOWN  ', tmp_current_direction_if_array, ' ',tmp_target_device_vport_if_array)
                        flag_exist_if_vport_array[1][1] = True
                        exit_if_vport_num_array[1] += 1

                        for tmp_update_l2_table_array in update_l2_table_array:
                            if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] == tmp_target_device_vport_if_array[0] \
                                and tmp_update_l2_table_array[4] != 'Switch (L2)':
                                exit_if_vport_num_l3_only_array[1] += 1

        if current_direction_if_array[3] != []:
            #print('Exsit IF RIGHT ',current_direction_if_array[3])
            flag_exist_if_vport_array[2][0] = True
            for tmp_current_direction_if_array in current_direction_if_array[3]:
                for tmp_target_device_vport_if_array in target_device_vport_if_array:
                    if tmp_current_direction_if_array in tmp_target_device_vport_if_array[1]:
                        #print('Exsit Vport RIGHT  ', tmp_current_direction_if_array, ' ',tmp_target_device_vport_if_array[0])
                        flag_exist_if_vport_array[2][1] = True
                        exit_if_vport_num_array[2] += 1

                        for tmp_update_l2_table_array in update_l2_table_array:
                            if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] == tmp_target_device_vport_if_array[0] \
                                and tmp_update_l2_table_array[4] != 'Switch (L2)':
                                exit_if_vport_num_l3_only_array[2] += 1

        if current_direction_if_array[4] != []:
            #print('Exsit IF LEFT ',current_direction_if_array[4])
            flag_exist_if_vport_array[3][0] = True
            for tmp_current_direction_if_array in current_direction_if_array[4]:
                for tmp_target_device_vport_if_array in target_device_vport_if_array:
                    if tmp_current_direction_if_array in tmp_target_device_vport_if_array[1]:
                        #print('Exsit Vport LEFT  ', tmp_current_direction_if_array, ' ',tmp_target_device_vport_if_array)
                        flag_exist_if_vport_array[3][1] = True
                        exit_if_vport_num_array[3] += 1

                        for tmp_update_l2_table_array in update_l2_table_array:
                            if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] == tmp_target_device_vport_if_array[0] \
                                and tmp_update_l2_table_array[4] != 'Switch (L2)':
                                exit_if_vport_num_l3_only_array[3] += 1

        #print('--- flag_exist_if_vport_array ,exit_if_vport_num_array ,exit_if_vport_num_l3_only_array up/down/right/left ---')
        #print(flag_exist_if_vport_array,exit_if_vport_num_array,exit_if_vport_num_l3_only_array)

        ''' count Virtual port that has not physical IF(include loopback)'''
        count_other_if = 0
        other_if_array = []
        for tmp_target_device_vport_if_array in target_device_vport_if_array:
            if tmp_target_device_vport_if_array[1] == ['']:
                count_other_if += 1
                other_if_array.append(tmp_target_device_vport_if_array[0])
        #print('--- target_device_name, other_if_array,count_other_if ,other_if_array---')
        #print(target_device_name, other_if_array,count_other_if,other_if_array)

        #extend device frame distance up/down/right/left by vport exist
        if flag_exist_if_vport_array[0][1] == True or count_other_if != 0:
            device_size_array[1] -= l2seg_size_margin
            device_size_array[3] += l2seg_size_margin
        if flag_exist_if_vport_array[1][1] == True:
            device_size_array[3] += l2seg_size_margin
        if flag_exist_if_vport_array[2][1] == True:
            device_size_array[2] += (l2seg_size_margin + l2seg_size_margin_left_right_add)
        if flag_exist_if_vport_array[3][1] == True:
            device_size_array[0] -= (l2seg_size_margin + l2seg_size_margin_left_right_add)
            device_size_array[2] += (l2seg_size_margin + l2seg_size_margin_left_right_add)

        #print('--- device_size_array (left, top , width, hight) at 2nd ---  ')
        #print(device_size_array)



        ''' extend device_size_array by left number of if and vport to downside'''
        need_top_distance_leftside = len(current_direction_if_array[4]) * (self.shape_hight_min + between_tag) + l2seg_size_array[0][1] + l2seg_size_array[0][3]
        need_top_distance_rightside = l2seg_size_array[-1][1] + l2seg_size_array[-1][3] + exit_if_vport_num_l3_only_array[2] * (self.shape_hight_min + between_tag) - (l2seg_size_margin * 0.75)
        downside_keep_distance = device_size_array[1] + device_size_array[3] - l2seg_size_margin - (self.shape_hight_min * 1.5)
        #print('--- need_top_distance_leftside,need_top_distance_rightside,downside_keep_distance --- ', need_top_distance_leftside,need_top_distance_rightside,downside_keep_distance)

        if need_top_distance_leftside > need_top_distance_rightside:
            if downside_keep_distance < need_top_distance_leftside:
                device_size_array[3] += (need_top_distance_leftside - downside_keep_distance)
        else:
            if downside_keep_distance < need_top_distance_rightside:
                device_size_array[3] += (need_top_distance_rightside - downside_keep_distance)

        ''' extend device_size_array by right number of if and vport to upside'''
        need_top_distance_rightside = len(current_direction_if_array[3]) * (self.shape_hight_min + between_tag) + l2seg_size_margin
        if flag_exist_if_vport_array[0][1] == True:
            need_top_distance_rightside += l2seg_size_margin

        need_top_distance_leftside = (l2seg_size_array[-1][1] - l2seg_size_array[0][1]) + exit_if_vport_num_l3_only_array[3] * (self.shape_hight_min + between_tag) + (l2seg_size_margin * 0.75)
        upside_keep_distance = l2seg_size_array[-1][1] - device_size_array[1]
        #print('--- need_top_distance_rightside,need_top_distance_leftside,upside_keep_distance --- ', need_top_distance_rightside,need_top_distance_leftside,upside_keep_distance)

        if need_top_distance_leftside > need_top_distance_rightside:
            if upside_keep_distance < need_top_distance_leftside:
                device_size_array[1] -= (need_top_distance_leftside - upside_keep_distance)
                device_size_array[3] += (need_top_distance_leftside - upside_keep_distance)
        else:
            if upside_keep_distance < need_top_distance_rightside:
                device_size_array[1] -= (need_top_distance_rightside - upside_keep_distance)
                device_size_array[3] += (need_top_distance_rightside - upside_keep_distance)

        #print('--- device_size_array (left, top , width, hight) at 3rd ---  ')
        #print(device_size_array)

        '''
        write physical if and vport tag
        '''
        tag_size_array = []
        '''write physical if of target shape [UP]'''
        tmp_up_tag_distance_sum = l2seg_size_array[0][0] + (l2seg_size_array[0][2] * 0.5)
        last_tag_left = tmp_up_tag_distance_sum
        last_tag_width = 0.0

        if flag_exist_if_vport_array[0][0] == True: # upside IF TAG
            #write up tag
            #if flag_exist_if_vport_array[1][0] == True:
                for tmp_current_direction_if_array in current_direction_if_array[1]:
                    tag_type = ''
                    for tmp_update_l2_table_array in update_l2_table_array:
                        if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[3] == tmp_current_direction_if_array:
                            if 'L2' in str(tmp_update_l2_table_array[2]):
                                tag_type = 'L2_TAG'
                                break
                            elif 'L3' in str(tmp_update_l2_table_array[2]):
                                tag_type = 'L3_TAG'
                                break

                    tmp_if_array = ns_def.adjust_portname(tmp_current_direction_if_array)
                    tmp_if_name = ns_def.get_tag_name_from_full_name(target_device_name, tmp_current_direction_if_array, self.position_line_tuple)
                    tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                    tag_left = tmp_up_tag_distance_sum
                    tag_top = device_size_array[1] - ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1] * 0.5
                    tag_width = tmp_if_distance
                    tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                    tag_name = tmp_if_name

                    self.shape = self.slide.shapes
                    ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                    tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name , tmp_current_direction_if_array])

                    tmp_up_tag_distance_sum += tmp_if_distance + between_tag

                last_tag_left = tag_left
                last_tag_width = tag_width
                # adjust device size width
                l2seg_rightside = l2seg_size_array[-1][0] + l2seg_size_array[-1][2] + l2seg_size_margin
                if_tag_leftside = tag_left + tag_width + l2seg_size_margin

                if if_tag_leftside > l2seg_rightside:
                    device_size_array[2] += (if_tag_leftside-l2seg_rightside)

        ''' write virtual if of target shape [UP]'''
        used_vport_name_array = []
        vport_with_l2seg_array = []
        offset_vport_L3 = 0.0

        if flag_exist_if_vport_array[0][1] == True or count_other_if != 0:
            tag_type = ''
            for tmp_update_l2_table_array in update_l2_table_array:
                if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[5] not in used_vport_name_array \
                        and (tmp_update_l2_table_array[3] in current_direction_if_array[1] or tmp_update_l2_table_array[3] == ''):
                    if 'L2' in str(tmp_update_l2_table_array[4]):
                        tag_type = 'L2_TAG'
                    elif 'Routed (L3)' in str(tmp_update_l2_table_array[4]):
                        tag_type = 'L3_TAG'
                    else:
                        tag_type = 'GRAY_TAG'

                    # print(tmp_update_l2_table_array)
                    tmp_if_array = ns_def.adjust_portname(tmp_update_l2_table_array[5])
                    tmp_if_name = str(tmp_if_array[0]) + ' ' + str(tmp_if_array[2])
                    tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                    for tmp_tag_size_array in tag_size_array:
                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                            if tag_type == 'L2_TAG':
                                tag_left = tmp_tag_size_array[0] + (tmp_tag_size_array[2] - tmp_if_distance) * 0.5
                                used_vport_name_array.append(tmp_update_l2_table_array[5])
                                vport_with_l2seg_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_update_l2_table_array[5], tmp_update_l2_table_array[6], tmp_update_l2_table_array[7]])
                                break

                            if tag_type == 'L3_TAG' or tag_type == 'GRAY_TAG':
                                tag_left = last_tag_left + last_tag_width + offset_vport_L3 + (between_tag * 0.5)
                                used_vport_name_array.append(tmp_update_l2_table_array[5])
                                offset_vport_L3 += tmp_if_distance + (between_tag * 0.5)

                    ### for other_if_array
                    for tmp_other_if_array in other_if_array:
                        if tmp_other_if_array == tmp_update_l2_table_array[5]:
                            tag_left = last_tag_left + last_tag_width + offset_vport_L3 + (between_tag * 0.5)
                            used_vport_name_array.append(tmp_other_if_array)
                            offset_vport_L3 += tmp_if_distance + (between_tag * 0.5)


                    tag_top = device_size_array[1] + l2seg_size_margin - ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1] * 0.5
                    tag_width = tmp_if_distance
                    tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                    tag_name = tmp_if_name

                    self.shape = self.slide.shapes
                    #print(tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)

                    ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                    tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_update_l2_table_array[5]])
                    vport_with_l2seg_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_update_l2_table_array[5], tmp_update_l2_table_array[6], tmp_update_l2_table_array[7]])

                    ''' write Directory L2 Segment name under virtual port tag [UP]'''
                    if tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[6] == '' and tmp_update_l2_table_array[7] != '':
                        tmp_l2seg = []
                        tmp_char = tmp_update_l2_table_array[7].replace(' ', '')  # [Temporary setting] replace ' ' to '' in l2 segment name
                        tmp_l2seg = tmp_char.split(',')

                        offset_hight = 0.0
                        offset_left = 0.05

                        for tmp_tmp_l2seg in tmp_l2seg:
                            offset_hight += tag_hight
                            # offset_left += 0.05
                            tag_width = ns_def.get_description_width_hight(self.shae_font_size,tmp_tmp_l2seg)[0]
                            self.shape = self.slide.shapes
                            ns_ddx_figure.extended.add_shape(self, 'L2SEG_TEXT', tag_left + offset_left, tag_top + offset_hight, tag_width, tag_hight, tmp_tmp_l2seg)

            # adjust device size left + width
            if (tag_left + tag_width + l2seg_size_margin) > (device_size_array[0] + device_size_array[2]):
                device_size_array[2] += ((tag_left + tag_width + l2seg_size_margin) - (device_size_array[0] + device_size_array[2]))

        '''write physical if of target shape [DOWN]'''
        tmp_down_tag_distance_sum = l2seg_size_array[0][0]

        if flag_exist_if_vport_array[1][0] == True: # Downside IF TAG
            #set left start point of if tag
            for tmp_current_direction_if_array in current_direction_if_array[2]:
                tmp_if_name = ns_def.get_tag_name_from_full_name(target_device_name, tmp_current_direction_if_array, self.position_line_tuple)
                tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]
                tag_left = tmp_down_tag_distance_sum
                tmp_down_tag_distance_sum += tmp_if_distance + between_tag

            #offset start point
            tag_offset = 0.0

            if tmp_down_tag_distance_sum > l2seg_size_array[-1][0]:
                #print(tmp_down_tag_distance_sum , l2seg_size_array[-1][0])
                tag_offset = tmp_down_tag_distance_sum - l2seg_size_array[-1][0]
                tmp_down_tag_distance_sum = l2seg_size_array[0][0] - tag_offset
                device_size_array[0] -= tag_offset
                device_size_array[2] += tag_offset
            else:
                tmp_down_tag_distance_sum = l2seg_size_array[0][0]

            #write down tag
            flag_down_tag_left = False
            if flag_exist_if_vport_array[1][0] == True:
                for tmp_current_direction_if_array in current_direction_if_array[2]:
                    tag_type = ''
                    for tmp_update_l2_table_array in update_l2_table_array:
                        if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[3] == tmp_current_direction_if_array:
                            if 'L2' in str(tmp_update_l2_table_array[2]):
                                tag_type = 'L2_TAG'
                                break
                            elif 'L3' in str(tmp_update_l2_table_array[2]):
                                tag_type = 'L3_TAG'
                                break

                    tmp_if_array = ns_def.adjust_portname(tmp_current_direction_if_array)
                    tmp_if_name = ns_def.get_tag_name_from_full_name(target_device_name, tmp_current_direction_if_array, self.position_line_tuple)
                    tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                    tag_left = tmp_down_tag_distance_sum
                    tag_top = device_size_array[1] + device_size_array[3] - ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1] * 0.5
                    tag_width = tmp_if_distance
                    tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                    tag_name = tmp_if_name

                    self.shape = self.slide.shapes
                    ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                    tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name , tmp_current_direction_if_array])

                    if flag_down_tag_left == False:
                        down_tag_left_edge = tag_left
                        flag_down_tag_left = True

                    tmp_down_tag_distance_sum += tmp_if_distance + between_tag

        ''' write virtual if of target shape [DOWN]'''
        vport_with_l2seg_array =[]
        offset_vport_L3 = 0.0
        if flag_exist_if_vport_array[1][1] == True:
            tag_type = ''
            for tmp_update_l2_table_array in  reversed(update_l2_table_array):
                if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[5] not in used_vport_name_array \
                        and tmp_update_l2_table_array[3] in current_direction_if_array[2]:
                    if 'L2' in str(tmp_update_l2_table_array[4]):
                        tag_type = 'L2_TAG'
                    elif 'Routed (L3)' in str(tmp_update_l2_table_array[4]):
                        tag_type = 'L3_TAG'
                    else:
                        tag_type = 'GRAY_TAG'

                    #print(tmp_update_l2_table_array)
                    tmp_if_array = ns_def.adjust_portname(tmp_update_l2_table_array[5])
                    tmp_if_name = str(tmp_if_array[0]) + ' ' + str(tmp_if_array[2])
                    tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                    for tmp_tag_size_array in tag_size_array:
                        if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                            if tag_type == 'L2_TAG':
                                tag_left = tmp_tag_size_array[0] + (tmp_tag_size_array[2] - tmp_if_distance) * 0.5
                                used_vport_name_array.append(tmp_update_l2_table_array[5])
                                vport_with_l2seg_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_update_l2_table_array[5], tmp_update_l2_table_array[6], tmp_update_l2_table_array[7]])
                                break

                    if tag_type == 'L3_TAG' or tag_type == 'GRAY_TAG':
                        tag_left = down_tag_left_edge - offset_vport_L3 - tmp_if_distance - (between_tag * 0.5)
                        used_vport_name_array.append(tmp_update_l2_table_array[5])
                        offset_vport_L3 += tmp_if_distance + (between_tag * 0.5)
                        device_size_array[0] = device_size_array[0] - tmp_if_distance - (between_tag * 0.5)
                        device_size_array[2] = device_size_array[2] + tmp_if_distance + (between_tag * 0.5)

                    tag_top = device_size_array[1] + device_size_array[3] - (ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1] * 0.5) - l2seg_size_margin
                    tag_width = tmp_if_distance
                    tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                    tag_name = tmp_if_name

                    self.shape = self.slide.shapes
                    ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                    tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_update_l2_table_array[5]])
                    vport_with_l2seg_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_update_l2_table_array[5], tmp_update_l2_table_array[6], tmp_update_l2_table_array[7]])

                    ''' write Directory L2 Segment name under virtual port tag [DOWN]'''
                    if tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[6] == '' and tmp_update_l2_table_array[7] != '':
                        #print(tag_name,tmp_update_l2_table_array[7])
                        tmp_l2seg = []
                        tmp_char = tmp_update_l2_table_array[7].replace(' ', '')  # [Temporary setting] replace ' ' to '' in l2 segment name
                        tmp_l2seg = tmp_char.split(',')

                        offset_hight = 0.0
                        offset_left = 0.05

                        for tmp_tmp_l2seg in tmp_l2seg:
                            offset_hight += tag_hight
                            #offset_left += 0.05
                            tag_width = ns_def.get_description_width_hight(self.shae_font_size,tmp_tmp_l2seg)[0]
                            self.shape = self.slide.shapes
                            ns_ddx_figure.extended.add_shape(self, 'L2SEG_TEXT', tag_left + offset_left, tag_top + offset_hight, tag_width, tag_hight, tmp_tmp_l2seg)

        '''write physical if of target shape [RIGHT]'''
        if flag_exist_if_vport_array[2][0] == True:
            offset_hight = 0.0
            for tmp_current_direction_if_array in reversed(current_direction_if_array[3]):
                tmp_if_array = ns_def.adjust_portname(tmp_current_direction_if_array)
                tmp_if_name = ns_def.get_tag_name_from_full_name(target_device_name, tmp_current_direction_if_array, self.position_line_tuple)
                tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                tag_left = device_size_array[0] + device_size_array[2] - (tmp_if_distance * 0.5)
                tag_top = l2seg_size_array[-1][1] - (l2seg_size_array[0][3] * 2) + offset_hight
                tag_width = tmp_if_distance
                tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                tag_name = tmp_if_name

                tag_type = 'GRAY_TAG'
                for tmp_update_l2_table_array in update_l2_table_array:
                    if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[3] == tmp_current_direction_if_array:
                        if 'Routed (L3)' == tmp_update_l2_table_array[2]:
                            tag_type = 'L3_TAG'
                        elif 'Switch (L2)' == tmp_update_l2_table_array[2]:
                            tag_type = 'L2_TAG'
                        else:
                            tag_type = 'GRAY_TAG'

                self.shape = self.slide.shapes
                ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_current_direction_if_array])

                offset_hight -= (self.shape_hight_min + between_tag)

        ''' write virtual if of target shape [RIGHT]'''
        offset_hight = 0
        if flag_exist_if_vport_array[2][1] == True:
            for tmp_update_l2_table_array in  update_l2_table_array:
                if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[5] not in used_vport_name_array \
                        and tmp_update_l2_table_array[3] in current_direction_if_array[3]:
                    tmp_if_array = []
                    for tmp_target_device_vport_if_array in target_device_vport_if_array:
                        if tmp_update_l2_table_array[3] in tmp_target_device_vport_if_array[1] and tmp_target_device_vport_if_array[0] not in used_vport_name_array:
                            used_vport_name_array.append(tmp_target_device_vport_if_array[0])

                            tmp_if_array = ns_def.adjust_portname(tmp_target_device_vport_if_array[0])
                            tmp_if_name = str(tmp_if_array[0]) + ' ' + str(tmp_if_array[2])
                            tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                            tag_left = device_size_array[0] + device_size_array[2] - (tmp_if_distance * 0.5) - l2seg_size_margin - l2seg_size_margin_left_right_add
                            tag_top = l2seg_size_array[-1][1] + (l2seg_size_array[-1][3] * 2 ) + offset_hight
                            tag_width = tmp_if_distance
                            tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                            tag_name = tmp_if_name

                            if 'Routed (L3)' == tmp_update_l2_table_array[4]:
                                tag_type = 'L3_TAG'
                                break
                            elif 'Switch (L2)' == tmp_update_l2_table_array[4]:
                                tag_type = 'L2_TAG'
                                for tmp_tag_size_array in tag_size_array:
                                    if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                        tag_top = tmp_tag_size_array[1]
                                        break
                                break
                            elif 'Loopback (L3)' == tmp_update_l2_table_array[4]:
                                tag_type = 'GRAY_TAG'
                                break

                    self.shape = self.slide.shapes
                    ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                    tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_target_device_vport_if_array[0]])

                    if tag_type != 'L2_TAG':
                        offset_hight += (ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1] + between_tag)

                        ''' write Directory L2 Segment name under virtual port tag [RIGHT]'''
                        if tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[6] == '' and tmp_update_l2_table_array[7] != '':
                            tmp_l2seg = []
                            tmp_char = tmp_update_l2_table_array[7].replace(' ', '')  # [Temporary setting] replace ' ' to '' in l2 segment name
                            tmp_l2seg = tmp_char.split(',')



                            half_num = math.floor(len(tmp_l2seg) * 0.5)
                            flag_over_half = False
                            half_count = 0
                            upside_l2seg = ''
                            downside_l2seg = ''
                            for tmp_tmp_l2seg in tmp_l2seg:
                                if half_num > half_count:
                                    upside_l2seg += (tmp_tmp_l2seg + ' ')
                                    half_count += 1
                                else:
                                    flag_over_half = True
                                    downside_l2seg += (tmp_tmp_l2seg + ' ')

                            offset_hight_l2seg = tag_hight
                            offset_RIGHT_l2seg  = 0.05

                            if upside_l2seg == '' and downside_l2seg != '':
                                upside_l2seg = downside_l2seg
                                downside_l2seg = ''

                            self.shape = self.slide.shapes
                            tag_width = ns_def.get_description_width_hight(self.shae_font_size,upside_l2seg)[0]
                            ns_ddx_figure.extended.add_shape(self, 'L2SEG_TEXT', tag_left + offset_RIGHT_l2seg, tag_top + offset_hight_l2seg, tag_width, tag_hight, upside_l2seg)

                            offset_hight_l2seg += ns_def.get_description_width_hight(self.shae_font_size,downside_l2seg)[1]
                            tag_width = ns_def.get_description_width_hight(self.shae_font_size,downside_l2seg)[0]

                            self.shape = self.slide.shapes
                            ns_ddx_figure.extended.add_shape(self, 'L2SEG_TEXT', tag_left + offset_RIGHT_l2seg, tag_top + offset_hight_l2seg, tag_width, tag_hight, downside_l2seg)


        '''write physical if of target shape [LEFT]'''
        if flag_exist_if_vport_array[3][0] == True:
            offset_hight = 0.0
            for tmp_current_direction_if_array in current_direction_if_array[4]:
                tmp_if_array = ns_def.adjust_portname(tmp_current_direction_if_array)
                tmp_if_name = ns_def.get_tag_name_from_full_name(target_device_name, tmp_current_direction_if_array, self.position_line_tuple)
                tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                tag_left = device_size_array[0] - (tmp_if_distance * 0.5)
                tag_top = l2seg_size_array[0][1] + (l2seg_size_array[0][3] * 2) + offset_hight
                tag_width = tmp_if_distance
                tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                tag_name = tmp_if_name

                tag_type = 'GRAY_TAG'
                for tmp_update_l2_table_array in update_l2_table_array:
                    if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[3] == tmp_current_direction_if_array:
                        if 'Routed (L3)' == tmp_update_l2_table_array[2]:
                            tag_type = 'L3_TAG'
                        elif 'Switch (L2)' == tmp_update_l2_table_array[2]:
                            tag_type = 'L2_TAG'
                        else:
                            tag_type = 'GRAY_TAG'

                self.shape = self.slide.shapes
                ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_current_direction_if_array])

                offset_hight += (ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1] + between_tag)

        ''' write virtual if of target shape [LEFT]'''
        offset_hight = 0.0
        if flag_exist_if_vport_array[3][1] == True:
            for tmp_update_l2_table_array in  update_l2_table_array:
                if tmp_update_l2_table_array[1] == target_device_name and tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[5] not in used_vport_name_array \
                        and tmp_update_l2_table_array[3] in current_direction_if_array[4]:
                    tmp_if_array = []
                    for tmp_target_device_vport_if_array in target_device_vport_if_array:
                        if tmp_update_l2_table_array[3] in tmp_target_device_vport_if_array[1] and tmp_target_device_vport_if_array[0] not in used_vport_name_array:
                            used_vport_name_array.append(tmp_target_device_vport_if_array[0])

                            tmp_if_array = ns_def.adjust_portname(tmp_target_device_vport_if_array[0])
                            tmp_if_name = str(tmp_if_array[0]) + ' ' + str(tmp_if_array[2])
                            tmp_if_distance = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[0]

                            tag_left = device_size_array[0] - (tmp_if_distance * 0.5) + l2seg_size_margin + l2seg_size_margin_left_right_add
                            tag_top = l2seg_size_array[0][1] - ((exit_if_vport_num_l3_only_array[3]) * (between_tag + l2seg_size_array[0][3])) + offset_hight
                            tag_width = tmp_if_distance
                            tag_hight = ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1]
                            tag_name = tmp_if_name

                            if 'Routed (L3)' == tmp_update_l2_table_array[4]:
                                tag_type = 'L3_TAG'
                                break
                            elif 'Switch (L2)' == tmp_update_l2_table_array[4]:
                                tag_type = 'L2_TAG'
                                for tmp_tag_size_array in tag_size_array:
                                    if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                        tag_top = tmp_tag_size_array[1]
                                        break
                                break
                            elif 'Loopback (L3)' == tmp_update_l2_table_array[4]:
                                tag_type = 'GRAY_TAG'
                                break

                    self.shape = self.slide.shapes
                    ns_ddx_figure.extended.add_shape(self, tag_type, tag_left, tag_top, tag_width, tag_hight, tag_name)
                    tag_size_array.append([tag_left, tag_top, tag_width, tag_hight, tag_name, tmp_target_device_vport_if_array[0]])

                    if tag_type != 'L2_TAG':
                        offset_hight += (ns_def.get_description_width_hight(self.shae_font_size,tmp_if_name)[1] + between_tag)

                        ''' write Directory L2 Segment name under virtual port tag [LEFT]'''
                        if tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[6] == '' and tmp_update_l2_table_array[7] != '':
                            tmp_l2seg = []
                            tmp_char = tmp_update_l2_table_array[7].replace(' ', '')  # [Temporary setting] replace ' ' to '' in l2 segment name
                            tmp_l2seg = tmp_char.split(',')

                            half_num = math.floor(len(tmp_l2seg) * 0.5)
                            flag_over_half = False
                            half_count = 0
                            upside_l2seg = ''
                            downside_l2seg = ''
                            for tmp_tmp_l2seg in tmp_l2seg:
                                if half_num > half_count:
                                    upside_l2seg += (tmp_tmp_l2seg + ' ')
                                    half_count += 1
                                else:
                                    flag_over_half = True
                                    downside_l2seg += (tmp_tmp_l2seg + ' ')

                            offset_hight_l2seg = tag_hight
                            offset_left_l2seg  = 0.05

                            if upside_l2seg == '' and downside_l2seg != '':
                                upside_l2seg = downside_l2seg
                                downside_l2seg = ''

                            self.shape = self.slide.shapes
                            tag_width = ns_def.get_description_width_hight(self.shae_font_size,upside_l2seg)[0]
                            ns_ddx_figure.extended.add_shape(self, 'L2SEG_TEXT', tag_left + offset_left_l2seg, tag_top + offset_hight_l2seg, tag_width, tag_hight, upside_l2seg)

                            offset_hight_l2seg += ns_def.get_description_width_hight(self.shae_font_size,downside_l2seg)[1]
                            tag_width = ns_def.get_description_width_hight(self.shae_font_size,downside_l2seg)[0]
                            self.shape = self.slide.shapes
                            ns_ddx_figure.extended.add_shape(self, 'L2SEG_TEXT', tag_left + offset_left_l2seg, tag_top + offset_hight_l2seg, tag_width, tag_hight, downside_l2seg)



        '''
        write lines
        '''
        used_vport_name_array = []
        inche_from_connect_x = 0.0
        inche_from_connect_y = 0.0
        inche_to_connect_x = 0.0
        inche_to_connect_y = 0.0
        #print('--- tag_size_array 2nd  (left, top , width, hight)--- ')
        #print(tag_size_array)
        for tmp_update_l2_table_array in update_l2_table_array:
            if tmp_update_l2_table_array[1] == target_device_name:
                ### Physical IF to Virtual port
                if tmp_update_l2_table_array[3] != '' and tmp_update_l2_table_array[5] != '':
                    #print('#LINE#  [From] ', tmp_update_l2_table_array[3], ' [To] ', tmp_update_l2_table_array[5])
                    if tmp_update_l2_table_array[3] in current_direction_if_array[1]:
                        #UP
                        for tmp_tag_size_array in tag_size_array:
                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                inche_from_connect_y = tmp_tag_size_array[1] + tmp_tag_size_array[3]
                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                inche_to_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                inche_to_connect_y = tmp_tag_size_array[1]

                    elif tmp_update_l2_table_array[3] in current_direction_if_array[2]:
                        #DOWN
                        for tmp_tag_size_array in tag_size_array:
                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                inche_from_connect_y = tmp_tag_size_array[1]
                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                inche_to_connect_x  = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                inche_to_connect_y = tmp_tag_size_array[1] + tmp_tag_size_array[3]

                    elif tmp_update_l2_table_array[3] in current_direction_if_array[3]:
                        #RIHGT
                        for tmp_tag_size_array in tag_size_array:
                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                inche_from_connect_x = tmp_tag_size_array[0]
                                inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)
                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                inche_to_connect_x = tmp_tag_size_array[0] + tmp_tag_size_array[2]
                                inche_to_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)

                    elif tmp_update_l2_table_array[3] in current_direction_if_array[4]:
                        #LEFT
                        for tmp_tag_size_array in tag_size_array:
                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                inche_from_connect_x = tmp_tag_size_array[0] + tmp_tag_size_array[2]
                                inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)
                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                inche_to_connect_x  = tmp_tag_size_array[0]
                                inche_to_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)

                    line_type = 'NORMAL'
                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                ### Virtual port to L2 Segment
                if tmp_update_l2_table_array[6] != '' and tmp_update_l2_table_array[5] != '' and tmp_update_l2_table_array[5] not in used_vport_name_array:
                    # print('#LINE#  [From] ', tmp_update_l2_table_array[3], ' [To] ', tmp_update_l2_table_array[5])
                    used_vport_name_array.append(tmp_update_l2_table_array[5])

                    # UP for other's L3 Virtual port (include loopback)
                    if tmp_update_l2_table_array[5] in other_if_array:
                        for tmp_tag_size_array in tag_size_array:
                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                inche_from_connect_y = tmp_tag_size_array[1] + tmp_tag_size_array[3]

                        modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                        for tmp_l2seg_size_array in l2seg_size_array:
                            for tmp_modify_update_l2_table_array in modify_update_l2_table_array:
                                if tmp_l2seg_size_array[4] == tmp_modify_update_l2_table_array:
                                    inche_to_connect_x = tmp_l2seg_size_array[0] + tmp_l2seg_size_array[2] - (tmp_l2seg_size_array[2] * 0.1)
                                    inche_to_connect_y = tmp_l2seg_size_array[1]
                                    line_type = 'NORMAL'
                                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                    if tmp_update_l2_table_array[3] in current_direction_if_array[1]:
                        # UP
                        for tmp_tag_size_array in tag_size_array:
                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                inche_from_connect_y = tmp_tag_size_array[1] + tmp_tag_size_array[3]

                        modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                        for tmp_l2seg_size_array in l2seg_size_array:
                            for tmp_modify_update_l2_table_array in modify_update_l2_table_array:
                                if tmp_l2seg_size_array[4] == tmp_modify_update_l2_table_array:
                                    inche_to_connect_x = tmp_l2seg_size_array[0] + tmp_l2seg_size_array[2] - (tmp_l2seg_size_array[2] * 0.1)
                                    inche_to_connect_y = tmp_l2seg_size_array[1]
                                    line_type = 'NORMAL'
                                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                    elif tmp_update_l2_table_array[3] in current_direction_if_array[2]:
                        #DOWN
                        for tmp_tag_size_array in tag_size_array:
                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                inche_from_connect_y = tmp_tag_size_array[1]

                        modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                        for tmp_l2seg_size_array in l2seg_size_array:
                            for tmp_modify_update_l2_table_array in modify_update_l2_table_array:
                                if tmp_l2seg_size_array[4] == tmp_modify_update_l2_table_array:
                                    inche_to_connect_x = tmp_l2seg_size_array[0] + (tmp_l2seg_size_array[2] * 0.1)
                                    inche_to_connect_y = tmp_l2seg_size_array[1] + tmp_l2seg_size_array[3]
                                    #used_vport_name_array.append(tmp_update_l2_table_array[5])
                                    line_type = 'NORMAL'
                                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                    elif tmp_update_l2_table_array[3] in current_direction_if_array[3]:
                        # RIHGT
                        for tmp_tag_size_array in tag_size_array:
                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                inche_from_connect_x = tmp_tag_size_array[0]
                                inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)

                        modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                        for tmp_l2seg_size_array in l2seg_size_array:
                            for tmp_modify_update_l2_table_array in modify_update_l2_table_array:
                                if tmp_l2seg_size_array[4] == tmp_modify_update_l2_table_array:
                                    inche_to_connect_x = tmp_l2seg_size_array[0] + tmp_l2seg_size_array[2]
                                    inche_to_connect_y = tmp_l2seg_size_array[1] + (tmp_l2seg_size_array[3] * 0.5)
                                    #used_vport_name_array.append(tmp_update_l2_table_array[5])
                                    line_type = 'NORMAL'
                                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)


                    elif tmp_update_l2_table_array[3] in current_direction_if_array[4]:
                        # LEFT
                        for tmp_tag_size_array in tag_size_array:
                            if tmp_tag_size_array[5] == tmp_update_l2_table_array[5]:
                                inche_from_connect_x = tmp_tag_size_array[0] + tmp_tag_size_array[2]
                                inche_from_connect_y = tmp_tag_size_array[1]+ (tmp_tag_size_array[3] * 0.5)

                        modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                        for tmp_l2seg_size_array in l2seg_size_array:
                            for tmp_modify_update_l2_table_array in modify_update_l2_table_array:
                                if tmp_l2seg_size_array[4] == tmp_modify_update_l2_table_array:
                                    inche_to_connect_x = tmp_l2seg_size_array[0]
                                    inche_to_connect_y = tmp_l2seg_size_array[1] + (tmp_l2seg_size_array[3] * 0.5)
                                    #used_vport_name_array.append(tmp_update_l2_table_array[5])
                                    line_type = 'NORMAL'
                                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                ### Physical IF to L2 Segment
                if tmp_update_l2_table_array[6] != '' and tmp_update_l2_table_array[5] == '' and tmp_update_l2_table_array[3] != '':
                    if tmp_update_l2_table_array[3] in current_direction_if_array[1]:
                        # UP
                        if flag_exist_if_vport_array[0][1] == True:
                            for tmp_tag_size_array in tag_size_array:
                                if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                    inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                    inche_from_connect_y = tmp_tag_size_array[1] + tmp_tag_size_array[3]
                                    inche_to_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                    inche_to_connect_y = device_size_array[1] + l2seg_size_margin
                                    line_type = 'NORMAL'
                                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                    inche_from_connect_x = inche_to_connect_x
                                    inche_from_connect_y = inche_to_connect_y

                        else:
                            for tmp_tag_size_array in tag_size_array:
                                if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                    inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                    inche_from_connect_y = tmp_tag_size_array[1] + tmp_tag_size_array[3]

                        modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                        for tmp_l2seg_size_array in l2seg_size_array:
                            for tmp_update_l2_table_array in modify_update_l2_table_array:
                                if tmp_l2seg_size_array[4] == tmp_update_l2_table_array:
                                    # print('tmp_tmp_update_l2_table_array   ',tmp_tmp_update_l2_table_array )
                                    inche_to_connect_x = tmp_l2seg_size_array[0] + tmp_l2seg_size_array[2] - (tmp_l2seg_size_array[2] * 0.1)
                                    inche_to_connect_y = tmp_l2seg_size_array[1]
                                    line_type = 'NORMAL'
                                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                    elif tmp_update_l2_table_array[3] in current_direction_if_array[2]:
                        # DOWN
                        if flag_exist_if_vport_array[1][1] == True:
                            for tmp_tag_size_array in tag_size_array:
                                if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                    inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                    inche_from_connect_y = tmp_tag_size_array[1]
                                    inche_to_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                    inche_to_connect_y = device_size_array[1] + device_size_array[3] - 0.1 - l2seg_size_margin
                                    line_type = 'NORMAL'
                                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                    inche_from_connect_x = inche_to_connect_x
                                    inche_from_connect_y = inche_to_connect_y

                        else:
                            for tmp_tag_size_array in tag_size_array:
                                if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                    inche_from_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5)
                                    inche_from_connect_y = tmp_tag_size_array[1]

                        modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                        for tmp_l2seg_size_array in l2seg_size_array:
                            for tmp_update_l2_table_array in modify_update_l2_table_array:
                                if tmp_l2seg_size_array[4] == tmp_update_l2_table_array:
                                    #print('tmp_tmp_update_l2_table_array   ',tmp_tmp_update_l2_table_array )
                                    inche_to_connect_x = tmp_l2seg_size_array[0] + (tmp_l2seg_size_array[2] * 0.1)
                                    inche_to_connect_y = tmp_l2seg_size_array[1] + tmp_l2seg_size_array[3]
                                    line_type = 'NORMAL'
                                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                    elif tmp_update_l2_table_array[3] in current_direction_if_array[3]:
                        #RIHGT
                        if flag_exist_if_vport_array[2][1] == True:
                            for tmp_tag_size_array in tag_size_array:
                                if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                    inche_from_connect_x = tmp_tag_size_array[0]
                                    inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)
                                    inche_to_connect_x = tmp_tag_size_array[0]  - l2seg_size_margin - l2seg_size_margin_left_right_add
                                    inche_to_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)
                                    line_type = 'NORMAL'
                                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                    inche_from_connect_x = inche_to_connect_x
                                    inche_from_connect_y = inche_to_connect_y

                        else:
                            for tmp_tag_size_array in tag_size_array:
                                if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                    inche_from_connect_x = tmp_tag_size_array[0]
                                    inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)

                        modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                        for tmp_l2seg_size_array in l2seg_size_array:
                            for tmp_update_l2_table_array in modify_update_l2_table_array:
                                if tmp_l2seg_size_array[4] == tmp_update_l2_table_array:
                                    # print('tmp_tmp_update_l2_table_array   ',tmp_tmp_update_l2_table_array )
                                    inche_to_connect_x = tmp_l2seg_size_array[0] + tmp_l2seg_size_array[2]
                                    inche_to_connect_y = tmp_l2seg_size_array[1] + (tmp_l2seg_size_array[3] * 0.5)
                                    line_type = 'NORMAL'
                                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)


                    elif tmp_update_l2_table_array[3] in current_direction_if_array[4]:
                        #LEFT
                        if flag_exist_if_vport_array[3][1] == True:
                            for tmp_tag_size_array in tag_size_array:
                                if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                    inche_from_connect_x = tmp_tag_size_array[0] + tmp_tag_size_array[2]
                                    inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)
                                    inche_to_connect_x = tmp_tag_size_array[0] + (tmp_tag_size_array[2] * 0.5) + l2seg_size_margin +l2seg_size_margin_left_right_add + 0.3
                                    inche_to_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)
                                    line_type = 'NORMAL'
                                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)

                                    inche_from_connect_x = inche_to_connect_x
                                    inche_from_connect_y = inche_to_connect_y

                        else:
                            for tmp_tag_size_array in tag_size_array:
                                if tmp_tag_size_array[5] == tmp_update_l2_table_array[3]:
                                    inche_from_connect_x = tmp_tag_size_array[0] + tmp_tag_size_array[2]
                                    inche_from_connect_y = tmp_tag_size_array[1] + (tmp_tag_size_array[3] * 0.5)

                        modify_update_l2_table_array = str(tmp_update_l2_table_array[6]).replace(' ', '').split(',')
                        for tmp_l2seg_size_array in l2seg_size_array:
                            for tmp_update_l2_table_array in modify_update_l2_table_array:
                                if tmp_l2seg_size_array[4] == tmp_update_l2_table_array:
                                    # print('tmp_tmp_update_l2_table_array   ',tmp_tmp_update_l2_table_array )
                                    inche_to_connect_x = tmp_l2seg_size_array[0]
                                    inche_to_connect_y = tmp_l2seg_size_array[1] + (tmp_l2seg_size_array[3] * 0.5)
                                    line_type = 'NORMAL'
                                    ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y, inche_to_connect_x, inche_to_connect_y)


        '''write device frame'''
        self.shape = self.slide.shapes
        if target_device_name in wp_list_array:
            tmp_device_type = 'WAY_POINT'
        else:
            tmp_device_type = 'DEVICE_FRAME'

        ns_ddx_figure.extended.add_shape(self, tmp_device_type, device_size_array[0], device_size_array[1], device_size_array[2], device_size_array[3], target_device_name)
        self.slide.shapes._spTree.remove(self.shape._element)   # move shape to back layer
        self.slide.shapes._spTree.insert(2, self.shape._element)  # move shape to back layer


        if action_type == 'RETURN_DEVICE_SIZE':
            return ([device_size_array[0], device_size_array[1],device_size_array[2],device_size_array[3]])

        return ()



