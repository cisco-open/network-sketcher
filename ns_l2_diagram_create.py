'''
SPDX-License-Identifier: Apache-2.0

Copyright 2023 Cisco Systems, Inc. and its affiliates

Complete cross-platform multiprocessing version with adaptive process count
'''

from pptx import *
import sys
import os
import re
import numpy as np
import math
import ns_def
import ns_ddx_figure
import openpyxl
from pptx import Presentation
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from pptx.dml.line import LineFormat
from pptx.shapes.connector import Connector
from pptx.util import Inches, Cm, Pt
import time
import json
from datetime import datetime
from multiprocessing import Pool, cpu_count
import traceback
import pickle
import platform

# ===================================================================
# Per-process worker context (initializer pattern for performance)
# ===================================================================

_WORKER_CTX = None
_WORKER_WP_LIST = None

def is_multiprocessing_supported():
    """
    Return True if multiprocessing Pool is supported in this runtime.
    If platform detection fails or unsupported, return False.
    """
    try:
        _ = cpu_count()

        sysname = platform.system()

        if sysname == 'Windows':
            import multiprocessing as mp
            try:
                return 'spawn' in mp.get_all_start_methods()
            except Exception:
                return False

        return True

    except Exception:
        return False

def worker_initializer(serialized_data, wp_list_array):
    """
    Pool initializer: runs once per worker process (cross-platform safe).
    Builds CompleteContext once per process, avoiding per-device overhead.
    """
    global _WORKER_CTX, _WORKER_WP_LIST

    # Force reinit (ignore any stale state)
    _WORKER_CTX = None
    _WORKER_WP_LIST = None

    try:
        class CompleteContext:
            def __init__(self, data_dict):
                self.position_style_shape_array = data_dict['position_style_shape_array']
                self.position_shape_array = data_dict['position_shape_array']
                self.position_line_array = data_dict['position_line_array']
                self.position_folder_array = data_dict['position_folder_array']
                self.position_tag_array = data_dict['position_tag_array']
                self.root_folder_array = data_dict['root_folder_array']
                self.folder_wp_name_array = data_dict['folder_wp_name_array']
                self.l2_table_array = data_dict['l2_table_array']

                self.position_style_shape_tuple = ns_def.convert_array_to_tuple(self.position_style_shape_array)
                self.position_shape_tuple = ns_def.convert_array_to_tuple(self.position_shape_array)
                self.position_line_tuple = ns_def.convert_array_to_tuple(self.position_line_array)
                self.position_folder_tuple = ns_def.convert_array_to_tuple(self.position_folder_array)
                self.position_tag_tuple = ns_def.convert_array_to_tuple(self.position_tag_array)
                self.root_folder_tuple = ns_def.convert_array_to_tuple(self.root_folder_array)

                self.ppt_edge_margin = 1.0
                self.ppt_width = 15.0
                self.ppt_hight = 10.0
                self.active_ppt = Presentation()

                for key, value in data_dict.items():
                    if not hasattr(self, key):
                        setattr(self, key, value)

                class DummyControl:
                    def __init__(self, value=''):
                        self.value = value
                    def get(self):
                        return self.value

                for key in list(data_dict.keys()):
                    if key.endswith('_value'):
                        control_name = key[:-6]
                        setattr(self, control_name, DummyControl(data_dict[key]))

            def _ensure_l2_shared_prepared(self):
                # If shared arrays are missing, rebuild even if flag says prepared
                if (not hasattr(self, "l2_shared_data_prepared") or not self.l2_shared_data_prepared
                        or not hasattr(self, "shared_new_l2_table_array")
                        or not hasattr(self, "shared_update_l2_table_array")
                        or not hasattr(self, "shared_new_direction_if_array")):
                    ns_ddx_figure.ns_ddx_figure_run.prepare_l2_shared_data(self)

        _WORKER_CTX = CompleteContext(serialized_data)

        # IMPORTANT: reset any cached/shared state from previous runs
        _WORKER_CTX.l2_shared_data_prepared = False
        for k in list(vars(_WORKER_CTX).keys()):
            if k.startswith("shared_"):
                delattr(_WORKER_CTX, k)

        _WORKER_WP_LIST = wp_list_array

    except Exception as e:
        _WORKER_CTX = None
        _WORKER_WP_LIST = None

# ===================================================================
# Cross-platform configuration
# ===================================================================

def get_platform_config():
    """Get platform-specific configuration"""
    try:
        system = platform.system()
    except Exception:
        system = "Unknown"

    try:
        cpun = cpu_count()
    except Exception:
        cpun = 1

    config = {
        'system': system,
        'is_windows': system == 'Windows',
        'is_macos': system == 'Darwin',
        'is_linux': system == 'Linux',
        'python_version': platform.python_version(),
        'cpu_count': cpun
    }
    return config


def calculate_optimal_processes(total_devices):
    """
    Calculate optimal number of processes
    Balances performance and efficiency based on system capabilities
    """
    try:
        cpu_cores = cpu_count()
    except Exception:
        cpu_cores = 1

    # Leave 2 cores for system
    base_processes = max(1, cpu_cores - 2)

    # Apply efficiency-based caps
    if cpu_cores <= 8:
        optimal = base_processes
        reason = "Small system - using most cores"
    elif cpu_cores <= 16:
        optimal = min(base_processes, 12)
        reason = "Medium system - capped at 12 for efficiency"
    elif cpu_cores <= 32:
        optimal = min(base_processes, 16)
        reason = "Large system - capped at 16 for efficiency"
    else:
        optimal = min(base_processes, 24)
        reason = "Very large system - capped at 24 for efficiency"

    # Don't use more processes than devices
    if optimal > total_devices:
        optimal = total_devices
        reason = f"Limited to device count ({total_devices})"

    optimal = max(1, optimal)

    # Estimate efficiency
    if optimal <= 4:
        est_eff = 60
    elif optimal <= 8:
        est_eff = 45
    elif optimal <= 12:
        est_eff = 30
    elif optimal <= 16:
        est_eff = 22
    elif optimal <= 24:
        est_eff = 18
    else:
        est_eff = 15

    return optimal, reason, est_eff


class PerformanceMonitor:
    """Performance monitoring and progress management class"""

    def __init__(self, enable_logging=True):
        self.logs = {}
        self.start_times = {}
        self.enable_logging = enable_logging
        self.parent_child_relationships = {}

    def start(self, section_name, parent=None):
        self.start_times[section_name] = time.time()
        if parent:
            self.parent_child_relationships[section_name] = parent
        if self.enable_logging:
            print(f"[START] {section_name}")

    def end(self, section_name, details=None):
        if section_name in self.start_times:
            elapsed = time.time() - self.start_times[section_name]
            self.logs[section_name] = {
                'elapsed_time': round(elapsed, 2),
                'details': details,
                'is_child': section_name in self.parent_child_relationships
            }
            if self.enable_logging:
                detail_str = f" ({details})" if details else ""
                print(f"[END] {section_name}: {elapsed:.2f}s{detail_str}")
            del self.start_times[section_name]

    def progress(self, current, total, message=""):
        if self.enable_logging and total > 0:
            percentage = (current / total) * 100
            print(f"  Progress: {current}/{total} ({percentage:.1f}%) {message}")

    def save_log(self, filepath='performance_log.json'):
        try:
            # Cross-platform path handling
            filepath = os.path.normpath(filepath)
            log_dir = os.path.dirname(filepath)
            if log_dir and not os.path.exists(log_dir):
                os.makedirs(log_dir, exist_ok=True)

            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(self.logs, f, indent=2, ensure_ascii=False)
            if self.enable_logging:
                print(f"[LOG SAVED] {filepath}")
        except Exception as e:
            print(f"Warning: Failed to save log: {e}")

    def print_summary(self):
        if not self.logs:
            return

        print("\n" + "=" * 80)
        print("Performance Summary")
        print("=" * 80)

        actual_time = self.logs.get('Total Processing', {}).get('elapsed_time', 0)
        sorted_logs = sorted(self.logs.items(), key=lambda x: x[1]['elapsed_time'], reverse=True)

        parent_sections = []
        child_sections = {}

        for section, log in sorted_logs:
            is_child = log.get('is_child', False)
            if is_child:
                parent = self.parent_child_relationships.get(section)
                if parent not in child_sections:
                    child_sections[parent] = []
                child_sections[parent].append((section, log))
            else:
                parent_sections.append((section, log))

        for section, log in parent_sections:
            elapsed = log['elapsed_time']
            percentage = (elapsed / actual_time * 100) if actual_time > 0 else 0
            details = log.get('details', '')
            detail_str = f" - {details}" if details else ""

            print(f"{section:45s}: {elapsed:8.2f}s ({percentage:5.1f}%){detail_str}")

            if section in child_sections:
                for child_section, child_log in child_sections[section]:
                    child_elapsed = child_log['elapsed_time']
                    child_percentage = (child_elapsed / actual_time * 100) if actual_time > 0 else 0
                    child_details = child_log.get('details', '')
                    child_detail_str = f" - {child_details}" if child_details else ""
                    print(f"  └─ {child_section:41s}: {child_elapsed:8.2f}s ({child_percentage:5.1f}%){child_detail_str}")

        print("-" * 80)
        print(f"{'Total wall clock time':45s}: {actual_time:8.2f}s")
        print("=" * 80)


def extract_serializable_attributes(obj):
    """Extract all serializable attributes from object"""
    attrs = {}

    for attr_name in dir(obj):
        if attr_name.startswith('_'):
            continue

        try:
            attr_value = getattr(obj, attr_name)

            if callable(attr_value):
                if hasattr(attr_value, 'get') and callable(attr_value.get):
                    try:
                        attrs[attr_name + '_value'] = attr_value.get()
                    except:
                        attrs[attr_name + '_value'] = ''
                continue

            try:
                pickle.dumps(attr_value)
                attrs[attr_name] = attr_value
            except:
                if isinstance(attr_value, (str, int, float, bool, type(None))):
                    attrs[attr_name] = attr_value
                elif isinstance(attr_value, (list, dict, tuple)):
                    try:
                        pickle.dumps(attr_value)
                        attrs[attr_name] = attr_value
                    except:
                        pass

        except:
            continue

    return attrs


# ===================================================================
# CRITICAL: Worker function must be at module level
# ===================================================================

def process_device_worker(device_name):
    """
    Worker function: device_name only.
    Uses per-process global context initialized by worker_initializer().
    MUST be at module level for cross-platform multiprocessing.
    """
    global _WORKER_CTX, _WORKER_WP_LIST

    try:
        if _WORKER_CTX is None:
            return (device_name, None, "WorkerContextNotInitialized")

        device_size_array = ns_ddx_figure.extended.l2_device_materials(
            _WORKER_CTX, 'RETURN_DEVICE_SIZE', device_name, [], _WORKER_WP_LIST
        )
        return (device_name, device_size_array, None)

    except Exception as e:
        return (device_name, None, f"{type(e).__name__}: {str(e)}")

class ns_l2_diagram_create():
    """L2 Network Diagram Creation Class (Cross-platform with Adaptive Processes)"""

    def __init__(self):
        # Get platform info
        platform_config = get_platform_config()

        print("=" * 80)
        print("Network Diagram Generation Started (Cross-platform Multiprocessing)")
        print(f"Platform: {platform_config['system']} | Python: {platform_config['python_version']}")
        print(f"CPU Cores: {platform_config['cpu_count']}")
        print(f"Start time: {datetime.now()}")
        print("=" * 80)

        self.perf_monitor = PerformanceMonitor(enable_logging=True)
        self.perf_monitor.start('Total Processing')

        try:
            '''STEP0: Load Master Data'''
            self.perf_monitor.start('STEP0: Load Master Data', parent='Total Processing')

            ws_name = 'Master_Data'
            ws_l2_name = 'Master_Data_L2'
            excel_maseter_file = self.inFileTxt_L2_3_1.get()

            # Cross-platform path normalization
            excel_maseter_file = os.path.normpath(excel_maseter_file)
            print(f"Excel file: {excel_maseter_file}")

            self.position_folder_array = ns_def.convert_master_to_array(ws_name, excel_maseter_file, '<<POSITION_FOLDER>>')
            self.position_shape_array = ns_def.convert_master_to_array(ws_name, excel_maseter_file, '<<POSITION_SHAPE>>')
            self.position_line_array = ns_def.convert_master_to_array(ws_name, excel_maseter_file, '<<POSITION_LINE>>')
            self.position_style_shape_array = ns_def.convert_master_to_array(ws_name, excel_maseter_file, '<<STYLE_SHAPE>>')
            self.position_tag_array = ns_def.convert_master_to_array(ws_name, excel_maseter_file, '<<POSITION_TAG>>')
            self.root_folder_array = ns_def.convert_master_to_array(ws_name, excel_maseter_file, '<<ROOT_FOLDER>>')

            self.position_folder_tuple = ns_def.convert_array_to_tuple(self.position_folder_array)
            self.position_shape_tuple = ns_def.convert_array_to_tuple(self.position_shape_array)
            self.position_line_tuple = ns_def.convert_array_to_tuple(self.position_line_array)
            self.position_style_shape_tuple = ns_def.convert_array_to_tuple(self.position_style_shape_array)
            self.position_tag_tuple = ns_def.convert_array_to_tuple(self.position_tag_array)
            self.root_folder_tuple = ns_def.convert_array_to_tuple(self.root_folder_array)

            self.folder_wp_name_array = ns_def.get_folder_wp_array_from_master(ws_name, excel_maseter_file)
            self.l2_table_array = ns_def.convert_master_to_array(ws_l2_name, excel_maseter_file, '<<L2_TABLE>>')

            total_data = (len(self.position_folder_array) + len(self.position_shape_array) +
                         len(self.position_line_array) + len(self.l2_table_array))
            self.perf_monitor.end('STEP0: Load Master Data', f'Data rows: {total_data}')

            '''STEP1: Create Device List'''
            self.perf_monitor.start('STEP1: Create Device List', parent='Total Processing')

            new_l2_table_array = []
            for tmp_l2_table_array in self.l2_table_array:
                if tmp_l2_table_array[0] != 1 and tmp_l2_table_array[0] != 2:
                    tmp_l2_table_array[1].extend(['', '', '', '', '', '', '', ''])
                    del tmp_l2_table_array[1][8:]
                    new_l2_table_array.append(tmp_l2_table_array)

            device_set = set()
            wp_set = set()
            device_list_array = []
            wp_list_array = []

            target_area = str(self.comboL2_3_6.get())

            for tmp_new_l2_table_array in new_l2_table_array:
                device_name = tmp_new_l2_table_array[1][1]
                area_name = tmp_new_l2_table_array[1][0]

                if area_name == 'N/A':
                    if device_name not in wp_set:
                        wp_set.add(device_name)
                        wp_list_array.append(device_name)
                else:
                    if target_area == area_name:
                        if device_name not in device_set:
                            device_set.add(device_name)
                            device_list_array.append(device_name)

            #print(f"Target area: {target_area}")
            #print(f"Device count: {len(device_list_array)}")
            #print(f"WP count: {len(wp_list_array)}")

            self.perf_monitor.end('STEP1: Create Device List',
                                 f'Devices: {len(device_list_array)}, WPs: {len(wp_list_array)}')

            '''STEP2: GET all device l2 size (OPTIMIZED DATA TRANSFER)'''
            self.perf_monitor.start('STEP2: Get Device Sizes (Optimized)', parent='Total Processing')

            self.ppt_edge_margin = 1.0
            self.ppt_width = 15.0
            self.ppt_hight = 10.0

            all_devices = device_list_array + wp_list_array
            total_devices = len(all_devices)

            # ===== DECISION: Use appropriate method based on device count =====
            USE_MULTIPROCESSING = is_multiprocessing_supported()
            #USE_MULTIPROCESSING = False   # for single-core testing

            global _WORKER_CTX, _WORKER_WP_LIST

            def _run_single_core():
                global _WORKER_CTX, _WORKER_WP_LIST
                # Single-core fallback (use same context model as multiprocessing)
                nonlocal wp_list_array, all_devices
                self.all_device_l2_size_array = []
                nonlocal success_count, failed_devices, start_time
                success_count = 0
                failed_devices = []
                start_time = time.time()

                print("Extracting attributes... (single-core)")
                all_attrs = extract_serializable_attributes(self)

                serialized_data = {
                    'position_style_shape_array': self.position_style_shape_array,
                    'position_shape_array': self.position_shape_array,
                    'position_line_array': self.position_line_array,
                    'position_folder_array': self.position_folder_array,
                    'position_tag_array': self.position_tag_array,
                    'root_folder_array': self.root_folder_array,
                    'folder_wp_name_array': self.folder_wp_name_array,
                    'l2_table_array': self.l2_table_array,
                }
                serialized_data.update(all_attrs)

                _WORKER_CTX = None
                _WORKER_WP_LIST = None
                worker_initializer(serialized_data, wp_list_array)

                if _WORKER_CTX is None:
                    raise RuntimeError("Single-core fallback failed: Worker context not initialized")

                for device_name in all_devices:
                    try:
                        device_size_array = ns_ddx_figure.extended.l2_device_materials(
                            _WORKER_CTX, 'RETURN_DEVICE_SIZE', device_name, [], wp_list_array
                        )
                        self.all_device_l2_size_array.append([device_name, device_size_array])
                        success_count += 1
                    except Exception as e:
                        failed_devices.append((device_name, f"{type(e).__name__}: {str(e)}"))

            if USE_MULTIPROCESSING:
                # Clear any stale worker context from previous runs
                _WORKER_CTX = None
                _WORKER_WP_LIST = None
                # Small dataset: Multiprocessing is effective
                #print(f"\n{'=' * 70}")
                #print(f"Using Multiprocessing (Optimal for <={256} devices)")
                #print(f"  Total devices: {total_devices}")

                num_processes, optimization_reason, estimated_efficiency = calculate_optimal_processes(total_devices)

                #print(f"  Processes: {num_processes}")
                #print(f"  Expected speedup: {num_processes * estimated_efficiency / 100:.1f}x")
                #print(f"{'=' * 70}\n")

                print("Extracting attributes...")
                all_attrs = extract_serializable_attributes(self)

                serialized_data = {
                    'position_style_shape_array': self.position_style_shape_array,
                    'position_shape_array': self.position_shape_array,
                    'position_line_array': self.position_line_array,
                    'position_folder_array': self.position_folder_array,
                    'position_tag_array': self.position_tag_array,
                    'root_folder_array': self.root_folder_array,
                    'folder_wp_name_array': self.folder_wp_name_array,
                    'l2_table_array': self.l2_table_array,
                }
                serialized_data.update(all_attrs)

                self.all_device_l2_size_array = []
                success_count = 0
                failed_devices = []
                start_time = time.time()

                print("Starting parallel processing...")

                try:
                    with Pool(
                            processes=num_processes,
                            initializer=worker_initializer,
                            initargs=(serialized_data, wp_list_array),
                    ) as pool:
                        chunksize = max(1, total_devices // (num_processes * 4))
                        results_iter = pool.imap_unordered(process_device_worker, all_devices, chunksize=chunksize)

                        progress_interval = max(1, total_devices // 50)

                        for idx, result in enumerate(results_iter, 1):
                            device_name, device_size_array, error = result

                            if error is None and device_size_array is not None:
                                self.all_device_l2_size_array.append([device_name, device_size_array])
                                success_count += 1
                            else:
                                failed_devices.append((device_name, error))
                                if len(failed_devices) <= 10:
                                    print(f"  Warning: {device_name} - {error}")

                            if idx % progress_interval == 0 or idx == total_devices:
                                elapsed = time.time() - start_time
                                rate = idx / elapsed if elapsed > 0 else 0
                                eta_min = ((total_devices - idx) / rate / 60) if rate > 0 else 0
                                speedup = rate / 1.02 if rate > 0 else 0

                                self.perf_monitor.progress(
                                    idx, total_devices,
                                    f"{device_name} | Rate: {rate:.2f} dev/s | Speedup: {speedup:.1f}x | "
                                    f"ETA: {eta_min:.1f}min | Success: {success_count}/{idx}"
                                )

                except Exception as e:
                    print(f"\n[ERROR] Multiprocessing failed: {e}")
                    traceback.print_exc()
                    print("[FALLBACK] Switching to single-core processing...")
                    USE_MULTIPROCESSING = False
                    _run_single_core()

            else:
                _run_single_core()

            # Common ending for both paths
            if failed_devices:
                print(f"\n{'=' * 70}")
                print(f"Warning: {len(failed_devices)} devices failed")
                for device_name, error in failed_devices[:10]:
                    print(f"  - {device_name}: {error}")
                if len(failed_devices) > 10:
                    print(f"  ... and {len(failed_devices) - 10} more")
                print(f"{'=' * 70}\n")

            self.device_size_dict = {item[0]: item[1] for item in self.all_device_l2_size_array}

            total_time = time.time() - start_time
            actual_rate = success_count / total_time if total_time > 0 else 0

            method_used = "Multiprocessing" if USE_MULTIPROCESSING else "Single-thread Batch"

            self.perf_monitor.end('STEP2: Get Device Sizes (Optimized)',
                                  f'Success: {success_count}/{total_devices}, Failed: {len(failed_devices)}, '
                                  f'Method: {method_used}, Rate: {actual_rate:.2f} dev/s')




            '''Create per area l2 ppt'''
            if self.click_value == 'L2-3-2':
                self.perf_monitor.start('STEP3: Create Area Diagram', parent='Total Processing')
                action_type = 'CREATE_L2_AREA'
                data_array = []
                ns_l2_diagram_create.l2_area_create(self, action_type, data_array)
                self.perf_monitor.end('STEP3: Create Area Diagram')

            '''Create per device l2 ppt'''
            if self.click_value == 'L2-3-3':
                self.perf_monitor.start('STEP3: Create Device Diagram', parent='Total Processing')

                print("Determining PPT size...")
                for device_name, size_info in self.device_size_dict.items():
                    if size_info[2] > self.ppt_width:
                        self.ppt_width = size_info[2]
                    if size_info[3] > self.ppt_hight:
                        self.ppt_hight = size_info[3]

                print(f"PPT size: width={self.ppt_width:.1f}in, height={self.ppt_hight:.1f}in")

                self.active_ppt = Presentation()

                total_devices = len(self.all_device_l2_size_array)
                print(f"Creating device diagrams: {total_devices} devices")

                progress_interval = max(1, total_devices // 10)

                for idx, tmp_all_device_size_array in enumerate(self.all_device_l2_size_array, 1):
                    if idx % progress_interval == 0 or idx == total_devices:
                        self.perf_monitor.progress(idx, total_devices, tmp_all_device_size_array[0])

                    input_device_name = tmp_all_device_size_array[0]
                    device_size_array = tmp_all_device_size_array[1]
                    action_type = 'WRITE_DEVICE_L2'
                    write_left_top_array = [3.0, 3.0, device_size_array]

                    try:
                        ns_ddx_figure.extended.l2_device_materials(
                            self, action_type, input_device_name, write_left_top_array, wp_list_array
                        )
                    except Exception as e:
                        print(f"  Warning: Failed to create diagram for '{input_device_name}': {e}")
                        continue

                # Cross-platform path handling for output
                output_file = os.path.normpath(self.output_ppt_file)
                print(f"\nSaving PPT file: {output_file}")
                try:
                    self.active_ppt.save(output_file)
                    print(f"Save completed: {output_file}")
                except Exception as e:
                    print(f"Error: Failed to save PPT file: {e}")
                    raise

                self.perf_monitor.end('STEP3: Create Device Diagram', f'Created: {total_devices}')
                return

        except Exception as e:
            print(f"\n[ERROR] Exception occurred: {e}")
            traceback.print_exc()

        finally:
            self.perf_monitor.end('Total Processing')
            self.perf_monitor.print_summary()

            # Cross-platform log path
            log_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'performance_log.json')
            #self.perf_monitor.save_log(log_path)

            print("\n" + "=" * 80)
            print(f"End time: {datetime.now()}")
            print("=" * 80)

    def l2_area_create(self, action_type, data_array):
        """Area-based L2 diagram creation"""
        #self.perf_monitor.start('Area Processing Detail', parent='STEP3: Create Area Diagram')

        ws_name = 'Master_Data'
        tmp_ws_name = '_tmp_'
        ppt_meta_file = self.inFileTxt_L2_3_1.get()

        # Initialize slide size attributes
        self.all_slide_max_width = 0
        self.all_slide_max_hight = 0

        #print("Converting L2 device style...")
        l2_position_style_shape_array = []

        for tmp_position_style_shape_array in self.position_style_shape_array:
            if tmp_position_style_shape_array[0] in [1, 2, 3]:
                l2_position_style_shape_array.append(tmp_position_style_shape_array)
            else:
                device_name = tmp_position_style_shape_array[1][0]
                if device_name in self.device_size_dict:
                    size_info = self.device_size_dict[device_name]
                    tmp_position_style_shape_array[1][1] = size_info[2]
                    tmp_position_style_shape_array[1][2] = size_info[3]
                l2_position_style_shape_array.append(tmp_position_style_shape_array)

        l2_position_style_shape_tuple = ns_def.convert_array_to_tuple(l2_position_style_shape_array)

        #print("Creating WayPoint-Folder mapping...")
        wp_with_folder_tuple = {}

        for tmp_wp_folder_name in self.folder_wp_name_array[1]:
            current_row = 1
            flag_start_row = False
            flag_end_row = False

            while flag_end_row == False:
                if str(self.position_shape_tuple[current_row, 1]) == tmp_wp_folder_name:
                    start_row = current_row
                    flag_start_row = True
                if flag_start_row == True and str(self.position_shape_tuple[current_row, 1]) == '<END>':
                    flag_end_row = True
                    end_row = current_row - 1
                current_row += 1

            for i in range(start_row, end_row + 1):
                flag_start_column = False
                current_column = 2
                while flag_start_column == False:
                    if str(self.position_shape_tuple[i, current_column]) != '<END>':
                        wp_with_folder_tuple[self.position_shape_tuple[i, current_column]] = tmp_wp_folder_name
                    else:
                        flag_start_column = True
                    current_column += 1

        tmp_folder_name = self.comboL2_3_6.get()
        print(f"Target folder: {tmp_folder_name}")

        current_row = 1
        flag_start_row = False
        flag_end_row = False

        while flag_end_row == False:
            if str(self.position_shape_tuple[current_row, 1]) == tmp_folder_name:
                start_row = current_row
                flag_start_row = True
            if flag_start_row == True and str(self.position_shape_tuple[current_row, 1]) == '<END>':
                flag_end_row = True
                end_row = current_row - 1
            current_row += 1

        tmp_folder_array = []
        for i in range(start_row, end_row + 1):
            flag_start_column = False
            current_column = 2
            while flag_start_column == False:
                if str(self.position_shape_tuple[i, current_column]) != '<END>':
                    tmp_folder_array.append(self.position_shape_tuple[i, current_column])
                else:
                    flag_start_column = True
                current_column += 1

        #print(f"Shapes in folder: {len(tmp_folder_array)}")

        connected_wp_folder_set = set()

        for tmp_shpae_name in tmp_folder_array:
            for tmp_position_line_tuple in self.position_line_tuple:
                if tmp_position_line_tuple[0] != 1:
                    if tmp_shpae_name == self.position_line_tuple[tmp_position_line_tuple[0], 1]:
                        for tmp_wp_with_folder_tuple in wp_with_folder_tuple:
                            if tmp_wp_with_folder_tuple == self.position_line_tuple[tmp_position_line_tuple[0], 2]:
                                connected_wp_folder_set.add(wp_with_folder_tuple[tmp_wp_with_folder_tuple])
                    if tmp_shpae_name == self.position_line_tuple[tmp_position_line_tuple[0], 2]:
                        for tmp_wp_with_folder_tuple in wp_with_folder_tuple:
                            if tmp_wp_with_folder_tuple == self.position_line_tuple[tmp_position_line_tuple[0], 1]:
                                connected_wp_folder_set.add(wp_with_folder_tuple[tmp_wp_with_folder_tuple])

        connected_wp_folder_array = list(connected_wp_folder_set)
        #print(f"Connected WP folders: {len(connected_wp_folder_array)}")

        extract_folder_tuple = {}
        for tmp_position_folder_tuple in self.position_folder_tuple:
            if self.position_folder_tuple[tmp_position_folder_tuple] == tmp_folder_name or \
               self.position_folder_tuple[tmp_position_folder_tuple] in connected_wp_folder_array:
                extract_folder_tuple[tmp_position_folder_tuple] = self.position_folder_tuple[tmp_position_folder_tuple]
                extract_folder_tuple[tmp_position_folder_tuple[0] - 1, tmp_position_folder_tuple[1]] = \
                    self.position_folder_tuple[tmp_position_folder_tuple[0] - 1, tmp_position_folder_tuple[1]]
                extract_folder_tuple[tmp_position_folder_tuple[0], 1] = \
                    self.position_folder_tuple[tmp_position_folder_tuple[0], 1]
                extract_folder_tuple[tmp_position_folder_tuple[0] - 1, 1] = \
                    self.position_folder_tuple[tmp_position_folder_tuple[0] - 1, 1]

        print(f"Copying Excel sheet: {ws_name} -> {tmp_ws_name}")
        ns_def.copy_excel_sheet(ws_name, ppt_meta_file, tmp_ws_name)

        clear_section_tuple = self.position_folder_tuple
        ns_def.clear_section_sheet(tmp_ws_name, ppt_meta_file, clear_section_tuple)

        convert_array = ns_def.convert_tuple_to_array(extract_folder_tuple)
        offset_row = convert_array[0][0] * -1 + 2

        current_y_grid_array = []
        flag_first_array = True
        if convert_array[0][1][0] == '<SET_WIDTH>':
            for tmp_array in convert_array:
                current_y_grid_array.append([tmp_array[0] + offset_row, tmp_array[1]])
        else:
            for tmp_array in convert_array:
                if flag_first_array == True:
                    current_y_grid_array.append([tmp_array[0] + offset_row - 1, tmp_array[1]])
                    flag_first_array = False
                else:
                    current_y_grid_array.append([tmp_array[0] + offset_row - 1, tmp_array[1]])

        convert_tuple = ns_def.convert_array_to_tuple(current_y_grid_array)

        if action_type == 'CREATE_L2_AREA':
            ns_def.write_excel_meta(convert_tuple, ppt_meta_file, tmp_ws_name, '<<POSITION_FOLDER>>', 0, 0)

        if action_type == 'CREATE_L2_AREA':
            ns_def.clear_section_sheet(tmp_ws_name, ppt_meta_file, self.position_style_shape_tuple)

        ns_def.write_excel_meta(l2_position_style_shape_tuple, ppt_meta_file, tmp_ws_name, '<<STYLE_SHAPE>>', 0, 0)

        print("Calculating folder size...")
        master_folder_size_array = ns_def.get_folder_width_size(
            convert_tuple, l2_position_style_shape_tuple, self.position_shape_tuple, 0.8
        )

        master_root_folder_tuple = ns_def.get_root_folder_tuple(self, master_folder_size_array, tmp_folder_name)

        clear_section_tuple = dict(master_root_folder_tuple)
        clear_section_tuple[1, 1] = '<<ROOT_FOLDER>>'
        if action_type == 'CREATE_L2_AREA':
            ns_def.clear_section_sheet(tmp_ws_name, ppt_meta_file, clear_section_tuple)
            ns_def.write_excel_meta(master_root_folder_tuple, ppt_meta_file, tmp_ws_name, '<<ROOT_FOLDER>>', 0, 0)

        import copy
        update_current_y_grid_array = copy.deepcopy(current_y_grid_array)

        tmp_sheet_position_folder_array = ns_def.convert_master_to_array(tmp_ws_name, ppt_meta_file, '<<POSITION_FOLDER>>')
        tmp_sheet_position_folder_tuple = ns_def.convert_array_to_tuple(tmp_sheet_position_folder_array)

        for tmp_tmp_sheet_position_folder_tuple in tmp_sheet_position_folder_tuple:
            for tmp_master_folder_size_array_2 in master_folder_size_array[2]:
                if tmp_sheet_position_folder_tuple[tmp_tmp_sheet_position_folder_tuple[0], tmp_tmp_sheet_position_folder_tuple[1]] == \
                   tmp_master_folder_size_array_2[1][0][0] and tmp_master_folder_size_array_2[1][0][0] != 10 and \
                   isinstance(tmp_sheet_position_folder_tuple[tmp_tmp_sheet_position_folder_tuple[0], tmp_tmp_sheet_position_folder_tuple[1]], str) and \
                   isinstance(tmp_master_folder_size_array_2[1][0][0], str):
                    tmp_num = 0
                    for tmp_current_y_grid_array in current_y_grid_array:
                        if tmp_current_y_grid_array[0] == tmp_tmp_sheet_position_folder_tuple[0] - 1 and \
                           tmp_master_folder_size_array_2[1][0][1] != 0:
                            update_current_y_grid_array[tmp_num][1][tmp_tmp_sheet_position_folder_tuple[1] - 1] = \
                                tmp_master_folder_size_array_2[1][0][1]
                            break
                        tmp_num += 1

        tmp_num = 0
        for tmp_current_y_grid_array in current_y_grid_array:
            if isinstance(tmp_current_y_grid_array[1][0], (int, float)):
                max_y_grid_current = tmp_current_y_grid_array[1][0]
                for tmp_tmp_current_y_grid_array in tmp_current_y_grid_array[1:]:
                    for tmp_master_folder_size_array_2 in master_folder_size_array[2]:
                        if tmp_master_folder_size_array_2[1][0][0] == tmp_tmp_current_y_grid_array[1]:
                            if max_y_grid_current < tmp_master_folder_size_array_2[1][0][2]:
                                max_y_grid_current = tmp_master_folder_size_array_2[1][0][2]
                                update_current_y_grid_array[tmp_num][1][0] = max_y_grid_current
                                break
            tmp_num += 1

        update_current_y_grid_tuple = ns_def.convert_array_to_tuple(update_current_y_grid_array)

        clear_sheet_position_folder_array = ns_def.convert_master_to_array(tmp_ws_name, ppt_meta_file, '<<POSITION_FOLDER>>')
        clear_sheet_position_folder_tuple = ns_def.convert_array_to_tuple(clear_sheet_position_folder_array)
        ns_def.clear_section_sheet(tmp_ws_name, ppt_meta_file, clear_sheet_position_folder_tuple)

        if master_root_folder_tuple[2, 7] > self.all_slide_max_width:
            self.all_slide_max_width = master_root_folder_tuple[2, 7]
        if master_root_folder_tuple[2, 8] > self.all_slide_max_hight:
            self.all_slide_max_hight = master_root_folder_tuple[2, 8]

        if action_type == 'CREATE_L2_AREA':
            ns_def.write_excel_meta(update_current_y_grid_tuple, ppt_meta_file, tmp_ws_name, '<<POSITION_FOLDER>>', 0, 0)

            self.root_width = self.all_slide_max_width
            self.root_hight = self.all_slide_max_hight

        self.output_diagram_path = self.outFileTxt_2_1.get()
        self.excel_file_path = ppt_meta_file
        self.worksheet_name = tmp_ws_name

        if action_type == 'CREATE_L2_AREA':
            self.l2_folder_name = tmp_folder_name
            self.all_tag_size_array = []
            print("Generating PPT diagram...")
            ns_ddx_figure.ns_ddx_figure_run.__init__(self)

        if action_type == 'CREATE_L2_AREA':
            #print(f"Removing temporary sheet: {tmp_ws_name}")
            ns_def.remove_excel_sheet(ppt_meta_file, tmp_ws_name)

        #self.perf_monitor.end('Area Processing Detail')

        return ([action_type, [self.all_slide_max_width, self.all_slide_max_hight]])


# ===================================================================
# CRITICAL: Cross-platform entry point guard
# This MUST be at the end of the file
# ===================================================================

if __name__ == '__main__':
    # Required for Windows multiprocessing
    # Prevents infinite process spawning
    print("\n" + "=" * 80)
    print("Note: This module should be imported by the main application.")
    print("If you see this message, the module is being run directly.")
    print("=" * 80)