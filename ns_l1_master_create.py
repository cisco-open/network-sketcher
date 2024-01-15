'''
SPDX-License-Identifier: Apache-2.0

Copyright 2023 Cisco Systems, Inc. and its affiliates

Licensed under the Apache License, Version 2.0 (the "License");
you may not use this file except in compliance with the License.
You may obtain a copy of the License at

http://www.apache.org/licenses/LICENSE-2.0

Unless required by applicable law or agreed to in writing, software
distributed under the License is distributed on an "AS IS" BASIS,
WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
See the License for the specific language governing permissions and
limitations under the License.
'''

from pptx import *
import sys, os, re
import numpy as np
import math
import ns_def
import openpyxl

class  ns_l1_master_create():
    def __init__(self):
        '''parameter '''
        flag_put_line_tag = True  #Write Tag Name in <<POSITION_LINE>>
        line_offset_value = 0.2  # in <<POSITION_LINE>> inches
        wp_roundness = 0.5 # in <<STYLE_SHAPE>> 0.0-1.0 * 100(%)
        wp_roundness_right_left = 0.2 # in <<STYLE_SHAPE>> 0.0-1.0 * 100(%)   Ver2.0 add fix wp roundness when right or left pattern
        shape_width_min = 0.5 # in <<STYLE_SHAPE>> inches
        shape_hight_min = 0.2 # in <<STYLE_SHAPE>> inches
        per_char_inchi = 0.16 # inches of per char count in shape
        color_wp = 'BLUE'  # in <<STYLE_SHAPE>>  ORANGE , BLUE , GREEN ,GRAY, empty is ''
        color_shape = 'GREEN' # in <<STYLE_SHAPE>>  ORANGE , BLUE , GREEN ,GRAY, empty is ''
        color_atmark = ''  # in <<STYLE_SHAPE>>  ORANGE , BLUE , GREEN ,GRAY, empty is ''
        tag_offet_inche = 0.02 # in <<POSITION_TAG>> inches of per char count in shape
        tag_name_prefix = 'GE 0/' # in <<POSITION_LINE>> TAG name prefix
        port_name_prefix = 'GigabitEthernet' # in <<POSITION_LINE>> port name prefix
        port_speed = 'Auto'  # in <<POSITION_LINE>> port speed
        port_duplex = 'Auto'  # in <<POSITION_LINE>> port speed
        port_type = '1000BASE-T'  # in <<POSITION_LINE>> port speed
        line_offset_ratio_wp = 3.0 # in <<POSITION_LINE>> increase ratio for WP , ver 1.12   reflect only up or down pattern
        shape_hight_offset_inches_wp = 0.1 # in <<STYLE_SHAPE>> increase hight inches for WP , ver 1.12
        self.shae_font_size = 6.0 #pt

        '''
        Create Excel file. And Write Template on Master_Data Sheet.
        '''
        ### Define File path ###
        if self.click_value == '92-3':
            # check : file is being opened
            if ns_def.check_file_open(str(self.inFileTxt_92_1.get()) + '__TMP__.xlsx') == True:
                return ()
            self.input_sketch_ppt = Presentation(self.inFileTxt_92_1.get())
            self.output_diagram_path = str(self.inFileTxt_92_1.get()) + '__TMP__.pptx'
            self.excel_file_path = self.inFileTxt_92_2.get()


        elif self.click_value == '1-4a':
            self.input_sketch_ppt = Presentation("./_tmp_tmp_tmp_.pptx")
            self.output_diagram_path = self.outFileTxt_1a_1.get()
            self.excel_file_path = self.outFileTxt_1a_2.get()
        elif self.click_value == '1-4b':
            self.input_sketch_ppt = Presentation("./_tmp_tmp_tmp_.pptx")
            #self.output_diagram_path = self.outFileTxt_1b_1.get()
            #self.excel_file_path = self.outFileTxt_1b_2.get()
        else:
            self.input_sketch_ppt = Presentation(self.inFileTxt_1_1.get())
            self.output_diagram_path = self.outFileTxt_1_1.get()
            self.excel_file_path = self.outFileTxt_1_2.get()

        if self.click_value == '92-3' and self.click_value_2nd != 'self.sub1_1_button_3':
            self.worksheet_name = 'Master_Data'
            wb = openpyxl.load_workbook(self.excel_file_path)
            del wb[self.worksheet_name]
            wb.create_sheet(index=0, title=self.worksheet_name )
            wb.save(self.excel_file_path)

        elif self.click_value == '92-3' and self.click_value_2nd == 'self.sub1_1_button_3':
            self.worksheet_name = 'Master_Data_tmp_'
            wb = openpyxl.load_workbook(self.excel_file_path)
            wb.create_sheet(index=0, title=self.worksheet_name )
            wb.save(self.excel_file_path)

        else:
            ### Create new data excel file
            self.worksheet_name = 'Master_Data'
            wb = openpyxl.Workbook()
            sheet = wb.active
            sheet.title = self.worksheet_name
            wb.save(self.excel_file_path)

        '''
        Write Template on Master_Data Sheet
        '''
        tmp_master_data_array = []
        tmp_master_data_array.append([1,['<<ROOT_FOLDER>>','Title Text','ratio_x(0.1-1.00)','ratio_y(0.1-1.00)','left(inches)','top(inches)','width(inches)','hight(inches)']])
        tmp_master_data_array.append([4, ['<<POSITION_FOLDER>>']])
        tmp_master_data_array.append([9, ['<<STYLE_FOLDER>>','Outline(YES/NO)','Text(NO/UP/DOWN)','Offset Upside Margin(<AUTO> or inches)','Offset Downside Margin(<AUTO> or inches)']])
        tmp_master_data_array.append([10, ['<DEFAULT>','YES','N/A','N/A','N/A']])
        tmp_master_data_array.append([11, ['<EMPTY>','NO','N/A','N/A','N/A']])
        tmp_master_data_array.append([15, ['<<POSITION_SHAPE>>']])
        tmp_master_data_array.append([19, ['<<STYLE_SHAPE>>','Width(Inches)','Hight(Inches)','Roundness(0.0-1.0)','Color(ORANGE/BLUE/GREEN/GRAY']])
        tmp_master_data_array.append([20, ['<DEFAULT>','0.5','0.3','0','N/A']])
        tmp_master_data_array.append([21, ['<EMPTY>','0.5','0.3','0','N/A']])
        tmp_master_data_array.append([25, ['<<POSITION_LINE>>']])
        tmp_master_data_array.append([26, ['From_Name','To_Name','From_Tag_Name','To_Tag_Name','From_Side(RIGHT/LEFT)','To_Side(RIGHT/LEFT)','Offset From_X (inches)','Offset From_Y (inches)','Offset To_X (inches)','Offset To_Y (inches)','Channel(inches)','Color(No) only)' \
            ,'From_Port_Name','From_Speed','From_Duplex','From_Port_Type','To__Port_Name','To_Speed','To_From_Duplex','To_From_Port_Type']])
        tmp_master_data_array.append([30, ['<<POSITION_TAG>>','Type(SHAPE/LINE)','Offset_SHAPE_X','Offset_SHAPE_Y','Offset_LINE(inches)','Adjust_LINE_Angle(YES/NO)']])
        tmp_master_data_array.append([31, ['<DEFAULT>','SHAPE','0','0','0.3','YES']])
        #print(tmp_master_data_array)

        template_master_data_tuple = {}
        template_master_data_tuple = ns_def.convert_array_to_tuple(tmp_master_data_array)

        #print('Create --- template_master_data_tuple---')
        #print(template_master_data_tuple)
        offset_row = 0
        offset_column = 0
        write_to_section = '_template_'
        ns_def.write_excel_meta(template_master_data_tuple, self.excel_file_path, self.worksheet_name, write_to_section, offset_row, offset_column)

        '''
        MAIN RUN
        '''
        # get number of slides
        for idx, slide in enumerate(self.input_sketch_ppt.slides):
            num_slide = idx + 1
        print('---- number of slides %d ----' %  num_slide)

        # get shapes folder waypoint
        tmp_shape_array = []
        master_shape_array = []
        folder_array = []
        wp_array = []
        tmp_line_array = []
        line_array = []
        exclude_folder_array = []
        tmp_shp = 0
        x_grid = 0
        y_grid = 0
        icon_num = 1
        for i, sld in enumerate(self.input_sketch_ppt.slides, start=1):
            #print(f'-- {i} --')
            current_folder_size = 0

            for shp in sld.shapes:
                ### GET Group(Icons),Picture and icon   # Add Ver 1.1
                if 'GROUP' in str(shp.shape_type)  or 'PICTURE' in str(shp.shape_type) :
                    tmp_shape_array.append(['@Icon@~'+str(icon_num)+'~', shp.left, shp.top, shp.width, shp.height, shp.rotation, i])
                    tmp_shp += 1
                    icon_num += 1

                ### GET Line ( connector )
                if 'LINE' in str(shp.shape_type):
                    tmp_line_array.append([shp.begin_x, shp.begin_y, shp.end_x, shp.end_y, i])
                    line_array.append([shp.begin_x, shp.begin_y, shp.end_x, shp.end_y, i])

                ### GET Shape
                if 'AUTO_SHAPE' in str(shp.shape_type) and str(shp.text) != '':
                    # check updated device '/n' and append array     ver 1.1
                    if '\n' in str(shp.text):
                        idx = str(shp.text).find('\n')
                        self.updated_name_array.append([str(shp.text)[:idx],str(shp.text)[idx + 1:] ])
                        tmp_shape_array.append([str(shp.text)[idx + 1:], shp.left, shp.top, shp.width, shp.height, shp.rotation, i ])
                    else:
                        tmp_shape_array.append([shp.text, shp.left, shp.top, shp.width, shp.height, shp.rotation, i])

                    ### Add multiple _AIR_   # Add Ver 1.11
                    if '_AIR_' in str(shp.text):
                        tmp_air = str(shp.text).replace('_AIR_','')
                        if tmp_air.isdigit() == True:
                            #print('### Add multiple _AIR_###')
                            tmp_air = int(tmp_air) - 1
                            for tmp_air_num in range(tmp_air):
                                tmp_shape_array.append([shp.text, shp.left + tmp_air_num , shp.top, shp.width, shp.height, shp.rotation, i])

                    tmp_shp += 1

                    '''Elect Folder array '''
                    if current_folder_size < (shp.width + shp.height):
                        current_folder_size = shp.width + shp.height
                        current_folder_array = tmp_shape_array[tmp_shp-1]

            ### add root folder if there is not root folder. Add ver 1.11
            tmp_count_include_shape = 0
            for tmp_check_shape in tmp_shape_array:
                if current_folder_array[1] < tmp_check_shape[1] and current_folder_array[2] < tmp_check_shape[2] and \
                        (current_folder_array[1] + current_folder_array[3]) > (tmp_check_shape[1] + tmp_check_shape[3]) and \
                        (current_folder_array[2] + current_folder_array[4]) > (tmp_check_shape[2] + tmp_check_shape[4]):
                    tmp_count_include_shape += 1
                    #print(current_folder_array,tmp_check_shape)

            #print('---tmp_count_include_shape---  \n', tmp_count_include_shape)

            if tmp_count_include_shape == 0:
                tmp_shape_array.append(['_tmp_', 1, 1, 999999999999, 999999999999, 0.0, i])
                current_folder_array = ['_tmp_','_tmp_', 1, 1, 999999999999, 999999999999, 0.0, i]

            '''### duplicate folder name check ###'''  #update ver 1.1
            duplicate_num = 1

            if str(current_folder_array[0])[-1] == "~":
                p = r'\~.*\~'
                r = re.findall(p, str(current_folder_array[0]))
                r = re.sub("\[\'\~", "", str(r))
                r = re.sub("\~\'\]", "", str(r))
                if r.isdecimal() == True and current_folder_array[0] not in exclude_folder_array:
                    exclude_folder_array.append(current_folder_array[0])
            #print(exclude_folder_array)

            for tmp_tmp_array in folder_array:
                if current_folder_array[0] == tmp_tmp_array[0]:
                    flag_num = True
                    while flag_num == True:
                        if str(current_folder_array[0]) + '~' + str(duplicate_num) + '~' in exclude_folder_array and '~' not in str(current_folder_array[0]):
                            duplicate_num += 1

                        else:
                            #print('current_folder_array[0]   ',current_folder_array[0])
                            flag_num = False
                            current_folder_array[0]= str(current_folder_array[0]) + '~' + str(duplicate_num) + '~'
                            exclude_folder_array.append(current_folder_array[0])
                            duplicate_num += 1

            folder_array.append(current_folder_array)

            ### insert Folder name ####
            for tmp_shape in tmp_shape_array:
                tmp_shape.insert(0,current_folder_array[0])
                master_shape_array.append(tmp_shape)
            tmp_shape_array = []
            tmp_line_array = []
            tmp_shp = 0

        ### Elect Way Point array ###
        for tmp_folder in folder_array:
            for tmp_master in master_shape_array:
                if tmp_folder[0] == tmp_master[0]:
                    if tmp_master[3] < tmp_folder[3]:
                        wp_array.append([tmp_master[0], tmp_master[1], 'UP'])
                    elif (tmp_master[3]+tmp_master[5]) > (tmp_folder[3]+tmp_folder[5]):
                        wp_array.append([tmp_master[0], tmp_master[1], 'DOWN'])
                    elif tmp_master[2] < tmp_folder[2]:
                        wp_array.append([tmp_master[0], tmp_master[1], 'LEFT'])
                    elif (tmp_master[2]+tmp_master[4]) > (tmp_folder[2]+tmp_folder[4]):
                        wp_array.append([tmp_master[0], tmp_master[1], 'RIGHT'])
        #print('master_shape_array    ',master_shape_array)
        ### MAKE Shape array ###
        pre_shape_array = []
        for tmp_shape in master_shape_array:
            tmp_include_bit = False
            # exclude folder and way point
            for tmp_folder in folder_array:
                if tmp_shape[1] == tmp_folder[1]:
                    tmp_include_bit = True
            for tmp_wp in wp_array:
                if tmp_shape[1] == tmp_wp[1]:
                    tmp_include_bit = True
            if  tmp_include_bit == False:
                pre_shape_array.append(tmp_shape)

        '''duplicate shape name check''' # update ver 1.1
        shape_array = []
        duplicate_num = 1
        #print('---pre_shape_array--- \n',pre_shape_array)

        # check include '~' and get exclude_pre_shape_array
        exclude_pre_shape_array = []
        for tmp_exclude_array in pre_shape_array:
            if str(tmp_exclude_array[1])[-1] == "~" and "~" in str(tmp_exclude_array[1]) :
                p = r'\~.*\~'
                r = re.findall(p, str(tmp_exclude_array[1]))
                r = re.sub("\[\'\~", "", str(r))
                r = re.sub("\~\'\]", "", str(r))
                if r.isdecimal() == True:
                    exclude_pre_shape_array.append(tmp_exclude_array[1])
        #print(exclude_pre_shape_array)

        # check duplicate
        for tmp_array in pre_shape_array:
            tmp_count = 0
            last_array = tmp_array
            for tmp_tmp_array in pre_shape_array:
                if tmp_array[1] == tmp_tmp_array[1] and "~" not in str(tmp_array[1]):
                    tmp_count += 1
                    if tmp_count >= 2:
                        flag_num = True
                        while flag_num == True:
                            if str(tmp_array[1]) + '~' + str(duplicate_num) + '~' in exclude_pre_shape_array:
                                duplicate_num += 1
                            else:
                                flag_num = False
                                tmp_array[1] = str(tmp_array[1]) + '~' + str(duplicate_num) + '~'
                                exclude_pre_shape_array.append(tmp_array[1])
                                duplicate_num += 1
                                last_array = tmp_array
            duplicate_num = 1
            shape_array.append(last_array)

        print('-- master_shape_array --')
        #print(master_shape_array)

        print('-- shape_array --')
        #print(shape_array)

        print('-- line_array --')
        #print(line_array)

        print('-- folder_array --')
        #print(folder_array)

        print('-- wp_array --')
        #print(wp_array)

        ### Add Way Point to Folder array ###
        new_wp_array = wp_array
        new_folder_wp_array =[]
        tmp_up_flag = ''
        tmp_down_flag = ''
        tmp_right_flag = ''
        tmp_left_flag = ''

        for tmp_folder in folder_array:
            for tmp_wp_array in new_wp_array:
                if tmp_folder[0] == tmp_wp_array[0]:
                    if tmp_wp_array[2] == 'UP':
                        tmp_up_flag = tmp_wp_array[1]
                    elif tmp_wp_array[2] == 'DOWN':
                        tmp_down_flag = tmp_wp_array[1]
                    elif tmp_wp_array[2] == 'LEFT':
                        tmp_left_flag = tmp_wp_array[1]
                    elif tmp_wp_array[2] == 'RIGHT':
                        tmp_right_flag = tmp_wp_array[1]


            new_folder_wp_array.append([tmp_folder[0],tmp_up_flag,tmp_down_flag,tmp_left_flag,tmp_right_flag,x_grid, y_grid])
            x_grid += 1
            tmp_up_flag = ''
            tmp_down_flag = ''
            tmp_right_flag = ''
            tmp_left_flag = ''

        print('-- new_folder_wp_array --')
        #print(new_folder_wp_array)

        ''' add connected way point'''
        up_flag = ''
        down_flag = ''
        left_flag = ''
        right_flag = ''
        connect_folder_wp_array = []
        for tmp_folder in new_folder_wp_array:
            for tmp_tmp_folder in new_folder_wp_array:
                if tmp_folder[1] == tmp_tmp_folder[2]:
                    up_flag = tmp_folder[1]
                if tmp_folder[2] == tmp_tmp_folder[1]:
                    down_flag = tmp_folder[2]
                if tmp_folder[3] == tmp_tmp_folder[4]:
                    left_flag = tmp_folder[3]
                if tmp_folder[4] == tmp_tmp_folder[3]:
                    right_flag = tmp_folder[4]

            tmp_array= []

            tmp_array.append(tmp_folder[0])
            if up_flag != '':
                tmp_array.append(up_flag)
            else:
                tmp_array.append('')
            if down_flag != '':
                tmp_array.append(down_flag)
            else:
                tmp_array.append('')
            if left_flag != '':
                tmp_array.append(left_flag)
            else:
                tmp_array.append('')
            if right_flag != '':
                tmp_array.append(right_flag)
            else:
                tmp_array.append('')
            up_flag = ''
            down_flag = ''
            left_flag = ''
            right_flag = ''

            connect_folder_wp_array.append(tmp_array)
        print('--connect_folder_wp_array--')
        #print(connect_folder_wp_array)

        ''' grouping x grid folder'''
        ### pre Kyuusai ###
        kyusai_a = []
        kyusai_b = []
        for tmp_folder in connect_folder_wp_array:
            if str(tmp_folder[3]) != '' and str(tmp_folder[4]) == '':
                kyusai_b.append(tmp_folder)
            else:
                kyusai_a.append(tmp_folder)

        for tmp_folder in kyusai_b:
            kyusai_a.append(tmp_folder)
        connect_folder_wp_array = kyusai_a

        tmp_array = []
        group_folder_wp_array = []
        for tmp_folder in connect_folder_wp_array:
            if tmp_folder[4] != '':
                for tmp_tmp_folder in connect_folder_wp_array:
                    if tmp_folder[4] == tmp_tmp_folder[3]:
                        if str(tmp_folder) in str(group_folder_wp_array) or str(tmp_tmp_folder) in str(group_folder_wp_array):
                            if str(tmp_folder) in str(group_folder_wp_array):
                                new_tmp_folder = tmp_folder
                            else:
                                new_tmp_folder = tmp_tmp_folder
                            ii = len(group_folder_wp_array)
                            for iii in range(ii):
                                if new_tmp_folder in group_folder_wp_array[iii]:
                                    iiii = iii
                            n = len(group_folder_wp_array[iiii])
                            for nn in range(n-1):
                                if new_tmp_folder == group_folder_wp_array[iiii][nn]:
                                    if nn == 0:
                                        if str(tmp_folder) not in str(group_folder_wp_array[iiii]):
                                            group_folder_wp_array[iiii].insert(0,tmp_folder)
                                    else:
                                        if str(tmp_tmp_folder) not in str(group_folder_wp_array[iiii]):
                                            group_folder_wp_array[iiii].append(tmp_tmp_folder)

                        else:
                            tmp_array.append(tmp_folder)
                            tmp_array.append(tmp_tmp_folder)

                if tmp_array != []:
                    group_folder_wp_array.append(tmp_array)
                tmp_array = []
            elif tmp_folder[3] == '':
                group_folder_wp_array.append(tmp_folder)
            else:
                ### Kyuusai-2 ###
                if str(tmp_folder) not in str(group_folder_wp_array):
                    #print('---  Kyusai-2  ---  ', tmp_folder,group_folder_wp_array)
                    tmp_i = 0
                    for tmp_folder2 in group_folder_wp_array:
                        if str(tmp_folder[3]+'\''+']]') in str(tmp_folder2):
                            break
                        tmp_i += 1
                    group_folder_wp_array[tmp_i].append(tmp_folder)

        print('--group_folder_wp_array--')
        #print(group_folder_wp_array)

        ''' set y grid folder'''
        num_folder_group = len(group_folder_wp_array)
        y_num_array = []
        for i in range(num_folder_group):
            y_num_array.append(1)

        #add [] to single folder group
        new_group_folder_array = []
        for tmp_num_folder in group_folder_wp_array:
            if ']]' not in str(tmp_num_folder):
                #single folder
                test_folder = []
                test_folder.append(tmp_num_folder)
                new_group_folder_array.append(test_folder)
            else:
                new_group_folder_array.append(tmp_num_folder)
        print('--new_group_folder_array  ,   add [] to single folder-')
        #print(new_group_folder_array)

        # change Y grid of y_num_array
        print(' -- move each group to downside - ')
        change_bit = True
        first_change_bit = True
        max_y_grid = 1
        while change_bit == True:
            pre_y_num_array = y_num_array
            for b in range(len(new_group_folder_array)):
                for i in range(len(new_group_folder_array[b])):
                    for bb in range(len(new_group_folder_array)):
                        for ii in range(len(new_group_folder_array[bb])):
                            #print(new_group_folder_array[b][i],new_group_folder_array[bb][ii],new_group_folder_array[b][i][2],new_group_folder_array[bb][ii][1],new_group_folder_array[b][i][2],y_num_array[b] , y_num_array[bb])
                            if new_group_folder_array[b][i] != new_group_folder_array[bb][ii] and new_group_folder_array[b][i][2] == new_group_folder_array[bb][ii][1] \
                                    and new_group_folder_array[b][i][2] != '' and y_num_array[b] >= y_num_array[bb]:
                                y_num_array[bb] += 1
                                if y_num_array[bb] > max_y_grid:
                                    max_y_grid = y_num_array[bb]
                                first_change_bit = True
                                #print(y_num_array)

                            '''a group which does not have UP and Down move to max downside'''
                            if new_group_folder_array[bb][ii][1] == '' and new_group_folder_array[bb][ii][2] == '' and max_y_grid >= y_num_array[bb]:
                                y_num_array[bb] = max_y_grid + 1
                                first_change_bit = True
                                #print(y_num_array)

            if pre_y_num_array == y_num_array:
                if first_change_bit == True:
                    first_change_bit = False
                else:
                    change_bit = False

        '''shapes convert to excel'''
        y_num_bit = True
        y_num_start = 1
        update_group_folder_array = []

        while y_num_bit == True:
            y_num_current = -1
            for tmp_yum in y_num_array:
                y_num_current +=1
                if tmp_yum == y_num_start or len(new_group_folder_array) == 1:
                    update_group_folder_array.append([new_group_folder_array[y_num_current] , tmp_yum])
            y_num_start +=1
            if y_num_start-1 == len(y_num_array):
                y_num_bit = False

        print(' -- Update group folder array for Excel position -- ')
        #print(update_group_folder_array)

        '''
        make excel format
        '''
        # merge merge_group_folder_array <-- update_group_folder_array
        folder_name_of_wp = [] # for put wp shapes in own wp folder
        merge_group_folder_array = update_group_folder_array
        for i in range(len(update_group_folder_array)):
            for ii in range(len(update_group_folder_array[i][0])):
                for tmp_new_folder_wp_array in new_folder_wp_array:
                    if update_group_folder_array[i][0][ii][0] == tmp_new_folder_wp_array[0]:
                        if tmp_new_folder_wp_array[1] != '':
                            merge_group_folder_array[i][0][ii][1] = tmp_new_folder_wp_array[1] + '_wp_'
                        if tmp_new_folder_wp_array[2] != '':
                            merge_group_folder_array[i][0][ii][2] = tmp_new_folder_wp_array[2] + '_wp_'
                        if tmp_new_folder_wp_array[3] != '':
                            merge_group_folder_array[i][0][ii][3] = tmp_new_folder_wp_array[3] + '_wp_'
                        if tmp_new_folder_wp_array[4] != '':
                            merge_group_folder_array[i][0][ii][4] = tmp_new_folder_wp_array[4] + '_wp_'
                        break
        print(' -- merge_group_folder_array -- ')
        #print(merge_group_folder_array)

        #Line up shapes
        row = 2       # Excel row
        column = 2    # Excel column
        offset_row = 0
        offset_column =0
        master_folder_tuple = {} # After this process, write the meta data to excel
        max_num_y_grid = merge_group_folder_array[-1][1]
        current_y_grid = 0

        #kyusai y grid == 1

        pre_y_num = merge_group_folder_array[0][1]
        flag_sigle_y_num = True
        for cur_merge_group_folder_array in merge_group_folder_array:
            if pre_y_num != cur_merge_group_folder_array[1]:
                flag_sigle_y_num = False
                break

        if flag_sigle_y_num == True:
            row -= 4

        ''' make master_folder_tuple'''
        flag_first_shape = True
        for i in range(1,max_num_y_grid+1):
            print('---- Y Grid %d ----  ' % i)
            flag_up_wp = False
            flag_down_wp = False
            flag_left_wp = False
            current_y_grid += 1

            # Check up flag and left and down flag
            for tmp_array in merge_group_folder_array:
                if tmp_array[1] == i:
                    for tmp_tmp_array in tmp_array[0]:
                        #check up or left side WP
                        if tmp_tmp_array[1] != '':
                            flag_up_wp = True
                        if tmp_tmp_array[2] != '':
                            flag_down_wp = True
                        if column == 2 and tmp_tmp_array[3] != '':
                            flag_left_wp = True

                        # kyusai Y
                        if flag_up_wp == False and i == 1:
                            flag_first_y_up = True
                        else:
                            flag_first_y_up = False

            #print('Flag | UP DOWN LEFT |',flag_up_wp,flag_down_wp,flag_left_wp )

            ### adjust pre down flag ###
            if i >=2:
                if pre_flag_down_wp == False and flag_up_wp == False and flag_sigle_y_num == False:
                    ## insert empty folder between folders that do not connect
                    master_folder_tuple[row + 1 -2 , 1] = '<SET_WIDTH>'
                    master_folder_tuple[row + 1 -2 , 2] = 10
                    master_folder_tuple[row + 2 -2 , 1] = 1

                # kyusai y grid == 1
                if flag_down_wp == False and flag_up_wp == False and flag_sigle_y_num == True:
                    row -= 1
                if flag_down_wp == True and flag_up_wp == False and flag_sigle_y_num == True:
                    row -= 1

            # make excel meta
            flag_row_column_2_2 = False   ### bug fix 001 ###
            #flag_row_column_onetime = False  ### bug fix 001 ###
            flag_check_leftup_wp = False ### bug fix 001 ###
            flag_check_leftup_wp_end = False ### bug fix 001 ###

            for tmp_array in merge_group_folder_array:
                write_up_wp = False

                if tmp_array[1] == current_y_grid:
                    for tmp_tmp_array in tmp_array[0]:
                        # check the same wp on upside
                        if flag_row_column_2_2 == True:             ### bug fix 001 ###
                            master_folder_tuple[row - 1, column] = 0.999 ### bug fix 001 ###

                        if flag_up_wp == True and tmp_tmp_array[1] != '':
                            tmp_i = 0
                            flag_write_up_wp = False
                            for num_wp_up in master_folder_tuple:
                                if num_wp_up[0] == row:
                                    tmp_i += 1
                                    if tmp_tmp_array[1] == master_folder_tuple[num_wp_up] and tmp_tmp_array[1] != '':
                                        flag_write_up_wp = True

                            if flag_write_up_wp == True:
                                write_up_wp = False
                            else:
                                write_up_wp = True

                            if write_up_wp == True and flag_check_leftup_wp == False:

                                ### bug fix 001
                                if row == 2 and column == 6:
                                    column -= 1
                                    flag_check_leftup_wp = True

                                master_folder_tuple[row, column] = tmp_tmp_array[1]
                                master_folder_tuple[row - 1, column] = 10
                                folder_name_of_wp.append(tmp_tmp_array[1])
                                #print(row, column,tmp_tmp_array[1])


                                ### bug fix 001 ###
                                if row == 2 and column == 2:
                                    flag_row_column_2_2 = True
                                    master_folder_tuple[row, column - 1] = 10

                        # check previous right side
                        flag_previous_right_side = False
                        for check_right in master_folder_tuple:
                            if check_right == (row + 2, column):
                                flag_previous_right_side = True

                        if flag_previous_right_side == True:
                            add_x = 1
                        else:
                            add_x = 0

                        ### check left side and wirte
                        if tmp_tmp_array[3] != '' or column == 2:
                            if write_up_wp == True and tmp_tmp_array[3] != '' and column != 3:
                                if flag_check_leftup_wp == False and (tmp_tmp_array[1] == '' or tmp_tmp_array[3] == ''):  ### bug fix 001
                                    column -= 1                                                                             ### bug fix 001
                                master_folder_tuple[row + 1, 1] = '<SET_WIDTH>'
                                master_folder_tuple[row + 2, 1] = 10
                                master_folder_tuple[row + 1, column + add_x] = 10
                                master_folder_tuple[row + 1, column + 1 + add_x] = 10
                                master_folder_tuple[row + 2, column + add_x] = tmp_tmp_array[3]
                                master_folder_tuple[row + 2, column + 1 + add_x] = tmp_tmp_array[0]
                                folder_name_of_wp.append(tmp_tmp_array[3])

                                if flag_row_column_2_2 == True and row == 2 and column == 2:  ### bug fix 001 ###
                                    master_folder_tuple[row - 1, column + 1] = 0.999  ### bug fix 001 ###

                            # kyusai y grid == 1
                            elif write_up_wp == False and tmp_tmp_array[3] != '' and column == 2:
                                master_folder_tuple[row + 1, 1] = '<SET_WIDTH>'
                                master_folder_tuple[row + 2, 1] = 10
                                master_folder_tuple[row + 1, column + add_x] = 10
                                master_folder_tuple[row + 1, column + 1 + add_x] = 10
                                master_folder_tuple[row + 2, column + add_x] = tmp_tmp_array[3]
                                master_folder_tuple[row + 2, column + 1 + add_x] = tmp_tmp_array[0]
                                folder_name_of_wp.append(tmp_tmp_array[3])

                            # kyusai y grid == 1
                            elif write_up_wp == False and flag_down_wp == True and tmp_tmp_array[3] == '' and column == 2 and flag_first_shape == True and flag_first_y_up == True:
                                if flag_first_shape == True:
                                    row -= 1
                                    flag_first_shape = False
                                master_folder_tuple[row + 1, 1] = '<SET_WIDTH>'
                                master_folder_tuple[row + 2, 1] = 10
                                master_folder_tuple[row + 1, column + add_x] = 10
                                master_folder_tuple[row + 2, column + add_x] = tmp_tmp_array[3]
                                master_folder_tuple[row + 2, column + add_x] = tmp_tmp_array[0]
                                folder_name_of_wp.append(tmp_tmp_array[3])

                            else:
                                master_folder_tuple[row + 1, 1] = '<SET_WIDTH>'
                                master_folder_tuple[row + 2, 1] = 10
                                master_folder_tuple[row + 1, column + add_x] = 10
                                master_folder_tuple[row + 2, column + add_x] = tmp_tmp_array[0]

                                if flag_row_column_2_2 == True and column != 2:  ### bug fix 001 ###
                                    master_folder_tuple[row - 1, column + add_x] = 0.999  ### bug fix 001 ###


                        else:
                            master_folder_tuple[row + 1, column + add_x] = 1
                            master_folder_tuple[row + 2, column + add_x] = ''
                            master_folder_tuple[row + 1, column + add_x + 1] = 10
                            master_folder_tuple[row + 2, column + add_x + 1] = tmp_tmp_array[0]

                        #if there is right wp
                        if tmp_tmp_array[4] != '':
                            if tmp_tmp_array[3] != '':
                                master_folder_tuple[row + 1, column + 2] = 10
                                master_folder_tuple[row + 2, column + 2] = tmp_tmp_array[4]
                                folder_name_of_wp.append(tmp_tmp_array[4])
                            else:
                                master_folder_tuple[row + 1, column + 1] = 10
                                master_folder_tuple[row + 2, column + 1] = tmp_tmp_array[4]
                                folder_name_of_wp.append(tmp_tmp_array[4])

                        #if there is down wp
                        if flag_down_wp == True and tmp_tmp_array[2] != '':
                            #check the same wp on downside
                            tmp_i = 0
                            write_down_wp = True
                            for num_wp_down in master_folder_tuple:
                                if num_wp_down[0] == row +4:
                                    tmp_i += 1
                                    if tmp_tmp_array[2] == master_folder_tuple[num_wp_down] and tmp_tmp_array[2] != '':
                                        write_down_wp = False

                            if write_down_wp == True:
                                master_folder_tuple[row + 4, column] = tmp_tmp_array[2]
                                master_folder_tuple[row + 3, column] = 10
                                master_folder_tuple[row + 3, 1] = '<SET_WIDTH>'
                                master_folder_tuple[row + 4, 1] = 10
                                folder_name_of_wp.append(tmp_tmp_array[2])
                                #print(row + 4, column, tmp_tmp_array[2])

                        #column move to right side
                        if tmp_tmp_array[3] != '':
                            column += 1
                        if tmp_tmp_array[3] == '' and flag_previous_right_side == True:
                            column += 1
                        column += 1

            row += 4
            column = 2

            ## previsious down flag ##
            pre_flag_down_wp = flag_down_wp

        print(' -- master_folder_tuple <<POSITION_FOLDER>> -- ')
        #print(master_folder_tuple)

        '''bug fix 001,  put 0.999 on empty cell(width)'''
        master_folder_tuple[1, 1] = '<<POSITION_FOLDER>>'
        bug_fix_master_folder_tuple = master_folder_tuple
        tmp_add_tuple = {}
        for tmp_bug_fix_master_folder_tuple in bug_fix_master_folder_tuple:
            if tmp_bug_fix_master_folder_tuple[1] == 1 and (bug_fix_master_folder_tuple[tmp_bug_fix_master_folder_tuple] == '<<POSITION_FOLDER>>' or bug_fix_master_folder_tuple[tmp_bug_fix_master_folder_tuple] == '<SET_WIDTH>'):
                #print(bug_fix_master_folder_tuple[tmp_bug_fix_master_folder_tuple])
                tmp_num_max_bug_fix_column = 1
                for tmp_tmp_bug_fix_master_folder_tuple in bug_fix_master_folder_tuple:
                    if tmp_tmp_bug_fix_master_folder_tuple[0] == tmp_bug_fix_master_folder_tuple[0]:
                        if tmp_num_max_bug_fix_column < tmp_tmp_bug_fix_master_folder_tuple[1]:
                            tmp_num_max_bug_fix_column = tmp_tmp_bug_fix_master_folder_tuple[1]

                #print(tmp_num_max_bug_fix_column)
                for tmp_num in range(2, tmp_num_max_bug_fix_column):
                    flag_tuple_exist = False
                    for tmp_tmp_bug_fix_master_folder_tuple in bug_fix_master_folder_tuple:
                        if tmp_tmp_bug_fix_master_folder_tuple[0] == tmp_bug_fix_master_folder_tuple[0] and tmp_tmp_bug_fix_master_folder_tuple[1] == tmp_num:
                            flag_tuple_exist = True

                    if flag_tuple_exist == False:
                        #print('### put 0.999 ###')
                        #print(tmp_bug_fix_master_folder_tuple[0],tmp_num)
                        tmp_add_tuple[tmp_bug_fix_master_folder_tuple[0],tmp_num] = 0.999

        for tmp_tmp_add_tuple in tmp_add_tuple:
            master_folder_tuple[tmp_tmp_add_tuple] = tmp_add_tuple[tmp_tmp_add_tuple]

        ### sort master_folder_tuple ###
        sort_master_folder_tuple = sorted(master_folder_tuple)
        #print('--- sort_master_folder_tuple ---  \n',sort_master_folder_tuple)
        i_master_folder_tuple = {}
        flag_exist_2_1 = False
        for tmp_sort_master_folder_tuple in sort_master_folder_tuple:
            if master_folder_tuple[tmp_sort_master_folder_tuple] != '<<POSITION_FOLDER>>':
                i_master_folder_tuple[tmp_sort_master_folder_tuple] = master_folder_tuple[tmp_sort_master_folder_tuple]

            if tmp_sort_master_folder_tuple[0] == 2 and tmp_sort_master_folder_tuple[1] == 1:
                flag_exist_2_1 = True

        ### kyusai bug fix 001 , add 10 to 2,1 in master_folder_tuple
        if flag_exist_2_1 == False:
            i_master_folder_tuple[2,1] = 10
        #print('--- i_master_folder_tuple ---   \n', i_master_folder_tuple)

        ### remove 0.999 in i_master_folder_tuple ###
        tmp_change_column = 1
        tmp_i_current_row = 1
        flag_previous_three_nine = False
        master_folder_tuple = {}
        used_master_folder_tuple = {}
        tmp_offset_column = 0

        for tmp_i_master_folder_tuple in i_master_folder_tuple:
            if tmp_i_master_folder_tuple[0] == 1 or i_master_folder_tuple[tmp_i_master_folder_tuple[0],1] == '<SET_WIDTH>':
                if tmp_i_current_row != tmp_i_master_folder_tuple[0]:
                    tmp_change_column = 1
                    flag_previous_three_nine = False
                    tmp_offset_column = 0

                tmp_i_current_row = tmp_i_master_folder_tuple[0]

                flag_used_folder = False
                for tmp_used_master_folder_tuple in used_master_folder_tuple:
                    if tmp_i_master_folder_tuple == tmp_used_master_folder_tuple:
                        flag_used_folder = True

                if i_master_folder_tuple[tmp_i_master_folder_tuple] == 0.999 and flag_previous_three_nine == False:
                    tmp_change_column = tmp_i_master_folder_tuple[1]
                    flag_previous_three_nine = True
                elif i_master_folder_tuple[tmp_i_master_folder_tuple] != 0.999 and flag_previous_three_nine == True and flag_used_folder == False:
                    master_folder_tuple[tmp_i_master_folder_tuple[0],tmp_change_column + tmp_offset_column] = i_master_folder_tuple[tmp_i_master_folder_tuple]
                    master_folder_tuple[tmp_i_master_folder_tuple[0] + 1, tmp_change_column + tmp_offset_column] = i_master_folder_tuple[tmp_i_master_folder_tuple[0] + 1 , tmp_i_master_folder_tuple[1]]
                    used_master_folder_tuple[tmp_i_master_folder_tuple] = i_master_folder_tuple[tmp_i_master_folder_tuple]
                    used_master_folder_tuple[tmp_i_master_folder_tuple[0] + 1,tmp_i_master_folder_tuple[1]] = i_master_folder_tuple[tmp_i_master_folder_tuple[0] + 1,tmp_i_master_folder_tuple[1]]
                    tmp_change_column = 1
                    flag_previous_three_nine = False
                    tmp_offset_column = tmp_offset_column - ( tmp_i_master_folder_tuple[1] - tmp_change_column)
                    #print(tmp_i_master_folder_tuple,used_master_folder_tuple)
                elif flag_used_folder == False and i_master_folder_tuple[tmp_i_master_folder_tuple] != 0.999:
                    master_folder_tuple[tmp_i_master_folder_tuple[0] + tmp_offset_column ,tmp_i_master_folder_tuple[1]] = i_master_folder_tuple[tmp_i_master_folder_tuple]

                    flag_exist_plus_one_row = False
                    for tmp_ii_master_folder_tuple in i_master_folder_tuple:
                        if tmp_ii_master_folder_tuple[0] == tmp_i_master_folder_tuple[0] + 1 and tmp_ii_master_folder_tuple[1] == tmp_i_master_folder_tuple[1]:
                            flag_exist_plus_one_row = True

                    if flag_exist_plus_one_row == True:
                        master_folder_tuple[tmp_i_master_folder_tuple[0] + 1 + tmp_offset_column, tmp_i_master_folder_tuple[1]] = i_master_folder_tuple[tmp_i_master_folder_tuple[0] + 1,tmp_i_master_folder_tuple[1]]

        ### kyusai bug fix 001 , add 10 to 2,1 in master_folder_tuple
        flag_exist_2_1_a = False
        for iii_master_folder_tuple in master_folder_tuple:
            if iii_master_folder_tuple[0] == 2 and iii_master_folder_tuple[1] == 1:
                flag_exist_2_1_a = True

        if flag_exist_2_1_a == False:
            master_folder_tuple[2,1] = 10

        print('--- master_folder_tuple <<POSITION_FOLDER>> -- ')
        #print(master_folder_tuple)

        '''For Debug, write excel meta'''
        #write_to_section = '<<POSITION_FOLDER>>'
        #ns_def.write_excel_meta(master_folder_tuple, self.excel_file_path, self.worksheet_name, write_to_section,offset_row, offset_column)
        #exit()

        '''
        Write <<STYLE_FOLDER>>
        '''
        master_style_folder_meta = {} # After this process, write the style folder meta data to excel
        style_folder_row = 1
        for tmp_array in folder_array:
            style_folder_row += 1
            master_style_folder_meta[style_folder_row, 1] = tmp_array[0]
            master_style_folder_meta[style_folder_row, 2] = 'YES'
            master_style_folder_meta[style_folder_row, 3] = 'UP'
            master_style_folder_meta[style_folder_row, 4] = '<AUTO>'
            master_style_folder_meta[style_folder_row, 5] = '<AUTO>'

        for meta in master_folder_tuple:
            flag_wp_array = False
            for tmp_wp in wp_array:
                if master_folder_tuple[meta[0], meta[1]] == tmp_wp[1] + '_wp_':
                    flag_wp_array = True

            if flag_wp_array == True:
                style_folder_row += 1
                master_style_folder_meta[style_folder_row, 1] = master_folder_tuple[meta[0], meta[1]]
                master_style_folder_meta[style_folder_row, 2] = 'NO'
                master_style_folder_meta[style_folder_row, 3] = 'NO'
                master_style_folder_meta[style_folder_row, 4] = '<AUTO>'
                master_style_folder_meta[style_folder_row, 5] = '<AUTO>'

        print('---- master_style_folder_meta ----')
        #print(master_style_folder_meta)

        write_to_section = '<<STYLE_FOLDER>>'
        offset_row = 2
        ns_def.write_excel_meta(master_style_folder_meta, self.excel_file_path, self.worksheet_name, write_to_section, offset_row, offset_column)
        offset_row = 0

        '''
        Write <<POSITION_SHAPE>>
        '''
        ### GET Shapes of each Folder
        start_row = 2
        master_shape_tuple = {}
        for each_folder_array in folder_array:
            flag_target_folder = False
            current_shape_array = []
            for each_shape_array in shape_array:
                if each_folder_array[0] == each_shape_array[0]:
                    flag_target_folder = True
                    current_shape_array.append(each_shape_array)

                if flag_target_folder == True and (each_folder_array[0] != each_shape_array[0] or each_shape_array[1] == shape_array[-1][1]):
                    current_shape_tuple = ns_def.return_shape_tuple(current_shape_array,start_row)

                    for tmp_tuple in current_shape_tuple:
                        master_shape_tuple[tmp_tuple] = current_shape_tuple[tmp_tuple]
                        last_row_num = tmp_tuple[0]
                    start_row = last_row_num + 1
                    break

        '''
        Write <<POSITION_SHAPE>> for WayPoint
        '''
        print('---- folder_name_of_wp ----')
        #print(folder_name_of_wp)
        pre_folder_name_of_wp = []
        for replace_folder_name_of_wp in folder_name_of_wp:
            pre_folder_name_of_wp.append(replace_folder_name_of_wp.replace('_wp_', ''))
        folder_name_of_wp = pre_folder_name_of_wp

        ### Check multiple WPs in a same folder
        multi_wp = []
        for tmp_wp in wp_array:
            if tmp_wp[1] not in folder_name_of_wp and tmp_wp[1] not in str(multi_wp):
                for tmp_tmp_wp in wp_array:
                    if tmp_wp[0] == tmp_tmp_wp[0] and tmp_wp[2] == tmp_tmp_wp[2] and tmp_wp[1] != tmp_tmp_wp[1]:
                        multi_wp.append([tmp_wp[1], tmp_wp[2],tmp_tmp_wp[1]])

        print('---- multi_wp ----')
        #print(multi_wp)

        for tmp_folder_name_of_wp in folder_name_of_wp:
            flag_multi_wp = False
            tmp_sort_hight = [tmp_folder_name_of_wp]

            for tmp_multi_wp in multi_wp:
                if tmp_folder_name_of_wp == tmp_multi_wp[2]:
                    flag_multi_wp = True
                    tmp_sort_hight.append(tmp_multi_wp[0])

            if flag_multi_wp == False:
                master_shape_tuple[start_row, 1] = tmp_folder_name_of_wp + '_wp_'
                master_shape_tuple[start_row, 2] = tmp_folder_name_of_wp
                master_shape_tuple[start_row, 3] = '<END>'

            elif flag_multi_wp == True:
                master_shape_tuple[start_row, 1] = tmp_folder_name_of_wp + '_wp_'
                #sort WPs with hight
                tmp_hight_wp = []
                tmp_sort_multi_wp = []
                current_multi_wp_tuple = {}
                for tmp_shape_array in master_shape_array:
                    if tmp_shape_array[1] in tmp_sort_hight and tmp_shape_array[1] not in str(tmp_hight_wp):
                        tmp_hight_wp.append([tmp_shape_array[1],tmp_shape_array[3]])
                        tmp_sort_multi_wp.append(tmp_shape_array)

                current_multi_wp_tuple = ns_def.return_shape_tuple(tmp_sort_multi_wp, start_row)
                print(' --- current_multi_wp_tuple --- ')
                #print(master_shape_tuple[start_row, 1],current_multi_wp_tuple)

                pre_multi_wp_hight = 0
                count_multi_wp_hight = 0
                flag_first = True
                for tmp_current_multi_wp_tuple in current_multi_wp_tuple:
                    if tmp_current_multi_wp_tuple[0] > pre_multi_wp_hight:
                        count_multi_wp_hight += 1

                    if flag_first == True:
                        master_shape_tuple[tmp_current_multi_wp_tuple] = tmp_folder_name_of_wp + '_wp_'
                        flag_first = False
                    else:
                        master_shape_tuple[tmp_current_multi_wp_tuple] = current_multi_wp_tuple[tmp_current_multi_wp_tuple]

                    pre_multi_wp_hight = tmp_current_multi_wp_tuple[0]

                start_row += count_multi_wp_hight

                master_shape_tuple[start_row - 1, 3] = ''
                start_row -= 2

            master_shape_tuple[start_row + 1, 1] = '<END>'
            start_row += 2

        print('---- master_shape_tuple ----   <<POSITION_SHAPE>>')
        #print(master_shape_tuple)
        write_to_section = '<<POSITION_SHAPE>>'
        ns_def.write_excel_meta(master_shape_tuple, self.excel_file_path, self.worksheet_name, write_to_section,offset_row, offset_column)

        '''
        Write <<POSITION_LINE>>
        '''
        numpy_shape_array = np.array(master_shape_array)
        numpy_line_array = np.array(line_array)
        np.line_point = []
        for i in range(1,num_slide + 1):
            tmp_numpy_shape_array =  numpy_shape_array[numpy_shape_array[:, 7] == str(i)]
            tmp_numpy_line_array = numpy_line_array[numpy_line_array[:, 4] == i]
            for tmp_numpy_line in tmp_numpy_line_array:
                begin_cal_x = 999999
                begin_cal_y = 999999
                end_cal_x = 999999
                end_cal_y = 999999
                for tmp_numpy_shape in tmp_numpy_shape_array:
                    ### begin_point
                    if abs(tmp_numpy_line[1] - int(tmp_numpy_shape[3])) <= abs(tmp_numpy_line[1] - (int(tmp_numpy_shape[3])+int(tmp_numpy_shape[5]))):
                        begin_near_y = abs(tmp_numpy_line[1] - int(tmp_numpy_shape[3]))
                    else:
                        begin_near_y = abs(tmp_numpy_line[1] - (int(tmp_numpy_shape[3])+int(tmp_numpy_shape[5])))

                    if abs(tmp_numpy_line[0] - int(tmp_numpy_shape[2])) <= abs(tmp_numpy_line[0] - (int(tmp_numpy_shape[2]) + int(tmp_numpy_shape[4]))):
                        begin_near_x = abs(tmp_numpy_line[0] - int(tmp_numpy_shape[2]))
                    else:
                        begin_near_x = abs(tmp_numpy_line[0] - (int(tmp_numpy_shape[2]) + int(tmp_numpy_shape[4])))

                    if (begin_cal_x + begin_cal_y) > (begin_near_x + begin_near_y):
                        begin_cal_y = begin_near_y
                        begin_cal_x = begin_near_x
                        begin_shape_name = tmp_numpy_shape[1]
                        begin_folder_name = tmp_numpy_shape[0]

                for tmp_numpy_shape in tmp_numpy_shape_array:
                    ### end_point
                    if abs(tmp_numpy_line[3] - int(tmp_numpy_shape[3])) <= abs(tmp_numpy_line[3] - (int(tmp_numpy_shape[3])+int(tmp_numpy_shape[5]))):
                        end_near_y = abs(tmp_numpy_line[3] - int(tmp_numpy_shape[3]))
                    else:
                        end_near_y = abs(tmp_numpy_line[3] - (int(tmp_numpy_shape[3])+int(tmp_numpy_shape[5])))

                    if abs(tmp_numpy_line[2] - int(tmp_numpy_shape[2])) <= abs(tmp_numpy_line[2] - (int(tmp_numpy_shape[2]) + int(tmp_numpy_shape[4]))):
                        end_near_x = abs(tmp_numpy_line[2] - int(tmp_numpy_shape[2]))
                    else:
                        end_near_x = abs(tmp_numpy_line[2] - (int(tmp_numpy_shape[2]) + int(tmp_numpy_shape[4])))

                    if (end_cal_x + end_cal_y) > (end_near_x + end_near_y) and begin_shape_name != tmp_numpy_shape[1]:
                        end_cal_y = end_near_y
                        end_cal_x = end_near_x
                        end_shape_name = tmp_numpy_shape[1]
                        end_folder_name = tmp_numpy_shape[0]

                np.line_point.append([begin_shape_name,begin_near_x,begin_near_y,end_shape_name,end_near_x,end_near_y,tmp_numpy_line[4]])
        print(' --- np.line_point ----')
        #print(np.line_point)

        ### mark RIGHT LEFT ###
        mark_lr = []
        for tmp_line in np.line_point:
            flag_begin_is_wp = False
            flag_end_is_wp = False
            tmp_begin_tuple = ''
            tmp_end_tuple = ''
            for tmp_shape_tuple in master_shape_tuple:
                if master_shape_tuple[tmp_shape_tuple] == tmp_line[0]:
                    tmp_begin_tuple = tmp_shape_tuple

                    for tmp_wp in wp_array:
                        if master_shape_tuple[tmp_shape_tuple] == tmp_wp[1]:
                            flag_begin_is_wp = True

                if master_shape_tuple[tmp_shape_tuple] == tmp_line[3]:
                    tmp_end_tuple = tmp_shape_tuple

                    for tmp_wp in wp_array:
                        if master_shape_tuple[tmp_shape_tuple] == tmp_wp[1]:
                            flag_end_is_wp = True

            if flag_begin_is_wp == True or flag_end_is_wp == True:
                ### get folder name of WP ###
                tmp_start_folder_name = master_shape_tuple[tmp_begin_tuple]
                tmp_end_folder_name = master_shape_tuple[tmp_end_tuple]
                if master_shape_tuple[tmp_begin_tuple] in folder_name_of_wp:
                    tmp_start_folder_name = master_shape_tuple[tmp_begin_tuple] + '_wp_'

                for tmp_multi_wp in multi_wp:
                    if master_shape_tuple[tmp_begin_tuple] == tmp_multi_wp[0]:
                        #print(tmp_multi_wp[2] + '_wp_')
                        tmp_start_folder_name = tmp_multi_wp[2] + '_wp_'

                if master_shape_tuple[tmp_end_tuple] in folder_name_of_wp:
                    tmp_end_folder_name = master_shape_tuple[tmp_end_tuple] + '_wp_'

                for tmp_multi_wp in multi_wp:
                    if master_shape_tuple[tmp_end_tuple] == tmp_multi_wp[0]:
                        #print(tmp_multi_wp[2] + '_wp_')
                        tmp_end_folder_name = tmp_multi_wp[2] + '_wp_'

                ### input foldername without wp ###
                if tmp_start_folder_name not in '_wp_':
                    for tmp_shape_array in shape_array:
                        if tmp_shape_array[1] == tmp_start_folder_name:
                            tmp_start_folder_name = tmp_shape_array[0]
                            break

                if tmp_end_folder_name not in '_wp_':
                    for tmp_shape_array in shape_array:
                        if tmp_shape_array[1] == tmp_end_folder_name:
                            tmp_end_folder_name = tmp_shape_array[0]
                            break

                ### print line and wp folder name ###
                #print(master_shape_tuple[tmp_begin_tuple], tmp_start_folder_name, master_shape_tuple[tmp_end_tuple], tmp_end_folder_name)

                ### Check wp RIGHT LEFT ###
                tmp_start_column = 0
                tmp_end_left = 0
                tmp_start_row = 0
                tmp_end_top = 0
                flag_wp = False
                for tmp_master_excel_meta in master_folder_tuple:
                    if master_folder_tuple[tmp_master_excel_meta] == tmp_start_folder_name:
                        tmp_start_column = tmp_master_excel_meta[1]
                        tmp_start_row = tmp_master_excel_meta[0]
                        flag_wp = True

                    if master_folder_tuple[tmp_master_excel_meta] == tmp_end_folder_name:
                        tmp_end_left = tmp_master_excel_meta[1]
                        tmp_end_top = tmp_master_excel_meta[0]
                        flag_wp = True

                #print(flag_wp , master_shape_tuple[tmp_begin_tuple], master_shape_tuple[tmp_end_tuple],tmp_start_column,tmp_end_left,tmp_start_row,tmp_end_top)

                if tmp_start_column > tmp_end_left and tmp_start_row == tmp_end_top and flag_wp == True:
                    mark_lr.append([master_shape_tuple[tmp_begin_tuple], master_shape_tuple[tmp_end_tuple], 'LEFT', 'RIGHT'])
                elif tmp_start_column < tmp_end_left and tmp_start_row == tmp_end_top and flag_wp == True:
                    mark_lr.append([master_shape_tuple[tmp_begin_tuple], master_shape_tuple[tmp_end_tuple], 'RIGHT', 'LEFT'])

            else:
                if tmp_begin_tuple[0] == tmp_end_tuple[0] and tmp_begin_tuple[1] > tmp_end_tuple[1]:
                    mark_lr.append([master_shape_tuple[tmp_begin_tuple],master_shape_tuple[tmp_end_tuple],'LEFT','RIGHT'])

                if tmp_begin_tuple[0] == tmp_end_tuple[0] and tmp_begin_tuple[1] < tmp_end_tuple[1]:
                    mark_lr.append([master_shape_tuple[tmp_begin_tuple], master_shape_tuple[tmp_end_tuple], 'RIGHT', 'LEFT'])

        print('---- mark_lr ----')
        #print(mark_lr)

        ### make tuple of line ###
        master_line_tuple = {}
        start_row = 0
        for tmp_line_point in np.line_point:
            start_row += 1

            ### write begin end shape name ###
            master_line_tuple[start_row, 1] = tmp_line_point[0]
            master_line_tuple[start_row, 2] = tmp_line_point[3]

            ### write LEFT RIGHT ###
            for tmp_lr in mark_lr:
                if tmp_lr[0] == tmp_line_point[0] and tmp_lr[1] == tmp_line_point[3]:
                    master_line_tuple[start_row, 5] = tmp_lr[2]
                    master_line_tuple[start_row, 6] = tmp_lr[3]

        '''
        change line offset values
        '''
        ### count connected lines on up down left right. ###
        finish_shape_list = []
        num_line = 0

        for tmp_count in master_line_tuple:
            if tmp_count[1] == 1:
                num_line += 1

        full_line_array = []
        for tmp_num_line in range (1,num_line + 1):
            tmp_target_list = []
            tmp_current_shape = master_line_tuple[tmp_num_line,1]
            for tmp_tmp_num_line in range (1,num_line + 1):
                if tmp_current_shape == master_line_tuple[tmp_tmp_num_line,1] and tmp_current_shape != master_line_tuple[tmp_tmp_num_line,2] and master_line_tuple[tmp_tmp_num_line,1] not in finish_shape_list:
                    tmp_target_list.append(master_line_tuple[tmp_tmp_num_line,2])
                if tmp_current_shape == master_line_tuple[tmp_tmp_num_line,2] and tmp_current_shape != master_line_tuple[tmp_tmp_num_line,1] and master_line_tuple[tmp_tmp_num_line,2] not in finish_shape_list:
                    tmp_target_list.append(master_line_tuple[tmp_tmp_num_line,1])
            finish_shape_list.append(tmp_current_shape)

            if tmp_target_list != []:
                full_line_array.append([tmp_current_shape, tmp_target_list])
                #print(tmp_current_shape, tmp_target_list)

        for tmp_num_line in range (1,num_line + 1):
            tmp_target_list = []
            tmp_current_shape = master_line_tuple[tmp_num_line,2]
            for tmp_tmp_num_line in range (1,num_line + 1):
                if tmp_current_shape == master_line_tuple[tmp_tmp_num_line,1] and tmp_current_shape != master_line_tuple[tmp_tmp_num_line,2] and master_line_tuple[tmp_tmp_num_line,1] not in finish_shape_list:
                    tmp_target_list.append(master_line_tuple[tmp_tmp_num_line,2])
                if tmp_current_shape == master_line_tuple[tmp_tmp_num_line,2] and tmp_current_shape != master_line_tuple[tmp_tmp_num_line,1] and master_line_tuple[tmp_tmp_num_line,2] not in finish_shape_list:
                    tmp_target_list.append(master_line_tuple[tmp_tmp_num_line,1])
            finish_shape_list.append(tmp_current_shape)

            if tmp_target_list != []:
                full_line_array.append([tmp_current_shape, tmp_target_list])
                #print(tmp_current_shape, tmp_target_list , '   Second')

        print('---- full_line_array ----')
        #print(full_line_array)

        ### mark to up down left right ###
        master_udlr_array = []
        for tmp_full_line_array in full_line_array:
            tmp_udlr_array = []
            for tmp_tmp_full_line_array in tmp_full_line_array[1]:
                ### check left right ###
                tmp_uplf = ''
                for tmp_mark_lr in mark_lr:
                    if tmp_mark_lr[0] == tmp_full_line_array[0] and tmp_mark_lr[1] == tmp_tmp_full_line_array:
                        tmp_uplf = tmp_mark_lr[2]

                    elif tmp_mark_lr[1] == tmp_full_line_array[0] and tmp_mark_lr[0] == tmp_tmp_full_line_array:
                        tmp_uplf = tmp_mark_lr[3]

                if tmp_uplf != '':
                    tmp_begin_top = 0
                    tmp_end_top = 0
                    tmp_begin_left = 0
                    tmp_end_left = 0

                    for i in range(1, num_slide + 1):
                        flag_a = False
                        flag_b = False
                        tmp_numpy_shape_array = numpy_shape_array[numpy_shape_array[:, 7] == str(i)]
                        for tmp_tmp_numpy_shape_array in tmp_numpy_shape_array:
                            if tmp_tmp_numpy_shape_array[1] == tmp_full_line_array[0]:
                                flag_b = True
                                tmp_begin_top = tmp_tmp_numpy_shape_array[3]
                                tmp_begin_left = tmp_tmp_numpy_shape_array[2]
                            if tmp_tmp_numpy_shape_array[1] == tmp_tmp_full_line_array:
                                flag_a = True
                                tmp_end_top = tmp_tmp_numpy_shape_array[3]
                                tmp_end_left = tmp_tmp_numpy_shape_array[2]

                        if flag_a == True and flag_b == True:
                            break

                    #### Kyusai -> WP's LINE belongs folder ####

                    for tmp_folder in master_folder_tuple:
                        #print(master_folder_tuple[tmp_folder])
                        tmp_current_folder_name = numpy_shape_array[numpy_shape_array[:, 1] == str(tmp_tmp_full_line_array)]
                        if master_folder_tuple[tmp_folder] == str(tmp_current_folder_name[0][0]):
                            tmp_begin_top = int(tmp_begin_top) + (100000000 * int(tmp_folder[0]))
                            tmp_end_top = int(tmp_end_top) + (100000000 * int(tmp_folder[0]))
                            tmp_begin_left = int(tmp_begin_left) + (100000000 * int(tmp_folder[1]))
                            tmp_end_left = int(tmp_end_left) + (100000000 * int(tmp_folder[1]))

                    #### Output ####
                    tmp_udlr_array.append([tmp_tmp_full_line_array, tmp_uplf,tmp_begin_top ,tmp_end_top,tmp_begin_left ,tmp_end_left])

                if tmp_uplf == '':
                    ##### UP or DOWN
                    tmp_begin_top = 0
                    tmp_end_top = 0
                    tmp_begin_left = 0
                    tmp_end_left = 0

                    for i in range(1, num_slide + 1):
                        flag_a = False
                        flag_b = False
                        tmp_numpy_shape_array = numpy_shape_array[numpy_shape_array[:, 7] == str(i)]
                        for tmp_tmp_numpy_shape_array in tmp_numpy_shape_array:
                            if tmp_tmp_numpy_shape_array[1] == tmp_full_line_array[0]:
                                flag_b = True
                                tmp_begin_top = tmp_tmp_numpy_shape_array[3]
                                tmp_begin_left = tmp_tmp_numpy_shape_array[2]
                            if tmp_tmp_numpy_shape_array[1] == tmp_tmp_full_line_array:
                                flag_a = True
                                tmp_end_top = tmp_tmp_numpy_shape_array[3]
                                tmp_end_left = tmp_tmp_numpy_shape_array[2]

                        if flag_a == True and flag_b == True:
                            break

                    #### Kyusai -> WP's LINE belongs folder ####

                    for tmp_folder in master_folder_tuple:
                        #print(master_folder_tuple[tmp_folder])
                        tmp_current_folder_name = numpy_shape_array[numpy_shape_array[:, 1] == str(tmp_tmp_full_line_array)]
                        if master_folder_tuple[tmp_folder] == str(tmp_current_folder_name[0][0]):
                            tmp_begin_top = int(tmp_begin_top) + (100000000 * int(tmp_folder[0]))
                            tmp_end_top = int(tmp_end_top) + (100000000 * int(tmp_folder[0]))
                            tmp_begin_left = int(tmp_begin_left) + (100000000 * int(tmp_folder[1]))
                            tmp_end_left = int(tmp_end_left) + (100000000 * int(tmp_folder[1]))

                    #### Output ####
                    if tmp_begin_top >= tmp_end_top:
                        tmp_udlr_array.append([tmp_tmp_full_line_array, 'UP', tmp_begin_top, tmp_end_top, tmp_begin_left, tmp_end_left])
                    else:
                        tmp_udlr_array.append([tmp_tmp_full_line_array, 'DOWN', tmp_begin_top, tmp_end_top, tmp_begin_left, tmp_end_left])
            udlr_array = []
            udlr_array.append([tmp_full_line_array[0],tmp_udlr_array])
            #print(udlr_array)

            pre_udlr_array = []
            ### append UP DOWN LEFT RIGHT ###
            for tmp_udlr_array in udlr_array:
                tmp_up_array = []
                tmp_down_array = []
                tmp_left_array = []
                tmp_right_array = []
                for tmp_tmp_udlr_array in tmp_udlr_array[1]:
                    tmp_tmp_udlr_array[2] = int(tmp_tmp_udlr_array[2])
                    tmp_tmp_udlr_array[3] = int(tmp_tmp_udlr_array[3])
                    tmp_tmp_udlr_array[4] = int(tmp_tmp_udlr_array[4])
                    tmp_tmp_udlr_array[5] = int(tmp_tmp_udlr_array[5])

                    if tmp_tmp_udlr_array[1] == 'UP':
                        tmp_up_array.append(tmp_tmp_udlr_array)
                    elif tmp_tmp_udlr_array[1] == 'DOWN':
                        tmp_down_array.append(tmp_tmp_udlr_array)
                    elif tmp_tmp_udlr_array[1] == 'LEFT':
                        tmp_left_array.append(tmp_tmp_udlr_array)
                    elif tmp_tmp_udlr_array[1] == 'RIGHT':
                        tmp_right_array.append(tmp_tmp_udlr_array)

                #sort up down left right
                tmp_up_array = sorted(tmp_up_array, reverse=False, key=lambda x: x[5])  # sort for tmp_end_top
                tmp_down_array = sorted(tmp_down_array, reverse=False, key=lambda x: x[5])  # sort for tmp_end_left
                tmp_left_array = sorted(tmp_left_array, reverse=False, key=lambda x: x[4])  # sort for tmp_end_top
                tmp_right_array = sorted(tmp_right_array, reverse=False, key=lambda x: x[4])  # sort for tmp_end_left

                pre_udlr_array.append([tmp_udlr_array[0],[tmp_up_array,tmp_down_array,tmp_left_array,tmp_right_array]])
                master_udlr_array.append(pre_udlr_array)

        print('---- master_udlr_array ----')
        #print(master_udlr_array)

        ### put offset value to master_line_tuple
        pre_master_line_tuple = {}
        finish_line_list = []

        for tmp_master_udlr_array in master_udlr_array:
            for tmp_tmp_master_udlr_array in tmp_master_udlr_array[0][1]:
                tmp_len_line = len(tmp_tmp_master_udlr_array)
                if tmp_len_line >= 2:
                    offset = ((tmp_len_line - 1) / 2) * -1 * line_offset_value
                    for b in range(tmp_len_line):
                        #print(b,tmp_master_udlr_array[0][0] ,tmp_tmp_master_udlr_array[b][0])
                        for tmp_master_line_tuple in master_line_tuple:
                            if (master_line_tuple[tmp_master_line_tuple[0],1] == tmp_master_udlr_array[0][0] and master_line_tuple[tmp_master_line_tuple[0],2] == tmp_tmp_master_udlr_array[b][0]) \
                                    or (master_line_tuple[tmp_master_line_tuple[0],1] == tmp_tmp_master_udlr_array[b][0] and master_line_tuple[tmp_master_line_tuple[0],2] == tmp_master_udlr_array[0][0]):
                                if master_line_tuple[tmp_master_line_tuple[0],1] == tmp_tmp_master_udlr_array[b][0] and master_line_tuple[tmp_master_line_tuple[0],2] == tmp_master_udlr_array[0][0]:
                                    tmp_offset_column = 2
                                else:
                                    tmp_offset_column = 0

                                if tmp_tmp_master_udlr_array[0][1] == 'UP' or tmp_tmp_master_udlr_array[0][1] == 'DOWN':
                                    if [tmp_master_line_tuple[0],7+tmp_offset_column] not in finish_line_list:
                                        pre_master_line_tuple[tmp_master_line_tuple[0], 7+tmp_offset_column] = offset
                                        offset += line_offset_value
                                        finish_line_list.append([tmp_master_line_tuple[0],7+tmp_offset_column])
                                        break

                                if tmp_tmp_master_udlr_array[0][1] == 'LEFT' or tmp_tmp_master_udlr_array[0][1] == 'RIGHT':
                                    if [tmp_master_line_tuple[0],8+tmp_offset_column] not in finish_line_list:
                                        pre_master_line_tuple[tmp_master_line_tuple[0], 8+tmp_offset_column] = offset
                                        offset += line_offset_value
                                        finish_line_list.append([tmp_master_line_tuple[0],8+tmp_offset_column])
                                        break

        print('---- pre_master_line_tuple ----')
        #print(pre_master_line_tuple)
        for tmp_pre_master_line_tuple in pre_master_line_tuple:
            master_line_tuple[tmp_pre_master_line_tuple] = pre_master_line_tuple[tmp_pre_master_line_tuple]

        '''
        Write Tag Name in <<POSITION_LINE>>
        '''
        pre_pre_master_line_tuple = {}
        if flag_put_line_tag == True:
            if_num = 0
            for tmp_master_line_tuple in master_line_tuple:
                if tmp_master_line_tuple[1] == 1:
                    #write port setting
                    pre_pre_master_line_tuple[tmp_master_line_tuple[0], 3] = tag_name_prefix + str(if_num)
                    pre_pre_master_line_tuple[tmp_master_line_tuple[0], 4] = tag_name_prefix + str(if_num)
                    pre_pre_master_line_tuple[tmp_master_line_tuple[0], 13] = port_name_prefix
                    pre_pre_master_line_tuple[tmp_master_line_tuple[0], 14] = port_speed
                    pre_pre_master_line_tuple[tmp_master_line_tuple[0], 15] = port_duplex
                    pre_pre_master_line_tuple[tmp_master_line_tuple[0], 16] = port_type
                    pre_pre_master_line_tuple[tmp_master_line_tuple[0], 17] = port_name_prefix
                    pre_pre_master_line_tuple[tmp_master_line_tuple[0], 18] = port_speed
                    pre_pre_master_line_tuple[tmp_master_line_tuple[0], 19] = port_duplex
                    pre_pre_master_line_tuple[tmp_master_line_tuple[0], 20] = port_type
                    if_num += 1

        #print('pre_pre_master_line_tuple')
        #print(pre_pre_master_line_tuple)
        for tmp_pre_master_line_tuple in pre_pre_master_line_tuple:
            master_line_tuple[tmp_pre_master_line_tuple] = pre_pre_master_line_tuple[tmp_pre_master_line_tuple]

        print('---- master_line_tuple ----')
        #print(master_line_tuple)

        ### increse offset value for WP , ver 1.12
        wp_name_array = []
        wp_width_offset_inche_array = []
        for i_wp_array in wp_array:
            wp_name_array.append(i_wp_array[1])

        i_master_line_tuple = master_line_tuple
        for tmp_i_master_line_tuple in i_master_line_tuple:
            if ns_def.check_tuple_num_exist(i_master_line_tuple, tmp_i_master_line_tuple[0], 1) == True or \
                    ns_def.check_tuple_num_exist(i_master_line_tuple, tmp_i_master_line_tuple[0], 2) == True:
                if tmp_i_master_line_tuple[1] == 1 and i_master_line_tuple[tmp_i_master_line_tuple[0],1] in wp_name_array:
                    if ns_def.check_tuple_num_exist(i_master_line_tuple, tmp_i_master_line_tuple[0], 7) == True:
                        wp_width_offset_inche_array.append([i_master_line_tuple[tmp_i_master_line_tuple[0], 1],i_master_line_tuple[tmp_i_master_line_tuple[0], 7] * (line_offset_ratio_wp - 1.0)])
                        master_line_tuple[tmp_i_master_line_tuple[0], 7] = i_master_line_tuple[tmp_i_master_line_tuple[0], 7] * line_offset_ratio_wp
                    if ns_def.check_tuple_num_exist(i_master_line_tuple, tmp_i_master_line_tuple[0], 8) == True:
                        wp_width_offset_inche_array.append([i_master_line_tuple[tmp_i_master_line_tuple[0], 1],i_master_line_tuple[tmp_i_master_line_tuple[0], 8] ])
                        master_line_tuple[tmp_i_master_line_tuple[0], 8] = i_master_line_tuple[tmp_i_master_line_tuple[0],8]
                elif tmp_i_master_line_tuple[1] == 2 and i_master_line_tuple[tmp_i_master_line_tuple[0],2] in wp_name_array:
                    if ns_def.check_tuple_num_exist(i_master_line_tuple, tmp_i_master_line_tuple[0], 9) == True:
                        wp_width_offset_inche_array.append([i_master_line_tuple[tmp_i_master_line_tuple[0], 2],i_master_line_tuple[tmp_i_master_line_tuple[0], 9] * (line_offset_ratio_wp - 1.0)])
                        master_line_tuple[tmp_i_master_line_tuple[0], 9] = i_master_line_tuple[tmp_i_master_line_tuple[0],9] * line_offset_ratio_wp
                    if ns_def.check_tuple_num_exist(i_master_line_tuple, tmp_i_master_line_tuple[0], 10) == True:
                        wp_width_offset_inche_array.append([i_master_line_tuple[tmp_i_master_line_tuple[0], 2],i_master_line_tuple[tmp_i_master_line_tuple[0], 10] ])
                        master_line_tuple[tmp_i_master_line_tuple[0], 10] = i_master_line_tuple[tmp_i_master_line_tuple[0],10]

        write_to_section = '<<POSITION_LINE>>'
        offset_row = 2
        ns_def.write_excel_meta(master_line_tuple, self.excel_file_path, self.worksheet_name, write_to_section, offset_row, offset_column)
        offset_row = 0

        '''
        Write <<STYLE_SHAPE>>
        '''
        master_style_shape_tuple = {}
        line_wh_array = []
        tmp_tuple_row = 1
        for tmp_master_udlr_array in master_udlr_array:
            for tmp_tmp_master_udlr_array in tmp_master_udlr_array[0]:
                if len(tmp_tmp_master_udlr_array) < 4: # Modify Ver 1.1
                    continue
                len_up = len(tmp_tmp_master_udlr_array[0])
                len_down = len(tmp_tmp_master_udlr_array[1])
                len_left = len(tmp_tmp_master_udlr_array[2])
                len_right = len(tmp_tmp_master_udlr_array[3])

            if len_up >= len_down:
                line_width = len_up
            else:
                line_width = len_down
            line_width = ((line_width - 1) * line_offset_value) * (1 + wp_roundness)

            if len_left >= len_right:
                line_hight = len_left
            else:
                line_hight = len_right
            line_hight = ((line_hight - 1) * line_offset_value) * (1 + wp_roundness)

            #### get num char
            num_char = ns_def.get_description_width_hight(self.shae_font_size,tmp_master_udlr_array[0][0])[0]
            if num_char > line_width:
                line_width = num_char
            if shape_width_min > line_width:
                line_width = shape_width_min
            if shape_hight_min > line_hight:
                line_hight = shape_hight_min

            ##input roundness
            tmp_roundness = 0
            tmp_color = color_shape

            # change color, way point
            for tmp_wp_array in wp_array:
                if tmp_wp_array[1] == tmp_master_udlr_array[0][0]:
                    tmp_roundness = wp_roundness
                    tmp_color = color_wp
                    if tmp_wp_array[2] == 'RIGHT' or tmp_wp_array[2] == 'LEFT':  # ver 2.0  fixed wp roundness when RIGHT or LEFT
                        tmp_roundness = wp_roundness_right_left

            # change color, @Icon@
            if '@Icon@' in str(tmp_master_udlr_array[0][0]):
                tmp_color = color_atmark

            line_wh_array.append([tmp_master_udlr_array[0][0],line_width,line_hight,tmp_roundness])

            master_style_shape_tuple[(tmp_tuple_row ,1)] = tmp_master_udlr_array[0][0]
            master_style_shape_tuple[(tmp_tuple_row, 2)] = line_width
            master_style_shape_tuple[(tmp_tuple_row, 3)] = line_hight
            master_style_shape_tuple[(tmp_tuple_row, 4)] = tmp_roundness
            master_style_shape_tuple[(tmp_tuple_row, 5)] = tmp_color
            tmp_tuple_row += 1

        ### increse offset value for WP , ver 1.12
        i_master_style_shape_tuple = master_style_shape_tuple
        for tmp_i_master_style_shape_tuple in i_master_style_shape_tuple:
            tmp_max_offset_width = 0
            if tmp_i_master_style_shape_tuple[1] == 1 and i_master_style_shape_tuple[tmp_i_master_style_shape_tuple[0], 1] in wp_name_array:
                for tmp_wp_width_offset_inche_array in wp_width_offset_inche_array:
                    if tmp_wp_width_offset_inche_array[0] == i_master_style_shape_tuple[tmp_i_master_style_shape_tuple[0], 1] and tmp_max_offset_width < tmp_wp_width_offset_inche_array[1]:
                        tmp_max_offset_width = tmp_wp_width_offset_inche_array[1]

                master_style_shape_tuple[tmp_i_master_style_shape_tuple[0], 2] = i_master_style_shape_tuple[tmp_i_master_style_shape_tuple[0], 2] + (tmp_max_offset_width * 2) + 0.2
                master_style_shape_tuple[tmp_i_master_style_shape_tuple[0], 3] = i_master_style_shape_tuple[ tmp_i_master_style_shape_tuple[0], 3] + shape_hight_offset_inches_wp

        #print('---- master_style_shape_tuple ---- ')
        #print(master_style_shape_tuple)

        write_to_section = '<<STYLE_SHAPE>>'
        offset_row = 3
        ns_def.write_excel_meta(master_style_shape_tuple, self.excel_file_path, self.worksheet_name, write_to_section, offset_row, offset_column)
        offset_row = 0

        '''
        Write <<POSITION_TAG>>
        '''
        master_line_tag_tuple = {}
        if flag_put_line_tag == True:

            for tmp_master_style_shape_tuple in master_style_shape_tuple:
                if tmp_master_style_shape_tuple[1] == 1:
                    ### check max char num of shape's tag name
                    max_len_num = 1
                    for tmp_master_line_tuple in master_line_tuple:
                        if master_line_tuple[tmp_master_line_tuple] == master_style_shape_tuple[tmp_master_style_shape_tuple]:
                            if tmp_master_line_tuple[1] == 1:
                                tmp_len_char = len(master_line_tuple[tmp_master_line_tuple[0], 3])
                            elif tmp_master_line_tuple[1] == 2:
                                tmp_len_char = len(master_line_tuple[tmp_master_line_tuple[0], 4])
                            if max_len_num < tmp_len_char:
                                max_len_num = tmp_len_char

                    master_line_tag_tuple[(tmp_master_style_shape_tuple[0], 1)] = master_style_shape_tuple[tmp_master_style_shape_tuple]
                    master_line_tag_tuple[(tmp_master_style_shape_tuple[0], 2)] = 'LINE'
                    master_line_tag_tuple[(tmp_master_style_shape_tuple[0], 5)] = max_len_num * tag_offet_inche
                    master_line_tag_tuple[(tmp_master_style_shape_tuple[0], 6)] = 'YES'

            print('---- master_line_tag_tuple ---- ')
            #print(master_line_tag_tuple)
            write_to_section = '<<POSITION_TAG>>'
            offset_row = 2
            ns_def.write_excel_meta(master_line_tag_tuple, self.excel_file_path, self.worksheet_name, write_to_section, offset_row, offset_column)
            offset_row = 0

        '''
        Calculate the best size of a Slide 
        '''
        # GET each folder size
        master_folder_size_array = []
        'parameter'
        min_tag_inches = 0.3  # inches,  between side of folder and eghe shape. left and right.

        #### GET best width size ####
        master_folder_size_array = ns_def.get_folder_width_size(master_folder_tuple,master_style_shape_tuple,master_shape_tuple,min_tag_inches)
        #print('-----master_folder_size_array-----  ',master_folder_size_array)  #[slide_max_width_inches, master_width_size_y_grid ,master_folder_size,slide_max_hight_inches,master_hight_size_y_grid]

        #print(master_folder_tuple)
        update_master_folder_tuple = {}
        for tmp_master_width_size_y_grid in master_folder_size_array[1]:
            for tmp_master_folder_size in master_folder_size_array[2]:
                if tmp_master_width_size_y_grid[0] == tmp_master_folder_size[0]:
                    if master_folder_size_array[0] == tmp_master_width_size_y_grid[1]: # check max width in the slide
                        for tmp_master_folder_tuple in master_folder_tuple:
                            if tmp_master_folder_tuple[0] == tmp_master_folder_size[0] and master_folder_tuple[tmp_master_folder_tuple] == tmp_master_folder_size[1][0][0]:
                                update_master_folder_tuple[tmp_master_folder_tuple[0]-1,tmp_master_folder_tuple[1]] = tmp_master_folder_size[1][0][1]
                                update_master_folder_tuple[tmp_master_folder_tuple[0], tmp_master_folder_tuple[1]] = master_folder_tuple[tmp_master_folder_tuple]
                            elif tmp_master_folder_tuple[0] == tmp_master_folder_size[0] and master_folder_tuple[tmp_master_folder_tuple] == '':
                                update_master_folder_tuple[tmp_master_folder_tuple[0] - 1, tmp_master_folder_tuple[1]] = tmp_master_width_size_y_grid[2]
                                update_master_folder_tuple[tmp_master_folder_tuple[0], tmp_master_folder_tuple[1]] = ''
                            elif tmp_master_folder_tuple[1] == 1: # write ALL column =1
                                update_master_folder_tuple[tmp_master_folder_tuple] = master_folder_tuple[tmp_master_folder_tuple]

                    else: # insert empty folder to left and right side
                        #print('-insert empty folder to left and right side - ', tmp_master_folder_size)
                        tmp_max_row = 0
                        tmp_max_column = 0

                        tmp_bothside_empty = (master_folder_size_array[0] - tmp_master_width_size_y_grid[1]) * 0.25
                        #print('tmp_bothside_empty----  ',tmp_master_min_size_y_grid)

                        for tmp_master_folder_tuple in master_folder_tuple:
                            ### set first column
                            if tmp_master_folder_tuple[0] == tmp_master_folder_size[0]:

                                #check to exist tuble 'master_folder_tuple[tmp_master_folder_tuple[0]-1,tmp_master_folder_tuple[1]-1]' # for bux fix 001
                                flag_exist_POSITION_FOLDER = False
                                flag_exist_SET_WIDTH = False
                                for tmp_bug_fix_tuple in master_folder_tuple:
                                    if tmp_bug_fix_tuple[0] == tmp_master_folder_tuple[0]-1 and tmp_bug_fix_tuple[1] == tmp_master_folder_tuple[1]-1:
                                        if master_folder_tuple[tmp_master_folder_tuple[0]-1,tmp_master_folder_tuple[1]-1] == '<<POSITION_FOLDER>>':
                                            flag_exist_POSITION_FOLDER = True
                                        if master_folder_tuple[tmp_master_folder_tuple[0]-1,tmp_master_folder_tuple[1]-1] == '<SET_WIDTH>':
                                            flag_exist_SET_WIDTH = True


                                if tmp_master_folder_tuple[0] != 1 and tmp_master_folder_tuple[1] == 2 and (flag_exist_SET_WIDTH == True or flag_exist_POSITION_FOLDER == True):
                                    update_master_folder_tuple[tmp_master_folder_tuple[0]-1, tmp_master_folder_tuple[1]] = tmp_bothside_empty
                                elif tmp_master_folder_tuple[0] != 1 and tmp_master_folder_tuple[1] == 2 and flag_exist_SET_WIDTH == True and flag_exist_POSITION_FOLDER == True:
                                    update_master_folder_tuple[tmp_master_folder_tuple[0], tmp_master_folder_tuple[1]] = ''

                            ### set body
                            if tmp_master_folder_tuple[0] == tmp_master_folder_size[0] and tmp_master_folder_tuple[1] != 1:
                                if tmp_master_folder_tuple[0] == tmp_master_folder_size[0] and master_folder_tuple[tmp_master_folder_tuple] == tmp_master_folder_size[1][0][0]:
                                    update_master_folder_tuple[tmp_master_folder_tuple[0] - 1, tmp_master_folder_tuple[1]+1] = tmp_master_folder_size[1][0][1]
                                    update_master_folder_tuple[tmp_master_folder_tuple[0], tmp_master_folder_tuple[1]+1] = master_folder_tuple[tmp_master_folder_tuple]
                                elif tmp_master_folder_tuple[0] == tmp_master_folder_size[0] and master_folder_tuple[tmp_master_folder_tuple] == '':
                                    update_master_folder_tuple[tmp_master_folder_tuple[0] - 1, tmp_master_folder_tuple[1]+1] = tmp_master_width_size_y_grid[2]
                                    update_master_folder_tuple[tmp_master_folder_tuple[0], tmp_master_folder_tuple[1]+1] = ''

                            ### set last coumn
                            if tmp_master_folder_tuple[0] == tmp_master_folder_size[0] and tmp_master_folder_tuple[0] > tmp_max_row:
                                tmp_max_row = tmp_master_folder_tuple[0]
                            if tmp_master_folder_tuple[0] == tmp_master_folder_size[0] and tmp_master_folder_tuple[1] > tmp_max_column:
                                tmp_max_column = tmp_master_folder_tuple[1]

                        if (tmp_max_column + 2) == 3:
                            #kyusai only empty row
                            update_master_folder_tuple[tmp_max_row -1, tmp_max_column + 1] = 10
                        else:
                            update_master_folder_tuple[tmp_max_row - 1, tmp_max_column + 2] = tmp_bothside_empty
                            update_master_folder_tuple[tmp_max_row, tmp_max_column + 2] = ''

        flag_wp_only = True
        for tmp_master_folder_size in master_folder_size_array[2]:
            if tmp_master_width_size_y_grid[0] == tmp_master_folder_size[0]:
                if '_wp_' not in str(tmp_master_folder_size[1][0][0]):
                    flag_wp_only = False
                    break

        ####only way point
        if flag_wp_only == True:
            for tmp_tmp_master_folder_tuple in master_folder_tuple:
                if tmp_master_width_size_y_grid[0] == tmp_tmp_master_folder_tuple[0] and tmp_tmp_master_folder_tuple[1] != 1:
                    #update_master_folder_tuple[tmp_tmp_master_folder_tuple ] = master_folder_tuple[tmp_tmp_master_folder_tuple]
                    if ns_def.check_tuple_num_exist(update_master_folder_tuple,tmp_tmp_master_folder_tuple[0]-1, tmp_tmp_master_folder_tuple[1]+2) == True:
                        update_master_folder_tuple[tmp_tmp_master_folder_tuple[0]-1,tmp_tmp_master_folder_tuple[1]] = update_master_folder_tuple[tmp_tmp_master_folder_tuple[0]-1,tmp_tmp_master_folder_tuple[1]+2]
                    else:
                        update_master_folder_tuple[tmp_tmp_master_folder_tuple[0] - 1, tmp_tmp_master_folder_tuple[1]] = 0.999
                elif master_folder_tuple[tmp_tmp_master_folder_tuple] == 0.999: ###bug fix 001###
                    update_master_folder_tuple[tmp_tmp_master_folder_tuple] = master_folder_tuple[tmp_tmp_master_folder_tuple]



        #### update best hight size ####
        pre_update_master_folder_tuple = update_master_folder_tuple
        update_master_folder_tuple = {}
        for tmp_pre_update_master_folder_tuple in pre_update_master_folder_tuple:
            if tmp_pre_update_master_folder_tuple[1] != 1 or pre_update_master_folder_tuple[tmp_pre_update_master_folder_tuple] == '<SET_WIDTH>':
                update_master_folder_tuple[tmp_pre_update_master_folder_tuple] = pre_update_master_folder_tuple[tmp_pre_update_master_folder_tuple]
            else:
                for tmp_master_folder_size_array in master_folder_size_array[4]:
                    if tmp_master_folder_size_array[0] == tmp_pre_update_master_folder_tuple[0]:
                        update_master_folder_tuple[tmp_pre_update_master_folder_tuple] = tmp_master_folder_size_array[1]

        print('update_master_folder_tuple-----  <<POSITION_FOLDER>> ')
        #print(update_master_folder_tuple)

        ### write excel meta ###
        write_to_section = '<<POSITION_FOLDER>>'
        offset_row = 0
        ns_def.write_excel_meta(update_master_folder_tuple, self.excel_file_path, self.worksheet_name, write_to_section,offset_row, offset_column)

        '''
        Write <<ROOT_FOLDER>>
        '''
        ppt_min_width = 8  #  inches
        ppt_min_hight = 3    #  inches
        ppt_max_width = 56  #  inches
        ppt_max_hight = 56   #  inches

        self.root_left = 0.28
        self.root_top  = 1.42
        self.root_width = math.ceil(master_folder_size_array[0] * 10) / 10
        self.root_hight = math.ceil(master_folder_size_array[3] * 10) / 10

        # width inches of root folder ( ver 1.1)
        if (self.root_width + self.root_left * 2) > ppt_max_width:
            self.root_width = ppt_max_width - (self.root_left * 2)
        # hight inches of root folder ( ver 1.1)
        if (self.root_hight + self.root_top * 2) > ppt_max_hight:
            self.root_hight = ppt_max_hight - (self.root_top * 2)


        master_root_folder_tuple = {}
        master_root_folder_tuple[2, 3] = 1
        master_root_folder_tuple[2, 4] = 1
        master_root_folder_tuple[2, 5] = self.root_left
        master_root_folder_tuple[2, 6] = self.root_top

        if self.root_width < (ppt_min_width  - (self.root_left * 2)):
            master_root_folder_tuple[2, 7] = (ppt_min_width  - (self.root_left * 2))
        else:
            master_root_folder_tuple[2, 7] = self.root_width

        if self.root_hight < (ppt_min_hight - (self.root_top * 1.5 )):
            master_root_folder_tuple[2, 8] = (ppt_min_hight - (self.root_top * 1.5 ))
        else:
            master_root_folder_tuple[2, 8] = self.root_hight

        master_root_folder_tuple[2, 2] = '[L1] All Areas'

        write_to_section = '<<ROOT_FOLDER>>'
        ns_def.write_excel_meta(master_root_folder_tuple, self.excel_file_path, self.worksheet_name, write_to_section, offset_row, offset_column)

        #print('---- <<ROOT_FOLDER>>  master_root_folder_tuple ---- ',master_root_folder_tuple)
        #print(master_root_folder_tuple)

if __name__ == '__main__':
    ns_l1_master_create()
