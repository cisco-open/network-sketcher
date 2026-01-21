'''
SPDX-License-Identifier: Apache-2.0

Copyright 2023 Cisco Systems, Inc. and its affiliates

Licensed under the Apache License, Version 2.0 (the "License");
you may not use this file except in compliance with the License.
You may obtain a copy of the License at

http://www.apache.org/licenses/LICENSE-2.0

Unless required by applicable law or agreed to in writing, software
distributed under the License is distributed on an "AS IS" BASIS,
WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
See the License for the specific language governing permissions and
limitations under the License.
'''
import ns_def, ns_egt_maker, ns_ddx_figure, ns_cli
from collections import Counter
import tkinter as tk ,tkinter.ttk , openpyxl
import ipaddress, sys, os, re, shutil
import numpy as np
import networkx as nx
import subprocess

#add at ver 2.5.2
class ai_context():
    def export_ai_context_file(self, dummy):
        print('--- export_ai_context_file ---')
        excel_maseter_file = self.inFileTxt_L2_3_1.get()
        iDir = os.path.abspath(os.path.dirname(excel_maseter_file))

        basename_without_ext = os.path.splitext(os.path.basename(excel_maseter_file))[0]
        self.outFileTxt_11_3.delete(0, tkinter.END)
        self.outFileTxt_11_3.insert(tk.END, iDir + ns_def.return_os_slash() + '[AI_Context]' + basename_without_ext.replace('[MASTER]', '') + '.txt')

        ## check file open
        ns_def.check_file_open(self.outFileTxt_11_3.get())

        # remove exist flow file
        if os.path.isfile(self.outFileTxt_11_3.get()) == True:
            os.remove(self.outFileTxt_11_3.get())

        self.ai_context_file = self.outFileTxt_11_3.get()
        print(self.ai_context_file)

        content_to_append = ''
        content_to_append += '\'\'\'' + '\n' + 'Basic response policy'+ '\n' + '\'\'\'' + '\n'
        content_to_append += '* You are a network specialist and technical consultant at Cisco.' + '\n'
        content_to_append += '* You provide specific, logical answers to broad and technical questions or consultations, including your reasoning, and you possess a high level of analytical ability.' + '\n'
        content_to_append += '* A customer has provided output from the OSS tool “Network Sketcher” using a show command.' + '\n' + '\n'

        content_to_append += '\'\'\'' + '\n' + 'All data in the master file'+ '\n' + '\'\'\'' + '\n'

        export_num = 0
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_area' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'area'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_area_device' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'area_device'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_area_location' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'area_location'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_attribute' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'attribute'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_device' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'device'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_device_interface' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'device_interface'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_device_location' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'device_location'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_l1_interface' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'l1_interface'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_l1_link' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'l1_link'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_l2_broadcast_domain' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'l2_broadcast_domain'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_l2_interface' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'l2_interface'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_l3_broadcast_domain' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'l3_broadcast_domain'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_l3_interface' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'l3_interface'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_waypoint' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'waypoint'])) + '\n'
        print('** Exporting ' + str(export_num := export_num + 1) + '/15')
        content_to_append += '** show_waypoint_interface' + '\n' + str(ns_cli.ns_cli_run.cli_show(self, excel_maseter_file, ['show', 'waypoint_interface'])) + '\n'+ '\n'

        # add commands's guide
        def resource_path(relative_path):
            try:
                base_path = sys._MEIPASS
            except Exception:
                base_path = os.path.abspath(os.path.dirname(__file__))

            return os.path.join(base_path, relative_path)

        file_path = resource_path('ns_extensions_cmd_list.txt')
        with open(file_path, 'r', encoding='utf-8') as f:
            content_to_append += f.read()


        try:
            # Open the file in append mode ('a'). If the file does not exist, it will be created automatically.
            with open(self.ai_context_file, 'a', encoding='utf-8') as file:
                # Write the content to the file, followed by a newline character
                file.write(content_to_append + '\n')
                # Print a success message indicating the file has been updated
            #print(f"'{content_to_append}' has been appended to {self.ai_context_file}.")
        except Exception as e:
            # Handle any errors that occur during the file operation
            print(f"An error occurred: {e}")

class flow_report():
    def create_device_flow_table(self,full_filepath_master,device_name):
        print('--- create_device_flow_table ---',full_filepath_master,device_name)
        ## check Flow_Data sheet exists in Master file
        input_excel_master = openpyxl.load_workbook(full_filepath_master)
        ws_list_master = input_excel_master.sheetnames
        input_excel_master.close()

        ws_flow_name = 'Flow_Data'
        if ws_flow_name in ws_list_master:
            master_flow_array = ns_def.convert_excel_to_array(ws_flow_name, full_filepath_master, 3)
            master_flow_array = master_flow_array[:-1]
            #print(master_flow_array)
            target_flow_array = []
            for tmp_master_flow_array in master_flow_array:
                tmp_routing_path = ''
                if tmp_master_flow_array[1][6] != '':
                    tmp_routing_path = tmp_master_flow_array[1][6]
                elif tmp_master_flow_array[1][6] == '' and tmp_master_flow_array[1][6] != ' ':
                    tmp_routing_path = tmp_master_flow_array[1][7]

                if tmp_master_flow_array[1][1] == device_name or tmp_master_flow_array[1][2] == device_name or device_name in tmp_routing_path:
                    target_flow_array.append([tmp_master_flow_array[1][0],tmp_master_flow_array[1][1],tmp_master_flow_array[1][2],tmp_master_flow_array[1][3],tmp_master_flow_array[1][4],tmp_master_flow_array[1][5],tmp_routing_path])
            #print(target_flow_array)


            '''
            export flow report
            '''
            excel_maseter_file = full_filepath_master
            iDir = os.path.abspath(os.path.dirname(excel_maseter_file))

            basename_without_ext = os.path.splitext(os.path.basename(excel_maseter_file))[0]
            self.outFileTxt_11_3.delete(0, tkinter.END)
            self.outFileTxt_11_3.insert(tk.END, iDir + ns_def.return_os_slash() + '[FLOW_REPORT]' + basename_without_ext.replace('[MASTER]', '') + '.xlsx')

            ## check file open
            ns_def.check_file_open(self.outFileTxt_11_3.get())

            # flag exist flow file
            flag_flow_table_exist = False
            if os.path.isfile(self.outFileTxt_11_3.get()) == True:
                #os.remove(self.outFileTxt_11_3.get())
                flag_flow_table_exist = True

            self.excel_flow_file = self.outFileTxt_11_3.get()

            ## check Flow_Data sheet exists in Master file
            input_excel_master = openpyxl.load_workbook(excel_maseter_file)
            ws_list_master = input_excel_master.sheetnames
            input_excel_master.close()

            '''
            MAKE Flows Table List
            '''
            master_device_table_tuple = {}
            flow_list_array = []
            egt_maker_width_array = ['5','25', '25','25', '25','15', '20', '25', '40']  # for Network Sketcher Ver 2.0
            flow_list_array.append([1, ['<RANGE>', '1','1', '1', '1', '1', '1', '1', '1', '1', '<END>']])
            flow_list_array.append([2, ['<HEADER>', 'No','Source Device Name', 'Destination Device Name','Source IP Address', 'Destination IP Address','TCP/UDP/ICMP','Service name(Port)',  'Max. bandwidth(Mbps)', 'Routing path settings', '<END>']])
            current_row_num = 3

            '''add flow table list'''
            #print(self.show_l3_interface)
            # Initialize a dictionary to hold devices and their IP addresses
            device_ips = {}

            # Iterate through the list of interfaces
            for tmp_show_l3_interface in self.show_l3_interface:
                tmp_device_name = tmp_show_l3_interface[0]
                tmp_ip_address = tmp_show_l3_interface[3]

                # Add the IP address to the corresponding device in the dictionary
                if tmp_device_name not in device_ips:
                    device_ips[tmp_device_name] = []
                device_ips[tmp_device_name].append(tmp_ip_address)

            for tmp_target_flow_array in target_flow_array:
                tmp_target_flow_array = list(map(str, tmp_target_flow_array))
                tmp_target_flow_array.insert(0, '')
                tmp_target_flow_array.append('<END>')
                source_ip_array = device_ips[tmp_target_flow_array[2]]
                destination_ip_array = device_ips[tmp_target_flow_array[3]]
                str_source_ip = ', '.join(map(str, source_ip_array ))
                str_destination_ip = ', '.join(map(str, destination_ip_array))
                tmp_target_flow_array.insert(4, str_source_ip)
                tmp_target_flow_array.insert(5, str_destination_ip)
                flow_list_array.append([current_row_num,tmp_target_flow_array])
                current_row_num += 1

            #add last <EMD>
            flow_list_array.append([current_row_num, ['<END>']])
            #print(flow_list_array)

            #print(flow_list_array)
            ### Convert to tuple
            master_device_table_tuple = ns_def.convert_array_to_tuple(flow_list_array)

            ''' 
            Create temp input data file
            '''
            # List of characters not allowed in Excel worksheet names
            forbidden_chars = [':', '\\', '/', '?', '*', '[', ']']
            # Remove forbidden characters using a list comprehension
            cleaned_device_name = ''.join(char for char in device_name if char not in forbidden_chars)


            ### Create the flow table excel file or add sheet
            self.worksheet_name = cleaned_device_name
            if flag_flow_table_exist == False:
                wb = openpyxl.Workbook()
                sheet = wb.active
                sheet.title = self.worksheet_name
                wb.save(self.excel_flow_file)
            else:
                wb = openpyxl.load_workbook(self.excel_flow_file)
                if self.worksheet_name in wb.sheetnames:
                    # Remove the existing worksheet
                    sheet_to_remove = wb[self.worksheet_name]
                    wb.remove(sheet_to_remove)

                wb.create_sheet(title=self.worksheet_name)
                wb.save(self.excel_flow_file)

            '''
            Create [FLOW_REPORT] file
            '''
            tmp_master_data_array = []
            tmp_master_data_array.append([1, [self.worksheet_name]])
            #print(tmp_master_data_array)

            template_master_data_tuple = {}
            template_master_data_tuple = ns_def.convert_array_to_tuple(tmp_master_data_array)

            #print('Create --- template_master_data_tuple---')
            #print(template_master_data_tuple)
            offset_row = 0
            offset_column = 0
            write_to_section = '_template_'
            ns_def.write_excel_meta(template_master_data_tuple, self.excel_flow_file, self.worksheet_name, write_to_section, offset_row, offset_column)

            ###
            input_excel_name = self.excel_flow_file
            output_excel_name = self.outFileTxt_11_3.get()
            if flag_flow_table_exist == False:
                NEW_OR_ADD = 'NEW'
            else:
                NEW_OR_ADD = 'ADD_OPTION1'
            ns_egt_maker.create_excel_gui_tree(input_excel_name,output_excel_name,NEW_OR_ADD, egt_maker_width_array)

            '''
            Add FLOW_List table from meta
            '''
            # Write normal tuple to excel
            tmp_ws_name = '_tmp_'
            master_excel_meta = master_device_table_tuple
            ppt_meta_file = output_excel_name
            excel_file_path = ppt_meta_file
            worksheet_name = tmp_ws_name
            section_write_to = '<<N/A>>'
            offset_row = 0
            offset_column = 0
            ns_def.create_excel_sheet(ppt_meta_file, tmp_ws_name)
            ns_def.write_excel_meta(master_excel_meta, excel_file_path, worksheet_name, section_write_to, offset_row, offset_column)

            #print(output_excel_name)
            self.input_tree_excel = openpyxl.load_workbook(output_excel_name)
            worksheet_name = cleaned_device_name
            start_row = 1
            start_column = 0
            custom_table_name = ppt_meta_file
            self.input_tree_excel = ns_egt_maker.insert_custom_excel_table(self.input_tree_excel, worksheet_name, start_row, start_column, custom_table_name)
            self.input_tree_excel.save(output_excel_name)

            # Remove _tmp_ sheet from excel master self.worksheet_name
            ns_def.remove_excel_sheet(ppt_meta_file, tmp_ws_name)


class flow():
    def add_routing_path_to_flow(self,full_filepath_master,flow_list_array):
        print('--- Routing path calculation ---')
        argv_array = ['show', 'l3_broadcast_domain']
        l3_broadcast_array = ns_cli.ns_cli_run.cli_show(self, full_filepath_master, argv_array)
        #print(l3_broadcast_array)

        G = nx.Graph()
        # Add nodes and edges to the graph
        for domain_info in l3_broadcast_array:
            broadcast_domain, device_interfaces = domain_info
            devices = [dev[0] for dev in device_interfaces]

            # Connect all devices in the broadcast domain (full graph)
            for i in range(len(devices)):
                for j in range(i + 1, len(devices)):
                    G.add_edge(devices[i], devices[j])

        for row in flow_list_array[2:]:
            data = row[1]
            while len(data) <= 7:
                data.append('')
            source = data[1]
            target = data[2]

            path = flow.get_shortest_path(G, source, target)
            #print(path)
            if 'is not in G' in path:
                continue

            if len(path) >= 2:
                path = path[1:-1]
            path2 = ', '.join([f"'{p}'" for p in path])
            if path2 == '':
                path2 = ' '

            data[7] = path2

        #print(flow_list_array)
        return (flow_list_array)

    def get_shortest_path(G,source, target):
        try:
            path = nx.shortest_path(G, source=source, target=target)
            #print(source,target,path, G)
            return path
        except nx.NetworkXNoPath:
            return f"The Path from {source} to {target} could not be found."
        except nx.NodeNotFound as e:
            return str(e)

    def append_flows_to_diagram(self,variable3_7_y_1,variable3_7_y_2,variable3_7_y_3): #add at ver 2.4.3
        print('--- append_flows_to_diagram ---')
        #print(variable3_7_y_1.get(), variable3_7_y_2.get(), variable3_7_y_3.get())
        #print(self.pptx_full_filepath)
        #print(self.full_filepath)

        ws_flow_name = 'Flow_Data'
        excel_maseter_file = self.full_filepath
        master_flow_array = ns_def.convert_excel_to_array(ws_flow_name, excel_maseter_file, 3)

        # Exclude the last element (['<<END_MARK>>'])
        filtered_master_flow = master_flow_array[:-1]
        #print(filtered_master_flow)

        # Exclude invalid lines
        filterd2_master_flow = []
        for element in filtered_master_flow:
            sublist = element[1]  # Take the sublist
            # Check if the second or third elements are not empty
            if len(sublist) >= 2:
                if sublist[1] and sublist[2]:
                    filterd2_master_flow.append(sublist)  # Append the sublist without the first number
        #print(filterd2_master_flow)

        # Filter lines by the target
        filtered_target_flow = []
        for element in filterd2_master_flow:
            # Check if the criteria are met
            if (element[1] == variable3_7_y_1.get() or variable3_7_y_1.get() == 'Any') and \
                    (element[2] == variable3_7_y_2.get() or variable3_7_y_2.get() == 'Any') and \
                    (element[4] == variable3_7_y_3.get() or variable3_7_y_3.get() == 'Any'):
                filtered_target_flow.append(element)

        #print(filtered_target_flow)

        '''read the pptx file and shapes data'''
        self.shape_name_grid_array = []

        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation(self.pptx_full_filepath)
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                if hasattr(shape, "adjustments"):
                    try:
                        if shape.adjustments[0] not in [0.99445, 0.50444, 0.30045, 0.00046, 0.15005, 0.00057, 0.2007] and shape.text.strip() != '': #exclude not (device or wp) shape. 0.0001 device, 0.0008 L3 instance device , 0.2002 way point, 0.2007 l3 instance in device
                            #print(shape.text.strip(), shape.adjustments[0] , shape.left, shape.top, shape.width, shape.height)
                            self.shape_name_grid_array.append([shape.text.strip(), shape.left, shape.top, shape.width, shape.height])

                    except IndexError:
                        pass

        #print(self.shape_name_grid_array)

        ''' deside routes and write flows '''
        if os.path.isfile(self.pptx_full_filepath) == True:
            self.active_ppt = Presentation(self.pptx_full_filepath)
            self.slide = self.active_ppt.slides[0]

        for tmp_filtered_target_flow in filtered_target_flow:

            # select routing path is auto or static
            selected_route_path = []
            if tmp_filtered_target_flow[6] == '' and tmp_filtered_target_flow[7] == ' ':
                selected_route_path = [tmp_filtered_target_flow[1], tmp_filtered_target_flow[2]]

            elif tmp_filtered_target_flow[6] == '' and tmp_filtered_target_flow[7] != ' ':
                selected_route_path = [element.strip().strip("'") for element in tmp_filtered_target_flow[7].split(',')]
                selected_route_path.insert(0, tmp_filtered_target_flow[1])
                selected_route_path.append(tmp_filtered_target_flow[2])

            elif tmp_filtered_target_flow[6] != '':
                selected_route_path = [element.strip().strip("'") for element in tmp_filtered_target_flow[6].split(',')]
                selected_route_path.insert(0, tmp_filtered_target_flow[1])
                selected_route_path.append(tmp_filtered_target_flow[2])

            #print(tmp_filtered_target_flow,selected_route_path)

            ''' write line'''
            if len(selected_route_path) == 2:
                # get source grid value
                source_grid = next(
                    (item for item in self.shape_name_grid_array if item[0] == tmp_filtered_target_flow[1]),
                    None  # Returns None if no match is found.
                )

                destination_grid = next(
                    (item for item in self.shape_name_grid_array if item[0] == tmp_filtered_target_flow[2]),
                    None  # Returns None if no match is found.
                )

                if source_grid == None or destination_grid == None:
                    continue

                #print(source_grid, destination_grid)

                line_type = 'FLOW0'
                inche_from_connect_x = (source_grid[1] + source_grid[3] * 0.25) / 914400
                inche_from_connect_y = (source_grid[2] + source_grid[4] * 0.5) / 914400
                inche_to_connect_x = (destination_grid[1] + destination_grid[3] * 0.75) / 914400
                inche_to_connect_y = (destination_grid[2] + destination_grid[4] * 0.5)  / 914400
                ns_ddx_figure.extended.add_line(self,line_type,inche_from_connect_x,inche_from_connect_y,inche_to_connect_x,inche_to_connect_y)

            elif len(selected_route_path) >= 3:
                #print(selected_route_path)
                for i in range(len(selected_route_path) - 1):
                    pair = [selected_route_path[i], selected_route_path[i + 1]]
                    #print(pair,i,len(selected_route_path) - 2 )

                    # get source grid value
                    source_grid = next(
                        (item for item in self.shape_name_grid_array if item[0] == pair[0]),
                        None  # Returns None if no match is found.
                    )

                    destination_grid = next(
                        (item for item in self.shape_name_grid_array if item[0] == pair[1]),
                        None  # Returns None if no match is found.
                    )

                    if source_grid == None or destination_grid == None:
                        continue

                    if i == 0:
                        line_type = 'FLOW1'
                        inche_from_connect_x = (source_grid[1] + source_grid[3] * 0.25) / 914400
                        inche_from_connect_y = (source_grid[2] + source_grid[4] * 0.5) / 914400
                        inche_to_connect_x = (destination_grid[1] + destination_grid[3] * 0.5) / 914400
                        inche_to_connect_y = (destination_grid[2] + destination_grid[4] * 0.5) / 914400
                        ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y,inche_to_connect_x, inche_to_connect_y)
                    elif i == len(selected_route_path) - 2:
                        line_type = 'FLOW2'
                        inche_from_connect_x = (source_grid[1] + source_grid[3] * 0.5) / 914400
                        inche_from_connect_y = (source_grid[2] + source_grid[4] * 0.5) / 914400
                        inche_to_connect_x = (destination_grid[1] + destination_grid[3] * 0.75) / 914400
                        inche_to_connect_y = (destination_grid[2] + destination_grid[4] * 0.5) / 914400
                        ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y,inche_to_connect_x, inche_to_connect_y)
                    else:
                        line_type = 'FLOW3'
                        inche_from_connect_x = (source_grid[1] + source_grid[3] * 0.5) / 914400
                        inche_from_connect_y = (source_grid[2] + source_grid[4] * 0.5) / 914400
                        inche_to_connect_x = (destination_grid[1] + destination_grid[3] * 0.5) / 914400
                        inche_to_connect_y = (destination_grid[2] + destination_grid[4] * 0.5) / 914400
                        ns_ddx_figure.extended.add_line(self, line_type, inche_from_connect_x, inche_from_connect_y,inche_to_connect_x, inche_to_connect_y)

        folder = os.path.dirname(self.pptx_full_filepath)
        filename = os.path.basename(self.pptx_full_filepath)
        modified_filepath = os.path.join(folder, f"Added_flows_{filename}")
        self.active_ppt.save(modified_filepath)

        #file open
        ns_def.messagebox_file_open(modified_filepath)

    def get_flow_item_list(self): # add at ver 2.4.3
        #print('--- get_flow_item_list ---')
        excel_maseter_file = self.inFileTxt_L2_3_1.get()
        ## check Flow_Data sheet exists in Master file
        input_excel_master = openpyxl.load_workbook(excel_maseter_file)
        ws_list_master = input_excel_master.sheetnames
        input_excel_master.close()

        ws_flow_name = 'Flow_Data'
        if ws_flow_name in ws_list_master:
            master_flow_array = ns_def.convert_excel_to_array(ws_flow_name, excel_maseter_file, 3)

            # Exclude the last element (['<<END_MARK>>'])
            filtered_master_flow = master_flow_array[:-1]

            # Group elements from the 2nd, 3rd, 4th, and 5th positions into separate lists
            category_wise_data = [[] for _ in range(4)]  # Prepare 4 category lists

            for entry in filtered_master_flow:
                data = entry[1]  # Extract the second element (list)
                for i in range(4):  # Process the 2nd, 3rd, 4th, and 5th elements (index 1 to 4)
                    if len(data) >= 4:
                        value = data[i + 1].strip()
                        if value and value not in category_wise_data[i]:  # Add only non-empty, non-duplicate values
                            category_wise_data[i].append(value)

            update_master_flow_array = category_wise_data

            return(update_master_flow_array)

        else:
            print('--- Master file does not have Flow_Data sheet ---')


    def export_flow_file(self,dummy):
        print('--- export_flow_file ---')

        excel_maseter_file = self.inFileTxt_L2_3_1.get()
        iDir = os.path.abspath(os.path.dirname(excel_maseter_file))

        basename_without_ext = os.path.splitext(os.path.basename(excel_maseter_file))[0]
        self.outFileTxt_11_3.delete(0, tkinter.END)
        self.outFileTxt_11_3.insert(tk.END, iDir + ns_def.return_os_slash() + '[FLOW]' + basename_without_ext.replace('[MASTER]', '') + '.xlsx')

        ## check file open
        ns_def.check_file_open(self.outFileTxt_11_3.get())

        # remove exist flow file
        if os.path.isfile(self.outFileTxt_11_3.get()) == True:
            os.remove(self.outFileTxt_11_3.get())

        self.excel_flow_file = self.outFileTxt_11_3.get()

        ## check Flow_Data sheet exists in Master file
        input_excel_master = openpyxl.load_workbook(excel_maseter_file)
        ws_list_master = input_excel_master.sheetnames
        input_excel_master.close()

        ws_flow_name = 'Flow_Data'
        flag_master_has_flow_sheet = False
        if ws_flow_name in ws_list_master:
            flag_master_has_flow_sheet = True
            master_flow_array = []
            master_flow_array = ns_def.convert_excel_to_array(ws_flow_name, excel_maseter_file, 3)
            if '<<END_MARK>>' in master_flow_array[-1][1]:
                master_flow_array = master_flow_array[:-1]
            #print(master_flow_array)

        '''
        MAKE Flows List
        '''
        master_device_table_tuple = {}
        flow_list_array = []
        egt_maker_width_array = ['5','25', '25','15', '20', '25', '40', '40']  # for Network Sketcher Ver 2.0
        flow_list_array.append([1, ['<RANGE>', '1','1', '1', '1', '1', '1', '1', '1', '<END>']])
        flow_list_array.append([2, ['<HEADER>', 'No','Source Device Name', 'Destination Device Name','TCP/UDP/ICMP','Service name(Port)',  'Max. bandwidth(Mbps)', 'Manually rouging path settings', 'Automatic rouging path settings', '<END>']])

        current_row_num = 3
        all_empty = False
        if flag_master_has_flow_sheet == True:
            # check last ten column = empty
            last_10_elements = [item[1] for item in master_flow_array[-1:]]
            all_empty = all(all(element == '' for element in item[1:7]) for item in last_10_elements)

            for tmp_master_flow_array in master_flow_array:
                for i in range(1, 8):
                    while len(tmp_master_flow_array[1]) <= 7:
                        tmp_master_flow_array[1].append('')

                    if tmp_master_flow_array[1][i] == '':
                        tmp_master_flow_array[1][i] = '<EMPTY>'

                flow_list_array.append([current_row_num, ['',str(current_row_num - 2), '>>' + str(tmp_master_flow_array[1][1]), '>>' + str(tmp_master_flow_array[1][2]), '>>' + str(tmp_master_flow_array[1][3]), '>>' + str(tmp_master_flow_array[1][4]), '>>' + str(tmp_master_flow_array[1][5]), '>>' + str(tmp_master_flow_array[1][6]),str(tmp_master_flow_array[1][7]), '<END>']])
                current_row_num += 1

        if flag_master_has_flow_sheet == True and all_empty == True:
            add_column_num = 0
        elif flag_master_has_flow_sheet == True and all_empty == False:
            add_column_num = 10
        else:
            add_column_num = 100

        current_row_max = add_column_num + current_row_num
        for n in range(current_row_num, current_row_max):
            flow_list_array.append([n, ['',str(n - 2), '<EMPTY>', '<EMPTY>', '<EMPTY>', '<EMPTY>', '<EMPTY>', '<EMPTY>',' ', '<END>']])

        flow_list_array.append([current_row_max, ['<END>']])

        #print(flow_list_array)
        ### Convert to tuple
        master_device_table_tuple = ns_def.convert_array_to_tuple(flow_list_array)

        ''' 
        Create temp input data file
        '''
        ### Create new data excel file
        self.worksheet_name = 'Flow_List'
        wb = openpyxl.Workbook()
        sheet = wb.active
        sheet.title = self.worksheet_name
        wb.save(self.excel_flow_file)

        '''
        Create [FLOW] file
        '''

        tmp_master_data_array = []
        tmp_master_data_array.append([1,[self.worksheet_name]])
        #tmp_master_data_array.append([2,[self.worksheet_name]])
        #print(tmp_master_data_array)


        template_master_data_tuple = {}
        template_master_data_tuple = ns_def.convert_array_to_tuple(tmp_master_data_array)

        #print('Create --- template_master_data_tuple---')
        #print(template_master_data_tuple)
        offset_row = 0
        offset_column = 0
        write_to_section = '_template_'
        ns_def.write_excel_meta(template_master_data_tuple, self.excel_flow_file, self.worksheet_name, write_to_section, offset_row, offset_column)

        ###
        input_excel_name = self.excel_flow_file
        output_excel_name = self.outFileTxt_11_3.get()
        NEW_OR_ADD = 'NEW'
        ns_egt_maker.create_excel_gui_tree(input_excel_name,output_excel_name,NEW_OR_ADD, egt_maker_width_array)

        '''
        Add FLOW_List table from meta
        '''
        # Write normal tuple to excel
        tmp_ws_name = '_tmp_'
        master_excel_meta = master_device_table_tuple
        ppt_meta_file = output_excel_name
        excel_file_path = ppt_meta_file
        worksheet_name = tmp_ws_name
        section_write_to = '<<N/A>>'
        offset_row = 0
        offset_column = 0
        ns_def.create_excel_sheet(ppt_meta_file, tmp_ws_name)
        ns_def.write_excel_meta(master_excel_meta, excel_file_path, worksheet_name, section_write_to, offset_row, offset_column)

        #print(output_excel_name)
        self.input_tree_excel = openpyxl.load_workbook(output_excel_name)
        worksheet_name = 'Flow_List'
        start_row = 1
        start_column = 0
        custom_table_name = ppt_meta_file
        self.input_tree_excel = ns_egt_maker.insert_custom_excel_table(self.input_tree_excel, worksheet_name, start_row, start_column, custom_table_name)
        self.input_tree_excel.save(output_excel_name)

        '''
        Add Drop list
        '''
        from openpyxl.worksheet.datavalidation import DataValidation
        # Load the Excel file
        wb = openpyxl.load_workbook(output_excel_name)
        ws = wb['Flow_List']  # Select the worksheet 'Flow_List'

        # Create a dropdown list (Enable in-cell dropdown)
        dv2 = DataValidation(type="list", formula1='"TCP,UDP,ICMP"', allow_blank=True, showDropDown=False)
        dv3 = DataValidation(type="list", formula1='"Any,FTP Data(20),FTP Control(21),SSH(22),Telnet(23),SMTP(25),DNS(53),DHCP Server(67),DHCP Client(68),HTTP(80),NNTP(119),NTP(123),IMAP(143),SNMP(161),SNMP Trap(162),BGP(179),HTTPS(443),SMB(445),SMTPS(465),SMTP(587),IMAPS(993),RDP(3389)"', allow_blank=True, showDropDown=False)

        # Apply data validation to cell C3
        row = 3
        for n in range(row, 103):
            column = 4
            ws.add_data_validation(dv2)
            dv2.add(ws.cell(row=n, column=column))
            column = 5
            ws.add_data_validation(dv3)
            dv3.add(ws.cell(row=n, column=column))

        # Save the flow file
        output_file = output_excel_name
        wb.save(output_file)

        print(f"Flow file is saved: {output_file}")

        # Remove _tmp_ sheet from excel master
        ns_def.remove_excel_sheet(ppt_meta_file, tmp_ws_name)

class  ip_report():
    def export_ip_report(self,dummy):
        print('--- export_ip_report ---')
        excel_maseter_file = self.inFileTxt_L2_3_1.get()
        iDir = os.path.abspath(os.path.dirname(excel_maseter_file))

        # SET IP Address report file patch
        basename_without_ext = os.path.splitext(os.path.basename(excel_maseter_file))[0]
        self.outFileTxt_11_3.delete(0, tkinter.END)
        self.outFileTxt_11_3.insert(tk.END, iDir + ns_def.return_os_slash() + '[IP_REPORT]' + basename_without_ext.replace('[MASTER]', '') + '.xlsx') #change IP_TABLE to IP_REPORT at ver 2.5.1
        self.excel_file_path = iDir + ns_def.return_os_slash() + '_template_[IP_REPORT]' + basename_without_ext.replace('[MASTER]', '') + '.xlsx'  #change IP_TABLE to IP_REPORT at ver 2.5.1

        ## check file open
        ns_def.check_file_open(self.outFileTxt_11_3.get())

        # remove exist ip table file
        if os.path.isfile(self.outFileTxt_11_3.get()) == True:
            os.remove(self.outFileTxt_11_3.get())

        self.excel_file_path = self.outFileTxt_11_3.get()

        '''
        MAKE IP Address List
        '''
        master_device_table_tuple = {}
        ip_address_list_array = []
        egt_maker_width_array = ['20', '20', '20', '20', '25', '15', '20']  # for Network Sketcher Ver 2.0
        ip_address_list_array.append([1, ['<RANGE>', '1', '1', '1', '1', '1', '1', '1', '<END>']])
        ip_address_list_array.append([2, ['<HEADER>', 'IP Address', 'Mask', 'Network Address', 'Device Name', 'L3 IF Name', 'L3 Instance', 'Area', '<END>']])

        current_row_num = 3

        kari_ip_address_list_array = []
        l3_segment_group_array = ns_def.get_l3_segments(self)
        #print(l3_segment_group_array)
        tmp_seg_array = []
        for tmp_l3_segment_group_array in l3_segment_group_array:
            #print(tmp_l3_segment_group_array)

            for tmp_tmp_l3_segment_group_array in tmp_l3_segment_group_array:
                tmp_seg_array.append([tmp_tmp_l3_segment_group_array[0],tmp_tmp_l3_segment_group_array[4]])

                ip_with_subnet = tmp_tmp_l3_segment_group_array[4]
                ip_address = '[None]'
                subnet_mask = '[None]'
                network_address = '[None]'

                L3_instance = ' '
                if tmp_tmp_l3_segment_group_array[3] != '':
                    L3_instance = tmp_tmp_l3_segment_group_array[3]

                if '/' in str(ip_with_subnet):
                    network = ipaddress.ip_network(ip_with_subnet, strict=False)
                    ip_interface = ipaddress.ip_interface(ip_with_subnet)
                    ip_address = str(ip_interface.ip)
                    subnet_mask = str(ip_interface.netmask)
                    ip_address_dummy, prefix_length = ip_with_subnet.split('/')
                    network_address = str(network.network_address) + str('/') + str(prefix_length)
                    numeric_sequence = ''.join(f'{int(octet):03}' for octet in ip_address.split('.'))
                if ip_address == '[None]':
                    numeric_sequence = str(255255255255)
                kari_ip_address_list_array.append([numeric_sequence,ip_address,subnet_mask,network_address,tmp_tmp_l3_segment_group_array[1],tmp_tmp_l3_segment_group_array[2],L3_instance,tmp_tmp_l3_segment_group_array[0], '<END>'])
        
        # Remove completely duplicate columns at ver 2.2.1(c)
        unique_tuples_set = set(tuple(item) for item in kari_ip_address_list_array)
        unique_list = [list(item) for item in unique_tuples_set]
        unique_array = np.array(unique_list)
        sorted_lists = sorted(unique_array, key=lambda x: x[0], reverse=False)
        
        #print(sorted_lists)

        for tmp_sorted_lists in sorted_lists:
            tmp_sorted_lists[0] = ''
            ip_address_list_array.append([current_row_num,tmp_sorted_lists])
            current_row_num += 1

        ip_address_list_array.append([current_row_num, ['<END>']])

        #print(ip_address_list_array)
        ### Convert to tuple
        master_device_table_tuple = ns_def.convert_array_to_tuple(ip_address_list_array)

        '''
        MAKE Summary
        '''
        summary_list_master_device_table_tuple = {}
        summary_list_array = []
        summary_list_array.append([1, ['<RANGE>', '1', '1','<END>']])
        summary_list_array.append([2, ['<HEADER>', 'Area', 'Summary(CIDR)',  '<END>']])

        area_list = ip_report.get_folder_list(self)
        current_row_num = 3

        #print(tmp_seg_array)
        get_folder = ip_report.get_folder_list(self)
        #print(get_folder)

        for tmp_get_folder in get_folder:
            kari_sum_array = []
            for tmp_tmp_seg_array in tmp_seg_array:
                if tmp_tmp_seg_array[0] == tmp_get_folder and tmp_tmp_seg_array[1] != '':
                    kari_sum_array.append(tmp_tmp_seg_array[1])

            #print(kari_sum_array)
            networks = [ipaddress.ip_network(cidr, strict=False) for cidr in kari_sum_array]

            # clac summary
            summary_address = ipaddress.collapse_addresses(networks)
            summary_address_list = [str(network) for network in summary_address]
            #print(tmp_get_folder,str(summary_address_list))

            first_area_flag = True
            for tmp_summary_address_list in summary_address_list:
                if first_area_flag == True:
                    summary_list_array.append([current_row_num, ['', tmp_get_folder, str(tmp_summary_address_list), '<END>']])
                    current_row_num += 1
                    first_area_flag = False
                else:
                    summary_list_array.append([current_row_num, ['', '', str(tmp_summary_address_list), '<END>']])
                    current_row_num += 1

        summary_list_array.append([current_row_num, ['<END>']])

        #print(summary_list_array)

        ### Convert to tuple
        master_summary_table_tuple = ns_def.convert_array_to_tuple(summary_list_array)

        ''' 
        Create temp input data file
        '''
        ### Create new data excel file
        self.worksheet_name = 'IP Address_List'
        wb = openpyxl.Workbook()
        sheet = wb.active
        sheet.title = self.worksheet_name
        wb.save(self.excel_file_path)

        '''
        Create [IP Address] file
        '''

        tmp_master_data_array = []
        tmp_master_data_array.append([1,['Summary']])
        tmp_master_data_array.append([2,[self.worksheet_name]])
        #print(tmp_master_data_array)


        template_master_data_tuple = {}
        template_master_data_tuple = ns_def.convert_array_to_tuple(tmp_master_data_array)

        #print('Create --- template_master_data_tuple---')
        #print(template_master_data_tuple)
        offset_row = 0
        offset_column = 0
        write_to_section = '_template_'
        ns_def.write_excel_meta(template_master_data_tuple, self.excel_file_path, self.worksheet_name, write_to_section, offset_row, offset_column)

        ###
        input_excel_name = self.excel_file_path
        output_excel_name = self.outFileTxt_11_3.get()
        NEW_OR_ADD = 'NEW'
        ns_egt_maker.create_excel_gui_tree(input_excel_name,output_excel_name,NEW_OR_ADD, egt_maker_width_array)

        '''
        Add IP Address_List table from meta
        '''
        # Write normal tuple to excel
        tmp_ws_name = '_tmp_'
        master_excel_meta = master_summary_table_tuple
        ppt_meta_file = output_excel_name
        excel_file_path = ppt_meta_file
        worksheet_name = tmp_ws_name
        section_write_to = '<<N/A>>'
        offset_row = 0
        offset_column = 0
        ns_def.create_excel_sheet(ppt_meta_file, tmp_ws_name)
        ns_def.write_excel_meta(master_excel_meta, excel_file_path, worksheet_name, section_write_to, offset_row, offset_column)

        #print(output_excel_name)
        self.input_tree_excel = openpyxl.load_workbook(output_excel_name)
        worksheet_name = 'Summary'
        start_row = 1
        start_column = 0
        custom_table_name = ppt_meta_file
        self.input_tree_excel = ns_egt_maker.insert_custom_excel_table(self.input_tree_excel, worksheet_name, start_row, start_column, custom_table_name)
        self.input_tree_excel.save(output_excel_name)

        # Remove _tmp_ sheet from excel master
        ns_def.remove_excel_sheet(ppt_meta_file, tmp_ws_name)

        '''
        Add Summary table from meta
        '''

        # Write normal tuple to excel
        tmp_ws_name = '_tmp_'
        master_excel_meta = master_device_table_tuple
        ppt_meta_file = output_excel_name
        excel_file_path = ppt_meta_file
        worksheet_name = tmp_ws_name
        section_write_to = '<<N/A>>'
        offset_row = 0
        offset_column = 0
        ns_def.create_excel_sheet(ppt_meta_file, tmp_ws_name)
        ns_def.write_excel_meta(master_excel_meta, excel_file_path, worksheet_name, section_write_to, offset_row, offset_column)

        #print(output_excel_name)
        self.input_tree_excel = openpyxl.load_workbook(output_excel_name)
        worksheet_name = 'IP Address_List'
        start_row = 1
        start_column = 0
        custom_table_name = ppt_meta_file
        self.input_tree_excel = ns_egt_maker.insert_custom_excel_table(self.input_tree_excel, worksheet_name, start_row, start_column, custom_table_name)
        self.input_tree_excel.save(output_excel_name)

        # Remove _tmp_ sheet from excel master
        ns_def.remove_excel_sheet(ppt_meta_file, tmp_ws_name)


    def get_folder_list(self):
        #print('--- get_folder_list ---')
        #parameter
        ws_name = 'Master_Data'
        excel_maseter_file = self.inFileTxt_L2_3_1.get()

        # GET Folder and wp name List
        self.folder_wp_name_array = ns_def.get_folder_wp_array_from_master(ws_name, excel_maseter_file)
        #print('---- folder_wp_name_array ----')
        #print(self.folder_wp_name_array)

        return_array = self.folder_wp_name_array[0]
        return_array.sort(reverse=False)
        #if len(self.folder_wp_name_array[1]) >= 1:
        #    return_array.append("_WAN(Way_Point)_")

        return return_array

class  auto_ip_addressing():
    def get_folder_list(self):
        #print('--- get_folder_list ---')
        #parameter
        ws_name = 'Master_Data'
        excel_maseter_file = self.inFileTxt_L2_3_1.get()

        # GET Folder and wp name List
        self.folder_wp_name_array = ns_def.get_folder_wp_array_from_master(ws_name, excel_maseter_file)
        #print('---- folder_wp_name_array ----')
        #print(self.folder_wp_name_array)

        return_array = self.folder_wp_name_array[0]
        return_array.sort(reverse=False)
        if len(self.folder_wp_name_array[1]) >= 1:
            return_array.append("_WAN(Way_Point)_")
        return return_array

    def get_auto_ip_param(self,target_area_name):
        #print('--- get_auto_ip_param ---')
        #print(target_area_name)

        if target_area_name == "_WAN(Way_Point)_":
            target_area_name = 'N/A'

        '''get values of Master Data'''
        #parameter
        ws_name = 'Master_Data'
        ws_l2_name = 'Master_Data_L2'
        ws_l3_name = 'Master_Data_L3'
        excel_maseter_file = self.inFileTxt_L3_3_1.get()

        self.result_get_l2_broadcast_domains =  ns_def.get_l2_broadcast_domains.run(self,excel_maseter_file)  ## 'self.update_l2_table_array, device_l2_boradcast_domain_array, device_l2_directly_l3vport_array, device_l2_other_array, marged_l2_broadcast_group_array'

        #print('--- self.update_l2_table_array ---')
        #print(self.result_get_l2_broadcast_domains[0])
        #print('--- self.device_l2_boradcast_domain_array ---')
        #print(self.result_get_l2_broadcast_domains[1])
        #self.device_l2_boradcast_domain_array = self.result_get_l2_broadcast_domains[1]
        #print('--- device_l2_directly_l3vport_array ---')
        #print(self.result_get_l2_broadcast_domains[2])
        #self.device_l2_directly_l3vport_array = self.result_get_l2_broadcast_domains[2]
        #print('--- device_l2_other_array ---')
        #print(self.result_get_l2_broadcast_domains[3])
        self.device_l2_other_array = self.result_get_l2_broadcast_domains[3]
        #print('--- marged_l2_broadcast_group_array ---')
        #print(self.result_get_l2_broadcast_domains[4])
        self.marged_l2_broadcast_group_array = self.result_get_l2_broadcast_domains[4]
        #print('--- self.target_l2_broadcast_group_array ---')
        #print(self.target_l2_broadcast_group_array)

        self.l3_table_array = ns_def.convert_master_to_array(ws_l3_name, excel_maseter_file, '<<L3_TABLE>>')
        #print('--- self.l3_table_array  ---')
        #print(self.l3_table_array )

        # check ip address exists in target area
        flag_no_ipaddress = True
        add_l3_table_array = []
        for index, tmp_l3_table_array  in enumerate(self.l3_table_array):
            str(tmp_l3_table_array).replace(' ', '')
            if index >= 2:
                if tmp_l3_table_array[1][0] == target_area_name:
                    #print(tmp_l3_table_array[1])
                    if len(tmp_l3_table_array[1]) == 5:
                        flag_no_ipaddress = False

                if len(tmp_l3_table_array[1]) == 5:
                    if ',' in str(tmp_l3_table_array[1][4]):
                        #print('--- tmp_l3_table_array ', str(tmp_l3_table_array))
                        tmp_tmp_l3_table_array= str(tmp_l3_table_array[1][4]).split(',')
                        for tmp_add_array in tmp_tmp_l3_table_array:
                            tmp_tmp_tmp_l3_table_array = tmp_l3_table_array
                            tmp_tmp_tmp_l3_table_array[1][4] = tmp_add_array
                            #print('--- tmp_tmp_tmp_l3_table_array ', tmp_tmp_tmp_l3_table_array)
                            self.l3_table_array.append([tmp_tmp_tmp_l3_table_array[0],[tmp_tmp_tmp_l3_table_array[1][0],tmp_tmp_tmp_l3_table_array[1][1],tmp_tmp_tmp_l3_table_array[1][2],tmp_tmp_tmp_l3_table_array[1][3],tmp_tmp_tmp_l3_table_array[1][4]]])

        #print ('flag_no_ipaddress', str(flag_no_ipaddress))

        outside_ip_address_list = []
        inside_ip_address_list = []
        full_ip_address_list = []
        for index, tmp_l3_table_array in enumerate(self.l3_table_array):
            if index >= 2:
                if len(tmp_l3_table_array[1]) == 5 and ',' not in str(tmp_l3_table_array[1][4]):
                    full_ip_address_list.append(str(tmp_l3_table_array[1][4]).replace(' ', ''))

                if tmp_l3_table_array[1][0] != target_area_name:
                    if len(tmp_l3_table_array[1]) == 5 and ',' not in str(tmp_l3_table_array[1][4]):
                        # print(tmp_l3_table_array[1][4])
                        first_octet = int(str(tmp_l3_table_array[1][4]).split('.')[0])
                        second_octet = int(str(tmp_l3_table_array[1][4]).split('.')[1])
                        third_octet = int(str(tmp_l3_table_array[1][4]).split('.')[2])

                        outside_ip_address_list.append(str(first_octet) + '.' + str(second_octet)+ '.' + str(third_octet) + '.0')
                else:
                    if len(tmp_l3_table_array[1]) == 5 and ',' not in str(tmp_l3_table_array[1][4]):
                        # print(tmp_l3_table_array[1][4])
                        first_octet = int(str(tmp_l3_table_array[1][4]).split('.')[0])
                        second_octet = int(str(tmp_l3_table_array[1][4]).split('.')[1])
                        third_octet = int(str(tmp_l3_table_array[1][4]).split('.')[2])

                        inside_ip_address_list.append(str(first_octet) + '.' + str(second_octet) + '.' + str(third_octet) + '.0')

        if flag_no_ipaddress == True:
            current_ip_address_list = outside_ip_address_list
        else:
            current_ip_address_list = inside_ip_address_list

        #print(current_ip_address_list )
        if current_ip_address_list != []:
            word_counts = Counter(current_ip_address_list)
            most_common_word, most_common_count = word_counts.most_common(1)[0]
            print(f"--- most_common_word: {most_common_word} (most_common_count: {most_common_count})")
        else:
            most_common_word = '10.0.0.0'

        use_network = ''
        '''get starting ip address'''
        # set starting ip address
        start_ip = ipaddress.IPv4Address(most_common_word)

        # count for 192 , 172 , 10
        if most_common_word.startswith('192.168.'):
            increase_count = 256
        elif most_common_word.startswith('172.'):
            for i in range(16, 32):
                if most_common_word.startswith('172.' + str(i) + '.'):
                    increase_count = 256 * 16
        elif most_common_word.startswith('10.'):
            increase_count = 256 * 256
        else:
            start_ip = ipaddress.IPv4Address('10.0.0.0')
            increase_count = 256 * 256

        # output start network address(CIDR)
        flag_1st_third_octet = True

        for _ in range(increase_count):
            # Convert IP address to byte array
            ip_bytes = start_ip.packed
            # Get the third octet and increase by 1
            third_octet = ip_bytes[2] + 1
            # If the third octet exceeds 255, the second octet is also increased
            if third_octet > 255:
                second_octet = ip_bytes[1] + 1
                third_octet = 0  # Omitted if the second octet exceeds 255
            else:
                second_octet = ip_bytes[1]

            if flag_1st_third_octet == True:
                third_octet -= 1
                flag_1st_third_octet = False

            # Build a new IP address
            new_ip_bytes = bytearray(ip_bytes)
            new_ip_bytes[1] = second_octet
            new_ip_bytes[2] = third_octet
            start_ip = ipaddress.IPv4Address(bytes(new_ip_bytes))
            #print(start_ip)

            if len(full_ip_address_list) == 0:
                use_network = '10.0.0.0/24'
                break

            flag_found_use_network = True
            for tmp_full_ip_address_list in full_ip_address_list:

                # Define a network to determine
                network1 = ipaddress.ip_network(str(start_ip) +'/24')
                network2 = ipaddress.ip_network(str(tmp_full_ip_address_list), strict=False)

                # Determine if network1 and network2 overlap
                if network1.overlaps(network2):
                    #print(f"{network1} overlaps with {network2}")
                    flag_found_use_network = False
                    break
                else:
                    #print(f"{network1} does not overlap with {network2}")
                    use_network = network1

            if flag_found_use_network == True:
                break

        if flag_no_ipaddress == False:
            use_network = str(most_common_word) + str('/24')
        #print(use_network)

        # set the value to GUI entry
        self.sub3_4_3_entry_1.delete(0, tkinter.END)
        self.sub3_4_3_entry_1.insert(0, use_network)

    def run_auto_ip(self,target_area_name):
        #print('--- run_auto_ip ---')
        l3_segment_group_array = ns_def.get_l3_segments(self)
        #print(l3_segment_group_array)

        '''Create existing IP address list'''
        exist_ip_list = []
        for tmp_l3_segment_group_array in l3_segment_group_array:
            for tmp_tmp_l3_segment_group_array in tmp_l3_segment_group_array:
                if tmp_tmp_l3_segment_group_array[4] != '':
                    exist_ip_list.append(tmp_tmp_l3_segment_group_array[4])
        #print('--- exist_ip_list ---')
        #print(exist_ip_list)

        # Set to store unique networks without duplicates
        unique_networks = set()

        # Extract networks from each CIDR notation and add them to the set
        for cidr in exist_ip_list:
            # Remove any whitespace (if present in the input)
            cidr = cidr.strip()
            # Create an IPv4Network object
            network = ipaddress.IPv4Network(cidr, strict=False)
            # Add the network to the set
            unique_networks.add(network)

        # Print the list of unique network addresses in CIDR notation
        #print("--- Unique network addresses in CIDR notation: ---")
        #for network in sorted(unique_networks):
        #    print(network.with_prefixlen)

        '''calc ip address'''
        ip_assigned_l3_segment_group_array = []
        ip_address_exists_array = []
        ip_address_exists_array_sub = []

        for tmp_l3_segment_group_array in l3_segment_group_array:
            flag_way_point_exists = False
            for tmp_tmp_l3_segment_group_array in tmp_l3_segment_group_array:
                if tmp_tmp_l3_segment_group_array[0] == 'N/A':
                    flag_way_point_exists = True

            if (tmp_l3_segment_group_array[0][0] == target_area_name and flag_way_point_exists == False) or (target_area_name == 'N/A' and flag_way_point_exists == True):
                #check to ip address exists
                flag_ip_address_exists = False

                for tmp_tmp_l3_segment_group_array in tmp_l3_segment_group_array:
                    if tmp_tmp_l3_segment_group_array[4] != '':
                        ip_address_exists = tmp_tmp_l3_segment_group_array[4]
                        ip_address_exists_array.append(tmp_tmp_l3_segment_group_array)
                        ip_address_exists_array_sub.append(tmp_tmp_l3_segment_group_array[4])
                        ip_address_exists_subnet = ip_address_exists.split('/')
                        flag_ip_address_exists = True

                #calc ip address subnet
                required_ip_num = len(tmp_l3_segment_group_array) + int(self.sub3_4_2_entry_1.get())
                # Add 2 for network address and broadcast address
                required_ip_num += 2
                # Find the number of bits in the subnet mask
                subnet_mask_bits = 32
                while True:
                    if 2 ** (32 - subnet_mask_bits) >= required_ip_num:
                        break
                    subnet_mask_bits -= 1

                # Display subnet mask in CIDR format
                cidr_notation = f"/{subnet_mask_bits}"
                #print(f"Required subnet mask in CIDR notation: {cidr_notation}")
                start_ip = str(self.sub3_4_3_entry_1.get()).split('/')[0]

                if flag_ip_address_exists == True:
                    network_exists = ipaddress.ip_network(ip_address_exists, strict=False)
                    start_ip = network_exists.network_address
                    cidr_notation = str('/') + str(ip_address_exists_subnet[1])

                '''check overlap'''
                # Check if a specific IP range overlaps with any of the unique networks
                #print(str(start_ip) + str(cidr_notation))
                specific_range = ipaddress.IPv4Network(str(start_ip) + str(cidr_notation))

                while True:
                    # Check for overlaps
                    overlap = any(specific_range.overlaps(network) for network in unique_networks)
                    if not overlap or flag_ip_address_exists == True:
                        break

                    # Calculate the broadcast address of the initial subnet
                    broadcast_address = specific_range.broadcast_address
                    # Calculate the first IP address of the next subnet by adding 1 to the broadcast address
                    next_subnet_network_address = broadcast_address + 1
                    # Create the next subnet based on the calculated network address
                    next_subnet = ipaddress.IPv4Network(f'{next_subnet_network_address}{str(cidr_notation)}')
                    specific_range = next_subnet
                #print(f'specific_range: {specific_range}')
                # Get all available hosts in the subnet
                subnet = ipaddress.IPv4Network(specific_range)
                available_hosts = list(subnet.hosts())

                ''' Descending order '''
                if str(self.combo3_4_4_1.get()) == "Descending order":
                    available_hosts = sorted(available_hosts, reverse=True)

                ''' Assign each host IP address from the subnet'''
                #for index, host in enumerate(available_hosts, start=1):
                #    print(f"Host {index}: {host}")

                ip_assign_num = 0
                pre_ip_assigned_l3_segment_group_array = []
                #print(ip_address_exists_array)


                for tmp_tmp_l3_segment_group_array in tmp_l3_segment_group_array:
                    if str(self.combo3_4_6_1.get()) == "Reassign within the same subnet":
                        #print("Reassign within the same subnet")
                        pre_ip_assigned_l3_segment_group_array.append([tmp_tmp_l3_segment_group_array[0],tmp_tmp_l3_segment_group_array[1],tmp_tmp_l3_segment_group_array[2],tmp_tmp_l3_segment_group_array[3],str(available_hosts[ip_assign_num]) + str(cidr_notation)])
                        ip_assign_num += 1
                    elif str(self.combo3_4_6_1.get()) == "Keep existing IP address":
                        #print("Keep existing IP address")
                        while True:
                            if tmp_tmp_l3_segment_group_array in ip_address_exists_array:
                                pre_ip_assigned_l3_segment_group_array.append(tmp_tmp_l3_segment_group_array)
                                ip_assign_num += 1
                                break
                            elif str(available_hosts[ip_assign_num]) + str(cidr_notation) in ip_address_exists_array_sub:
                                ip_assign_num += 1
                            else:
                                pre_ip_assigned_l3_segment_group_array.append([tmp_tmp_l3_segment_group_array[0], tmp_tmp_l3_segment_group_array[1], tmp_tmp_l3_segment_group_array[2], tmp_tmp_l3_segment_group_array[3], str(available_hosts[ip_assign_num]) + str(cidr_notation)])
                                ip_assign_num += 1
                                break

                ip_assigned_l3_segment_group_array.append(pre_ip_assigned_l3_segment_group_array)
                #print(f"pre_ip_assigned_l3_segment_group_array: {pre_ip_assigned_l3_segment_group_array}")

                unique_networks.add(subnet)

        #print('--- ip_assigned_l3_segment_group_array ---')
        #print(ip_assigned_l3_segment_group_array)

        '''
        Update to the Master file
        '''
        ws_l3_name = 'Master_Data_L3'
        excel_maseter_file = self.inFileTxt_L3_3_1.get()
        self.l3_table_array = ns_def.convert_master_to_array(ws_l3_name, excel_maseter_file, '<<L3_TABLE>>')

        updated_l3_table_array = []
        for tmp_l3_table_array in self.l3_table_array:
            #print(tmp_l3_table_array)
            if tmp_l3_table_array[1][0] != target_area_name and (tmp_l3_table_array[0] == 1 or tmp_l3_table_array[0] == 2):
                updated_l3_table_array.append(tmp_l3_table_array)
            else:
                flag_l3_if_match = False
                for tmp_ip_assigned_l3_segment_group_array in ip_assigned_l3_segment_group_array:
                    for tmp_tmp_tmp_ip_assigned_l3_segment_group_array in tmp_ip_assigned_l3_segment_group_array:

                        if tmp_tmp_tmp_ip_assigned_l3_segment_group_array[1] == tmp_l3_table_array[1][1]  and tmp_tmp_tmp_ip_assigned_l3_segment_group_array[2] == tmp_l3_table_array[1][2]:
                            flag_l3_if_match = True
                            if len(tmp_l3_table_array) != 5:
                                updated_l3_table_array.append([tmp_l3_table_array[0],tmp_tmp_tmp_ip_assigned_l3_segment_group_array])
                            else:
                                updated_l3_table_array.append(tmp_l3_table_array)
                if flag_l3_if_match == False:
                    updated_l3_table_array.append(tmp_l3_table_array)

        #print('--- updated_l3_table_array ---')
        #print(updated_l3_table_array)

        '''
        write the Master file
        '''
        last_l3_table_tuple = {}
        last_l3_table_tuple = ns_def.convert_array_to_tuple(updated_l3_table_array)

        excel_master_ws_name_l3 = 'Master_Data_L3'
        # delete L3 Table sheet
        ns_def.remove_excel_sheet(excel_maseter_file, excel_master_ws_name_l3)
        # create L3 Table sheet
        ns_def.create_excel_sheet(excel_maseter_file, excel_master_ws_name_l3)
        # write tuple to excel master data
        ns_def.write_excel_meta(last_l3_table_tuple, excel_maseter_file, excel_master_ws_name_l3, '_template_', 0, 0)

class  summary_diagram():

    def export_summary_diagram(self,dummy): # add at ver 2.3.4
        print('--- export_summary_diagram ---')
        iDir = os.path.abspath(os.path.dirname(self.full_filepath ))

        # SET IP Address report file patch
        basename_without_ext = os.path.splitext(os.path.basename(self.inFileTxt_L2_3_1.get()))[0]

        self.outFileTxt_Lx_1_1 = tk.Entry()
        self.outFileTxt_Lx_1_1.delete(0, tkinter.END)
        self.outFileTxt_Lx_1_1.insert(tk.END, iDir + ns_def.return_os_slash() + basename_without_ext.replace('[MASTER]','__TMP__[MASTER]') + '.xlsx')
        self.excel_maseter_file_backup  = self.outFileTxt_Lx_1_1.get()

        # remove exist _TMP_ file
        if os.path.isfile(str(self.excel_maseter_file_backup)) == True:
            os.remove(str(self.excel_maseter_file_backup))

        shutil.copy(str(self.inFileTxt_L2_3_1.get()), str(self.excel_maseter_file_backup))

        #GET backup master file parameter
        # parameter
        ws_name = 'Master_Data'
        tmp_ws_name = '_tmp_'
        ppt_meta_file = str(self.excel_maseter_file_backup)

        # convert from master to array and convert to tuple
        self.position_folder_array = ns_def.convert_master_to_array(ws_name, ppt_meta_file, '<<POSITION_FOLDER>>')
        self.position_shape_array = ns_def.convert_master_to_array(ws_name, ppt_meta_file, '<<POSITION_SHAPE>>')
        self.root_folder_array = ns_def.convert_master_to_array(ws_name, ppt_meta_file, '<<ROOT_FOLDER>>')
        self.positoin_line_array = ns_def.convert_master_to_array(ws_name, ppt_meta_file, '<<POSITION_LINE>>')
        self.position_folder_tuple = ns_def.convert_array_to_tuple(self.position_folder_array)
        self.position_shape_tuple = ns_def.convert_array_to_tuple(self.position_shape_array)
        self.root_folder_tuple = ns_def.convert_array_to_tuple(self.root_folder_array)
        self.positoin_line_tuple = ns_def.convert_array_to_tuple(self.positoin_line_array)

        #print('---- self.position_folder_tuple ----')
        #print(self.position_folder_tuple)
        #print('---- self.position_folder_array ----')
        #print(self.position_folder_array)
        #print('---- self.position_shape_tuple ----')
        #print(self.position_shape_tuple)
        #print('--- self.position_shape_array ---')
        #print(self.position_shape_array)

        # GET Folder and wp name List
        folder_wp_name_array = ns_def.get_folder_wp_array_from_master(ws_name, ppt_meta_file)
        #print('---- folder_wp_name_array ----')
        #print(folder_wp_name_array)

        '''delete shape in <<POSITION_SHAPE>>'''
        new_position_shape_array = []
        counter = 1

        for item in self.position_shape_array:
            number = item[0]
            elements = item[1]

            # Rule 1: If the first element is 1, do not change it
            if number == 1:
                new_position_shape_array.append([counter, elements])
                counter += 1
                continue

            # Additional Rule: If the first item contains '_wp_', do not change it
            if '_wp_' in elements[0]:
                new_position_shape_array.append([counter, elements])
                counter += 1
                continue

            # Rule 2: If the first element of the second item is '', remove the entire entry
            if elements[0] == '':
                continue

            # Rule 3: Otherwise, set the second element to '<END>' and remove the rest
            new_entry = [elements[0], '<END>']

            # Assign a new ascending order number
            new_position_shape_array.append([counter, new_entry])
            counter += 1

        # Output the new array
        #print('--- new_position_shape_array ---')
        #print(new_position_shape_array)

        new_position_shape_tuple = ns_def.convert_array_to_tuple(new_position_shape_array)
        #print(new_position_shape_tuple)

        # SET new <<POSITION_SHAPE>>
        write_to_section = '<<POSITION_SHAPE>>'
        offset_row = 0
        offset_column = 0
        ns_def.clear_section_sheet('Master_Data', self.excel_maseter_file_backup, self.position_shape_tuple)
        ns_def.write_excel_meta(new_position_shape_tuple, self.excel_maseter_file_backup, 'Master_Data',write_to_section, offset_row, offset_column)

        ''' change value of <<POSITION_FOLDER>>'''
        self_position_folder_tuple = self.position_folder_tuple

        # Initialize new tuple
        new_position_folder_tuple = {}

        # Ratio values
        ratio_x = 0.35
        ratio_y = 0.25

        offset_sum_x = 0.0
        max_offset_sum_x = 0.0
        offset_sum_y = 0.0

        current_y_axis = 1

        # Iterate through the original tuple to apply the transformations
        for (y, x), value in self_position_folder_tuple.items():
            # Check the conditions for transformation -- X-axis
            if (
                    x >= 2 and
                    self_position_folder_tuple.get((y, 1)) not in ['<SET_WIDTH>', '<<POSITION_FOLDER>>'] and
                    '_wp_' not in str(value) and str(value) != ''
            ):
                # Calculate the new key and transformed value
                new_key = (y - 1, x)
                new_value = self_position_folder_tuple.get((y - 1 , x))

                # Add to the new tuple
                new_position_folder_tuple[new_key] = new_value * ratio_x

                if current_y_axis == y:
                    offset_sum_x += (new_value - new_value * ratio_x)
                    if max_offset_sum_x < offset_sum_x:
                        max_offset_sum_x = offset_sum_x
                else:
                    offset_sum_x = (new_value - new_value * ratio_x)
                    current_y_axis = y
                    if max_offset_sum_x < offset_sum_x:
                        max_offset_sum_x = offset_sum_x

            # Check the conditions for transformation -- Y-axis
            if (
                    x == 1 and
                    self_position_folder_tuple.get((y, 1)) not in ['<SET_WIDTH>', '<<POSITION_FOLDER>>']
            ):
                flag_wp_include = False
                for (yy, xx), value in self_position_folder_tuple.items():
                    # Check the conditions for transformation 2nd
                    if (
                            y == yy and
                            isinstance(self_position_folder_tuple[yy, xx], str) and  # Check if the value is a string
                            '_wp_' not in str(self_position_folder_tuple[yy, xx]) and  # Check if '_wp_' is in the string
                            self_position_folder_tuple[yy, xx] != ''
                    ):
                        flag_wp_include = True

                if flag_wp_include == True :
                    new_position_folder_tuple[y,x] = self_position_folder_tuple[y,x] * ratio_y
                    offset_sum_y += (self_position_folder_tuple[y,x] - (self_position_folder_tuple[y,x] * ratio_y))

        # Output the new tuple
        #print(new_position_folder_tuple)

        # Combine dictionaries, prioritizing new_position_folder_tuple
        combined_position_folder_tuple = self_position_folder_tuple.copy()  # Start with a copy of the original
        combined_position_folder_tuple.update(new_position_folder_tuple)  # Update with new values, overwriting where keys overlap

        # Output the combined result
        #print(combined_position_folder_tuple)

        # SET new <<FOLDER_SHAPE>>
        write_to_section = '<<POSITION_FOLDER>>'
        offset_row = 0
        offset_column = 0
        ns_def.clear_section_sheet('Master_Data', self.excel_maseter_file_backup, self.position_folder_tuple)
        ns_def.write_excel_meta(combined_position_folder_tuple, self.excel_maseter_file_backup, 'Master_Data',write_to_section, offset_row, offset_column)

        ''' change value of <<ROOT_FOLDER>>'''
        # Create a new tuple with the modifications
        new_root_folder_tuple = self.root_folder_tuple.copy()

        # Modify the specific values
        new_root_folder_tuple[(2, 2)] = 'Summary Diagram'
        new_root_folder_tuple[(2, 7)] = self.root_folder_tuple[(2, 7)] - max_offset_sum_x
        new_root_folder_tuple[(2, 8)] = offset_sum_y

        # SET new <<ROOT_SHAPE>>
        write_to_section = '<<ROOT_FOLDER>>'
        offset_row = 0
        offset_column = 0
        ns_def.clear_section_sheet('Master_Data', self.excel_maseter_file_backup, self.root_folder_tuple)
        ns_def.write_excel_meta(new_root_folder_tuple, self.excel_maseter_file_backup, 'Master_Data',write_to_section, offset_row, offset_column)

        ''' change value of <<POSITION_LINE>>'''
        new_get_shape_folder_tuple =  ns_def.get_shape_folder_tuple(self.position_shape_tuple)
        #print(new_get_shape_folder_tuple)
        #print(self.positoin_line_tuple)
        self_positoin_line_tuple = self.positoin_line_tuple
        last_positoin_line_tuple = self_positoin_line_tuple.copy()

        oppsite_wp_array = []
        for (y, x), value in self_positoin_line_tuple.items():
            if value in new_get_shape_folder_tuple: # for ns-011 bug of tmp workaround
                if x in [1, 2] and y >= 3 and '_wp_' in new_get_shape_folder_tuple[value]:
                    if x == 1:
                        last_positoin_line_tuple[(y, x + 1)] = new_get_shape_folder_tuple[self_positoin_line_tuple[y, (x + 1)]]
                    if x == 2:
                        last_positoin_line_tuple[(y, x - 1)] = new_get_shape_folder_tuple[self_positoin_line_tuple[y, (x - 1)]]

        # Output the result
        #print(last_positoin_line_tuple)

        # SET new <<FOLDER_SHAPE>>
        write_to_section = '<<POSITION_LINE>>'
        offset_row = 0
        offset_column = 0
        ns_def.clear_section_sheet('Master_Data', self.excel_maseter_file_backup, self.positoin_line_tuple)
        ns_def.write_excel_meta(last_positoin_line_tuple, self.excel_maseter_file_backup, 'Master_Data',write_to_section, offset_row, offset_column)


class cli_on_gui():
    """
    CLI command input GUI for Network Sketcher
    """

    def __init__(self):
        """
        Initialize CLI GUI handler
        """
        self.root = None
        self.full_filepath = None
        self.cli_window = None
        self.cli_text = None
        self.from_master_panel = False  # Flag to indicate if opened from Master Panel


    def load_cli_from_file(self):
        """
        Load CLI commands from a file
        """
        fTyp = [("Text files", "*.txt"), ("All files", "*.*")]
        iDir = os.path.abspath(os.path.dirname(sys.argv[0]))
        file_path = tk.filedialog.askopenfilename(filetypes=fTyp, initialdir=iDir)

        if file_path:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                self.cli_text.delete('1.0', 'end')
                self.cli_text.insert('1.0', content)
                self.cli_text.config(fg='black')
            except Exception as e:
                tkinter.messagebox.showerror('Error', f'Failed to load file:\n{str(e)}')

    def execute_cli_commands(self):
        """
        Execute the inputted CLI commands using subprocess
        """
        cli_commands = self.cli_text.get('1.0', 'end-1c').strip()

        if not cli_commands:
            tkinter.messagebox.showwarning('Warning', 'Please enter valid CLI commands.')
            return

        try:
            # Split CLI commands by line and filter out comments and empty lines
            command_lines = []
            for line in cli_commands.split('\n'):
                line = line.strip()
                # Skip empty lines and comment lines (starting with #)
                if line and not line.startswith('#'):
                    command_lines.append(line)

            if not command_lines:
                tkinter.messagebox.showwarning('Warning', 'No valid commands found. All lines are comments or empty.')
                return

            # Prepare Master file path (use the path from the main window entry)
            master_file = self.full_filepath

            # Get network_sketcher.py path
            if getattr(sys, 'frozen', False):
                # Running as compiled executable
                sketcher_path = sys.executable
            else:
                # Running as script
                sketcher_path = os.path.abspath(sys.argv[0])

            # Get python executable path
            python_executable = sys.executable

            # Execute each command
            success_count = 0
            error_count = 0
            error_messages = []
            executed_commands = []
            command_outputs = []

            for cmd_index, cmd in enumerate(command_lines, 1):
                try:
                    print(f"\n[{cmd_index}/{len(command_lines)}] Executing: {cmd}")

                    # Parse command line more carefully using shlex for proper quote handling
                    import shlex

                    # First, try to parse with shlex to handle quotes properly
                    try:
                        parts = shlex.split(cmd)
                    except ValueError:
                        # If shlex fails, fall back to simple split
                        parts = cmd.split()

                    # Filter out network_sketcher.py path and --master arguments
                    filtered_parts = []
                    skip_next = False
                    i = 0
                    while i < len(parts):
                        part = parts[i]

                        if skip_next:
                            skip_next = False
                            i += 1
                            continue

                        if part.lower() in ['python', 'python.exe']:
                            i += 1
                            continue

                        if 'network_sketcher.py' in part.lower():
                            i += 1
                            continue

                        if part in ['--master', '-master']:
                            skip_next = True  # Skip the next argument (master file path)
                            i += 1
                            continue

                        filtered_parts.append(part)
                        i += 1

                    if not filtered_parts:
                        print(f"[SKIP] Empty command after filtering")
                        continue

                    # Check if this is an 'export ai_context_file' command
                    # Show security warning dialog when executed from cli_on_gui
                    if 'export' in filtered_parts and 'ai_context_file' in filtered_parts:
                        # Check if --accept-security-risk flag is already present
                        if '--accept-security-risk' not in filtered_parts:
                            # Show security warning dialog to user
                            result = tkinter.messagebox.askyesno(
                                "Warning",
                                "The exported AI Context file contains data from the master file, "
                                "which includes all configuration information for the network (NW). "
                                "Please be aware that there is a risk of data leakage if the exported file "
                                "is loaded into a Large Language Model (LLM). "
                                "Do you fully understand and accept this risk before proceeding with the export?"
                            )

                            if not result:
                                # User declined - skip this command
                                print(f"[SKIP] User declined ai_context_file export")
                                command_outputs.append({
                                    'command': ' '.join(filtered_parts),
                                    'stdout': '[Info] Export cancelled by user via dialog',
                                    'stderr': '',
                                    'returncode': 0
                                })
                                success_count += 1
                                executed_commands.append(' '.join(filtered_parts) + ' (cancelled by user)')
                                continue

                            # User accepted - add the flag to bypass CLI confirmation
                            filtered_parts.append('--accept-security-risk')
                            print(f"[INFO] User accepted security risk, added --accept-security-risk flag")

                    # Construct command array for subprocess
                    # Format: [python_path, sketcher_path, command, args..., '--master', master_file_path]
                    subprocess_cmd = [python_executable, sketcher_path] + filtered_parts + ['--master', master_file]

                    print(f"Subprocess command: {subprocess_cmd}")

                    # Execute command using subprocess
                    # stdin=subprocess.DEVNULL prevents blocking on input() calls
                    result = subprocess.run(
                        subprocess_cmd,
                        capture_output=True,
                        text=True,
                        timeout=300,  # 5 minutes timeout
                        encoding='utf-8',
                        errors='replace',
                        stdin=subprocess.DEVNULL
                    )

                    # Check return code
                    if result.returncode == 0:
                        success_count += 1
                        executed_commands.append(' '.join(filtered_parts))
                        command_outputs.append({
                            'command': ' '.join(filtered_parts),
                            'stdout': result.stdout,
                            'stderr': result.stderr,
                            'returncode': 0
                        })
                        print(f"[SUCCESS] Command executed: {' '.join(filtered_parts)}")
                        if result.stdout:
                            print(f"Output:\n{result.stdout}")
                    else:
                        error_count += 1
                        error_msg = f"Command: {' '.join(filtered_parts)}\nReturn code: {result.returncode}"
                        if result.stderr:
                            error_msg += f"\nError output:\n{result.stderr}"
                        if result.stdout:
                            error_msg += f"\nStandard output:\n{result.stdout}"
                        error_messages.append(error_msg)

                        # Also add to command_outputs for display
                        command_outputs.append({
                            'command': ' '.join(filtered_parts),
                            'stdout': result.stdout,
                            'stderr': result.stderr,
                            'returncode': result.returncode
                        })

                        print(f"[ERROR] Command failed: {' '.join(filtered_parts)}")
                        print(f"Return code: {result.returncode}")
                        if result.stderr:
                            print(f"Error output:\n{result.stderr}")
                        if result.stdout:
                            print(f"Standard output:\n{result.stdout}")

                except subprocess.TimeoutExpired:
                    error_count += 1
                    error_msg = f"Command: {' '.join(filtered_parts) if filtered_parts else cmd}\nError: Command timeout (exceeded 5 minutes)"
                    error_messages.append(error_msg)
                    print(f"[ERROR] Command timeout: {cmd}")

                except Exception as cmd_error:
                    error_count += 1
                    error_msg = f"Command: {' '.join(filtered_parts) if filtered_parts else cmd}\nError: {str(cmd_error)}"
                    error_messages.append(error_msg)
                    print(f"[ERROR] Command: {cmd}")
                    print(f"Error: {str(cmd_error)}")
                    import traceback
                    traceback.print_exc()

            print(f"\n=== Execution Summary ===")
            print(f"Total commands: {len(command_lines)}")
            print(f"Success: {success_count}")
            print(f"Failed: {error_count}")

            # Show result summary with proper error handling
            try:
                if error_count == 0:
                    # Show output details in a separate window
                    try:
                        self.show_command_output_window(command_outputs)
                    except Exception as output_error:
                        print(f"[WARNING] Could not show output window: {output_error}")

                    # Close CLI window
                    try:
                        self.cli_window.destroy()
                    except:
                        pass

                    # Only open Master Panel if NOT opened from Master Panel
                    if not self.from_master_panel:
                        # ★★★ NEW: Check if master file has data before opening Master Panel ★★★
                        try:
                            import ns_def

                            # Check if master file has any areas/devices
                            has_data = False
                            try:
                                style_folder_array = ns_def.convert_master_to_array('Master_Data', master_file, '<<STYLE_FOLDER>>')
                                # Check if there are any areas (rows beyond header rows 1, 2, 3)
                                for item in style_folder_array:
                                    if item[0] not in [1, 2, 3] and len(item[1]) > 0 and item[1][0]:
                                        has_data = True
                                        break
                            except Exception as check_error:
                                print(f"[WARNING] Could not check master file data: {check_error}")
                                has_data = False

                            if not has_data:
                                # Master file has no data - don't open Master Panel
                                print(f"[INFO] Master file has no areas/devices. Master Panel not opened.")
                                '''tkinter.messagebox.showinfo('Complete',
                                                            'CLI commands executed successfully!\n\n'
                                                            'The master file has been updated but contains no areas or devices.\n'
                                                            'Please add areas and devices using CLI commands or reload the master file.')'''
                            else:
                                # Master file has data - open Master Panel
                                # Automatically reopen the Master Panel
                                print(f"[INFO] Automatically opening Master Panel...")

                                # Import network_sketcher to access ns_front_run
                                import network_sketcher

                                # Create a minimal mock instance with required methods
                                class TempInstance:
                                    def __init__(self, root, full_filepath, cli_gui_instance):
                                        self.root = root
                                        self.full_filepath = full_filepath
                                        self.filename = os.path.basename(full_filepath)
                                        self.click_value = ''
                                        self.click_value_2nd = ''
                                        self.click_value_3rd = ''
                                        self.click_value_VPN = ''
                                        self.cli_flag_no_export = False
                                        self.flag_summary_diagram = False
                                        self.keep_root_width = 0.0
                                        self.keep_root_hight = 0.0
                                        self.cli_gui_instance = cli_gui_instance

                                    def on_combobox_select(self, event):
                                        """Handle combobox selection"""
                                        try:
                                            import ns_def
                                            self.attribute_tuple1_1 = ns_def.get_global_attribute_tuple(
                                                self.inFileTxt_L2_3_1.get(),
                                                self.comboATTR_1_1.get()
                                            )
                                        except Exception as e:
                                            print(f"[WARNING] Could not update attribute tuple: {e}")

                                    def click_action_sub(self, click_value, target_area_name):
                                        """Delegate to network_sketcher.ns_front_run.click_action_sub"""
                                        import network_sketcher
                                        network_sketcher.ns_front_run.click_action_sub(self, click_value, target_area_name)

                                temp_instance = TempInstance(self.root, self.full_filepath, self)

                                file_type_array = ['EXCEL_MASTER', 'EXCEL_MASTER']
                                network_sketcher.ns_front_run.sub_excel_master_1(temp_instance, file_type_array)

                                print(f"[SUCCESS] Master Panel opened successfully")

                        except Exception as panel_error:
                            print(f"[WARNING] Could not open Master Panel automatically: {panel_error}")
                            import traceback
                            traceback.print_exc()
                            # Fallback: Show manual instruction
                            tkinter.messagebox.showinfo('Next Step',
                                                        'CLI commands executed successfully!\n\n'
                                                        'Please click "Submit" button in the main window\n'
                                                        'to open the Master Panel with the updated data.')
                    else:
                        # Opened from Master Panel - just show completion message without opening Master Panel
                        print(f"[INFO] CLI commands executed from Master Panel. Master Panel not reopened.")

                elif success_count > 0:
                    error_detail = '\n\n'.join(error_messages[:3])  # Show first 3 errors
                    if len(error_messages) > 3:
                        error_detail += f'\n\n... and {len(error_messages) - 3} more errors'

                    # Show output details in a separate window
                    try:
                        if command_outputs:
                            self.show_command_output_window(command_outputs)
                    except Exception as output_error:
                        print(f"[WARNING] Could not show output window: {output_error}")

                    # Close CLI window
                    try:
                        self.cli_window.destroy()
                    except:
                        pass

                    # Different message based on where CLI window was opened from
                    if not self.from_master_panel:
                        result = tkinter.messagebox.askyesno('Partial Success',
                                                             f'Executed: {success_count} commands\n'
                                                             f'Failed: {error_count} commands\n\n'
                                                             f'Error details:\n{error_detail}\n\n'
                                                             f'Do you want to open the Master Panel?')

                        if result:
                            # ★★★ NEW: Check if master file has data before opening Master Panel ★★★
                            try:
                                import ns_def

                                # Check if master file has any areas/devices
                                has_data = False
                                try:
                                    style_folder_array = ns_def.convert_master_to_array('Master_Data', master_file, '<<STYLE_FOLDER>>')
                                    for item in style_folder_array:
                                        if item[0] not in [1, 2, 3] and len(item[1]) > 0 and item[1][0]:
                                            has_data = True
                                            break
                                except Exception as check_error:
                                    print(f"[WARNING] Could not check master file data: {check_error}")
                                    has_data = False

                                if not has_data:
                                    tkinter.messagebox.showinfo('Info',
                                                                'The master file contains no areas or devices.\n'
                                                                'Please add areas and devices using CLI commands first.')
                                else:
                                    # Open Master Panel
                                    import network_sketcher

                                    class TempInstance:
                                        def __init__(self, root, full_filepath, cli_gui_instance):
                                            self.root = root
                                            self.full_filepath = full_filepath
                                            self.filename = os.path.basename(full_filepath)
                                            self.click_value = ''
                                            self.click_value_2nd = ''
                                            self.click_value_3rd = ''
                                            self.click_value_VPN = ''
                                            self.cli_flag_no_export = False
                                            self.flag_summary_diagram = False
                                            self.keep_root_width = 0.0
                                            self.keep_root_hight = 0.0
                                            self.cli_gui_instance = cli_gui_instance

                                        def on_combobox_select(self, event):
                                            """Handle combobox selection"""
                                            try:
                                                import ns_def
                                                self.attribute_tuple1_1 = ns_def.get_global_attribute_tuple(
                                                    self.inFileTxt_L2_3_1.get(),
                                                    self.comboATTR_1_1.get()
                                                )
                                            except Exception as e:
                                                print(f"[WARNING] Could not update attribute tuple: {e}")

                                        def click_action_sub(self, click_value, target_area_name):
                                            """Delegate to network_sketcher.ns_front_run.click_action_sub"""
                                            import network_sketcher
                                            network_sketcher.ns_front_run.click_action_sub(self, click_value, target_area_name)

                                    temp_instance = TempInstance(self.root, self.full_filepath, self)

                                    file_type_array = ['EXCEL_MASTER', 'EXCEL_MASTER']
                                    network_sketcher.ns_front_run.sub_excel_master_1(temp_instance, file_type_array)

                            except Exception as panel_error:
                                print(f"[WARNING] Could not open Master Panel: {panel_error}")
                                import traceback
                                traceback.print_exc()
                                tkinter.messagebox.showinfo('Next Step',
                                                            'Please click "Submit" button in the main window\n'
                                                            'to open the Master Panel with the updated data.')
                    else:
                        # Opened from Master Panel - just show completion message with errors
                        tkinter.messagebox.showwarning('Partial Success',
                                                       f'Executed: {success_count} commands\n'
                                                       f'Failed: {error_count} commands\n\n'
                                                       f'Error details:\n{error_detail}\n\n'
                                                       f'The Master file has been partially updated.\n'
                                                       f'Please refresh or reopen the Master Panel if needed.')

                else:
                    error_detail = '\n\n'.join(error_messages[:3])  # Show first 3 errors
                    if len(error_messages) > 3:
                        error_detail += f'\n\n... and {len(error_messages) - 3} more errors'

                    tkinter.messagebox.showerror('All Commands Failed',
                                                 f'All {error_count} commands failed.\n\n'
                                                 f'Error details:\n{error_detail}')

                    # Close CLI window
                    try:
                        self.cli_window.destroy()
                    except:
                        pass

            except Exception as summary_error:
                print(f"[ERROR] Error showing summary: {summary_error}")
                import traceback
                traceback.print_exc()
                # Show a simple error message
                tkinter.messagebox.showinfo('Execution Complete',
                                            f'Commands executed: {success_count}\n'
                                            f'Commands failed: {error_count}')

        except Exception as e:
            # This except block catches errors in the outer try block
            error_message = f'Failed to execute CLI commands:\n{str(e)}'
            print(f"[FATAL ERROR] {error_message}")
            import traceback
            traceback.print_exc()
            tkinter.messagebox.showerror('Error', error_message)




    def show_cli_input_window(self):
        """
        Display CLI input window when Master file
        """
        self.cli_window = tk.Toplevel(self.root)
        self.cli_window.title('CLI Input')
        self.cli_window.geometry("900x650+150+150")

        # Main frame
        main_frame = tk.LabelFrame(
            self.cli_window,
            text="Run CLI Commands on the Master File",
            font=("", 14),
            background="#FFF3CD",
            padx=10,
            pady=10
        )
        main_frame.pack(fill='both', expand=True, padx=10, pady=10)

        # Information label
        info_label = tk.Label(
            main_frame,
            text="You can enter and execute multiple CLI commands below.\n - There is no need to specify the path to network.sketcher.py or the master file.\n - File paths entered in the CLI will be ignored.",
            font=("", 11),
            background="#FFF3CD",
            justify='left'
        )
        info_label.pack(pady=(0, 10))

        # CLI input area
        cli_frame = tk.LabelFrame(
            main_frame,
            text="CLI Commands",
            font=("", 12),
            background="#FFFFFF"
        )
        cli_frame.pack(fill='both', expand=True, pady=5)

        # Text area with scrollbar
        scrollbar = tk.Scrollbar(cli_frame)
        scrollbar.pack(side='right', fill='y')

        self.cli_text = tk.Text(
            cli_frame,
            font=("Courier", 10),
            wrap='word',
            yscrollcommand=scrollbar.set
        )
        self.cli_text.pack(fill='both', expand=True, padx=5, pady=5)
        scrollbar.config(command=self.cli_text.yview)

        # Placeholder text - Updated to show proper syntax without quotes around arrays
        placeholder_text = """# Enter CLI commands here (one per line)\n# show ...\n# add ...\n# delete ...\n# rename ...\n# export ...\n"""
        self.cli_text.insert('1.0', placeholder_text)
        self.cli_text.config(fg='gray')

        # Placeholder clear handling
        def on_focus_in(event):
            if self.cli_text.get('1.0', 'end-1c') == placeholder_text.strip():
                self.cli_text.delete('1.0', 'end')
                self.cli_text.config(fg='black')

        def on_focus_out(event):
            if not self.cli_text.get('1.0', 'end-1c').strip():
                self.cli_text.insert('1.0', placeholder_text)
                self.cli_text.config(fg='gray')

        self.cli_text.bind('<FocusIn>', on_focus_in)
        self.cli_text.bind('<FocusOut>', on_focus_out)

        # Button frame
        button_frame = tk.Frame(main_frame, background="#FFF3CD")
        button_frame.pack(pady=10)

        # Execute button
        execute_button = tk.Button(
            button_frame,
            text="Execute CLI Commands",
            font=("", 12),
            command=self.execute_cli_commands,
            bg="#28a745",
            fg="white",
            padx=20,
            pady=5,
        )
        execute_button.pack(side='left', padx=5)

        '''
        # Cancel button
        cancel_button = tk.Button(
            button_frame,
            text="Cancel",
            font=("", 12),
            command=self.cli_window.destroy,
            bg="#dc3545",
            fg="white",
            padx=20,
            pady=5
        )
        cancel_button.pack(side='left', padx=5)

        # Load from file button
        load_button = tk.Button(
            button_frame,
            text="Load from File",
            font=("", 12),
            command=self.load_cli_from_file,
            padx=20,
            pady=5
        )
        load_button.pack(side='left', padx=5)'''

    def show_command_output_window(self, command_outputs):
        """
        Display command execution output in a separate window
        """
        if not command_outputs:
            print("[INFO] No command outputs to display")
            return

        try:
            output_window = tk.Toplevel(self.root)
            output_window.title('Command Execution Output')
            output_window.geometry("800x600+200+200")

            # Main frame
            main_frame = tk.Frame(output_window)
            main_frame.pack(fill='both', expand=True, padx=10, pady=10)

            # Title
            title_label = tk.Label(
                main_frame,
                text=f"Command Execution Results ({len(command_outputs)} commands)",
                font=("", 12, "bold")
            )
            title_label.pack(pady=(0, 10))

            # Text area with scrollbar
            text_frame = tk.Frame(main_frame)
            text_frame.pack(fill='both', expand=True)

            scrollbar = tk.Scrollbar(text_frame)
            scrollbar.pack(side='right', fill='y')

            output_text = tk.Text(
                text_frame,
                font=("Courier", 9),
                wrap='word',
                yscrollcommand=scrollbar.set
            )
            output_text.pack(fill='both', expand=True)
            scrollbar.config(command=output_text.yview)

            # Error indicators list (case-insensitive matching)
            ERROR_INDICATORS = [
                # Error keywords
                '[error]',
                'error:',
                'error ',
                # Failure keywords
                'failed',
                'failure',
                'failing',
                # Exception keywords
                'exception',
                'traceback',
                # Warning/Critical keywords (optional - can be removed if too strict)
                'critical',
                'fatal',
                # Denial keywords
                'denied',
                'not found',
                'not exist',
                'does not exist',
                'cannot ',
                'could not',
                'unable to',
                # Invalid keywords
                'invalid',
                'illegal',
                'unsupported',
                # Abort keywords
                'aborted',
                'abort:',
                'terminated',
                # Permission keywords
                'permission denied',
                'access denied',
                'unauthorized',
                # Timeout keywords
                'timeout',
                'timed out',
            ]

            # Helper function to check if output contains error
            def is_error_output(output):
                """Check if the command output indicates an error"""
                # Check return code first
                returncode = output.get('returncode', 0)
                if returncode != 0:
                    return True

                # Check stdout for error indicators
                stdout = output.get('stdout', '')
                if stdout:
                    stdout_lower = stdout.lower()
                    for indicator in ERROR_INDICATORS:
                        if indicator in stdout_lower:
                            return True

                # Check stderr for any content (stderr usually indicates issues)
                stderr = output.get('stderr', '')
                if stderr and stderr.strip():
                    stderr_lower = stderr.lower()
                    for indicator in ERROR_INDICATORS:
                        if indicator in stderr_lower:
                            return True

                return False

            # ========== Add Summary Section at TOP ==========
            try:
                # Count success and failed commands using enhanced error detection
                success_count = sum(1 for o in command_outputs if not is_error_output(o))
                failed_count = len(command_outputs) - success_count

                # Add summary header
                output_text.insert('end', f"{'=' * 80}\n", 'separator')
                output_text.insert('end', "EXECUTION SUMMARY\n", 'summary_header')
                output_text.insert('end', f"{'=' * 80}\n", 'separator')

                # Add statistics
                output_text.insert('end', f"Total Commands: {len(command_outputs)}\n", 'summary_text')
                output_text.insert('end', f"Success: ", 'summary_text')
                output_text.insert('end', f"{success_count}\n", 'success')
                output_text.insert('end', f"Failed: ", 'summary_text')
                output_text.insert('end', f"{failed_count}\n", 'failed' if failed_count > 0 else 'success')
                output_text.insert('end', f"{'=' * 80}\n", 'separator')

                # Add command list summary
                for idx, output in enumerate(command_outputs, 1):
                    has_error = is_error_output(output)
                    status = "FAILED" if has_error else "SUCCESS"
                    status_tag = 'failed' if has_error else 'success'

                    output_text.insert('end', f"Command {idx}: ", 'summary_text')
                    output_text.insert('end', f"[{status}] ", status_tag)
                    output_text.insert('end', f"{output['command']}\n", 'summary_text')

                output_text.insert('end', f"{'=' * 80}\n", 'separator')
                output_text.insert('end', '\n\n')

            except Exception as summary_error:
                print(f"[WARNING] Error inserting summary: {summary_error}")

            # ========== Add Detailed Output Section ==========
            output_text.insert('end', "DETAILED OUTPUT\n", 'summary_header')
            output_text.insert('end', f"{'=' * 80}\n", 'separator')
            output_text.insert('end', '\n')

            # Insert command outputs
            for idx, output in enumerate(command_outputs, 1):
                try:
                    output_text.insert('end', f"{'=' * 80}\n", 'separator')

                    # Add return code indicator with enhanced error detection
                    has_error = is_error_output(output)
                    status = "FAILED" if has_error else "SUCCESS"
                    status_tag = 'failed' if has_error else 'success'

                    output_text.insert('end', f"Command {idx}: ", 'command')
                    output_text.insert('end', f"[{status}] ", status_tag)
                    output_text.insert('end', f"{output['command']}\n", 'command')
                    output_text.insert('end', f"{'=' * 80}\n", 'separator')

                    if output.get('stdout'):
                        output_text.insert('end', f"\nOutput:\n", 'header')
                        output_text.insert('end', f"{output['stdout']}\n", 'output')

                    if output.get('stderr'):
                        output_text.insert('end', f"\nErrors/Warnings:\n", 'error_header')
                        output_text.insert('end', f"{output['stderr']}\n", 'error')

                    output_text.insert('end', '\n\n')
                except Exception as insert_error:
                    print(f"[WARNING] Error inserting output for command {idx}: {insert_error}")

            # Configure tags for colored text
            output_text.tag_config('separator', foreground='gray')
            output_text.tag_config('command', foreground='blue', font=("Courier", 9, "bold"))
            output_text.tag_config('success', foreground='green', font=("Courier", 9, "bold"))
            output_text.tag_config('failed', foreground='red', font=("Courier", 9, "bold"))
            output_text.tag_config('header', foreground='green', font=("Courier", 9, "bold"))
            output_text.tag_config('error_header', foreground='red', font=("Courier", 9, "bold"))
            output_text.tag_config('output', foreground='black')
            output_text.tag_config('error', foreground='red')
            output_text.tag_config('summary_header', foreground='blue', font=("Courier", 11, "bold"))
            output_text.tag_config('summary_text', foreground='black', font=("Courier", 9))

            output_text.config(state='disabled')

            # Close button
            close_button = tk.Button(
                main_frame,
                text="Close",
                font=("", 11),
                command=output_window.destroy,
                padx=20,
                pady=5
            )
            close_button.pack(pady=(10, 0))

            print(f"[INFO] Output window created successfully with {len(command_outputs)} commands")

        except Exception as e:
            print(f"[ERROR] Failed to create output window: {e}")
            import traceback
            traceback.print_exc()

